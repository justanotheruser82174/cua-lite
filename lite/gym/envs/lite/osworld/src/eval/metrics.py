"""Format-aware eval metrics for lite_osworld synth tasks.

OSWorld upstream's `compare_docx_files` only compares `paragraph.text` (plain
strings), missing all formatting and non-paragraph content (tables, images,
footers). This causes ~50 task families in writer.py to trivially pass even
when the agent does nothing — verified empirically: tasks that hit
`max_steps=15`, never saved, never even attempted the task can still get
reward=1.0.

This module provides `compare_docx_strict` that checks:
  - Paragraph text + per-character formatting (bold/italic/underline/strike,
    font name/size/color, highlight)
  - Paragraph-level format (alignment, line_spacing, style.name, indent)
  - Tables (count + dimensions + cell text)
  - Inline images / shapes (count + byte equality)
  - Section margins / page size
  - Header + footer text and digit presence (covers page-numbers tasks)

Dispatched by lite/gym/envs/lite/osworld/src/eval/runner.py before
falling back to OSWorld upstream metrics.

Usage in evaluators:
  {"func": "compare_docx_strict", "result": {...}, "expected": {...}}

literal_match override
----------------------
Upstream `literal_match` does `str(result) == str(expected)`. When the
evaluator specifies `expected = {"type": "rule", "rules": {"expected": "..."}}`,
the lite_osworld getter returns the inner `rules` dict, so the upstream
comparison stringifies the dict and never matches. The override here unwraps
`{"expected": ...}` before delegating to upstream.

GIMP overrides
--------------
check_brightness_decrease_and_structure_sim and check_contrast_increase_and_structure_sim
are kept ONLY as thin mirrors of upstream at UPSTREAM's thresholds (MSE 0.03 / SSIM 0.65).
A loosened threshold (0.03→0.15 / 0.65→0.40) is not version-forced: a container probe
proved lite GIMP == osworld VM GIMP == 2.10.30-1ubuntu0.1 (identical), and on the real
task assets the default non-legacy op clears the upstream thresholds anyway
(brightness-decrease MSE ≈0.00001; max contrast SSIM 0.757). A loosened band is also
reward-hackable (admits posterize/solarize/wash-then-darken) — see each function's
docstring + devs/envs/lite.osworld/bridge/AGENTS.md (REIMPL guard-audit). Tasks whose
UPSTREAM evaluator explicitly sets a threshold (e.g. multi_apps_4c26e3f3 → 0.15) still pass that
value through: that is upstream parity, not a lite override.
"""

from __future__ import annotations

import logging
from typing import Any

logger = logging.getLogger(__name__)


def _normal_default_font(doc) -> str | None:
    """The document's Normal-style default font name (or None)."""
    try:
        return doc.styles["Normal"].font.name
    except (KeyError, AttributeError):
        return None


def _effective_font_name(run, paragraph, normal_font: str | None) -> str | None:
    """Resolve a run's EFFECTIVE font name through the inheritance cascade:
    run.font.name → run char-style → paragraph style → Normal default.

    LibreOffice's docx save normaliser hoists a document-wide font to the
    paragraph-style / Normal default, leaving per-run ``font.name=None`` even
    though the rendered text IS that font. Reading raw ``run.font.name`` then
    false-fails a correct edit. Shared by ``compare_font_names_loose`` and
    ``_char_format_signature`` so both honour the cascade identically.
    """
    try:
        if run.font.name:
            return run.font.name
        rs = getattr(run, "style", None)
        if rs is not None and rs.font.name:
            return rs.font.name
        ps = getattr(paragraph, "style", None)
        if ps is not None and ps.font.name:
            return ps.font.name
        return normal_font
    except AttributeError:
        return None


def _normal_default_size(doc) -> float | None:
    """The document's Normal-style default font size in pt (or None)."""
    try:
        sz = doc.styles["Normal"].font.size
        return sz.pt if sz is not None else None
    except (KeyError, AttributeError):
        return None


def _normal_default_color(doc) -> str | None:
    """The document's Normal-style default font color as RGB hex (or None)."""
    try:
        c = doc.styles["Normal"].font.color
        return str(c.rgb) if c is not None and c.rgb is not None else None
    except (KeyError, AttributeError):
        return None


def _effective_font_size(run, paragraph, normal_size: float | None) -> float | None:
    """Resolve a run's EFFECTIVE font size (pt) through the inheritance cascade:
    run.font.size → run char-style → paragraph style → Normal default.

    Parallels ``_effective_font_name``: LO's docx save hoists a document-wide size
    to the paragraph/Normal style, leaving per-run ``font.size=None`` even though
    the rendered size is unchanged. Reading raw ``run.font.size`` then false-fails a
    correct size edit whose value equals the style default. Both docs resolve
    identically, so a genuinely-different explicit size still mismatches.
    """
    try:
        if run.font.size is not None:
            return run.font.size.pt
        rs = getattr(run, "style", None)
        if rs is not None and rs.font.size is not None:
            return rs.font.size.pt
        ps = getattr(paragraph, "style", None)
        if ps is not None and ps.font.size is not None:
            return ps.font.size.pt
        return normal_size
    except AttributeError:
        return None


def _effective_color(run, paragraph, normal_color: str | None) -> str | None:
    """Resolve a run's EFFECTIVE font color (RGB hex) through the inheritance cascade:
    run.font.color → run char-style → paragraph style → Normal default.

    Parallels ``_effective_font_name``: LO hoists a document-wide color to the
    paragraph/Normal style on save, leaving per-run ``font.color`` unset even though
    the rendered color is unchanged. Theme/scheme colors expose ``rgb is None`` and
    fall through to the normal default (both sides do so identically). A genuinely-
    different explicit RGB still mismatches.
    """
    try:
        c = run.font.color
        if c is not None and c.rgb is not None:
            return str(c.rgb)
        rs = getattr(run, "style", None)
        if rs is not None and rs.font.color is not None and rs.font.color.rgb is not None:
            return str(rs.font.color.rgb)
        ps = getattr(paragraph, "style", None)
        if ps is not None and ps.font.color is not None and ps.font.color.rgb is not None:
            return str(ps.font.color.rgb)
        return normal_color
    except AttributeError:
        return None


_W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_W_FILL = f"{{{_W_NS}}}fill"
_W_VAL = f"{{{_W_NS}}}val"
_YELLOW_HIGHLIGHT_VALUE = 7
_YELLOW_SHADING_FILLS = {"FFFF00", "FFFFFF00", "YELLOW"}


def _run_highlight_value(run) -> int | None:
    """Return the effective yellow text-background value for a run.

    LibreOffice may save a yellow text highlight made through the UI as
    run-level Word shading (`w:shd w:fill="FFFF00"`) instead of
    `w:highlight w:val="yellow"`. They render the same for these tasks, so
    normalize yellow run shading to the python-docx yellow highlight enum.
    """
    if run.font.highlight_color is not None:
        return run.font.highlight_color.value

    rpr = getattr(run._element, "rPr", None)
    if rpr is None:
        return None
    shd = rpr.find(f"{{{_W_NS}}}shd")
    if shd is None:
        return None
    fill = str(shd.get(_W_FILL) or "").strip().upper()
    val = str(shd.get(_W_VAL) or "").strip().lower()
    if fill in _YELLOW_SHADING_FILLS and val in {"", "clear", "solid"}:
        return _YELLOW_HIGHLIGHT_VALUE
    return None


def _char_format_signature(paragraph, *,
                           examine_font_name: bool = True,
                           examine_font_size: bool = True,
                           examine_color: bool = True,
                           examine_highlight: bool = True,
                           normal_font: str | None = None,
                           normal_size: float | None = None,
                           normal_color: str | None = None) -> list[tuple]:
    """Per-character format signature.

    For each char of `paragraph.text`, emit a tuple of formatting attrs.  Two
    paragraphs are format-equivalent iff their signatures are equal even if
    they were tokenized into runs differently (LibreOffice often splits or
    merges runs after edits).

    The `examine_*` flags let callers relax specific checks. LibreOffice's
    docx round-trip can normalize font name/size/color on untouched runs,
    causing false-fails for writer perturbs even when the agent did the
    right thing.
    """
    sig: list[tuple] = []
    for run in paragraph.runs:
        text = run.text or ""
        # python-docx tri-state: True / False / None.  Treat None as False
        # (the default — inherit from style).
        attrs = (
            bool(run.bold),
            bool(run.italic),
            bool(run.underline),
            bool(run.font.strike),
            # Resolve the EFFECTIVE font via the inheritance cascade (not raw
            # run.font.name) so LO's hoist-to-style normalisation doesn't false-
            # fail a correct Ctrl+A→font edit. Both docs resolve identically, so
            # a genuinely-wrong font still mismatches. (validation writer-font bug)
            (_effective_font_name(run, paragraph, normal_font) if examine_font_name else None),
            # Resolve EFFECTIVE size via the SAME style cascade as font name (LO hoists a
            # doc-wide size to the style on save → raw run.font.size=None false-fails a
            # correct edit whose size equals the style default; a wrong size still mismatches).
            (_effective_font_size(run, paragraph, normal_size) if examine_font_size else None),
            # Highlight: WD_COLOR_INDEX enum value or equivalent LO run shading.
            (_run_highlight_value(run) if examine_highlight else None),
            # Color: EFFECTIVE RGB hex via the style cascade (parallels font name/size); LO
            # hoists a doc-wide color to the style on save → raw run.font.color=None false-
            # fails a correct color edit equal to the style default. A wrong color still mismatches.
            (_effective_color(run, paragraph, normal_color) if examine_color else None),
        )
        sig.extend([attrs] * len(text))
    # writer_0e763496: LibreOffice leaves a TRAILING-whitespace run
    # unformatted (e.g. an unstruck trailing "  ") while the gold struck the whole
    # paragraph, so the trailing-space chars' per-char format tuple differs -> FN.
    # Neutralize the format tuples of trailing (visually invisible) whitespace chars
    # ONLY. Text is still compared separately by the caller, and an INTERIOR unstruck
    # word still mismatches -> proven no-FP.
    full = paragraph.text
    n_trailing = len(full) - len(full.rstrip())
    if n_trailing and sig:
        neutral = (None,) * len(sig[0])
        for k in range(len(sig) - n_trailing, len(sig)):
            sig[k] = neutral
    return sig


def _paragraph_break_count(paragraph) -> tuple[int, int]:
    """Count (page_break, line_break) elements in this paragraph.

    Page breaks are `<w:br w:type="page"/>` inside runs and don't contribute
    to `paragraph.text` or per-char signatures (often inside an empty-text
    run).  We count them here so a paragraph that gained a page break is
    distinguishable from one that didn't.
    """
    page = 0
    line = 0
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    for run in paragraph.runs:
        for br in run._element.findall(f".//{ns}br"):
            t = br.get(f"{ns}type")
            if t == "page":
                page += 1
            else:
                # default (textWrapping) and "column" both treated as line break
                line += 1
    return (page, line)


def _paragraph_format_signature(paragraph) -> tuple:
    """Paragraph-level format (style, alignment, line_spacing, indent).

    Excludes `space_before` / `space_after`: LibreOffice's docx save normalizes
    these to a concrete value inherited from the paragraph's style (e.g., Heading 1
    has 24pt space_before by default), while python-docx-written paragraphs leave
    them as None (inherit-at-render-time). Both are semantically equivalent. If we
    compare them strict-equal we get false fails on every successfully-completed
    task that goes through LibreOffice's save (i.e., all of them, since the env's
    postconfig forces a ctrl+s before evaluation).

    Tasks that target space_before/after specifically are rare in our pool; we
    accept that those few become degenerate as a trade-off for not false-failing
    the common format tasks (bold, font, alignment, line_spacing, etc.).
    """
    pf = paragraph.paragraph_format

    def _effective(attr, hard_default):
        """Resolve a paragraph_format attr through the style cascade (own value →
        paragraph style → base styles), falling back to the format's HARD default.
        LO's save drops a REDUNDANT explicit value (one equal to the inherited
        default) to None, while the non-round-tripped python-docx gold keeps it
        explicit — comparing raw then false-fails a font-only edit (gold=LEFT/1.0,
        LO-result=None). Resolving the EFFECTIVE value (None → style → hard default)
        on both sides matches when the attr is untouched, yet still distinguishes a
        genuine NON-default alignment/line-spacing edit LO keeps explicit.
        (validation: 936321ce/e528b65e — font-only edits false-failed on incidental
        alignment/line_spacing that LO normalised to inherit.)
        """
        own = getattr(pf, attr)
        if own is not None:
            return own
        style = paragraph.style
        while style is not None:
            sval = getattr(style.paragraph_format, attr, None)
            if sval is not None:
                return sval
            style = getattr(style, "base_style", None)
        return hard_default

    # alignment hard default = LEFT (unaligned text renders left); line_spacing = 1.0 (single).
    eff_align = _effective("alignment", None)

    def _eff_indent_pt(attr):
        # left/right/first_line indent through the SAME style cascade as alignment/
        # line_spacing. LO drops a redundant explicit indent (equal to the style
        # default) to None on save while the python-docx gold keeps it explicit →
        # raw compare false-fails a styled paragraph (Heading/List/Quote) whose indent
        # is inherited. hard default None = no indent; a genuinely-different explicit
        # indent LO keeps still mismatches → no false positive.
        v = _effective(attr, None)
        return v.pt if v is not None else None

    return (
        paragraph.style.name if paragraph.style is not None else None,
        eff_align.value if eff_align is not None else 0,
        _effective("line_spacing", 1.0),
        _eff_indent_pt("left_indent"),
        _eff_indent_pt("right_indent"),
        _eff_indent_pt("first_line_indent"),
    )


def _table_signature(table) -> tuple:
    """Cell text grid — captures row/col count and content."""
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cells.append(cell.text)
        rows.append(tuple(cells))
    return tuple(rows)


def _section_signature(section) -> tuple:
    """Page geometry + header/footer plain text."""
    def _emu_or_none(x):
        return x.pt if x is not None else None

    header_text = ""
    footer_text = ""
    if section.header is not None and section.header.paragraphs:
        header_text = "\n".join(p.text for p in section.header.paragraphs)
    if section.footer is not None and section.footer.paragraphs:
        footer_text = "\n".join(p.text for p in section.footer.paragraphs)
    # Boolean flag: does footer/header contain a digit (proxy for page-number field).
    # python-docx doesn't render fields, so we also check the underlying XML for
    # a `<w:fldChar>` referencing PAGE.
    return (
        _emu_or_none(section.page_height),
        _emu_or_none(section.page_width),
        _emu_or_none(section.left_margin),
        _emu_or_none(section.right_margin),
        _emu_or_none(section.top_margin),
        _emu_or_none(section.bottom_margin),
        header_text,
        footer_text,
        _section_has_page_field(section),
    )


def _section_has_page_field(section) -> bool:
    """Detect Word PAGE field in header or footer (page numbers).

    python-docx exposes paragraph.text but not field codes, so look at the
    raw XML for `<w:instrText>PAGE` or similar.
    """
    try:
        for region in (section.header, section.footer):
            if region is None:
                continue
            xml = region._element.xml
            if "PAGE" in xml and "instrText" in xml:
                return True
    except Exception:
        pass
    return False


def _images_signature(doc) -> list[bytes]:
    """Sorted hashes of inline image bytes.

    Retained for callers that want strict byte equality. `compare_docx_strict`
    now uses `_images_match_tolerant` instead.
    """
    import hashlib
    blobs = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blobs.append(hashlib.sha256(rel.target_part.blob).hexdigest())
            except Exception:
                pass
    return sorted(blobs)


def _extract_image_blobs(doc) -> list[bytes]:
    """Return raw image blobs from a docx in relationship-iteration order."""
    blobs: list[bytes] = []
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                blobs.append(rel.target_part.blob)
            except Exception:
                pass
    return blobs


def _decode_blob_to_rgb_array(blob: bytes):
    """Decode an image blob (JPEG/PNG/…) into an (H, W, 3) uint8 numpy array.

    Returns None if PIL can't open the blob.
    """
    import io
    from PIL import Image
    import numpy as np
    try:
        img = Image.open(io.BytesIO(blob))
        if img.mode != "RGB":
            img = img.convert("RGB")
        return np.asarray(img, dtype=np.uint8)
    except Exception as e:
        logger.warning("_decode_blob_to_rgb_array: PIL failed: %s", e)
        return None


def _pair_similarity(arr_a, arr_b) -> float:
    """SSIM-or-PSNR similarity between two RGB uint8 arrays.

    Resizes the larger image down to the smaller one's shape if they differ
    (LO can decode-then-re-encode at a slightly different size on some
    pipelines; the gold and result both went through LO once so they normally
    match, but we handle the rare drift defensively).

    Returns a normalized score in [0, 1]:
      - SSIM directly if skimage is available.
      - Else PSNR mapped via min(1.0, psnr / 40.0). PSNR≥30 dB → score ≥0.75.
    """
    import numpy as np
    from PIL import Image

    if arr_a is None or arr_b is None:
        return 0.0
    if arr_a.shape != arr_b.shape:
        # Resize larger → smaller via PIL.
        h = min(arr_a.shape[0], arr_b.shape[0])
        w = min(arr_a.shape[1], arr_b.shape[1])
        try:
            arr_a = np.asarray(
                Image.fromarray(arr_a).resize((w, h), Image.BILINEAR), dtype=np.uint8
            )
            arr_b = np.asarray(
                Image.fromarray(arr_b).resize((w, h), Image.BILINEAR), dtype=np.uint8
            )
        except Exception as e:
            logger.warning("_pair_similarity: resize failed: %s", e)
            return 0.0

    # Try SSIM first (preferred — perceptual + bounded in [0, 1]).
    try:
        from skimage.metrics import structural_similarity as ssim_fn
        min_dim = min(arr_a.shape[0], arr_a.shape[1])
        win = 7 if min_dim >= 7 else (min_dim if min_dim % 2 == 1 else min_dim - 1)
        if win < 3:
            # Image too small for SSIM; fall through to PSNR.
            raise RuntimeError("image too small for SSIM")
        try:
            score = ssim_fn(arr_a, arr_b, win_size=win, channel_axis=2)
        except TypeError:
            score = ssim_fn(arr_a, arr_b, win_size=win, multichannel=True)
        return float(score)
    except Exception:
        pass

    # Fallback: PSNR mapped to [0, 1]
    diff = arr_a.astype(np.float64) - arr_b.astype(np.float64)
    mse = float(np.mean(diff * diff))
    if mse == 0:
        return 1.0
    psnr = 10.0 * np.log10((255.0 * 255.0) / mse)
    return min(1.0, max(0.0, psnr / 40.0))


def _images_match_tolerant(doc1, doc2, threshold: float = 0.85) -> bool:
    """Pixel-tolerant docx image comparison.

    Replaces strict byte-hash equality of embedded image blobs (which fails on
    JPEG re-encode through LO's Ctrl+S). For each pair of
    images at the same iteration index:
      1. Decode both blobs via PIL into RGB uint8 arrays.
      2. Resize to a common shape if necessary.
      3. Compute SSIM (or PSNR fallback). Reject if < threshold.

    Hash-equal blobs short-circuit to True (cheap path). Images must match in
    count.
    """
    import hashlib

    blobs1 = _extract_image_blobs(doc1)
    blobs2 = _extract_image_blobs(doc2)
    if len(blobs1) != len(blobs2):
        logger.debug("image count mismatch: %d vs %d", len(blobs1), len(blobs2))
        return False

    # Order-agnostic: try to pair via byte-hash multiset first (catches the
    # zero-edit case cheaply), then fall back to greedy pairing by similarity.
    h1 = sorted(hashlib.sha256(b).hexdigest() for b in blobs1)
    h2 = sorted(hashlib.sha256(b).hexdigest() for b in blobs2)
    if h1 == h2:
        return True

    # Greedy max-similarity matching (n is tiny in practice — ≤4 images per
    # docx in the photo_to_docx family). For each image in doc1, find the
    # best-similarity unmatched image in doc2; require every pair ≥ threshold.
    arrs1 = [_decode_blob_to_rgb_array(b) for b in blobs1]
    arrs2 = [_decode_blob_to_rgb_array(b) for b in blobs2]
    used = [False] * len(arrs2)
    for a in arrs1:
        best_j = -1
        best_s = -1.0
        for j, b in enumerate(arrs2):
            if used[j]:
                continue
            s = _pair_similarity(a, b)
            if s > best_s:
                best_s = s
                best_j = j
        if best_j == -1 or best_s < threshold:
            logger.debug(
                "image-tolerant: best pair similarity %.3f < threshold %.2f",
                best_s, threshold,
            )
            return False
        used[best_j] = True
    return True


def compare_docx_strict(file1, file2, **options) -> int:
    """Strict comparison of two .docx files.

    Returns 1 if file1 and file2 match on:
      - Paragraph count + each paragraph's text
      - Per-character formatting (bold/italic/underline/strike/font/highlight/color)
      - Paragraph-level format (style/alignment/spacing/indent)
      - Table count + each table's cell text
      - Image count + image byte-hashes (multiset equality)
      - Section count + page geometry + header/footer text + PAGE field presence
    Returns 0 on any mismatch or on parse error.

    Options (default True except where noted):
      examine_font_name, examine_font_size, examine_color, examine_highlight:
        char-format relaxations. Disable individual fields to skip strict
        equality on them. LibreOffice's docx round-trip can normalize these
        on untouched runs even when the agent's edit was correct.
      examine_images: pixel-tolerant SSIM equality of embedded images.
        Short-circuits to True on byte-hash multiset equality (cheap zero-
        edit case). Otherwise decodes blobs via PIL, greedy-pairs by SSIM,
        and requires every pair ≥ image_ssim_threshold. This replaces the
        old byte-hash semantics that false-failed on LO's JPEG re-encode
        during validation.
      image_ssim_threshold (default 0.85): min per-pair SSIM under
        examine_images. Lower for noisier images; raise to make stricter.
      allow_paragraph_diff (default False): if True, paragraphs may differ
        in text content too; only structure-level checks (count, etc.) apply.
    """
    if not file1 or not file2:
        return 0
    if not (file1.endswith(".docx") and file2.endswith(".docx")):
        logger.warning("compare_docx_strict only supports .docx; got %s, %s", file1, file2)
        return 0

    examine_font_name = options.get("examine_font_name", True)
    examine_font_size = options.get("examine_font_size", True)
    examine_color = options.get("examine_color", True)
    examine_highlight = options.get("examine_highlight", True)
    examine_images = options.get("examine_images", True)

    try:
        from docx import Document
        d1 = Document(file1)
        d2 = Document(file2)
    except Exception as e:
        logger.error("compare_docx_strict: failed to open: %s", e)
        return 0

    # Paragraphs: count + text + char-format + para-format
    if len(d1.paragraphs) != len(d2.paragraphs):
        logger.debug("paragraph count mismatch: %d vs %d", len(d1.paragraphs), len(d2.paragraphs))
        return 0
    # Per-doc Normal-style default font, for the effective-font cascade in
    # _char_format_signature (so LO's hoist-to-style font doesn't false-fail).
    _n1 = _normal_default_font(d1)
    _n2 = _normal_default_font(d2)
    _sz1 = _normal_default_size(d1)
    _sz2 = _normal_default_size(d2)
    _nc1 = _normal_default_color(d1)
    _nc2 = _normal_default_color(d2)
    for i, (p1, p2) in enumerate(zip(d1.paragraphs, d2.paragraphs)):
        if p1.text != p2.text:
            logger.debug("paragraph %d text mismatch", i)
            return 0
        sig1 = _char_format_signature(
            p1, examine_font_name=examine_font_name, examine_font_size=examine_font_size,
            examine_color=examine_color, examine_highlight=examine_highlight,
            normal_font=_n1, normal_size=_sz1, normal_color=_nc1,
        )
        sig2 = _char_format_signature(
            p2, examine_font_name=examine_font_name, examine_font_size=examine_font_size,
            examine_color=examine_color, examine_highlight=examine_highlight,
            normal_font=_n2, normal_size=_sz2, normal_color=_nc2,
        )
        if sig1 != sig2:
            logger.debug("paragraph %d char-format mismatch", i)
            return 0
        if _paragraph_format_signature(p1) != _paragraph_format_signature(p2):
            logger.debug("paragraph %d para-format mismatch", i)
            return 0
        if _paragraph_break_count(p1) != _paragraph_break_count(p2):
            logger.debug("paragraph %d break-count mismatch", i)
            return 0

    # Tables
    if len(d1.tables) != len(d2.tables):
        logger.debug("table count mismatch: %d vs %d", len(d1.tables), len(d2.tables))
        return 0
    for i, (t1, t2) in enumerate(zip(d1.tables, d2.tables)):
        if _table_signature(t1) != _table_signature(t2):
            logger.debug("table %d content mismatch", i)
            return 0

    # Images: pixel-tolerant (SSIM-based). LO's Ctrl+S re-encodes JPEG blobs,
    # so raw byte-hash equality false-fails on visually-identical content
    # (29/37 photo_to_docx rollouts in validation scored 0). The
    # tolerant path falls back to byte-hash multiset equality first (cheap)
    # and only invokes PIL/SSIM when hashes differ. Threshold 0.85 still
    # rejects unrelated/wrong images (SSIM≈0) while accepting JPEG re-encode
    # drift (typically SSIM≥0.98 for the same source image).
    if examine_images:
        image_threshold = float(options.get("image_ssim_threshold", 0.85))
        if not _images_match_tolerant(d1, d2, threshold=image_threshold):
            logger.debug("image content mismatch (tolerant SSIM)")
            return 0

    # Sections (page geometry + header/footer)
    if len(d1.sections) != len(d2.sections):
        logger.debug("section count mismatch: %d vs %d", len(d1.sections), len(d2.sections))
        return 0
    for i, (s1, s2) in enumerate(zip(d1.sections, d2.sections)):
        if _section_signature(s1) != _section_signature(s2):
            logger.debug("section %d mismatch", i)
            return 0

    return 1


def check_gitignore_has_entries(result_path: str, **options) -> float:
    """Check that a .gitignore file at *result_path* contains every entry in
    `options["must_have"]` as a standalone non-comment line.

    Rationale (validation fix for vs_code_create_gitignore_*):
    The prior eval was `compare_text_file` byte-equal vs a 200+-line canonical
    GitHub gitignore template — impossible to recall verbatim from instruction.
    The skill we actually want to test is "produce a sensible .gitignore for
    project type X", so we relax the eval to a small set of MUST-HAVE entries
    (e.g. ``node_modules/`` + ``*.log`` for Node, ``__pycache__/`` + ``*.pyc``
    for Python). Comments and blank lines are stripped; entries are matched
    exactly against the remaining lines (after strip()) so a trailing-slash
    convention (``dir/`` vs ``dir``) is honoured.

    Options:
      must_have (list[str], required): pattern strings that must each appear
        as a standalone line in the file.
    """
    if not result_path:
        return 0.0
    must_have: list[str] = options.get("must_have", []) or []
    if not must_have:
        return 0.0
    try:
        with open(result_path, "r", encoding="utf-8", errors="replace") as fh:
            raw = fh.read()
    except OSError as e:
        logger.warning("check_gitignore_has_entries: open failed: %s", e)
        return 0.0
    # Keep only non-comment, non-blank lines.
    lines = {
        ln.strip()
        for ln in raw.splitlines()
        if ln.strip() and not ln.lstrip().startswith("#")
    }
    missing = [m for m in must_have if m not in lines]
    if missing:
        logger.debug(
            "check_gitignore_has_entries: missing entries %s in %s",
            missing, result_path,
        )
        return 0.0
    return 1.0


def check_thunderbird_prefs_loose(result, rule) -> float:
    """Lite-side wrapper around upstream `check_thunderbird_prefs` that
    accepts BOTH `mail.server.default.X` and `mail.server.serverN.X` keys
    interchangeably when matching the expected rules.

    Rationale (Validation, F_TB_28/29/30 set_pref cluster):
    Thunderbird's Account Settings GUI writes per-account `serverN.X` prefs;
    `mail.server.default.X` is only set via about:config (Edit → Config Editor).
    Eval rules generated by `_pref_evaluator` use `default.X` (which matches
    the `_pref_oracle` heredoc that writes there programmatically), but
    agents performing the natural UI path land in `serverN.X` → strict
    exact-key match returns 0. This wrapper accepts either namespace, so
    UI-edit AND about:config-edit both score 1.0 when the value matches.

    Only used by train-side synth tasks (eval.jsonl tasks still use upstream
    strict `check_thunderbird_prefs`).
    """
    from desktop_env.evaluators.metrics.thunderbird import _pref_pattern, _match_pref
    if result is None:
        return 0.0
    expect_rules = rule.get("expect", {})
    unexpect_rules = rule.get("unexpect", {})

    # Build accept-set per key: for keys starting with mail.server.default.,
    # also accept mail.server.<anything>.<suffix>.
    def _aliases_for(key: str) -> set[str]:
        aliases = {key}
        prefix = "mail.server.default."
        if key.startswith(prefix):
            suffix = key[len(prefix):]
            # accept mail.server.serverN.<suffix>, e.g. server1, server2, server3
            for n in range(1, 21):
                aliases.add(f"mail.server.server{n}.{suffix}")
        return aliases

    # Reverse map: for each prefs.js key encountered, find which expected key (if any) it satisfies.
    expect_alias_map = {}  # alias -> original_expected_key
    for k in expect_rules:
        for alias in _aliases_for(k):
            expect_alias_map[alias] = k
    unexpect_alias_map = {}
    for k in unexpect_rules:
        for alias in _aliases_for(k):
            unexpect_alias_map[alias] = k

    expect_metrics = {k: False for k in expect_rules}
    unexpect_metric = True
    try:
        with open(result) as f:
            for line in f:
                m = _pref_pattern.match(line.strip())
                if m is None: continue
                key = m.group("key")
                try:
                    import json as _json
                    value = _json.loads(m.group("val"))
                except Exception:
                    continue
                if key in expect_alias_map:
                    orig_key = expect_alias_map[key]
                    if _match_pref(value, expect_rules[orig_key]):
                        expect_metrics[orig_key] = True
                elif key in unexpect_alias_map:
                    orig_key = unexpect_alias_map[key]
                    if _match_pref(value, unexpect_rules[orig_key]):
                        unexpect_metric = False
    except Exception as e:
        logger.warning("check_thunderbird_prefs_loose: read failed: %s", e)
        return 0.0
    return 1.0 if (all(expect_metrics.values()) and unexpect_metric) else 0.0


def compare_font_names_loose(docx_file, rules) -> float:
    """Lite-side wrapper around upstream `compare_font_names` that resolves
    paragraph-style and document-default font inheritance before comparing.

    Upstream `compare_font_names` (`desktop_env/evaluators/metrics/docs.py:555`)
    reads `run.font.name` raw — which is `None` for any run whose `rPr/rFonts`
    is inherited rather than explicit. LO Writer's save normaliser preferentially
    moves font.name to the paragraph style / document-default style, leaving
    individual runs with `font.name=None`. Result: an agent that correctly
    Ctrl+A → Format → Character → Font Name → "Times New Roman" reaches a
    document whose RENDERED font is TNR everywhere, but whose RAW run.font.name
    is `None` on most runs → upstream returns 0.

    This helper resolves the inheritance cascade:
        run.font.name
        → run.style.font.name        (character style)
        → paragraph.style.font.name  (paragraph style)
        → docx.styles['Normal'].font.name  (document default)
    and compares the EFFECTIVE font name against `rules["font_name"]`. Only
    used by train-side synth/perturb tasks (synth/libreoffice_writer.py's
    `_build_font_names_evaluator`); eval.jsonl tasks continue to use upstream
    strict `compare_font_names` so OSWorld benchmark semantics are preserved.

    Skips empty/whitespace-only runs (they have no visible glyphs to honour
    any font name, so their font.name is irrelevant for the user-visible
    correctness of the edit).
    """
    if not docx_file:
        return 0.0
    try:
        from docx import Document  # python-docx
        doc = Document(docx_file)
    except Exception as e:
        logger.warning("compare_font_names_loose: open failed: %s", e)
        return 0.0

    expected_font = rules["font_name"]

    def _normal_default() -> str | None:
        try:
            return doc.styles["Normal"].font.name
        except (KeyError, AttributeError):
            return None

    normal_font = _normal_default()

    def _effective_font(run, paragraph) -> str | None:
        try:
            if run.font.name:
                return run.font.name
            run_style = getattr(run, "style", None)
            if run_style is not None and run_style.font.name:
                return run_style.font.name
            para_style = getattr(paragraph, "style", None)
            if para_style is not None and para_style.font.name:
                return para_style.font.name
            return normal_font
        except AttributeError:
            return None

    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            # Skip whitespace-only / empty runs — they don't render glyphs so
            # their font is irrelevant to the visible result.
            if not (run.text or "").strip():
                continue
            if _effective_font(run, paragraph) != expected_font:
                return 0.0
    return 1.0


def compare_references(file1: str, file2: str, **options) -> float:
    """Override of upstream compare_references that returns 0 when neither
    document has the References heading.

    Upstream returns 1.0 when both documents lack the reference indicator
    (vacuous equality). For tasks where the agent must ADD a References
    section, this causes a trivial pass before the agent does anything.
    We return 0.0 instead so the task only scores 1.0 when the result
    document actually contains the References section.
    """
    from desktop_env.evaluators.metrics.docs import compare_references as _upstream
    from docx import Document as _Document

    if not file1 or not file2:
        return 0.0

    reference_indicator = options.get("reference_indicator", "References")
    try:
        doc_result = _Document(file1)
        paragraphs_result = [p.text for p in doc_result.paragraphs]
    except Exception:
        return 0.0

    ref_result_idx = (
        paragraphs_result.index(reference_indicator)
        if reference_indicator in paragraphs_result
        else -1
    )

    # If result has no References heading, the task was not completed.
    if ref_result_idx == -1:
        return 0.0

    # Result has a References section — delegate the full comparison upstream.
    return _upstream(file1, file2, **options)


# ---------------------------------------------------------------------------
# GIMP overrides: looser thresholds for GIMP's non-legacy B-C algorithm
# ---------------------------------------------------------------------------

def check_brightness_decrease_and_structure_sim(src_path: str, tgt_path: str, threshold: float = 0.03) -> float:
    """Thin mirror of upstream check_brightness_decrease_and_structure_sim at upstream's
    default threshold (0.03). Pure one-directional wrapper: names the threshold, then delegates.

    Kept at upstream's 0.03, not loosened: lite GIMP == osworld VM GIMP == 2.10.30
    (identical), so there is no version-forced reason to diverge from upstream.
    Empirically, on the real asset a legit non-legacy brightness DECREASE scores
    MSE ≈ 0.00001 (upstream normalizes both images to brightness 128, so a pure
    brightness change cancels out) — far under 0.03, so upstream's 0.03 does NOT
    false-fail the task. A looser band such as [0.03,0.15) is reward-hackable —
    posterize/solarize/wash-then-darken score 0.033-0.116 and would pass 0.15
    while destroying structure.
    NOTE: tasks whose UPSTREAM evaluator explicitly sets options.threshold (e.g.
    multi_apps_4c26e3f3 → 0.15) still get that value passed through — that is upstream parity,
    NOT a lite override, so it is left intact. See bridge/AGENTS.md (REIMPL guard-audit).
    """
    from desktop_env.evaluators.metrics.gimp import (
        check_brightness_decrease_and_structure_sim as _upstream,
    )
    return _upstream(src_path, tgt_path, threshold=threshold)


def check_contrast_increase_and_structure_sim(src_path: str, tgt_path: str) -> float:
    """Thin mirror of upstream check_contrast_increase_and_structure_sim at upstream's
    threshold (0.65). Pure one-directional wrapper: names the threshold, then delegates.

    Kept at upstream's 0.65, not loosened: a container probe showed lite GIMP ==
    osworld VM GIMP == 2.10.30-1ubuntu0.1 (identical), so upstream's 0.65 behaves
    the same on both substrates — there is no version-forced reason to diverge.
    On the actual eval asset (f723c744/berries.png) the DEFAULT non-legacy op
    scores SSIM 0.757 even at MAX contrast, so 0.65 does NOT false-fail a legit
    edit. A looser band such as [0.40,0.65) admits wrong-but-std↑ ops
    (posterize/brightness-with-clip/hist-eq: SSIM 0.45-0.62) — reward-hackable
    since this metric feeds train. See devs/envs/lite.osworld/bridge/AGENTS.md
    (REIMPL guard-audit).
    """
    from desktop_env.evaluators.metrics.gimp import (
        check_contrast_increase_and_structure_sim as _upstream,
    )
    return _upstream(src_path, tgt_path)


# ---------------------------------------------------------------------------
# Impress eval-tolerance overrides (deferred queue, validation)
# ---------------------------------------------------------------------------
#
# The two helpers below wrap upstream `compare_pptx_files` by copying both the
# result and expected .pptx files to scratch paths, snapping the offending
# property to a quantization grid in BOTH copies, then delegating the
# (now-tolerant) comparison to upstream. Because the same snap is applied
# symmetrically, equal-modulo-tolerance values become byte-identical and
# upstream's strict equality returns 1 on real successes.
#
# Implementation choice: we mutate copies via the python-pptx model and let
# python-pptx round-trip the OOXML. The original files are NOT touched, so
# repeated eval calls (and oracle replays) remain idempotent.


def _quantize_pptx_colors(src_path: str, dst_path: str, rgb_tolerance: int) -> bool:
    """Copy src_path → dst_path with every explicit RGB run color snapped to
    a `2*rgb_tolerance+1`-wide grid. Returns True on success.

    The snap maps any byte v ∈ [0..255] to `round(v / step) * step` where
    `step = 2*rgb_tolerance + 1`. Two original bytes that differ by ≤
    rgb_tolerance MAY still snap to neighboring bins (boundary case), but in
    practice gold colors are picked at clean 5/10/20-step values so
    rgb_tolerance≈2 is plenty of slack vs the Custom-Color picker drift of
    1-3 bytes per channel.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    step = max(1, 2 * rgb_tolerance + 1)

    def _snap_byte(b: int) -> int:
        return min(255, max(0, round(b / step) * step))

    def _snap_rgb(rgb: RGBColor) -> RGBColor:
        r, g, b = rgb[0], rgb[1], rgb[2]
        return RGBColor(_snap_byte(r), _snap_byte(g), _snap_byte(b))

    def _walk(shape):
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    color = run.font.color
                    try:
                        if color.rgb is not None:
                            color.rgb = _snap_rgb(color.rgb)
                    except (AttributeError, TypeError):
                        # Theme/scheme colors (color.type == MSO_THEME_COLOR)
                        # don't expose `.rgb` as RGBColor — leave them (upstream
                        # compares enum identity, stable across LO round-trip).
                        # A color-LESS run (color.type is None) must NOT compare
                        # EQUAL to a colored run via upstream's hasattr-guarded
                        # skip, so stamp a reserved sentinel RGB no gold uses.
                        try:
                            if color.type is None:
                                color.rgb = RGBColor(0x01, 0x02, 0x03)  # "unset" sentinel
                        except Exception:
                            pass
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            _c = run.font.color
                            try:
                                if _c.rgb is not None:
                                    _c.rgb = _snap_rgb(_c.rgb)
                            except (AttributeError, TypeError):
                                # See run-branch note: stamp the "unset" sentinel
                                # on a color-less cell run so it can't vacuously
                                # equal a colored gold; theme colors untouched.
                                try:
                                    if _c.type is None:
                                        _c.rgb = RGBColor(0x01, 0x02, 0x03)
                                except Exception:
                                    pass
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                _walk(sub)

    try:
        prs = Presentation(src_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                _walk(shape)
        prs.save(dst_path)
        return True
    except Exception as e:
        logger.warning("_quantize_pptx_colors failed for %s: %s", src_path, e)
        return False


def _quantize_pptx_positions(src_path: str, dst_path: str, cm_tolerance: float) -> bool:
    """Copy src_path → dst_path with every shape (left, top, width, height)
    snapped to a `step` EMU grid where step = round(cm_tolerance * 360000 * 2)+1.

    Snapping symmetrically in both files turns sub-mm round-trip drift on
    the gold (which traveled through LO save once) and the result (which
    traveled through LO save once) into byte-identical EMU values that pass
    upstream's `is_approximately_equal` (0.5%-relative-tolerance) trivially.
    """
    from pptx import Presentation
    # 1 cm = 360000 EMU. Use a step ≈ 2*tolerance so any two values within
    # tolerance of each other map to the same bin (modulo boundary effects).
    step = max(1, int(round(cm_tolerance * 360000)))

    def _snap(v):
        if v is None:
            return v
        # ROUND to nearest bin, not floor. pptx positions are clean cm
        # multiples (360000 EMU) and step (18000) divides them, so gold sits ON a
        # floor bin-edge; a negative-drift result then floors to the bin BELOW and
        # fails upstream's 0.5% compare (and floor could split a would-pass pair
        # 359900/360100 across two bins). round() matches _quantize_pptx_colors and
        # keeps a real ≥1cm move (≥20 bins apart) failing → no false positive.
        return round(v / step) * step

    def _walk(shape):
        # python-pptx exposes .left/.top/.width/.height as int (EMU). Group
        # shapes don't allow direct assignment of left/top, so guard with
        # try/except.
        try:
            if shape.left is not None:
                shape.left = _snap(shape.left)
            if shape.top is not None:
                shape.top = _snap(shape.top)
            if shape.width is not None:
                shape.width = _snap(shape.width)
            if shape.height is not None:
                shape.height = _snap(shape.height)
        except (AttributeError, TypeError, ValueError):
            pass
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                _walk(sub)

    try:
        prs = Presentation(src_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                _walk(shape)
        prs.save(dst_path)
        return True
    except Exception as e:
        logger.warning("_quantize_pptx_positions failed for %s: %s", src_path, e)
        return False


def compare_pptx_files_color_tolerant(result, expected, **options) -> int:
    """RGB-tolerant variant of upstream compare_pptx_files.

    Handles the title_color cluster (>=15 impress variants): the Custom-Color
    picker rounds to the display gamut, so eval-side ARGB differs by 1-3 bytes
    per channel from the gold's exact ARGB.

    Approach: copy both pptx files to scratch paths with run RGB colors
    snapped to a (2*rgb_tolerance+1)-byte grid in BOTH copies, then delegate
    to upstream `compare_pptx_files`. Theme/scheme colors are passed through
    unchanged (upstream compares them by enum identity, which is round-trip
    stable). All other comparisons (text, fonts, positions, …) are
    unaffected because the snap touches only the SrgbClr values.
    """
    import os
    import tempfile

    from desktop_env.evaluators.metrics.slides import compare_pptx_files as _upstream

    rgb_tolerance = int(options.pop("rgb_tolerance", 2))

    if not result or not expected:
        return 0
    if not (str(result).endswith(".pptx") and str(expected).endswith(".pptx")):
        return _upstream(result, expected, **options)

    tmpdir = tempfile.mkdtemp(prefix="pptx_color_tolerant_")
    snap_result = os.path.join(tmpdir, "result.pptx")
    snap_expected = os.path.join(tmpdir, "expected.pptx")

    ok1 = _quantize_pptx_colors(result, snap_result, rgb_tolerance)
    ok2 = _quantize_pptx_colors(expected, snap_expected, rgb_tolerance)
    if not (ok1 and ok2):
        # Snap failed — fall back to upstream on the originals so we never
        # silently report 1 on a broken eval.
        return _upstream(result, expected, **options)

    try:
        return _upstream(snap_result, snap_expected, **options)
    finally:
        try:
            import shutil
            shutil.rmtree(tmpdir, ignore_errors=True)
        except Exception:
            pass


def compare_pptx_files_position_tolerant(result, expected, **options) -> int:
    """Cm-tolerant variant of upstream compare_pptx_files.

    Wires the position-compare cluster (d_imp_59 + d_imp_60 x3 = 5 members
    per cluster escalation growth #1, deferred queue, trigger
    O — pptx round-trip via LO emits sub-mm drift on (left, top, width,
    height) values that fail strict EMU equality).

    Approach: copy both pptx files to scratch paths with every shape's
    (left, top, width, height) snapped to a `cm_tolerance` Cm grid in BOTH
    copies, then delegate to upstream. The default `cm_tolerance=0.05`
    (≈0.05 Cm = 18000 EMU) is well below the smallest semantic move in our
    gold set (≥1 Cm) and above the empirical round-trip drift (≤0.01 Cm).
    """
    import os
    import tempfile

    from desktop_env.evaluators.metrics.slides import compare_pptx_files as _upstream

    cm_tolerance = float(options.pop("cm_tolerance", 0.05))

    if not result or not expected:
        return 0
    if not (str(result).endswith(".pptx") and str(expected).endswith(".pptx")):
        return _upstream(result, expected, **options)

    tmpdir = tempfile.mkdtemp(prefix="pptx_pos_tolerant_")
    snap_result = os.path.join(tmpdir, "result.pptx")
    snap_expected = os.path.join(tmpdir, "expected.pptx")

    ok1 = _quantize_pptx_positions(result, snap_result, cm_tolerance)
    ok2 = _quantize_pptx_positions(expected, snap_expected, cm_tolerance)
    if not (ok1 and ok2):
        return _upstream(result, expected, **options)

    try:
        return _upstream(snap_result, snap_expected, **options)
    finally:
        try:
            import shutil
            shutil.rmtree(tmpdir, ignore_errors=True)
        except Exception:
            pass


# Impress geometry reconcile. One reconcile helper, driven by the
# catalog's position_mode. POSITION_DRIFT_EMU is > every observed LibreOffice
# round-trip + manual drift (<=27360 EMU) and < the smallest semantic move (171720
# EMU). A single flat tolerance is provably impossible (lane_C6 counter-proof): #21's
# correct top-drag needs tol>=~160000 but #23's centered no-op sits at 171540 -> any
# flat tol that passes #21 false-positives #23. Hence region + per-axis tolerance.
POSITION_DRIFT_EMU = 50000


def _reconcile_pptx_geometry(res_path, gold_path, dst, *, tol_emu, region, ignore_geom):
    """Copy res_path->dst, snapping each result shape's (left, top, width, height) ONTO
    the matching gold shape wherever the difference is semantically acceptable, then let
    upstream compare the snapped copy. Snapping toward gold (never away) means a
    genuinely-wrong shape stays different -> upstream returns 0.

      - ignore_geom : set of (slide_idx, shape_type) whose geometry is unspecified
                      (e.g. an inserted table) -> snap all four attrs unconditionally.
      - region      : {"slide": i, "shape_type": t, "edge": "top|bottom|left|right"} for a
                      "move to edge" task. shape_type may be the sentinel "title" (resolved
                      by placeholder idx-0, mirroring slide.shapes.title) or a numeric MSO
                      type (13 picture, 19 table). Position snaps when BOTH centers land in
                      the named third (a confirmed move) OR -- for a near-full-span shape
                      whose center cannot leave the middle third (#23) -- when the residual
                      is within tol (round-trip drift). Size snaps only within tol (a move
                      must not resize).
      - tol_emu     : every OTHER (unmoved) shape gets each axis snapped when the residual
                      is within tol, absorbing global round-trip drift so upstream's strict
                      examine_shape gate does not false-fail the untouched shapes."""
    from pptx import Presentation

    rp, gp = Presentation(res_path), Presentation(gold_path)
    W, H = rp.slide_width, rp.slide_height

    def in_edge(sh, edge):
        # Bias-past-tolerance, not center-in-third. A near-full-width /
        # near-full-height shape's center can never leave the middle third, so the old
        # `band()==edge` predicate never fired for it (it rode on the tol fallback only,
        # which a realistic drag overshoots). Instead require the center displaced from
        # SLIDE CENTER toward the target edge by >= tol_emu. This is a strict superset of
        # the third test for the compact shapes (#21 title / #22 picture): center in the
        # bottom third => cy > 2/3 H > H/2 + tol, so those still fire; and it additionally
        # fires for a full-width table nudged right past tol (#23). Requiring it for BOTH
        # result and gold keeps every neg control 0 (centered => displacement 0 < tol).
        cx, cy = sh.left + sh.width / 2, sh.top + sh.height / 2
        return {"top": (H / 2 - cy) >= tol_emu, "bottom": (cy - H / 2) >= tol_emu,
                "left": (W / 2 - cx) >= tol_emu, "right": (cx - W / 2) >= tol_emu}[edge]

    def is_title(sh):
        # python-pptx `slide.shapes.title` is the placeholder whose idx == 0.
        return sh.is_placeholder and sh.placeholder_format.idx == 0

    def is_slide_number_placeholder(sh):
        try:
            text = (sh.text or "").strip()
            return sh.is_placeholder and (text == "<number>" or text.isdigit())
        except Exception:
            return False

    def is_region_target(shr, shg):
        st = region["shape_type"]
        if st == "title":
            return is_title(shr) and is_title(shg)
        return shr.shape_type == st

    for si, (sr, sg) in enumerate(zip(rp.slides, gp.slides)):
        g = list(sg.shapes)
        for idx, shr in enumerate(sr.shapes):
            if idx >= len(g):
                break
            shg = g[idx]

            def cp(attrs, _shr=shr, _shg=shg):
                for a in attrs:
                    try:
                        setattr(_shr, a, getattr(_shg, a))
                    except Exception:
                        pass

            def within(a, _shr=shr, _shg=shg):
                rv, gv = getattr(_shr, a), getattr(_shg, a)
                return rv is not None and gv is not None and abs(rv - gv) <= tol_emu

            if (si, shr.shape_type) in ignore_geom:
                cp(("left", "top", "width", "height"))
                continue
            # LO headless normalization can collapse page-number placeholders
            # to top=0/height=0 on one copy but not the other. Position-mode
            # tasks never grade page-number geometry, so absorb this drift.
            if is_slide_number_placeholder(shr) and is_slide_number_placeholder(shg):
                cp(("left", "top", "width", "height"))
                continue
            if region and region["slide"] == si and is_region_target(shr, shg):
                if in_edge(shr, region["edge"]) and in_edge(shg, region["edge"]):
                    cp(("left", "top"))            # confirmed move -> position matches gold
                else:
                    for a in ("left", "top"):      # tight move / unmoved axis: absorb drift
                        if within(a):
                            cp((a,))
                for a in ("width", "height"):
                    if within(a):
                        cp((a,))
                continue
            for a in ("left", "top", "width", "height"):
                if within(a):
                    cp((a,))
    rp.save(dst)
    return True


def compare_pptx_files(result, expected, **options) -> int:
    """Lite override: perturb impress color tasks call the bare
    `compare_pptx_files` with `examine_color_rgb:true`, hitting upstream's
    exact-byte color compare — which false-fails on the GUI Custom-Color picker's
    1-3 byte gamut rounding. The synth generator routes ITS color tasks to
    `compare_pptx_files_color_tolerant` (99 rows), but perturb (24 color rows) +
    eval use the bare name. Route only color-examining tasks to the tolerant
    variant; everything else — the 243 synth bare NON-color tasks, eval — is
    upstream verbatim (so blast radius is exactly the color tasks).
    """
    if options.get("examine_color_rgb"):
        return compare_pptx_files_color_tolerant(result, expected, **options)
    # validation (impress reposition cluster 2b94c692/ac1b39ff/15aece23 etc.):
    # examine_shape compares (left,top,width,height) at upstream's 0.5% tolerance.
    # The pptx round-trip through LO emits sub-mm drift on these EMU values, which
    # can fail strict-ish equality even when the move is semantically correct, so
    # route shape-examining tasks through the cm-grid-snapping tolerant variant.
    # NOTE: this absorbs ROUND-TRIP DRIFT only (default 0.05cm grid). It does NOT
    # make vague directional drags ("move to the right") or unspecified-position
    # inserts (insert_table) pass — those golds use a fixed EMU target the manual
    # placement won't hit within 0.05cm. Fixing those needs a directional/existence
    # eval (compare "moved toward edge?" / "table of N×M exists?"), not a tolerance
    # bump — tracked separately; a larger tolerance would overlap the ≥1cm smallest
    # semantic move and cause false positives.
    if options.get("position_mode"):  # {"region"|"ignore"|"tolerance"} carried by the catalog
        import os
        import shutil
        import tempfile

        from desktop_env.evaluators.metrics.slides import compare_pptx_files as _upstream

        td = tempfile.mkdtemp(prefix="pptx_geom_")
        try:
            snap = os.path.join(td, "res.pptx")
            _reconcile_pptx_geometry(
                result, expected, snap,
                tol_emu=int(options.pop("position_tolerance_emu", POSITION_DRIFT_EMU)),
                region=options.pop("position_region", None),
                ignore_geom={tuple(x) for x in options.pop("ignore_shape_geometry", [])},
            )
            options.pop("position_mode", None)
            # Strict geometry gate on the SNAPPED copy. Upstream defaults examine_shape True,
            # but the generator's _BASE_OPTS carries it False, so set it True explicitly here
            # -- otherwise a wrong-region / wrong-size result would skip upstream's dimension
            # gate and false-pass (breaking every region/size neg control).
            options["examine_shape"] = True
            return _upstream(snap, expected, **options)
        finally:
            shutil.rmtree(td, ignore_errors=True)
    if options.get("examine_shape"):
        return compare_pptx_files_position_tolerant(result, expected, **options)
    from desktop_env.evaluators.metrics.slides import compare_pptx_files as _upstream
    return _upstream(result, expected, **options)


def compare_pptx_appended_titles(result, expected, num_appended: int = 0, **options) -> float:
    """Tier-A1 Chrome->pptx *append-slides* comparator (multi_apps
    perturb family 4c26e3f3 / 47f7c0ce / 778efd0a / bb83cab4).

    The task APPENDS N slides whose only meaningful content is the title TEXT.
    Upstream ``compare_pptx_files`` cannot score it: (1) its per-slide shape-count
    check (``slides.py`` shape-count gate) is UNCONDITIONAL (runs even with
    examine_shape=False) — the gold builder materializes the last layout's
    placeholders (~8 shapes) while an LO-Impress agent appends a title-only
    (~1-shape) slide -> mismatch -> 0; and (2) the source template carries a slide
    python-pptx cannot parse (missing <p:spTree>), which makes upstream RAISE.

    We require equal slide count and, for the APPENDED slides only (the last
    ``num_appended``), equal per-slide non-empty TEXT gathered across all shapes
    (the title lives in a textbox, not a ``shapes.title`` placeholder). The PREFIX
    slides are NOT compared: they are unchanged by the task and LO-normalized on
    save (one is unreadable), so any prefix compare would false-fail a correct
    agent. A genuinely wrong / mis-ordered / short append still fails.
    """
    from pptx import Presentation

    def _slide_texts(slide) -> list[str]:
        texts: list[str] = []

        def walk(shape):
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
            if hasattr(shape, "shapes"):
                for sub in shape.shapes:
                    walk(sub)

        for shape in slide.shapes:
            walk(shape)
        return sorted(texts)

    res_slides = list(Presentation(result).slides)
    exp_slides = list(Presentation(expected).slides)
    if len(res_slides) != len(exp_slides):
        return 0.0
    if num_appended <= 0:
        return 1.0
    for res_slide, exp_slide in zip(res_slides[-num_appended:], exp_slides[-num_appended:]):
        if _slide_texts(res_slide) != _slide_texts(exp_slide):
            return 0.0
    return 1.0


# ---------------------------------------------------------------------------
# Calc chart-eval override (deferred queue, validation)
# ---------------------------------------------------------------------------


def compare_calc_chart_type(result: str, expected: str, **options) -> float:
    """Series-ref-agnostic variant of upstream `compare_table` chart rule.

    Wires the calc chart cluster (f_calc_55__chart_box_office_bar +
    f_calc_77__chart_ticket_price, deferred queue). Root cause:
    upstream `load_charts` keys `chart_set` by the series-reference STRING
    (e.g. `'BoxOffice!$C$2:$C$17,BoxOffice!$A$2:$A$17'`). When the agent
    inserts a chart with the same type/dimensions but LO normalizes the
    range reference (e.g. drops a sheet prefix, or expands the range), the
    dict-equality `charts1 == charts2` fails even though the chart_props
    (type, direction, etc.) match.

    Approach: re-implement the `compare_table` driver loop, but for
    `type=="chart"` rules compare a SORTED LIST OF info dicts (dropping the
    series-ref key) instead of upstream's dict-by-series. All other rule
    types delegate to a clone of upstream's logic by re-using
    `compare_table` with the chart rules stripped out.
    """
    import os

    import openpyxl
    from desktop_env.evaluators.metrics.table import compare_table as _upstream_compare_table
    from desktop_env.evaluators.metrics.utils import load_charts

    if result is None or expected is None:
        return 0.0
    if not os.path.exists(result) or not os.path.exists(expected):
        logger.error("compare_calc_chart_type: missing file (result=%s, expected=%s)", result, expected)
        return 0.0

    rules = options.get("rules", [])
    chart_rules = [r for r in rules if r.get("type") == "chart"]
    non_chart_rules = [r for r in rules if r.get("type") != "chart"]

    # Delegate non-chart rules to upstream verbatim (preserves all of
    # upstream's sheet_data / style / sparkline / etc. semantics).
    if non_chart_rules:
        non_chart_opts = {**options, "rules": non_chart_rules}
        upstream_score = _upstream_compare_table(result, expected, **non_chart_opts)
        # `compare_table` returns 1.0 on full pass, 0.0 otherwise.
        if not upstream_score:
            return 0.0

    # Chart-rule handling: load both books with openpyxl, then compare each
    # chart rule via load_charts(sorted-list-of-info-dicts) rather than
    # upstream's dict-keyed-by-series equality.
    try:
        wb_r = openpyxl.load_workbook(filename=result)
        wb_e = openpyxl.load_workbook(filename=expected)
    except Exception as e:
        logger.error("compare_calc_chart_type: failed to load workbooks: %s", e)
        return 0.0

    def _resolve_sheet(book, sheet_idx):
        # Mirror _parse_sheet_idx's int → sheet name path. `EI<int>` /
        # `RI<int>` / `EN<name>` / `RN<name>` strings are also accepted.
        if isinstance(sheet_idx, int):
            names = book.sheetnames
            return book, names[sheet_idx] if 0 <= sheet_idx < len(names) else ""
        if not isinstance(sheet_idx, str):
            return book, ""
        if sheet_idx.startswith("RI"):
            names = wb_r.sheetnames
            try:
                return wb_r, names[int(sheet_idx[2:])]
            except (ValueError, IndexError):
                return wb_r, ""
        if sheet_idx.startswith("RN"):
            return wb_r, sheet_idx[2:]
        if sheet_idx.startswith("EI"):
            names = wb_e.sheetnames
            try:
                return wb_e, names[int(sheet_idx[2:])]
            except (ValueError, IndexError):
                return wb_e, ""
        if sheet_idx.startswith("EN"):
            return wb_e, sheet_idx[2:]
        return book, ""

    for r in chart_rules:
        book_a, sheet_a = _resolve_sheet(wb_r, r["sheet_idx0"])
        book_b, sheet_b = _resolve_sheet(wb_e, r["sheet_idx1"])
        charts_a = load_charts(book_a, sheet_a, **r)
        charts_b = load_charts(book_b, sheet_b, **r)
        if r.get("require_value_range"):
            # calc chart perturb family: opt-in range-awareness.
            # Dropping the WHOLE series key accepts a chart pointing at the WRONG
            # value range (a real FP — proven: default path scores 1.0 for a col-B
            # chart when gold charts col C). Re-key each chart by a sheet-prefix-
            # normalized VALUE ref (drop only the trailing category ref) so LO's
            # sheet-prefix drop is tolerated (0->1) while a wrong value-range still
            # fails. Synth f_calc_55/f_calc_77 omit the flag -> unchanged.
            def _valkey(series_key):
                return ";".join(
                    seg.split(",", 1)[0].split("!", 1)[-1]
                    for seg in series_key.split(";")
                )
            list_a = sorted(
                ((_valkey(k), v) for k, v in charts_a.items()),
                key=lambda kv: repr((kv[0], sorted(kv[1].items()))),
            )
            list_b = sorted(
                ((_valkey(k), v) for k, v in charts_b.items()),
                key=lambda kv: repr((kv[0], sorted(kv[1].items()))),
            )
        else:
            # Default (unchanged): drop the series-ref key by taking only
            # `.values()` then sorting by a stable serialization of each info dict.
            list_a = sorted(charts_a.values(), key=lambda d: repr(sorted(d.items())))
            list_b = sorted(charts_b.values(), key=lambda d: repr(sorted(d.items())))
        if list_a != list_b:
            logger.debug(
                "compare_calc_chart_type: chart info mismatch: %s vs %s",
                list_a, list_b,
            )
            return 0.0

    return 1.0


def _sheet_data_dtype_tolerant(result: str, expected: str, rule: dict) -> bool:
    """Lite override: upstream `sheet_data` does `s1.round(4).equals(s2)`, which is
    dtype-STRICT — a CORRECT sheet whose columns differ only in dtype (int-vs-
    float, or a NaN forcing object dtype) after the LibreOffice round-trip vs the
    openpyxl-built gold false-fails. Re-check the same two sheets with pandas
    `assert_frame_equal(check_dtype=False)`: same value tolerance as
    upstream (so a WRONG value still fails, and string/datetime cells still
    compare exactly) but ignoring dtype. Used ONLY as an OR-fallback AFTER the
    exact upstream check already rejected, so it can only ADD acceptances of
    value-equal sheets — never relax value-strictness. Reuses upstream's own
    sheet loaders so the two paths can't diverge. Returns False on ANY error
    (conservative — never a spurious pass).
    """
    import pandas as pd
    from desktop_env.evaluators.metrics.table import _load_sheet, _parse_sheet_idx
    try:
        pdr = pd.ExcelFile(result)
        pde = pd.ExcelFile(expected)
        rnames, enames = pdr.sheet_names, pde.sheet_names
        b0, i0 = _parse_sheet_idx(rule["sheet_idx0"], pdr, pde, rnames, enames)
        b1, i1 = _parse_sheet_idx(rule["sheet_idx1"], pdr, pde, rnames, enames)
        s1, s2 = _load_sheet(b0, i0), _load_sheet(b1, i1)
        if not isinstance(s1, pd.DataFrame) or not isinstance(s2, pd.DataFrame):
            return False
        prec = int(rule.get("precision", 4))
        s1 = s1.round(prec).reset_index(drop=True)
        s2 = s2.round(prec).reset_index(drop=True)
        pd.testing.assert_frame_equal(
            s1, s2, check_dtype=False, check_exact=False,
            atol=10 ** (-prec), rtol=1e-6, check_like=False,  # rel-tol so a 2dp-displayed total matches a full-precision stored gold (add-acceptance-only; OR-fallback)
        )
        return True
    except Exception:
        return False


def _sheet_data_time_tolerant(result: str, expected: str, rule: dict) -> bool:
    """calc_15 filter_by_origin: LibreOffice stores pasted time cells as
    datetime.time while the synth gold stores time STRINGS ("06:30") -> DataFrame.equals
    is False -> FN (not rescued by _sheet_data_dtype_tolerant). Canonicalize
    datetime.time/datetime and time-looking strings (^\\d{1,2}:\\d{2}(:\\d{2})?$) to
    HH:MM:SS on BOTH sides, then value-compare. OR-fallback ONLY (after upstream strict
    already rejected) so it can only ADD acceptances of value-equal sheets -- a wrong
    time/value still fails -> no FP. Do NOT broaden to numeric coercion (the rejected C2
    dtype-coercion). Returns False on ANY error (conservative)."""
    import datetime as _dt
    import re as _re
    import pandas as pd
    from desktop_env.evaluators.metrics.table import _load_sheet, _parse_sheet_idx
    _TIME_RE = _re.compile(r"^\d{1,2}:\d{2}(:\d{2})?$")
    try:
        pdr = pd.ExcelFile(result)
        pde = pd.ExcelFile(expected)
        rnames, enames = pdr.sheet_names, pde.sheet_names
        b0, i0 = _parse_sheet_idx(rule["sheet_idx0"], pdr, pde, rnames, enames)
        b1, i1 = _parse_sheet_idx(rule["sheet_idx1"], pdr, pde, rnames, enames)
        s1, s2 = _load_sheet(b0, i0), _load_sheet(b1, i1)
        if not isinstance(s1, pd.DataFrame) or not isinstance(s2, pd.DataFrame):
            return False
        if s1.shape != s2.shape:
            return False

        def _canon(v):
            if isinstance(v, (_dt.time, _dt.datetime)):
                return v.strftime("%H:%M:%S")
            if isinstance(v, str) and _TIME_RE.match(v.strip()):
                parts = [int(p) for p in v.strip().split(":")]
                while len(parts) < 3:
                    parts.append(0)
                return "%02d:%02d:%02d" % tuple(parts[:3])
            return v

        c1 = s1.reset_index(drop=True).map(_canon)
        c2 = s2.reset_index(drop=True).map(_canon)
        return bool(c1.equals(c2))
    except Exception:
        return False


def compare_table_numfmt_tolerant(result: str, expected: str = None, **options) -> float:
    """Locale/quote-tolerant variant of upstream `compare_table` for the
    `number_format` style rule.

    Root cause: the gold is openpyxl-written
    (``cell.number_format = '$#,##0.00'``), but when the agent applies the SAME
    visible format via LibreOffice it is stored as a locale-tagged / re-quoted
    string — e.g. the toolbar Currency button yields `[$$-409]#,##0.00`, and LO
    re-quotes `"$"#,##0`. Upstream does exact `styles1 == styles2`, so a
    correctly-formatted result false-fails.

    Fix: split out `props==["number_format"]` style rules, delegate every other
    rule to upstream `compare_table` verbatim, and for the number_format rules
    compare `load_xlsx_styles` dicts AFTER canonicalizing each format string —
    stripping only the `[$<sym>-<locale>]` locale id (keeping the currency
    symbol) and `"`-quoting. Digits / decimals / `%` are untouched, so a wrong
    format (e.g. `0.0` vs `0.00`) still fails, and an unformatted result cell
    (`General`) still mismatches the gold → no trivial pass.
    """
    import os
    import re

    import openpyxl
    from desktop_env.evaluators.metrics.table import compare_table as _upstream_compare_table
    from desktop_env.evaluators.metrics.utils import load_xlsx_styles

    if result is None:
        return 0.0
    if not os.path.exists(result):
        logger.error("compare_table_numfmt_tolerant: missing result file %s", result)
        return 0.0

    rules = options.get("rules", [])
    nf_rules = [r for r in rules
                if r.get("type") == "style" and (r.get("props") or []) == ["number_format"]]
    other_rules = [r for r in rules if r not in nf_rules]
    sd_rules = [r for r in other_rules if r.get("type") == "sheet_data"]
    non_sd_rules = [r for r in other_rules if r.get("type") != "sheet_data"]

    # `expected` is the gold xlsx; it's None for check_cell-only rules (the ref
    # is baked into the rule, no gold file). Those flow entirely through the
    # non_sd_rules → upstream path below (upstream `compare_table` accepts
    # expected=None). nf/sd rules DO need a gold file — if any are present
    # without an expected, that's a malformed task, fail explicitly. Mirrors
    # upstream's `expected: str = None` signature. (validation: calc_357ef137
    # crashed with a TypeError on its expected:null check_cell rule.)
    if expected is None:
        if nf_rules or sd_rules:
            logger.error("compare_table_numfmt_tolerant: nf/sd rules require a gold file but expected is None")
            return 0.0
    elif not os.path.exists(expected):
        logger.error("compare_table_numfmt_tolerant: missing expected file %s", expected)
        return 0.0

    # Non-sheet_data, non-number_format rules flow through upstream verbatim.
    if non_sd_rules:
        if not _upstream_compare_table(result, expected, **{**options, "rules": non_sd_rules}):
            return 0.0
    # sheet_data rules: run upstream's exact dtype-STRICT check FIRST,
    # so anything it accepts still passes (no regression — the eval oracle's
    # gold-vs-gold compare is byte-identical and keeps passing here). Only when
    # upstream rejects do we retry with a dtype-TOLERANT compare, to rescue a
    # correct sheet whose columns differ only in dtype (int-vs-float, or a NaN
    # forcing object dtype) after the LibreOffice round-trip vs the openpyxl gold.
    # The fallback keeps upstream's value tolerance, so a WRONG value
    # still fails — it can only ADD acceptances of value-equal sheets.
    for r in sd_rules:
        if _upstream_compare_table(result, expected, **{**options, "rules": [r]}):
            continue
        if not _sheet_data_dtype_tolerant(result, expected, r) \
                and not _sheet_data_time_tolerant(result, expected, r):  # time OR-fallback
            return 0.0
    if not nf_rules:
        return 1.0

    try:
        wb_r = openpyxl.load_workbook(filename=result)
        wb_e = openpyxl.load_workbook(filename=expected)
    except Exception as e:
        logger.error("compare_table_numfmt_tolerant: failed to load workbooks: %s", e)
        return 0.0

    _LOCALE_CURRENCY = re.compile(r"\[\$([^\]\-]*)(?:-[0-9A-Fa-f]+)?\]")
    # LO's Currency toolbar button always appends a red-negative subformat
    # (`;[RED]-$#,##0.00`) that the openpyxl-built gold (positive-only) lacks.
    # Strip just that `;[RED]…` clause so the positive format compares equal.
    _RED_NEG = re.compile(r";\[red\][^;]*", re.IGNORECASE)

    def _canon(fmt):
        if fmt is None:
            return None
        f = _LOCALE_CURRENCY.sub(lambda m: m.group(1), str(fmt))  # [$$-409] -> $
        f = _RED_NEG.sub("", f)                                    # drop ;[RED]-… negative clause
        # The openpyxl-written gold escapes the literal currency symbol
        # as `\$#,##0.00`, but LO stores it unescaped (canon -> `$#,##0.00`). A `\`
        # in a number format only escapes the NEXT char to be a literal (same visual),
        # so dropping it normalizes `\$`==`$` without changing the displayed format.
        f = f.replace("\\", "")                                    # \$#,##0 -> $#,##0
        return f.replace('"', "")                                  # "$"#,##0 -> $#,##0

    def _resolve_sheet(book, sheet_idx):
        # Mirror compare_calc_chart_type._resolve_sheet: int -> positional name;
        # RI/RN -> result book; EI/EN -> expected book.
        if isinstance(sheet_idx, int):
            names = book.sheetnames
            return book, names[sheet_idx] if 0 <= sheet_idx < len(names) else ""
        if not isinstance(sheet_idx, str):
            return book, ""
        if sheet_idx.startswith("RI"):
            names = wb_r.sheetnames
            try:
                return wb_r, names[int(sheet_idx[2:])]
            except (ValueError, IndexError):
                return wb_r, ""
        if sheet_idx.startswith("RN"):
            return wb_r, sheet_idx[2:]
        if sheet_idx.startswith("EI"):
            names = wb_e.sheetnames
            try:
                return wb_e, names[int(sheet_idx[2:])]
            except (ValueError, IndexError):
                return wb_e, ""
        if sheet_idx.startswith("EN"):
            return wb_e, sheet_idx[2:]
        return book, ""

    def _fix_formula_numfmt(styles, book, sheet_name):
        # Upstream `_read_cell_style` returns None for a FORMULA cell's
        # number_format because it requires `cell.data_type == "n"` (plain number).
        # A derived column the agent built with `=ROUND(...)` formulas (data_type
        # "f") therefore reads None and false-fails the number_format rule even
        # when correctly formatted, while the openpyxl gold's plain-value cells
        # read fine — the dominant calc-derived false-negative. Re-read the REAL
        # `cell.number_format` for such formula cells. A wrong format still fails
        # (we read the actual format), so this only rescues correct outputs.
        try:
            ws = book[sheet_name]
        except KeyError:
            return styles
        for coord, vals in styles.items():
            if vals and vals[0] is None:
                c = ws[coord]
                if c.value is not None and c.data_type == "f":
                    vals[0] = c.number_format
        return styles

    for r in nf_rules:
        book_a, sheet_a = _resolve_sheet(wb_r, r["sheet_idx0"])
        book_b, sheet_b = _resolve_sheet(wb_e, r["sheet_idx1"])
        styles_a = _fix_formula_numfmt(load_xlsx_styles(book_a, sheet_a, result, **r), book_a, sheet_a)
        styles_b = _fix_formula_numfmt(load_xlsx_styles(book_b, sheet_b, expected, **r), book_b, sheet_b)
        canon_a = {k: [_canon(v) for v in vals] for k, vals in styles_a.items()}
        canon_b = {k: [_canon(v) for v in vals] for k, vals in styles_b.items()}
        # LibreOffice often saves a stored dimension one row/col past the
        # data (phantom empty cells, value=None → numfmt None), so the agent's sheet
        # gets extra `A18..E18 = [None]` keys the openpyxl gold lacks → dict
        # inequality on otherwise-identical formats. Drop keys whose every style
        # value is None (an empty, unformatted cell) from BOTH sides before
        # comparing — symmetric, so a wrong format on a real DATA cell still fails.
        canon_a = {k: v for k, v in canon_a.items() if any(x is not None for x in v)}
        canon_b = {k: v for k, v in canon_b.items() if any(x is not None for x in v)}
        if canon_a != canon_b:
            logger.debug("compare_table_numfmt_tolerant: numfmt mismatch: %s vs %s", canon_a, canon_b)
            return 0.0

    return 1.0


# Global lite override of `compare_table`: the runner resolves metric
# funcs against this module FIRST (runner.py:150 getattr(custom_metrics, ...)), so
# aliasing the upstream name here makes the number_format locale-tolerance apply
# UNIFORMLY across synth + perturb + eval — wherever a task's evaluator func is the
# bare "compare_table" (perturb uses it directly; synth chart/freeze route elsewhere).
# Safe: for any rule set WITHOUT a number_format style rule, the tolerant impl
# delegates 100% to upstream `compare_table` verbatim, so behavior is identical.
compare_table = compare_table_numfmt_tolerant


# --- Lite patches over upstream-OSWorld evaluator false negatives -------------
# These INTENTIONALLY relax verbatim-upstream verifiers that false-fail a correct
# agent. They
# only ADD acceptances of provably-equivalent results; a genuinely-wrong answer
# still fails. The runner resolves funcs against this module first, so aliasing
# the upstream name here applies the patch uniformly across synth/perturb/eval.

# §N-1: upstream `compare_font_names` (docs.py:555) reads `run.font.name` RAW, so a
# correct Ctrl+A→font edit false-fails when LibreOffice's save hoists the font to
# the paragraph/Normal style (per-run name → None). `compare_font_names_loose`
# resolves the same inheritance cascade; a genuinely-wrong font still mismatches.
compare_font_names = compare_font_names_loose


def compare_docx_files(file1, file2, **options):
    """§N-2 lite patch over upstream `compare_docx_files` (docs.py:159): the
    exact-compare branch (`ignore_blanks=False`) does `if p1 != p2: return 0` on
    raw paragraph text, so a result that strips a TRAILING SPACE the gold kept
    (`'… lectures. '` vs `'… lectures.'`) false-fails — trailing whitespace is
    visually/semantically void. Delegate to upstream first; only if it returns 0
    AND both files are .docx, retry with each paragraph's text strip'd (honouring
    ignore_case; strip() also absorbs a sentence-per-line LEADING space). A
    real content difference survives strip → still fails."""
    from desktop_env.evaluators.metrics.docs import compare_docx_files as _upstream
    score = _upstream(file1, file2, **options)
    if score or not file1 or not file2:
        return score
    if not (str(file1).endswith(".docx") and str(file2).endswith(".docx")):
        return score
    # strip-tolerant fallback — only widens the exact-compare path.
    if options.get("content_only") or options.get("fuzzy_match") or options.get("ignore_order"):
        return score  # those branches aren't the trailing-space bug; don't touch
    try:
        from docx import Document
        p1 = [p.text for p in Document(file1).paragraphs]
        p2 = [p.text for p in Document(file2).paragraphs]
    except Exception:
        return score
    if options.get("ignore_blanks", True):
        p1 = [t for t in p1 if t.strip()]
        p2 = [t for t in p2 if t.strip()]
    if len(p1) != len(p2):
        return score
    ic = options.get("ignore_case", False)
    # LibreOffice Writer autocorrect rewrites typed ASCII punctuation on entry:
    # " - " -> en-dash, "--" -> em-dash, straight -> curly quotes, "..." -> "…".
    # The python-docx gold keeps the literal ASCII the source showed, so a
    # verbatim-typed paragraph false-fails on the glyph alone (FN on
    # multi_apps_7ff48d5b "–" vs "-"). Fold this LO-autocorrect class back to
    # ASCII before compare. Same OR-fallback shape as the strip rule below:
    # runs only after upstream exact-match failed and requires equal para count,
    # so it can only ADD acceptances of typography-equal paragraphs — content
    # diffs still fail.
    _AUTOCORRECT = {"–": "-", "—": "-", "‘": "'", "’": "'",
                    "“": '"', "”": '"', "…": "..."}
    def _norm(t):
        # strip() not rstrip(): splitting sentences onto separate lines leaves the
        # space after '.' as LEADING whitespace on the next paragraph, a benign
        # artifact of the transform (sentence-per-line FN). Fallback only runs
        # after upstream fails and requires equal para count + per-para equality
        # modulo strip, so genuine content diffs still fail (negative-control checked).
        t = t.strip()
        for _k, _v in _AUTOCORRECT.items():
            t = t.replace(_k, _v)
        return t.lower() if ic else t
    return 1 if all(_norm(a) == _norm(b) for a, b in zip(p1, p2)) else score


def compare_line_spacing(docx_file1, docx_file2) -> int:
    """Lite override: upstream compares `paragraph_format.line_spacing`
    with `!=`, but LibreOffice exports a DEFAULT single-spaced paragraph with no
    explicit spacing → python-docx reads `None`, while the gold may carry explicit
    `1.0`. Normalize `None`→`1.0` (single IS the default) before comparing, so a
    correct single-spacing edit isn't false-failed. A non-single gold still requires
    the agent to set it (no trivial pass).
    """
    from docx import Document
    from desktop_env.evaluators.metrics.docs import compare_docx_files as _cmp_docx
    if not docx_file1 or not docx_file2:
        return 0
    if not _cmp_docx(docx_file1, docx_file2):
        return 0
    try:
        doc1, doc2 = Document(docx_file1), Document(docx_file2)
    except Exception as e:
        logger.error("compare_line_spacing: %s", e)
        return 0
    if len(doc1.paragraphs) != len(doc2.paragraphs):
        return 0

    def _sp(p):
        s = p.paragraph_format.line_spacing
        return 1.0 if s is None else s

    for p1, p2 in zip(doc1.paragraphs, doc2.paragraphs):
        if _sp(p1) != _sp(p2):
            return 0
    return 1


# ---------------------------------------------------------------------------
# EPUB comparison (validation deferred) — local override
# ---------------------------------------------------------------------------
# Upstream `compare_epub` (desktop_env/.../others.py:72) extracts toc.ncx,
# content.opf, and *.html from each epub, then diffs each pair via
# SequenceMatcher. Its `process_epub` only strips `dc:identifier` from
# content.opf and `navPoint` from toc.ncx — leaving `<dcterms:modified>`,
# pandoc generator meta, and other runtime-varying tags. Even with
# SOURCE_DATE_EPOCH baked into both gold + agent commands, multi_epub_md_*
# tasks still scored below the 0.5 pass threshold.
#
# This override widens the strip set so OPF/NCX noise can't drag the ratio
# below 0.5 for byte-identical-modulo-noise epubs.


def _process_epub_strict(filename: str) -> list[str]:
    """Like upstream process_epub but with broader noise stripping."""
    import logging as _lg
    import os as _os
    import zipfile as _zf

    import lxml.html as _lh

    log = _lg.getLogger("lite_osworld.compare_epub")
    base_dir = filename + ".dir"
    _os.makedirs(base_dir, exist_ok=True)
    file_list: list[str] = []

    _OPF_NOISE = (
        "dc:identifier",      # UUID
        "dcterms:modified",   # pandoc wallclock (sometimes leaks past SOURCE_DATE_EPOCH)
        "dc:date",            # pandoc emits when env var absent
        'meta name="generator"',
        'meta name="cover"',
        '<meta property="dcterms:modified"',
    )
    _NCX_NOISE = ("navPoint", "<meta name=\"dtb:uid\"")

    # Pandoc EPUB3 emits `EPUB/content.opf`, `EPUB/toc.ncx`, and `.xhtml` files
    # (not bare `content.opf` / `toc.ncx` / `.html` that the older EPUB2 spec
    # assumed). Match by basename + extension so we cover both layouts.
    try:
        with _zf.ZipFile(filename, "r") as zf:
            names = zf.namelist()
            ncx = next((n for n in names if _os.path.basename(n).lower() == "toc.ncx"), None)
            opf = next((n for n in names if n.lower().endswith(".opf")), None)
            html_like = [n for n in names if n.lower().endswith((".html", ".xhtml"))]
            if ncx:
                out_path = _os.path.join(base_dir, "toc.ncx")
                with zf.open(ncx) as in_f, open(out_path, "w") as out_f:
                    for line in in_f.read().decode().splitlines():
                        if not any(p in line for p in _NCX_NOISE):
                            out_f.write(line + "\n")
                file_list.append(out_path)
            if opf:
                out_path = _os.path.join(base_dir, "content.opf")
                with zf.open(opf) as in_f, open(out_path, "w") as out_f:
                    for line in in_f.read().decode().splitlines():
                        if not any(p in line for p in _OPF_NOISE):
                            out_f.write(line + "\n")
                file_list.append(out_path)
            for fn in html_like:
                # Flatten nested paths to a unique basename inside base_dir
                # so the listing order is comparable across runs.
                out_path = _os.path.join(base_dir, _os.path.basename(fn))
                with zf.open(fn) as in_f, open(out_path, "w") as out_f:
                    raw = in_f.read().decode()
                    cleaned = "".join(c for c in raw if c not in "\n\r").encode()
                    html = _lh.fromstring(cleaned)
                    out_f.write(_lh.tostring(html, pretty_print=True, encoding="unicode"))
                file_list.append(out_path)
        log.debug("%s: %s", filename, file_list)
        return sorted(file_list)
    except _zf.BadZipFile:
        return []


def _epub_chapter_text(filename: str) -> str:
    """Concatenated, layout-independent CHAPTER text of an epub.

    Pulls every .html/.xhtml entry's visible text (tags stripped, whitespace
    collapsed) in spine order and joins it — so the comparison does not depend
    on how many files the prose is split across or what they're named.
    """
    import os as _os
    import zipfile as _zf

    import lxml.html as _lh

    try:
        with _zf.ZipFile(filename, "r") as zf:
            html_like = sorted(n for n in zf.namelist()
                               if n.lower().endswith((".html", ".xhtml")))
            parts: list[str] = []
            for fn in html_like:
                # skip pure-nav/cover scaffolding that carries no prose
                if _os.path.basename(fn).lower() in ("nav.xhtml", "toc.xhtml"):
                    continue
                # lxml.html.fromstring REJECTS a `str` that carries an XML
                # encoding declaration ("<?xml … encoding=…?>", which pandoc's
                # EPUB3 xhtml always has) — pass the raw BYTES so it parses.
                raw = zf.read(fn)
                try:
                    txt = _lh.fromstring(raw).text_content()
                except Exception:
                    continue
                txt = " ".join(txt.split())
                if txt:
                    parts.append(txt)
            return "\n".join(parts)
    except _zf.BadZipFile:
        return ""


def compare_epub(result: str, expected: str, **options) -> float:
    """Local override of upstream compare_epub.

    Behavioural delta vs upstream: strips more known-noisy OPF/NCX lines
    (`dcterms:modified`, `meta name="generator"`, etc.) before diffing.
    Same return type / threshold semantics — returns mean SequenceMatcher
    ratio across extracted files; runner.py applies the 0.5 pass cutoff.

    §N-3 lite patch: the per-file `zip()` pairs files by SORTED BASENAME, which
    misaligns when the agent's tool emits a different file layout than the gold
    (pandoc EPUB3 `ch001.xhtml`+`nav` vs an EPUB2 `file0.html`+`cover` gold) —
    chapter prose then gets diffed against `content.opf`/`nav` → false 0 even
    when the prose is byte-identical. So if the positional diff falls below the
    0.5 pass cutoff, fall back to a layout-INDEPENDENT compare of the
    concatenated chapter TEXT (tags stripped); take the better ratio. A genuinely
    different book still scores low on both. (validation: multi_apps 42d25c08.)
    """
    import difflib
    if result is None:
        return 0.0
    result_files = _process_epub_strict(result)
    expected_files = _process_epub_strict(expected)
    if not result_files:
        return 0.0
    total = 0.0
    for f1, f2 in zip(result_files, expected_files):
        with open(f1) as fh:
            a = fh.read().splitlines()
        with open(f2) as fh:
            b = fh.read().splitlines()
        total += difflib.SequenceMatcher(a=a, b=b).ratio()
    positional = total / len(result_files)
    if positional >= 0.5:
        return positional
    # layout-independent chapter-text fallback. Compare at WORD granularity (the
    # concatenated prose is whitespace-collapsed and may be a single line, so
    # splitlines() would yield one element and a 0/1 ratio).
    rt = _epub_chapter_text(result)
    et = _epub_chapter_text(expected)
    if not rt or not et:
        return positional
    content = difflib.SequenceMatcher(a=rt.split(), b=et.split()).ratio()
    return max(positional, content)


# ---------------------------------------------------------------------------
# VLC compare_videos duration-aware override
# ---------------------------------------------------------------------------
#
# Upstream `desktop_env/evaluators/metrics/vlc.py:compare_videos` iterates
# at most `max_frames_to_check=100` frames at `fps=25`, so the loop exits
# after ~4 seconds of footage regardless of true container duration. Any
# video ≥4s passes the implicit "same length" check trivially, even when
# the agent's clip is the full source instead of the requested trim. This
# defeats the F_VLC_16 trim_5s/trim_3s tasks: the source video and a no-op copy share pHash
# and both run >4s, so upstream returns 1.0 on both sides.
#
# This helper:
#   1. Reads container duration of both files via ffprobe (no frame loop).
#   2. Rejects when |dur_a − dur_b| ≥ 0.5 s.
#   3. Only then defers to upstream `compare_videos` for the pHash check,
#      with `max_frames_to_check` raised so long clips are sampled across
#      their entire length (cheap on 5-10 s clips).
# Returns 1.0 iff BOTH duration and pHash agree, 0.0 otherwise.


def _video_duration(path: str) -> float | None:
    """Container duration in seconds.

    Probes via OpenCV first (already a hard dep of upstream compare_videos,
    so no PATH/shell concerns). Falls back to ``ffprobe`` only if cv2 can't
    open the file or reports a zero/negative fps. Returns None on failure.
    """
    try:
        import cv2  # pylint: disable=import-outside-toplevel

        cap = cv2.VideoCapture(path)
        try:
            if not cap.isOpened():
                return None
            fps = cap.get(cv2.CAP_PROP_FPS) or 0.0
            frames = cap.get(cv2.CAP_PROP_FRAME_COUNT) or 0.0
            if fps > 0 and frames > 0:
                return float(frames) / float(fps)
        finally:
            cap.release()
    except Exception as e:
        logger.debug("_video_duration cv2 path failed for %s: %s", path, e)

    # ffprobe fallback. Try /usr/bin/ffprobe explicitly before falling back
    # to PATH — some host envs ship a conda ffprobe with broken libiconv
    # linkage that fails before producing duration.
    import os as _os
    import subprocess

    for binary in ("/usr/bin/ffprobe", "ffprobe"):
        if binary != "ffprobe" and not _os.path.exists(binary):
            continue
        try:
            out = subprocess.run(
                [
                    binary, "-v", "error", "-select_streams", "v:0",
                    "-show_entries", "format=duration",
                    "-of", "default=nokey=1:noprint_wrappers=1", path,
                ],
                capture_output=True, text=True, timeout=15,
            )
            raw = (out.stdout or "").strip()
            if raw:
                return float(raw)
        except (subprocess.SubprocessError, OSError, ValueError) as e:
            logger.debug("_video_duration ffprobe (%s) failed: %s", binary, e)
    return None


def compare_videos_full_duration(
    video_path1: str,
    video_path2: str,
    *,
    duration_tolerance_s: float = 0.5,
    max_frames_to_check: int = 600,
    threshold: int = 5,
    **_unused,
) -> float:
    """Duration-aware variant of upstream compare_videos.

    Returns 1.0 only when:
      * container durations agree within `duration_tolerance_s` (default 0.5 s);
      * AND upstream `compare_videos` returns 1.0 with a raised
        `max_frames_to_check` so the pHash loop spans the FULL clip rather
        than just the first ~4 s window that the upstream default samples.

    Returns 0.0 on any duration mismatch, missing file, or pHash divergence.
    """
    if not video_path1 or not video_path2:
        return 0.0
    import os as _os

    if not (_os.path.isfile(video_path1) and _os.path.isfile(video_path2)):
        logger.warning(
            "compare_videos_full_duration: missing file (%s | %s)",
            video_path1, video_path2,
        )
        return 0.0

    dur_a = _video_duration(video_path1)
    dur_b = _video_duration(video_path2)
    if dur_a is None or dur_b is None:
        return 0.0
    if abs(dur_a - dur_b) >= duration_tolerance_s:
        logger.debug(
            "compare_videos_full_duration: duration mismatch %.3fs vs %.3fs"
            " (tolerance %.3fs)", dur_a, dur_b, duration_tolerance_s,
        )
        return 0.0

    # Duration agrees — defer to upstream's pHash check with widened sampling
    # so long clips are inspected end-to-end rather than first-4-seconds-only.
    from desktop_env.evaluators.metrics.vlc import compare_videos as _upstream
    return float(_upstream(
        video_path1, video_path2,
        max_frames_to_check=max_frames_to_check,
        threshold=threshold,
    ))


# ---------------------------------------------------------------------------
# Calc freeze-pane direct probe
# ---------------------------------------------------------------------------
#
# `_LO_NORMALIZE_TAIL` re-saves the openpyxl-written gold xlsx via
# `soffice --headless --convert-to xlsx`. That conversion drops the
# `<sheetView><pane state="frozen"/></sheetView>` element produced by
# openpyxl, so the post-normalize gold has `ws.freeze_panes is None`.
# Upstream's `compare_table` "freeze" rule then compares `None == None`
# and trivially passes BEFORE the agent has done anything.
#
# This helper bypasses the LO normalize path entirely: it opens the
# agent's result xlsx directly with openpyxl (the runner downloads it
# raw from the container before normalize) and checks for the requested
# `freeze_panes` value. To remain robust against runner normalize, the
# generator's `result` getter must read pre-normalize bytes, OR (the
# simpler path) we rely on the agent's interactive
# Ctrl+S which produces an LO-saved xlsx with `<pane state="frozen"/>`
# intact when freeze was actually applied via View → Freeze Cells.
# Empirically, the interactive LO save preserves
# freeze_panes; only the headless --convert-to xlsx drops it.


def check_xlsx_freeze_pane(
    result_path: str,
    *,
    expected_first_row: int = 2,
    expected_first_col_letter: str = "A",
    **_unused,
) -> float:
    """Probe a result xlsx for an `A<row>`-style frozen pane via openpyxl.

    Returns 1.0 iff `ws.freeze_panes` on sheet 0 equals
    `f"{expected_first_col_letter}{expected_first_row}"` (default `"A2"`,
    i.e. row 1 frozen). Returns 0.0 on missing file, missing sheet, parse
    error, or freeze-cell mismatch.

    Note: this MUST be called against the agent's LO-interactive-saved
    xlsx (not the runner's LO-normalize-rewritten copy). For the
    F_CALC_88 freeze_header_row template the agent saves via Ctrl+S, which
    keeps the `<pane state="frozen"/>` sheetView element intact; only the
    headless `soffice --convert-to xlsx` strips it. The runner's docx /
    pptx normalize pass does NOT touch xlsx, so this probe sees the
    interactive save verbatim.
    """
    if not result_path:
        return 0.0
    import os as _os

    if not _os.path.isfile(result_path):
        logger.warning("check_xlsx_freeze_pane: missing %s", result_path)
        return 0.0

    expected = f"{expected_first_col_letter}{expected_first_row}"
    try:
        import openpyxl
        wb = openpyxl.load_workbook(result_path, read_only=False, data_only=False)
        ws = wb.worksheets[0]
        actual = ws.freeze_panes
        # openpyxl's freeze_panes returns only topLeftCell and DISCARDS xSplit,
        # so a row+column freeze (xSplit=1 topLeftCell="B2") reads back identically
        # to a cursor-artifact row-only freeze (xSplit=0 topLeftCell="B2"). Read
        # xSplit off the pane so the row-only relaxation can require NO column
        # freeze. The relaxation must not pass a "freeze first row" result that
        # also froze column A.
        _pane = ws.sheet_view.pane
        actual_xsplit = getattr(_pane, "xSplit", None) if _pane is not None else None
        wb.close()
    except Exception as e:
        logger.warning("check_xlsx_freeze_pane: load failed for %s: %s",
                       result_path, e)
        return 0.0

    if actual == expected:
        return 1.0
    # Row-only freeze (expected col == "A", i.e. no column freeze intended):
    # accept any top-left COLUMN as long as the frozen ROW matches. LO writes
    # the freeze top-left cell using the active cell's column, so "freeze the
    # first row" with the cursor parked in column B yields "B2" — row 1 IS still
    # frozen (ySplit correct); the column is a cursor-position artifact, not a
    # wrong freeze. (validation: d681960f false-failed; sibling froze identically
    # and passed only because its cursor happened to sit in column A.)
    if expected_first_col_letter == "A" and isinstance(actual, str):
        import re as _re
        m = _re.match(r"^[A-Z]+(\d+)$", actual)
        # Row matches AND no column is also frozen (xSplit absent/0). A non-zero
        # xSplit means the agent ALSO froze a column — a genuinely-different freeze
        # state, not a cursor artifact, so it must NOT pass a row-only task.
        if m and int(m.group(1)) == expected_first_row and not actual_xsplit:
            return 1.0
    logger.debug(
        "check_xlsx_freeze_pane: freeze_panes=%r expected=%r", actual, expected,
    )
    return 0.0


def check_page_number_colors(pptx_file: str, rules: dict, **_options) -> float:
    """Slide-number color check, scoped to the `sldNum` placeholder.

    OSWorld upstream `check_page_number_colors`
    (`desktop_env/evaluators/metrics/slides.py`) tries the right XPath first
    (`.//p:ph[@type="sldNum"]//a:solidFill//a:srgbClr`) but on miss falls
    through to a reverse-scan over **all** `.//a:solidFill//a:srgbClr` in
    `slideMaster1.xml` and picks the first non-`000000`. In typical
    LibreOffice-saved masters that deterministically lands on a theme accent
    (often `FFFFFF`), so the check fails no matter what the agent did.

    This override only trusts the placeholder-scoped lookup: find every
    `p:sp` whose `p:nvSpPr/p:nvPr/p:ph[@type="sldNum"]` is present, then
    read the first `a:solidFill//a:srgbClr@val` underneath that shape (so
    either the `a:fld`/`a:rPr` run-color or a body-level color counts).
    Returns 1.0 iff that color matches the named bucket (`red` / `blue` /
    `green` / `black`) using upstream's RGB heuristic.
    """
    import xml.etree.ElementTree as ET
    import zipfile

    if not pptx_file:
        return 0.0
    color_name = (rules or {}).get("color")
    if not color_name:
        logger.warning("check_page_number_colors: no color in rules=%r", rules)
        return 0.0

    ns = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }

    try:
        with zipfile.ZipFile(pptx_file, "r") as zf:
            with zf.open("ppt/slideMasters/slideMaster1.xml") as f:
                root = ET.parse(f).getroot()
    except (zipfile.BadZipFile, KeyError, FileNotFoundError, ET.ParseError) as e:
        logger.warning("check_page_number_colors: cannot open %s: %s", pptx_file, e)
        return 0.0

    found_val: str | None = None
    for sp in root.findall(".//p:sp", ns):
        ph = sp.find("./p:nvSpPr/p:nvPr/p:ph", ns)
        if ph is None or ph.get("type") != "sldNum":
            continue
        clr = sp.find(".//a:solidFill/a:srgbClr", ns)
        if clr is not None and clr.get("val"):
            found_val = clr.get("val")
            break

    if found_val is None:
        logger.info("check_page_number_colors: no sldNum-scoped color found in %s",
                    pptx_file)
        return 0.0

    try:
        r = int(found_val[0:2], 16)
        g = int(found_val[2:4], 16)
        b = int(found_val[4:6], 16)
    except (ValueError, IndexError):
        logger.warning("check_page_number_colors: bad color val %r", found_val)
        return 0.0

    th = 50
    is_red = r > g + th and r > b + th
    is_blue = b > g + th and b > r + th
    is_green = g > r + th and g > b + th
    is_black = r < th and g < th and b < th
    ok = (
        (color_name == "red" and is_red)
        or (color_name == "blue" and is_blue)
        or (color_name == "green" and is_green)
        or (color_name == "black" and is_black)
    )
    logger.info(
        "check_page_number_colors: file=%s rgb=#%s expect=%s -> %s",
        pptx_file, found_val, color_name, ok,
    )
    return 1.0 if ok else 0.0


def literal_match(result: Any, expected: Any, **options) -> float:
    """Byte-equal string compare with `rule`-getter unwrap.

    Upstream OSWorld `literal_match` does `str(result) == str(expected)`. When
    `expected` comes from a `{"type": "rule", "rules": {"expected": "..."}}`
    config, the runner's getter returns the inner `rules` dict (not the string),
    so the comparison stringifies the dict (e.g. ``"{'expected': '...'}"``) and
    never matches. Unwrap the dict here, then delegate to upstream so the
    `type=list` / `ignore_case` options keep working unchanged.
    """
    if isinstance(expected, dict) and "expected" in expected:
        expected = expected["expected"]
    from desktop_env.evaluators.metrics.general import literal_match as _upstream
    return _upstream(result, expected, **options)


# ---------------------------------------------------------------------------
# GIMP memsize serialization override.
# Upstream check_config_status exact-string-compares the value token, but GIMP
# serializes memsize configs (tile-cache-size, undo-size, ...) with unit
# suffixes ('2G', '1024M') via gimp_memsize_serialize — never the raw byte int
# the task rule carries. Every real GUI agent hits '2G' and fails; the oracle
# passes only because it bypasses GIMP and writes the raw int. Compare
# numerically when both sides parse as memsize. Recovers the whole gimp
# memsize-config family. Validated 0->1 with negative controls (512M/1024M stay 0).
# ---------------------------------------------------------------------------
_GIMP_MEMSIZE_UNIT = {"": 1, "b": 1, "k": 1024, "m": 1024**2, "g": 1024**3}


def _gimp_memsize_to_bytes(tok: str) -> int | None:
    import re
    m = re.fullmatch(r"\s*(\d+)\s*([bkmgBKMG]?)\s*", str(tok or ""))
    return int(m.group(1)) * _GIMP_MEMSIZE_UNIT[m.group(2).lower()] if m else None


# VS Code omits a setting from settings.json when it is set to its shipped
# default via the Settings UI (the UI removes the key rather than writing the
# default). An absent whitelisted key therefore reads as "left at default";
# accept it IFF expected == that default. NARROW (2 keys); every other key/value
# is compared exactly — byte-identical to upstream everywhere else.
# ⚠ Landing this override ALONE is a trivial-pass FP: the base image ships
# settings.json WITHOUT these keys and the perturb setup does not seed them, so a
# no-op agent scores 1.0 whenever expected == the default. It is landed PAIRED
# with the perturb generator seeding the non-default starting value
# (gen/train/perturb/vs_code.py::perturb_vscode_bool_setting) so a no-op starts
# off-target -> 0, while a real UI-set-to-default (key removed) still recovers.
_VSCODE_EFFECTIVE_DEFAULTS = {
    "debug.focusEditorOnBreak": True,
    "workbench.editor.wrapTabs": False,
}


def check_json_settings(actual: str, expected: str, **options) -> float:
    """Default-aware override of upstream check_json_settings. Present keys
    exact; an ABSENT expected key passes only if it is a whitelisted VS Code
    effective-default AND expected == that default. Strict superset of upstream —
    the only verdict that changes is {whitelisted key absent AND exp == default}."""
    import json
    if not actual:
        return 0.0
    try:
        with open(actual) as f:
            data = json.load(f)
    except Exception:
        return 0.0
    _MISSING = object()
    for key, value in expected["expected"].items():
        if key in data:
            if data[key] != value:
                return 0.0
        elif _VSCODE_EFFECTIVE_DEFAULTS.get(key, _MISSING) != value:
            return 0.0
    return 1.0


def check_config_status(actual_config_path, rule):
    if actual_config_path is None:
        return 0.0
    with open(actual_config_path) as f:
        content = f.readlines()

    def _val_match(actual: str, expected: str) -> bool:
        if actual == expected:
            return True
        a, b = _gimp_memsize_to_bytes(actual), _gimp_memsize_to_bytes(expected)
        return a is not None and b is not None and a == b

    for line in content:
        if line.startswith("#") or line == "\n":
            continue
        items = line.strip().lstrip("(").rstrip(")\n").split()
        if not items:
            continue
        if isinstance(rule["key"], str):
            if items[0] == rule["key"] and _val_match(items[-1], rule["value"]):
                return 1.0
        elif isinstance(rule["key"], list) and len(rule["key"]) == 2:
            if items[0] == rule["key"][0] and items[1] == rule["key"][1] \
                    and _val_match(items[-1], rule["value"]):
                return 1.0
    return 0.0


# ---------------------------------------------------------------------------
# Thunderbird filter-action alias override.
# The real TB "Add Star" UI serializes action="Mark flagged" (with a space)
# into msgFilterRules.dat, while the catalog gold invented "MarkFlagged" (no
# space) — same semantic action, different serialized spelling -> exact match 0.
# Normalize the action through an alias table before delegating to upstream.
# (An oracle replay cannot reproduce this — its oracle_actions write MarkFlagged
# directly; this is a real-UI-only FN.) Validated 0->1, 4 controls hold.
# ---------------------------------------------------------------------------
def check_thunderbird_filter(result, rules):
    import copy
    import re
    import tempfile
    from desktop_env.evaluators.metrics.thunderbird import (
        check_thunderbird_filter as _upstream,
    )
    if result is None:
        return 0.0
    def _canon(a):
        # tb-filter class (generalizes the 2-entry alias table): Thunderbird
        # serializes filter actions CamelCase (MarkFlagged, MoveToFolder,
        # CopyToFolder, ChangePriority) but the gold uses spaced/underscored labels
        # ("Mark flagged"). Normalize BOTH sides by casefold + strip spaces/underscores
        # -> the whole action family, not a hand-listed alias set. Genuinely-different
        # actions stay distinct (MoveToFolder != Delete), so no new false-positive.
        return re.sub(r"[ _]+", "", a.strip().lower()) if isinstance(a, str) else a

    text = re.sub(r'action="([^"]*)"',
                  lambda m: 'action="%s"' % _canon(m.group(1)),
                  open(result).read())
    rules2 = copy.deepcopy(rules)
    for bucket in ("expect", "unexpect"):
        for r in rules2.get(bucket, []):
            if "action" in r:
                r["action"] = _canon(r["action"])
    import os
    with tempfile.NamedTemporaryFile("w", suffix=".dat", delete=False) as f:
        f.write(text)
        tmp = f.name
    try:
        return _upstream(tmp, rules2)
    finally:
        os.unlink(tmp)


def check_list(result, rules):
    """Override upstream check_list for the thunderbird_a10b69e1 class (17
    `.msf` rows): Thunderbird auto-generates `.msf` index sidecars with its own
    title-case (e.g. `Drafts.msf` for a folder named DRAFTS), while the eval
    `expect` regexes are uppercase (`\\bDRAFTS\\.msf\\b`). Match `.msf`-sidecar
    patterns case-INSENSITIVELY; leave folder/mbox-name patterns (user-controlled)
    case-sensitive. A pattern targets a `.msf` sidecar iff it POSITIVELY matches
    `.msf` — contains `\\.msf` AND not inside a negative lookahead `(?!...)` (the
    mbox patterns `\\bDRAFTS/?(?!\\.msf)` contain `.msf` only in a lookahead, so a
    naive `"msf" in ptt` would wrongly relax them and break the lowercase-mbox
    control)."""
    import re
    if result is None:
        return 0.0

    def _compile(ptt):
        is_msf = (r"\.msf" in ptt) and ("(?!" not in ptt)
        return re.compile(ptt, re.IGNORECASE if is_msf else 0)

    expect = [_compile(p) for p in rules.get("expect", [])]
    unexpect = [_compile(p) for p in rules.get("unexpect", [])]
    expect_metrics = [False] * len(expect)
    unexpect_metric = True
    with open(result) as f:
        for line in f:
            for i, r in enumerate(expect):
                expect_metrics[i] = expect_metrics[i] or (r.search(line) is not None)
            unexpect_metric = unexpect_metric and all(r.search(line) is None for r in unexpect)
    return float(all(expect_metrics) and unexpect_metric)
