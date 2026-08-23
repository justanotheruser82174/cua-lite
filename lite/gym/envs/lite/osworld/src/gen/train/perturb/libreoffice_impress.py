"""Impress domain per-task perturbation (Track B, TYPE_1 + TYPE_2).

14 op types with deterministic oracle (python-pptx). TYPE_1 keeps eval op
list, resamples (slide_idx, param). TYPE_2 uses structurally-different variants.

Usage:
    uv run python -m lite.gym.envs.lite.osworld.src.gen.train --track perturb --domain libreoffice_impress
"""

from __future__ import annotations

import copy
import json
import random
from functools import partial
from pathlib import Path

from lite.gym.envs.lite.osworld.src.gen.common import (
    LO_SAVE_POSTCONFIG,
)
from lite.gym.envs.lite.osworld.src.gen.train.perturb._utils import (
    make_perturb_row,
)

PERTURB_TYPE_1 = "type1"
PERTURB_TYPE_2 = "type2"
# TYPE_3 = atomic eval-evaluator perturb. Targets one of 8 impress-specific
# evaluators (`check_presenter_console_disable`, `check_auto_saving_time`,
# `evaluate_presentation_fill_to_rgb_distance`, `compare_images`,
# `check_image_stretch_and_center`, `check_page_number_colors`,
# `compare_audios`, `check_slide_orientation_Portrait`). Each TYPE_3 row uses
# the eval task's downloaded source file as base, mutates initial state via a
# config_step (so reward=0 before agent acts), and supplies an oracle that
# writes the gold state via raw zip/xml/python-pptx manipulation. The
# evaluator is the actual eval-side evaluator with rules resampled where the
# evaluator accepts parameter rules; binary evaluators (presenter_console,
# image_stretch_and_center, slide_orientation_Portrait) keep eval rules
# unchanged but vary which source pptx the oracle constructs from. See
# `devs/envs/lite.osworld/perturb/libreoffice_impress.md` §"TYPE_3".
PERTURB_TYPE_3 = "type3"

# ---------------------------------------------------------------------------
# Pools
# ---------------------------------------------------------------------------

_COLOR_POOL = [
    ("red", 255, 0, 0), ("blue", 0, 0, 255), ("green", 0, 128, 0),
    ("yellow", 255, 255, 0), ("black", 0, 0, 0), ("white", 255, 255, 255),
    ("orange", 255, 165, 0), ("purple", 128, 0, 128), ("pink", 255, 192, 203),
]
_BG_COLOR_POOL = [
    ("red", 255, 0, 0), ("blue", 0, 0, 255), ("green", 0, 128, 0),
    ("yellow", 255, 255, 0), ("purple", 128, 0, 128), ("orange", 255, 165, 0),
    ("cyan", 0, 255, 255), ("white", 255, 255, 255), ("black", 0, 0, 0),
]
_RGB_TO_NAME: dict[tuple, str] = {
    (c[1], c[2], c[3]): c[0]
    for c in _COLOR_POOL + _BG_COLOR_POOL
}
_SIZE_POOL = [10, 12, 14, 16, 18, 20, 24, 28, 32, 36, 40, 44, 48, 60, 72]
_STYLE_POOL = ["bold", "italic", "underline", "strikethrough"]
# Validation (perturb validation): Calibri / Georgia not installed in LO VM
# (per writer validation; symptom: agent reports infeasible).
_FONT_POOL = ["Arial", "Times New Roman", "Courier New", "Verdana"]  # validation: drop Trebuchet MS too (validation edb61b14 agent reports not installed)
_TITLE_POOL = ["Introduction", "Summary", "Overview", "Conclusion", "Highlights",
               "References", "Agenda", "Results", "Discussion", "Background"]
_NOTE_POOL = [
    "Key point here.", "Remember to elaborate.", "Add more context.",
    "Important note.", "Follow up needed.", "Check data sources.",
    "Expand with examples.", "Review before presenting.",
]
_ALIGN_POOL = ["left", "center", "right", "justify"]
_TRANSITION_POOL = ["dissolve", "fade", "wipe", "push"]  # uncover: LO normalization drops <p:uncover>, oracle always scores 0
_PIC_CM_POOL = [5, 8, 10, 12, 15, 18, 20, 25]
_TABLE_POOL = [(3, 2), (4, 3), (5, 2), (6, 4)]

# ---------------------------------------------------------------------------
# Analysis loader
# ---------------------------------------------------------------------------

# Committed, reproducible slide-analysis (which slides have text/pics/tables per
# impress source). Was an EPHEMERAL /tmp/impress_full.json built by a manual
# script and ABSENT at regen -> every TYPE_1 op fell back to slide 1, which for
# several sources has no text/runs -> a no-op perturb -> expected==source ->
# trivial pass (eval=1.0 before the oracle). Baking it next to the generator
# makes regen deterministic and routes ops to real content slides. Rebuild via
# devs/envs/lite.osworld/perturb/libreoffice_impress.md when sources change.
# (#154 oracle-validation impress trivial_pass, + ~35 latent slide-1 no-ops.)
_ANALYSIS_PATH = Path(__file__).with_name("impress_analysis.json")
_ANALYSIS: dict | None = None


def _load_analysis() -> dict:
    global _ANALYSIS
    if _ANALYSIS is None:
        _ANALYSIS = json.loads(_ANALYSIS_PATH.read_text()) if _ANALYSIS_PATH.exists() else {}
    return _ANALYSIS


def _info(eval_row: dict) -> dict:
    return _load_analysis().get(eval_row["task_id"].split("_")[-1], {})


def _slide_list(info: dict, kind: str) -> list[int]:
    """1-based slide indices having the given kind (texts/pics/tables)."""
    return [s["idx"] for s in info.get("slides", []) if s.get(kind)]


# ---------------------------------------------------------------------------
# Core helpers
# ---------------------------------------------------------------------------

def _get_file_path(eval_row: dict) -> str | None:
    for step in eval_row["metadata"]["config"]:
        if step.get("type") == "download":
            files = step["parameters"].get("files", [])
            if files:
                return files[0]["path"]
        if step.get("type") == "open":
            return step["parameters"].get("path")
    return None


def _lo_normalize_cmd(path: str, fmt: str) -> str:
    return (
        f"tmpd=$(mktemp -d) && "
        f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
        f"--convert-to {fmt} --outdir \"$tmpd\" '{path}' 2>/dev/null && "
        f"[ -f \"$tmpd/$(basename '{path}')\" ] && "
        f"cp \"$tmpd/$(basename '{path}')\" '{path}'; "
        f"rm -rf \"$tmpd\"; true"
    )


def _make_config_step(py_code: str) -> dict:
    return {"type": "execute", "parameters": {
        "command": f"python3 << 'PYEOF'\n{py_code}\nPYEOF", "shell": True,
    }}


def _standard_oracle(expected_py: str, file_path: str, expected_path: str) -> list[dict]:
    # LO normalize is NOT idempotent (pptx shape geometry shrinks ~360 EMU per
    # pass). Build expected (step 1), normalize it once (step 2), then cp to
    # file_path (step 3) as a byte-identical copy. The eval result-getter adds
    # exactly one more normalize pass to BOTH sides -> result = expected = f(f1).
    # The old 4th action (re-normalize file_path) gave result an extra pass ->
    # drift -> spurious 0. Removed.
    return [
        {"type": "execute", "parameters": {"command": f"python3 << 'PYEOF'\n{expected_py}\nPYEOF", "shell": True}},
        {"type": "execute", "parameters": {"command": _lo_normalize_cmd(expected_path, "pptx"), "shell": True}},
        {"type": "execute", "parameters": {"command": f"cp '{expected_path}' '{file_path}'", "shell": True}},
    ]


# ---------------------------------------------------------------------------
# Evaluator builder
# ---------------------------------------------------------------------------

_BASE_OPTS: dict = {
    "examine_shape": False,
    "examine_font_name": False,
    "examine_font_size": False,
    "examine_alignment": False,
    "examine_run_count": False,
    "examine_color_rgb": False,
    "examine_indent": False,
    "examine_strike_through": False,
    "examine_font_underline": False,
    "examine_font_bold": False,
    "examine_font_italic": False,
    "examine_bullets": False,
    "examine_background_color": False,
    "examine_note": False,
}

_STYLE_FLAG = {
    "bold": "examine_font_bold",
    "italic": "examine_font_italic",
    "underline": "examine_font_underline",
    "strikethrough": "examine_strike_through",
}


def _op_flags(op: tuple) -> dict:
    t = op[0]
    if t == "set_font_color":
        return {"examine_color_rgb": True}
    if t == "set_font_size":
        return {"examine_font_size": True}
    if t == "set_font_style":
        return {_STYLE_FLAG.get(op[2], "examine_font_bold"): True}
    if t == "set_font_name":
        return {"examine_font_name": True}
    if t == "set_background_color":
        return {"examine_background_color": True}
    if t == "set_text_alignment":
        return {"examine_alignment": True}
    if t == "set_picture_size":
        return {"examine_modify_height": True}
    if t == "move_object":
        # #155 §C #21/#22/#23: emit the reconcile mode. Direction (edge) and object are
        # already decided at op-build time (_expected_snippet L478); the comparator resolves
        # "title" by placeholder idx-0 and snaps position via region-OR-tolerance, so a tight
        # wide-shape move (#23) that cannot leave the middle third still passes -- no distinct
        # mode needed. All emitted directions are edges (top/bottom/left/right); never "center".
        edge = op[2]
        obj = op[3] if len(op) > 3 else "picture"
        shape_type = {"title": "title", "picture": 13, "table": 19}[obj]
        return {
            "position_mode": "region",
            "position_region": {"slide": op[1] - 1, "shape_type": shape_type, "edge": edge},
            "position_tolerance_emu": 50000,
        }
    if t == "insert_table":
        # #155 §C #5: instruction fixes no geometry -> ignore the new table's L/T/W/H; the
        # table row/col + cell-text checks (upstream, unconditional) carry correctness.
        return {
            "position_mode": "ignore",
            "ignore_shape_geometry": [[op[1] - 1, 19]],
        }
    if t == "add_speaker_note":
        return {"examine_note": True}
    return {}


def _build_evaluator(ops: list, file_path: str, expected_path: str) -> dict:
    file_name = file_path.rsplit("/", 1)[-1]
    t_ops = [op for op in ops if op[0] == "set_slide_transition"]
    o_ops = [op for op in ops if op[0] != "set_slide_transition"]

    options = dict(_BASE_OPTS)
    for op in o_ops:
        options.update(_op_flags(op))

    if not t_ops:
        return {
            "func": "compare_pptx_files",
            "result": {"type": "vm_file", "path": file_path, "dest": file_name},
            "expected": {"type": "vm_file", "path": expected_path, "dest": "expected_file"},
            "options": options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    slide_0, transition = t_ops[0][1] - 1, t_ops[0][2]
    if not o_ops:
        return {
            "func": "check_transition",
            "result": {"type": "vm_file", "path": file_path, "dest": file_name},
            "expected": {"type": "rule", "rules": {"slide_idx": slide_0, "transition_type": transition}},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return {
        "func": ["check_transition", "compare_pptx_files"],
        "result": [
            {"type": "vm_file", "path": file_path, "dest": file_name},
            {"type": "vm_file", "path": file_path, "dest": "result_cmp.pptx"},
        ],
        "expected": [
            {"type": "rule", "rules": {"slide_idx": slide_0, "transition_type": transition}},
            {"type": "vm_file", "path": expected_path, "dest": "expected_file"},
        ],
        "options": [{}, options],
        "postconfig": LO_SAVE_POSTCONFIG,
    }


# ---------------------------------------------------------------------------
# Per-op code snippet generators
# ---------------------------------------------------------------------------

def _slide_run_op_lines(s: int, run_stmt: str) -> list[str]:
    """Gold/config lines applying ``run_stmt`` (a statement operating on ``_r``)
    to EVERY run in EVERY text frame AND EVERY table cell on slide ``s``.

    pptx tables have ``has_text_frame=False``, so a ``has_text_frame``-only loop
    silently skips all table text — making a "whole slide / every text run" gold
    UNDER-specify the table (it leaves cells at their original color/size). The
    evaluator (compare_pptx_files) DOES walk table cells, so an agent that
    correctly recolours the whole slide then mismatches the table-blind gold and
    false-fails. Recursing ``has_table`` here keeps gold and agent consistent;
    table-less slides hit neither table branch, so this is a no-op for them.
    (validation: 15aece23_59c09dd7 — blue "whole slide" left the table white.)
    """
    return [
        f"for _sh in prs.slides[{s}].shapes:",
        f"  if _sh.has_text_frame:",
        f"    for _p in _sh.text_frame.paragraphs:",
        f"      for _r in _p.runs: {run_stmt}",
        f"  elif _sh.has_table:",
        f"    for _row in _sh.table.rows:",
        f"      for _cell in _row.cells:",
        f"        for _p in _cell.text_frame.paragraphs:",
        f"          for _r in _p.runs: {run_stmt}",
    ]


def _config_snippet(op: tuple) -> tuple[set[str], list[str]]:
    t, s1 = op[0], op[1]
    s = s1 - 1  # 0-based

    if t == "set_font_color":
        return (
            {"from pptx.dml.color import RGBColor"},
            _slide_run_op_lines(s, "_r.font.color.rgb = RGBColor(128,128,128)"),
        )
    if t == "set_font_size":
        pt = op[2]; cfg = 14 if pt == 12 else 12
        return (
            {"from pptx.util import Pt"},
            _slide_run_op_lines(s, f"_r.font.size = Pt({cfg})"),
        )
    if t == "set_font_style":
        style = op[2]
        if style == "bold":
            return (set(), _slide_run_op_lines(s, "_r.font.bold = False"))
        if style == "italic":
            return (set(), _slide_run_op_lines(s, "_r.font.italic = False"))
        if style == "underline":
            return (set(), _slide_run_op_lines(s, "_r.font.underline = False"))
        # strikethrough
        return (set(), _slide_run_op_lines(
            s, "_r._r.get_or_add_rPr().attrib.pop('strike', None)"))
    if t == "set_font_name":
        font = op[2]; cfg = "Arial" if font != "Arial" else "Verdana"
        return (set(), _slide_run_op_lines(s, f"_r.font.name = '{cfg}'"))
    if t == "set_background_color":
        return (
            {"from pptx.dml.color import RGBColor"},
            [f"_bg = prs.slides[{s}].background",
             f"_bg.fill.solid()",
             f"_bg.fill.fore_color.rgb = RGBColor(128,128,128)"],
        )
    if t == "set_title_text":
        return (set(), [f"_tsh = prs.slides[{s}].shapes.title",
                        f"if _tsh: _tsh.text = 'TempTitle'"])
    if t == "set_text_alignment":
        return (
            {"from pptx.enum.text import PP_ALIGN"},
            [f"for _sh in prs.slides[{s}].shapes:",
             f"  if _sh.has_text_frame:",
             f"    for _p in _sh.text_frame.paragraphs: _p.alignment = PP_ALIGN.LEFT"],
        )
    if t == "set_picture_size":
        h = op[2]; cfg_h = 8 if h != 8 else 10
        return (set(), [
            f"for _sh in prs.slides[{s}].shapes:",
            f"  if _sh.shape_type == 13:",
            f"    _oh={{}};_oh['v']=_sh.height; _nh=int({cfg_h}*360000)",
            f"    if _oh['v']>0: _sh.width=int(_sh.width*_nh/_oh['v'])",
            f"    _sh.height=_nh; break",
        ])
    if t == "add_speaker_note":
        return (set(), [f"prs.slides[{s}].notes_slide.notes_text_frame.text = ''"])
    if t == "set_slide_transition":
        return (
            {"from pptx.oxml.ns import qn as _qn"},
            [f"_sld = prs.slides[{s}]._element",
             f"for _tr in list(_sld.findall(_qn('p:transition'))): _sld.remove(_tr)"],
        )
    if t == "edit_table_cell":
        row_sel = op[2]
        return (set(), [
            f"for _sh in prs.slides[{s}].shapes:",
            f"  if hasattr(_sh,'table'):",
            f"    _ri={{'first':0,'second':1,'last':-1}}['{row_sel}']",
            f"    if _ri<0: _ri=len(_sh.table.rows)+_ri",
            f"    for _ci,_c in enumerate(_sh.table.rows[_ri].cells): _c.text=chr(65+_ci)",
            f"    break",
        ])
    if t == "move_object":
        obj = op[3] if len(op) > 3 else "picture"
        if obj == "title":
            return (set(), [
                f"_sw=prs.slide_width or 9144000; _sh2=prs.slide_height or 6858000",
                f"_tsh=prs.slides[{s}].shapes.title",
                f"if _tsh:",
                f"  _tsh.left=int((_sw-_tsh.width)/2); _tsh.top=int((_sh2-_tsh.height)/2)",
            ])
        shape_t = 13 if obj == "picture" else 19
        return (set(), [
            f"_sw=prs.slide_width or 9144000; _sh2=prs.slide_height or 6858000",
            f"for _sh in prs.slides[{s}].shapes:",
            f"  if _sh.shape_type=={shape_t}:",
            f"    _sh.left=int((_sw-_sh.width)/2); _sh.top=int((_sh2-_sh.height)/2); break",
        ])
    # insert_table, reorder_slides, add_bullet_point: no initial state change
    return (set(), [])


_NS_PPTX = "http://schemas.openxmlformats.org/presentationml/2006/main"
_TR_INNER = {
    "dissolve": "<p:dissolve/>",
    "fade":     "<p:fade/>",
    "wipe":     "<p:wipe/>",
    "push":     "<p:push/>",
}


def _expected_snippet(op: tuple) -> tuple[set[str], list[str]]:
    t, s1 = op[0], op[1]
    s = s1 - 1  # 0-based

    if t == "set_font_color":
        r, g, b = op[2]
        return (
            {"from pptx.dml.color import RGBColor"},
            _slide_run_op_lines(s, f"_r.font.color.rgb = RGBColor({r},{g},{b})"),
        )
    if t == "set_font_size":
        pt = op[2]
        return (
            {"from pptx.util import Pt"},
            _slide_run_op_lines(s, f"_r.font.size = Pt({pt})"),
        )
    if t == "set_font_style":
        style = op[2]
        if style == "bold":
            return (set(), _slide_run_op_lines(s, "_r.font.bold = True"))
        if style == "italic":
            return (set(), _slide_run_op_lines(s, "_r.font.italic = True"))
        if style == "underline":
            return (set(), _slide_run_op_lines(s, "_r.font.underline = True"))
        return (set(), _slide_run_op_lines(
            s, "_r._r.get_or_add_rPr().set('strike','sngStrike')"))
    if t == "set_font_name":
        font = op[2].replace("'", "\\'")
        return (set(), _slide_run_op_lines(s, f"_r.font.name = '{font}'"))
    if t == "set_background_color":
        r, g, b = op[2]
        return (
            {"from pptx.dml.color import RGBColor"},
            [f"_bg = prs.slides[{s}].background",
             f"_bg.fill.solid()",
             f"_bg.fill.fore_color.rgb = RGBColor({r},{g},{b})"],
        )
    if t == "set_title_text":
        text = op[2].replace("'", "\\'")
        return (set(), [
            f"_sl={{}};_sl['s']=prs.slides[{s}]; _tsh=_sl['s'].shapes.title",
            f"if _tsh: _tsh.text='{text}'",
            f"else:",
            f"  for _sh in _sl['s'].shapes:",
            f"    if _sh.has_text_frame: _sh.text='{text}'; break",
        ])
    if t == "set_text_alignment":
        pp = {"left": "LEFT", "center": "CENTER", "right": "RIGHT", "justify": "JUSTIFY"}[op[2]]
        return (
            {"from pptx.enum.text import PP_ALIGN"},
            [f"for _sh in prs.slides[{s}].shapes:",
             f"  if _sh.has_text_frame:",
             f"    for _p in _sh.text_frame.paragraphs: _p.alignment = PP_ALIGN.{pp}"],
        )
    if t == "set_picture_size":
        h_cm = op[2]
        return (set(), [
            f"for _sh in prs.slides[{s}].shapes:",
            f"  if _sh.shape_type == 13:",
            f"    _oh={{}};_oh['v']=_sh.height; _nh=int({h_cm}*360000)",
            f"    if _oh['v']>0: _sh.width=int(_sh.width*_nh/_oh['v'])",
            f"    _sh.height=_nh; break",
        ])
    if t == "add_speaker_note":
        note = op[2].replace("'", "\\'")
        return (set(), [f"prs.slides[{s}].notes_slide.notes_text_frame.text = '{note}'"])
    if t == "set_slide_transition":
        tr = op[2]
        xml = f'<p:transition xmlns:p="{_NS_PPTX}">{_TR_INNER[tr]}</p:transition>'
        return (
            {"from lxml import etree as _etree", "from pptx.oxml.ns import qn as _qn"},
            [
                f"_sld = prs.slides[{s}]._element",
                f"for _tr in list(_sld.findall(_qn('p:transition'))): _sld.remove(_tr)",
                f"_sld.append(_etree.fromstring({repr(xml)}))",
            ],
        )
    if t == "insert_table":
        rows, cols = op[2], op[3]
        return (
            {"from pptx.util import Inches"},
            [
                f"_tbl=prs.slides[{s}].shapes.add_table({rows},{cols},Inches(1),Inches(2),Inches(8),Inches(3)).table",
                f"for _r in _tbl.rows:",
                f"  for _c in _r.cells: _c.text=''",
            ],
        )
    if t == "edit_table_cell":
        row_sel = op[2]; cells = op[3]
        return (set(), [
            f"for _sh in prs.slides[{s}].shapes:",
            f"  if hasattr(_sh,'table'):",
            f"    _ri={{'first':0,'second':1,'last':-1}}['{row_sel}']",
            f"    if _ri<0: _ri=len(_sh.table.rows)+_ri",
            f"    _nc={repr(cells)}",
            f"    for _ci,_c in enumerate(_sh.table.rows[_ri].cells):",
            f"      if _ci<len(_nc): _c.text=_nc[_ci]",
            f"    break",
        ])
    if t == "reorder_slides":
        src_0 = op[1] - 1; dst_0 = op[2] - 1
        return (set(), [
            f"_lst=prs.slides._sldIdLst",
            f"_e=_lst[{src_0}]; _lst.remove(_e); _lst.insert({dst_0},_e)",
        ])
    if t == "move_object":
        direction = op[2]
        obj = op[3] if len(op) > 3 else "picture"
        pos_lines = {
            "top":    ["    _sh.top = 457200"],
            "bottom": ["    _sh.top = _sh2 - _sh.height - 457200"],
            "left":   ["    _sh.left = 457200"],
            "right":  ["    _sh.left = _sw - _sh.width - 457200"],
            "center": ["    _sh.left=int((_sw-_sh.width)/2); _sh.top=int((_sh2-_sh.height)/2)"],
        }[direction]
        if obj == "title":
            return (set(), [
                f"_sw=prs.slide_width or 9144000; _sh2=prs.slide_height or 6858000",
                f"_sh=prs.slides[{s}].shapes.title",
                f"if _sh:",
            ] + [f"  {ln.strip()}" for ln in pos_lines])
        shape_t = 13 if obj == "picture" else 19
        return (set(), [
            f"_sw=prs.slide_width or 9144000; _sh2=prs.slide_height or 6858000",
            f"for _sh in prs.slides[{s}].shapes:",
            f"  if _sh.shape_type=={shape_t}:",
        ] + pos_lines + ["    break"])
    if t == "add_bullet_point":
        text = op[2].replace("'", "\\'")
        return (set(), [
            f"_sl=prs.slides[{s}]",
            f"for _sh in _sl.shapes:",
            # #155 impress_f23acfd2 (GOLD): a bullet belongs in the body/OBJECT
            # placeholder, not the title. Skip title (1) + centered-title (3)
            # placeholders so add_paragraph lands in the body text frame -- matching
            # what an agent adding a bullet point does. General over all bullet tasks.
            f"  if _sh.has_text_frame and not (_sh.is_placeholder and _sh.placeholder_format.type in (1, 3)):",
            f"    _tf=_sh.text_frame",
            f"    _para=_tf.add_paragraph()",
            f"    _para.text='{text}'",
            f"    break",
        ])
    return (set(), [])


# ---------------------------------------------------------------------------
# Script builder
# ---------------------------------------------------------------------------

def _build_py(ops: list, base_path: str, save_path: str, snippet_fn) -> str:
    imports: set[str] = {"from pptx import Presentation"}
    body: list[str] = [f"prs = Presentation('{base_path}')"]
    for op in ops:
        extra, lines = snippet_fn(op)
        imports |= extra
        body.extend(lines)
    body.append(f"prs.save('{save_path}')")
    return "\n".join(sorted(imports)) + "\n" + "\n".join(body)


# ---------------------------------------------------------------------------
# Instruction builder
# ---------------------------------------------------------------------------

def _op_instr(op: tuple, rng: random.Random) -> str:
    """Per-op instruction paraphrase pool (5 paraphrases per op type).

    Pools mix imperative / context-then-action / action-with-rationale styles
    with varied length (8-28 words). Eval baseline is ~22.6 avg_words, ~17%
    polite, so the longer paraphrases keep avg_words on target while
    `_build_instruction` wraps with a polite prefix at ~17-25% realised rate.
    """
    t, s1 = op[0], op[1]
    if t == "set_font_color":
        name = _RGB_TO_NAME.get(op[2], str(op[2])) if isinstance(op[2], tuple) else str(op[2])
        # Validation note: agent picked "Indigo" when instruction said
        # "blue" (8979838c trace) — color name ambiguity. Append explicit RGB
        # hex spec so the agent can disambiguate against LO's color picker
        # (which has named shades like Indigo/Light Blue/Navy under the broad
        # "blue" umbrella). evaluator (compare_pptx_files with
        # examine_color_rgb=True) checks exact RGB so the hex is what matters.
        if isinstance(op[2], tuple):
            r, g, b = op[2]
            hex_spec = f"#{r:02X}{g:02X}{b:02X}"
            color_str = f"{name} (RGB {hex_spec})"
        else:
            color_str = name
        return rng.choice([
            f"Change the font color of all text on slide {s1} to {color_str}",
            f"Set all text on slide {s1} to {color_str}",
            f"Make the font color {color_str} for every text element on slide {s1}",
            f"On slide {s1}, switch every text run's color over to {color_str} so the whole slide reads in one color",
            f"For slide {s1}, recolor each text element to {color_str} — every textbox should pick up the same shade",
        ])
    if t == "set_font_size":
        return rng.choice([
            f"Set the font size of all text on slide {s1} to {op[2]} pt",
            f"Change the text size on slide {s1} to {op[2]} points",
            f"Resize all text on slide {s1} to {op[2]} pt",
            f"On slide {s1}, bump every text run up to {op[2]} pt so the whole slide reads at one size",
            f"Make the text on slide {s1} display at {op[2]} pt across every textbox and placeholder",
        ])
    if t == "set_font_style":
        style = op[2]
        return rng.choice([
            f"Apply {style} formatting to all text on slide {s1}",
            f"Make all text on slide {s1} {style}",
            f"On slide {s1}, mark every text run as {style} across all the textboxes",
            f"Turn the text on slide {s1} {style} so each line picks up the {style} styling",
            f"For slide {s1}, switch every text element over to {style} formatting",
        ])
    if t == "set_font_name":
        return rng.choice([
            f"Change the font of all text on slide {s1} to {op[2]}",
            f"Set the typeface of all text on slide {s1} to {op[2]}",
            f"Apply the {op[2]} font to all text on slide {s1}",
            f"On slide {s1}, switch every text element over to the {op[2]} typeface for a consistent look",
            f"Re-typeset the text on slide {s1} in {op[2]} so all of the textboxes share one font family",
        ])
    if t == "set_background_color":
        name = _RGB_TO_NAME.get(op[2], str(op[2])) if isinstance(op[2], tuple) else str(op[2])
        # See validation set_font_color comment: append explicit RGB hex to
        # disambiguate against LO's color picker shades (Indigo / Light Blue /
        # Navy under broad "blue").
        if isinstance(op[2], tuple):
            r, g, b = op[2]
            color_str = f"{name} (RGB #{r:02X}{g:02X}{b:02X})"
        else:
            color_str = name
        return rng.choice([
            f"Change the background color of slide {s1} to {color_str}",
            f"Set slide {s1}'s background to {color_str}",
            f"Make the background of slide {s1} {color_str}",
            f"For slide {s1}, repaint the slide background in solid {color_str}",
            f"On slide {s1}, swap the background fill over to a {color_str} color so the slide reads {color_str} behind the content",
        ])
    if t == "set_title_text":
        return rng.choice([
            f'Change the title of slide {s1} to "{op[2]}"',
            f'Set the title on slide {s1} to "{op[2]}"',
            f'Update slide {s1}\'s title text to "{op[2]}"',
            f'On slide {s1}, replace whatever\'s in the title placeholder with "{op[2]}"',
            f'Rewrite the title heading on slide {s1} to read "{op[2]}"',
        ])
    if t == "set_text_alignment":
        return rng.choice([
            f"Set the text alignment on slide {s1} to {op[2]}",
            f"Align all text on slide {s1} to the {op[2]}",
            f"Change the paragraph alignment on slide {s1} to {op[2]}",
            f"On slide {s1}, realign every paragraph in every textbox to {op[2]}",
            f"For slide {s1}, switch the text alignment over to {op[2]} across all of the text frames",
        ])
    if t == "set_picture_size":
        return rng.choice([
            f"Resize the picture on slide {s1} to have a height of {op[2]} cm",
            f"Set the image height on slide {s1} to {op[2]} cm",
            f"Change the picture's height on slide {s1} to {op[2]} centimeters",
            f"On slide {s1}, scale the picture so its height comes out to exactly {op[2]} cm while keeping the aspect ratio",
            f"Resize the image on slide {s1} so it stands {op[2]} cm tall — adjust the width proportionally to match",
        ])
    if t == "add_speaker_note":
        return rng.choice([
            f'Add the speaker note "{op[2]}" to slide {s1}',
            f'Set the notes for slide {s1} to "{op[2]}"',
            f'Insert "{op[2]}" as the speaker note on slide {s1}',
            f'On slide {s1}, drop "{op[2]}" into the speaker notes pane so it shows up during presenter mode',
            f'For slide {s1}, write "{op[2]}" into the notes section underneath the slide',
        ])
    if t == "set_slide_transition":
        return rng.choice([
            f"Set the slide transition of slide {s1} to {op[2]}",
            f"Apply a {op[2]} transition to slide {s1}",
            f"Change the transition effect on slide {s1} to {op[2]}",
            f"On slide {s1}, configure the entry transition to use the {op[2]} effect",
            f"For slide {s1}, switch the slide transition over to {op[2]} so it animates in with that effect",
        ])
    if t == "insert_table":
        return rng.choice([
            f"Insert a {op[2]}-row by {op[3]}-column table on slide {s1}",
            f"Add a table with {op[2]} rows and {op[3]} columns to slide {s1}",
            f"Place a {op[2]}×{op[3]} table on slide {s1}",
            f"On slide {s1}, drop in a new empty table sized {op[2]} rows by {op[3]} columns",
            f"For slide {s1}, insert a fresh {op[2]} row, {op[3]} column table so we can fill it in later",
        ])
    if t == "edit_table_cell":
        cells_str = ", ".join(f'"{c}"' for c in op[3])
        return rng.choice([
            f"In the table on slide {s1}, set the {op[2]} row cells to {cells_str}",
            f"Update the {op[2]} row of the table on slide {s1} with the values {cells_str}",
            f"Fill the {op[2]} row of the table on slide {s1} with {cells_str}",
            f"On slide {s1}, edit the table's {op[2]} row so the cells read {cells_str} from left to right",
            f"For the table on slide {s1}, replace the contents of the {op[2]} row with {cells_str}",
        ])
    if t == "reorder_slides":
        return rng.choice([
            f"Move slide {s1} to position {op[2]} in the presentation",
            f"Reorder the slides so that slide {s1} becomes slide {op[2]}",
            f"Place slide {s1} at index {op[2]} in the slide deck",
            f"In the slide panel, drag slide {s1} into position {op[2]} so it moves to that spot in the deck",
            f"Reorder the deck by moving slide {s1} to the {op[2]} position in the running order",
        ])
    if t == "move_object":
        obj = op[3] if len(op) > 3 else "picture"
        return rng.choice([
            f"Move the {obj} on slide {s1} to the {op[2]} of the slide",
            f"Reposition the {obj} on slide {s1} to the {op[2]}",
            f"Place the {obj} on slide {s1} at the {op[2]} of the slide",
            f"On slide {s1}, drag the {obj} over to the {op[2]} edge of the slide canvas",
            f"For slide {s1}, shift the {obj} so it sits at the {op[2]} of the slide rather than where it is now",
        ])
    if t == "add_bullet_point":
        return rng.choice([
            f'Add the bullet point "{op[2]}" to slide {s1}',
            f'Insert a new bullet "{op[2]}" on slide {s1}',
            f'Append the bullet point "{op[2]}" to the text on slide {s1}',
            f'On slide {s1}, tack on a new bullet that reads "{op[2]}" at the end of the existing list',
            f'For slide {s1}, append "{op[2]}" as an additional bullet underneath the existing items',
        ])
    return f"Perform {t} on slide {s1}"


# Polite prefixes — applied at ~22% rate; eval baseline is ~17%.
_POLITE_PREFIXES_SINGLE = (
    "Could you help me",
    "Please",
    "I need to",
    "Can you",
    "I want to",
    "I'd like to",
)
_POLITE_PREFIXES_MULTI = (
    "Please",
    "I need to",
    "I want to",
    "I'd like to",
)
# Action verbs that look natural after a polite prefix ("Please <verb> ...").
# Excludes prepositions like "on"/"for"/"in" deliberately — polite-wrapping a
# preposition-led paraphrase ("Please on slide N, ...") reads ungrammatically.
# This means roughly half of paraphrases are not eligible for polite wrapping,
# which together with the random gate keeps realised polite% near eval baseline.
_ACTION_VERBS = frozenset({
    "change", "set", "apply", "move", "resize", "reposition", "place",
    "add", "insert", "update", "fill", "make", "align", "reorder",
    "turn", "switch", "drop", "drag", "shift", "rewrite", "scale", "bump",
    "tack", "append", "edit", "configure", "rewrite",
})


def _build_instruction(ops: list, rng: random.Random) -> str:
    """Build a natural-English instruction from an ops list.

    Multi-op fusion strategy:
      - 1 op  -> single sentence (most common path)
      - 2 ops -> "<op1>, and <op2>." single-sentence "and"-fusion (~60%);
                 "First <op1>, then <op2>." multi-step ordered fusion (~25%);
                 rare ". " fallback (~15%)
      - 3+ ops -> ". " (capital-letter new sentences, ~50%);
                  ordered "First X. Then, Y. Finally, Z." (~25%);
                  rare ". Also,/Next,/Additionally," connectors (~25%)

    Eval baseline (impress) shows ~21% of instructions carry multi-step
    sequence keywords (then / next / first / once / after / before / finally),
    so the 2-op "First X, then Y." path and the 3+-op ordered path push
    perturb multi-step% toward that band without inflating eval-rare
    connectors (Also/Additionally) above their ~2% share.
    """
    parts = [_op_instr(op, rng) for op in ops]

    if len(parts) == 1:
        core = parts[0] + "."
    elif len(parts) == 2:
        r = rng.random()
        if r < 0.60:
            # Single-sentence "and"-fusion (grammatically safe when both
            # parts start with imperative verbs; lowercase-fold below works
            # because paraphrases start with capitalised verbs/prepositions).
            p0 = parts[0].rstrip(".")
            p1 = parts[1][0].lower() + parts[1][1:]
            core = f"{p0}, and {p1}."
        elif r < 0.85:
            # Multi-step ordered fusion: "First <op1>, then <op2>." adds
            # explicit sequence keywords (first/then) that match eval's
            # multi-step phrasing without inflating Also/Additionally.
            p0 = parts[0][0].lower() + parts[0][1:]
            p1 = parts[1][0].lower() + parts[1][1:]
            core = f"First {p0.rstrip('.')}, then {p1}."
        else:
            # Bare ". " (two new sentences). Eval-rare connector path.
            core = parts[0] + ". " + parts[1] + "."
    else:
        r = rng.random()
        if r < 0.50:
            # Bare ". " (new sentences) — neutral, eval-rare-connector-free.
            core = ". ".join(parts) + "."
        elif r < 0.75:
            # Ordered multi-step: "First X. Then, Y. Finally, Z." — adds
            # sequence keywords (first/then/next/finally) consistent with
            # eval's ~21% multi-step ratio.
            seq_connectors = [". Then, ", ". Next, ", ". Finally, "]
            joined = ["First " + parts[0][0].lower() + parts[0][1:]]
            for i, p in enumerate(parts[1:]):
                lower_p = p[0].lower() + p[1:]
                # Last part gets "Finally,"; middle parts get "Then,/Next,".
                if i == len(parts) - 2:
                    joined.append(". Finally, " + lower_p)
                else:
                    joined.append(rng.choice(seq_connectors[:2]) + lower_p)
            core = "".join(joined) + "."
        else:
            # Eval-rare ". Also,/Additionally," connectors (kept ≤25% of 3+).
            connector = rng.choice([". Also, ", ". Additionally, "])
            joined = [parts[0]]
            for p in parts[1:]:
                joined.append(p[0].lower() + p[1:])
            core = connector.join(joined) + "."

    # Polite wrapping at ~30%, gated by verb-start (eval baseline 17%; many
    # paraphrases begin with prepositions like "On slide N, ..." which the
    # verb check skips, so realised polite% lands ~17-25%).
    is_single_sentence = len(parts) == 1 or (len(parts) == 2 and ", and " in core)
    first_word = core.split(maxsplit=1)[0].lower().rstrip(",")
    if first_word in _ACTION_VERBS and rng.random() < 0.40:
        first_lower = core[0].lower() + core[1:]
        prefix_pool = _POLITE_PREFIXES_SINGLE if is_single_sentence else _POLITE_PREFIXES_MULTI
        return f"{rng.choice(prefix_pool)} {first_lower}"
    return core


# ---------------------------------------------------------------------------
# Row builder
# ---------------------------------------------------------------------------

def _build_row(
    eval_row: dict,
    ops: list,
    file_path: str,
    expected_path: str,
    perturb_type: str,
    rng: random.Random | None = None,
) -> dict | None:
    if not ops:
        return None

    config_py = _build_py(ops, file_path, file_path, _config_snippet)
    expected_py = _build_py(ops, file_path, expected_path, _expected_snippet)

    # Only inject config step when at least one op modifies initial state
    _no_config_ops = {"insert_table", "reorder_slides", "add_bullet_point"}
    has_config = any(op[0] not in _no_config_ops for op in ops)

    oracle = _standard_oracle(expected_py, file_path, expected_path)
    evaluator = _build_evaluator(ops, file_path, expected_path)
    _rng = rng if rng is not None else random.Random()
    instruction = _build_instruction(ops, _rng)

    modified_eval_row = copy.deepcopy(eval_row)
    cfg = modified_eval_row["metadata"]["config"]
    open_idx = next((i for i, s in enumerate(cfg) if s.get("type") == "open"), len(cfg))
    # NOTE: expected_py was previously
    # only embedded in oracle_actions (the oracle-replay path) — during real
    # agent rollouts oracle never runs, so /tmp/perturb_expected_<short>.pptx
    # was never created. The eval reads it as None → score=0 unconditionally,
    # regardless of whether the agent did the right thing. Mirror calc's
    # `perturb_config_step=_make_config_step(gold_py)` pattern: append the
    # expected-file writer to config so it runs at agent-rollout setup time.
    # Order: config_py (initial-state setup, mutates file_path in place) MUST
    # run before expected_py (reads file_path post-config_py, writes
    # expected_path). Build the list bottom-up: insert expected_py first, then
    # config_py before it — both end up before the original `open` step.
    cfg.insert(open_idx, _make_config_step(expected_py))
    if has_config:
        cfg.insert(open_idx, _make_config_step(config_py))

    knobs = {
        "perturb_type": perturb_type,
        "ops": "|".join(op[0] for op in ops),
        "slide0": ops[0][1],
    }

    return make_perturb_row(
        eval_row=modified_eval_row,
        knob_assignment=knobs,
        new_instruction=instruction,
        new_oracle=oracle,
        new_evaluator=evaluator,
        oracle_after_postconfig=True,
    )


# ---------------------------------------------------------------------------
# Resample helpers
# ---------------------------------------------------------------------------

def _rc(rng: random.Random, exclude=None):
    pool = [c for c in _COLOR_POOL if c != exclude]
    return rng.choice(pool or _COLOR_POOL)


def _rb(rng: random.Random, exclude=None):
    pool = [c for c in _BG_COLOR_POOL if c != exclude]
    return rng.choice(pool or _BG_COLOR_POOL)


def _rs(rng: random.Random, candidates: list, exclude=None):
    pool = [x for x in candidates if x != exclude]
    return rng.choice(pool or candidates)


def _rn(rng: random.Random, exclude=None):
    pool = [f for f in _FONT_POOL if f != exclude]
    return rng.choice(pool or _FONT_POOL)


def _rt(rng: random.Random, exclude=None):
    pool = [t for t in _TITLE_POOL if t != exclude]
    return rng.choice(pool or _TITLE_POOL)


# ---------------------------------------------------------------------------
# TYPE_1 per-task functions  (returns ops list or [])
# ---------------------------------------------------------------------------

def _t1_04578141(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: 3 set_font_color on s1 → (yellow, red, green); exclude per slot so no slot
    # can reproduce its eval color, guaranteeing at least one (slide,color) always differs.
    # Sample WITHOUT replacement so the 3 ops never target the same slide with
    # contradictory colors.
    n = min(3, len(txt))
    slides = rng.sample(txt, n)
    colors = [
        _rc(rng, ("yellow", 255, 255, 0)),
        _rc(rng, ("red",    255, 0,   0)),
        _rc(rng, ("green",  0,   128, 0)),
    ]
    return [("set_font_color", slides[i], (colors[i][1], colors[i][2], colors[i][3])) for i in range(n)]


def _t1_05dd4c1d(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: 3 set_text_alignment; resample 3 slides+alignments. Sample slides
    # WITHOUT replacement so the 3 ops never target the same slide with
    # contradictory alignments (e.g. slide 6 = center AND justify).
    n = min(3, len(txt))
    slides = rng.sample(txt, n)
    aligns = rng.sample(_ALIGN_POOL, min(n, len(_ALIGN_POOL)))
    return [("set_text_alignment", slides[i], aligns[i % len(aligns)]) for i in range(n)]


def _t1_08aced46(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_title_text("Note") + set_text_alignment(right) on s1
    s = rng.choice(txt)
    title = _rt(rng, "Note")
    align = _rs(rng, _ALIGN_POOL, "right")
    return [("set_title_text", s, title), ("set_text_alignment", s, align)]


def _t1_15aece23(rng, info):
    # eval: move title s2 → bottom; resample direction
    directions = ["top", "right"]
    d = rng.choice(directions)
    return [("move_object", 2, d, "title")]


def _t1_21760ecb(rng, info):
    txt = _slide_list(info, "texts") or list(range(1, 25))
    # eval: set_slide_transition(s1, dissolve)
    s = rng.choice(txt)
    tr = _rs(rng, _TRANSITION_POOL, "dissolve")
    return [("set_slide_transition", s, tr)]


def _t1_2b94c692(rng, info):
    # eval: move image s2 → right
    directions = ["top", "bottom", "left"]
    d = rng.choice(directions)
    return [("move_object", 2, d, "picture")]


def _t1_3161d64e(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: 2x set_font_size on s14. Sample WITHOUT replacement so the 2 ops
    # never target the same slide with contradictory sizes.
    n = min(2, len(txt))
    slides = rng.sample(txt, n)
    pt1 = _rs(rng, _SIZE_POOL, 60)
    pt2 = _rs(rng, _SIZE_POOL, 28)
    pts = [pt1, pt2]
    return [("set_font_size", slides[i], pts[i]) for i in range(n)]


def _t1_358aa0a7(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_name all slides
    s = rng.choice(txt)
    font = _rn(rng)
    return [("set_font_name", s, font)]


def _t1_39be0d19(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: insert_table(5r×2c, s3)
    s = rng.choice(txt)
    rows, cols = rng.choice([t for t in _TABLE_POOL if t != (5, 2)] or _TABLE_POOL)
    return [("insert_table", s, rows, cols)]


def _t1_3b27600c(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_background_color(all → blue)
    s = rng.choice(txt)
    c = _rb(rng, ("blue", 0, 0, 255))
    return [("set_background_color", s, (c[1], c[2], c[3]))]


def _t1_4ed5abd0(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_color(s2,3,5→black) + set_font_style(underline)
    s1, s2 = rng.choice(txt), rng.choice(txt)
    c = _rc(rng, ("black", 0, 0, 0))
    style = _rs(rng, _STYLE_POOL, "underline")
    return [("set_font_color", s1, (c[1], c[2], c[3])), ("set_font_style", s2, style)]


def _t1_550ce7e7(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_style(strikethrough)
    s = rng.choice(txt)
    style = _rs(rng, _STYLE_POOL, "strikethrough")
    return [("set_font_style", s, style)]


def _t1_57667013(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_color(s5, yellow)
    s = rng.choice(txt)
    c = _rc(rng, ("yellow", 255, 255, 0))
    return [("set_font_color", s, (c[1], c[2], c[3]))]


def _t1_5c1a6c3d(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_style(bold,s1) + set_font_size(44pt,s1) + set_font_style(underline,s1)
    # The two set_font_style ops must NOT land on the same slide with
    # different styles (contradiction). Sample s1 and s3 without replacement.
    if len(txt) >= 2:
        s1, s3 = rng.sample(txt, 2)
    else:
        s1 = s3 = txt[0]
    s2 = rng.choice(txt)
    style1 = _rs(rng, _STYLE_POOL, "bold")
    pt = _rs(rng, _SIZE_POOL, 44)
    style3 = _rs(rng, _STYLE_POOL, "underline")
    return [("set_font_style", s1, style1), ("set_font_size", s2, pt), ("set_font_style", s3, style3)]


def _t1_5cfb9197(rng, info):
    tbl = _slide_list(info, "tables") or [4]
    # eval: edit_table_cell(s4, first row)
    s = rng.choice(tbl)
    # Validation: edit_table_cell relies on `info.tbl` (slides containing
    # tables) but the spec source pptx may have an empty table or rows with
    # values that coincide with the perturb target (any string we pick risks
    # a vacuous-pass). Switched to set_background_color which always
    # mutates the slide property.
    c = _rb(rng, ("orange", 255, 165, 0))
    return [("set_background_color", s, (c[1], c[2], c[3]))]


def _t1_7ae48c60(rng, info):
    pic = _slide_list(info, "pics") or [3, 4, 6]
    # eval: set_picture_size(s3→20cm, s4→30cm, s6→25cm); pics are s3,s4,s6
    # use s4 or s6 to avoid 7 small icons on s3
    safe_slides = [s for s in pic if s in [4, 6]] or pic
    # NOTE: the previous
    # `[rng.choice(safe_slides) for _ in range(2)]` picked WITH replacement —
    # both ops could land on the same slide with different target heights
    # (e.g. "set s6 to 8cm AND set s6 to 20cm") → instruction self-contradicts
    # → agent reports infeasible. Use rng.sample for distinct slides; fall
    # back to a single op if fewer than 2 slides are available.
    if len(safe_slides) >= 2:
        slides = rng.sample(safe_slides, 2)
        h1 = _rs(rng, _PIC_CM_POOL, 30)
        h2 = _rs(rng, _PIC_CM_POOL, 25)
        return [("set_picture_size", slides[0], h1), ("set_picture_size", slides[1], h2)]
    h1 = _rs(rng, _PIC_CM_POOL, 30)
    return [("set_picture_size", safe_slides[0], h1)]


def _t1_7dbc52a6(rng, info):
    txt = _slide_list(info, "texts") or [2]
    # eval: add_speaker_note(=title) + set_font_style(bold)
    s = rng.choice(txt)
    note = rng.choice(_NOTE_POOL)
    style = _rs(rng, _STYLE_POOL, "bold")
    return [("add_speaker_note", s, note), ("set_font_style", s, style)]


def _t1_841b50aa(rng, info):
    # eval: add_speaker_note("APP") + set_background_color(purple)
    note = rng.choice(_NOTE_POOL)
    c = _rb(rng, ("purple", 128, 0, 128))
    return [("add_speaker_note", 1, note), ("set_background_color", 1, (c[1], c[2], c[3]))]


def _t1_8979838c(rng, info):
    # eval: set_background_color(purple) + add_speaker_note(title)
    c = _rb(rng, ("purple", 128, 0, 128))
    note = rng.choice(_NOTE_POOL)
    return [("set_background_color", 1, (c[1], c[2], c[3])), ("add_speaker_note", 1, note)]


def _t1_986fc832(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_style(underline) + set_font_color(incl table)
    s = rng.choice(txt)
    style = _rs(rng, _STYLE_POOL, "underline")
    c = _rc(rng)
    return [("set_font_style", s, style), ("set_font_color", s, (c[1], c[2], c[3]))]


def _t1_9cf05d24(rng, info):
    txt = _slide_list(info, "texts") or [2]
    # eval: set_background_color(s1→green)
    s = rng.choice(txt)
    c = _rb(rng, ("green", 0, 128, 0))
    return [("set_background_color", s, (c[1], c[2], c[3]))]


def _t1_a434992a(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_size(12) + set_font_color(orange) + set_background_color(?)
    s = rng.choice(txt)
    pt = _rs(rng, _SIZE_POOL, 12)
    c_font = _rc(rng, ("orange", 255, 165, 0))
    c_bg = _rb(rng)
    return [
        ("set_font_size", s, pt),
        ("set_font_color", s, (c_font[1], c_font[2], c_font[3])),
        ("set_background_color", s, (c_bg[1], c_bg[2], c_bg[3])),
    ]


def _t1_a53f80cd(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_color(s2-3 title→black) + set_font_style(bold)
    s1, s2 = rng.choice(txt), rng.choice(txt)
    c = _rc(rng, ("black", 0, 0, 0))
    style = _rs(rng, _STYLE_POOL, "bold")
    return [("set_font_color", s1, (c[1], c[2], c[3])), ("set_font_style", s2, style)]


def _t1_ac1b39ff(rng, info):
    # eval: move table s3 → bottom; keep table bottom (only option with evaluator)
    # resample to "right" position instead
    directions = ["top", "right"]
    d = rng.choice(directions)
    return [("move_object", 3, d, "table")]


def _t1_ac9bb6cb(rng, info):
    # eval: set_font_color(footer→red); footer manipulation complex, skip TYPE_1
    return []


def _t1_af2d657a(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_title_text("Happy Family") + set_font_name("Microsoft JhengHei")
    s = rng.choice(txt)
    title = _rt(rng, "Happy Family")
    font = _rn(rng, "Microsoft JhengHei")
    return [("set_title_text", s, title), ("set_font_name", s, font)]


def _t1_b8adbc24(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_title_text(s2→"Online Shopping")
    s = rng.choice(txt)
    title = _rt(rng, "Online Shopping")
    return [("set_title_text", s, title)]


def _t1_e4ef0baf(rng, info):
    pic = _slide_list(info, "pics") or [1]
    txt = _slide_list(info, "texts") or [1]
    # eval: set_picture_size(s3→20cm) + set_font_size(s6→40pt)
    sp = rng.choice(pic)
    st = rng.choice(txt)
    h = _rs(rng, _PIC_CM_POOL, 20)
    pt = _rs(rng, _SIZE_POOL, 40)
    return [("set_picture_size", sp, h), ("set_font_size", st, pt)]


def _t1_ed43c15f(rng, info):
    txt = _slide_list(info, "texts") or [1]
    pic = _slide_list(info, "pics") or [2]
    # eval: move_object(pic s2→top) + set_font_style(underline, s1+s2)
    sp = rng.choice(pic)
    st = rng.choice(txt)
    d = _rs(rng, ["right", "bottom"], "top")
    style = _rs(rng, _STYLE_POOL, "underline")
    return [("move_object", sp, d, "picture"), ("set_font_style", st, style)]


def _t1_edb61b14(rng, info):
    txt = _slide_list(info, "texts") or [1]
    # eval: set_font_name(last slide → "Times New Roman")
    s = rng.choice(txt)
    font = _rn(rng, "Times New Roman")
    return [("set_font_name", s, font)]


def _t1_f23acfd2(rng, info):
    # eval: add_bullet_point(slide 1); resample bullet text
    bullets = ["Key finding", "Action item", "Next steps", "Important point",
               "Review needed", "Follow up", "Open question"]
    text = rng.choice(bullets)
    return [("add_bullet_point", 1, text)]


# ---------------------------------------------------------------------------
# TYPE_3: atomic eval-evaluator perturb (P3-4-impress)
#
# Each TYPE_3 fn returns a list of fully-built rows (not ops), since the
# evaluator structure is task-specific (rules vary per evaluator). Each row
# is constructed independently with its own config_step, oracle, and
# evaluator that exactly mirrors one of 8 impress-domain eval evaluators.
# ---------------------------------------------------------------------------

# LO user registry path inside the OSWorld VM. Both
# `check_presenter_console_disable` and `check_auto_saving_time` read this
# file and grep for `EnablePresenterScreen` / `AutoSaveTimeIntervall`.
_LO_REGISTRY_PATH = "/home/user/.config/libreoffice/4/user/registrymodifications.xcu"


def _registry_xcu_template(items_xml: str) -> str:
    """Build a minimal valid registrymodifications.xcu wrapping `items_xml`.

    The evaluators only need the matching `<item oor:path=...>` element to
    parse correctly; surrounding namespaces stay constant. `items_xml` must
    be a sequence of `<item ...>...</item>` elements rendered as a string.
    """
    return (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<oor:items xmlns:oor="http://openoffice.org/2001/registry" '
        'xmlns:xs="http://www.w3.org/2001/XMLSchema" '
        'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">\n'
        f'{items_xml}\n'
        '</oor:items>\n'
    )


def _xcu_presenter_item(enabled: bool) -> str:
    val = "true" if enabled else "false"
    return (
        '  <item oor:path="/org.openoffice.Office.Impress/Misc/Start">'
        '<prop oor:name="EnablePresenterScreen" oor:op="fuse">'
        f'<value>{val}</value>'
        '</prop></item>'
    )


def _xcu_autosave_item(minutes: int, enabled: bool = True) -> str:
    en = "true" if enabled else "false"
    return (
        '  <item oor:path="/org.openoffice.Office.Common/Save/Document">'
        '<prop oor:name="AutoSave" oor:op="fuse">'
        f'<value>{en}</value></prop>'
        '<prop oor:name="AutoSaveTimeIntervall" oor:op="fuse">'
        f'<value>{minutes}</value></prop>'
        '</item>'
    )


def _write_xcu_step(items_xml: str) -> dict:
    """Emit an execute step that writes the registry xcu file with `items_xml`."""
    xcu = _registry_xcu_template(items_xml)
    # Write via python heredoc so quoting stays clean (xcu has many quotes).
    py = (
        "import os, pathlib\n"
        f"p = pathlib.Path({_LO_REGISTRY_PATH!r})\n"
        "p.parent.mkdir(parents=True, exist_ok=True)\n"
        f"p.write_text({xcu!r}, encoding='utf-8')\n"
    )
    return {"type": "execute", "parameters": {
        "command": f"python3 << 'PYEOF'\n{py}\nPYEOF", "shell": True,
    }}


# ----- _t3 per-task fns ---------------------------------------------------
# Each fn signature: (rng, info, eval_row, file_path) -> list[dict]
# Returns a list of fully-built perturb rows.
# ----------------------------------------------------------------------------

def _make_t3_row(
    eval_row: dict,
    instruction: str,
    config_steps: list[dict],
    oracle_actions: list[dict],
    evaluator: dict,
    knob_tag: str,
) -> dict:
    """Assemble a TYPE_3 row: inject config_steps before `open`, set evaluator."""
    modified_eval_row = copy.deepcopy(eval_row)
    if config_steps:
        cfg = modified_eval_row["metadata"]["config"]
        # Validation note: when eval has no `open` step (T3 base launches
        # via xdg-open / fresh soffice), the previous fallback `len(cfg)` would
        # APPEND xcu-seeds after LibreOffice already launched — the running
        # process has cached its registry defaults and ignores the post-launch
        # file write (caught on impress 2cd43775: agent's AutoSave changes
        # never persisted because xcu was overwritten by the cached registry).
        # Inject BEFORE the first `open` OR `launch`, whichever comes first.
        anchor_idx = next(
            (i for i, s in enumerate(cfg) if s.get("type") in ("open", "launch")),
            len(cfg),
        )
        for step in reversed(config_steps):
            cfg.insert(anchor_idx, step)
    knobs = {"perturb_type": PERTURB_TYPE_3, "ops": knob_tag}
    return make_perturb_row(
        eval_row=modified_eval_row,
        knob_assignment=knobs,
        new_instruction=instruction,
        new_oracle=oracle_actions,
        new_evaluator=evaluator,
        oracle_after_postconfig=True,
    )


def _t3_0f84bef9(rng, info, eval_row, file_path) -> list[dict]:
    """check_presenter_console_disable — disable presenter console in xcu.

    Initial: write xcu with EnablePresenterScreen=true (eval returns 0).
    Oracle: rewrite xcu with EnablePresenterScreen=false (eval returns 1).
    """
    init_step = _write_xcu_step(_xcu_presenter_item(enabled=True))
    oracle_step = _write_xcu_step(_xcu_presenter_item(enabled=False))
    instr_pool = [
        "Disable the LibreOffice Impress presenter console so it shows the slideshow on a single monitor.",
        "Turn off the presenter view in Impress — I want one screen with just the actual slideshow.",
        "I want to use only one display when presenting; please disable the presenter screen feature in Impress.",
        "Configure Impress so that running a slideshow uses one monitor only (no presenter console).",
        # multi-step (first/then) — sequences the open-options + disable steps
        "First open the LibreOffice Impress options panel, then disable the presenter console so the slideshow runs on a single monitor.",
    ]
    evaluator = {
        "func": "check_presenter_console_disable",
        "result": {"type": "vm_file", "path": _LO_REGISTRY_PATH, "dest": "registrymodifications.xcu"},
        "expected": {},
        "options": {},
    }
    return [_make_t3_row(
        eval_row, rng.choice(instr_pool), [init_step], [oracle_step],
        evaluator, "presenter_console_disable",
    )]


def _t3_2cd43775(rng, info, eval_row, file_path) -> list[dict]:
    """check_auto_saving_time — set LO autosave interval.

    Eval rule: minutes (integer). Resample minutes ∈ {3, 5, 10, 15} excluding
    eval's value=3 (eval baseline says "every 3min"). Build 2 variants.
    """
    rows = []
    for minutes in rng.sample([5, 10, 15], 2):
        # Initial state: a different minutes value (so reward=0).
        init_minutes = 1 if minutes != 1 else 2
        init_step = _write_xcu_step(_xcu_autosave_item(init_minutes))
        oracle_step = _write_xcu_step(_xcu_autosave_item(minutes))
        instr_pool = [
            f"Enable auto-save every {minutes} minutes in LibreOffice Impress so I don't have to keep hitting Ctrl+S.",
            f"Turn on auto-save with a {minutes}-minute interval in LibreOffice settings.",
            f"Configure LibreOffice to auto-save documents every {minutes} minutes.",
            f"Set the autosave interval in LibreOffice to {minutes} min so my work gets saved automatically.",
            # multi-step (first/then) — open settings then enable autosave
            f"First open the LibreOffice general options, then turn on auto-save and set the interval to {minutes} minutes.",
        ]
        evaluator = {
            "func": "check_auto_saving_time",
            "result": {"type": "vm_file", "path": _LO_REGISTRY_PATH, "dest": "registrymodifications.xcu"},
            "expected": {"type": "rule", "rules": {"minutes": minutes}},
            "options": {},
        }
        rows.append(_make_t3_row(
            eval_row, rng.choice(instr_pool), [init_step], [oracle_step],
            evaluator, f"autosave_{minutes}min",
        ))
    return rows


def _t3_3b27600c(rng, info, eval_row, file_path) -> list[dict]:
    """evaluate_presentation_fill_to_rgb_distance — set slide bg via raw RGB.

    Eval rule: rgb=[0,0,255] (blue). Resample to red/green/yellow/purple.
    Two variants. Initial state: gray fill (config_py mutates source pptx).
    """
    targets = [(255, 0, 0, "red"), (0, 128, 0, "green"),
               (255, 255, 0, "yellow"), (128, 0, 128, "purple")]
    chosen = rng.sample(targets, 2)
    rows = []
    for r, g, b, name in chosen:
        # config_py: mutate downloaded pptx so background is gray, not target.
        config_py = (
            "from pptx import Presentation\n"
            "from pptx.dml.color import RGBColor\n"
            f"prs = Presentation('{file_path}')\n"
            "for sld in prs.slides:\n"
            "    bg = sld.background\n"
            "    bg.fill.solid()\n"
            "    bg.fill.fore_color.rgb = RGBColor(128, 128, 128)\n"
            f"prs.save('{file_path}')\n"
        )
        # oracle: rewrite same pptx with target rgb.
        oracle_py = (
            "from pptx import Presentation\n"
            "from pptx.dml.color import RGBColor\n"
            f"prs = Presentation('{file_path}')\n"
            "for sld in prs.slides:\n"
            "    bg = sld.background\n"
            "    bg.fill.solid()\n"
            f"    bg.fill.fore_color.rgb = RGBColor({r}, {g}, {b})\n"
            f"prs.save('{file_path}')\n"
        )
        config_step = _make_config_step(config_py)
        oracle_step = {"type": "execute", "parameters": {
            "command": f"python3 << 'PYEOF'\n{oracle_py}\nPYEOF", "shell": True,
        }}
        instr_pool = [
            f"Please make the background {name} on all my slides.",
            f"Set every slide's background to a solid {name} color.",
            f"Change the background color across all slides in this presentation to {name}.",
            f"I'd like every slide to have a {name} background — apply that to the whole deck.",
            # multi-step (once/then) — open then apply
            f"Once the deck is open, set the slide background to {name} and then apply it to every slide in the presentation.",
        ]
        evaluator = {
            "func": "evaluate_presentation_fill_to_rgb_distance",
            "result": {"type": "vm_file", "path": file_path,
                       "dest": file_path.rsplit("/", 1)[-1]},
            "expected": {"type": "rule", "rules": {
                "rgb": [r, g, b], "original_rgb": [128, 128, 128],
            }},
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }
        rows.append(_make_t3_row(
            eval_row, rng.choice(instr_pool), [config_step], [oracle_step],
            evaluator, f"bg_fill_rgb_{name}",
        ))
    return rows


def _t3_455d3c66(rng, info, eval_row, file_path) -> list[dict]:
    """compare_images — export pptx to png and compare against gold png.

    Strategy: setup uses LibreOffice headless `--convert-to png` to render a
    hidden expected image under /tmp, then the oracle copies that image to the
    Desktop result path. Initial state: leave the result path absent so the
    pre-oracle compare_images score is exactly 0.0 instead of a tiny SSIM
    residual from an unrelated dummy image.
    """
    out_png = "/home/user/Desktop/res.png"
    expected_png = "/tmp/res_gold.png"
    # init: keep the result image absent; compare_images treats that as 0.0.
    init_step = {"type": "execute", "parameters": {
        "command": f"rm -f '{out_png}'",
        "shell": True,
    }}
    # SETUP gold: render pptx during config so a normal rollout has
    # the expected image to compare against (previously res_gold.png was produced
    # ONLY by the oracle → absent in a real rollout → compare_images returns 0).
    # Runs headless BEFORE the first `open`, so no GUI soffice is contended.
    render_gold_cmd = (
        f"tmpd=$(mktemp -d) && "
        f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
        f"--convert-to png --outdir \"$tmpd\" '{file_path}' 2>/dev/null && "
        f"cp \"$tmpd\"/*.png '{expected_png}'; "
        f"rm -rf \"$tmpd\"; true"
    )
    gold_step = {"type": "execute", "parameters": {"command": render_gold_cmd, "shell": True}}
    # oracle now only does the trivial copy of the setup-rendered gold.
    oracle_step = {"type": "execute", "parameters": {
        "command": f"cp '{expected_png}' '{out_png}'", "shell": True}}
    instr_pool = [
        "Could you help me export this Impress file to a PNG image and save it as res.png on the Desktop?",
        "Please export the current presentation to a .png image file named res.png on my Desktop.",
        "Render the slideshow as res.png on the Desktop using the default export settings.",
        "Save an image export of this Impress deck as res.png on the Desktop.",
        # multi-step (first/then) — open the deck, then export to PNG on Desktop
        "First open the presentation, then export it as a PNG image named res.png on the Desktop.",
    ]
    evaluator = {
        "func": "compare_images",
        "result": {"type": "vm_file", "path": out_png, "dest": "res.png"},
        "expected": {"type": "vm_file", "path": expected_png, "dest": "res_gold.png"},
        "options": {},
    }
    return [_make_t3_row(
        eval_row, rng.choice(instr_pool), [init_step, gold_step], [oracle_step],
        evaluator, "export_png",
    )]


def _t3_5d901039(rng, info, eval_row, file_path) -> list[dict]:
    """check_image_stretch_and_center — stretch slide-1 image to slide bounds.

    Eval expected = original pptx (cloud_file). The evaluator compares modified
    against original by matching image blobs and asserting modified image's
    width/height/center within Inches(0.5) of the slide. Oracle: load source,
    enlarge first picture on slide 1 to slide width/height, save.
    """
    expected_path = file_path.replace(".pptx", "_orig.pptx")
    # init: shrink the picture to a tiny size (so it's not centered/stretched).
    config_py = (
        "from pptx import Presentation\n"
        "from pptx.util import Inches\n"
        f"prs = Presentation('{file_path}')\n"
        "for sh in prs.slides[0].shapes:\n"
        "    if sh.shape_type == 13:\n"
        "        sh.left = Inches(0); sh.top = Inches(0)\n"
        "        sh.width = Inches(2); sh.height = Inches(2)\n"
        "        break\n"
        f"prs.save('{file_path}')\n"
    )
    # oracle: copy file to expected_path BEFORE stretch (= original), then
    # stretch in-place on file_path. Evaluator compares post-postconfig
    # `file_path` against `expected_path` (saved as original snapshot).
    # We need the snapshot to be saved BEFORE we mutate the result. Use
    # oracle_after_postconfig=True so oracle runs after agent saves.
    oracle_py = (
        "from pptx import Presentation\n"
        f"prs = Presentation('{file_path}')\n"
        "slide = prs.slides[0]\n"
        "sw = prs.slide_width; sh = prs.slide_height\n"
        "for shape in slide.shapes:\n"
        "    if shape.shape_type == 13:\n"
        "        shape.width = sw; shape.height = sh\n"
        "        shape.left = 0; shape.top = 0\n"
        "        break\n"
        f"prs.save('{file_path}')\n"
    )
    # The init step needs to mutate AFTER the file is downloaded but BEFORE
    # any oracle snapshot. The snapshot in oracle_py copies file_path (already
    # init-mutated) to expected_path so the evaluator sees that as original.
    # Then oracle stretches file_path. Evaluator compares modified pptx vs
    # init-mutated snapshot — both have the same image blob (only size/pos
    # differs), which is exactly what `check_image_stretch_and_center` tests.
    config_step = _make_config_step(config_py)
    # SETUP: snapshot the (shrunk) file as _orig so the evaluator's expected
    # original exists during a normal rollout. Same image blob the checker
    # matches; runs after shrink, before `open`. §C #25 (5d901039).
    snapshot_step = _make_config_step(
        "import shutil\n"
        f"shutil.copy('{file_path}', '{expected_path}')\n"
    )
    oracle_step = {"type": "execute", "parameters": {
        "command": f"python3 << 'PYEOF'\n{oracle_py}\nPYEOF", "shell": True,
    }}
    instr_pool = [
        "Stretch the first image on slide 1 to fill the entire slide while keeping it centered.",
        "Resize the picture on the first slide so it covers the whole slide and is centered on the page.",
        "Make the image on slide 1 a full-page cover image — stretch it to the slide bounds and center it.",
        "I'd like the picture on the first slide to fill the slide — please stretch it to the page size and center it.",
        # multi-step (first/then) — stretch then center
        "First stretch the picture on slide 1 to the full slide width and height, then center it on the page so it sits flush with the slide edges.",
    ]
    evaluator = {
        "func": "check_image_stretch_and_center",
        "result": {"type": "vm_file", "path": file_path, "dest": file_path.rsplit("/", 1)[-1]},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected_original.pptx"},
        "options": {},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return [_make_t3_row(
        eval_row, rng.choice(instr_pool), [config_step, snapshot_step], [oracle_step],
        evaluator, "image_stretch_center",
    )]


def _t3_ac9bb6cb(rng, info, eval_row, file_path) -> list[dict]:
    """check_page_number_colors — set slide-master page-number color.

    Eval rule: color ∈ {red, blue, green, black}. Resample 2 variants. Oracle:
    edit slideMaster1.xml to set the slide-number placeholder srgbClr to a
    color matching the rule's classifier (e.g. red threshold). Init: set to
    a different color so reward=0.
    """
    # Hex values that pass the evaluator's red/blue/green/black classifier.
    color_hex = {"red": "FF0000", "blue": "0000FF", "green": "00B050", "black": "000000"}
    init_hex = "808080"  # gray — passes none of the classifiers
    rows = []
    for color in rng.sample(["red", "blue", "green"], 2):
        target_hex = color_hex[color]
        # config_py: rewrite slideMaster1.xml so srgbClr val=init_hex (initial).
        # oracle_py: rewrite slideMaster1.xml so srgbClr val=target_hex.
        # Use a shared helper that rebuilds the zip with the patched master.
        config_py = _build_master_color_patch(file_path, init_hex)
        oracle_py = _build_master_color_patch(file_path, target_hex)
        instr_pool = [
            f"The slide number is barely visible — change its color to {color} so it stands out.",
            f"Make the page number on every slide {color} so I can read it more easily.",
            f"Could you change the slide-number color to {color}?",
            f"Please update the page number color across all slides to {color}.",
            # multi-step (first/then) — open master then recolor
            f"First open the slide master, then change the page number color to {color} so every slide picks up the new color.",
        ]
        evaluator = {
            "func": "check_page_number_colors",
            "result": {"type": "vm_file", "path": file_path,
                       "dest": file_path.rsplit("/", 1)[-1]},
            "expected": {"type": "rule", "rules": {"color": color}},
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }
        config_step = {"type": "execute", "parameters": {
            "command": f"python3 << 'PYEOF'\n{config_py}\nPYEOF", "shell": True,
        }}
        oracle_step = {"type": "execute", "parameters": {
            "command": f"python3 << 'PYEOF'\n{oracle_py}\nPYEOF", "shell": True,
        }}
        rows.append(_make_t3_row(
            eval_row, rng.choice(instr_pool), [config_step], [oracle_step],
            evaluator, f"page_number_{color}",
        ))
    return rows


def _build_master_color_patch(pptx_path: str, hex_color: str) -> str:
    """Return python code that patches slideMaster1.xml's first srgbClr val.

    The evaluator searches sldNum/ftr/dt placeholders first, then text-run
    rPr solidFill, then any non-black srgbClr. To target the most common
    code path, we set ALL non-black srgbClr in slideMaster1 to `hex_color`.
    """
    return (
        "import zipfile, shutil, os, re\n"
        f"src = '{pptx_path}'\n"
        f"tmpd = '/tmp/_master_patch_{hex_color}'\n"
        "if os.path.exists(tmpd):\n"
        "    shutil.rmtree(tmpd)\n"
        "os.makedirs(tmpd)\n"
        "with zipfile.ZipFile(src, 'r') as z:\n"
        "    z.extractall(tmpd)\n"
        "mp = os.path.join(tmpd, 'ppt/slideMasters/slideMaster1.xml')\n"
        "if os.path.exists(mp):\n"
        "    with open(mp, 'r', encoding='utf-8') as f:\n"
        "        xml = f.read()\n"
        f"    xml = re.sub(r'<a:srgbClr val=\"[0-9A-Fa-f]{{6}}\"', '<a:srgbClr val=\"{hex_color}\"', xml)\n"
        "    with open(mp, 'w', encoding='utf-8') as f:\n"
        "        f.write(xml)\n"
        "os.remove(src)\n"
        "with zipfile.ZipFile(src, 'w', zipfile.ZIP_DEFLATED) as zf:\n"
        "    for root, _, files in os.walk(tmpd):\n"
        "        for fn in files:\n"
        "            fp = os.path.join(root, fn)\n"
        "            zf.write(fp, os.path.relpath(fp, tmpd))\n"
        "shutil.rmtree(tmpd)\n"
    )


def _t3_c59742c0(rng, info, eval_row, file_path) -> list[dict]:
    """compare_audios — embed Baseball.mp3 into pptx slide1.

    Eval downloads both pptx and Baseball.mp3 to Desktop. The evaluator
    extracts the audio embedded in slide-1 of pptx and compares it to the
    Desktop mp3. Oracle: extract pptx, copy mp3 into ppt/media/, register
    the relationship, repackage. Mirror synth `audio_insert` oracle.
    """
    audio_name = "Baseball.mp3"
    mp3_path = "/home/user/Desktop/Baseball.mp3"
    # Initial state: the eval download already provides the mp3 next to the
    # pptx, but the pptx doesn't have it embedded yet. So reward=0 by default
    # — no init mutation needed beyond what eval config provides.
    embed_py = (
        "import zipfile, shutil, os\n"
        "import xml.etree.ElementTree as ET\n"
        f"pptx_path = '{file_path}'\n"
        f"mp3_path = '{mp3_path}'\n"
        f"audio_name = '{audio_name}'\n"
        "tmp_dir = '/tmp/_pptx_audio_embed'\n"
        "if os.path.exists(tmp_dir):\n"
        "    shutil.rmtree(tmp_dir)\n"
        "os.makedirs(tmp_dir)\n"
        "with zipfile.ZipFile(pptx_path, 'r') as z:\n"
        "    z.extractall(tmp_dir)\n"
        "media_dir = os.path.join(tmp_dir, 'ppt/media')\n"
        "os.makedirs(media_dir, exist_ok=True)\n"
        "shutil.copy2(mp3_path, os.path.join(media_dir, audio_name))\n"
        "rels_path = os.path.join(tmp_dir, 'ppt/slides/_rels/slide1.xml.rels')\n"
        "if os.path.exists(rels_path):\n"
        "    tree = ET.parse(rels_path); root = tree.getroot()\n"
        "    max_id = 0\n"
        "    for r in root:\n"
        "        rid = r.get('Id', '')\n"
        "        if rid.startswith('rId'):\n"
        "            try: max_id = max(max_id, int(rid[3:]))\n"
        "            except ValueError: pass\n"
        "    new_rid = f'rId{max_id + 1}'\n"
        "    ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/relationships')\n"
        "    rel = ET.SubElement(root, 'Relationship')\n"
        "    rel.set('Id', new_rid)\n"
        "    rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio')\n"
        "    rel.set('Target', f'../media/{audio_name}')\n"
        "    tree.write(rels_path, xml_declaration=True, encoding='UTF-8')\n"
        "ct_path = os.path.join(tmp_dir, '[Content_Types].xml')\n"
        "ct_tree = ET.parse(ct_path); ct_root = ct_tree.getroot()\n"
        "ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/content-types')\n"
        "if not any(e.get('Extension') == 'mp3' for e in ct_root):\n"
        "    e = ET.SubElement(ct_root, 'Default')\n"
        "    e.set('Extension', 'mp3'); e.set('ContentType', 'audio/mpeg')\n"
        "ct_tree.write(ct_path, xml_declaration=True, encoding='UTF-8')\n"
        "os.remove(pptx_path)\n"
        "with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:\n"
        "    for root_dir, _, files in os.walk(tmp_dir):\n"
        "        for f in files:\n"
        "            fp = os.path.join(root_dir, f)\n"
        "            zf.write(fp, os.path.relpath(fp, tmp_dir))\n"
        "shutil.rmtree(tmp_dir)\n"
    )
    oracle_step = {"type": "execute", "parameters": {
        "command": f"python3 << 'PYEOF'\n{embed_py}\nPYEOF", "shell": True,
    }}
    instr_pool = [
        f'Add the audio file "{audio_name}" from the Desktop into the first slide of this presentation.',
        f"Insert {audio_name} (it's on my Desktop) into slide 1 of this Impress deck.",
        f'Embed the "{audio_name}" file into the first slide so it plays during the presentation.',
        f"I'd like to add the {audio_name} audio from the Desktop into the opening slide.",
        # multi-step (first/then) — locate then embed
        f'First locate "{audio_name}" on the Desktop, then embed it into slide 1 of this presentation so it plays once the slide is shown.',
    ]
    evaluator = {
        "func": "compare_audios",
        "result": {
            "type": "audio_in_slide",
            "ppt_file_path": file_path,
            "slide_index": 0,
            "dest": audio_name,
        },
        "expected": {"type": "vm_file", "path": mp3_path, "dest": "Baseball_to_be_placed.mp3"},
        "options": {},
    }
    return [_make_t3_row(
        eval_row, rng.choice(instr_pool), [], [oracle_step],
        evaluator, "audio_embed_slide1",
    )]


def _t3_ce88f674(rng, info, eval_row, file_path) -> list[dict]:
    """check_slide_orientation_Portrait — flip slide width<height.

    Init: ensure landscape (width>=height — typically the source default).
    Oracle: swap slide_width/slide_height so width<height.
    """
    config_py = (
        "from pptx import Presentation\n"
        f"prs = Presentation('{file_path}')\n"
        "# force landscape\n"
        "if prs.slide_width < prs.slide_height:\n"
        "    prs.slide_width, prs.slide_height = prs.slide_height, prs.slide_width\n"
        f"prs.save('{file_path}')\n"
    )
    oracle_py = (
        "from pptx import Presentation\n"
        f"prs = Presentation('{file_path}')\n"
        "# force portrait\n"
        "if prs.slide_width >= prs.slide_height:\n"
        "    prs.slide_width, prs.slide_height = prs.slide_height, prs.slide_width\n"
        f"prs.save('{file_path}')\n"
    )
    config_step = _make_config_step(config_py)
    oracle_step = {"type": "execute", "parameters": {
        "command": f"python3 << 'PYEOF'\n{oracle_py}\nPYEOF", "shell": True,
    }}
    instr_pool = [
        "Set the slides to portrait orientation instead of landscape.",
        "Switch the slide layout from landscape to portrait orientation.",
        "Please change the page setup to portrait so the slides are taller than they are wide.",
        "Rotate the slide orientation to portrait — I want them upright, not sideways.",
        # multi-step (first/then) — open page setup then flip
        "First open the slide properties dialog, then switch the orientation from landscape to portrait so the slides stand taller than they are wide.",
    ]
    evaluator = {
        "func": "check_slide_orientation_Portrait",
        "result": {"type": "vm_file", "path": file_path,
                   "dest": file_path.rsplit("/", 1)[-1]},
        "expected": {},
        "options": {},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return [_make_t3_row(
        eval_row, rng.choice(instr_pool), [config_step], [oracle_step],
        evaluator, "orientation_portrait",
    )]


_TYPE3_FNS: dict[str, callable] = {
    "0f84bef9": _t3_0f84bef9,
    "2cd43775": _t3_2cd43775,
    "3b27600c": _t3_3b27600c,
    "455d3c66": _t3_455d3c66,
    "5d901039": _t3_5d901039,
    "ac9bb6cb": _t3_ac9bb6cb,
    "c59742c0": _t3_c59742c0,
    "ce88f674": _t3_ce88f674,
}


_TYPE1_FNS: dict[str, callable] = {
    "04578141": _t1_04578141,
    "05dd4c1d": _t1_05dd4c1d,
    "08aced46": _t1_08aced46,
    "15aece23": _t1_15aece23,
    "21760ecb": _t1_21760ecb,
    "2b94c692": _t1_2b94c692,
    "3161d64e": _t1_3161d64e,
    "358aa0a7": _t1_358aa0a7,
    "39be0d19": _t1_39be0d19,
    "3b27600c": _t1_3b27600c,
    "4ed5abd0": _t1_4ed5abd0,
    "550ce7e7": _t1_550ce7e7,
    "57667013": _t1_57667013,
    "5c1a6c3d": _t1_5c1a6c3d,
    "5cfb9197": _t1_5cfb9197,
    "7ae48c60": _t1_7ae48c60,
    "7dbc52a6": _t1_7dbc52a6,
    "841b50aa": _t1_841b50aa,
    "8979838c": _t1_8979838c,
    "986fc832": _t1_986fc832,
    "9cf05d24": _t1_9cf05d24,
    "9ec204e4": lambda rng, info: [],  # duplicate_slides: too complex
    "a434992a": _t1_a434992a,
    "a53f80cd": _t1_a53f80cd,
    "a669ef01": lambda rng, info: [],  # indent_adjust: not in 14-op pool
    "ac1b39ff": _t1_ac1b39ff,
    "ac9bb6cb": _t1_ac9bb6cb,
    "af2d657a": _t1_af2d657a,
    "b8adbc24": _t1_b8adbc24,
    "e4ef0baf": _t1_e4ef0baf,
    "ed43c15f": _t1_ed43c15f,
    "edb61b14": _t1_edb61b14,
    "f23acfd2": _t1_f23acfd2,
    "73c99fb9": lambda rng, info: [],  # edit_textbox_content: not in pool
}


# ---------------------------------------------------------------------------
# TYPE_2 per-task variant lists
# Note: slide indices are 1-based
# ---------------------------------------------------------------------------

def _t2_variants(short_id: str, rng: random.Random) -> list[list[tuple]]:
    note = rng.choice(_NOTE_POOL)
    c_red = (255, 0, 0)
    c_blue = (0, 0, 255)
    c_green = (0, 128, 0)
    c_yellow = (255, 255, 0)
    c_purple = (128, 0, 128)
    c_orange = (255, 165, 0)

    _V: dict[str, list[list]] = {
        "04578141": [
            [("edit_table_cell", 18, "first", ["Col A", "Col B"])],
            [("reorder_slides", 11, 7)],
            [("add_speaker_note", 3, note), ("set_font_name", 5, "Verdana")],
            # D2: replaced set_slide_transition+set_font_size with text-only ops.
            [("set_font_size", 4, 28), ("set_font_color", 5, c_purple)],
        ],
        "05dd4c1d": [
            # txt slides only on s2/s3/s6: keep font ops on those; bg/pic/note
            # are slide-property ops and run on any slide.
            [("set_picture_size", 2, 15), ("add_speaker_note", 6, note)],
            [("set_font_name", 3, "Verdana"), ("set_picture_size", 5, 12)],
            # D2: replaced set_slide_transition with set_font_style.
            [("set_font_style", 2, "italic"), ("set_font_style", 6, "underline")],
        ],
        "08aced46": [
            [("set_font_color", 1, c_red)],
            [("set_background_color", 2, c_blue), ("set_font_style", 1, "underline")],
            [("set_font_name", 1, "Arial"), ("set_font_size", 1, 28)],
        ],
        "15aece23": [
            [("set_font_size", 2, 24)],
            [("set_background_color", 1, c_purple)],
            [("set_font_color", 2, c_blue), ("set_font_size", 2, 20)],
        ],
        "21760ecb": [
            [("set_font_style", 3, "bold"), ("set_font_color", 5, c_blue)],
            [("reorder_slides", 11, 15)],
            [("insert_table", 2, 3, 2), ("set_font_name", 4, "Verdana")],
            [("set_font_color", 6, c_red), ("set_font_size", 8, 20)],
        ],
        "2b94c692": [
            [("set_font_name", 2, "Times New Roman")],
            [("set_text_alignment", 2, "center")],
            [("set_background_color", 1, c_blue), ("set_picture_size", 2, 12)],
        ],
        "3161d64e": [
            [("reorder_slides", 7, 4), ("set_font_style", 5, "italic")],
            # Validation: slide 18 has no table in source pptx → edit_table_cell
            # was a no-op → expected.pptx == source → trivial_pass. Switched to
            # guaranteed-mutating set_background_color on slide 1.
            [("set_background_color", 1, (255, 165, 0))],
            # D2: replaced set_slide_transition with set_font_color.
            [("add_speaker_note", 3, note), ("set_background_color", 7, c_green)],
            [("set_picture_size", 10, 12), ("set_font_name", 14, "Times New Roman")],
        ],
        "358aa0a7": [
            # D2: replaced set_slide_transition with set_font_style.
            [("set_font_style", 5, "italic"), ("set_font_color", 10, c_red)],
            [("reorder_slides", 14, 10)],
            [("add_speaker_note", 3, note), ("set_font_size", 7, 20)],
            [("set_picture_size", 6, 15), ("set_font_color", 12, c_red)],
        ],
        "39be0d19": [
            [("set_title_text", 1, "Overview"), ("set_background_color", 2, c_blue)],
            # D2: removed standalone set_slide_transition.
            [("set_font_style", 3, "bold")],
            [("set_font_style", 2, "underline"), ("set_font_color", 3, c_red)],
        ],
        "3b27600c": [
            # D2: replaced set_slide_transition with set_text_alignment.
            [("set_font_size", 3, 22), ("set_font_size", 7, 24)],
            [("reorder_slides", 6, 10), ("set_font_color", 5, c_blue)],
            [("insert_table", 2, 4, 3), ("set_font_name", 4, "Verdana")],
            [("set_font_style", 10, "italic"), ("set_font_style", 12, "underline")],
        ],
        "4ed5abd0": [
            [("set_title_text", 6, "Summary"), ("set_font_size", 4, 20)],
            # D2: replaced set_slide_transition with set_text_alignment.
            [("set_font_color", 2, c_blue), ("set_background_color", 5, c_yellow)],
            [("set_picture_size", 3, 12), ("set_font_name", 1, "Arial")],
        ],
        "550ce7e7": [
            [("set_font_name", 4, "Times New Roman"), ("insert_table", 6, 3, 2)],
            [("set_picture_size", 3, 15), ("add_speaker_note", 5, note)],
            # D2: replaced set_slide_transition with set_font_style.
            [("set_font_style", 2, "bold"), ("set_font_color", 7, c_red)],
            [("add_speaker_note", 1, note), ("set_font_size", 4, 24)],
        ],
        "57667013": [
            # txt slides only on s3/s5: keep font ops on those.
            [("set_text_alignment", 3, "right"), ("set_picture_size", 2, 12)],
            [("set_font_style", 3, "underline"), ("set_font_size", 5, 20)],
            [("set_font_name", 3, "Verdana"), ("set_background_color", 1, c_green)],
        ],
        "5c1a6c3d": [
            # txt on s1/s2/s3/s5 (not s4); pics on s1-s5. Keep font ops off s4.
            [("set_font_style", 3, "bold"), ("set_font_color", 5, c_blue)],
            [("set_picture_size", 2, 15), ("set_background_color", 5, c_yellow)],
            [("set_title_text", 3, "Results"), ("set_font_name", 5, "Arial")],
        ],
        "5cfb9197": [
            [("set_font_name", 3, "Verdana"), ("set_font_color", 5, c_red)],
            [("set_font_size", 2, 22), ("set_font_size", 3, 22)],
            [("insert_table", 1, 3, 2), ("set_font_color", 2, c_blue)],
        ],
        "73c99fb9": [
            [("set_font_color", 1, c_red)],
            [("set_font_color", 1, c_blue), ("set_background_color", 2, c_green)],
            [("set_font_style", 1, "bold")],
        ],
        "7ae48c60": [
            [("set_font_size", 2, 24), ("set_font_color", 5, c_red)],
            [("add_speaker_note", 4, note), ("set_font_name", 1, "Verdana")],
            # D2: replaced set_slide_transition with set_font_size.
            [("set_font_size", 3, 24), ("set_background_color", 5, c_purple)],
        ],
        "7dbc52a6": [
            [("set_font_color", 2, c_red), ("set_font_size", 2, 24)],
            [("set_slide_transition", 1, "fade")],
        ],
        # 841b50aa: s1 has only a table (no text shapes) → set_font_color silently
        # no-ops. Use set_background_color (slide-property) and edit_table_cell
        # which always apply regardless of shape content. D2: removed standalone
        # set_slide_transition.
        "841b50aa": [
            [("set_background_color", 1, c_blue)],
            # Validation (2026-05-16): slide 1 has no table (edit_table_cell
            # no-op) AND probably no editable text shapes (set_font_color
            # also no-op → trivial). Switched to add_speaker_note which
            # creates a new note element via python-pptx — guaranteed
            # mutation independent of source slide structure.
            [("add_speaker_note", 1, "Action item: review next quarter.")],
        ],
        "8979838c": [
            [("set_font_style", 1, "bold")],
            [("set_font_name", 1, "Verdana")],
            [("set_font_color", 1, c_blue)],
        ],
        "986fc832": [
            [("set_font_color", 1, c_red)],
            # Validation: slide 1 has no table → edit_table_cell no-op →
            # trivial. Switched to set_background_color (slide property).
            [("set_background_color", 1, (0, 128, 255))],
            [("set_font_style", 1, "underline")],
        ],
        "9cf05d24": [
            [("set_font_size", 2, 20), ("set_picture_size", 2, 12)],
            [("edit_table_cell", 2, "first", ["Item 1", "Item 2"]), ("set_font_color", 2, c_blue)],
            [("set_font_size", 2, 22)],
        ],
        "9ec204e4": [
            [("reorder_slides", 9, 13), ("set_font_color", 3, c_red)],
            # D2: replaced set_slide_transition with set_font_style.
            [("set_font_style", 5, "underline"), ("set_background_color", 10, c_blue)],
            [("set_font_style", 7, "underline"), ("set_font_name", 12, "Verdana")],
            [("set_picture_size", 5, 15), ("set_font_size", 3, 20)],
        ],
        "a434992a": [
            [("set_font_style", 1, "underline")],
            [("set_font_size", 1, 28)],
        ],
        "a53f80cd": [
            [("set_picture_size", 4, 15), ("set_font_size", 2, 20)],
            # D2: replaced set_slide_transition with set_font_color.
            [("set_font_color", 5, c_purple), ("set_font_name", 2, "Times New Roman")],
            [("set_font_style", 3, "bold"), ("set_font_color", 4, c_red)],
        ],
        "a669ef01": [
            [("insert_table", 4, 4, 3), ("set_font_color", 6, c_red)],
            # D2: replaced set_slide_transition with set_text_alignment.
            [("set_font_size", 2, 24), ("set_background_color", 7, c_yellow)],
            [("set_picture_size", 3, 12), ("set_font_size", 5, 24)],
        ],
        "ac1b39ff": [
            # txt only on s2/s3: keep alignment/font ops on those slides.
            # D2: removed standalone set_slide_transition.
            [("set_font_size", 2, 24)],
            [("edit_table_cell", 3, "first", ["R1", "R2", "R3"]), ("set_font_color", 2, c_blue)],
            [("reorder_slides", 3, 2), ("set_background_color", 2, c_red)],
        ],
        "ac9bb6cb": [
            [("set_picture_size", 5, 15), ("set_font_size", 3, 20)],
            [("set_font_color", 7, c_red), ("set_font_color", 10, c_red)],
            [("reorder_slides", 8, 4), ("set_font_name", 8, "Verdana")],
            [("insert_table", 2, 3, 2), ("set_font_style", 4, "underline")],
        ],
        "af2d657a": [
            # All slides have empty title placeholders (no text-frame runs):
            # font ops on standalone slides silently no-op → trivial_pass risk.
            # Use slide-property ops (bg/note/insert_table) only.
            [("insert_table", 3, 3, 2)],
            [("set_background_color", 2, c_blue), ("add_speaker_note", 4, note)],
            # D2: removed standalone set_slide_transition.
            # Replaced standalone set_font_color (no text-frames anywhere) with
            # set_background_color which is a slide-property op.
            [("set_background_color", 1, c_orange)],
        ],
        "b8adbc24": [
            [("edit_table_cell", 1, "first", ["A", "B", "C"]), ("set_font_style", 1, "underline")],
            [("set_font_color", 1, c_red), ("set_picture_size", 1, 12)],
        ],
        "e4ef0baf": [
            [("set_font_style", 3, "underline"), ("set_background_color", 5, c_purple)],
            # D2: replaced set_slide_transition with set_text_alignment.
            [("set_font_size", 1, 22), ("set_font_name", 6, "Times New Roman")],
            [("set_font_style", 4, "underline"), ("set_font_color", 6, c_blue)],
            [("set_font_color", 4, c_red), ("set_font_size", 6, 24)],
        ],
        "ed43c15f": [
            [("edit_table_cell", 4, "first", ["T1", "T2"]), ("set_font_color", 3, c_blue)],
            [("reorder_slides", 6, 2), ("set_background_color", 5, c_green)],
            [("set_font_style", 6, "bold"), ("set_font_size", 6, 20)],
        ],
        "edb61b14": [
            # D2: replaced set_slide_transition with set_text_alignment.
            [("set_title_text", 3, "Summary"), ("set_font_size", 5, 24)],
            [("edit_table_cell", 4, "first", ["V1", "V2", "V3"]), ("set_font_color", 2, c_red)],
            [("set_picture_size", 2, 12), ("set_font_color", 5, c_red)],
        ],
        "f23acfd2": [
            [("set_font_color", 1, c_red)],
            [("set_font_color", 1, c_red)],
            [("set_font_style", 1, "bold")],
        ],
    }
    return _V.get(short_id, [])


# ---------------------------------------------------------------------------
# Main perturb function
# ---------------------------------------------------------------------------

# Tasks legitimately excluded from all perturbation (no file ops or external
# deps). 0f84bef9 was previously here (Presenter console UI toggle) but is now
# perturbable via TYPE_3 (`_t3_0f84bef9`) which writes the registry xcu file
# directly, matching the eval evaluator `check_presenter_console_disable`.
_EXCLUDED_TASKS: frozenset[str] = frozenset()


def perturb_impress_per_task(
    eval_row: dict,
    rng: random.Random,
    perturb_types: tuple[str, ...] = (PERTURB_TYPE_1, PERTURB_TYPE_2, PERTURB_TYPE_3),
    max_type1: int = 1,
    max_type2: int = 2,
    max_type3: int = 2,
) -> list[dict]:
    if eval_row["task_id"] in _EXCLUDED_TASKS:
        return []

    # file_path may be None for tasks that don't download a pptx (e.g. 2cd43775
    # which only launches LO Impress). TYPE_1/TYPE_2 require file_path; TYPE_3
    # tasks that target registry-only evaluators (presenter_console, autosave)
    # can run without it.
    file_path = _get_file_path(eval_row)

    tid = eval_row["task_id"]
    short_id = tid.split("_")[-1]
    info = _info(eval_row)
    expected_path = f"/tmp/perturb_expected_{short_id}.pptx"

    rows: list[dict] = []

    if PERTURB_TYPE_1 in perturb_types and file_path is not None:
        t1_fn = _TYPE1_FNS.get(short_id)
        for _ in range(max_type1):
            if t1_fn is None:
                break
            try:
                ops = t1_fn(rng, info)
            except Exception:
                break
            if not ops:
                break
            row = _build_row(eval_row, ops, file_path, expected_path, PERTURB_TYPE_1, rng)
            if row:
                rows.append(row)

    if PERTURB_TYPE_2 in perturb_types and file_path is not None:
        variants = _t2_variants(short_id, rng)
        if variants:
            sampled = rng.sample(variants, min(max_type2, len(variants)))
            for ops in sampled:
                try:
                    row = _build_row(eval_row, ops, file_path, expected_path, PERTURB_TYPE_2, rng)
                    if row:
                        rows.append(row)
                except Exception:
                    continue

    if PERTURB_TYPE_3 in perturb_types:
        t3_fn = _TYPE3_FNS.get(short_id)
        if t3_fn is not None:
            try:
                t3_rows = t3_fn(rng, info, eval_row, file_path)
            except Exception:
                t3_rows = []
            for r in t3_rows[:max_type3]:
                if r:
                    rows.append(r)

    return rows
