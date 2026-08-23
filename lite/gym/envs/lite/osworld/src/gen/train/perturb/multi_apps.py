"""Multi_apps domain per-task perturbation.

Two archetypes, in dispatch priority:

1. **Tier A1** (chrome → LO sink): genuine multi-app — agent reads source
   content from a local-HTML page in Chrome and transcribes it into the LO
   sink (xlsx/docx). See `_TIER_A1_TASKS` for per-base specs. Used when the
   eval base has both a chrome_open_tabs step and an editable LO file.

2. **Legacy file-op** (xlsx/docx/pptx single-app): same op pool as the
   single-domain perturbs (sort, bold, title change, etc.). Applies to bases
   with an editable LO file but no Tier A1 spec. These rows are single-app
   under a multi_apps label and are queued for replacement by future Tier
   A1/A2 work; they are kept until that work lands.

Tasks lacking both a Tier A1 spec and an editable LO file produce 0 perturb
rows (no rephrase fallback — single-app rephrase rows under multi_apps were
audited as low-value training signal and dropped).

Usage:
    from lite.gym.envs.lite.osworld.src.gen.train.perturb.multi_apps import MULTI_APPS_PERTURB_FNS
"""

from __future__ import annotations

import base64
import copy
import html
import json
import random
import re
import textwrap
from pathlib import Path

from lite.gym.envs.lite.osworld.src.gen.common import (
    LO_SAVE_POSTCONFIG,
)
from lite.gym.envs.lite.osworld.src.gen.train.perturb._utils import (
    make_perturb_row,
)


def _find_project_root() -> Path:
    p = Path(__file__).resolve().parent
    for _ in range(15):
        if (p / ".osworld").exists() or (p / ".git").exists():
            return p
        p = p.parent
    return Path(__file__).resolve().parent


# Tasks that cannot be perturbed due to persistent failures:
# - 51f5801c: PPTX has scheme-colored slide backgrounds → compare_pptx_files crashes on .rgb
# - da922383: import fitz fails
# - 69acbb55: pip install InstructorEmbedding succeeds but module not importable at eval time
# - eb303e01: PPTX contains ink annotations that LibreOffice strips on save → compare_pptx_files fails
# - c7c1e4c3 / b337d106: Google reCAPTCHA blocks search route non-deterministically
_BROKEN_MULTI_TASKS = frozenset({
    # 51f5801c, d1acdb87 removed during validation: the broken-setup reasons applied
    # to the original Tier A1 dispatcher (which reused eval row's config).
    # Our new Tier A4/A6 archetypes (`_strip_eval_config` + own pre-config)
    # don't share that state — see _TIER_A23_TASKS specs for these tids.
    "osworld_multi_apps_da922383",
    "osworld_multi_apps_69acbb55",
    "osworld_multi_apps_eb303e01",
    # c7c1e4c3 removed during validation: original chrome+xlsx state used Google
    # reCAPTCHA. New a4_xlsx_table_to_docx perturb is calc+writer only —
    # no chrome involvement, no reCAPTCHA dependency.
    "osworld_multi_apps_b337d106",
    # f8cfa149 removed during validation: original chrome+xls-search used Google
    # reCAPTCHA. New a4_xlsx_table_to_docx perturb is calc+writer only —
    # agent reads cell_search.xls in Calc and transcribes records to docx.
})

def _make_config_step(py_code: str) -> dict:
    return {"type": "execute", "parameters": {
        "command": f"python3 << 'PYEOF'\n{py_code}\nPYEOF", "shell": True,
    }}



# ---------------------------------------------------------------------------
# Tier A1: Chrome (local HTML) → LO sink (xlsx/docx/pptx)
#
# Genuine multi-app perturbations: each row REQUIRES the agent to read content
# from a Chrome tab (loaded from /tmp/perturb_<short>.html that we write) and
# transcribe it into the LO sink file. Gold is generated deterministically by
# applying the same transcription via openpyxl/python-docx/python-pptx.
#
# Per-task spec (`_TIER_A1_TASKS`) lists the sink path, schema, and 1+ variants
# (each variant = source HTML content + instruction + transcription rows).
# Schema between source HTML and sink file is strictly aligned to remove
# agent ambiguity.
# ---------------------------------------------------------------------------

def _render_source_html(title: str, headers: list[str], rows: list[list]) -> str:
    """Render a tiny self-contained HTML page with one table.

    Constraints:
    - Uses ASCII-only fonts/punctuation (avoids copy-paste pitfalls).
    - Single-screen size: ≤ 8 rows, ≤ 5 cols, no scrolling at default zoom.
    - Header text matches sink-file header text exactly.
    """
    th_cells = "".join(f"<th>{html.escape(str(h))}</th>" for h in headers)
    tr_lines = []
    for row in rows:
        td_cells = "".join(f"<td>{html.escape(str(c))}</td>" for c in row)
        tr_lines.append(f"<tr>{td_cells}</tr>")
    body_rows = "\n      ".join(tr_lines)
    return textwrap.dedent(f"""\
        <!DOCTYPE html>
        <html><head>
        <meta charset="utf-8">
        <title>{html.escape(title)}</title>
        <style>
          body {{ font-family: sans-serif; padding: 32px; font-size: 20px; color: #222; }}
          h2 {{ margin: 0 0 18px 0; }}
          table {{ border-collapse: collapse; }}
          th, td {{ border: 1px solid #888; padding: 10px 18px; text-align: left; }}
          th {{ background: #eee; font-weight: bold; }}
        </style>
        </head><body>
        <h2>{html.escape(title)}</h2>
        <table>
          <tr>{th_cells}</tr>
          {body_rows}
        </table>
        </body></html>
        """)


def _build_html_write_step(html_path: str, html_content: str) -> dict:
    """Pre-config step that writes the source HTML to /tmp/ via base64.

    Base64 is used so embedded quotes/specials survive shell quoting.
    """
    b64 = base64.b64encode(html_content.encode()).decode()
    return {"type": "execute", "parameters": {
        "command": f"echo {b64} | base64 -d > {html_path}",
        "shell": True,
    }}


def _ensure_chrome_with_url(eval_row: dict, new_url: str) -> dict:
    """Deep-copy eval_row and ensure chrome_open_tabs points at new_url.

    Two paths:
    - If chrome_open_tabs already exists: replace its first URL with new_url.
    - Else: inject `google-chrome --remote-debugging-port=1337` launch +
      socat forwarder + chrome_open_tabs steps before the first
      download/launch step. Both chrome and socat are skipped if already
      present elsewhere in the config.

    The injected chrome stack mirrors the standard OSWorld pattern used by
    bases that already use chrome (e.g. 3e3fc409 / 67890eb6).
    """
    er = copy.deepcopy(eval_row)
    cfg = er["metadata"]["config"]

    # Path 1: chrome_open_tabs already there — just replace URL
    for step in cfg:
        if step.get("type") == "chrome_open_tabs":
            params = step.setdefault("parameters", {})
            urls = params.get("urls_to_open", [])
            params["urls_to_open"] = [new_url] + (urls[1:] if urls else [])
            return er

    # Path 2: no chrome_open_tabs — inject the missing chrome-stack pieces.
    # Order: chrome launch → socat launch → chrome_open_tabs. If chrome and/or
    # socat already exist, only insert what's missing, and place
    # chrome_open_tabs immediately AFTER the last existing chrome-stack
    # launch so chrome+socat are running before the tab-open call.
    chrome_idx = -1
    socat_idx = -1
    for i, s in enumerate(cfg):
        if s.get("type") != "launch":
            continue
        cmd = s.get("parameters", {}).get("command", [])
        if not (isinstance(cmd, list) and cmd):
            continue
        if "google-chrome" in cmd[0] or "chromium" in cmd[0]:
            chrome_idx = i
        elif "socat" in cmd[0]:
            socat_idx = i

    new_steps: list[dict] = []
    if chrome_idx == -1:
        new_steps.append({"type": "launch", "parameters": {
            "command": ["google-chrome", "--remote-debugging-port=1337"],
        }})
    if socat_idx == -1:
        new_steps.append({"type": "launch", "parameters": {
            "command": ["socat", "tcp-listen:9222,fork", "tcp:localhost:1337"],
        }})
    new_steps.append({"type": "chrome_open_tabs", "parameters": {
        "urls_to_open": [new_url],
    }})

    last_existing = max(chrome_idx, socat_idx)
    if last_existing >= 0:
        insert_at = last_existing + 1
    else:
        insert_at = next(
            (i for i, s in enumerate(cfg) if s.get("type") in ("download", "launch")),
            0,
        )
    cfg[insert_at:insert_at] = new_steps
    return er


def _drop_download_for_path(eval_row: dict, target_path: str) -> dict:
    """Deep-copy eval_row and remove any `download` step that targets
    `target_path`. Used by sink_starts_empty=True specs where we want the
    pre-config gold-py (which writes a fresh empty workbook to that path)
    to be the file that LO opens, not the upstream populated copy.
    """
    er = copy.deepcopy(eval_row)
    cfg = er["metadata"]["config"]
    new_cfg = []
    for step in cfg:
        if step.get("type") == "download":
            files = step.get("parameters", {}).get("files", [])
            files = [f for f in files if f.get("path") != target_path]
            if not files:
                continue  # drop the download step entirely
            step = copy.deepcopy(step)
            step["parameters"]["files"] = files
        new_cfg.append(step)
    er["metadata"]["config"] = new_cfg
    return er


def _build_xlsx_append_gold_py(
    initial_path: str,
    append_rows: list,
    expected_path: str,
    *,
    prepend_headers: list | None = None,
) -> str:
    """Python heredoc that loads initial xlsx, appends rows, saves to expected_path.

    If `prepend_headers` is given (used for empty-sink specs):
    - First REPLACES the sink xlsx at `initial_path` with a fresh empty
      workbook (the upstream HF-cached xlsx contains unrelated pre-populated
      data despite the "empty" label, which would mislead an agent following
      the "headers in row 1" instruction).
    - Then writes [prepend_headers, *append_rows] to expected_path so the
      gold matches what an agent following the "into the empty xlsx —
      headers in row 1 and data rows below" instruction would produce.

    Without `prepend_headers`, the sink is left as-is (original data
    preserved) and expected = original + appended rows.
    """
    all_rows = ([prepend_headers] + list(append_rows)) if prepend_headers else list(append_rows)
    if prepend_headers is None:
        return textwrap.dedent(f"""\
            import openpyxl
            # data_only=True: bake cached formula values so an existing formula
            # cell (e.g. Email `=LOWER(A2)&"@someuniversity.edu"`) lands in the
            # gold as its COMPUTED string, matching the agent's LO-saved result,
            # not re-emitted as a bare formula that pandas reads NaN (#155
            # multi_apps_f5c13cdd FN). Pareto-safe: formula-free sinks (67890eb6)
            # return identical literals -> no-op; a formula sink is already
            # broken (formula-string vs value) without this.
            wb = openpyxl.load_workbook({initial_path!r}, data_only=True)
            ws = wb.worksheets[0]
            # #155 multi_apps_67890eb6 (GOLD): compact trailing all-blank rows so
            # append() lands right after the last POPULATED row (natural agent
            # behavior), not after styled/blank filler rows that ws.max_row counts.
            while ws.max_row > 1 and all(c.value is None for c in ws[ws.max_row]):
                ws.delete_rows(ws.max_row)
            for row in {all_rows!r}:
                ws.append(row)
            wb.save({expected_path!r})
            """)
    # sink_starts_empty=True: replace sink with empty xlsx, then build expected.
    return textwrap.dedent(f"""\
        import openpyxl, os
        os.makedirs(os.path.dirname({initial_path!r}), exist_ok=True)
        empty_wb = openpyxl.Workbook()
        # Workbook() creates a default 'Sheet' with one empty row context;
        # active sheet has no data — perfect for "empty xlsx" semantics.
        empty_wb.save({initial_path!r})
        wb = openpyxl.load_workbook({initial_path!r})
        ws = wb.worksheets[0]
        for row in {all_rows!r}:
            ws.append(row)
        wb.save({expected_path!r})
        """)


def _build_docx_append_gold_py(initial_path: str, append_paragraphs: list[str], expected_path: str) -> str:
    """Python heredoc that loads initial docx, appends paragraphs, saves."""
    return textwrap.dedent(f"""\
        from docx import Document
        doc = Document({initial_path!r})
        for text in {append_paragraphs!r}:
            doc.add_paragraph(text)
        doc.save({expected_path!r})
        """)


def _build_pptx_append_slides_py(initial_path: str, slide_titles: list[str], expected_path: str) -> str:
    """Python heredoc that loads initial pptx, appends N blank-layout slides
    each containing one centered title text box, saves to expected_path.

    Uses a layout-agnostic approach (manually-positioned text box on a blank
    layout) so the result is deterministic regardless of which slide layouts
    the original presentation defines.
    """
    return textwrap.dedent(f"""\
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation({initial_path!r})
        # Use the last layout (typically Blank) — most reliable across templates.
        blank_layout = prs.slide_layouts[-1]
        for title_text in {slide_titles!r}:
            slide = prs.slides.add_slide(blank_layout)
            tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
            tf = tx.text_frame
            tf.text = title_text
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(36)
        prs.save({expected_path!r})
        """)


# ---------------------------------------------------------------------------
# Tier-A1 instruction pool / stylization
#
# Each variant carries:
#   - html_title: short topic phrase (also rendered as the source HTML <h2>).
#   - rows: data rows transcribed from Chrome to sink.
# The dispatcher derives a 3-paraphrase `instr_pool` per variant via
# `_build_instr_pool`, then picks one with `rng.choice` and applies
# `_stylize_multi_apps_instruction` (polite prefix ~38%, narrative motivation
# pad ~50%, trailing context fluff ~37%) to land in V3 distribution targets:
#
#     polite ratio   30-40% (eval baseline 31%)
#     save mention   0%     (NEVER append "save the file" — saving is a
#                            postconfig step, not part of the instruction)
#     avg_words      35-45  (eval baseline 40.6)
#     multi_sep      ~3%    (multi-step paraphrases include ". Also,/Then,")
#
# The pool is not stored on the variant — generated inline so adding a new
# variant only requires (html_title, rows). Pool entries are semantically
# equivalent (same xlsx-append / docx-append / pptx-slide-append op).
# ---------------------------------------------------------------------------

_A1_POLITE_PREFIXES = [
    "Could you ",
    "Please ",
    "Can you ",
    "I'd like you to ",
    "I'd like to ",
    "I want to ",
    "I need to ",
    "Would you mind helping me ",
]

# Per-extension verb tuples used by paraphrase templates.
_A1_VERBS = {
    "xlsx": [("Append", "to the bottom of"), ("Add", "as new rows at the end of"), ("Copy", "into")],
    "docx": [("Append", "as new paragraphs at the end of"), ("Add", "as new paragraphs to"), ("Paste", "as new paragraphs into")],
    "pptx": [("Append", "as new slides at the end of"), ("Add", "as new slides to"), ("Insert", "as new slides at the bottom of")],
}

# Source-tab phrasings — used to vary the "the Chrome tab shows ..." opener.
_A1_SOURCE_OPENERS = [
    "The Chrome tab shows {topic}",
    "{topic_cap} are listed in the Chrome page",
    "I have {topic} in the open Chrome tab",
]

# Narrative motivation prefixes — give each row a why/context. Two flavors:
#   - "lead-in" (ends with "— "): next sentence continues lowercase.
#   - "stand-alone" (ends with ". "): next sentence starts a new sentence.
# Empty entries skip motivation (~50% no-prefix mix for tonal variety).
_A1_MOTIVATION_LEAD_IN = [
    "I'm putting together a quick {topic_kind} and want this in one place — ",
    "I rely on my {sink_kind} as my main {topic_kind}, so ",
    "I'm tidying up the {sink_kind} before {milestone}, so ",
]
_A1_MOTIVATION_STANDALONE = [
    "These came in today and I want to track them alongside the rest of the {topic_kind}. ",
    "I noticed these still need to be logged before {milestone}. ",
    "I'd like the {sink_kind} to stay current as my single source of truth. ",
]
_A1_MOTIVATION_TEMPLATES = (
    [("lead", t) for t in _A1_MOTIVATION_LEAD_IN]
    + [("stand", t) for t in _A1_MOTIVATION_STANDALONE]
    + [("", "")] * 6  # ~50% skip-motivation
)

# Trailing context fluff — adds 8-15 words of benign reminder. ~35% non-empty.
# Two pools: ext-specific (xlsx says "rows/columns", docx/pptx says
# "paragraphs/slides"). Empty slots = no fluff (~50% mix).
_A1_FLUFF_SUFFIXES_XLSX = [
    " Keep the original column order and don't touch any of the existing rows.",
    " Leave the rest of the file untouched — only the new rows should be added.",
    " Don't reformat anything else; only the new rows should be added.",
    "",
    "",
    "",
    "",
    "",
]
_A1_FLUFF_SUFFIXES_DOCX = [
    " Don't change any of the existing paragraphs — only the new ones should be added.",
    " Leave the rest of the document untouched; only the new paragraphs should be added.",
    " Keep the existing text exactly as it is.",
    "",
    "",
    "",
    "",
    "",
]
_A1_FLUFF_SUFFIXES_PPTX = [
    " Don't modify any of the existing slides — only the new ones should be added.",
    " Leave the rest of the deck untouched; only the new slides should be added.",
    " Keep the existing slide content exactly as it is.",
    "",
    "",
    "",
    "",
    "",
]

# Multi-step second-sentence pad — drives V3 multi_sep ratio (~3% target).
# Most slots are empty so the V3 `\. (Also|Additionally|Then|Next),` regex stays
# near eval baseline.
_A1_MULTI_SEP_PADS = [
    ". Also, double-check that the entries match the source",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
]

# Multi-step keyword pads — second-sentence sequencing using
# then/next/first/once/after/before/finally. Eval has multi-step keyword
# coverage 19.8% (broader than the V3 `Also|Then|Next` regex), perturb sat at
# ~12.9%. These pads bring the keyword ratio up while the V3 multi_sep regex
# stays at baseline because most pads use sequencing words (Once / After /
# Before / Finally) that the V3 regex doesn't match. ~10% non-empty slots.
_A1_MULTI_STEP_KEYWORD_PADS = [
    " Once the new entries are in place, scan the column order one more time so nothing got swapped",
    " After you've copied the rows over, give the file a quick visual pass to confirm the order matches the tab",
    " Before wrapping up, double-check that no existing rows were edited along the way",
    " Finally, glance at the last few entries to confirm they line up with what's shown in the tab",
    " Once the data is in, verify each new row reads exactly the same as in the Chrome page",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
]

# Long narrative pads — full multi-step backstory (50-90 words) that frames the
# whole task as a sequenced workflow ("Tomorrow we have X — first do A, then B,
# finally C"). Used to boost p75 / max word count toward eval distribution
# (eval p75=49, max=101; perturb was p75=43, max=69). Triggers infrequently
# (~6% of rows) so mean-words drift stays ±5pp of baseline.
#
# Each pad takes a leading kwarg context ({sink_kind}, {topic}, {milestone})
# and replaces motivation+fluff (the stylizer skips both when a long pad fires)
# so a single row carries one coherent narrative rather than three stacked
# fluff layers.
_A1_LONG_NARRATIVE_PADS = [
    (
        "Tomorrow morning we have {milestone} and the {topic} reference is "
        "split across two places at the moment — first open up the Chrome tab "
        "I left running, then take {rows_phrase} from it and bring them into "
        "the {sink_kind} so everything lives in one file before the meeting. "
        "Once that's done, the {sink_kind} should be the single source of truth. "
    ),
    (
        "I'm getting the {sink_kind} ready ahead of {milestone}, and the only "
        "missing piece is the {topic} listed in the open Chrome page. First "
        "check that the columns line up, next pull {rows_phrase} from the tab, "
        "and after that drop them into the {sink_kind} in the same order they "
        "appear on screen. Before you finish, give it a quick scan so nothing "
        "got dropped. "
    ),
    (
        "Quick errand before {milestone}: I've staged the {topic} in a Chrome "
        "tab and I want them folded into the {sink_kind} in one pass. First "
        "look at the table headers in the tab, then take {rows_phrase} across "
        "in the same order, and finally leave the existing entries exactly as "
        "they were so the diff is clean. "
    ),
    (
        "Could you help me consolidate the {topic} for {milestone}? I've got "
        "the data in a Chrome tab and the {sink_kind} open side by side. "
        "First, copy the headers across if the file is empty, then bring "
        "{rows_phrase} into the {sink_kind} below them, and once that's "
        "done, take a moment to sanity-check that the column order on screen "
        "matches the file. "
    ),
    (
        "I'm putting together the {sink_kind} for {milestone} and the only "
        "outstanding item is the {topic} in the Chrome window. Please first "
        "skim the tab so you know the layout, next transfer {rows_phrase} "
        "into the {sink_kind} matching the column order one-for-one, and "
        "after that leave the rest of the file alone — the existing entries "
        "should not change at all. "
    ),
    (
        "Before {milestone}, I want the {sink_kind} to be the single place "
        "I check for {topic}. The Chrome tab has the new rows. First make "
        "sure no existing data is touched, then bring {rows_phrase} into the "
        "{sink_kind} in the same column order, and finally close out without "
        "renaming or reformatting anything else. "
    ),
    (
        "Heads-up: we have {milestone} on the calendar and the {topic} I "
        "need are still only in a Chrome tab. Could you first read through "
        "the tab to confirm the schema, then bring {rows_phrase} across into "
        "the {sink_kind} below the existing data, and once that's done leave "
        "every other row untouched so the file diff is minimal? "
    ),
    (
        "Tidying up before {milestone} — I've got the {topic} sitting in a "
        "Chrome tab and the {sink_kind} open and ready. First glance at the "
        "tab to note the column order, next copy {rows_phrase} into the "
        "{sink_kind} keeping that order intact, and after that let the file "
        "sit as-is so I can review it before the meeting. "
    ),
]


def _topic_from_html_title(html_title: str) -> str:
    """Derive a topic phrase from html_title (used inside paraphrases).

    Strategy: lowercase the title, strip parentheticals, drop leading
    quantifier words ("New", "Two more", "Three"). Fall back to the original
    if the result is empty.
    """
    t = html_title.strip()
    # Drop trailing parenthetical e.g. "(June)"
    if "(" in t:
        t = t[: t.index("(")].strip()
    # Drop leading quantifier prefixes
    lowered = t.lower()
    for prefix in (
        "two more ", "three more ", "two ", "three ", "four ", "five ",
        "new ", "additional ", "more ", "older ", "recent ",
    ):
        if lowered.startswith(prefix):
            t = t[len(prefix):]
            break
    return t.strip().lower() or html_title.lower()


def _sink_kind_phrase(sink_ext: str, sink_filename: str) -> str:
    """Short colloquial phrase for the sink file (used in narrative pads)."""
    if sink_ext == "xlsx":
        return f"{sink_filename} spreadsheet"
    if sink_ext == "docx":
        return f"{sink_filename} document"
    if sink_ext == "pptx":
        return f"{sink_filename} deck"
    return sink_filename


def _columns_phrase(headers: list[str]) -> str:
    """Render headers as a comma-separated, instruction-grade column list."""
    return ", ".join(headers)


def _build_instr_pool(spec: dict, variant: dict) -> list[str]:
    """Generate 3 semantically-equivalent paraphrases for a variant.

    All three describe the same xlsx-append / docx-append / pptx-slide-append
    op against the same sink file with the same data rows. They differ only in
    surface form: opener, verb, column-order phrase. None of them mention
    "save" — saving is handled by the eval postconfig.
    """
    sink_path = spec["sink_path"]
    sink_ext = spec["sink_ext"]
    sink_headers = spec["sink_headers"]
    sink_starts_empty = bool(spec.get("sink_starts_empty"))
    sink_filename = sink_path.rsplit("/", 1)[-1]
    n = len(variant["rows"])
    topic = _topic_from_html_title(variant["html_title"])
    cols = _columns_phrase(sink_headers)

    # Count phrasing
    if n == 1:
        count_word = "one"
        rows_phrase = "the row"
        entry_word = "entry"
    elif n == 2:
        count_word = "two"
        rows_phrase = "both rows"
        entry_word = "entries"
    elif n == 3:
        count_word = "three"
        rows_phrase = "all three rows"
        entry_word = "entries"
    else:
        count_word = str(n)
        rows_phrase = f"all {n} rows"
        entry_word = "entries"

    if sink_ext == "xlsx":
        if sink_starts_empty:
            # Empty sink: instruction must say "headers in row 1, data below".
            # All three paraphrases start with an action verb so the polite-
            # prefix injector in `_stylize_multi_apps_instruction` can hit any
            # of them (V3 polite-ratio target depends on this).
            p1 = (
                f"Copy the small table from the Chrome tab — columns {cols} "
                f"plus {count_word} {entry_word} — into the empty "
                f"{sink_filename}, placing the headers in row 1 and the data "
                f"rows directly below them, preserving the column order."
            )
            p2 = (
                f"Transfer the whole table from the open Chrome tab "
                f"(header row {cols} plus {count_word} data "
                f"row{'s' if n != 1 else ''}) into the empty {sink_filename}, "
                f"starting at cell A1 with headers on the first row."
            )
            p3 = (
                f"Take the {count_word} {entry_word} listed in Chrome with "
                f"the column titles {cols} and copy the headers plus those "
                f"{count_word} data row{'s' if n != 1 else ''} into the empty "
                f"{sink_filename} so it ends up with one header row followed "
                f"by the {count_word} data row{'s' if n != 1 else ''}."
            )
            return [p1, p2, p3]
        # Non-empty xlsx — append below existing data. All three start with
        # an action verb (Append / Add / Take) so the polite-prefix injector
        # can apply to any choice.
        p1 = (
            f"Append {rows_phrase} from the Chrome tab ({topic}, "
            f"{count_word} new {entry_word}) to the bottom of "
            f"{sink_filename}, preserving the column order ({cols})."
        )
        p2 = (
            f"Add the {count_word} {entry_word} from the open Chrome page "
            f"covering {topic} as new rows at the end of {sink_filename}, "
            f"keeping the column order ({cols})."
        )
        pronoun = "it" if n == 1 else "them"
        p3 = (
            f"Take the {count_word} {entry_word} from the Chrome tab "
            f"({topic}) and append {pronoun} to {sink_filename}. Use the same "
            f"column order as the existing data: {cols}."
        )
        return [p1, p2, p3]

    if sink_ext == "docx":
        # Docx variants are pure text-append — gold renders rows as paragraphs.
        # Don't mention column order; instead say "preserve text exactly".
        # All three start with an action verb (Append / Add / Take).
        p1 = (
            f"Append the {count_word} {entry_word} from the Chrome tab "
            f"({topic}) as new paragraphs at the end of {sink_filename}, "
            f"preserving the text exactly."
        )
        p2 = (
            f"Add the {count_word} {entry_word} listed in the Chrome page "
            f"about {topic} as new paragraphs at the bottom of "
            f"{sink_filename}, keeping the wording verbatim."
        )
        p3 = (
            f"Take the {count_word} {entry_word} shown in Chrome ({topic}) "
            f"and paste each as a new paragraph at the end of "
            f"{sink_filename}, keeping the text exactly as written in the tab."
        )
        return [p1, p2, p3]

    if sink_ext == "pptx":
        # Pptx variants append N new slides — each slide title from the row.
        # All three start with an action verb (Insert / Append / Take).
        p1 = (
            f"Insert {count_word} new slide{'s' if n != 1 else ''} at the "
            f"end of {sink_filename}, setting each slide's title to the "
            f"corresponding text from the Chrome tab "
            f"(which shows {count_word} title{'s' if n != 1 else ''} on {topic})."
        )
        p2 = (
            f"Append {count_word} new slide{'s' if n != 1 else ''} at the "
            f"bottom of {sink_filename}, each titled with the matching text "
            f"from the {count_word} title{'s' if n != 1 else ''} listed in "
            f"the Chrome page about {topic}."
        )
        p3 = (
            f"Take the {count_word} title{'s' if n != 1 else ''} from the "
            f"Chrome tab ({topic}) and add a new slide for each at the end "
            f"of {sink_filename}, with that text as the slide title."
        )
        return [p1, p2, p3]

    raise ValueError(f"unsupported sink_ext={sink_ext}")


_A1_ACTION_VERBS = frozenset({
    "append", "add", "copy", "take", "insert", "transfer", "paste",
})


def _stylize_multi_apps_instruction(
    instr: str,
    rng: random.Random,
    *,
    spec: dict,
    variant: dict,
) -> str:
    """Stylize a paraphrase: polite-prefix injection + narrative motivation
    pad + trailing context fluff. Targets V3 distribution metrics (see
    section docstring above).

    Defensive: also strips any trailing "save"-mentioning phrase if it leaks
    into a paraphrase, so V3 save-ratio stays at 0%.
    """
    # Defensive save-strip — paraphrases shouldn't contain "save" but guard
    # against future edits leaking it.
    save_re = re.compile(
        r"\s*(save the (file|spreadsheet|document|deck)|save when done|save\.)\s*\.?\s*$",
        re.IGNORECASE,
    )
    while save_re.search(instr):
        instr = save_re.sub(".", instr).rstrip()
        if not instr.endswith("."):
            instr += "."

    sink_ext = spec["sink_ext"]
    sink_filename = spec["sink_path"].rsplit("/", 1)[-1]
    sink_kind = _sink_kind_phrase(sink_ext, sink_filename)
    topic = _topic_from_html_title(variant["html_title"])
    n = len(variant["rows"])
    if n == 1:
        rows_phrase_long = "the new entry"
    elif n == 2:
        rows_phrase_long = "the two new entries"
    elif n == 3:
        rows_phrase_long = "the three new entries"
    else:
        rows_phrase_long = f"the {n} new entries"

    # Long narrative pad (~6% of rows) — replaces motivation + fluff with a
    # full multi-step backstory (50-90 words) so p75 / max word count moves
    # toward eval distribution. Mutually exclusive with the standalone
    # motivation/fluff/multi-step pads below to avoid stacking and to keep
    # mean-words drift inside ±5pp.
    long_narrative_applied = False
    if rng.random() < 0.06:
        template = rng.choice(_A1_LONG_NARRATIVE_PADS)
        long_pad = template.format(
            sink_kind=sink_kind,
            topic=topic,
            rows_phrase=rows_phrase_long,
            milestone=rng.choice([
                "tomorrow's review", "Friday's sync", "the leadership review",
                "the next stakeholder check-in", "the deadline",
                "tomorrow's standup",
            ]),
        )
        # Lead-in style: lowercase first letter of the original paraphrase so
        # the stitched narrative reads like one continuous request.
        if instr and instr[0].isupper():
            instr_body = instr[0].lower() + instr[1:]
        else:
            instr_body = instr
        instr = long_pad + instr_body
        long_narrative_applied = True

    # Polite-prefix injection (~32% target). Apply BEFORE motivation so the
    # polite token can appear at sentence-1 start (V3's `^polite` regex
    # matches position 0). Only when paraphrase starts with an action verb.
    # Skip when long narrative already provides framing.
    already_polite = any(instr.lower().startswith(p.lower().strip()) for p in _A1_POLITE_PREFIXES)
    first_word = instr.split(" ", 1)[0].lower().rstrip(".,") if instr else ""
    starts_with_verb = first_word in _A1_ACTION_VERBS
    polite_applied = False
    if not long_narrative_applied and starts_with_verb and not already_polite and rng.random() < 0.38:
        prefix = rng.choice(_A1_POLITE_PREFIXES)
        instr = prefix + instr[0].lower() + instr[1:]
        polite_applied = True

    # Narrative motivation prefix. ~50% of rows pick a non-empty motivation,
    # but skip if we already applied a polite prefix (would push polite token
    # out of position 0 and break V3 regex match) or a long narrative pad.
    if not polite_applied and not long_narrative_applied:
        topic_kind_choices = ["set of notes", "running record", "tracker", "log"]
        flavor, motivation_template = rng.choice(_A1_MOTIVATION_TEMPLATES)
        if motivation_template:
            motivation = motivation_template.format(
                topic_kind=rng.choice(topic_kind_choices),
                sink_kind=sink_kind,
                milestone=rng.choice(["the deadline", "tomorrow's review", "the next sync", "Friday"]),
            )
            if flavor == "lead":
                # Lead-in ends with "— "; next clause continues lowercase.
                if instr and instr[0].isupper():
                    instr_lc = instr[0].lower() + instr[1:]
                else:
                    instr_lc = instr
                instr = motivation + instr_lc
            else:
                # Stand-alone motivation is its own sentence; keep next
                # sentence capitalized.
                instr = motivation + instr

    # Trailing context fluff (~35% non-empty); ext-aware so the reminder uses
    # the right vocabulary (rows/columns vs paragraphs vs slides).
    # Skip when long narrative already supplies the closing context.
    if not long_narrative_applied:
        fluff_pool = {
            "xlsx": _A1_FLUFF_SUFFIXES_XLSX,
            "docx": _A1_FLUFF_SUFFIXES_DOCX,
            "pptx": _A1_FLUFF_SUFFIXES_PPTX,
        }[sink_ext]
        fluff = rng.choice(fluff_pool)
        if fluff:
            if not instr.rstrip().endswith("."):
                instr = instr.rstrip() + "."
            instr = instr + fluff

    # Multi-step second-sentence pad (~3% non-empty) — V3 `Also|Then|Next`
    # regex driver. Mutually exclusive with long narrative.
    if not long_narrative_applied:
        multi_pad = rng.choice(_A1_MULTI_SEP_PADS)
        if multi_pad:
            if instr.endswith("."):
                instr = instr[:-1]
            instr = instr + multi_pad + "."

    # Multi-step keyword pad (~10% non-empty) — uses sequencing words
    # (Once / After / Before / Finally) that don't trigger V3's narrow
    # multi_sep regex but do show up under the broader multi-step keyword
    # check (then/next/first/once/after/before/finally). Skipped under the
    # long narrative, which already has its own sequencing.
    if not long_narrative_applied:
        kw_pad = rng.choice(_A1_MULTI_STEP_KEYWORD_PADS)
        if kw_pad:
            if not instr.rstrip().endswith("."):
                instr = instr.rstrip() + "."
            instr = instr + kw_pad

    if not instr.endswith("."):
        instr = instr + "."
    return instr


# Tier A1 specs — per-task, hand-written for highest correctness/quality.
# Each entry: sink file metadata + N variants (each = source HTML + topic +
# transcription data). Schema headers MUST match the actual file's header row.
#
# Each variant supplies (html_title, rows) — instruction text is generated
# from these by `_build_instr_pool` (3-paraphrase pool) + dispatcher's
# `_stylize_multi_apps_instruction` (polite prefix / motivation pad / fluff).
# Legacy `instr` field is preserved on existing variants for reference but
# IGNORED by the dispatcher.
#
# All "variants" are independent perturb rows. There is no TYPE_1/TYPE_2 split;
# every row is a multi-app transcription task.
_TIER_A1_TASKS: dict[str, dict] = {
    # 3e3fc409: chrome+IMDB → movies.xlsx (151r×4c)
    "osworld_multi_apps_3e3fc409": {
        "sink_path": "/home/user/Desktop/movies.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["title", "release year", "ratings", "description"],
        "variants": [
            {
                "html_title": "New movies to add",
                "rows": [
                    ["Stalker", 1979, 8.2, "Sci-fi by Andrei Tarkovsky"],
                    ["Persona", 1966, 8.1, "Drama by Ingmar Bergman"],
                    ["Tokyo Story", 1953, 8.2, "Family drama by Yasujiro Ozu"],
                ],
                "instr": (
                    "The Chrome tab shows a small table of three movies. "
                    "Append all three rows to the open movies.xlsx, preserving "
                    "the column order (title, release year, ratings, description). "
                    "Save the file when done."
                ),
            },
            {
                "html_title": "Recent releases to log",
                "rows": [
                    ["Past Lives", 2023, 7.8, "Romance drama by Celine Song"],
                    ["The Holdovers", 2023, 7.9, "Comedy drama by Alexander Payne"],
                ],
                "instr": (
                    "Two recent movies are listed in the Chrome browser. "
                    "Add them as new rows at the bottom of movies.xlsx, keeping "
                    "the same column order (title, release year, ratings, "
                    "description). Save the spreadsheet."
                ),
            },
            {
                "html_title": "International films",
                "rows": [
                    ["Parasite", 2019, 8.6, "Korean dark comedy by Bong Joon-ho"],
                    ["Roma", 2018, 7.7, "Mexican drama by Alfonso Cuaron"],
                    ["Amour", 2012, 7.9, "French drama by Michael Haneke"],
                ],
                "instr": (
                    "Three international films are listed in the Chrome page. "
                    "Add them as new rows at the bottom of movies.xlsx in the "
                    "given column order: title, release year, ratings, description. "
                    "Save."
                ),
            },
            {
                "html_title": "Animated picks",
                "rows": [
                    ["Spirited Away", 2001, 8.6, "Animated fantasy by Hayao Miyazaki"],
                    ["WALL-E", 2008, 8.4, "Pixar science-fiction animation"],
                ],
                "instr": (
                    "Two animated films are shown in the Chrome tab. Append "
                    "both rows to movies.xlsx (title, release year, ratings, "
                    "description). Save."
                ),
            },
            {
                "html_title": "Documentary additions",
                "rows": [
                    ["RBG", 2018, 7.5, "Documentary on Justice Ruth Bader Ginsburg"],
                    ["Won't You Be My Neighbor", 2018, 8.4, "Mister Rogers documentary"],
                    ["Sound of Metal", 2019, 7.7, "Drama on hearing loss"],
                ],
                "instr": (
                    "Three documentary entries are listed in the Chrome page. "
                    "Append all three rows at the bottom of movies.xlsx in "
                    "the order title, release year, ratings, description. Save."
                ),
            },
        ],
    },
    # 3f05f3b9: chrome+Picard → tally_book.xlsx (Service/Month/Amount, 6 rows)
    "osworld_multi_apps_3f05f3b9": {
        "sink_path": "/home/user/Documents/Finance/tally_book.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Service", "Month", "Amount"],
        "variants": [
            {
                "html_title": "New April expenses",
                "rows": [
                    ["Netflix", "April", 15.99],
                    ["Spotify", "April", 9.99],
                    ["Gym Membership", "April", 50.00],
                ],
                "instr": (
                    "The Chrome tab lists three new expense entries. "
                    "Append all three rows to the bottom of tally_book.xlsx "
                    "(columns: Service, Month, Amount). Save the file."
                ),
            },
            {
                "html_title": "May utilities",
                "rows": [
                    ["Electricity", "May", 88.40],
                    ["Internet", "May", 60.00],
                ],
                "instr": (
                    "Two utility bills are shown in the Chrome page. "
                    "Add them to tally_book.xlsx as new rows in the order "
                    "Service, Month, Amount. Save the spreadsheet."
                ),
            },
            {
                "html_title": "Cloud services (June)",
                "rows": [
                    ["GCP Compute", "June", 215.40],
                    ["Cloudflare Workers", "June", 5.00],
                    ["DNS Provider", "June", 12.00],
                ],
                "instr": (
                    "The Chrome tab shows three cloud-service expenses. "
                    "Append all three rows to tally_book.xlsx (columns: "
                    "Service, Month, Amount). Save."
                ),
            },
            {
                "html_title": "Marketing tools (July)",
                "rows": [
                    ["Mailchimp", "July", 29.00],
                    ["HubSpot", "July", 50.00],
                ],
                "instr": (
                    "Two marketing-tool expenses are listed in the Chrome "
                    "page. Add them as new rows at the bottom of "
                    "tally_book.xlsx, in the order Service, Month, Amount. "
                    "Save the file."
                ),
            },
        ],
    },
    # 42d25c08: chrome+GitHub → tally_book.xlsx (same Service/Month/Amount)
    "osworld_multi_apps_42d25c08": {
        "sink_path": "/home/user/Documents/Finance/tally_book.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Service", "Month", "Amount"],
        "variants": [
            {
                "html_title": "Subscription renewals (June)",
                "rows": [
                    ["Cloud Storage", "June", 12.50],
                    ["VPN Service", "June", 7.99],
                    ["News Site", "June", 4.50],
                ],
                "instr": (
                    "Three subscription renewals are listed in the Chrome tab. "
                    "Append all three rows to tally_book.xlsx (Service, Month, "
                    "Amount). Save."
                ),
            },
            {
                "html_title": "Q3 office expenses",
                "rows": [
                    ["Coffee", "July", 22.00],
                    ["Office Supplies", "July", 35.75],
                ],
                "instr": (
                    "Two office expenses are displayed in the Chrome page. "
                    "Add them as new rows to tally_book.xlsx in the order "
                    "Service, Month, Amount. Save the file."
                ),
            },
            {
                "html_title": "Team welfare (August)",
                "rows": [
                    ["Yoga Class", "August", 90.00],
                    ["Team Lunch", "August", 145.50],
                    ["Massage", "August", 60.00],
                ],
                "instr": (
                    "The Chrome tab shows three team-welfare expenses. "
                    "Append all three rows to tally_book.xlsx (columns: "
                    "Service, Month, Amount). Save the file."
                ),
            },
            {
                "html_title": "Insurance premiums (September)",
                "rows": [
                    ["Medical Insurance", "September", 240.00],
                    ["Dental Plan", "September", 45.00],
                ],
                "instr": (
                    "Two insurance premium entries are listed in the Chrome "
                    "page. Add them as new rows at the bottom of "
                    "tally_book.xlsx, preserving the column order Service, "
                    "Month, Amount. Save."
                ),
            },
        ],
    },
    # 67890eb6: chrome+aclanthology → best_awards_acl.xlsx (5r×4c, NB original
    # header has typo 'tile' instead of 'title' — preserve verbatim).
    "osworld_multi_apps_67890eb6": {
        "sink_path": "/home/user/Desktop/best_awards_acl.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["tile", "year", "author list", "PDF link"],
        "variants": [
            {
                "html_title": "Recent ACL best long papers",
                "rows": [
                    [
                        "Do Androids Laugh at Electric Sheep",
                        2023,
                        "Alice Smith, Bob Lee",
                        "https://aclanthology.org/2023.acl-best.pdf",
                    ],
                    [
                        "Cross-Lingual Distillation",
                        2024,
                        "Carol Wang, Dan Park",
                        "https://aclanthology.org/2024.acl-best.pdf",
                    ],
                ],
                "instr": (
                    "The Chrome tab shows two ACL best long paper entries. "
                    "Append both rows to best_awards_acl.xlsx, preserving the "
                    "column order (tile, year, author list, PDF link). Save."
                ),
            },
            {
                "html_title": "ACL 2024 highlight",
                "rows": [
                    [
                        "Reflective Reasoning in LLMs",
                        2024,
                        "Eve Chen, Frank Park",
                        "https://aclanthology.org/2024.acl-highlight.pdf",
                    ],
                ],
                "instr": (
                    "One ACL 2024 paper entry is shown in the Chrome page. "
                    "Add it as a new row at the bottom of best_awards_acl.xlsx "
                    "(columns: tile, year, author list, PDF link). Save."
                ),
            },
            {
                "html_title": "Older ACL entries",
                "rows": [
                    [
                        "Bridging Syntax and Semantics",
                        2019,
                        "G. Kim, H. Lee",
                        "https://aclanthology.org/2019.acl-best.pdf",
                    ],
                    [
                        "Memory-Efficient Transformer Pretraining",
                        2020,
                        "I. Singh, J. Liu",
                        "https://aclanthology.org/2020.acl-best.pdf",
                    ],
                    [
                        "Compositional Generalization Benchmarks",
                        2021,
                        "K. Tan, M. Brown",
                        "https://aclanthology.org/2021.acl-best.pdf",
                    ],
                ],
                "instr": (
                    "The Chrome tab lists three older ACL paper entries "
                    "(2019-2021). Append all three rows to "
                    "best_awards_acl.xlsx in the order tile, year, author "
                    "list, PDF link. Save."
                ),
            },
            {
                "html_title": "Findings track papers",
                "rows": [
                    [
                        "Adaptive Curriculum Learning for NLP",
                        2023,
                        "N. Reyes, P. Watson",
                        "https://aclanthology.org/2023.findings-acl.pdf",
                    ],
                    [
                        "Sparse Attention Variants",
                        2024,
                        "Q. Zhao, R. Diaz",
                        "https://aclanthology.org/2024.findings-acl.pdf",
                    ],
                ],
                "instr": (
                    "Two findings-track ACL papers are listed in the Chrome "
                    "page. Add them as new rows at the bottom of "
                    "best_awards_acl.xlsx, preserving the column order tile, "
                    "year, author list, PDF link. Save."
                ),
            },
            {
                "html_title": "Outstanding paper awards",
                "rows": [
                    [
                        "Pretraining Data Curation at Scale",
                        2024,
                        "S. Anand, T. Liu",
                        "https://aclanthology.org/2024.acl-outstanding.pdf",
                    ],
                    [
                        "Translation Quality Estimation Methods",
                        2023,
                        "U. Vega, V. Park",
                        "https://aclanthology.org/2023.acl-outstanding.pdf",
                    ],
                    [
                        "Robustness in Multilingual Models",
                        2022,
                        "W. Xu, X. Tan",
                        "https://aclanthology.org/2022.acl-outstanding.pdf",
                    ],
                ],
                "instr": (
                    "Three outstanding-paper-award entries are listed in the "
                    "Chrome page. Append all three rows to "
                    "best_awards_acl.xlsx, preserving the column order tile, "
                    "year, author list, PDF link. Save."
                ),
            },
        ],
    },
    # 236833a3: chrome+Huggingface → paper_reading_2024_03_01.docx (9 paragraphs)
    "osworld_multi_apps_236833a3": {
        "sink_path": "/home/user/Desktop/paper_reading_2024_03_01.docx",
        "sink_ext": "docx",
        "sink_headers": ["Paper entry"],
        "variants": [
            {
                "html_title": "Daily papers (additional entries)",
                "rows": [
                    [
                        "MM-Reasoner: A Multi-Modal Reasoning Benchmark "
                        "(arxiv 2403.00010, by A. Lin et al.)"
                    ],
                    [
                        "Sparse Mixture-of-Experts at Scale "
                        "(arxiv 2403.00011, by R. Khosla et al.)"
                    ],
                    [
                        "Long-Context Retrieval with Hierarchical Memory "
                        "(arxiv 2403.00012, by S. Park et al.)"
                    ],
                ],
                "instr": (
                    "The Chrome tab lists three paper entries. Append each "
                    "entry as a new paragraph at the end of paper_reading_2024_03_01.docx, "
                    "preserving the entry text exactly. Save the document."
                ),
            },
            {
                "html_title": "Two more papers to record",
                "rows": [
                    [
                        "Efficient Inference via KV-Cache Compression "
                        "(arxiv 2403.00013, by D. Han et al.)"
                    ],
                    [
                        "Self-Refining Code Generation Agents "
                        "(arxiv 2403.00014, by L. Wei et al.)"
                    ],
                ],
                "instr": (
                    "Two additional paper entries are shown in the Chrome page. "
                    "Add each as a new paragraph at the end of "
                    "paper_reading_2024_03_01.docx. Save the document."
                ),
            },
            {
                "html_title": "Vision-language papers",
                "rows": [
                    [
                        "Multi-Image Instruction Tuning at Scale "
                        "(arxiv 2403.00015, by T. Brooks et al.)"
                    ],
                    [
                        "Open-Vocabulary Visual Grounding "
                        "(arxiv 2403.00016, by U. Vasquez et al.)"
                    ],
                ],
                "instr": (
                    "Two vision-language paper entries are shown in the Chrome "
                    "page. Append each as a new paragraph at the end of "
                    "paper_reading_2024_03_01.docx. Save."
                ),
            },
        ],
    },
    # 5990457f: NO chrome in original config (chrome injected by Tier A1).
    # Sink: researchers.xlsx (Name / All citations / h-index / i10-index /
    # Top Cited Work / PDF Link of Top Cited Work — 5 rows).
    "osworld_multi_apps_5990457f": {
        "sink_path": "/home/user/Desktop/researchers.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": [
            "Name", "All citations", "h-index", "i10-index",
            "Top Cited Work", "PDF Link of Top Cited Work",
        ],
        "variants": [
            {
                "html_title": "Researchers to add to the table",
                "rows": [
                    [
                        "John Doe", 12500, 45, 110,
                        "Sparse Reasoning at Scale",
                        "https://example.org/papers/sparse-reasoning.pdf",
                    ],
                    [
                        "Jane Smith", 8800, 38, 92,
                        "Verifier-Guided Decoding",
                        "https://example.org/papers/verifier-decoding.pdf",
                    ],
                ],
                "instr": (
                    "The Chrome tab shows a small table of two researchers. "
                    "Append both rows to the bottom of researchers.xlsx, "
                    "preserving the column order (Name, All citations, "
                    "h-index, i10-index, Top Cited Work, PDF Link of Top "
                    "Cited Work). Save the file."
                ),
            },
            {
                "html_title": "One more researcher entry",
                "rows": [
                    [
                        "Alice Park", 21300, 60, 145,
                        "Long-Context Retrieval",
                        "https://example.org/papers/long-context.pdf",
                    ],
                ],
                "instr": (
                    "One researcher entry is shown in the Chrome page. "
                    "Add it as a new row at the bottom of researchers.xlsx "
                    "(columns: Name, All citations, h-index, i10-index, "
                    "Top Cited Work, PDF Link of Top Cited Work). Save."
                ),
            },
            {
                "html_title": "Industry researchers",
                "rows": [
                    [
                        "Maya Singh", 15600, 52, 130,
                        "Tool-Augmented Inference",
                        "https://example.org/papers/tool-aug.pdf",
                    ],
                    [
                        "Hiroshi Tanaka", 19200, 58, 142,
                        "Memory-Efficient Pretraining",
                        "https://example.org/papers/memory-efficient.pdf",
                    ],
                    [
                        "Nia Okafor", 7400, 32, 78,
                        "Robustness in Code LLMs",
                        "https://example.org/papers/code-robustness.pdf",
                    ],
                ],
                "instr": (
                    "The Chrome tab shows three industry-researcher entries. "
                    "Append all three rows to researchers.xlsx in the order "
                    "Name, All citations, h-index, i10-index, Top Cited Work, "
                    "PDF Link of Top Cited Work. Save."
                ),
            },
            {
                "html_title": "Two academia entries",
                "rows": [
                    [
                        "Prof. Liang Wei", 33000, 78, 200,
                        "Causal Reasoning in NLP",
                        "https://example.org/papers/causal-nlp.pdf",
                    ],
                    [
                        "Dr. Carla Mendez", 11200, 41, 105,
                        "Adversarial Robustness Survey",
                        "https://example.org/papers/adv-robust.pdf",
                    ],
                ],
                "instr": (
                    "Two academic researcher entries are listed in the Chrome "
                    "page. Add them as new rows at the bottom of "
                    "researchers.xlsx, preserving the column order (Name, "
                    "All citations, h-index, i10-index, Top Cited Work, PDF "
                    "Link of Top Cited Work). Save."
                ),
            },
            {
                "html_title": "Cross-disciplinary entries",
                "rows": [
                    [
                        "Dr. Eva Romanov", 9500, 34, 88,
                        "Computational Social Science",
                        "https://example.org/papers/cs-society.pdf",
                    ],
                    [
                        "Prof. Hassan Ali", 26800, 65, 170,
                        "AI for Sustainability",
                        "https://example.org/papers/ai-sust.pdf",
                    ],
                    [
                        "Dr. Yuki Sato", 5200, 27, 60,
                        "Neuro-Symbolic Reasoning",
                        "https://example.org/papers/neuro-symbolic.pdf",
                    ],
                ],
                "instr": (
                    "Three cross-disciplinary researcher entries are shown in "
                    "the Chrome page. Append all three rows to "
                    "researchers.xlsx in the order Name, All citations, "
                    "h-index, i10-index, Top Cited Work, PDF Link of Top "
                    "Cited Work. Save."
                ),
            },
        ],
    },
    # 415ef462: NO chrome — tally_book.xlsx (Service/Month/Amount, 6 rows)
    "osworld_multi_apps_415ef462": {
        "sink_path": "/home/user/Documents/Finance/tally_book.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Service", "Month", "Amount"],
        "variants": [
            {
                "html_title": "August AWS expenses",
                "rows": [
                    ["AWS EC2", "August", 145.20],
                    ["AWS S3", "August", 22.10],
                    ["AWS Lambda", "August", 6.50],
                ],
                "instr": (
                    "The Chrome tab lists three new AWS expenses. "
                    "Append all three rows to tally_book.xlsx (columns: "
                    "Service, Month, Amount). Save."
                ),
            },
            {
                "html_title": "Software licenses (September)",
                "rows": [
                    ["IDE License", "September", 99.00],
                    ["Database License", "September", 250.00],
                ],
                "instr": (
                    "Two software license entries are shown in the Chrome page. "
                    "Add both as new rows at the bottom of tally_book.xlsx, "
                    "keeping the column order Service, Month, Amount. Save."
                ),
            },
            {
                "html_title": "Database tools (October)",
                "rows": [
                    ["MongoDB Atlas", "October", 75.00],
                    ["Postgres Cloud", "October", 50.00],
                    ["Redis Cluster", "October", 35.00],
                ],
                "instr": (
                    "The Chrome tab shows three database-tool expenses. "
                    "Append all three rows to tally_book.xlsx (Service, Month, "
                    "Amount). Save."
                ),
            },
            {
                "html_title": "Security audits (November)",
                "rows": [
                    ["Pen-Test Engagement", "November", 1800.00],
                    ["Vulnerability Scan", "November", 220.00],
                ],
                "instr": (
                    "Two security audit expenses are shown in the Chrome page. "
                    "Add both as new rows in tally_book.xlsx in the order "
                    "Service, Month, Amount. Save the spreadsheet."
                ),
            },
        ],
    },
    # 48c46dc7: NO chrome — tally_book.xlsx
    "osworld_multi_apps_48c46dc7": {
        "sink_path": "/home/user/Documents/Finance/tally_book.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Service", "Month", "Amount"],
        "variants": [
            {
                "html_title": "Travel expenses (October)",
                "rows": [
                    ["Flight", "October", 320.00],
                    ["Hotel", "October", 180.00],
                    ["Taxi", "October", 45.50],
                ],
                "instr": (
                    "The Chrome tab shows three travel expense entries. "
                    "Append all three rows to tally_book.xlsx (Service, "
                    "Month, Amount). Save the file."
                ),
            },
            {
                "html_title": "Conference fees",
                "rows": [
                    ["Registration", "November", 450.00],
                    ["Workshop", "November", 95.00],
                ],
                "instr": (
                    "Two conference fee entries are listed in the Chrome page. "
                    "Add them to tally_book.xlsx as new rows in the order "
                    "Service, Month, Amount. Save."
                ),
            },
            {
                "html_title": "Catering (December)",
                "rows": [
                    ["Holiday Dinner", "December", 480.00],
                    ["Office Snacks", "December", 65.00],
                    ["Coffee Bar", "December", 38.00],
                ],
                "instr": (
                    "The Chrome tab shows three catering expenses. Append "
                    "all three rows to tally_book.xlsx (columns: Service, "
                    "Month, Amount). Save."
                ),
            },
            {
                "html_title": "Shuttle service (January)",
                "rows": [
                    ["Airport Shuttle", "January", 75.00],
                    ["Local Taxi", "January", 22.50],
                ],
                "instr": (
                    "Two shuttle service entries are listed in the Chrome "
                    "page. Add them as new rows at the bottom of "
                    "tally_book.xlsx in the order Service, Month, Amount. "
                    "Save the file."
                ),
            },
        ],
    },
    # 788b3701: NO chrome — tally_book.xlsx
    "osworld_multi_apps_788b3701": {
        "sink_path": "/home/user/Documents/Finance/tally_book.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Service", "Month", "Amount"],
        "variants": [
            {
                "html_title": "Hardware purchases (December)",
                "rows": [
                    ["Monitor", "December", 280.00],
                    ["Keyboard", "December", 90.00],
                    ["Webcam", "December", 60.00],
                ],
                "instr": (
                    "The Chrome tab lists three hardware purchases. "
                    "Append all three rows to tally_book.xlsx (Service, "
                    "Month, Amount). Save the file."
                ),
            },
            {
                "html_title": "Annual renewals",
                "rows": [
                    ["Domain", "January", 12.00],
                    ["Hosting", "January", 96.00],
                ],
                "instr": (
                    "Two annual renewal entries are shown in the Chrome page. "
                    "Add both rows at the bottom of tally_book.xlsx, keeping "
                    "the column order Service, Month, Amount. Save."
                ),
            },
            {
                "html_title": "Maintenance (February)",
                "rows": [
                    ["Cleaning Service", "February", 220.00],
                    ["Repair Work", "February", 95.00],
                    ["Pest Control", "February", 75.00],
                ],
                "instr": (
                    "The Chrome tab shows three maintenance expenses. "
                    "Append all three rows to tally_book.xlsx (Service, "
                    "Month, Amount). Save."
                ),
            },
            {
                "html_title": "Training expenses (March)",
                "rows": [
                    ["Online Course", "March", 199.00],
                    ["Workshop Fee", "March", 350.00],
                ],
                "instr": (
                    "Two training expenses are listed in the Chrome page. "
                    "Add them as new rows at the bottom of tally_book.xlsx "
                    "in the order Service, Month, Amount. Save the file."
                ),
            },
        ],
    },
    # 6f4073b8: NO chrome — ConferenceCity.xlsx (Year/Conference/City, 22 rows)
    "osworld_multi_apps_6f4073b8": {
        "sink_path": "/home/user/Desktop/ConferenceCity.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Year", "Conference", "City"],
        "variants": [
            {
                "html_title": "Future conferences",
                "rows": [
                    [2025, "ICLR", "Singapore"],
                    [2025, "ACL", "Vienna"],
                    [2025, "EMNLP", "Suzhou"],
                ],
                "instr": (
                    "The Chrome tab shows three future conference entries. "
                    "Append all three rows to ConferenceCity.xlsx (columns: "
                    "Year, Conference, City). Save."
                ),
            },
            {
                "html_title": "Historical conferences",
                "rows": [
                    [2018, "NeurIPS", "Montreal"],
                    [2017, "ICML", "Sydney"],
                ],
                "instr": (
                    "Two historical conference entries are listed in the "
                    "Chrome page. Add them as new rows in ConferenceCity.xlsx "
                    "in the order Year, Conference, City. Save the file."
                ),
            },
            {
                "html_title": "Regional workshops",
                "rows": [
                    [2024, "EACL", "Malta"],
                    [2024, "AACL", "Bali"],
                    [2024, "NAACL", "Mexico City"],
                ],
                "instr": (
                    "The Chrome tab shows three regional ACL-family entries. "
                    "Append all three rows to ConferenceCity.xlsx (columns: "
                    "Year, Conference, City). Save."
                ),
            },
            {
                "html_title": "Online editions",
                "rows": [
                    [2020, "ACL Online", "Virtual"],
                    [2021, "EMNLP Online", "Virtual"],
                ],
                "instr": (
                    "Two online conference editions are listed in the Chrome "
                    "page. Add them as new rows at the bottom of "
                    "ConferenceCity.xlsx in the order Year, Conference, City. "
                    "Save."
                ),
            },
        ],
    },
    # 7f35355e: NO chrome — stock.xlsx (Symbol/Company/Stock Price/Revenue/Market Cap)
    "osworld_multi_apps_7f35355e": {
        "sink_path": "/home/user/Desktop/stock.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Symbol", "Company", "Stock Price", "Revenue", "Market Cap"],
        "variants": [
            {
                "html_title": "Tech stock additions",
                "rows": [
                    ["NVDA", "NVIDIA Corp", 480.50, 60900, 1180000],
                    ["AMD", "Advanced Micro Devices", 142.30, 22680, 230000],
                ],
                "instr": (
                    "The Chrome tab lists two tech stock entries. Append both "
                    "rows to stock.xlsx, preserving the column order (Symbol, "
                    "Company, Stock Price, Revenue, Market Cap). Save."
                ),
            },
            {
                "html_title": "Financial stocks",
                "rows": [
                    ["JPM", "JPMorgan Chase", 192.10, 158000, 555000],
                    ["GS", "Goldman Sachs", 425.80, 47330, 142000],
                ],
                "instr": (
                    "Two financial stock entries are shown in the Chrome page. "
                    "Add them as new rows at the bottom of stock.xlsx in the "
                    "order Symbol, Company, Stock Price, Revenue, Market Cap. Save."
                ),
            },
            {
                "html_title": "Consumer-staples stocks",
                "rows": [
                    ["WMT", "Walmart Inc", 60.20, 638790, 489000],
                    ["COST", "Costco Wholesale", 705.50, 242290, 312000],
                    ["PG", "Procter & Gamble", 159.80, 82000, 376000],
                ],
                "instr": (
                    "The Chrome tab shows three consumer-staples stock "
                    "entries. Append all three rows to stock.xlsx in the "
                    "given column order: Symbol, Company, Stock Price, "
                    "Revenue, Market Cap. Save."
                ),
            },
            {
                "html_title": "Healthcare stocks",
                "rows": [
                    ["PFE", "Pfizer Inc", 28.40, 58500, 161000],
                    ["MRK", "Merck & Co", 124.30, 60100, 314000],
                ],
                "instr": (
                    "Two healthcare stock entries are listed in the Chrome "
                    "page. Add both as new rows at the bottom of stock.xlsx, "
                    "preserving the column order (Symbol, Company, Stock "
                    "Price, Revenue, Market Cap). Save."
                ),
            },
            {
                "html_title": "Energy sector stocks",
                "rows": [
                    ["XOM", "ExxonMobil Corp", 110.50, 344582, 446000],
                    ["CVX", "Chevron Corp", 152.40, 196913, 280000],
                    ["BP", "BP p.l.c.", 36.20, 213032, 96000],
                ],
                "instr": (
                    "Three energy-sector stock entries are listed in the "
                    "Chrome page. Append all three rows to stock.xlsx in "
                    "the order Symbol, Company, Stock Price, Revenue, "
                    "Market Cap. Save."
                ),
            },
        ],
    },
    # 881deb30: NO chrome — supported_rate.xlsx (Year + 8 HK universities)
    "osworld_multi_apps_881deb30": {
        "sink_path": "/home/user/Documents/Fundings/supported_rate.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": [
            "Year", "CityU", "HKBU", "LU", "CUHK",
            "HKIEd (EduHK)", "PolyU", "HKUST", "HKU",
        ],
        "variants": [
            {
                "html_title": "Recent year support rates",
                "rows": [
                    [2023, 0.42, 0.31, 0.28, 0.45, 0.30, 0.40, 0.48, 0.50],
                    [2024, 0.44, 0.33, 0.29, 0.46, 0.32, 0.41, 0.49, 0.52],
                ],
                "instr": (
                    "The Chrome tab shows two new yearly support rate rows. "
                    "Append both rows to supported_rate.xlsx (columns: Year, "
                    "CityU, HKBU, LU, CUHK, HKIEd (EduHK), PolyU, HKUST, HKU). "
                    "Save."
                ),
            },
            {
                "html_title": "Projection for next year",
                "rows": [
                    [2025, 0.45, 0.34, 0.30, 0.47, 0.33, 0.42, 0.50, 0.53],
                ],
                "instr": (
                    "One projected support rate row is shown in the Chrome "
                    "page. Add it as a new row at the bottom of "
                    "supported_rate.xlsx (Year, CityU, HKBU, LU, CUHK, "
                    "HKIEd (EduHK), PolyU, HKUST, HKU). Save the file."
                ),
            },
            {
                "html_title": "Three projection years",
                "rows": [
                    [2026, 0.46, 0.35, 0.31, 0.48, 0.34, 0.43, 0.51, 0.54],
                    [2027, 0.47, 0.36, 0.31, 0.49, 0.34, 0.43, 0.52, 0.55],
                    [2028, 0.48, 0.36, 0.32, 0.50, 0.35, 0.44, 0.52, 0.55],
                ],
                "instr": (
                    "The Chrome tab shows three projection-year rows. Append "
                    "all three rows to supported_rate.xlsx in the order "
                    "Year, CityU, HKBU, LU, CUHK, HKIEd (EduHK), PolyU, "
                    "HKUST, HKU. Save."
                ),
            },
            {
                "html_title": "Older years to backfill",
                "rows": [
                    [2010, 0.32, 0.22, 0.20, 0.36, 0.21, 0.31, 0.39, 0.41],
                    [2011, 0.33, 0.23, 0.21, 0.37, 0.22, 0.32, 0.40, 0.42],
                ],
                "instr": (
                    "Two older yearly rows are listed in the Chrome page. "
                    "Add both as new rows at the bottom of supported_rate.xlsx "
                    "in the order Year, CityU, HKBU, LU, CUHK, HKIEd (EduHK), "
                    "PolyU, HKUST, HKU. Save."
                ),
            },
            {
                "html_title": "Mid-range years",
                "rows": [
                    [2015, 0.37, 0.27, 0.25, 0.41, 0.27, 0.36, 0.43, 0.46],
                    [2016, 0.38, 0.28, 0.26, 0.42, 0.27, 0.37, 0.44, 0.47],
                    [2017, 0.39, 0.29, 0.27, 0.43, 0.28, 0.38, 0.45, 0.48],
                ],
                "instr": (
                    "Three mid-range yearly rows are listed in the Chrome "
                    "page. Append all three rows to supported_rate.xlsx in "
                    "the order Year, CityU, HKBU, LU, CUHK, HKIEd (EduHK), "
                    "PolyU, HKUST, HKU. Save."
                ),
            },
        ],
    },
    # f5c13cdd: NO chrome — tuition_payment.xlsx (Name/Email/Payment, 27 rows)
    "osworld_multi_apps_f5c13cdd": {
        "sink_path": "/home/user/Documents/Departments/finance/tuition_payment.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Name", "Email", "Payment"],
        "variants": [
            {
                "html_title": "Late payments to record",
                "rows": [
                    ["Liam Brown", "liam.b@example.edu", "Pending"],
                    ["Olivia Lin", "olivia.l@example.edu", "Pending"],
                    ["Noah Patel", "noah.p@example.edu", "Pending"],
                ],
                "instr": (
                    "The Chrome tab shows three new tuition records. Append "
                    "all three rows to tuition_payment.xlsx (columns: Name, "
                    "Email, Payment). Save."
                ),
            },
            {
                "html_title": "Newly cleared payments",
                "rows": [
                    ["Sophia Reed", "sophia.r@example.edu", "Paid"],
                    ["Mason Wu", "mason.w@example.edu", "Paid"],
                ],
                "instr": (
                    "Two newly cleared tuition records are listed in the "
                    "Chrome page. Add them as new rows at the bottom of "
                    "tuition_payment.xlsx (Name, Email, Payment). Save."
                ),
            },
            {
                "html_title": "Scholarship recipients",
                "rows": [
                    ["Daria Cole", "daria.c@example.edu", "Waived"],
                    ["Kenji Ito", "kenji.i@example.edu", "Waived"],
                    ["Priya Rao", "priya.r@example.edu", "Waived"],
                ],
                "instr": (
                    "The Chrome tab shows three scholarship-recipient records. "
                    "Append all three rows to tuition_payment.xlsx (columns: "
                    "Name, Email, Payment). Save the file."
                ),
            },
            {
                "html_title": "International student records",
                "rows": [
                    ["Aoi Yamada", "aoi.y@example.edu", "Pending"],
                    ["Carlos Mendez", "carlos.m@example.edu", "Paid"],
                ],
                "instr": (
                    "Two international-student records are listed in the "
                    "Chrome page. Add them as new rows in tuition_payment.xlsx, "
                    "preserving the column order Name, Email, Payment. Save."
                ),
            },
        ],
    },
    # 00fa164e: NO chrome — awe_desk_env.docx (31 paragraphs, research environment)
    "osworld_multi_apps_00fa164e": {
        "sink_path": "/home/user/Documents/awesome-desktop/awe_desk_env.docx",
        "sink_ext": "docx",
        "sink_headers": ["Note"],
        "variants": [
            {
                "html_title": "Notes to append",
                "rows": [
                    ["Reviewed setup with the lab on 2024-09-12; no blockers."],
                    ["Migrated nightly snapshots to the secondary disk array."],
                ],
                "instr": (
                    "The Chrome tab lists two short notes. Append each note "
                    "as a new paragraph at the end of awe_desk_env.docx, "
                    "preserving the text exactly. Save the document."
                ),
            },
            {
                "html_title": "Three more notes",
                "rows": [
                    ["Confirmed CUDA driver pinned to v535 across all nodes."],
                    ["Monitoring agents now report to the central dashboard."],
                    ["Next maintenance window scheduled for 2024-10-04."],
                ],
                "instr": (
                    "Three additional notes are shown in the Chrome page. Add "
                    "each as a new paragraph at the end of awe_desk_env.docx. "
                    "Save."
                ),
            },
            {
                "html_title": "Backup and recovery notes",
                "rows": [
                    ["Tested restore from the off-site mirror within the SLA window."],
                    ["Verified that the rollback playbook completes in under twenty minutes."],
                ],
                "instr": (
                    "The Chrome tab shows two backup-and-recovery notes. "
                    "Append each as a new paragraph at the end of "
                    "awe_desk_env.docx, preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Security follow-ups",
                "rows": [
                    ["Confirmed least-privilege roles on the build cluster."],
                    ["Enrolled the new on-call rotation into the alerting tier."],
                    ["Quarterly secrets rotation completed without downtime."],
                ],
                "instr": (
                    "Three security follow-up notes are shown in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "awe_desk_env.docx. Save the document."
                ),
            },
        ],
    },
    # 02ce9a50: NO chrome — top-10-linux-commands-for-newbies.docx (175 paragraphs)
    "osworld_multi_apps_02ce9a50": {
        "sink_path": "/home/user/Desktop/top-10-linux-commands-for-newbies.docx",
        "sink_ext": "docx",
        "sink_headers": ["Tip"],
        "variants": [
            {
                "html_title": "Two extra tips",
                "rows": [
                    ["Tip: use `man <command>` to read the manual page for any tool."],
                    ["Tip: append `--help` to a command to see its quick usage summary."],
                ],
                "instr": (
                    "The Chrome tab shows two extra Linux tips. Append each "
                    "tip as a new paragraph at the end of "
                    "top-10-linux-commands-for-newbies.docx, preserving the "
                    "text exactly. Save."
                ),
            },
            {
                "html_title": "Three more tips",
                "rows": [
                    ["Tip: combine `find` with `xargs` to act on many files."],
                    ["Tip: pipe `history` into `grep` to recall a recent command."],
                    ["Tip: use `tar -czvf out.tar.gz dir/` to compress a folder."],
                ],
                "instr": (
                    "Three additional Linux tips are listed in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "top-10-linux-commands-for-newbies.docx. Save."
                ),
            },
            {
                "html_title": "Networking tips",
                "rows": [
                    ["Tip: use `ss -tulpn` to list listening sockets and owning processes."],
                    ["Tip: `curl -I url` shows just the response headers."],
                ],
                "instr": (
                    "The Chrome tab lists two networking tips. Append each "
                    "as a new paragraph at the end of "
                    "top-10-linux-commands-for-newbies.docx, preserving the "
                    "text exactly. Save."
                ),
            },
            {
                "html_title": "Process management tips",
                "rows": [
                    ["Tip: use `pgrep -f <pattern>` to find process IDs by command line."],
                    ["Tip: `kill -SIGTERM <pid>` requests graceful shutdown; `-9` forces it."],
                    ["Tip: `nohup <cmd> &` runs a command detached from the terminal."],
                ],
                "instr": (
                    "Three process-management tips are shown in the Chrome "
                    "page. Append each as a new paragraph at the end of "
                    "top-10-linux-commands-for-newbies.docx. Save."
                ),
            },
            {
                "html_title": "Disk and storage tips",
                "rows": [
                    ["Tip: `du -sh dir/` shows the total size of a directory."],
                    ["Tip: `df -h` reports filesystem disk usage in human-readable form."],
                ],
                "instr": (
                    "The Chrome tab shows two disk-and-storage tips. Append "
                    "each as a new paragraph at the end of "
                    "top-10-linux-commands-for-newbies.docx, preserving the "
                    "text exactly. Save."
                ),
            },
        ],
    },
    # 1f18aa87: NO chrome — Grammer test 1.docx (58 paragraphs)
    "osworld_multi_apps_1f18aa87": {
        "sink_path": "/home/user/Desktop/Grammer test 1.docx",
        "sink_ext": "docx",
        "sink_headers": ["Sample sentence"],
        "variants": [
            {
                "html_title": "Two sample sentences",
                "rows": [
                    ["She has been working on the project since June, hasn't she?"],
                    ["By the time we arrived, they had already finished the meal."],
                ],
                "instr": (
                    "The Chrome tab shows two sample grammar sentences. "
                    "Append each as a new paragraph at the end of "
                    "'Grammer test 1.docx', preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Three more sentences",
                "rows": [
                    ["If I had known earlier, I would have told you immediately."],
                    ["The book that I borrowed from the library is overdue."],
                    ["Neither the manager nor the staff were aware of the change."],
                ],
                "instr": (
                    "Three additional sample sentences are listed in the "
                    "Chrome page. Add each as a new paragraph at the end of "
                    "'Grammer test 1.docx'. Save."
                ),
            },
            {
                "html_title": "Conditional sentences",
                "rows": [
                    ["Should you need help, please contact the support desk."],
                    ["Were he here, he would solve this in a minute."],
                ],
                "instr": (
                    "The Chrome tab shows two conditional-form sentences. "
                    "Append each as a new paragraph at the end of "
                    "'Grammer test 1.docx', preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Reported speech examples",
                "rows": [
                    ["She said that she had finished her homework before dinner."],
                    ["He told me that the meeting would be postponed by an hour."],
                    ["They asked whether we could join them for lunch tomorrow."],
                ],
                "instr": (
                    "Three reported-speech examples are listed in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "'Grammer test 1.docx'. Save."
                ),
            },
        ],
    },
    # 20236825: NO chrome — Bubble_Sort_tutorial.docx (71 paragraphs)
    "osworld_multi_apps_20236825": {
        "sink_path": "/home/user/Desktop/Bubble_Sort_tutorial.docx",
        "sink_ext": "docx",
        "sink_headers": ["Note"],
        "variants": [
            {
                "html_title": "Two practice notes",
                "rows": [
                    ["Note: bubble sort has worst-case O(n^2) time complexity."],
                    ["Note: it is stable but rarely used for large datasets."],
                ],
                "instr": (
                    "The Chrome tab shows two practice notes. Append each as "
                    "a new paragraph at the end of Bubble_Sort_tutorial.docx, "
                    "preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Three implementation reminders",
                "rows": [
                    ["Reminder: swap two elements only when the left is greater than the right."],
                    ["Reminder: an early-exit optimization breaks when no swaps occur."],
                    ["Reminder: always test on already-sorted and reverse-sorted arrays."],
                ],
                "instr": (
                    "Three additional reminders are listed in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "Bubble_Sort_tutorial.docx. Save."
                ),
            },
            {
                "html_title": "Comparison with other sorts",
                "rows": [
                    ["Compared with insertion sort, bubble sort makes more swaps on average."],
                    ["Selection sort uses fewer swaps than bubble sort but the same number of comparisons."],
                ],
                "instr": (
                    "The Chrome tab shows two comparison notes. Append each "
                    "as a new paragraph at the end of "
                    "Bubble_Sort_tutorial.docx, preserving the text exactly. "
                    "Save."
                ),
            },
            {
                "html_title": "Common pitfalls",
                "rows": [
                    ["Pitfall: forgetting to reset the swap flag at the start of each pass."],
                    ["Pitfall: using strict `<` when ties matter for stability."],
                    ["Pitfall: ignoring the cost of repeated swaps on linked structures."],
                ],
                "instr": (
                    "Three common-pitfall notes are listed in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "Bubble_Sort_tutorial.docx. Save the document."
                ),
            },
        ],
    },
    # 2b9493d7: NO chrome — 15-MB-docx-file-download.docx (66 paragraphs)
    "osworld_multi_apps_2b9493d7": {
        "sink_path": "/home/user/Desktop/15-MB-docx-file-download.docx",
        "sink_ext": "docx",
        "sink_headers": ["Note"],
        "variants": [
            {
                "html_title": "Two appended notes",
                "rows": [
                    ["Editorial note: revised section three on 2024-09-15."],
                    ["Editorial note: cross-references checked against source."],
                ],
                "instr": (
                    "The Chrome tab shows two editorial notes. Append each "
                    "as a new paragraph at the end of "
                    "15-MB-docx-file-download.docx, preserving the text "
                    "exactly. Save."
                ),
            },
            {
                "html_title": "Three more notes",
                "rows": [
                    ["Page numbering verified end-to-end."],
                    ["Bibliography order matches the citation index."],
                    ["File metadata refreshed before final upload."],
                ],
                "instr": (
                    "Three additional editorial notes are shown in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "15-MB-docx-file-download.docx. Save."
                ),
            },
            {
                "html_title": "Reviewer feedback",
                "rows": [
                    ["Reviewer A: tighten the executive summary by one paragraph."],
                    ["Reviewer B: clarify the methodology footnote on page 4."],
                ],
                "instr": (
                    "The Chrome tab shows two reviewer-feedback notes. "
                    "Append each as a new paragraph at the end of "
                    "15-MB-docx-file-download.docx. Save."
                ),
            },
            {
                "html_title": "Distribution checklist",
                "rows": [
                    ["Verified target list against the latest CRM export."],
                    ["Confirmed embargoed items remain unpublished."],
                    ["Final PDF and DOCX hashes archived for audit."],
                ],
                "instr": (
                    "Three distribution-checklist items are listed in the "
                    "Chrome page. Add each as a new paragraph at the end of "
                    "15-MB-docx-file-download.docx. Save."
                ),
            },
        ],
    },
    # 8df7e444: NO chrome — Recruitment_and_retention...docx (250 paragraphs)
    "osworld_multi_apps_8df7e444": {
        "sink_path": "/home/user/Recruitment_and_retention_of_health_professionals_across_Europe.docx",
        "sink_ext": "docx",
        "sink_headers": ["Annotation"],
        "variants": [
            {
                "html_title": "Two reviewer annotations",
                "rows": [
                    ["Reviewer note: align terminology with the WHO 2023 health workforce report."],
                    ["Reviewer note: cross-check Table 4 against the OECD figures from Q3 2023."],
                ],
                "instr": (
                    "The Chrome tab shows two reviewer annotations. Append "
                    "each as a new paragraph at the end of the open "
                    "Recruitment_and_retention_of_health_professionals_"
                    "across_Europe.docx, preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Three reviewer notes",
                "rows": [
                    ["Reviewer note: consider adding a recent example from Eastern European regions."],
                    ["Reviewer note: figure captions should reference the source dataset."],
                    ["Reviewer note: appendix needs a glossary of abbreviations."],
                ],
                "instr": (
                    "Three additional reviewer notes are listed in the "
                    "Chrome page. Add each as a new paragraph at the end of "
                    "the recruitment-and-retention docx. Save."
                ),
            },
            {
                "html_title": "Editor questions",
                "rows": [
                    ["Editor question: can the abstract better foreground policy implications?"],
                    ["Editor question: should we add a sub-section on telehealth recruitment?"],
                ],
                "instr": (
                    "The Chrome tab shows two editor questions. Append each "
                    "as a new paragraph at the end of the recruitment-and-"
                    "retention docx, preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Methodology clarifications",
                "rows": [
                    ["Methodology: include selection criteria for the longitudinal study cohort."],
                    ["Methodology: state inter-coder reliability for the qualitative analysis."],
                    ["Methodology: list ethical-approval references in a numbered list."],
                ],
                "instr": (
                    "Three methodology clarifications are shown in the Chrome "
                    "page. Add each as a new paragraph at the end of the "
                    "recruitment-and-retention docx. Save."
                ),
            },
            {
                "html_title": "Stakeholder-feedback notes",
                "rows": [
                    ["Stakeholder note: ministerial reviewers requested clearer policy recommendations."],
                    ["Stakeholder note: hospital association asked for region-specific case studies."],
                ],
                "instr": (
                    "The Chrome tab shows two stakeholder-feedback notes. "
                    "Append each as a new paragraph at the end of the "
                    "recruitment-and-retention docx, preserving the text "
                    "exactly. Save."
                ),
            },
        ],
    },
    # 09a37c51: NO chrome — requirment.docx (6 paragraphs, image-edit request doc)
    "osworld_multi_apps_09a37c51": {
        "sink_path": "/home/user/Desktop/requirment.docx",
        "sink_ext": "docx",
        "sink_headers": ["Comment"],
        "variants": [
            {
                "html_title": "Two follow-up comments",
                "rows": [
                    ["Follow-up: please match the brand color to hex #2C5DA8."],
                    ["Follow-up: keep the original aspect ratio when resizing."],
                ],
                "instr": (
                    "The Chrome tab shows two follow-up comments. Append "
                    "each as a new paragraph at the end of requirment.docx, "
                    "preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Three more comments",
                "rows": [
                    ["Final draft due 2024-10-01."],
                    ["Use the provided palette for backgrounds only."],
                    ["Send the layered source file alongside the export."],
                ],
                "instr": (
                    "Three additional comments are listed in the Chrome page. "
                    "Add each as a new paragraph at the end of requirment.docx. "
                    "Save."
                ),
            },
            {
                "html_title": "Asset-handover checklist",
                "rows": [
                    ["Asset checklist: include both PNG and SVG exports."],
                    ["Asset checklist: name files using lowercase and dashes."],
                ],
                "instr": (
                    "The Chrome tab shows two asset-handover checklist items. "
                    "Append each as a new paragraph at the end of "
                    "requirment.docx, preserving the text exactly. Save."
                ),
            },
        ],
    },
    # 716a6079: NO chrome — secret.docx (2 paragraphs, hidden in /home/user/Data3/List3/)
    "osworld_multi_apps_716a6079": {
        "sink_path": "/home/user/Data3/List3/secret.docx",
        "sink_ext": "docx",
        "sink_headers": ["Note"],
        "variants": [
            {
                "html_title": "New entries for the secret notebook",
                "rows": [
                    ["Backup of the encryption key was rotated on 2024-09-30."],
                    ["Off-site copy verified by the on-call engineer."],
                ],
                "instr": (
                    "The Chrome tab shows two new notebook entries. Append "
                    "each as a new paragraph at the end of "
                    "/home/user/Data3/List3/secret.docx, preserving the text "
                    "exactly. Save the document."
                ),
            },
            {
                "html_title": "Three more entries",
                "rows": [
                    ["Audit log retention extended to 90 days."],
                    ["MFA enforced across all admin accounts."],
                    ["Quarterly access review scheduled for 2024-12-15."],
                ],
                "instr": (
                    "Three more notebook entries are listed in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "/home/user/Data3/List3/secret.docx. Save the file."
                ),
            },
        ],
    },
    # 7ff48d5b: NO chrome — AllLocations.docx (1 paragraph)
    "osworld_multi_apps_7ff48d5b": {
        "sink_path": "/home/user/Desktop/AllLocations.docx",
        "sink_ext": "docx",
        "sink_headers": ["Location entry"],
        "variants": [
            {
                "html_title": "New location entries",
                "rows": [
                    ["Hong Kong - Asia / Special Administrative Region."],
                    ["Lisbon - Portugal / coastal capital."],
                    ["Reykjavik - Iceland / North Atlantic capital."],
                ],
                "instr": (
                    "The Chrome tab shows three new location entries. Append "
                    "each as a new paragraph at the end of AllLocations.docx, "
                    "preserving the text exactly. Save the document."
                ),
            },
            {
                "html_title": "Two more locations",
                "rows": [
                    ["Vienna - Austria / Central European capital."],
                    ["Helsinki - Finland / Baltic capital."],
                ],
                "instr": (
                    "Two more location entries are shown in the Chrome page. "
                    "Add each as a new paragraph at the end of "
                    "AllLocations.docx. Save the file."
                ),
            },
        ],
    },
    # 873cafdd: NO chrome — Recommended_plugin_list.docx (6 paragraphs)
    "osworld_multi_apps_873cafdd": {
        "sink_path": "/home/user/Desktop/Recommended_plugin_list.docx",
        "sink_ext": "docx",
        "sink_headers": ["Plugin entry"],
        "variants": [
            {
                "html_title": "New plugin recommendations",
                "rows": [
                    ["Plugin: Better Color Picker - improved palette tools."],
                    ["Plugin: Quick Slice - one-click image slicing."],
                ],
                "instr": (
                    "The Chrome tab shows two new plugin recommendations. "
                    "Append each as a new paragraph at the end of "
                    "Recommended_plugin_list.docx, preserving the text "
                    "exactly. Save."
                ),
            },
            {
                "html_title": "Three additional plugins",
                "rows": [
                    ["Plugin: Layer Pro - layer state preservation across sessions."],
                    ["Plugin: Smart Crop - aspect-aware crop suggestions."],
                    ["Plugin: Batch Export - exports many sizes in one click."],
                ],
                "instr": (
                    "Three additional plugins are listed in the Chrome page. "
                    "Add each as a new paragraph at the end of "
                    "Recommended_plugin_list.docx. Save."
                ),
            },
            {
                "html_title": "Productivity plugins",
                "rows": [
                    ["Plugin: Quick Switcher - keyboard-only window navigation."],
                    ["Plugin: Snippet Vault - personal snippet manager."],
                ],
                "instr": (
                    "The Chrome tab shows two productivity plugins. Append "
                    "each as a new paragraph at the end of "
                    "Recommended_plugin_list.docx, preserving the text "
                    "exactly. Save."
                ),
            },
        ],
    },
    # df67aebb: chrome+dblp → references.docx (11 paragraphs)
    "osworld_multi_apps_df67aebb": {
        "sink_path": "/home/user/Desktop/references.docx",
        "sink_ext": "docx",
        "sink_headers": ["Reference entry"],
        "variants": [
            {
                "html_title": "New references to cite",
                "rows": [
                    [
                        "Smith, J. and Lee, K. 2023. "
                        "Robust Prompting for Code LLMs. ACL 2023, pp. 100-115."
                    ],
                    [
                        "Wang, M. 2024. "
                        "Verifier-Guided Decoding. EMNLP 2024, pp. 200-215."
                    ],
                ],
                "instr": (
                    "The Chrome tab shows two new references. Append each as "
                    "a new paragraph at the end of references.docx, preserving "
                    "the reference text verbatim. Save the document."
                ),
            },
            {
                "html_title": "Three more references",
                "rows": [
                    [
                        "Park, S. and Kim, J. 2024. "
                        "Multi-Agent Coordination Benchmarks. NeurIPS 2024."
                    ],
                    [
                        "Chen, R. 2023. "
                        "Sparse Attention for Long Documents. ICLR 2023, pp. 50-65."
                    ],
                    [
                        "Garcia, A. 2024. "
                        "Tool-Augmented Reasoning Survey. arXiv 2402.12345."
                    ],
                ],
                "instr": (
                    "Three additional references are listed in the Chrome page. "
                    "Add each as a new paragraph at the end of references.docx. "
                    "Save the file."
                ),
            },
            {
                "html_title": "Workshop references",
                "rows": [
                    [
                        "Ito, S. and Brown, P. 2023. "
                        "Curriculum Reinforcement Learning. ICLR Workshop 2023."
                    ],
                    [
                        "Diaz, R. 2024. "
                        "Vision-Language Pretraining Survey. CVPR Workshop 2024."
                    ],
                ],
                "instr": (
                    "Two workshop references are listed in the Chrome page. "
                    "Append each as a new paragraph at the end of "
                    "references.docx, preserving the reference text verbatim. "
                    "Save."
                ),
            },
            {
                "html_title": "Older foundational references",
                "rows": [
                    [
                        "Vaswani, A. et al. 2017. "
                        "Attention Is All You Need. NeurIPS 2017."
                    ],
                    [
                        "Devlin, J. et al. 2019. "
                        "BERT: Pre-training of Deep Bidirectional Transformers. NAACL 2019."
                    ],
                    [
                        "Brown, T. et al. 2020. "
                        "Language Models are Few-Shot Learners. NeurIPS 2020."
                    ],
                ],
                "instr": (
                    "Three foundational references are shown in the Chrome "
                    "page. Add each as a new paragraph at the end of "
                    "references.docx. Save."
                ),
            },
        ],
    },
    # 68a25bd4: empty xlsx (rsc-ebook-collection-2023.xlsx, 0r×0c)
    # `sink_starts_empty=True` → gold writes our defined headers as row 1.
    "osworld_multi_apps_68a25bd4": {
        "sink_path": "/home/user/Desktop/rsc-ebook-collection-2023.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Title", "Author", "Year", "Link"],
        "sink_starts_empty": True,
        "variants": [
            {
                "html_title": "Books to add to the collection",
                "rows": [
                    ["The Quantum Code", "A. Lee", 2022, "https://example.org/q.pdf"],
                    ["Synthetic Minds", "B. Park", 2023, "https://example.org/s.pdf"],
                    ["Algorithmic Histories", "C. Diaz", 2021, "https://example.org/a.pdf"],
                ],
                "instr": (
                    "The Chrome tab shows a small table with the columns "
                    "Title, Author, Year, Link and three book rows. Copy the "
                    "table into the empty rsc-ebook-collection-2023.xlsx — "
                    "place the headers in row 1 and the three data rows below "
                    "them, preserving the column order. Save."
                ),
            },
            {
                "html_title": "Two more books",
                "rows": [
                    ["Foundations of Reasoning", "D. Han", 2020, "https://example.org/f.pdf"],
                    ["Patterns in Practice", "E. Park", 2024, "https://example.org/p.pdf"],
                ],
                "instr": (
                    "The Chrome tab shows a header row (Title, Author, Year, "
                    "Link) and two book rows. Copy the entire table (headers "
                    "in row 1, data in rows 2-3) into the empty "
                    "rsc-ebook-collection-2023.xlsx. Save."
                ),
            },
            {
                "html_title": "International edition",
                "rows": [
                    ["Estructuras Datos", "F. Mendez", 2022, "https://example.org/ed.pdf"],
                    ["Reseaux Numeriques", "G. Lambert", 2023, "https://example.org/rn.pdf"],
                ],
                "instr": (
                    "Two international book entries are shown in the Chrome "
                    "page. Copy the headers (Title, Author, Year, Link) and "
                    "the two rows into the empty rsc-ebook-collection-2023.xlsx, "
                    "starting at cell A1. Save."
                ),
            },
            {
                "html_title": "Older classics",
                "rows": [
                    ["Foundations of Computing", "I. Park", 2010, "https://example.org/fc.pdf"],
                    ["Numerical Recipes", "J. Anand", 2007, "https://example.org/nr.pdf"],
                    ["Discrete Structures", "K. Romero", 2012, "https://example.org/ds.pdf"],
                ],
                "instr": (
                    "Three older-classic book entries are listed in the "
                    "Chrome page with headers (Title, Author, Year, Link) "
                    "and three rows. Copy the entire table into the empty "
                    "rsc-ebook-collection-2023.xlsx with headers in row 1 "
                    "and data in rows 2-4. Save."
                ),
            },
        ],
    },
    # 869de13e: empty xlsx (2023_validation_<uuid>.xlsx)
    "osworld_multi_apps_869de13e": {
        "sink_path": "/home/user/Desktop/2023_validation_7bd855d8-463d-4ed5-93ca-5fe35145f733.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Item", "Status", "Note"],
        "sink_starts_empty": True,
        "variants": [
            {
                "html_title": "Validation items",
                "rows": [
                    ["Module A", "Pass", "Reviewed by team"],
                    ["Module B", "Fail", "Investigate logs"],
                    ["Module C", "Pass", "No issues"],
                ],
                "instr": (
                    "The Chrome tab shows a 3-column validation table. Copy "
                    "the entire table (Item / Status / Note) into the empty "
                    "spreadsheet — headers in row 1 and the three data rows "
                    "below them. Save."
                ),
            },
            {
                "html_title": "Additional validations",
                "rows": [
                    ["Module D", "Pending", "Awaiting CI"],
                    ["Module E", "Pass", "Verified"],
                ],
                "instr": (
                    "Two more validation items are listed in the Chrome page. "
                    "Copy the header row (Item, Status, Note) and the two "
                    "data rows into the empty xlsx. Save."
                ),
            },
            {
                "html_title": "Audit results",
                "rows": [
                    ["Module F", "Pass", "Audit OK"],
                    ["Module G", "Fail", "Re-test required"],
                    ["Module H", "Pass", "Accepted"],
                ],
                "instr": (
                    "The Chrome tab shows three audit-result rows with "
                    "headers Item, Status, Note. Copy the whole table (1 "
                    "header row + 3 data rows) into the empty 2023 "
                    "validation xlsx. Save the file."
                ),
            },
            {
                "html_title": "Sign-off entries",
                "rows": [
                    ["Module I", "Pass", "Signed off by lead"],
                    ["Module J", "Pending", "Awaiting QA review"],
                ],
                "instr": (
                    "The Chrome tab shows two sign-off validation rows. "
                    "Copy the headers (Item, Status, Note) and the two "
                    "data rows into the empty 2023 validation xlsx. Save."
                ),
            },
        ],
    },
    # da52d699: empty xlsx (2023_validation_Book_Reading_Rate.xlsx)
    "osworld_multi_apps_da52d699": {
        "sink_path": "/home/user/Desktop/2023_validation_Book_Reading_Rate.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Book", "Reader", "Pages Read"],
        "sink_starts_empty": True,
        "variants": [
            {
                "html_title": "Reading log entries",
                "rows": [
                    ["Brave New World", "Avery", 120],
                    ["The Stranger", "Beth", 85],
                    ["Klara and the Sun", "Cleo", 200],
                ],
                "instr": (
                    "The Chrome tab shows a reading-log table with columns "
                    "Book, Reader, Pages Read. Copy the entire table (1 "
                    "header row + 3 data rows) into the empty "
                    "2023_validation_Book_Reading_Rate.xlsx, starting at A1. "
                    "Save."
                ),
            },
            {
                "html_title": "Two more entries",
                "rows": [
                    ["Project Hail Mary", "Dean", 320],
                    ["Educated", "Eliza", 150],
                ],
                "instr": (
                    "Two more reading-log entries are listed in the Chrome "
                    "page. Copy the header row (Book, Reader, Pages Read) "
                    "and the two data rows into the empty xlsx. Save."
                ),
            },
            {
                "html_title": "Childrens-section reading",
                "rows": [
                    ["Where the Wild Things Are", "Finn", 32],
                    ["Charlotte's Web", "Greta", 184],
                ],
                "instr": (
                    "The Chrome tab shows two children's-section reading "
                    "entries with column titles Book, Reader, Pages Read. "
                    "Copy the headers and the two rows into the empty "
                    "2023_validation_Book_Reading_Rate.xlsx. Save."
                ),
            },
            {
                "html_title": "Reading group records",
                "rows": [
                    ["Beloved", "Helen", 95],
                    ["The Underground Railroad", "Ian", 168],
                    ["Pachinko", "June", 240],
                ],
                "instr": (
                    "The Chrome tab shows three reading-group entries with "
                    "headers Book, Reader, Pages Read. Copy the headers and "
                    "the three rows into the empty 2023 reading-rate xlsx. "
                    "Save."
                ),
            },
        ],
    },
    # deec51c9: empty xlsx (New Large Language Models.xlsx)
    "osworld_multi_apps_deec51c9": {
        "sink_path": "/home/user/Desktop/New Large Language Models.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["Model", "Org", "Params (B)", "Released"],
        "sink_starts_empty": True,
        "variants": [
            {
                "html_title": "New LLM releases",
                "rows": [
                    ["Llama-3", "Meta", 70, "2024-04"],
                    ["Mixtral 8x22B", "Mistral", 141, "2024-04"],
                    ["Qwen2-72B", "Alibaba", 72, "2024-06"],
                ],
                "instr": (
                    "The Chrome tab shows a table with columns Model, Org, "
                    "Params (B), Released and three LLM releases. Copy the "
                    "entire table (1 header row + 3 data rows) into the empty "
                    "New Large Language Models.xlsx, starting at cell A1. Save."
                ),
            },
            {
                "html_title": "Two more LLMs",
                "rows": [
                    ["Phi-3-mini", "Microsoft", 4, "2024-04"],
                    ["Gemma-7B", "Google", 7, "2024-02"],
                ],
                "instr": (
                    "Two more LLM entries are listed in the Chrome page. "
                    "Copy the headers (Model, Org, Params (B), Released) "
                    "and the two data rows into the empty xlsx. Save."
                ),
            },
            {
                "html_title": "Open-source LLMs",
                "rows": [
                    ["DeepSeek-V2", "DeepSeek", 236, "2024-05"],
                    ["Yi-34B", "01.AI", 34, "2023-11"],
                ],
                "instr": (
                    "The Chrome tab shows two open-source LLM entries with "
                    "column titles Model, Org, Params (B), Released. Copy "
                    "the headers and the two data rows into the empty "
                    "New Large Language Models.xlsx. Save."
                ),
            },
            {
                "html_title": "Multilingual LLMs",
                "rows": [
                    ["Bloom", "BigScience", 176, "2022-07"],
                    ["mT5-XXL", "Google", 13, "2021-03"],
                    ["NLLB-200", "Meta", 54, "2022-07"],
                ],
                "instr": (
                    "The Chrome tab shows three multilingual LLM entries "
                    "with headers Model, Org, Params (B), Released. Copy "
                    "the headers and three data rows into the empty "
                    "New Large Language Models.xlsx, starting at row 1. "
                    "Save."
                ),
            },
        ],
    },
    # 3a93cae4: NO chrome — Zheng He.docx (129 paragraphs, students work)
    "osworld_multi_apps_3a93cae4": {
        "sink_path": "/home/user/Desktop/students work/Zheng He .docx",
        "sink_ext": "docx",
        "sink_headers": ["Reviewer comment"],
        "variants": [
            {
                "html_title": "Two reviewer comments",
                "rows": [
                    ["Reviewer comment: tighten the timeline section to focus on Zheng He's later voyages."],
                    ["Reviewer comment: cite primary sources for the 1421 expedition details."],
                ],
                "instr": (
                    "The Chrome tab shows two reviewer comments. Append each "
                    "as a new paragraph at the end of the open Zheng He.docx, "
                    "preserving the text exactly. Save the document."
                ),
            },
            {
                "html_title": "Three more comments",
                "rows": [
                    ["Reviewer comment: clarify the role of the Yongle Emperor in funding the voyages."],
                    ["Reviewer comment: include a brief paragraph on the social impact in port cities."],
                    ["Reviewer comment: cross-check ship-size estimates against modern reconstructions."],
                ],
                "instr": (
                    "Three more reviewer comments are listed in the Chrome "
                    "page. Append each as a new paragraph at the end of the "
                    "Zheng He.docx case-study. Save."
                ),
            },
            {
                "html_title": "Editorial suggestions",
                "rows": [
                    ["Editorial suggestion: split the trade-routes paragraph into two for readability."],
                    ["Editorial suggestion: add a short conclusion summarizing the lasting legacy."],
                ],
                "instr": (
                    "The Chrome tab shows two editorial suggestions. Add each "
                    "as a new paragraph at the end of the Zheng He.docx. Save."
                ),
            },
        ],
    },
    # 2c1ebcd7: same Zheng He.docx file (129p) — different eval base, separate task_id
    "osworld_multi_apps_2c1ebcd7": {
        "sink_path": "/home/user/Desktop/students work/Zheng He .docx",
        "sink_ext": "docx",
        "sink_headers": ["Annotation"],
        "variants": [
            {
                "html_title": "Two factual annotations",
                "rows": [
                    ["Fact: Zheng He commanded seven major voyages between 1405 and 1433."],
                    ["Fact: the treasure ships were among the largest wooden vessels ever built."],
                ],
                "instr": (
                    "The Chrome tab shows two factual annotations. Append "
                    "each as a new paragraph at the end of the Zheng He.docx, "
                    "preserving the text exactly. Save."
                ),
            },
            {
                "html_title": "Bibliography pointers",
                "rows": [
                    ["See also: Levathes (1994) - When China Ruled the Seas."],
                    ["See also: Dreyer (2007) - Zheng He: China and the Oceans in the Early Ming Dynasty."],
                    ["See also: Sen (2016) - The Impact of Zheng He's Expeditions on Indian Ocean Interactions."],
                ],
                "instr": (
                    "Three bibliography pointers are listed in the Chrome "
                    "page. Add each as a new paragraph at the end of the "
                    "Zheng He.docx. Save."
                ),
            },
            {
                "html_title": "Discussion prompts",
                "rows": [
                    ["Discussion: how did Ming maritime policy shift after the Yongle reign?"],
                    ["Discussion: compare the voyages with contemporaneous European exploration."],
                ],
                "instr": (
                    "The Chrome tab shows two discussion prompts. Append "
                    "each as a new paragraph at the end of the Zheng He.docx. "
                    "Save the document."
                ),
            },
        ],
    },
    # 3680a5ee: file1.xlsx (5001 rows × 1 col, single column "First Name")
    "osworld_multi_apps_3680a5ee": {
        "sink_path": "/home/user/Desktop/file1.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": ["First Name"],
        "variants": [
            {
                "html_title": "First names to add",
                "rows": [
                    ["Aria"],
                    ["Bryce"],
                    ["Cleo"],
                ],
                "instr": (
                    "The Chrome tab shows three first names in a single-"
                    "column table. Append all three rows to the bottom of "
                    "file1.xlsx (column: First Name). Save."
                ),
            },
            {
                "html_title": "Two more names",
                "rows": [
                    ["Dax"],
                    ["Elia"],
                ],
                "instr": (
                    "Two first names are shown in the Chrome page. Append "
                    "both as new rows at the bottom of file1.xlsx in the "
                    "First Name column. Save the file."
                ),
            },
        ],
    },
    # 4c26e3f3: NO chrome — PPT-Template_widescreen.pptx (16 slides)
    "osworld_multi_apps_4c26e3f3": {
        "sink_path": "/home/user/Desktop/PPT-Template_widescreen.pptx",
        "sink_ext": "pptx",
        "sink_headers": ["Slide title"],
        "variants": [
            {
                "html_title": "Slide titles to add",
                "rows": [
                    ["Future Outlook"],
                    ["Conclusion"],
                ],
                "instr": (
                    "The Chrome tab shows two slide titles in a single-column "
                    "table. Insert two new slides at the end of "
                    "PPT-Template_widescreen.pptx, set each new slide's "
                    "title to the corresponding text from Chrome, and save."
                ),
            },
            {
                "html_title": "Three closing slides",
                "rows": [
                    ["Q&A"],
                    ["Acknowledgements"],
                    ["Thank You"],
                ],
                "instr": (
                    "Three closing-slide titles are listed in the Chrome page. "
                    "Append three new slides at the end of the presentation, "
                    "each with the matching title text. Save."
                ),
            },
            {
                "html_title": "Section headers",
                "rows": [
                    ["Background"],
                    ["Methodology"],
                ],
                "instr": (
                    "The Chrome tab shows two section-header titles. Insert "
                    "two new slides at the end of "
                    "PPT-Template_widescreen.pptx, with each new slide's "
                    "title set to the corresponding text. Save."
                ),
            },
            {
                "html_title": "Risk-and-mitigation slides",
                "rows": [
                    ["Risks"],
                    ["Mitigation Plan"],
                    ["Open Questions"],
                ],
                "instr": (
                    "Three risk-and-mitigation titles are listed in the "
                    "Chrome page. Append three new slides at the end of "
                    "PPT-Template_widescreen.pptx, each titled with the "
                    "matching text. Save."
                ),
            },
        ],
    },
    # 47f7c0ce: NO chrome — Robotic_Workshop_Infographics.pptx (20 slides)
    "osworld_multi_apps_47f7c0ce": {
        "sink_path": "/home/user/Desktop/Robotic_Workshop_Infographics.pptx",
        "sink_ext": "pptx",
        "sink_headers": ["Slide title"],
        "variants": [
            {
                "html_title": "Robotics topics to add",
                "rows": [
                    ["Sensor Fusion Basics"],
                    ["Kinematic Models"],
                ],
                "instr": (
                    "The Chrome tab shows two robotics-topic slide titles. "
                    "Insert two new slides at the end of "
                    "Robotic_Workshop_Infographics.pptx, setting each new "
                    "slide's title to the corresponding text from Chrome. "
                    "Save."
                ),
            },
            {
                "html_title": "Practical sessions",
                "rows": [
                    ["Hands-on: Camera Calibration"],
                    ["Hands-on: Path Planning"],
                    ["Hands-on: Gripper Control"],
                ],
                "instr": (
                    "Three hands-on session titles are listed in the Chrome "
                    "page. Append three new slides at the end of the "
                    "presentation, each with the matching title text. Save."
                ),
            },
            {
                "html_title": "Closing items",
                "rows": [
                    ["Lessons Learned"],
                    ["Next Steps"],
                ],
                "instr": (
                    "The Chrome tab shows two closing-item titles. Insert "
                    "two new slides at the end of "
                    "Robotic_Workshop_Infographics.pptx, each titled with "
                    "the corresponding text. Save."
                ),
            },
            {
                "html_title": "Safety topics",
                "rows": [
                    ["Workspace Safety"],
                    ["Emergency Stop Protocol"],
                ],
                "instr": (
                    "The Chrome tab shows two safety-topic titles. Append "
                    "two new slides at the end of "
                    "Robotic_Workshop_Infographics.pptx, each titled with "
                    "the matching text from Chrome. Save."
                ),
            },
        ],
    },
    # 778efd0a: NO chrome — Minimalist_Business_Slides.pptx (16 slides)
    "osworld_multi_apps_778efd0a": {
        "sink_path": "/home/user/Desktop/Minimalist_Business_Slides.pptx",
        "sink_ext": "pptx",
        "sink_headers": ["Slide title"],
        "variants": [
            {
                "html_title": "Quarter review titles",
                "rows": [
                    ["Q3 Highlights"],
                    ["Q4 Outlook"],
                ],
                "instr": (
                    "The Chrome tab shows two quarter-review slide titles. "
                    "Insert two new slides at the end of "
                    "Minimalist_Business_Slides.pptx, with each new slide's "
                    "title set to the corresponding text. Save."
                ),
            },
            {
                "html_title": "Strategy section",
                "rows": [
                    ["Market Position"],
                    ["Competitive Landscape"],
                    ["Strategic Priorities"],
                ],
                "instr": (
                    "Three strategy-section titles are listed in the Chrome "
                    "page. Append three new slides at the end of the "
                    "presentation, each with the matching title text. Save."
                ),
            },
            {
                "html_title": "Closing slides",
                "rows": [
                    ["Open Discussion"],
                    ["Action Items"],
                ],
                "instr": (
                    "The Chrome tab shows two closing-slide titles. Insert "
                    "two new slides at the end of "
                    "Minimalist_Business_Slides.pptx, each titled with the "
                    "corresponding text. Save."
                ),
            },
            {
                "html_title": "Operations review",
                "rows": [
                    ["Cost Structure"],
                    ["Operational Efficiency"],
                    ["Vendor Partnerships"],
                ],
                "instr": (
                    "Three operations-review titles are shown in the Chrome "
                    "page. Append three new slides at the end of "
                    "Minimalist_Business_Slides.pptx, each titled with the "
                    "matching text. Save."
                ),
            },
        ],
    },
    # bb83cab4: NO chrome — Unlocking-the-Power-of-ChatGPT.pptx (28 slides)
    "osworld_multi_apps_bb83cab4": {
        "sink_path": "/home/user/Desktop/Unlocking-the-Power-of-ChatGPT.pptx",
        "sink_ext": "pptx",
        "sink_headers": ["Slide title"],
        "variants": [
            {
                "html_title": "Advanced prompting topics",
                "rows": [
                    ["Chain-of-Thought Patterns"],
                    ["Tool-Use Prompting"],
                ],
                "instr": (
                    "The Chrome tab shows two advanced-prompting topic "
                    "titles. Insert two new slides at the end of "
                    "Unlocking-the-Power-of-ChatGPT.pptx, with each new "
                    "slide's title set to the corresponding text. Save."
                ),
            },
            {
                "html_title": "Integration patterns",
                "rows": [
                    ["RAG Pipelines"],
                    ["Agentic Loops"],
                    ["Function Calling"],
                ],
                "instr": (
                    "Three integration-pattern titles are listed in the "
                    "Chrome page. Append three new slides at the end of the "
                    "presentation, each with the matching title text. Save."
                ),
            },
            {
                "html_title": "Closing items",
                "rows": [
                    ["Limitations & Risks"],
                    ["Looking Ahead"],
                ],
                "instr": (
                    "The Chrome tab shows two closing-item titles. Insert "
                    "two new slides at the end of "
                    "Unlocking-the-Power-of-ChatGPT.pptx, each titled with "
                    "the corresponding text from Chrome. Save."
                ),
            },
            {
                "html_title": "Use-case demonstrations",
                "rows": [
                    ["Demo: Customer Support"],
                    ["Demo: Code Review Assistant"],
                    ["Demo: Knowledge-Base Search"],
                ],
                "instr": (
                    "Three use-case demonstration titles are listed in the "
                    "Chrome page. Append three new slides at the end of "
                    "Unlocking-the-Power-of-ChatGPT.pptx, each titled with "
                    "the matching text. Save."
                ),
            },
        ],
    },
    # e135df7c: chrome+aclanthology → annual-enterprise-survey-2021... (32×10)
    # Headers verbatim (NZSIOC schema). Test that long header strings copy faithfully.
    "osworld_multi_apps_e135df7c": {
        "sink_path": "/home/user/Desktop/annual-enterprise-survey-2021-financial-year-provisional.xlsx",
        "sink_ext": "xlsx",
        "sink_headers": [
            "Year", "Industry_aggregation_NZSIOC", "Industry_code_NZSIOC",
            "Industry_name_NZSIOC", "Units", "Variable_code", "Variable_name",
            "Variable_category", "Value", "Industry_code_ANZSIC06",
        ],
        "variants": [
            {
                "html_title": "Industry survey rows (additional)",
                "rows": [
                    [2022, "Level 1", "99999", "All Industries", "Dollars (millions)",
                     "H01", "Total income", "Financial performance", 458200, "ANZSIC06 ZZ11"],
                    [2022, "Level 1", "99999", "All Industries", "Dollars (millions)",
                     "H05", "Total expenditure", "Financial performance", 410500, "ANZSIC06 ZZ11"],
                ],
                "instr": (
                    "The Chrome tab shows two additional industry survey "
                    "rows with the full ten-column schema. Append both rows "
                    "to the bottom of the annual-enterprise-survey-2021... "
                    "spreadsheet, preserving the column order (Year, "
                    "Industry_aggregation_NZSIOC, Industry_code_NZSIOC, "
                    "Industry_name_NZSIOC, Units, Variable_code, "
                    "Variable_name, Variable_category, Value, "
                    "Industry_code_ANZSIC06). Save."
                ),
            },
            {
                "html_title": "Manufacturing rows",
                "rows": [
                    [2022, "Level 2", "C", "Manufacturing", "Dollars (millions)",
                     "H01", "Total income", "Financial performance", 81200, "ANZSIC06 C"],
                ],
                "instr": (
                    "One manufacturing-industry row is shown in the Chrome "
                    "page. Add it as a new row at the bottom of the "
                    "annual-enterprise-survey-2021 spreadsheet, preserving "
                    "the ten-column order Year, Industry_aggregation_NZSIOC, "
                    "Industry_code_NZSIOC, Industry_name_NZSIOC, Units, "
                    "Variable_code, Variable_name, Variable_category, Value, "
                    "Industry_code_ANZSIC06. Save."
                ),
            },
            {
                "html_title": "Services rows",
                "rows": [
                    [2022, "Level 2", "K", "Financial and Insurance Services", "Dollars (millions)",
                     "H01", "Total income", "Financial performance", 47200, "ANZSIC06 K"],
                    [2022, "Level 2", "K", "Financial and Insurance Services", "Dollars (millions)",
                     "H05", "Total expenditure", "Financial performance", 28600, "ANZSIC06 K"],
                ],
                "instr": (
                    "Two financial-services rows are listed in the Chrome "
                    "page. Append both as new rows in the spreadsheet, "
                    "keeping the ten-column order. Save."
                ),
            },
        ],
    },
}


# ---------------------------------------------------------------------------
# Tier A2: Single-app cross-format conversions / file-level operations.
# Tier A3: Genuine cross-app coordination beyond the chrome→LO pattern.
#
# These specs are fully self-contained training tasks. Unlike Tier A1, which
# reuses the eval row's downloads/launches and only tweaks chrome_open_tabs,
# Tier A2/A3 STRIP the eval row's config entirely and provide their own
# pre-config setup (mkdir / write source files / launch apps). The eval base's
# task_id is the only thing reused — used as the perturb-row anchor.
#
# Each spec lists 2+ variants (V4d uniqueness) with paraphrased instruction
# pools. Schema:
#   {
#       "archetype": "<short identifier — drives dispatcher branch>",
#       "tier": "a2" | "a3",
#       "variants": [
#           {
#               "params": {...},                   # archetype-specific data
#               "instr_pool": [paraphrase1, ...],  # 3+ paraphrases
#           },
#       ],
#   }
#
# The dispatcher `_perturb_a23` looks up the archetype, builds setup/oracle/
# evaluator from the variant params, picks one paraphrase, and stylizes it via
# the shared `_stylize_a23_instruction` (similar polite-prefix / fluff logic
# tuned for these archetypes — separate from the A1 stylizer).
# ---------------------------------------------------------------------------

# Generic polite prefix list reused across A2/A3 stylization. Mirrors the
# A1 list but kept independent so future tweaks don't accidentally couple.
_A23_POLITE_PREFIXES = list(_A1_POLITE_PREFIXES)

# A2/A3 fluff pool — short, generic, archetype-agnostic. ~38% non-empty
# matches eval baseline.
_A23_FLUFF_SUFFIXES = [
    " Don't change anything else in the source files.",
    " Leave the rest of the project layout intact.",
    " Keep the original content untouched besides the requested change.",
    " The result file should match the expected output exactly.",
    " Write the output to the path I gave above.",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
    "",
]

# A2/A3 multi-step pad. ~3% non-empty matches eval ~3% multi_sep ratio.
_A23_MULTI_SEP_PADS = (
    [". Then, double-check that the new file exists at that path"]
    + [""] * 32
)


def _stylize_a23_instruction(
    instr: str,
    rng: random.Random,
) -> str:
    """Stylize a Tier-A2/A3 paraphrase: polite-prefix injection (~38%) +
    trailing fluff (~38%) + multi-step pad (~3%). Defensive save-strip.
    """
    save_re = re.compile(
        r"\s*(save the (file|spreadsheet|document|deck|archive|image|pdf)|"
        r"save when done|save\.)\s*\.?\s*$",
        re.IGNORECASE,
    )
    while save_re.search(instr):
        instr = save_re.sub(".", instr).rstrip()
        if not instr.endswith("."):
            instr += "."

    # Polite prefix (~38%) — only when paraphrase begins with an action verb
    # in imperative mood (capitalized first word, alpha-only). Avoids
    # double-polite by skipping when instr already starts with one of the
    # polite-prefix forms.
    already_polite = any(
        instr.lower().startswith(p.lower().strip())
        for p in _A23_POLITE_PREFIXES
    )
    first_word = instr.split(" ", 1)[0].rstrip(".,") if instr else ""
    starts_with_verb = bool(first_word) and first_word[0].isupper() and first_word.isalpha()
    if starts_with_verb and not already_polite and rng.random() < 0.38:
        prefix = rng.choice(_A23_POLITE_PREFIXES)
        instr = prefix + instr[0].lower() + instr[1:]

    fluff = rng.choice(_A23_FLUFF_SUFFIXES)
    if fluff:
        if not instr.rstrip().endswith("."):
            instr = instr.rstrip() + "."
        instr = instr + fluff

    multi_pad = rng.choice(_A23_MULTI_SEP_PADS)
    if multi_pad:
        if instr.endswith("."):
            instr = instr[:-1]
        instr = instr + multi_pad + "."

    if not instr.endswith("."):
        instr = instr + "."
    return instr


def _strip_eval_config(eval_row: dict) -> dict:
    """Deep-copy eval_row with metadata.config replaced by an empty list.

    Used by Tier A2/A3 so the perturb row's setup is fully owned by
    `pre_config_steps` — eval base downloads/launches are not reused.
    """
    er = copy.deepcopy(eval_row)
    er["metadata"]["config"] = []
    return er


def _shell_step(cmd: str) -> dict:
    """Build an `execute` step that runs a shell command."""
    return {"type": "execute", "parameters": {"command": cmd, "shell": True}}


def _lo_normalize_cmd(path: str, fmt: str) -> str:
    """Round-trip `path` through `soffice --headless --convert-to <fmt>` so its
    XML structure matches what the evaluator runner produces when it normalizes
    expected files on the eval side. Byte-identical to the helper in
    `perturb/libreoffice_writer.py` / `perturb/libreoffice_impress.py` so the
    LO-normalize asymmetry fix is uniform across domains.
    """
    return (
        f"tmpd=$(mktemp -d) && "
        f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
        f"--convert-to {fmt} --outdir \"$tmpd\" '{path}' 2>/dev/null && "
        f"[ -f \"$tmpd/$(basename '{path}')\" ] && "
        f"cp \"$tmpd/$(basename '{path}')\" '{path}'; "
        f"rm -rf \"$tmpd\"; true"
    )


def _build_oracle_lo(sink: str, expected: str, fmt: str = "docx") -> list[dict]:
    """3-step oracle that fixes the LO-normalize asymmetry bug for docx/pptx
    file-op evaluators (compare_docx_files / compare_docx_strict /
    compare_docx_tables / compare_docx_images / compare_pptx_files):

      ① normalize gold (so result = normalize(normalize(expected)))
      ② plant: cp gold → sink
      ③ normalize sink (so the sink matches what eval-side LO normalize produces)

    Without step ①, fields LO mutates during normalization (e.g. font fallback,
    XML namespace ordering) can leave expected and sink-after-normalize out of
    sync → false-fail on otherwise-passing tasks. Mirror of `_build_oracle` in
    `perturb/libreoffice_writer.py` (docx) and `_standard_oracle` minus the
    expected_py step in `perturb/libreoffice_impress.py` (pptx).
    """
    return [
        _shell_step(_lo_normalize_cmd(expected, fmt)),
        _shell_step(f"cp '{expected}' '{sink}'"),
        _shell_step(_lo_normalize_cmd(sink, fmt)),
    ]


def _write_text_file_step(path: str, content: str) -> dict:
    """Write a text file via base64 (safe across shell quoting)."""
    b64 = base64.b64encode(content.encode()).decode()
    return _shell_step(f"mkdir -p $(dirname {path}) && echo {b64} | base64 -d > {path}")


def _write_binary_via_python_step(path: str, py_code: str) -> dict:
    """Run a python heredoc that writes a binary file at `path`. The caller
    is responsible for putting the write logic in `py_code`.
    """
    return _make_config_step(py_code)


# ---------------------------------------------------------------------------
# Tier A2/A3 archetype builders. Each `_build_<archetype>(short, params)`
# returns (pre_config_steps, perturb_config_step, oracle, evaluator). The
# `short` is the 8-char eval-task suffix used in /tmp paths.
# ---------------------------------------------------------------------------

def _build_a2_xlsx_to_csv(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent opens an xlsx and exports its data as CSV.

    Setup: write a minimal xlsx via openpyxl, launch LO Calc on it.
    Oracle: regenerate expected CSV server-side via python.
    Evaluator: compare_csv on result CSV.
    """
    src_xlsx = f"/tmp/perturb_a2_{short}_src.xlsx"
    result_csv = f"/home/user/Desktop/perturb_a2_{short}.csv"
    expected_csv = f"/tmp/perturb_a2_{short}_expected.csv"

    headers = params["headers"]
    rows = params["rows"]
    delim = params.get("delim", ",")

    pre_py = textwrap.dedent(f"""\
        import openpyxl, os, csv
        os.makedirs(os.path.dirname({src_xlsx!r}), exist_ok=True)
        wb = openpyxl.Workbook()
        ws = wb.worksheets[0]
        ws.append({headers!r})
        for row in {rows!r}:
            ws.append(row)
        wb.save({src_xlsx!r})
        # Pre-build the expected CSV so oracle is a simple cp.
        with open({expected_csv!r}, "w", newline="") as fh:
            w = csv.writer(fh, delimiter={delim!r})
            w.writerow({headers!r})
            for row in {rows!r}:
                w.writerow(row)
        # Make sure result_csv does NOT exist initially.
        try:
            os.remove({result_csv!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["soffice", "--calc", src_xlsx],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_csv}' '{result_csv}'")]
    evaluator = {
        "func": "compare_csv",
        "result": {"type": "vm_file", "path": result_csv,
                   "dest": result_csv.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_csv,
                     "dest": "expected_file"},
        "options": {"strict": False},
    }
    return pre, None, oracle, evaluator


def _build_a2_xlsx_to_pdf(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent opens xlsx, exports/prints to PDF.

    Setup: write xlsx, launch LO Calc.
    Oracle: LO headless --convert-to pdf to produce expected.
    Evaluator: check_pdf_pages (pages == N) — robust to layout drift.
    """
    src_xlsx = f"/tmp/perturb_a2_{short}_src.xlsx"
    result_pdf = f"/home/user/Desktop/perturb_a2_{short}.pdf"
    headers = params["headers"]
    rows = params["rows"]
    expected_pages = params.get("expected_pages", 1)

    pre_py = textwrap.dedent(f"""\
        import openpyxl, os
        os.makedirs(os.path.dirname({src_xlsx!r}), exist_ok=True)
        wb = openpyxl.Workbook()
        ws = wb.worksheets[0]
        ws.append({headers!r})
        for row in {rows!r}:
            ws.append(row)
        wb.save({src_xlsx!r})
        try:
            os.remove({result_pdf!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["soffice", "--calc", src_xlsx],
        }},
    ]
    out_dir = "/home/user/Desktop"
    # LO is already running (--calc launched in pre); use a separate
    # -env:UserInstallation to avoid the singleton-lock conflict that would
    # cause --convert-to to silently no-op.
    oracle_cmd = (
        f"_lotmpd=$(mktemp -d) && _louser=$(mktemp -d) && "
        f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
        f"-env:UserInstallation=\"file://$_louser\" "
        f"--convert-to pdf --outdir \"$_lotmpd\" '{src_xlsx}' >/dev/null 2>&1; "
        f"_pdf=\"$_lotmpd/$(basename '{src_xlsx}' .xlsx).pdf\"; "
        f"[ -f \"$_pdf\" ] && cp \"$_pdf\" '{result_pdf}'; "
        f"rm -rf \"$_lotmpd\" \"$_louser\"; true"
    )
    oracle = [_shell_step(oracle_cmd)]
    # check_pdf_pages expects rules={"relation": <op>, "ref_value": <int>};
    # use ge so 1+ page PDFs pass (LO may render slightly differently across
    # versions but always produces at least `expected_pages` pages here).
    evaluator = {
        "func": "check_pdf_pages",
        "result": {"type": "vm_file", "path": result_pdf,
                   "dest": result_pdf.split("/")[-1]},
        "expected": {"type": "rule", "rules": {
            "relation": "ge", "ref_value": expected_pages,
        }},
    }
    return pre, None, oracle, evaluator


def _build_a2_docx_to_pdf(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent exports a docx as PDF. Mirror of xlsx_to_pdf."""
    src_docx = f"/tmp/perturb_a2_{short}_src.docx"
    result_pdf = f"/home/user/Desktop/perturb_a2_{short}.pdf"
    expected_pdf = f"/tmp/perturb_a2_{short}_expected.pdf"
    paragraphs = params["paragraphs"]

    pre_py = textwrap.dedent(f"""\
        from docx import Document
        import os
        os.makedirs(os.path.dirname({src_docx!r}), exist_ok=True)
        doc = Document()
        for p in {paragraphs!r}:
            doc.add_paragraph(p)
        doc.save({src_docx!r})
        try:
            os.remove({result_pdf!r})
        except FileNotFoundError:
            pass
        # Pre-build expected via LO headless so evaluator has a deterministic
        # reference (same conversion path as oracle).
        import subprocess, shutil, tempfile
        td = tempfile.mkdtemp()
        try:
            r = subprocess.run([
                'soffice', '--headless', '--norestore', '--nofirststartwizard',
                '--convert-to', 'pdf', '--outdir', td, {src_docx!r},
            ], env={{**os.environ, 'DISPLAY': ':1'}}, capture_output=True)
            cand = os.path.join(td, os.path.basename({src_docx!r}).replace('.docx', '.pdf'))
            if os.path.isfile(cand):
                shutil.copy(cand, {expected_pdf!r})
        finally:
            shutil.rmtree(td, ignore_errors=True)
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["soffice", "--writer", src_docx],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_pdf}' '{result_pdf}'")]
    evaluator = {
        "func": "compare_pdfs",
        "result": {"type": "vm_file", "path": result_pdf,
                   "dest": result_pdf.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_pdf,
                     "dest": "expected_file"},
    }
    return pre, None, oracle, evaluator


def _build_a2_image_format_convert(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: convert source image (png) to target format (jpg).

    Setup: PIL writes a small synthetic image as `src_path`; opens GIMP on it.
    Oracle: PIL re-converts source to target format at result path.
    Evaluator: compare_images on result vs gold.
    """
    fmt_in = params.get("fmt_in", "png")
    fmt_out = params.get("fmt_out", "jpg")
    size = params.get("size", (160, 120))
    color = params.get("color", (90, 140, 200))
    src_path = f"/home/user/Desktop/perturb_a2_{short}_src.{fmt_in}"
    result_path = f"/home/user/Desktop/perturb_a2_{short}.{fmt_out}"
    expected_path = f"/tmp/perturb_a2_{short}_expected.{fmt_out}"

    pre_py = textwrap.dedent(f"""\
        from PIL import Image
        import os
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        Image.new('RGB', {tuple(size)!r}, {tuple(color)!r}).save({src_path!r})
        # Pre-build expected (PIL JPEG default quality).
        Image.open({src_path!r}).convert('RGB').save({expected_path!r})
        try:
            os.remove({result_path!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["gimp", src_path],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_path}' '{result_path}'")]
    evaluator = {
        "func": "compare_images",
        "result": {"type": "vm_file", "path": result_path,
                   "dest": result_path.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_path,
                     "dest": "expected_file"},
    }
    return pre, None, oracle, evaluator


def _build_a2_archive_zip(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: zip a small directory of text files.

    Setup: write N text files in /home/user/Desktop/perturb_a2_<short>/.
    Oracle: zip the directory to result_zip; pre-build expected zip too.
    Evaluator: compare_archive (file_type=text).
    """
    files = params["files"]  # {filename: content}
    src_dir = f"/home/user/Desktop/perturb_a2_{short}_dir"
    result_zip = f"/home/user/Desktop/perturb_a2_{short}.zip"
    expected_zip = f"/tmp/perturb_a2_{short}_expected.zip"

    pre_py = textwrap.dedent(f"""\
        import os, zipfile, shutil
        if os.path.isdir({src_dir!r}):
            shutil.rmtree({src_dir!r})
        os.makedirs({src_dir!r}, exist_ok=True)
        for name, content in {list(files.items())!r}:
            with open(os.path.join({src_dir!r}, name), 'w') as fh:
                fh.write(content)
        # Pre-build expected zip.
        if os.path.isfile({expected_zip!r}):
            os.remove({expected_zip!r})
        with zipfile.ZipFile({expected_zip!r}, 'w', zipfile.ZIP_DEFLATED) as zf:
            for name, _ in {list(files.items())!r}:
                zf.write(os.path.join({src_dir!r}, name), arcname=name)
        try:
            os.remove({result_zip!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["nautilus", src_dir],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_zip}' '{result_zip}'")]
    evaluator = {
        "func": "compare_archive",
        "result": {"type": "vm_file", "path": result_zip,
                   "dest": result_zip.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_zip,
                     "dest": "expected_file"},
        "options": {"file_type": "text"},
    }
    return pre, None, oracle, evaluator


def _build_a2_text_file_create(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent writes a small text file (e.g. notes / config).

    Setup: write a "template" file with placeholder lines; launch gedit/text app.
    Oracle: replace with the expected final text.
    Evaluator: compare_text_file.
    """
    template = params["template"]
    expected = params["expected"]
    template_path = f"/home/user/Desktop/perturb_a2_{short}_template.txt"
    result_path = f"/home/user/Desktop/perturb_a2_{short}.txt"
    expected_path = f"/tmp/perturb_a2_{short}_expected.txt"

    pre_py = textwrap.dedent(f"""\
        import os
        os.makedirs(os.path.dirname({template_path!r}), exist_ok=True)
        with open({template_path!r}, 'w') as fh:
            fh.write({template!r})
        with open({expected_path!r}, 'w') as fh:
            fh.write({expected!r})
        try:
            os.remove({result_path!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        # gedit is NOT installed in the OSWorld VM (mousepad is the
        # available GUI text editor; gedit was removed when XFCE replaced
        # GNOME). Launching gedit silently fails (Popen FileNotFoundError
        # → server 500 → dispatch ignores). Use mousepad instead so the
        # agent actually sees an editor window during rollouts.
        {"type": "launch", "parameters": {
            "command": ["mousepad", template_path],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_path}' '{result_path}'")]
    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": result_path,
                   "dest": result_path.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_path,
                     "dest": "expected_file"},
        "options": {"ignore_blanks": True},
    }
    return pre, None, oracle, evaluator


def _build_a2_python_test_suite(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent fixes a small python file. Test suite validates result.

    Setup: write a buggy `module.py` and a test_module.py that imports it and
    runs `test()` returning bool.
    Oracle: overwrite module.py with the fixed version.
    Evaluator: check_python_file_by_test_suite on test_module.py.
    """
    buggy = params["buggy_module"]
    fixed = params["fixed_module"]
    test_code = params["test_code"]
    proj_dir = f"/home/user/Desktop/perturb_a2_{short}_proj"
    module_path = f"{proj_dir}/module.py"
    test_path = f"{proj_dir}/test_module.py"
    fixed_path = f"/tmp/perturb_a2_{short}_fixed.py"

    pre_py = textwrap.dedent(f"""\
        import os, shutil
        if os.path.isdir({proj_dir!r}):
            shutil.rmtree({proj_dir!r})
        os.makedirs({proj_dir!r}, exist_ok=True)
        with open({module_path!r}, 'w') as fh:
            fh.write({buggy!r})
        with open({test_path!r}, 'w') as fh:
            fh.write({test_code!r})
        with open({fixed_path!r}, 'w') as fh:
            fh.write({fixed!r})
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["code", proj_dir],
        }},
    ]
    oracle = [_shell_step(f"cp '{fixed_path}' '{module_path}'")]
    evaluator = {
        "func": "check_python_file_by_test_suite",
        # `check_python_file_by_test_suite(actual_files, test_file)` upstream
        # expects actual_files as a list (multi=True; first entry is the
        # module the test suite imports). test_file is the test module's
        # path (single vm_file, copied to env.cache_dir/test_module.py).
        "result": {"type": "vm_file",
                   "path": [module_path],
                   "dest": ["module.py"],
                   "multi": True},
        "expected": {"type": "vm_file", "path": test_path,
                     "dest": "test_module.py"},
        "options": {"test_function_name": "test"},
    }
    return pre, None, oracle, evaluator


def _build_a2_json_export(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent fills out a JSON config file. check_json validates keys.

    Setup: write a stub JSON with placeholder values; launch text editor.
    Oracle: write the final JSON.
    Evaluator: check_json with expect rules.
    """
    stub = params["stub_json"]
    final = params["final_json"]
    rules = params["rules"]
    stub_path = f"/home/user/Desktop/perturb_a2_{short}_stub.json"
    result_path = f"/home/user/Desktop/perturb_a2_{short}.json"
    expected_path = f"/tmp/perturb_a2_{short}_expected.json"

    pre_py = textwrap.dedent(f"""\
        import os, json
        os.makedirs(os.path.dirname({stub_path!r}), exist_ok=True)
        with open({stub_path!r}, 'w') as fh:
            json.dump({stub!r}, fh, indent=2)
        with open({expected_path!r}, 'w') as fh:
            json.dump({final!r}, fh, indent=2)
        try:
            os.remove({result_path!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        # gedit is NOT installed in the OSWorld VM (mousepad is the
        # available GUI text editor). Use mousepad so the agent sees an
        # actual editor window during rollouts.
        {"type": "launch", "parameters": {
            "command": ["mousepad", stub_path],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_path}' '{result_path}'")]
    evaluator = {
        "func": "check_json",
        "result": {"type": "vm_file", "path": result_path,
                   "dest": result_path.split("/")[-1]},
        "expected": {"type": "rule", "rules": rules},
    }
    return pre, None, oracle, evaluator


def _build_a2_check_include_exclude(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A2: agent removes a file or runs a small shell op; result is a
    text dump (e.g. `ls -R`) checked with check_include_exclude.

    Validation note: the evaluator must NOT
    depend on the oracle creating an intermediate file. Oracle only runs
    in validate.py, never in training rollout (utils/setup.py:42-44 only
    dispatches metadata.config). Earlier the evaluator did
    `cat '<result_log>'` where result_log was created by the oracle's
    `ls -R > result_log` — file never exists during training, eval reads
    empty / errors, all variants of 17 multi_apps A2 bases scored 0
    (2c9fc0de, 510f64c8, 937087b6, 26150609, 2fe4b718, etc.). Now the
    evaluator runs `ls -R` directly on the directory at eval time, which
    is independent of any oracle action. Oracle still kept for validate.py
    parity (so oracle-replay reproduces the same final state).
    """
    files = params["files"]  # {filename: content}
    files_to_remove = params["remove"]  # subset of files
    include = params.get("include", [])
    exclude = params.get("exclude", [])
    src_dir = f"/home/user/Desktop/perturb_a2_{short}_dir"

    pre_py = textwrap.dedent(f"""\
        import os, shutil
        if os.path.isdir({src_dir!r}):
            shutil.rmtree({src_dir!r})
        os.makedirs({src_dir!r}, exist_ok=True)
        for name, content in {list(files.items())!r}:
            with open(os.path.join({src_dir!r}, name), 'w') as fh:
                fh.write(content)
        """)
    # Oracle removes the requested files (kept for validate.py replay parity).
    rm_cmds = " && ".join(
        f"rm -f '{src_dir}/{name}'" for name in files_to_remove
    )
    oracle_cmd = rm_cmds
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["nautilus", src_dir],
        }},
    ]
    oracle = [_shell_step(oracle_cmd)]
    # Run ls -R directly at eval time, independent of any oracle action.
    # Empty output (missing dir) yields no include matches → eval=0 cleanly.
    eval_cmd = f"cd '{src_dir}' 2>/dev/null && ls -R || echo MISSING_DIR"
    evaluator = {
        "func": "check_include_exclude",
        "result": {"type": "vm_command_line",
                   "command": eval_cmd, "shell": True},
        "expected": {"type": "rule", "rules": {
            "include": include, "exclude": exclude,
        }},
    }
    return pre, None, oracle, evaluator


# ---------- Tier A3 archetypes (genuine cross-app) ----------

def _build_a3_chrome_to_writer(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A3: agent reads text from a chrome tab (local HTML) and writes it
    into a docx via LibreOffice Writer.

    This differs from Tier A1's chrome→writer pattern in evaluator (uses
    compare_text_file on a server-extracted .txt rather than compare_docx).
    """
    paragraphs = params["paragraphs"]
    title = params.get("title", "Article")
    sink_docx = f"/home/user/Desktop/perturb_a3_{short}.docx"
    expected_docx = f"/tmp/perturb_a3_{short}_expected.docx"
    html_path = f"/tmp/perturb_a3_{short}.html"

    html_body = "".join(f"<p>{html.escape(p)}</p>" for p in paragraphs)
    html_content = textwrap.dedent(f"""\
        <!DOCTYPE html><html><head><meta charset="utf-8">
        <title>{html.escape(title)}</title>
        <style>body {{ font-family: sans-serif; padding: 32px; font-size: 18px; }}</style>
        </head><body><h2>{html.escape(title)}</h2>{html_body}</body></html>
        """)

    pre_py = textwrap.dedent(f"""\
        from docx import Document
        import os
        os.makedirs(os.path.dirname({sink_docx!r}), exist_ok=True)
        # Empty-ish initial doc.
        doc = Document()
        doc.add_paragraph("")
        doc.save({sink_docx!r})
        # Expected doc has all paragraphs.
        gold = Document()
        for p in {paragraphs!r}:
            gold.add_paragraph(p)
        gold.save({expected_docx!r})
        """)
    pre = [
        _build_html_write_step(html_path, html_content),
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["google-chrome", "--remote-debugging-port=1337"],
        }},
        {"type": "launch", "parameters": {
            "command": ["socat", "tcp-listen:9222,fork", "tcp:localhost:1337"],
        }},
        {"type": "chrome_open_tabs", "parameters": {
            "urls_to_open": [f"file://{html_path}"],
        }},
        {"type": "launch", "parameters": {
            "command": ["soffice", "--writer", sink_docx],
        }},
    ]
    oracle = _build_oracle_lo(sink_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": sink_docx,
                   "dest": sink_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx,
                     "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a3_image_to_archive(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A3: agent processes images (e.g. resize) in GIMP, then zips them
    into an archive via file manager.
    """
    n_images = params.get("n_images", 3)
    out_size = params.get("out_size", (80, 60))
    src_dir = f"/home/user/Desktop/perturb_a3_{short}_imgs"
    result_zip = f"/home/user/Desktop/perturb_a3_{short}.zip"
    expected_zip = f"/tmp/perturb_a3_{short}_expected.zip"

    pre_py = textwrap.dedent(f"""\
        from PIL import Image
        import os, zipfile, shutil
        if os.path.isdir({src_dir!r}):
            shutil.rmtree({src_dir!r})
        os.makedirs({src_dir!r}, exist_ok=True)
        for i in range({n_images!r}):
            Image.new('RGB', (320, 240), (40 + 60 * i, 80, 200 - 30 * i)).save(
                os.path.join({src_dir!r}, f'image_{{i}}.png'))
        # Pre-build expected: resize each + zip.
        td = os.path.join({src_dir!r}, '_resized')
        os.makedirs(td, exist_ok=True)
        names = []
        for i in range({n_images!r}):
            src = os.path.join({src_dir!r}, f'image_{{i}}.png')
            dst = os.path.join(td, f'image_{{i}}.png')
            Image.open(src).resize({tuple(out_size)!r}).save(dst)
            names.append((dst, f'image_{{i}}.png'))
        if os.path.isfile({expected_zip!r}):
            os.remove({expected_zip!r})
        with zipfile.ZipFile({expected_zip!r}, 'w', zipfile.ZIP_DEFLATED) as zf:
            for src, arc in names:
                zf.write(src, arcname=arc)
        # Remove the _resized/ scratch directory so the agent only sees the
        # original images in src_dir (otherwise nautilus would show a stale
        # _resized/ folder that could confuse "zip the resized images" tasks).
        shutil.rmtree(td, ignore_errors=True)
        try:
            os.remove({result_zip!r})
        except FileNotFoundError:
            pass
        """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["gimp", src_dir],
        }},
        {"type": "launch", "parameters": {
            "command": ["nautilus", src_dir],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_zip}' '{result_zip}'")]
    evaluator = {
        "func": "compare_archive",
        "result": {"type": "vm_file", "path": result_zip,
                   "dest": result_zip.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_zip,
                     "dest": "expected_file"},
        "options": {"file_type": "image"},
    }
    return pre, None, oracle, evaluator


def _build_a3_vscode_filemanager(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A3: agent edits files in VS Code, then moves them via file
    manager. Final state inspected by check_include_exclude on `ls -R`.

    Validation note: same fix as a2_check_include_exclude. Run
    ls -RF + cat file contents directly at eval time, NOT via oracle output.
    Oracle never runs in training rollout (only validate.py), so the prior
    `cat '{result_log}'` always read empty/missing → 6 multi_apps A3 bases
    (0c825995, 91190194, 5bc63fb9, d68204bf, etc.) scored 0 unconditionally.
    Oracle still kept for validate.py replay parity.
    """
    files = params["files"]  # {filename: content}
    move_pairs = params["moves"]  # [(src_rel, dst_rel)]
    edits = params.get("edits", {})  # {filename: new_content}
    include = params.get("include", [])
    exclude = params.get("exclude", [])
    proj_dir = f"/home/user/Desktop/perturb_a3_{short}_proj"

    pre_py = textwrap.dedent(f"""\
        import os, shutil
        if os.path.isdir({proj_dir!r}):
            shutil.rmtree({proj_dir!r})
        os.makedirs({proj_dir!r}, exist_ok=True)
        for name, content in {list(files.items())!r}:
            full = os.path.join({proj_dir!r}, name)
            os.makedirs(os.path.dirname(full) or {proj_dir!r}, exist_ok=True)
            with open(full, 'w') as fh:
                fh.write(content)
        """)
    # Oracle: apply edits + moves (kept for validate.py replay parity).
    oracle_cmds = []
    for fname, new in edits.items():
        b64 = base64.b64encode(new.encode()).decode()
        oracle_cmds.append(
            f"echo {b64} | base64 -d > '{proj_dir}/{fname}'"
        )
    for src, dst in move_pairs:
        oracle_cmds.append(
            f"mkdir -p \"$(dirname '{proj_dir}/{dst}')\" && "
            f"mv '{proj_dir}/{src}' '{proj_dir}/{dst}'"
        )
    oracle_cmds.append(
        f"find '{proj_dir}' -mindepth 1 -type d -empty -delete"
    )
    oracle = [_shell_step(" && ".join(oracle_cmds))]
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["code", proj_dir],
        }},
        {"type": "launch", "parameters": {
            "command": ["nautilus", proj_dir],
        }},
    ]
    # Run ls -RF + cat file contents at eval time. -F adds trailing '/'
    # for dir markers (so "archive/" excludes work).
    eval_cmd = (
        f"cd '{proj_dir}' 2>/dev/null && "
        f"ls -RF 2>/dev/null && find . -type f -exec cat {{}} \\; 2>/dev/null"
        f" || echo MISSING_DIR"
    )
    evaluator = {
        "func": "check_include_exclude",
        "result": {"type": "vm_command_line",
                   "command": eval_cmd, "shell": True},
        "expected": {"type": "rule", "rules": {
            "include": include, "exclude": exclude,
        }},
    }
    return pre, None, oracle, evaluator


def _build_a4_pptx_notes_to_docx(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A4: extract speaker notes from real eval pptx → docx.

    Mirrors eval task 51f5801c (Dickinson_Slides.pptx → notes). Notes are
    pre-extracted at generation time and embedded as the canonical
    expected-docx content. Agent reads the same pptx in Impress, copies
    notes into the empty target docx.
    """
    src_url = params["src_url"]
    src_basename = params["src_basename"]
    target_basename = params["target_basename"]
    notes = params["notes"]  # pre-extracted from real eval source
    src_app = params.get("src_app", "--impress")  # --impress for pptx, --calc for xlsx, etc.
    src_path = f"/home/user/Desktop/{src_basename}"
    target_docx = f"/home/user/Desktop/{target_basename}"
    expected_docx = f"/tmp/perturb_a4_{short}_expected.docx"
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        if not os.path.exists({src_path!r}):
            urllib.request.urlretrieve({src_url!r}, {src_path!r})
        from docx import Document
        Document().save({target_docx!r})
        gold = Document()
        for note in {notes!r}:
            gold.add_paragraph(note)
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["soffice", src_app, src_path]}},
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", target_docx]}},
    ]
    oracle = _build_oracle_lo(target_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": target_docx, "dest": target_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a4_xlsx_table_to_docx(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A4: transcribe real eval xlsx table content → docx paragraphs.

    Mirrors eval task 81c425f5 (OSP_Envelope_Price-List → docx table). The
    canonical extracted rows (first N cols × M rows of real xlsx, joined by
    ' | ') are pre-computed and embedded. Agent opens xlsx in Calc, reads
    cells, writes to docx.
    """
    src_url = params["src_url"]
    src_basename = params["src_basename"]
    target_basename = params["target_basename"]
    rows = params["rows"]  # list[list[str]] pre-extracted from real xlsx
    src_path = f"/home/user/Desktop/{src_basename}"
    target_docx = f"/home/user/Desktop/{target_basename}"
    expected_docx = f"/tmp/perturb_a4_{short}_expected.docx"
    expected_paragraphs = [" | ".join(r) for r in rows]
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        if not os.path.exists({src_path!r}):
            urllib.request.urlretrieve({src_url!r}, {src_path!r})
        from docx import Document
        Document().save({target_docx!r})
        gold = Document()
        for p in {expected_paragraphs!r}:
            gold.add_paragraph(p)
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["soffice", "--calc", src_path]}},
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", target_docx]}},
    ]
    oracle = _build_oracle_lo(target_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": target_docx, "dest": target_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a4_xlsx_to_docx_realtable(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A4 (NEW signature: compare_docx_tables): real xlsx → docx with
    real Word table (rows×cols). Oracle inserts a python-docx table with
    each cell from xlsx. Agent's natural action: copy xlsx range → paste
    as table in writer. Skill match for compare_docx_tables eval.
    """
    src_url = params["src_url"]
    src_basename = params["src_basename"]
    target_basename = params["target_basename"]
    rows = params["rows"]  # list[list[str]] pre-extracted
    src_path = f"/home/user/Desktop/{src_basename}"
    target_docx = f"/home/user/Desktop/{target_basename}"
    expected_docx = f"/tmp/perturb_a4_{short}_expected.docx"
    n_rows, n_cols = len(rows), max(len(r) for r in rows) if rows else 0
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        if not os.path.exists({src_path!r}):
            urllib.request.urlretrieve({src_url!r}, {src_path!r})
        from docx import Document
        Document().save({target_docx!r})
        gold = Document()
        rows_data = {rows!r}
        n_r = {n_rows}
        n_c = {n_cols}
        if n_r > 0 and n_c > 0:
            t = gold.add_table(rows=n_r, cols=n_c)
            for i, r in enumerate(rows_data):
                for j, cell in enumerate(r):
                    t.cell(i, j).text = str(cell) if cell is not None else ''
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["soffice", "--calc", src_path]}},
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", target_docx]}},
    ]
    oracle = _build_oracle_lo(target_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_tables",
        "result": {"type": "vm_file", "path": target_docx, "dest": target_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a4_txt_to_xlsx_column(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A4: read items from real txt → write into xlsx column.

    Mirrors eval task d1acdb87 (restaurants.txt → MUST_VISIT.xlsx). Both
    real source files are downloaded; oracle writes pre-extracted names
    into A2:A_N of the xlsx (real workbook structure preserved). Eval
    evaluator: compare_table.
    """
    txt_url = params["txt_url"]
    txt_basename = params["txt_basename"]
    xlsx_url = params["xlsx_url"]
    xlsx_basename = params["xlsx_basename"]
    items = params["restaurant_names"]  # pre-extracted from real txt
    txt_path = f"/home/user/Desktop/{txt_basename}"
    xlsx_path = f"/home/user/Desktop/{xlsx_basename}"
    expected_xlsx = f"/tmp/perturb_a4_{short}_expected.xlsx"
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request, shutil
        os.makedirs(os.path.dirname({xlsx_path!r}), exist_ok=True)
        if not os.path.exists({txt_path!r}):
            urllib.request.urlretrieve({txt_url!r}, {txt_path!r})
        if not os.path.exists({xlsx_path!r}):
            urllib.request.urlretrieve({xlsx_url!r}, {xlsx_path!r})
        # Build expected: copy of original xlsx with items written into A2:A_N.
        shutil.copy({xlsx_path!r}, {expected_xlsx!r})
        from openpyxl import load_workbook
        wb = load_workbook({expected_xlsx!r})
        ws = wb.active
        for i, name in enumerate({items!r}):
            ws.cell(row=2 + i, column=1).value = name
        wb.save({expected_xlsx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["soffice", "--calc", xlsx_path]}},
    ]
    oracle = [_shell_step(f"cp '{expected_xlsx}' '{xlsx_path}'")]
    evaluator = {
        "func": "compare_table",
        "result": {"type": "vm_file", "path": xlsx_path,
                   "dest": xlsx_path.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_xlsx,
                     "dest": "expected_file"},
        "options": {"rules": [{"type": "sheet_data",
                               "sheet_idx0": 0, "sheet_idx1": "EI0"}]},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a6_merge_txt_to_docx(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A6: concat multiple .txt files into one docx in alphabetical order.

    Pattern: agent opens a folder with N .txt files, reads each in
    alphabetical order, writes content as paragraphs in target docx.
    Real eval analogue: 98e8e339 (concat .txt from vscode project) — but
    that's covered as a3_chrome_to_writer. This archetype uses synthetic
    .txt files with realistic content (notes, log entries).
    """
    txt_files = params["txt_files"]  # {filename: content (str)}
    src_dirname = params["src_dirname"]
    target_basename = params["target_basename"]
    src_dir = f"/home/user/Desktop/{src_dirname}"
    target_docx = f"/home/user/Desktop/{target_basename}"
    expected_docx = f"/tmp/perturb_a6_{short}_expected.docx"
    pre_py = textwrap.dedent(f"""\
        import os, shutil
        if os.path.isdir({src_dir!r}):
            shutil.rmtree({src_dir!r})
        os.makedirs({src_dir!r}, exist_ok=True)
        for fname, content in {txt_files!r}.items():
            with open(os.path.join({src_dir!r}, fname), 'w') as fh:
                fh.write(content)
        from docx import Document
        Document().save({target_docx!r})
        # Expected = concat of files in alphabetical filename order, each
        # non-empty line as its own paragraph.
        gold = Document()
        for fname in sorted({list(txt_files)!r}):
            with open(os.path.join({src_dir!r}, fname)) as fh:
                for line in fh.read().splitlines():
                    if line.strip():
                        gold.add_paragraph(line)
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["nautilus", src_dir]}},
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", target_docx]}},
    ]
    oracle = _build_oracle_lo(target_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": target_docx, "dest": target_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a10_image_to_pdf(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A10: convert real eval image → PDF.

    Mirrors eval task a503b07f (OIP.jpg receipt → PDF). Agent opens image
    in GIMP (or any viewer), exports as PDF. Oracle: PIL converts the
    real downloaded image to PDF using same dpi/page-size convention.
    """
    src_url = params["src_url"]
    src_basename = params["src_basename"]
    target_basename = params["target_basename"]
    src_path = f"/home/user/Desktop/{src_basename}"
    result_pdf = f"/home/user/Desktop/{target_basename}"
    expected_pdf = f"/tmp/perturb_a10_{short}_expected.pdf"
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        if not os.path.exists({src_path!r}):
            urllib.request.urlretrieve({src_url!r}, {src_path!r})
        # Validation (2026-05-16): upstream `compare_pdfs(None, valid_pdf)`
        # returns 1.0 because `fitz.open(None)` opens a fresh empty document
        # rather than raising, and `fuzz.ratio('', '')` on image-only PDFs
        # is 100. To force pre-eval non-trivial we MUST have a result file
        # whose text differs from expected. Plant an empty (0-byte) file so
        # `fitz.open(empty_file)` raises (verified: compare_pdfs(empty, valid)
        # returns 0.0 on host).
        try:
            os.remove({result_pdf!r})
        except FileNotFoundError:
            pass
        # Single null byte forces compare_pdfs to raise on result side.
        # Empty file (0 bytes) won't work because `_download_from_container`
        # treats empty `data` as a failure and returns None — then
        # `compare_pdfs(None, valid_pdf)` hits the upstream bug and
        # returns 1.0. A 1-byte file downloads as 1 byte; fitz.open then
        # raises FileDataError → score=0.0 → non-trivial.
        with open({result_pdf!r}, 'wb') as _f:
            _f.write(b'\\0')
        # Build expected via PIL (deterministic image→PDF conversion).
        from PIL import Image
        img = Image.open({src_path!r}).convert('RGB')
        img.save({expected_pdf!r}, 'PDF', resolution=100.0)
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["gimp", src_path]}},
    ]
    oracle = [_shell_step(f"cp '{expected_pdf}' '{result_pdf}'")]
    evaluator = {
        "func": "compare_pdfs",
        "result": {"type": "vm_file", "path": result_pdf, "dest": result_pdf.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_pdf, "dest": "expected_file"},
    }
    return pre, None, oracle, evaluator


def _build_a8_modify_section(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A8: modify-in-place specific section of a docx by heading anchor.

    Pattern: source docx has multiple sections each with a heading
    paragraph followed by body paragraphs. Agent edits a SPECIFIC section
    (identified by heading text) — replacing its body content with new
    paragraphs. Heading anchor avoids ordinal ambiguity (per writer
    perturb's _ORDINAL_SAFE_BASES lessons).
    """
    sections = params["sections"]  # list[(heading: str, body_paragraphs: list[str])]
    section_to_modify = params["section_heading"]  # exact heading text
    new_body = params["new_body"]  # replacement body paragraphs
    target_basename = params["target_basename"]
    target_docx = f"/home/user/Desktop/{target_basename}"
    expected_docx = f"/tmp/perturb_a8_{short}_expected.docx"
    pre_py = textwrap.dedent(f"""\
        from docx import Document
        import os
        os.makedirs(os.path.dirname({target_docx!r}), exist_ok=True)
        # Source: write all sections as Heading 2 + body paragraphs.
        doc = Document()
        for heading, body in {sections!r}:
            p = doc.add_paragraph(heading)
            p.style = doc.styles['Heading 2']
            for body_p in body:
                doc.add_paragraph(body_p)
        doc.save({target_docx!r})
        # Expected: same sections but `section_heading` body replaced.
        gold = Document()
        for heading, body in {sections!r}:
            p = gold.add_paragraph(heading)
            p.style = gold.styles['Heading 2']
            replacement = {new_body!r} if heading == {section_to_modify!r} else body
            for body_p in replacement:
                gold.add_paragraph(body_p)
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", target_docx]}},
    ]
    oracle = _build_oracle_lo(target_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": target_docx, "dest": target_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a10_batch_jpg_to_png(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A10: download N real jpg images, agent converts each to png.

    Mirrors eval task ce2b64a2 (3 real mountain jpgs). Original eval task
    asks for image identification (vision/external knowledge → OUT-OF-SCOPE).
    Perturb simplifies to format conversion: agent converts each jpg → png
    in the same Pictures folder. Oracle uses PIL for deterministic conversion.
    """
    src_urls = params["src_urls"]  # list[(url, basename)]
    pic_dirname = params["pic_dirname"]
    pic_dir = f"/home/user/Desktop/{pic_dirname}"
    expected_dir = f"/tmp/perturb_a10_{short}_expected"
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request, shutil
        os.makedirs({pic_dir!r}, exist_ok=True)
        for url, name in {src_urls!r}:
            dst = os.path.join({pic_dir!r}, name)
            if not os.path.exists(dst):
                urllib.request.urlretrieve(url, dst)
        # Build expected: same names but .png extension, deterministic PIL convert.
        if os.path.isdir({expected_dir!r}):
            shutil.rmtree({expected_dir!r})
        os.makedirs({expected_dir!r}, exist_ok=True)
        from PIL import Image
        for url, name in {src_urls!r}:
            jpg_path = os.path.join({pic_dir!r}, name)
            png_name = os.path.splitext(name)[0] + '.png'
            Image.open(jpg_path).convert('RGBA').save(
                os.path.join({expected_dir!r}, png_name), 'PNG'
            )
    """)
    # Oracle: copy each expected png into pic_dir.
    oracle_cmd = (
        f"for f in {expected_dir}/*.png; do "
        f"  cp \"$f\" {pic_dir}/; "
        f"done; true"
    )
    # Evaluator: check that all expected png files exist with matching content.
    # Use check_include_exclude on `ls` since per-file image hash is fragile.
    expected_pngs = [name.rsplit(".", 1)[0] + ".png" for _, name in src_urls]
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["nautilus", pic_dir]}},
    ]
    oracle = [_shell_step(oracle_cmd)]
    # Validation note: same oracle-dependency fix as a2/a3; eval
    # must directly compute `ls` at eval time, not depend on oracle creating
    # an intermediate file (oracle never runs in training rollout).
    eval_cmd = f"ls '{pic_dir}' 2>/dev/null || echo MISSING_DIR"
    evaluator = {
        "func": "check_include_exclude",
        "result": {"type": "vm_command_line",
                   "command": eval_cmd, "shell": True},
        "expected": {"type": "rule", "rules": {
            "include": expected_pngs, "exclude": [],
        }},
    }
    return pre, None, oracle, evaluator


def _build_a10_xcf_to_docx_image(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A10: open real .xcf in GIMP, export as png, insert into docx.

    Mirrors eval task 227d2f97 (xcf on Desktop → image.docx with the image).
    Oracle: download xcf, run GIMP batch-mode to export to png, then build
    expected.docx with that image inserted via python-docx.
    """
    src_url = params["src_url"]
    src_basename = params["src_basename"]
    target_basename = params["target_basename"]
    src_path = f"/home/user/Desktop/{src_basename}"
    target_docx = f"/home/user/Desktop/{target_basename}"
    expected_docx = f"/tmp/perturb_a10_{short}_expected.docx"
    intermediate_png = f"/tmp/perturb_a10_{short}_image.png"
    pre_py = textwrap.dedent(f"""\
        import os, urllib.request, subprocess
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        if not os.path.exists({src_path!r}):
            urllib.request.urlretrieve({src_url!r}, {src_path!r})
        # GIMP batch-mode: flatten xcf and save as png.
        try:
            os.remove({intermediate_png!r})
        except FileNotFoundError:
            pass
        subprocess.run([
            'gimp', '-i', '-b',
            f"(let* ((image (car (gimp-file-load RUN-NONINTERACTIVE \\\"{src_path}\\\" \\\"{src_basename}\\\")))) (gimp-image-flatten image) (file-png-save RUN-NONINTERACTIVE image (car (gimp-image-get-active-drawable image)) \\\"{intermediate_png}\\\" \\\"png\\\" 0 9 1 1 1 1 1))",
            '-b', '(gimp-quit 0)',
        ], capture_output=True, timeout=60)
        from docx import Document
        from docx.shared import Inches
        # Empty target docx for agent to fill.
        Document().save({target_docx!r})
        # Expected: docx with the png inserted (matches what agent's
        # GIMP-export + Writer-Insert-Image flow produces).
        gold = Document()
        if os.path.exists({intermediate_png!r}):
            gold.add_picture({intermediate_png!r}, width=Inches(5))
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["gimp", src_path]}},
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", target_docx]}},
    ]
    oracle = _build_oracle_lo(target_docx, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_images",
        "result": {"type": "vm_file", "path": target_docx, "dest": target_docx.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


def _build_a4_real_docx_to_pdf(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A4 (P1): real-source docx → PDF export.

    Unlike `a2_docx_to_pdf` (which generates a synthetic source docx),
    this archetype downloads the SAME source docx as the eval task. The
    expected.pdf is built by running LO headless conversion on the real
    docx — agent's PDF must match this conversion's output. Drops the
    eval task's Google Drive upload step (out-of-scope per plan).

    Targets eval bases like 22a4636f (Meeting-Agenda.docx),
    897e3b53 (form.docx) where the original instruction was
    docx→pdf+upload-to-Drive but only the docx→pdf part is reproducible.
    """
    src_url = params["src_url"]
    src_basename = params["src_basename"]
    target_basename = params["target_basename"]
    src_path = f"/home/user/Desktop/{src_basename}"
    result_pdf = f"/home/user/Desktop/{target_basename}"
    expected_pdf = f"/tmp/perturb_a4_{short}_expected.pdf"

    pre_py = textwrap.dedent(f"""\
        import os, subprocess, shutil, tempfile, urllib.request
        os.makedirs(os.path.dirname({src_path!r}), exist_ok=True)
        if not os.path.exists({src_path!r}):
            urllib.request.urlretrieve({src_url!r}, {src_path!r})
        try:
            os.remove({result_pdf!r})
        except FileNotFoundError:
            pass
        # Pre-build expected via LO headless conversion of the real source —
        # same conversion path the agent's File→Export As PDF would invoke,
        # so compare_pdfs against this reference is deterministic.
        td = tempfile.mkdtemp()
        try:
            subprocess.run([
                'soffice', '--headless', '--norestore', '--nofirststartwizard',
                '--convert-to', 'pdf', '--outdir', td, {src_path!r},
            ], env={{**os.environ, 'DISPLAY': ':1'}}, capture_output=True, timeout=30)
            cand = os.path.join(td, os.path.basename({src_path!r}).rsplit('.', 1)[0] + '.pdf')
            if os.path.isfile(cand):
                shutil.copy(cand, {expected_pdf!r})
        finally:
            shutil.rmtree(td, ignore_errors=True)
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["soffice", "--writer", src_path],
        }},
    ]
    oracle = [_shell_step(f"cp '{expected_pdf}' '{result_pdf}'")]
    evaluator = {
        "func": "compare_pdfs",
        "result": {"type": "vm_file", "path": result_pdf,
                   "dest": result_pdf.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_pdf,
                     "dest": "expected_file"},
    }
    return pre, None, oracle, evaluator


def _build_a23_xlsx_to_html_chrome_view(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A23: agent converts xlsx → html via LO Calc and opens it in Chrome.

    Mirrors eval e135df7c structure (cross-app: Calc + Chrome) with the same
    compound evaluator [is_expected_tabs, compare_htmls]. Each variant changes
    only the output filename and decoy-tab set so the cross-app skill is
    preserved while the reward target shifts.

    Params:
        xlsx_url: HF URL for source xlsx
        xlsx_path: VM path
        out_basename: html filename to produce (without dir)
        decoy_tabs: 3 pre-opened decoy tabs
    """
    xlsx_url = params["xlsx_url"]
    xlsx_path = params["xlsx_path"]
    out_basename = params["out_basename"]
    decoy_tabs = params["decoy_tabs"]
    out_dir = params.get("out_dir", "/home/user/Desktop")
    out_html = f"{out_dir}/{out_basename}"
    expected_html = f"/tmp/perturb_a23_{short}_expected.html"

    pre_py = textwrap.dedent(f"""\
        import os, urllib.request
        os.makedirs(os.path.dirname({xlsx_path!r}), exist_ok=True)
        if not os.path.exists({xlsx_path!r}):
            urllib.request.urlretrieve({xlsx_url!r}, {xlsx_path!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {
            "command": ["google-chrome", "--remote-debugging-port=1337"],
        }},
        {"type": "launch", "parameters": {
            "command": ["socat", "tcp-listen:9222,fork", "tcp:localhost:1337"],
        }},
        {"type": "chrome_open_tabs", "parameters": {"urls_to_open": decoy_tabs}},
        {"type": "open", "parameters": {"path": xlsx_path}},
    ]

    # Oracle: soffice --convert-to html (creates expected.html), cp to result,
    # then launch chrome with file:// URL so it becomes a tab. Also clear chrome
    # session caches like eval f8cfa149 pattern.
    oracle = [
        _shell_step(
            f"_lotmpd=$(mktemp -d) && "
            f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
            f"--convert-to html --outdir \"$_lotmpd\" '{xlsx_path}' 2>/dev/null && "
            f"_genhtml=\"$_lotmpd/$(basename '{xlsx_path}' .xlsx).html\" && "
            f"[ -f \"$_genhtml\" ] && cp \"$_genhtml\" '{expected_html}' && "
            f"cp \"$_genhtml\" '{out_html}'; "
            f"rm -rf \"$_lotmpd\"; true"
        ),
        _shell_step('rm -f "/home/user/chrome-data/Default/Last Session" '
                    '"/home/user/chrome-data/Default/Last Tabs" 2>/dev/null; true'),
        {"type": "launch", "parameters": {
            "command": ["google-chrome", "--remote-debugging-port=1337", f"file://{out_html}"],
        }},
        {"type": "sleep", "parameters": {"seconds": 5}},
    ]

    expected_urls = decoy_tabs + [f"file://{out_html}"]
    evaluator = {
        "func": ["is_expected_tabs", "compare_htmls"],
        "result": [
            {"type": "open_tabs_info"},
            {"type": "vm_file", "path": out_html, "dest": out_basename},
        ],
        "expected": [
            {"type": "rule", "rules": {"type": "url", "urls": expected_urls}},
            {"type": "vm_file", "path": expected_html, "dest": "expected_file"},
        ],
    }
    return pre, None, oracle, evaluator


def _build_a23_xlsx_row_to_docx_table(short: str, params: dict) -> tuple[list[dict], dict | None, list[dict], dict]:
    """Tier A23: agent extracts specific row(s)/columns from an xlsx and inserts
    them as a table in a pre-existing docx (cross-app calc + writer).

    Mirrors eval 00fa164e structure (extract GPT-4 row from expe-results.xlsx,
    insert as table in awe_desk_env.docx Main Results section). Each variant
    selects a different (rows, cols) subset so agent must navigate Calc to
    different cells and re-paste into Writer.

    Params:
        xlsx_url: HF URL for source xlsx (verbatim from eval)
        xlsx_path: VM path (e.g. "/home/user/Documents/awesome-desktop/expe-results.xlsx")
        docx_url: HF URL for pre-existing docx
        docx_path: VM path (e.g. "/home/user/Documents/awesome-desktop/awe_desk_env.docx")
        extracted_rows: list[list[str]] — the table to insert (pre-extracted)
        extracted_label: human-readable label for instruction (e.g. "GPT-4 row")
    """
    xlsx_url = params["xlsx_url"]
    xlsx_path = params["xlsx_path"]
    docx_url = params["docx_url"]
    docx_path = params["docx_path"]
    extracted_rows = params["extracted_rows"]
    expected_docx = f"/tmp/perturb_a23_{short}_expected.docx"
    n_rows, n_cols = len(extracted_rows), max(len(r) for r in extracted_rows) if extracted_rows else 0

    pre_py = textwrap.dedent(f"""\
        import os, urllib.request
        os.makedirs(os.path.dirname({xlsx_path!r}), exist_ok=True)
        os.makedirs(os.path.dirname({docx_path!r}), exist_ok=True)
        if not os.path.exists({xlsx_path!r}):
            urllib.request.urlretrieve({xlsx_url!r}, {xlsx_path!r})
        if not os.path.exists({docx_path!r}):
            urllib.request.urlretrieve({docx_url!r}, {docx_path!r})
        # Build expected docx: load source docx + append table at end.
        from docx import Document
        gold = Document({docx_path!r})
        rows_data = {extracted_rows!r}
        n_r = {n_rows}; n_c = {n_cols}
        if n_r > 0 and n_c > 0:
            t = gold.add_table(rows=n_r, cols=n_c)
            for i, r in enumerate(rows_data):
                for j, cell in enumerate(r):
                    t.cell(i, j).text = str(cell) if cell is not None else ''
        gold.save({expected_docx!r})
    """)
    pre = [
        _make_config_step(pre_py),
        {"type": "launch", "parameters": {"command": ["soffice", "--calc", xlsx_path]}},
        {"type": "launch", "parameters": {"command": ["soffice", "--writer", docx_path]}},
    ]
    oracle = _build_oracle_lo(docx_path, expected_docx, fmt="docx")
    evaluator = {
        "func": "compare_docx_tables",
        "result": {"type": "vm_file", "path": docx_path, "dest": docx_path.split("/")[-1]},
        "expected": {"type": "vm_file", "path": expected_docx, "dest": "expected_file"},
        "postconfig": LO_SAVE_POSTCONFIG,
    }
    return pre, None, oracle, evaluator


# Map archetype name -> builder function.
_A23_BUILDERS = {
    "a2_xlsx_to_csv": _build_a2_xlsx_to_csv,
    "a2_xlsx_to_pdf": _build_a2_xlsx_to_pdf,
    "a2_docx_to_pdf": _build_a2_docx_to_pdf,
    "a2_image_format_convert": _build_a2_image_format_convert,
    "a2_archive_zip": _build_a2_archive_zip,
    "a2_text_file_create": _build_a2_text_file_create,
    "a2_python_test_suite": _build_a2_python_test_suite,
    "a2_json_export": _build_a2_json_export,
    "a2_check_include_exclude": _build_a2_check_include_exclude,
    "a3_chrome_to_writer": _build_a3_chrome_to_writer,
    "a3_image_to_archive": _build_a3_image_to_archive,
    "a3_vscode_filemanager": _build_a3_vscode_filemanager,
    "a4_real_docx_to_pdf": _build_a4_real_docx_to_pdf,
    "a4_pptx_notes_to_docx": _build_a4_pptx_notes_to_docx,
    "a4_xlsx_table_to_docx": _build_a4_xlsx_table_to_docx,
    "a6_merge_txt_to_docx": _build_a6_merge_txt_to_docx,
    "a10_image_to_pdf": _build_a10_image_to_pdf,
    "a8_modify_section": _build_a8_modify_section,
    "a10_xcf_to_docx_image": _build_a10_xcf_to_docx_image,
    "a4_xlsx_to_docx_realtable": _build_a4_xlsx_to_docx_realtable,
    "a10_batch_jpg_to_png": _build_a10_batch_jpg_to_png,
    "a4_txt_to_xlsx_column": _build_a4_txt_to_xlsx_column,
    "a23_xlsx_row_to_docx_table": _build_a23_xlsx_row_to_docx_table,
    "a23_xlsx_to_html_chrome_view": _build_a23_xlsx_to_html_chrome_view,
}


# Per-task Tier A2/A3 specs, keyed by eval task_id. Each spec specifies an
# archetype (one of `_A23_BUILDERS`) and 2+ variants. Each variant is
# `{"params": {...}, "instr_pool": [3 paraphrases]}`. The dispatcher passes
# `params` to the archetype builder and selects + stylizes one paraphrase.
#
# Coverage rationale (uncovered eval bases targeted):
#   A2 archive  : 5df7b33a, 0e5303d4, f7dfbef3 (compress / extract / format)
#   A2 csv      : c867c42d, ee9a3c83
#   A2 pdf      : 337d318b, b52b40a5
#   A2 image    : 2fe4b718, c2751594
#   A2 text     : acb0f96b, f918266a
#   A2 python   : 26150609, 9219480b
#   A2 json     : e2392362
#   A2 incl/excl: 2c9fc0de, 937087b6, 510f64c8
#   A3 chrome→writer: 0c825995, 98e8e339, aad10cd7
#   A3 image→archive: 46407397, 82e3c869
#   A3 vscode→files : 91190194, d68204bf, 5bc63fb9
_TIER_A23_TASKS: dict[str, dict] = {
    # NOTE: 00fa164e (xlsx-row → docx-table cross-app) was drafted as
    # _build_a23_xlsx_row_to_docx_table but the base is already registered
    # in _TIER_A1_TASKS (chrome HTML → docx) which takes dispatch priority.
    # The A23 builder remains in _A23_BUILDERS for future bases that have
    # an xlsx + pre-existing docx WITHOUT chrome.

    # NOTE: e135df7c (xlsx → html → Chrome) drafted as
    # _build_a23_xlsx_to_html_chrome_view but the base is registered in
    # _TIER_A1_TASKS too (chrome HTML → xlsx). Builder kept in _A23_BUILDERS.

    # ---------- Tier A2: archive ----------
    "osworld_multi_apps_5df7b33a": {
        "archetype": "a2_archive_zip",
        "tier": "a2",
        "variants": [
            {
                "params": {"files": {
                    "chapter_1.txt": "Chapter 1: Spectral graph foundations.\nDefinitions and notation.\n",
                    "chapter_2.txt": "Chapter 2: The Laplacian eigenvalues.\nKey lemmas.\n",
                    "chapter_3.txt": "Chapter 3: Random walks on graphs.\nMixing time bounds.\n",
                }},
                "instr_pool": [
                    "Bundle the three chapter text files in the perturb_a2_5df7b33a_dir folder on the Desktop into a single zip archive saved as perturb_a2_5df7b33a.zip on the Desktop, preserving the original filenames inside the archive.",
                    "Compress chapter_1.txt, chapter_2.txt, and chapter_3.txt from the perturb_a2_5df7b33a_dir folder into one zip file named perturb_a2_5df7b33a.zip on the Desktop, keeping each file at the archive root.",
                    "Create a zip archive perturb_a2_5df7b33a.zip on the Desktop that contains the three chapter txt files from the perturb_a2_5df7b33a_dir folder, using the same filenames as inside the source folder.",
                ],
            },
            {
                "params": {"files": {
                    "intro.txt": "Welcome to the workshop.\n",
                    "agenda.txt": "Day 1: setup\nDay 2: practice\n",
                }},
                "instr_pool": [
                    "Zip intro.txt and agenda.txt from the perturb_a2_5df7b33a_dir folder into perturb_a2_5df7b33a.zip on the Desktop with each file at the archive root.",
                    "Pack both txt files in perturb_a2_5df7b33a_dir into a single archive perturb_a2_5df7b33a.zip on the Desktop, keeping the filenames intact.",
                    "Compress intro.txt and agenda.txt found in the perturb_a2_5df7b33a_dir folder into perturb_a2_5df7b33a.zip on the Desktop.",
                ],
            },
        ],
    },
    "osworld_multi_apps_0e5303d4": {
        "archetype": "a2_archive_zip",
        "tier": "a2",
        "variants": [
            {
                "params": {"files": {
                    "week_01.txt": "Lecture week 1: introduction to python\n",
                    "week_02.txt": "Lecture week 2: control flow\n",
                    "week_03.txt": "Lecture week 3: collections\n",
                    "week_04.txt": "Lecture week 4: functions\n",
                }},
                "instr_pool": [
                    "Bundle the four lecture text files in the perturb_a2_0e5303d4_dir folder on the Desktop into a zip archive named perturb_a2_0e5303d4.zip on the Desktop, keeping the filenames as they are.",
                    "Combine week_01.txt through week_04.txt from the perturb_a2_0e5303d4_dir folder into one zip file perturb_a2_0e5303d4.zip on the Desktop, preserving the names of each file.",
                    "Create a zip archive perturb_a2_0e5303d4.zip on the Desktop containing the four week txt files from perturb_a2_0e5303d4_dir, with each entry at the archive root.",
                ],
            },
            {
                "params": {"files": {
                    "notes_a.txt": "Lecture A notes.\n",
                    "notes_b.txt": "Lecture B notes.\n",
                }},
                "instr_pool": [
                    "Zip the two notes files from perturb_a2_0e5303d4_dir into perturb_a2_0e5303d4.zip on the Desktop, keeping the filenames at the archive root.",
                    "Pack notes_a.txt and notes_b.txt from the perturb_a2_0e5303d4_dir folder into one archive perturb_a2_0e5303d4.zip on the Desktop.",
                    "Combine both notes txt files in perturb_a2_0e5303d4_dir into a single zip file perturb_a2_0e5303d4.zip on the Desktop with the names unchanged.",
                ],
            },
        ],
    },
    "osworld_multi_apps_f7dfbef3": {
        "archetype": "a2_archive_zip",
        "tier": "a2",
        "variants": [
            {
                "params": {"files": {
                    "doc_one.txt": "First document body text.\n",
                    "doc_two.txt": "Second document body text.\n",
                    "doc_three.txt": "Third document body text.\n",
                }},
                "instr_pool": [
                    "Compress all three doc_*.txt files from perturb_a2_f7dfbef3_dir into one archive perturb_a2_f7dfbef3.zip on the Desktop, keeping the names intact.",
                    "Bundle doc_one.txt, doc_two.txt, and doc_three.txt from the perturb_a2_f7dfbef3_dir folder into perturb_a2_f7dfbef3.zip on the Desktop, with each file at the archive root.",
                    "Create the zip file perturb_a2_f7dfbef3.zip on the Desktop containing the three doc txt files from perturb_a2_f7dfbef3_dir, using the same filenames inside the archive.",
                ],
            },
            {
                "params": {"files": {
                    "memo.txt": "Internal memo body.\n",
                    "letter.txt": "Cover letter body.\n",
                }},
                "instr_pool": [
                    "Pack memo.txt and letter.txt from perturb_a2_f7dfbef3_dir into a single zip archive perturb_a2_f7dfbef3.zip on the Desktop.",
                    "Bundle the two text files in perturb_a2_f7dfbef3_dir into perturb_a2_f7dfbef3.zip on the Desktop, with the original filenames.",
                    "Compress memo.txt and letter.txt from the perturb_a2_f7dfbef3_dir folder into perturb_a2_f7dfbef3.zip on the Desktop.",
                ],
            },
        ],
    },
    # ---------- Tier A2: csv export ----------
    "osworld_multi_apps_c867c42d": {
        "archetype": "a2_xlsx_to_csv",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "headers": ["Name", "Email", "Phone"],
                    "rows": [
                        ["Alice Park", "alice@example.com", "+1-555-0101"],
                        ["Bruno Diaz", "bruno@example.com", "+1-555-0102"],
                        ["Caleb Lin", "caleb@example.com", "+1-555-0103"],
                    ],
                },
                "instr_pool": [
                    "Open the perturb_a2_c867c42d_src.xlsx contact list in LibreOffice Calc and export it as a CSV file at /home/user/Desktop/perturb_a2_c867c42d.csv, preserving the column order Name, Email, Phone with a comma separator.",
                    "Export the perturb_a2_c867c42d_src.xlsx contact list to a comma-separated CSV file at /home/user/Desktop/perturb_a2_c867c42d.csv, keeping the same Name / Email / Phone column order.",
                    "Convert the perturb_a2_c867c42d_src.xlsx contacts into a CSV file at /home/user/Desktop/perturb_a2_c867c42d.csv using a comma as the field delimiter and the same column order as the spreadsheet.",
                ],
            },
            {
                "params": {
                    "headers": ["product", "qty", "price"],
                    "rows": [
                        ["bolt", 100, 0.12],
                        ["nut", 200, 0.05],
                        ["washer", 500, 0.02],
                    ],
                },
                "instr_pool": [
                    "Export the perturb_a2_c867c42d_src.xlsx inventory table to /home/user/Desktop/perturb_a2_c867c42d.csv as comma-separated values, preserving the header row product, qty, price.",
                    "Convert the perturb_a2_c867c42d_src.xlsx parts list to a CSV file at /home/user/Desktop/perturb_a2_c867c42d.csv with a comma delimiter and the original column order.",
                    "Save the perturb_a2_c867c42d_src.xlsx inventory as /home/user/Desktop/perturb_a2_c867c42d.csv using comma separators and keeping the product, qty, price columns in the same order.",
                ],
            },
        ],
    },
    "osworld_multi_apps_ee9a3c83": {
        "archetype": "a2_xlsx_to_csv",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "headers": ["city", "country", "population"],
                    "rows": [
                        ["Helsinki", "Finland", 658864],
                        ["Oslo", "Norway", 697549],
                        ["Reykjavik", "Iceland", 131136],
                    ],
                },
                "instr_pool": [
                    "Convert the perturb_a2_ee9a3c83_src.xlsx city table into a CSV file at /home/user/Desktop/perturb_a2_ee9a3c83.csv with a comma separator and the same city, country, population column order.",
                    "Export the perturb_a2_ee9a3c83_src.xlsx data to /home/user/Desktop/perturb_a2_ee9a3c83.csv as comma-separated values, preserving the header row.",
                    "Save the perturb_a2_ee9a3c83_src.xlsx city list as a CSV file at /home/user/Desktop/perturb_a2_ee9a3c83.csv keeping the columns in the order city, country, population.",
                ],
            },
            {
                "params": {
                    "headers": ["item", "category", "stock"],
                    "rows": [
                        ["apple", "fruit", 24],
                        ["broccoli", "vegetable", 15],
                    ],
                },
                "instr_pool": [
                    "Export the perturb_a2_ee9a3c83_src.xlsx grocery sheet to /home/user/Desktop/perturb_a2_ee9a3c83.csv as a comma-separated CSV in the same column order item, category, stock.",
                    "Convert the perturb_a2_ee9a3c83_src.xlsx data into a CSV file at /home/user/Desktop/perturb_a2_ee9a3c83.csv with comma delimiters and the original header row.",
                    "Save the perturb_a2_ee9a3c83_src.xlsx grocery table as /home/user/Desktop/perturb_a2_ee9a3c83.csv preserving the item, category, stock column order.",
                ],
            },
        ],
    },
    # ---------- Tier A2: pdf export ----------
    "osworld_multi_apps_337d318b": {
        "archetype": "a2_xlsx_to_pdf",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "headers": ["Invoice", "Vendor", "Amount", "Status"],
                    "rows": [
                        ["INV-100", "Acme", 120.0, "paid"],
                        ["INV-101", "Beta", 240.0, "open"],
                        ["INV-102", "Gamma", 360.0, "paid"],
                    ],
                    "expected_pages": 1,
                },
                "instr_pool": [
                    "Open the perturb_a2_337d318b_src.xlsx invoice sheet in LibreOffice Calc and export it as a single-page PDF at /home/user/Desktop/perturb_a2_337d318b.pdf.",
                    "Export the perturb_a2_337d318b_src.xlsx workbook to a one-page PDF saved at /home/user/Desktop/perturb_a2_337d318b.pdf using LibreOffice Calc's export feature.",
                    "Convert the perturb_a2_337d318b_src.xlsx file into a PDF document at /home/user/Desktop/perturb_a2_337d318b.pdf, fitting the table on a single page.",
                ],
            },
            {
                "params": {
                    "headers": ["client", "balance"],
                    "rows": [["Olivia", 50.5], ["Pavel", 120.0]],
                    "expected_pages": 1,
                },
                "instr_pool": [
                    "Export the perturb_a2_337d318b_src.xlsx balance sheet to /home/user/Desktop/perturb_a2_337d318b.pdf as a one-page PDF.",
                    "Save the perturb_a2_337d318b_src.xlsx data as a single-page PDF at /home/user/Desktop/perturb_a2_337d318b.pdf.",
                    "Print the perturb_a2_337d318b_src.xlsx contents to a one-page PDF file at /home/user/Desktop/perturb_a2_337d318b.pdf.",
                ],
            },
        ],
    },
    "osworld_multi_apps_b52b40a5": {
        "archetype": "a2_docx_to_pdf",
        "tier": "a2",
        "variants": [
            {
                "params": {"paragraphs": [
                    "Paper Recommendation Summary",
                    "We have reviewed three recent papers on graph neural networks.",
                    "All three papers are attached as separate PDFs in the source bundle.",
                    "Please archive this summary along with the source PDFs.",
                ]},
                "instr_pool": [
                    "Open the perturb_a2_b52b40a5_src.docx document in LibreOffice Writer and export the text content to a PDF file at /home/user/Desktop/perturb_a2_b52b40a5.pdf.",
                    "Export the perturb_a2_b52b40a5_src.docx writer document to a PDF saved at /home/user/Desktop/perturb_a2_b52b40a5.pdf, keeping the paragraph order intact.",
                    "Convert the perturb_a2_b52b40a5_src.docx file into a PDF document at /home/user/Desktop/perturb_a2_b52b40a5.pdf using LibreOffice Writer's export feature.",
                ],
            },
            {
                "params": {"paragraphs": [
                    "Workshop Schedule",
                    "Day 1 covers fundamentals.",
                    "Day 2 focuses on hands-on exercises.",
                ]},
                "instr_pool": [
                    "Save the perturb_a2_b52b40a5_src.docx schedule as a PDF at /home/user/Desktop/perturb_a2_b52b40a5.pdf using LibreOffice Writer.",
                    "Export the perturb_a2_b52b40a5_src.docx workshop document to /home/user/Desktop/perturb_a2_b52b40a5.pdf as a PDF, preserving the paragraph order.",
                    "Convert the perturb_a2_b52b40a5_src.docx contents into a PDF file at /home/user/Desktop/perturb_a2_b52b40a5.pdf.",
                ],
            },
        ],
    },
    # P1 validation: tailored to real eval source files. Each base
    # uses the eval task's actual download URL — content extracted at
    # generation time from the real downloaded file (`/tmp/multi_apps_files/`).

    # 51f5801c — Dickinson_Slides.pptx has 9 slides; 8 have non-empty notes.
    # Notes pre-extracted via python-pptx; oracle writes them as docx
    # paragraphs in slide order.
    "osworld_multi_apps_51f5801c": {
        "archetype": "a4_pptx_notes_to_docx",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/51f5801c-18b3-4f25-b0c3-02f85507a078/Dickinson_Slides.pptx",
                    "src_basename": "Dickinson_Slides.pptx",
                    "target_basename": "Dickinson_Slides_notes.docx",
                    "notes": [
                        "This is opening slide.",
                        "Cover slide option #1",
                        "Cover slide option #3",
                        "This is a graph.",
                        "This is a table.",
                        "This is item lists.",
                        "This is an inserted image.",
                        "Blank ending slide",
                    ],
                },
                "instr_pool": [
                    "Open Dickinson_Slides.pptx in LibreOffice Impress, read the speaker note from each slide, and copy them into Dickinson_Slides_notes.docx as separate paragraphs in slide order. Skip any slide whose notes are empty.",
                    "Extract the speaker notes from Dickinson_Slides.pptx (one per slide, in slide order, skipping empty notes) and write each as a paragraph in Dickinson_Slides_notes.docx.",
                    "Read the speaker notes from every slide in Dickinson_Slides.pptx and append each non-empty note as a paragraph to Dickinson_Slides_notes.docx, preserving the original wording.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/51f5801c-18b3-4f25-b0c3-02f85507a078/Dickinson_Slides.pptx",
                    "src_basename": "Dickinson_Slides.pptx",
                    "target_basename": "Dickinson_first_three_notes.docx",
                    # Subset: only first 3 non-empty notes
                    "notes": [
                        "This is opening slide.",
                        "Cover slide option #1",
                        "Cover slide option #3",
                    ],
                },
                "instr_pool": [
                    "Take only the speaker notes from the FIRST 3 slides of Dickinson_Slides.pptx and copy them into Dickinson_first_three_notes.docx as paragraphs (one per slide, in slide order).",
                    "Open Dickinson_Slides.pptx in Impress, read the speaker notes for slides 1-3 only, and write each as a paragraph in Dickinson_first_three_notes.docx.",
                    "Extract the speaker notes from slides 1, 2, and 3 of Dickinson_Slides.pptx and put each into Dickinson_first_three_notes.docx as its own paragraph in the same order.",
                ],
            },
        ],
    },
    # 81c425f5 — OSP_Envelope_Price-List_2023_5000.xlsx is 23r×1024c. Take
    # the first 8 rows × 5 columns of the price-list region (envelope codes
    # MCC 150072..150066 with quantity-price tiers); join cells with ' | '
    # for paragraph form. Kept compact to fit max_steps budget.
    # 81c425f5 — eval evaluator is compare_docx_tables (expects a real Word
    # table in the docx). Switched validation from a4_xlsx_table_to_docx
    # (paragraphs) to a4_xlsx_to_docx_realtable (real table) for eval-
    # signature match. Real xlsx is 23r×1024c; we extract the first 5 cols
    # × 6 rows of the price-list region into a 6×5 docx table.
    "osworld_multi_apps_81c425f5": {
        "archetype": "a4_xlsx_to_docx_realtable",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/81c425f5-78f3-4771-afd6-3d2973825947/OSP_Envelope_Price-List_2023_5000.xlsx",
                    "src_basename": "OSP_Envelope_Price-List_2023_5000.xlsx",
                    "target_basename": "envelope_price_extract.docx",
                    # number-format DISPLAY strings (col A `#,##0`; cols B-E `$#,##0`,
                    # half-up per calc_102/69c7aedd Decimal precedent). A Calc->Writer
                    # copy-paste yields displayed text, so the gold must match display,
                    # not the raw cell value. §C #7 (81c425f5).
                    "rows": [
                        ["", "MCC 150072", "MCC 150006", "MCC 150063", "MCC 150066"],
                        ["5,000", "$617", "$645", "$623", "$664"],
                        ["6,000", "$656", "$690", "$664", "$713"],
                        ["7,000", "$696", "$735", "$704", "$762"],
                        ["7,500", "$716", "$758", "$725", "$786"],
                        ["10,000", "$814", "$870", "$826", "$908"],
                    ],
                },
                "instr_pool": [
                    "Open OSP_Envelope_Price-List_2023_5000.xlsx in LibreOffice Calc, select the first 5 columns of the first 6 rows, copy that range, and paste it as a 6-row × 5-column table into envelope_price_extract.docx.",
                    "Copy the rows 1-6 × columns A-E region from OSP_Envelope_Price-List_2023_5000.xlsx and paste it as a Word table in envelope_price_extract.docx, preserving cell values exactly.",
                    "Transfer the 6×5 region (rows 1-6, columns A-E) of OSP_Envelope_Price-List_2023_5000.xlsx into envelope_price_extract.docx as a real table — cell-for-cell match with the source.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/81c425f5-78f3-4771-afd6-3d2973825947/OSP_Envelope_Price-List_2023_5000.xlsx",
                    "src_basename": "OSP_Envelope_Price-List_2023_5000.xlsx",
                    "target_basename": "envelope_price_top4.docx",
                    # Same source region but a 5×4 subset (drop the last column,
                    # keep only the first 4 price tiers — narrower window).
                    "rows": [
                        ["", "MCC 150072", "MCC 150006", "MCC 150063"],
                        ["5,000", "$617", "$645", "$623"],
                        ["6,000", "$656", "$690", "$664"],
                        ["7,000", "$696", "$735", "$704"],
                        ["7,500", "$716", "$758", "$725"],
                    ],
                },
                "instr_pool": [
                    "Open OSP_Envelope_Price-List_2023_5000.xlsx and copy the first 4 columns of the first 5 rows into envelope_price_top4.docx as a 5×4 Word table, preserving the cell values exactly.",
                    "From OSP_Envelope_Price-List_2023_5000.xlsx, take the rows 1-5 × columns A-D region and paste it as a real table into envelope_price_top4.docx (5 rows × 4 columns).",
                    "Extract the upper-left 5×4 region of OSP_Envelope_Price-List_2023_5000.xlsx (rows 1-5, columns A-D) into envelope_price_top4.docx as a Word table — cell-for-cell match.",
                ],
            },
        ],
    },
    # c7c1e4c3 — Professor_Contact.xlsx has 3 HKU CS professors with their
    # universities, majors, websites (real values pre-extracted). Eval task
    # was chrome+xlsx state-heavy; perturb simplifies to "extract data → docx".
    "osworld_multi_apps_c7c1e4c3": {
        "archetype": "a4_xlsx_table_to_docx",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/c7c1e4c3-9e92-4eba-a4b8-689953975ea4/Professor_Contact.xlsx",
                    "src_basename": "Professor_Contact.xlsx",
                    "target_basename": "professor_directory.docx",
                    "rows": [
                        ["No.", "Professor", "University", "Major", "Website"],
                        ["1", "Qi Liu", "HKU", "Computer Science", "https://leuchine.github.io/"],
                        ["2", "Tao Yu", "HKU", "Computer Science", "https://taoyds.github.io/"],
                        ["3", "Lingpeng Kong", "HKU", "Computer Science", "https://ikekonglp.github.io/"],
                    ],
                },
                "instr_pool": [
                    "Open Professor_Contact.xlsx in LibreOffice Calc. Extract the first 5 columns (No., Professor, University, Major, Website) for the header and 3 professor rows, and write each row as a paragraph in professor_directory.docx joining cells with ' | '.",
                    "Read the professor records from Professor_Contact.xlsx (header row + 3 data rows, columns A-E excluding Email) and transcribe each row as a paragraph in professor_directory.docx, separating cells by ' | '.",
                    "Transcribe the professor table from Professor_Contact.xlsx (rows 2-5, columns A-E) into professor_directory.docx — one paragraph per row, ' | ' between cells.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/c7c1e4c3-9e92-4eba-a4b8-689953975ea4/Professor_Contact.xlsx",
                    "src_basename": "Professor_Contact.xlsx",
                    "target_basename": "professor_names_only.docx",
                    # Subset: only Professor + University columns (drop No., Major, Website).
                    "rows": [
                        ["Professor", "University"],
                        ["Qi Liu", "HKU"],
                        ["Tao Yu", "HKU"],
                        ["Lingpeng Kong", "HKU"],
                    ],
                },
                "instr_pool": [
                    "Open Professor_Contact.xlsx and extract just the Professor name and University columns (header row + 3 professor rows) into professor_names_only.docx — each row a paragraph with ' | ' between the two cells.",
                    "From Professor_Contact.xlsx, transcribe only the 'Professor' and 'University' columns (header + 3 data rows) into professor_names_only.docx as paragraphs joined by ' | '.",
                    "Take only the professor name and university columns from Professor_Contact.xlsx (header + 3 records) and write each row as a paragraph in professor_names_only.docx, separating the two cells with ' | '.",
                ],
            },
        ],
    },
    # 8e116af7 — my_bookkeeping.xlsx has a "Bookkeeping simple" header row +
    # column headers + 5 transactions. Eval task was complex (extract from
    # receipt images); perturb simplifies to "transcribe transactions → docx".
    "osworld_multi_apps_8e116af7": {
        "archetype": "a4_xlsx_table_to_docx",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/8e116af7-7db7-4e35-a68b-b0939c066c78/my_bookkeeping.xlsx",
                    "src_basename": "my_bookkeeping.xlsx",
                    "target_basename": "bookkeeping_summary.docx",
                    "rows": [
                        ["Description", "Category", "Type", "Amount", "Balance"],
                        ["Office Supplies Purchase", "Office Supplies", "Expense", "-150", "850"],
                        ["Client Payment Received", "Sales", "Income", "500", "1350"],
                        ["Internet Bill", "Utilities", "Expense", "-60", "1290"],
                    ],
                },
                "instr_pool": [
                    "Open my_bookkeeping.xlsx in LibreOffice Calc. Skip the first row ('Bookkeeping simple' title), then transcribe the column-header row and the first 3 transaction rows into bookkeeping_summary.docx — each row as a paragraph, columns separated by ' | '.",
                    "Read rows 2-5 of my_bookkeeping.xlsx (header + 3 transactions, columns A-E) and write each row as a paragraph in bookkeeping_summary.docx using ' | ' to separate cells.",
                    "Transcribe the header (Description, Category, Type, Amount, Balance) and the first 3 transactions from my_bookkeeping.xlsx into bookkeeping_summary.docx — one paragraph per row, columns delimited by ' | '.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/8e116af7-7db7-4e35-a68b-b0939c066c78/my_bookkeeping.xlsx",
                    "src_basename": "my_bookkeeping.xlsx",
                    "target_basename": "bookkeeping_amounts.docx",
                    # Subset: only Description + Amount columns, first 3 transactions.
                    "rows": [
                        ["Description", "Amount"],
                        ["Office Supplies Purchase", "-150"],
                        ["Client Payment Received", "500"],
                        ["Internet Bill", "-60"],
                    ],
                },
                "instr_pool": [
                    "Open my_bookkeeping.xlsx, skip the title row, and extract only the Description and Amount columns (header + first 3 transactions) into bookkeeping_amounts.docx — each row as a paragraph with ' | ' between the two cells.",
                    "From my_bookkeeping.xlsx (skip the 'Bookkeeping simple' row), take just the Description and Amount columns of the next 4 rows and write each as a paragraph in bookkeeping_amounts.docx joined by ' | '.",
                    "Transcribe only the Description and Amount fields of the first 3 transactions (with the header row above them) from my_bookkeeping.xlsx into bookkeeping_amounts.docx — one paragraph per row, two cells per paragraph.",
                ],
            },
        ],
    },
    # f8cfa149 — real eval source cell_search.xls has 51r×8c of person/customer
    # records: row-index col, First Name, Last Name, Gender, Country, Age,
    # Date, Id. Eval task was chrome-search-by-name (B6 ='Nereida') with
    # is_expected_search_query against active_tab_info — OFFICIALLY EXCLUDED
    # via metadata.others.exclude_reason: "Google reCAPTCHA redirects
    # /search?q=nereida to /sorry/index?continue=... non-deterministically,
    # breaking is_expected_search_query regex match". Validation
    # cross-app rebuild reused the chrome-search skeleton — same reCAPTCHA
    # failure mode applied to perturb. REVERTED to a4_xlsx_table_to_docx:
    # agent reads cell_search.xls in Calc and transcribes records to docx —
    # no chrome involvement, no reCAPTCHA dependency. Two variants exercise
    # different row/column subsets of the verified xls content.
    "osworld_multi_apps_f8cfa149": {
        "archetype": "a4_xlsx_table_to_docx",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/f8cfa149-d1c1-4215-8dac-4a0932bad3c2/file.xls",
                    "src_basename": "cell_search.xls",
                    "target_basename": "cell_search_records.docx",
                    # Pre-verified content: drop column A (row-index) per
                    # eval convention. Header + first 4 records.
                    "rows": [
                        ["First Name", "Last Name", "Gender", "Country", "Age"],
                        ["Dulce", "Abril", "Female", "United States", "32"],
                        ["Mara", "Hashimoto", "Female", "Great Britain", "25"],
                        ["Philip", "Gent", "Male", "France", "36"],
                        ["Kathleen", "Hanner", "Female", "United States", "25"],
                    ],
                },
                "instr_pool": [
                    "Open cell_search.xls in LibreOffice Calc. Read the columns First Name, Last Name, Gender, Country, Age (the header row + the first 4 data rows) and write each row as a paragraph in cell_search_records.docx, joining cells with ' | '. Skip the leading row-index column.",
                    "From cell_search.xls, transcribe columns B-F (header + first 4 records) into cell_search_records.docx as paragraphs separated by ' | ' — drop the row-index column at A and the Date / Id columns at G-H.",
                    "Take the first 4 records from cell_search.xls (First Name through Age, columns B-F) plus the header, and write each row as a paragraph in cell_search_records.docx using ' | ' between cells.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/f8cfa149-d1c1-4215-8dac-4a0932bad3c2/file.xls",
                    "src_basename": "cell_search.xls",
                    "target_basename": "cell_search_names.docx",
                    # Subset variant: only First Name + Last Name + Country
                    # for the first 5 records (different rows + columns).
                    "rows": [
                        ["First Name", "Last Name", "Country"],
                        ["Dulce", "Abril", "United States"],
                        ["Mara", "Hashimoto", "Great Britain"],
                        ["Philip", "Gent", "France"],
                        ["Kathleen", "Hanner", "United States"],
                        ["Nereida", "Magwood", "United States"],
                    ],
                },
                "instr_pool": [
                    "Open cell_search.xls in LibreOffice Calc. Extract just the First Name, Last Name, and Country columns for the header and the first 5 records into cell_search_names.docx — each row as a paragraph with ' | ' between the three cells.",
                    "From cell_search.xls, take only the First Name, Last Name, Country columns (header + first 5 records) and transcribe each row as a paragraph in cell_search_names.docx joined by ' | '.",
                    "Pull the First Name + Last Name + Country fields of the first 5 people from cell_search.xls (with the header above them) into cell_search_names.docx — one paragraph per row, three cells per row joined by ' | '.",
                ],
            },
        ],
    },
    # 185f29bd — Employee Performance Evaluation Summary.xlsx has 7 employees
    # × 26 columns. Take first 5 columns (Name, ID, Date, Position, Dept) for
    # 4 employees + header. Date col formatted as YYYY-MM-DD via openpyxl.
    "osworld_multi_apps_185f29bd": {
        "archetype": "a4_xlsx_table_to_docx",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/185f29bd-5da0-40a6-b69c-ba7f4e0324ef/Employee%20Performance%20Evaluation%20Summary.xlsx",
                    "src_basename": "Employee Performance Evaluation Summary.xlsx",
                    "target_basename": "employee_directory.docx",
                    "rows": [
                        ["EMPLOYEE NAME", "EMPLOYEE ID", "DATE OF CURRENT REVIEW", "POSITION HELD", "DEPARTMENT"],
                        ["John Doe", "12345", "2023-04-01", "Analyst", "Finance"],
                        ["Emily Johnson", "67890", "2023-03-25", "Project Manager", "Marketing"],
                        ["Michael Brown", "11223", "2023-03-15", "Software Engineer", "IT"],
                        ["Linda Green", "44556", "2023-04-05", "Sales Representative", "Sales"],
                    ],
                },
                "instr_pool": [
                    "Open Employee Performance Evaluation Summary.xlsx in LibreOffice Calc. Take the first 5 columns (EMPLOYEE NAME, EMPLOYEE ID, DATE OF CURRENT REVIEW, POSITION HELD, DEPARTMENT) of the header row and the first 4 employee rows, transcribe each as a paragraph in employee_directory.docx, joining cells with ' | '. Format dates as YYYY-MM-DD.",
                    "Read columns A-E of rows 1-5 from Employee Performance Evaluation Summary.xlsx and write each row as a paragraph in employee_directory.docx — columns separated by ' | ', dates rendered as YYYY-MM-DD.",
                    "Extract the first 5 columns × first 5 rows (header + 4 employees) of Employee Performance Evaluation Summary.xlsx into employee_directory.docx, one paragraph per row with ' | ' between cells; render the date column as YYYY-MM-DD strings.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/185f29bd-5da0-40a6-b69c-ba7f4e0324ef/Employee%20Performance%20Evaluation%20Summary.xlsx",
                    "src_basename": "Employee Performance Evaluation Summary.xlsx",
                    "target_basename": "employee_positions.docx",
                    # Subset: name + position + department only (drop ID, date), first 3 employees.
                    "rows": [
                        ["EMPLOYEE NAME", "POSITION HELD", "DEPARTMENT"],
                        ["John Doe", "Analyst", "Finance"],
                        ["Emily Johnson", "Project Manager", "Marketing"],
                        ["Michael Brown", "Software Engineer", "IT"],
                    ],
                },
                "instr_pool": [
                    "Open Employee Performance Evaluation Summary.xlsx and extract only the EMPLOYEE NAME, POSITION HELD, and DEPARTMENT columns (header row + first 3 employee rows) into employee_positions.docx — each row a paragraph with ' | ' between cells.",
                    "From Employee Performance Evaluation Summary.xlsx, take just the name + position + department columns of the first 3 employees (plus the header row) and write each row as a paragraph in employee_positions.docx joined by ' | '.",
                    "Transcribe only the EMPLOYEE NAME, POSITION HELD, DEPARTMENT fields of the first 3 employees (with the header row above them) from Employee Performance Evaluation Summary.xlsx into employee_positions.docx — one paragraph per row, three cells joined by ' | '.",
                ],
            },
        ],
    },
    # a503b07f — OIP.jpg (real receipt image) → PDF. Eval task says
    # "transform receipt image into PDF". Oracle uses PIL for deterministic
    # image→PDF conversion; agent typically does GIMP File→Export As PDF.
    "osworld_multi_apps_a503b07f": {
        "archetype": "a10_image_to_pdf",
        "tier": "a10",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/a503b07f-9119-456b-b75d-f5146737d24f/OIP.jpg",
                    "src_basename": "OIP.jpg",
                    "target_basename": "OIP.pdf",
                },
                "instr_pool": [
                    "Convert OIP.jpg (the receipt image on the Desktop) to a PDF saved at /home/user/Desktop/OIP.pdf, fitting the image to a single page.",
                    "Open OIP.jpg in GIMP and export it as a PDF at /home/user/Desktop/OIP.pdf.",
                    "Save the OIP.jpg receipt image as a PDF document at /home/user/Desktop/OIP.pdf.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/a503b07f-9119-456b-b75d-f5146737d24f/OIP.jpg",
                    "src_basename": "OIP.jpg",
                    "target_basename": "receipt_archive.pdf",
                },
                "instr_pool": [
                    "Take OIP.jpg from the Desktop and produce a single-page PDF of it at /home/user/Desktop/receipt_archive.pdf.",
                    "Convert the OIP.jpg receipt image to a PDF named receipt_archive.pdf on the Desktop, image fitted to one page.",
                    "I want to archive OIP.jpg as a PDF — save it as /home/user/Desktop/receipt_archive.pdf, fitting the photo to a single page.",
                ],
            },
        ],
    },
    # d1acdb87 — real eval sources are restaurants.txt (5 restaurant names,
    # one per line numbered) + MUST_VISIT.xlsx (1000-row template, header
    # "Restaurant Name | Address | Contact Number | Website" then all empty).
    # Eval evaluator is compare_table. Perturb: agent reads the 5 names
    # from restaurants.txt and writes them into Restaurant Name column
    # (A2:A6) of MUST_VISIT.xlsx.
    "osworld_multi_apps_d1acdb87": {
        "archetype": "a4_txt_to_xlsx_column",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "txt_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/d1acdb87-bb67-4f30-84aa-990e56a09c92/restaurants.txt",
                    "txt_basename": "restaurants.txt",
                    "xlsx_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/d1acdb87-bb67-4f30-84aa-990e56a09c92/MUST_VISIT.xlsx",
                    "xlsx_basename": "MUST_VISIT.xlsx",
                    # Pre-extracted from real restaurants.txt (numbered '1. Ming Pavilion', etc.)
                    # Strip the "N. " prefix to get bare names.
                    "restaurant_names": [
                        "Ming Pavilion",
                        "Cristal Room by Anne-Sophie Pic",
                        "Leela",
                        "Nobu",
                        "Niras",
                    ],
                },
                "instr_pool": [
                    "Open restaurants.txt on the Desktop, take the 5 restaurant names (strip the '1.', '2.', ... numbering prefix) and enter them into the Restaurant Name column (column A, rows 2-6) of MUST_VISIT.xlsx in the same order they appear in the txt file.",
                    "Read the 5 numbered restaurant names from restaurants.txt and fill them into MUST_VISIT.xlsx's Restaurant Name column starting at A2, preserving the order; remove the leading 'N.' numbering when copying.",
                    "Transfer the 5 restaurants listed in restaurants.txt (one per line) into the MUST_VISIT.xlsx spreadsheet — write them into column A (under the 'Restaurant Name' header) in rows 2-6, dropping the numeric prefix.",
                ],
            },
            {
                "params": {
                    "txt_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/d1acdb87-bb67-4f30-84aa-990e56a09c92/restaurants.txt",
                    "txt_basename": "restaurants.txt",
                    "xlsx_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/d1acdb87-bb67-4f30-84aa-990e56a09c92/MUST_VISIT.xlsx",
                    "xlsx_basename": "MUST_VISIT.xlsx",
                    # Top 3 only (subset of the full 5).
                    "restaurant_names": [
                        "Ming Pavilion",
                        "Cristal Room by Anne-Sophie Pic",
                        "Leela",
                    ],
                },
                "instr_pool": [
                    "Take only the FIRST 3 restaurants from restaurants.txt and write them into MUST_VISIT.xlsx column A rows 2-4 (drop the 'N.' numbering). Leave A5-A6 blank.",
                    "Open restaurants.txt and copy the top 3 entries (numbered 1, 2, 3) into the Restaurant Name column of MUST_VISIT.xlsx (rows A2-A4), stripping the '1.'/'2.'/'3.' prefixes.",
                    "Read the first three restaurant names from restaurants.txt and put them in MUST_VISIT.xlsx column A under the header (A2-A4) without the leading '#.' numbering.",
                ],
            },
        ],
    },
    # P1 validation archetype A4_real_docx_to_pdf: use eval's real
    # downloaded source docx instead of synthesizing one. Targets eval
    # tasks where original instruction was docx→pdf + Drive upload (Drive
    # part dropped per plan).
    # ce2b64a2 — 3 real mountain jpg images. Eval task asks image identification
    # (OUT-OF-SCOPE). Perturb redirects to file-format conversion: agent
    # converts each jpg → png in the same folder. check_include_exclude
    # checks that all expected png filenames exist (per-byte image match
    # would be fragile across encoders).
    "osworld_multi_apps_ce2b64a2": {
        "archetype": "a10_batch_jpg_to_png",
        "tier": "a10",
        "variants": [
            {
                "params": {
                    "pic_dirname": "Pictures_ce2b64a2",
                    "src_urls": [
                        ["https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/ce2b64a2-ddc1-4f91-8c7d-a88be7121aac/picture1.jpg", "picture1.jpg"],
                        ["https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/ce2b64a2-ddc1-4f91-8c7d-a88be7121aac/picture2.jpg", "picture2.jpg"],
                        ["https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/ce2b64a2-ddc1-4f91-8c7d-a88be7121aac/picture3.jpg", "picture3.jpg"],
                    ],
                },
                "instr_pool": [
                    "In the Pictures_ce2b64a2 folder on the Desktop, convert each .jpg file (picture1.jpg, picture2.jpg, picture3.jpg) to a .png file with the same base name in the same folder. Keep the original jpg files intact.",
                    "Take the 3 jpg images in Desktop/Pictures_ce2b64a2 and produce a .png copy of each (picture1.png, picture2.png, picture3.png) in the same folder.",
                    "For each jpg in the Pictures_ce2b64a2 folder on the Desktop, save a .png version with matching base name alongside the original.",
                ],
            },
            {
                "params": {
                    "pic_dirname": "Photos_ce2b64a2_v2",
                    "src_urls": [
                        ["https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/ce2b64a2-ddc1-4f91-8c7d-a88be7121aac/picture1.jpg", "shot_a.jpg"],
                        ["https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/ce2b64a2-ddc1-4f91-8c7d-a88be7121aac/picture2.jpg", "shot_b.jpg"],
                        ["https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/ce2b64a2-ddc1-4f91-8c7d-a88be7121aac/picture3.jpg", "shot_c.jpg"],
                    ],
                },
                "instr_pool": [
                    "In the Photos_ce2b64a2_v2 folder on the Desktop, convert each .jpg file (shot_a.jpg, shot_b.jpg, shot_c.jpg) to .png with the same base name. Originals stay.",
                    "Produce .png copies of all 3 jpgs (shot_a/b/c) in Desktop/Photos_ce2b64a2_v2 — same base name, .png extension.",
                    "Save .png versions of shot_a.jpg, shot_b.jpg, shot_c.jpg in Desktop/Photos_ce2b64a2_v2 alongside the originals.",
                ],
            },
        ],
    },
    # bc2b57f3 — real eval sources are workbook-with-sample-database.xlsx
    # (10 sheets: 'Mkt Ouallam', 'Sorghum', 'Millet spatial integration',
    # 'Millet', 'ReadMe', 'Mkt Gotheye', 'Raw_data', 'Beans', 'Mkt Mangaize',
    # 'Mkt Tera') and reminder.docx (assignment instructions). Eval task
    # asks the agent to modify the workbook per reminder; that requires
    # multi-step interpretation. Perturb simplifies to "list the workbook's
    # sheet names into a docx in their existing order".
    "osworld_multi_apps_bc2b57f3": {
        "archetype": "a4_pptx_notes_to_docx",  # reuses the "list-of-strings → docx paragraphs" pattern
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/bc2b57f3-686d-4ec9-87ce-edf850b7e442/workbook-with-sample-database.xlsx",
                    "src_basename": "workbook-with-sample-database.xlsx",
                    "target_basename": "workbook_sheet_names.docx",
                    "src_app": "--calc",
                    "notes": [
                        "Mkt Ouallam",
                        "Sorghum",
                        "Millet spatial integration",
                        "Millet",
                        "ReadMe",
                        "Mkt Gotheye",
                        "Raw_data",
                        "Beans",
                        "Mkt Mangaize",
                        "Mkt Tera",
                    ],
                },
                "instr_pool": [
                    "Open workbook-with-sample-database.xlsx in LibreOffice Calc, read each sheet tab name in tab order from left to right, and write each name as a paragraph in workbook_sheet_names.docx — preserve the exact tab labels and ordering.",
                    "List every sheet tab name from workbook-with-sample-database.xlsx (in tab-bar left-to-right order) as separate paragraphs in workbook_sheet_names.docx, keeping the names verbatim.",
                    "Read the sheet tab names from workbook-with-sample-database.xlsx (in their tab-bar order) and transcribe each as its own paragraph into workbook_sheet_names.docx on the Desktop.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/bc2b57f3-686d-4ec9-87ce-edf850b7e442/workbook-with-sample-database.xlsx",
                    "src_basename": "workbook-with-sample-database.xlsx",
                    "target_basename": "first_five_tabs.docx",
                    "src_app": "--calc",
                    # Subset of first 5 sheet names only.
                    "notes": [
                        "Mkt Ouallam",
                        "Sorghum",
                        "Millet spatial integration",
                        "Millet",
                        "ReadMe",
                    ],
                },
                "instr_pool": [
                    "Open workbook-with-sample-database.xlsx and list ONLY THE FIRST 5 sheet tab names (left-to-right order) into first_five_tabs.docx as separate paragraphs.",
                    "Take the names of the first 5 sheet tabs from workbook-with-sample-database.xlsx and write each as a paragraph in first_five_tabs.docx — keep the order and exact wording.",
                    "From workbook-with-sample-database.xlsx, extract the first 5 sheet tab labels (in tab-bar order) and put each as its own paragraph in first_five_tabs.docx.",
                ],
            },
        ],
    },
    # 227d2f97 — real .xcf file → docx with image embedded.
    # Eval expects a docx with the (flattened) xcf content as an image.
    "osworld_multi_apps_227d2f97": {
        "archetype": "a10_xcf_to_docx_image",
        "tier": "a10",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/227d2f97-562b-4ccb-ae47-a5ec9e142fbb/QTdHniCqfJbBLJe3L3nijU-1200-80.xcf",
                    "src_basename": "QTdHniCqfJbBLJe3L3nijU-1200-80.xcf",
                    "target_basename": "image.docx",
                },
                "instr_pool": [
                    "Open QTdHniCqfJbBLJe3L3nijU-1200-80.xcf in GIMP, flatten + export it as a PNG, then create image.docx on the Desktop containing that PNG as an embedded image.",
                    "Use GIMP to flatten the QTdHniCqfJbBLJe3L3nijU-1200-80.xcf file and export it to PNG, then insert that PNG into a LibreOffice Writer document saved as image.docx on the Desktop.",
                    "Take the .xcf on the Desktop (QTdHniCqfJbBLJe3L3nijU-1200-80.xcf), flatten it and export to PNG via GIMP, then paste the PNG into a Writer document saved as image.docx.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/227d2f97-562b-4ccb-ae47-a5ec9e142fbb/QTdHniCqfJbBLJe3L3nijU-1200-80.xcf",
                    "src_basename": "QTdHniCqfJbBLJe3L3nijU-1200-80.xcf",
                    "target_basename": "graphic_export.docx",
                },
                "instr_pool": [
                    "Flatten the .xcf on the Desktop in GIMP, export it as PNG, and embed that PNG into a Writer document at /home/user/Desktop/graphic_export.docx.",
                    "Convert QTdHniCqfJbBLJe3L3nijU-1200-80.xcf to PNG via GIMP (flatten first), then create graphic_export.docx on the Desktop with the PNG inserted.",
                    "Take the .xcf graphic, flatten + export to PNG in GIMP, and place the PNG inside a new Writer document saved as /home/user/Desktop/graphic_export.docx.",
                ],
            },
        ],
    },
    "osworld_multi_apps_22a4636f": {
        "archetype": "a4_real_docx_to_pdf",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/22a4636f-8179-4357-8e87-d1743ece1f81/Meeting-Agenda.docx",
                    "src_basename": "Meeting-Agenda.docx",
                    "target_basename": "Meeting-Agenda.pdf",
                },
                "instr_pool": [
                    "Open Meeting-Agenda.docx in LibreOffice Writer and export it as a PDF saved at /home/user/Desktop/Meeting-Agenda.pdf, preserving the document content exactly.",
                    "Convert the Meeting-Agenda.docx file (already on the Desktop) into a PDF at /home/user/Desktop/Meeting-Agenda.pdf using LibreOffice Writer's PDF export.",
                    "Export Meeting-Agenda.docx to /home/user/Desktop/Meeting-Agenda.pdf as a PDF; keep the original document untouched.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/22a4636f-8179-4357-8e87-d1743ece1f81/Meeting-Agenda.docx",
                    "src_basename": "Meeting-Agenda.docx",
                    "target_basename": "agenda_distribution.pdf",
                },
                "instr_pool": [
                    "Convert Meeting-Agenda.docx into a PDF named /home/user/Desktop/agenda_distribution.pdf using LibreOffice Writer.",
                    "Export Meeting-Agenda.docx to PDF for distribution — save the result as /home/user/Desktop/agenda_distribution.pdf.",
                    "I need a PDF copy of Meeting-Agenda.docx to share with attendees; export it as /home/user/Desktop/agenda_distribution.pdf via LibreOffice Writer.",
                ],
            },
        ],
    },
    "osworld_multi_apps_897e3b53": {
        "archetype": "a4_real_docx_to_pdf",
        "tier": "a4",
        "variants": [
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/897e3b53-5d4d-444b-85cb-2cdc8a97d903/form.docx",
                    "src_basename": "form.docx",
                    "target_basename": "form.pdf",
                },
                "instr_pool": [
                    "Convert form.docx (on the Desktop) to a PDF saved at /home/user/Desktop/form.pdf using LibreOffice Writer.",
                    "Open form.docx in LibreOffice Writer and export it to a PDF file at /home/user/Desktop/form.pdf.",
                    "Save form.docx as a PDF at /home/user/Desktop/form.pdf using LibreOffice Writer's export feature.",
                ],
            },
            {
                "params": {
                    "src_url": "https://huggingface.co/datasets/xlangai/ubuntu_osworld_file_cache/resolve/main/multi_apps/897e3b53-5d4d-444b-85cb-2cdc8a97d903/form.docx",
                    "src_basename": "form.docx",
                    "target_basename": "form_signed_copy.pdf",
                },
                "instr_pool": [
                    "Export form.docx as a PDF named form_signed_copy.pdf on the Desktop — use LibreOffice Writer's File > Export As > PDF flow.",
                    "Save form.docx as /home/user/Desktop/form_signed_copy.pdf via LibreOffice Writer's PDF export.",
                    "I need a PDF copy of form.docx for filing — export it from LibreOffice Writer to /home/user/Desktop/form_signed_copy.pdf.",
                ],
            },
        ],
    },
    # ---------- Tier A2: image format convert ----------
    "osworld_multi_apps_2fe4b718": {
        "archetype": "a2_image_format_convert",
        "tier": "a2",
        "variants": [
            {
                "params": {"fmt_in": "png", "fmt_out": "jpg",
                           "size": (200, 150), "color": (60, 120, 180)},
                "instr_pool": [
                    "Open perturb_a2_2fe4b718_src.png in GIMP and export it as a JPEG file at /home/user/Desktop/perturb_a2_2fe4b718.jpg, keeping the image content the same.",
                    "Convert the perturb_a2_2fe4b718_src.png image to a JPEG file saved at /home/user/Desktop/perturb_a2_2fe4b718.jpg using GIMP, preserving the visible content.",
                    "Save perturb_a2_2fe4b718_src.png as /home/user/Desktop/perturb_a2_2fe4b718.jpg in JPEG format using GIMP.",
                ],
            },
            {
                "params": {"fmt_in": "png", "fmt_out": "jpg",
                           "size": (320, 200), "color": (220, 90, 90)},
                "instr_pool": [
                    "Export perturb_a2_2fe4b718_src.png to JPEG format at /home/user/Desktop/perturb_a2_2fe4b718.jpg using GIMP, keeping the dimensions and content unchanged.",
                    "Convert the perturb_a2_2fe4b718_src.png file into a JPEG image at /home/user/Desktop/perturb_a2_2fe4b718.jpg without changing what the image looks like.",
                    "Save the perturb_a2_2fe4b718_src.png picture as /home/user/Desktop/perturb_a2_2fe4b718.jpg in JPEG format with GIMP.",
                ],
            },
        ],
    },
    "osworld_multi_apps_c2751594": {
        "archetype": "a2_image_format_convert",
        "tier": "a2",
        "variants": [
            {
                "params": {"fmt_in": "png", "fmt_out": "jpg",
                           "size": (240, 160), "color": (40, 160, 90)},
                "instr_pool": [
                    "Open perturb_a2_c2751594_src.png in GIMP and save a JPEG copy at /home/user/Desktop/perturb_a2_c2751594.jpg with the same image content.",
                    "Convert perturb_a2_c2751594_src.png to a JPEG file at /home/user/Desktop/perturb_a2_c2751594.jpg using GIMP, preserving the visible image.",
                    "Export perturb_a2_c2751594_src.png as JPEG to /home/user/Desktop/perturb_a2_c2751594.jpg using GIMP without modifying the picture content.",
                ],
            },
            {
                "params": {"fmt_in": "png", "fmt_out": "jpg",
                           "size": (160, 240), "color": (200, 140, 70)},
                "instr_pool": [
                    "Save perturb_a2_c2751594_src.png as /home/user/Desktop/perturb_a2_c2751594.jpg in JPEG format using GIMP.",
                    "Export the perturb_a2_c2751594_src.png picture to a JPEG file at /home/user/Desktop/perturb_a2_c2751594.jpg with GIMP.",
                    "Convert perturb_a2_c2751594_src.png into a JPEG image saved at /home/user/Desktop/perturb_a2_c2751594.jpg.",
                ],
            },
        ],
    },
    # ---------- Tier A2: text file create ----------
    "osworld_multi_apps_acb0f96b": {
        "archetype": "a2_text_file_create",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "template": "TODO: write project README\n",
                    "expected": "Project README\nThis project is a tiny demo for the Perturb training set.\n",
                },
                "instr_pool": [
                    "Replace the contents of perturb_a2_acb0f96b_template.txt and save the new text as /home/user/Desktop/perturb_a2_acb0f96b.txt with the README body 'Project README' on the first line and the demo description on the second line.",
                    "Write the project README into /home/user/Desktop/perturb_a2_acb0f96b.txt with two lines, the first being the title and the second describing it as a tiny Perturb training demo.",
                    "Create /home/user/Desktop/perturb_a2_acb0f96b.txt with two lines, the title 'Project README' and a one-sentence description that the project is a tiny demo for the Perturb training set.",
                ],
            },
            {
                "params": {
                    "template": "TODO: notes for the meeting\n",
                    "expected": "Meeting notes\nDiscuss the new schedule and the next deliverables.\n",
                },
                "instr_pool": [
                    "Write the meeting notes into /home/user/Desktop/perturb_a2_acb0f96b.txt with the title 'Meeting notes' and a second line summarising the topics: the new schedule and the next deliverables.",
                    "Save the meeting notes as /home/user/Desktop/perturb_a2_acb0f96b.txt containing the title and a one-sentence summary of the schedule and deliverables discussion.",
                    "Create /home/user/Desktop/perturb_a2_acb0f96b.txt holding the meeting notes, with 'Meeting notes' on line 1 and the discussion summary on line 2.",
                ],
            },
        ],
    },
    "osworld_multi_apps_f918266a": {
        "archetype": "a2_text_file_create",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "template": "calculator log placeholder\n",
                    "expected": "result = 42\nstatus = ok\n",
                },
                "instr_pool": [
                    "Write the calculator output to /home/user/Desktop/perturb_a2_f918266a.txt with two lines: 'result = 42' on the first line and 'status = ok' on the second line.",
                    "Save the calculator log as /home/user/Desktop/perturb_a2_f918266a.txt containing the line 'result = 42' followed by 'status = ok'.",
                    "Create /home/user/Desktop/perturb_a2_f918266a.txt with the calculator output: line 1 is 'result = 42' and line 2 is 'status = ok'.",
                ],
            },
            {
                "params": {
                    "template": "log placeholder\n",
                    "expected": "result = 100\nstatus = ok\n",
                },
                "instr_pool": [
                    "Write /home/user/Desktop/perturb_a2_f918266a.txt with the calculator log lines 'result = 100' and 'status = ok'.",
                    "Save the calculator output as /home/user/Desktop/perturb_a2_f918266a.txt containing 'result = 100' on line 1 and 'status = ok' on line 2.",
                    "Create /home/user/Desktop/perturb_a2_f918266a.txt holding the two lines 'result = 100' and 'status = ok' in that order.",
                ],
            },
        ],
    },
    # ---------- Tier A2: python test suite ----------
    "osworld_multi_apps_26150609": {
        "archetype": "a2_python_test_suite",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "buggy_module": (
                        "def add(a, b):\n"
                        "    # BUG: subtracting instead of adding\n"
                        "    return a - b\n"
                    ),
                    "fixed_module": (
                        "def add(a, b):\n"
                        "    return a + b\n"
                    ),
                    "test_code": (
                        "from module import add\n"
                        "def test():\n"
                        "    return add(2, 3) == 5 and add(-1, 1) == 0\n"
                    ),
                },
                "instr_pool": [
                    "Open the perturb_a2_26150609_proj folder in VS Code and fix the add function in module.py so that calling add(a, b) returns the actual sum of the two arguments instead of their difference.",
                    "Edit module.py inside the perturb_a2_26150609_proj project so that add(a, b) correctly returns a + b instead of subtracting the inputs.",
                    "Fix the buggy add(a, b) function in module.py inside the perturb_a2_26150609_proj folder so that it returns the sum of its two arguments.",
                ],
            },
            {
                "params": {
                    "buggy_module": (
                        "def is_even(n):\n"
                        "    # BUG: returning the wrong parity\n"
                        "    return n % 2 != 0\n"
                    ),
                    "fixed_module": (
                        "def is_even(n):\n"
                        "    return n % 2 == 0\n"
                    ),
                    "test_code": (
                        "from module import is_even\n"
                        "def test():\n"
                        "    return is_even(4) and not is_even(5) and is_even(0)\n"
                    ),
                },
                "instr_pool": [
                    "Open module.py in VS Code from the perturb_a2_26150609_proj folder and fix is_even(n) so that it returns True when n is even and False otherwise.",
                    "Edit module.py inside perturb_a2_26150609_proj to fix the is_even function so that it returns the correct parity for the input integer.",
                    "Fix the broken is_even(n) implementation in module.py inside perturb_a2_26150609_proj so that it returns True for even integers and False for odd integers.",
                ],
            },
        ],
    },
    "osworld_multi_apps_9219480b": {
        "archetype": "a2_python_test_suite",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "buggy_module": (
                        "def factorial(n):\n"
                        "    # BUG: missing base case, recurses forever for n=0\n"
                        "    return n * factorial(n - 1)\n"
                    ),
                    "fixed_module": (
                        "def factorial(n):\n"
                        "    if n <= 1:\n"
                        "        return 1\n"
                        "    return n * factorial(n - 1)\n"
                    ),
                    "test_code": (
                        "from module import factorial\n"
                        "def test():\n"
                        "    return factorial(0) == 1 and factorial(5) == 120\n"
                    ),
                },
                "instr_pool": [
                    "Open module.py from the perturb_a2_9219480b_proj folder in VS Code and fix the factorial function so it has a base case for n <= 1 and returns the right factorial for non-negative integers.",
                    "Edit factorial(n) inside module.py in perturb_a2_9219480b_proj so that it terminates for n = 0 and returns the correct factorial value.",
                    "Fix the recursive factorial implementation in module.py inside perturb_a2_9219480b_proj by adding a base case so factorial(0) returns 1 and factorial(5) returns 120.",
                ],
            },
            {
                "params": {
                    "buggy_module": (
                        "def fib(n):\n"
                        "    # BUG: off-by-one base case\n"
                        "    if n == 0:\n"
                        "        return 1\n"
                        "    if n == 1:\n"
                        "        return 2\n"
                        "    return fib(n - 1) + fib(n - 2)\n"
                    ),
                    "fixed_module": (
                        "def fib(n):\n"
                        "    if n == 0:\n"
                        "        return 0\n"
                        "    if n == 1:\n"
                        "        return 1\n"
                        "    return fib(n - 1) + fib(n - 2)\n"
                    ),
                    "test_code": (
                        "from module import fib\n"
                        "def test():\n"
                        "    return fib(0) == 0 and fib(1) == 1 and fib(6) == 8\n"
                    ),
                },
                "instr_pool": [
                    "Edit fib(n) in module.py inside the perturb_a2_9219480b_proj folder so that fib(0) is 0 and fib(1) is 1, fixing the off-by-one base case.",
                    "Fix the broken Fibonacci function in module.py inside perturb_a2_9219480b_proj so it returns the standard sequence starting with 0 and 1.",
                    "Open module.py in the perturb_a2_9219480b_proj project and correct the base case of fib(n) so that fib(0) = 0 and fib(1) = 1.",
                ],
            },
        ],
    },
    # ---------- Tier A2: json export ----------
    "osworld_multi_apps_e2392362": {
        "archetype": "a2_json_export",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "stub_json": {"name": "TODO", "title": "TODO", "email": "todo@example.com"},
                    "final_json": {"name": "Sam Carter", "title": "Researcher", "email": "sam@example.com"},
                    "rules": {"expect": [
                        {"key": ["name"], "method": "eq", "ref": "Sam Carter"},
                        {"key": ["title"], "method": "eq", "ref": "Researcher"},
                        {"key": ["email"], "method": "eq", "ref": "sam@example.com"},
                    ]},
                },
                "instr_pool": [
                    "Edit perturb_a2_e2392362_stub.json so that the personal config saved at /home/user/Desktop/perturb_a2_e2392362.json has name set to 'Sam Carter', title set to 'Researcher', and email set to 'sam@example.com'.",
                    "Replace the TODO placeholders in perturb_a2_e2392362_stub.json and write the populated JSON to /home/user/Desktop/perturb_a2_e2392362.json with name 'Sam Carter', title 'Researcher', and email 'sam@example.com'.",
                    "Fill the perturb_a2_e2392362_stub.json template and save the result as /home/user/Desktop/perturb_a2_e2392362.json with the values name='Sam Carter', title='Researcher', email='sam@example.com'.",
                ],
            },
            {
                "params": {
                    "stub_json": {"name": "TODO", "title": "TODO", "email": "todo@example.com"},
                    "final_json": {"name": "Lin Wei", "title": "Engineer", "email": "lin@example.com"},
                    "rules": {"expect": [
                        {"key": ["name"], "method": "eq", "ref": "Lin Wei"},
                        {"key": ["title"], "method": "eq", "ref": "Engineer"},
                        {"key": ["email"], "method": "eq", "ref": "lin@example.com"},
                    ]},
                },
                "instr_pool": [
                    "Edit perturb_a2_e2392362_stub.json and save the populated config to /home/user/Desktop/perturb_a2_e2392362.json with name='Lin Wei', title='Engineer', and email='lin@example.com'.",
                    "Fill in the stub at perturb_a2_e2392362_stub.json and write /home/user/Desktop/perturb_a2_e2392362.json with the JSON values name 'Lin Wei', title 'Engineer', and email 'lin@example.com'.",
                    "Replace the placeholders in perturb_a2_e2392362_stub.json and save the result as /home/user/Desktop/perturb_a2_e2392362.json with name 'Lin Wei', title 'Engineer', email 'lin@example.com'.",
                ],
            },
        ],
    },
    # ---------- Tier A2: include/exclude ----------
    "osworld_multi_apps_2c9fc0de": {
        "archetype": "a2_check_include_exclude",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "files": {
                        "keep_a.txt": "keep this file\n",
                        "keep_b.txt": "keep this file too\n",
                        "tmp_drop.txt": "delete me\n",
                    },
                    "remove": ["tmp_drop.txt"],
                    "include": ["keep_a.txt", "keep_b.txt"],
                    "exclude": ["tmp_drop.txt"],
                },
                "instr_pool": [
                    "Delete the file tmp_drop.txt from the perturb_a2_2c9fc0de_dir folder on the Desktop while keeping keep_a.txt and keep_b.txt intact, then leave the folder layout that way.",
                    "Remove tmp_drop.txt from the perturb_a2_2c9fc0de_dir directory on the Desktop, leaving keep_a.txt and keep_b.txt as they are.",
                    "Tidy up perturb_a2_2c9fc0de_dir on the Desktop by deleting tmp_drop.txt; keep both keep_a.txt and keep_b.txt where they are.",
                ],
            },
            {
                "params": {
                    "files": {
                        "report.txt": "final report\n",
                        "scratch_1.txt": "scratch\n",
                        "scratch_2.txt": "scratch\n",
                    },
                    "remove": ["scratch_1.txt", "scratch_2.txt"],
                    "include": ["report.txt"],
                    "exclude": ["scratch_1.txt", "scratch_2.txt"],
                },
                "instr_pool": [
                    "Delete both scratch_1.txt and scratch_2.txt from perturb_a2_2c9fc0de_dir on the Desktop while keeping report.txt where it is.",
                    "Remove the two scratch_*.txt files from the perturb_a2_2c9fc0de_dir folder; report.txt should stay in place.",
                    "Clean up perturb_a2_2c9fc0de_dir on the Desktop by deleting scratch_1.txt and scratch_2.txt and keeping report.txt.",
                ],
            },
        ],
    },
    "osworld_multi_apps_937087b6": {
        "archetype": "a2_check_include_exclude",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "files": {
                        "alpha.log": "alpha log\n",
                        "beta.log": "beta log\n",
                        "gamma.tmp": "temp\n",
                    },
                    "remove": ["gamma.tmp"],
                    "include": ["alpha.log", "beta.log"],
                    "exclude": ["gamma.tmp"],
                },
                "instr_pool": [
                    "Delete gamma.tmp from perturb_a2_937087b6_dir on the Desktop, keeping alpha.log and beta.log untouched.",
                    "Remove the temporary file gamma.tmp from perturb_a2_937087b6_dir; both alpha.log and beta.log should remain.",
                    "Clean up perturb_a2_937087b6_dir by deleting gamma.tmp and keeping the two log files.",
                ],
            },
            {
                "params": {
                    "files": {
                        "doc.md": "doc\n",
                        "old.bak": "old backup\n",
                    },
                    "remove": ["old.bak"],
                    "include": ["doc.md"],
                    "exclude": ["old.bak"],
                },
                "instr_pool": [
                    "Delete old.bak from perturb_a2_937087b6_dir on the Desktop while keeping doc.md.",
                    "Remove the old.bak backup from the perturb_a2_937087b6_dir folder; the doc.md file should stay.",
                    "Clean up perturb_a2_937087b6_dir by deleting old.bak and leaving doc.md in place.",
                ],
            },
        ],
    },
    "osworld_multi_apps_510f64c8": {
        "archetype": "a2_check_include_exclude",
        "tier": "a2",
        "variants": [
            {
                "params": {
                    "files": {
                        "src_main.py": "print('main')\n",
                        "src_util.py": "print('util')\n",
                        "stale.pyc": "compiled\n",
                    },
                    "remove": ["stale.pyc"],
                    "include": ["src_main.py", "src_util.py"],
                    "exclude": ["stale.pyc"],
                },
                "instr_pool": [
                    "Delete stale.pyc from perturb_a2_510f64c8_dir on the Desktop while keeping src_main.py and src_util.py.",
                    "Remove the stale.pyc compiled artifact from perturb_a2_510f64c8_dir; both src_main.py and src_util.py should stay where they are.",
                    "Clean up perturb_a2_510f64c8_dir by removing stale.pyc and leaving the two .py source files intact.",
                ],
            },
            {
                "params": {
                    "files": {
                        "data.csv": "a,b\n1,2\n",
                        "schema.json": "{\"x\": 1}\n",
                        "snapshot.tmp": "tmp\n",
                    },
                    "remove": ["snapshot.tmp"],
                    "include": ["data.csv", "schema.json"],
                    "exclude": ["snapshot.tmp"],
                },
                "instr_pool": [
                    "Delete snapshot.tmp from perturb_a2_510f64c8_dir on the Desktop, keeping data.csv and schema.json.",
                    "Remove snapshot.tmp from the perturb_a2_510f64c8_dir folder; data.csv and schema.json should stay in place.",
                    "Clean up perturb_a2_510f64c8_dir by deleting snapshot.tmp and leaving the data and schema files intact.",
                ],
            },
        ],
    },
    # ---------- Tier A3: chrome → writer ----------
    "osworld_multi_apps_0c825995": {
        "archetype": "a3_chrome_to_writer",
        "tier": "a3",
        "variants": [
            {
                "params": {
                    "title": "Environmental Policy Brief",
                    "paragraphs": [
                        "Lower carbon-intensity in transportation reduces urban smog levels.",
                        "Investment in public transit yields the largest emissions reduction per dollar.",
                        "Cycling and walking corridors complement transit by handling last-mile demand.",
                    ],
                },
                "instr_pool": [
                    "Read the three sentences shown in the Environmental Policy Brief Chrome tab and append each as its own paragraph into the perturb_a3_0c825995.docx open in LibreOffice Writer, in the same order, preserving the wording exactly.",
                    "Copy the three sentences from the Chrome tab titled 'Environmental Policy Brief' into the perturb_a3_0c825995.docx Writer document as three separate paragraphs in the order shown.",
                    "Take the three policy sentences shown in the Chrome tab and add each as a new paragraph in perturb_a3_0c825995.docx, preserving the original wording and the order on the page.",
                ],
            },
            {
                "params": {
                    "title": "Quarterly Report Highlights",
                    "paragraphs": [
                        "Revenue grew by twelve percent year-over-year.",
                        "Operating margin remained stable at eighteen percent.",
                    ],
                },
                "instr_pool": [
                    "Take the two highlight sentences from the Chrome tab and append each as a new paragraph in perturb_a3_0c825995.docx, keeping the wording verbatim.",
                    "Copy both sentences from the 'Quarterly Report Highlights' Chrome tab into perturb_a3_0c825995.docx as new paragraphs in the same order.",
                    "Read the two report highlights shown in Chrome and add each as its own paragraph in perturb_a3_0c825995.docx, preserving the text exactly.",
                ],
            },
        ],
    },
    "osworld_multi_apps_98e8e339": {
        "archetype": "a3_chrome_to_writer",
        "tier": "a3",
        "variants": [
            {
                "params": {
                    "title": "Project Notes",
                    "paragraphs": [
                        "Note one covers the project setup and dependency installation.",
                        "Note two covers the data ingestion pipeline and validation steps.",
                        "Note three covers the evaluation harness and reporting layout.",
                    ],
                },
                "instr_pool": [
                    "Read the three project notes from the Chrome tab and append each as its own paragraph into perturb_a3_98e8e339.docx, preserving the wording and the original order.",
                    "Copy the three notes from the 'Project Notes' Chrome tab into perturb_a3_98e8e339.docx as separate paragraphs in the same order they appear.",
                    "Take the three project-note sentences shown in Chrome and add each one as a new paragraph in perturb_a3_98e8e339.docx, keeping the text verbatim.",
                ],
            },
            {
                "params": {
                    "title": "Reading List",
                    "paragraphs": [
                        "Item one is the introduction chapter on linear algebra basics.",
                        "Item two is the chapter on probability theory fundamentals.",
                    ],
                },
                "instr_pool": [
                    "Take the two reading items from the Chrome tab and append each as a new paragraph in perturb_a3_98e8e339.docx, preserving the wording.",
                    "Copy the two list items from the 'Reading List' Chrome page into perturb_a3_98e8e339.docx as separate paragraphs in the same order.",
                    "Read the two items in the Chrome tab and add each as its own paragraph at the end of perturb_a3_98e8e339.docx.",
                ],
            },
        ],
    },
    "osworld_multi_apps_aad10cd7": {
        "archetype": "a3_chrome_to_writer",
        "tier": "a3",
        "variants": [
            {
                "params": {
                    "title": "Search Guidelines",
                    "paragraphs": [
                        "Provide informative placeholder text inside every empty search field.",
                        "Use a single search field per page wherever practical.",
                        "Always announce search results in a region that screen readers monitor.",
                    ],
                },
                "instr_pool": [
                    "Copy the three search-guideline sentences shown in the Chrome tab into perturb_a3_aad10cd7.docx as three separate paragraphs in the same order, keeping the wording exactly.",
                    "Read the three guidelines from the 'Search Guidelines' Chrome page and append each as its own paragraph in perturb_a3_aad10cd7.docx, preserving the text verbatim.",
                    "Take the three guideline sentences from Chrome and add each as a new paragraph in perturb_a3_aad10cd7.docx in the same order they appear.",
                ],
            },
            {
                "params": {
                    "title": "Accessibility Tips",
                    "paragraphs": [
                        "Provide visible focus styles for every interactive element.",
                        "Pair each input with a programmatically associated label.",
                    ],
                },
                "instr_pool": [
                    "Copy both accessibility tips from the Chrome tab into perturb_a3_aad10cd7.docx as two new paragraphs in the order shown, keeping the wording verbatim.",
                    "Read the two tips on the 'Accessibility Tips' Chrome page and append each as a separate paragraph in perturb_a3_aad10cd7.docx.",
                    "Take the two tip sentences from Chrome and add each as its own paragraph in perturb_a3_aad10cd7.docx, preserving the original text.",
                ],
            },
        ],
    },
    # ---------- Tier A3: image → archive ----------
    "osworld_multi_apps_46407397": {
        "archetype": "a3_image_to_archive",
        "tier": "a3",
        "variants": [
            {
                "params": {"n_images": 3, "out_size": (80, 60)},
                "instr_pool": [
                    "Open the three images in perturb_a3_46407397_imgs in GIMP, resize each one to 80 by 60 pixels, then bundle the resized copies into a zip archive at /home/user/Desktop/perturb_a3_46407397.zip with the same filenames.",
                    "Resize all three image_*.png files in perturb_a3_46407397_imgs to 80x60 pixels using GIMP, then zip the resized images into /home/user/Desktop/perturb_a3_46407397.zip preserving the filenames.",
                    "Take the three images in perturb_a3_46407397_imgs, downscale each to 80x60 pixels with GIMP, and pack the resized copies into a single zip archive at /home/user/Desktop/perturb_a3_46407397.zip.",
                ],
            },
            {
                "params": {"n_images": 4, "out_size": (100, 80)},
                "instr_pool": [
                    "Resize each of the four image_*.png files in perturb_a3_46407397_imgs to 100 by 80 pixels using GIMP and pack the resized copies into a zip archive at /home/user/Desktop/perturb_a3_46407397.zip with the original filenames.",
                    "Open the four images in perturb_a3_46407397_imgs in GIMP, scale each one to 100x80 pixels, then bundle the resized images into /home/user/Desktop/perturb_a3_46407397.zip.",
                    "Take all four images in perturb_a3_46407397_imgs, resize each to 100x80 pixels with GIMP, and zip the resized files into /home/user/Desktop/perturb_a3_46407397.zip preserving the filenames.",
                ],
            },
        ],
    },
    "osworld_multi_apps_82e3c869": {
        "archetype": "a3_image_to_archive",
        "tier": "a3",
        "variants": [
            {
                "params": {"n_images": 3, "out_size": (120, 90)},
                "instr_pool": [
                    "Resize each image in perturb_a3_82e3c869_imgs to 120 by 90 pixels using GIMP and pack them into a zip archive at /home/user/Desktop/perturb_a3_82e3c869.zip with the original filenames.",
                    "Open the three images in perturb_a3_82e3c869_imgs in GIMP, scale each one to 120x90 pixels, and bundle the resized copies into /home/user/Desktop/perturb_a3_82e3c869.zip.",
                    "Take the three picture files in perturb_a3_82e3c869_imgs, downscale each to 120x90 pixels with GIMP, then zip them into /home/user/Desktop/perturb_a3_82e3c869.zip preserving the filenames.",
                ],
            },
            {
                "params": {"n_images": 2, "out_size": (160, 120)},
                "instr_pool": [
                    "Resize the two image_*.png files in perturb_a3_82e3c869_imgs to 160x120 pixels using GIMP and zip the resized copies into /home/user/Desktop/perturb_a3_82e3c869.zip preserving filenames.",
                    "Open both images in perturb_a3_82e3c869_imgs in GIMP, scale each one to 160x120 pixels, then bundle the resized copies into /home/user/Desktop/perturb_a3_82e3c869.zip.",
                    "Take the two picture files in perturb_a3_82e3c869_imgs, resize each to 160x120 pixels with GIMP, and pack them into /home/user/Desktop/perturb_a3_82e3c869.zip.",
                ],
            },
        ],
    },
    # ---------- Tier A3: vscode + filemanager ----------
    "osworld_multi_apps_91190194": {
        "archetype": "a3_vscode_filemanager",
        "tier": "a3",
        "variants": [
            {
                "params": {
                    "files": {
                        "main.py": "print('hello')\n",
                        "drafts/scratch.txt": "draft notes\n",
                    },
                    "edits": {
                        "main.py": "print('hello world')\n",
                    },
                    "moves": [
                        ("drafts/scratch.txt", "archive/scratch.txt"),
                    ],
                    "include": ["main.py", "archive", "scratch.txt", "hello world"],
                    # validation: a `mv drafts/scratch.txt archive/` correctly empties
                    # drafts/ but leaves the (now-empty) dir behind, so `ls -RF` still
                    # prints its `drafts:` header. The task only asks to MOVE the file,
                    # not delete the folder — so don't exclude `drafts:`. Excluding the
                    # old file path is enough to verify the move. (siblings already do
                    # this — they exclude only the old path.)
                    "exclude": ["drafts/scratch.txt"],
                },
                "instr_pool": [
                    "Open the perturb_a3_91190194_proj folder in VS Code, change main.py so the printed message becomes 'hello world', then move drafts/scratch.txt into a new archive/ folder using the file manager.",
                    "In the perturb_a3_91190194_proj project, edit main.py in VS Code to print 'hello world' instead of 'hello', then move drafts/scratch.txt to archive/scratch.txt with the file manager.",
                    "Update main.py inside perturb_a3_91190194_proj so it prints 'hello world', then relocate drafts/scratch.txt into an archive/ subfolder using the file manager.",
                ],
            },
            {
                "params": {
                    "files": {
                        "config.json": "{\"value\": 1}\n",
                        "old/data.txt": "old data\n",
                    },
                    "edits": {
                        "config.json": "{\"value\": 2}\n",
                    },
                    "moves": [
                        ("old/data.txt", "data/data.txt"),
                    ],
                    "include": ["config.json", "\"value\": 2", "data/", "data.txt"],
                    "exclude": ["\"value\": 1", "old/data.txt"],
                },
                "instr_pool": [
                    "Open perturb_a3_91190194_proj in VS Code, set config.json's value to 2, then move old/data.txt into a new data/ folder via the file manager.",
                    "In the perturb_a3_91190194_proj project, change config.json so 'value' is 2 in VS Code, then relocate old/data.txt to data/data.txt with the file manager.",
                    "Update config.json inside perturb_a3_91190194_proj so the JSON value becomes 2, and move old/data.txt to a new data/ subfolder.",
                ],
            },
        ],
    },
    "osworld_multi_apps_d68204bf": {
        "archetype": "a3_vscode_filemanager",
        "tier": "a3",
        "variants": [
            {
                "params": {
                    "files": {
                        "notes.md": "# old\n",
                        "tmp/log.txt": "tmp log\n",
                    },
                    "edits": {
                        "notes.md": "# release notes\n",
                    },
                    "moves": [
                        ("tmp/log.txt", "logs/log.txt"),
                    ],
                    "include": ["notes.md", "release notes", "logs/", "log.txt"],
                    "exclude": ["# old", "tmp/log.txt"],
                },
                "instr_pool": [
                    "Open perturb_a3_d68204bf_proj in VS Code, change notes.md so the heading becomes 'release notes', then move tmp/log.txt into a new logs/ folder.",
                    "In the perturb_a3_d68204bf_proj project, replace the heading in notes.md with 'release notes' using VS Code, then relocate tmp/log.txt to logs/log.txt via the file manager.",
                    "Update notes.md inside perturb_a3_d68204bf_proj so it reads '# release notes', and move tmp/log.txt into a new logs/ subfolder.",
                ],
            },
            {
                "params": {
                    "files": {
                        "src_main.txt": "todo\n",
                        "drafts/old.txt": "old draft\n",
                    },
                    "edits": {
                        "src_main.txt": "completed\n",
                    },
                    "moves": [
                        ("drafts/old.txt", "archive/old.txt"),
                    ],
                    "include": ["src_main.txt", "completed", "archive/", "old.txt"],
                    "exclude": ["todo", "drafts/old.txt"],
                },
                "instr_pool": [
                    "Open perturb_a3_d68204bf_proj in VS Code, change src_main.txt so its line reads 'completed', then move drafts/old.txt into an archive/ folder.",
                    "Edit src_main.txt inside perturb_a3_d68204bf_proj to say 'completed' using VS Code, then relocate drafts/old.txt to archive/old.txt with the file manager.",
                    "Update src_main.txt inside perturb_a3_d68204bf_proj to read 'completed', and move drafts/old.txt into a new archive/ subfolder.",
                ],
            },
        ],
    },
    "osworld_multi_apps_5bc63fb9": {
        "archetype": "a3_vscode_filemanager",
        "tier": "a3",
        "variants": [
            {
                "params": {
                    "files": {
                        "data.json": "{\"answers\": [\"a\", \"b\", \"c\"]}\n",
                        "drafts/scratch.txt": "scratch\n",
                    },
                    "edits": {
                        "data.json": "{\"answers\": [\"x\", \"y\"]}\n",
                    },
                    "moves": [
                        ("drafts/scratch.txt", "archive/scratch.txt"),
                    ],
                    "include": ["data.json", "\"x\"", "\"y\"", "archive/", "scratch.txt"],
                    "exclude": ["\"a\"", "drafts/scratch.txt"],
                },
                "instr_pool": [
                    "Open perturb_a3_5bc63fb9_proj in VS Code, replace data.json's answers list with ['x', 'y'], then move drafts/scratch.txt into an archive/ folder.",
                    "In the perturb_a3_5bc63fb9_proj project, edit data.json so its answers field becomes ['x', 'y'] using VS Code, then relocate drafts/scratch.txt to archive/scratch.txt.",
                    "Update data.json inside perturb_a3_5bc63fb9_proj so 'answers' is the list ['x', 'y'], and move drafts/scratch.txt to a new archive/ subfolder.",
                ],
            },
            {
                "params": {
                    "files": {
                        "results.json": "{\"score\": 0}\n",
                        "tmp/note.txt": "tmp\n",
                    },
                    "edits": {
                        "results.json": "{\"score\": 100}\n",
                    },
                    "moves": [
                        ("tmp/note.txt", "logs/note.txt"),
                    ],
                    "include": ["results.json", "\"score\": 100", "logs/", "note.txt"],
                    "exclude": ["\"score\": 0", "tmp/note.txt"],
                },
                "instr_pool": [
                    "Open perturb_a3_5bc63fb9_proj in VS Code, change results.json so 'score' becomes 100, then move tmp/note.txt to a new logs/ folder.",
                    "In the perturb_a3_5bc63fb9_proj project, edit results.json so the score is 100 using VS Code, then relocate tmp/note.txt to logs/note.txt with the file manager.",
                    "Update results.json inside perturb_a3_5bc63fb9_proj so 'score' is 100, and move tmp/note.txt into a new logs/ subfolder.",
                ],
            },
        ],
    },
}


def _perturb_a23(eval_row: dict, rng: random.Random) -> list[dict]:
    """Tier A2/A3 dispatcher: emit one perturb row per variant.

    Each row's setup is fully self-contained — `eval_row.config` is stripped
    and replaced with the archetype's `pre_config_steps`. The only thing kept
    from the eval row is the `task_id` (used to anchor the perturb row and
    derive the `/tmp/perturb_a{2,3}_<short>_*` filename namespace).
    """
    tid = eval_row["task_id"]
    spec = _TIER_A23_TASKS.get(tid)
    if spec is None:
        return []

    archetype = spec["archetype"]
    builder = _A23_BUILDERS[archetype]
    short = tid[-8:]
    rows: list[dict] = []
    er_stripped = _strip_eval_config(eval_row)
    for i, variant in enumerate(spec["variants"]):
        pre, perturb_cfg, oracle, evaluator = builder(short, variant["params"])
        instr_raw = rng.choice(variant["instr_pool"])
        instr = _stylize_a23_instruction(instr_raw, rng)
        rows.append(make_perturb_row(
            eval_row=er_stripped,
            knob_assignment={"archetype": archetype, "v": i},
            new_instruction=instr,
            new_oracle=oracle,
            new_evaluator=evaluator,
            pre_config_steps=pre,
            perturb_config_step=perturb_cfg,
            oracle_after_postconfig=False,
        ))
    return rows


# Validation (P0): Tier A1 trim to address +23pp Δfeas multi_apps Layer A
# (archetype monoculture). Pre-trim: 36 bases × 3-5 variants ≈ 139 rows
# dominated multi_apps perturb (75% of 189) — perturb-trained agent saw
# only "Chrome→LO append" workflow. Trim cap reduces over-rep without
# losing pattern coverage; redundant bases (same sink schema) further
# trimmed. Net effect: 139 → ~62 Tier A1 rows; multi_apps total 189 → ~110.
_TIER_A1_MAX_VARIANTS_PER_BASE = 2
_TIER_A1_DROP_BASES: frozenset[str] = frozenset({
    # tally_book.xlsx schema repeated across 5 bases — keep one (3f05f3b9)
    "osworld_multi_apps_415ef462",
    "osworld_multi_apps_42d25c08",
    "osworld_multi_apps_48c46dc7",
    "osworld_multi_apps_788b3701",
    # Zheng He.docx repeated — keep 2c1ebcd7
    "osworld_multi_apps_3a93cae4",
})


def _perturb_chrome_html_to_lo(eval_row: dict, rng: random.Random) -> list[dict]:
    """Tier A1 dispatcher: chrome (local HTML) → LO sink (xlsx/docx).

    Returns one perturb row per variant in spec (capped at
    ``_TIER_A1_MAX_VARIANTS_PER_BASE``). Each row has:
    - pre_config_steps: write /tmp/perturb_<short>.html
    - eval_row.config (modified): chrome_open_tabs URL points to file://...
    - perturb_config_step: generate gold sink file via openpyxl/python-docx
    - oracle: cp gold over result file
    - oracle_after_postconfig=True: oracle runs after agent's Ctrl+S
    """
    tid = eval_row["task_id"]
    if tid in _TIER_A1_DROP_BASES:
        return []
    spec = _TIER_A1_TASKS.get(tid)
    if spec is None:
        return []

    short = tid[-8:]
    sink_path = spec["sink_path"]
    sink_ext = spec["sink_ext"]
    sink_headers = spec["sink_headers"]
    expected_path = f"/tmp/perturb_expected_{short}.{sink_ext}"
    html_path = f"/tmp/perturb_{short}.html"
    file_url = f"file://{html_path}"

    rows = []
    for variant in spec["variants"][:_TIER_A1_MAX_VARIANTS_PER_BASE]:
        if sink_ext == "xlsx":
            html_content = _render_source_html(
                variant["html_title"], sink_headers, variant["rows"],
            )
            # For empty-sink specs, the gold writes the sink_headers as row 1
            # before the data rows (since the original xlsx has no header row).
            prepend = sink_headers if spec.get("sink_starts_empty") else None
            gold_py = _build_xlsx_append_gold_py(
                sink_path, variant["rows"], expected_path,
                prepend_headers=prepend,
            )
            evaluator = {
                "func": "compare_table",
                "result": {"type": "vm_file", "path": sink_path,
                           "dest": sink_path.split("/")[-1]},
                "expected": {"type": "vm_file", "path": expected_path,
                             "dest": "expected_file"},
                "options": {"rules": [{"type": "sheet_data",
                                       "sheet_idx0": 0, "sheet_idx1": "EI0"}]},
                "postconfig": LO_SAVE_POSTCONFIG,
            }
        elif sink_ext == "docx":
            html_content = _render_source_html(
                variant["html_title"], sink_headers, variant["rows"],
            )
            paragraphs = [
                " | ".join(str(c) for c in row) for row in variant["rows"]
            ]
            gold_py = _build_docx_append_gold_py(
                sink_path, paragraphs, expected_path,
            )
            # Tier-A1 docx variants are pure text-append: instructions only
            # ask to "preserve text exactly", never bold/italic/font/etc. Use
            # OSWorld's text-only `compare_docx_files` (paragraph-text match,
            # `examine_*` format fields N/A here) to avoid `compare_docx_strict`
            # tripping on style.name renames that LO Writer's docx round-trip
            # applies to both new and existing paragraphs (e.g., python-docx
            # `add_paragraph` writes style=None which LO renames to
            # "Body Text" on save — false-fail despite identical text).
            evaluator = {
                "func": "compare_docx_files",
                "result": {"type": "vm_file", "path": sink_path,
                           "dest": sink_path.split("/")[-1]},
                "expected": {"type": "vm_file", "path": expected_path,
                             "dest": "expected_file"},
                "postconfig": LO_SAVE_POSTCONFIG,
            }
        elif sink_ext == "pptx":
            # pptx variants use single-column "Slide title" rows. HTML shows
            # the titles to copy; oracle/gold appends one blank slide per
            # title with that text in a centered text box.
            html_content = _render_source_html(
                variant["html_title"], sink_headers, variant["rows"],
            )
            slide_titles = [str(row[0]) for row in variant["rows"]]
            gold_py = _build_pptx_append_slides_py(
                sink_path, slide_titles, expected_path,
            )
            evaluator = {
                "func": "compare_pptx_appended_titles",
                "result": {"type": "vm_file", "path": sink_path,
                           "dest": sink_path.split("/")[-1]},
                "expected": {"type": "vm_file", "path": expected_path,
                             "dest": "expected_file"},
                "postconfig": LO_SAVE_POSTCONFIG,
                # #155 §C #11: the gold builder materializes the last layout's
                # placeholders (several empty shapes) while an LO-Impress agent
                # appends a title-only slide -> upstream compare_pptx_files'
                # UNCONDITIONAL per-slide shape-count gate false-fails the correct
                # title (and the template has a python-pptx-unreadable slide that
                # makes upstream raise). Route to the append-aware comparator:
                # slide count + APPENDED-slide title TEXT only (prefix is
                # LO-normalized and not part of the task, so it is not compared).
                "options": {"num_appended": len(slide_titles)},
            }
        else:
            raise ValueError(f"Tier A1: unsupported sink_ext={sink_ext}")

        pre_step = _build_html_write_step(html_path, html_content)
        gold_step = _make_config_step(gold_py)
        eval_row_mod = _ensure_chrome_with_url(eval_row, file_url)

        # For sink_starts_empty=True, the gold-py overwrites the sink with an
        # empty workbook. Two related fixups vs the default flow:
        #   1. drop the eval-row download step that fetches the populated
        #      upstream sink — otherwise it would clobber the empty file we
        #      write in pre-config.
        #   2. run gold-py in pre_config_steps (before the eval-row `open`
        #      step) so LO opens an already-empty file. Otherwise LO opens the
        #      populated file, keeps it in memory, and the agent's later
        #      Ctrl+S writes the stale populated buffer back over disk.
        # Smoking gun without these two fixes: the agent ends up saving a
        # workbook with the upstream tab (e.g. `2023 ebook list`) plus a new
        # `Sheet2` of their own data; eval reads sheet_idx0=0 → mismatch.
        sink_starts_empty = bool(spec.get("sink_starts_empty"))
        if sink_starts_empty:
            eval_row_mod = _drop_download_for_path(eval_row_mod, sink_path)
            pre_steps_for_row = [pre_step, gold_step]
            perturb_config_step_for_row = None
        else:
            pre_steps_for_row = [pre_step]
            perturb_config_step_for_row = gold_step

        # Oracle: cp gold over result. For docx/pptx, use the 3-step LO-normalize
        # pattern (normalize expected → cp → normalize sink) so the sink's XML
        # structure matches what the eval runner's LO normalization produces on
        # the expected file. Mirrors `_build_oracle` in
        # `perturb/libreoffice_writer.py` and `_standard_oracle` in
        # `perturb/libreoffice_impress.py`. For non-LO sinks (csv/xlsx), the
        # evaluator does not LO-normalize, so plain cp is correct.
        if sink_ext in ("docx", "pptx"):
            oracle = _build_oracle_lo(sink_path, expected_path, fmt=sink_ext)
        else:
            oracle = [_shell_step(f"cp '{expected_path}' '{sink_path}'")]

        # Build paraphrase pool from variant template, pick one, stylize.
        # See `_build_instr_pool` / `_stylize_multi_apps_instruction` docstrings
        # for distribution targets (polite ~32%, save 0%, avg_words 35-45).
        instr_pool = _build_instr_pool(spec, variant)
        instr = _stylize_multi_apps_instruction(
            rng.choice(instr_pool), rng, spec=spec, variant=variant,
        )

        rows.append(make_perturb_row(
            eval_row=eval_row_mod,
            knob_assignment={"variant": variant["html_title"][:30]},
            new_instruction=instr,
            new_oracle=oracle,
            new_evaluator=evaluator,
            pre_config_steps=pre_steps_for_row,
            perturb_config_step=perturb_config_step_for_row,
            oracle_after_postconfig=True,
        ))
    return rows


# ---------------------------------------------------------------------------
# Main dispatcher
# ---------------------------------------------------------------------------

def perturb_multi_apps_per_task(
    eval_row: dict,
    rng: random.Random,
) -> list[dict]:
    # Validation: multi_apps has no TYPE_1/TYPE_2 split (per the
    # _TIER_A1_TASKS docstring above) — every variant is an independent
    # multi-app row, so the legacy `max_type1` arg never had a callsite that
    # used it. The dispatcher (`apply_structural_perturbation`) filters kwargs
    # via `inspect.signature`, so callers passing `max_type1=...` are no-op
    # safe after this removal.
    tid = eval_row["task_id"]
    if tid in _BROKEN_MULTI_TASKS:
        return []

    # Highest-quality policy: only Tier A1 (chrome → LO sink) rows are emitted
    # under the multi_apps label. Bases without a Tier A1 spec produce 0 rows.
    # Rationale: legacy single-app file-op rows (sort/freeze/zoom/bold/font/...)
    # don't exercise multi-app coordination — they belong to the calc/writer/
    # impress single-domain perturbs, and emitting them under multi_apps would
    # dilute the training signal for "must use multiple apps".
    if tid in _TIER_A1_TASKS:
        return _perturb_chrome_html_to_lo(eval_row, rng)
    if tid in _TIER_A23_TASKS:
        return _perturb_a23(eval_row, rng)
    return []
