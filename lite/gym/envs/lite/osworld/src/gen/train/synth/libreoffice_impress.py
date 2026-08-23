"""LibreOffice Impress synth generator (Track A — host-heredoc design).

Per AGENTS.md / libreoffice_impress.md: every row builds BOTH the source
pptx AND the gold pptx via python-pptx heredocs in `pre_config_steps`
(Hard Constraint #2: gold MUST land in metadata.config, not only oracle —
fixes validation 41/41-fail bug). Oracle = `cp gold source`. Evaluator =
`compare_pptx_files` with a single non-default `examine_<field>=True`
(or `check_transition` for transition rows); other examine_* set False
to avoid LO round-trip false-fails on untouched runs. Postconfig =
`LO_SAVE_POSTCONFIG`. `oracle_after_postconfig=True`. Per-row `open_command`
launches LibreOffice Impress on the source pptx.

Implemented row groups (99/99 PASS — `set_title_font_slide1_arial` switched to
`source_font="DejaVu Sans"` so LO's `Liberation Sans → Arial` substitution
on PPTX export doesn't make source's slide-0 font.name match gold's; and
`title_to_bottom` retitled the target slide to the literal "Product Comparison"
required by the eval branch, restructured so gold = original-position pptx +
oracle moves source's slide-2 title down via python-pptx so post-oracle
`result.top > expected.top`):
- Per-shape font color: title (#1-class) and body (#2-class)
                                                     → compare_pptx_files examine_color=True
                                                       (`impress_set_title_color_slide1_navy`,
                                                        `impress_set_title_color_slide2_red`,
                                                        `impress_set_body_color_slide1_green`)
- Per-shape font size                                → compare_pptx_files examine_font_size=True
                                                       (`impress_set_body_size_slide2_24pt`,
                                                        `impress_set_title_size_slide1_32pt`,
                                                        `impress_set_body_size_slide1_24pt`)
- Per-shape font name                                → compare_pptx_files examine_font_name=True
                                                       (`impress_set_title_font_slide1_arial`,
                                                        `impress_set_title_font_slide2_times`)
- Center body text                                   → compare_pptx_files examine_alignment=True
                                                       (`impress_center_body_slide2`)
- Per-run bold / underline / italic                  → compare_pptx_files examine_bold/underline/italic=True
                                                       (`impress_bold_body_slide4`,
                                                        `impress_underline_body_slide4`,
                                                        `impress_italic_body_slide3`)
- Slide bg color                                     → compare_pptx_files examine_background=True
                                                       (`impress_bg_lightblue_slide3`,
                                                        `impress_bg_darkblue_slide1`,
                                                        `impress_bg_pastelgreen_slide2`)
- Speaker notes set                                  → compare_pptx_files examine_note=True
                                                       (`impress_note_office_hours_slide5`,
                                                        `impress_note_qa_reserved_slide5`,
                                                        `impress_note_coach_intro_slide1`)
- Slide reorder (per `rng.sample` distinct slot)     → compare_pptx_files (default)
                                                       (`impress_swap_slides_3_4_training`,
                                                        `impress_swap_slides_1_4_pitch`)
- Insert N×M table                                   → compare_pptx_files (table-presence + dims)
                                                       (`impress_insert_table_grading_5x3`,
                                                        `impress_insert_table_schedule_6x2`,
                                                        `impress_insert_table_measurements_4x3`)
- Image resize (existing inline picture)             → compare_pptx_files examine_image_size=True
                                                       (`impress_image_resize_slide1_4cm`,
                                                        `impress_image_resize_slide3_6cm`)
- Title-to-bottom reposition                         → compare_pptx_files examine_position=True
                                                       (`title_to_bottom`)
- Slide transition (`fade` / `dissolve`)             → check_transition
                                                       (`impress_transition_fade_slide2`,
                                                        `impress_transition_dissolve_slide3`)
- Real-photo skill-templates (Cartesian: 5 topic families × 7 skill ops = 35 templates):
  - image-resize    → compare_pptx_files examine_modify_height=True   (`impress_topic_<topic>_image_resize`)
  - embed-shrink    → compare_pptx_files examine_image_size=True      (`impress_topic_<topic>_embed_single`)
  - set-title       → compare_pptx_files examine_text=True            (`impress_topic_<topic>_set_title`)
  - gallery (color) → compare_pptx_files examine_color_rgb=True       (`impress_topic_<topic>_gallery_4_color_red`)
  - gallery (bold)  → compare_pptx_files examine_font_bold=True       (`impress_topic_<topic>_gallery_4_bold`)
  - caption-size    → compare_pptx_files examine_font_size=True       (`impress_topic_<topic>_caption_size`)
  - reposition      → compare_pptx_files examine_shape=True           (`impress_topic_<topic>_reposition`)
  Topics: see `_TOPIC_FAMILIES` (TopicTheme list). Each theme has its own
  `photo_dirs` whitelist sampled from `photos/<dir>/`. Same template + different
  seed = different photo, but multi-slide decks stay coherent.
- Body-paragraph alignment (center/right/justify)    → compare_pptx_files examine_alignment=True
                                                       (`impress_align_center_body_climate`,
                                                        `impress_align_right_body_postmortem`,
                                                        `impress_align_justify_body_research`)
- Speaker notes on real-world decks                  → compare_pptx_files examine_note=True
                                                       (`impress_note_climate_report_slide1`,
                                                        `impress_note_product_launch_qa`,
                                                        `impress_note_research_methodology_dataset`)
- Image reposition (move existing photo)             → covered by Cartesian
                                                       (`impress_topic_<topic>_reposition`, see above)
- Paragraph indent level                             → compare_pptx_files examine_indent=True
                                                       (`impress_indent_body_healthcare_slide2`)
- Wipe / push transition variants                    → check_transition
                                                       (`impress_transition_wipe_slide1`,
                                                        `impress_transition_push_slide2`,
                                                        `impress_transition_wipe_slide3_postmortem`)

Deck diversification (2026-05-10): existing 6 templates were retargeted
onto richer per-domain decks (Mars Mission, Q1 Postmortem, Wildlife Conservation,
Climate Report, Product Launch, Customer Stories) so the same edit op is
exercised across distinct semantic domains rather than reused boilerplate.
Template_ids unchanged; only `titles=`/`body_lines=` updated. The Deck
pool is declared after `_DECK_TECH` (`_DECK_MARS_MISSION`, `_DECK_POSTMORTEM`,
`_DECK_WILDLIFE_CONSERVATION`, `_DECK_CLIMATE_REPORT`, `_DECK_PRODUCT_LAUNCH`,
`_DECK_CUSTOMER_STORIES`, `_DECK_AI_ETHICS`, `_DECK_HEALTHCARE_ANALYTICS`,
`_DECK_RESEARCH_METHODOLOGY`).

Deferred (per devs/envs/lite.osworld/synth/libreoffice_impress.md `## Implementation status`):
- Chart insert rows — chart XML doesn't round-trip cleanly through LO save
  (F11); perturb has no synth-side chart precedent either.
- Audio embed rows (`compare_audios`) — needs real .mp3/.wav asset.
- PDF / image export rows (`compare_pdfs` / `compare_images`) — agent-side
  LO-Export-As GUI flow is fragile and has no oracle-stable variant on the
  synth side yet.
- Theme apply / master-slide edit / animation rows — gold is ill-defined
  via python-pptx alone; perturb side covers via TYPE_3 wireframes only.

Re-enable plan: slide-image rows now wired via `_stage_asset` (host_push);
investigate LO-stable chart XML construction; use `soffice --convert-to pdf`
for PDF-export gold.

Usage:
    uv run python -m lite.gym.envs.lite.osworld.src.gen.train \\
        --track synth --domain libreoffice_impress
"""

from __future__ import annotations

import random
import re
import textwrap
from dataclasses import dataclass, field
from typing import Callable

from lite.gym.envs.lite.osworld.src.gen.common import (
    LO_SAVE_POSTCONFIG,
)
from lite.gym.envs.lite.osworld.src.gen.train.synth._utils import (
    SynthTemplate,
    _ASSET_ROOT,
    _stage_asset,
    _stable_hash,
)


# ---------------------------------------------------------------------------
# LO post-launch settle steps (validation — Trigger F race fix). Owned by the
# LO domain per the per-domain setup convention; previously lived in
# common.py and was applied via a `template.domain == "libreoffice_*"`
# conditional.
# ---------------------------------------------------------------------------
_LO_POSTLAUNCH_SETTLE: list[dict] = [
    {"type": "sleep", "parameters": {"seconds": 6}},
    {"type": "activate_window", "parameters": {"window_name": "LibreOffice"}},
]


# Body-focus post-open step (libreoffice_impress Axis A H-trigger):
# after LO Impress finishes launching, force keyboard focus into the slide
# canvas and dismiss any startup banner / Tip-of-the-Day dialog so the
# subsequent agent keystrokes (e.g. Tab to body placeholder, F5 slideshow,
# Ctrl+S save dialog) land in the slide editor rather than in a stray
# transient overlay. Escape twice dismisses tip dialogs + clears any
# auto-focused placeholder; F6 cycles focus back to the document area.
# Used by the 6 base title/body-style templates that are particularly
# brittle to focus drift: imp_04 / imp_05 / imp_12 / imp_25 / imp_30 / imp_32.
_IMPRESS_BODY_FOCUS_STEPS: list[dict] = [
    *_LO_POSTLAUNCH_SETTLE,
    {"type": "execute", "parameters": {
        "command": (
            "WID=$(xdotool search --name 'LibreOffice Impress' 2>/dev/null | head -1); "
            "if [ -n \"$WID\" ]; then xdotool windowactivate --sync \"$WID\"; fi; true"
        ),
        "shell": True,
    }},
    {"type": "key", "parameters": {"key": "Escape"}},
    {"type": "sleep", "parameters": {"seconds": 0.2}},
    {"type": "key", "parameters": {"key": "Escape"}},
    {"type": "key", "parameters": {"key": "F6"}},
    {"type": "sleep", "parameters": {"seconds": 0.3}},
]


# ---------------------------------------------------------------------------
# Heredoc / step helpers
# ---------------------------------------------------------------------------

_DESKTOP = "/home/user/Desktop"


def _execute(command: str, *, shell: bool = True) -> dict:
    return {"type": "execute", "parameters": {"command": command, "shell": shell}}


def _py_step(py_code: str) -> dict:
    """Wrap a Python heredoc as an `execute` config step."""
    return _execute(f"python3 << 'PYEOF'\n{py_code}\nPYEOF")


# Header preamble shared by every source/gold builder.
_PPTX_PREAMBLE = textwrap.dedent("""\
    from pptx import Presentation
    from pptx.util import Pt, Cm, Inches, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from lxml import etree
""")


def _build_pptx_step(out_path: str, builder_body: str) -> dict:
    """Run a python-pptx heredoc that constructs a Presentation and saves to out_path."""
    py = (
        _PPTX_PREAMBLE
        + textwrap.dedent(f"""
            prs = Presentation()
            # 16:9 widescreen by default; matches LO Impress default since 7.x.
            prs.slide_width = Cm(25.4)
            prs.slide_height = Cm(14.29)
            path = {out_path!r}
        """)
        + "\n"
        + textwrap.dedent(builder_body)
        + "\nprs.save(path)\n"
    )
    return _py_step(py)


# ---------------------------------------------------------------------------
# Oracle helpers — symmetric LO normalize around the cp-plant. Mirrors perturb's
# `_lo_normalize_cmd` / `_standard_oracle` at perturb/libreoffice_impress.py:124
# (and synth/libreoffice_writer.py's `_build_oracle`).
# ---------------------------------------------------------------------------

def _lo_normalize_cmd(path: str, fmt: str = "pptx") -> str:
    """Headless LO round-trip a pptx in place. Byte-identical to perturb's
    `_lo_normalize_cmd` at perturb/libreoffice_impress.py:124 (just default
    fmt). Used by `_build_oracle_pptx` to symmetrize the eval-time normalize
    chain so the cp-planted expected matches the agent's LO_SAVE_POSTCONFIG'd
    result on round-trip-mutated fields (font name canonicalisation,
    paragraph reformatting, image byte-encoding)."""
    return (
        f"tmpd=$(mktemp -d) && "
        f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
        f"--convert-to {fmt} --outdir \"$tmpd\" '{path}' 2>/dev/null && "
        f"[ -f \"$tmpd/$(basename '{path}')\" ] && "
        f"cp \"$tmpd/$(basename '{path}')\" '{path}'; "
        f"rm -rf \"$tmpd\"; true"
    )


def _build_oracle_pptx(out_path: str, expected_path: str) -> list[dict]:
    """Oracle: ① normalize gold, ② plant at sink.

    oracle_after_postconfig=True kills LO before running oracle actions and
    sets _postconfig_done=True so LO_SAVE_POSTCONFIG (Ctrl+S) is skipped after
    oracle. Result = cp of norm(gold). Expected = norm(gold). They are
    byte-identical without a third normalize step.  A third normalize would
    produce result = norm(norm(gold)) which diverges when LO normalize is
    non-idempotent (font-name, size, colour, alignment properties in pptx).
    """
    return [
        _execute(_lo_normalize_cmd(expected_path, "pptx")),
        _execute(f"cp '{expected_path}' '{out_path}'"),
    ]


# ---------------------------------------------------------------------------
# Source-deck builders — each emits python-pptx body that adds slides with a
# title text frame and (optionally) a body text frame. Source decks cap text-
# frames at 1-2/slide (F9 mitigation).
# ---------------------------------------------------------------------------

def _slide_deck_body(
    titles: list[str], body_lines: list[str], *, source_font: str = "Liberation Sans",
) -> str:
    """Build a python-pptx body that adds slides with title + body text frames.

    `source_font` is the font.name written for every run. Default 'Liberation
    Sans' is the LO Linux default. For tests where the gold mutates a slide's
    font.name to a Latin font that LO substitutes from Liberation Sans on PPTX
    export (e.g. 'Arial' — LO writes "Arial" to the typeface attribute when
    exporting Liberation Sans for cross-platform compatibility), pass a
    non-substituted font (e.g. 'DejaVu Sans') so source != gold post-LO-save.
    """
    assert len(titles) == len(body_lines), "titles/body_lines length must match"
    py_lines = ["blank = prs.slide_layouts[6]"]
    for title, body in zip(titles, body_lines):
        py_lines.append("slide = prs.slides.add_slide(blank)")
        py_lines.append("tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(22.0), Cm(2.5))")
        py_lines.append("tf = tb.text_frame")
        py_lines.append("tf.word_wrap = True")
        py_lines.append("tf.text = " + repr(title))
        py_lines.append("tf.paragraphs[0].runs[0].font.size = Pt(28)")
        py_lines.append(f"tf.paragraphs[0].runs[0].font.name = {source_font!r}")
        # Explicit BLACK baseline so compare_pptx_files' per-run color check
        # ENGAGES. Upstream skips the color compare when hasattr(run.font.color,
        # "rgb") is False (a color-less run), which lets a color-less start deck
        # vacuously equal a colored gold -> trivial_pass. Gold decks replay this
        # body then override the target run's color via gold_mutate, so start
        # (black) != gold (colored). RGBColor already imported by _PPTX_PREAMBLE.
        py_lines.append("tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)")
        py_lines.append("bb = slide.shapes.add_textbox(Cm(1.0), Cm(4.0), Cm(22.0), Cm(8.0))")
        py_lines.append("bf = bb.text_frame")
        py_lines.append("bf.word_wrap = True")
        py_lines.append("bf.text = " + repr(body))
        # Validation fix (CRITICAL): empty body string yields paragraphs[0]
        # with NO runs, so .runs[0] raises IndexError. Guard the per-run
        # font-set lines so gallery templates (which pass body_lines=[""]) and
        # any future title-only deck don't crash at heredoc execution time.
        if " " not in repr(body) and repr(body) in ("''", '""'):
            # body is empty — set a single-char placeholder so eval still has
            # a run to inspect when examine_text is on.
            py_lines.append("bf.text = '\\xa0'  # NBSP placeholder for empty-body slide")
        py_lines.append("if bf.paragraphs[0].runs:")
        py_lines.append("    bf.paragraphs[0].runs[0].font.size = Pt(18)")
        py_lines.append(f"    bf.paragraphs[0].runs[0].font.name = {source_font!r}")
        py_lines.append("    bf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)")
    return "\n".join(py_lines)


# ---------------------------------------------------------------------------
# examine_<field> options bundle for `compare_pptx_files` — picks ONE skill
# dimension and suppresses the rest so LO round-trip artifacts on untouched
# runs don't false-fail. Mirrors perturb's `_BASE_OPTS` shape.
# ---------------------------------------------------------------------------

# Canonical valid `examine_*` keys read by
# `desktop_env.evaluators.metrics.slides.compare_pptx_files`. Anything not in
# this set is a SILENT NO-OP at eval time. The whitelist is enforced by
# `_examine_options` (raises on unknown keys) and surfaces typos / stale
# eval_class strings before they corrupt the train jsonl.
_VALID_EXAMINE_FIELDS: frozenset[str] = frozenset({
    "examine_number_of_slides", "examine_shape", "examine_text",
    "examine_indent", "examine_font_name", "examine_font_size",
    "examine_font_bold", "examine_font_italic", "examine_color_rgb",
    "examine_font_underline", "examine_strike_through", "examine_alignment",
    "examine_title_bottom_position", "examine_table_bottom_position",
    "examine_run_count", "examine_right_position", "examine_top_position",
    "examine_shape_for_shift_size", "examine_image_size",
    "examine_modify_height", "examine_bullets", "examine_background_color",
    "examine_note",
})


# Validation landmine fix:
#
# Before this fix, many `Param.examine_field` slots carried synth's eval_class
# strings (e.g. `bold_underline_text`, `set_font_color`, `edit_title`,
# `change_bg_color`) — none of which are real `compare_pptx_files` keys, so
# they were SILENT no-ops at eval time. Tasks passed only because gold mutated
# exactly the right attribute and source matched everywhere else under the
# strict-default `examine_*` fan-out.
#
# Each Param's `examine_field` now MUST be a canonical `examine_*` key
# (whitelist `_VALID_EXAMINE_FIELDS`). The mapping below documents the
# eval_class → canonical-field correspondence for the four invalid keys
# observed during validation (informational; per-Param rewrites use the
# gold-mutator function as the actual ground truth — `_gold_set_*_font_color`
# → `examine_color_rgb`, `_gold_set_*_font_size` → `examine_font_size`, etc.).
EVAL_CLASS_TO_EXAMINE_FIELD: dict[str, tuple[str, ...]] = {
    # set_font_color: gold mutates run.font.color.rgb on title or body shapes.
    "set_font_color": ("examine_color_rgb",),
    # edit_title: covers font-size / font-name / speaker-note / title-text edits.
    # Per-Param choice depends on the gold helper used; canonical default is
    # examine_text but specific helpers map to font_size / font_name / note.
    "edit_title": ("examine_text", "examine_font_size", "examine_font_name", "examine_note"),
    # change_bg_color: gold mutates slide.background fill.
    "change_bg_color": ("examine_background_color",),
    # bold_underline_text: format-op bucket — covers bold / underline / italic /
    # alignment depending on the gold helper.
    "bold_underline_text": ("examine_font_bold", "examine_font_underline",
                             "examine_font_italic", "examine_alignment"),
    # compound_pptx: paired-mutation params — opens 2+ flags via extra_examine.
    "compound_pptx": ("examine_font_bold", "examine_font_underline",
                      "examine_color_rgb", "examine_background_color",
                      "examine_alignment"),
    # image_stretch: gold resizes a picture; eval reads picture height delta.
    "image_stretch": ("examine_modify_height",),
    # reorder_slides: gold swaps slide indices; eval reads per-slide text order.
    "reorder_slides": ("examine_text",),
    # add_slide: D-IMP-51 inserts a table shape; eval reads shape presence.
    "add_slide": ("examine_shape",),
}


def _examine_options(field: str | tuple[str, ...]) -> dict:
    """Examine ONE skill dimension; suppress the rest so LO round-trip artifacts
    on untouched runs don't false-fail. Mirrors perturb's `_BASE_OPTS` shape.

    `field` is a canonical `examine_*` key from `_VALID_EXAMINE_FIELDS` (or a
    tuple of such keys for compound-mutation Params). Any unknown key raises
    `AssertionError` so eval_class typos can't silently regress to no-op."""
    opts = {
        "examine_shape": False,
        "examine_run_count": False,
        "examine_indent": False,
        "examine_strike_through": False,
        "examine_bullets": False,
        "examine_text": True,
    }
    if not field:
        return opts
    fields = (field,) if isinstance(field, str) else tuple(field)
    for f in fields:
        assert f in _VALID_EXAMINE_FIELDS, (
            f"_examine_options: {f!r} is not a valid compare_pptx_files key. "
            f"Valid keys: {sorted(_VALID_EXAMINE_FIELDS)}"
        )
        opts[f] = True
    return opts


# ---------------------------------------------------------------------------
# Targeted gold-mutation snippets — applied to `prs.slides[idx]` in the gold
# heredoc. Mirror perturb's `_expected_snippet` shape. Used by §I FileTasks.
# ---------------------------------------------------------------------------

def _gold_set_font_color(idx: int, rgb: tuple[int, int, int]) -> str:
    r, g, b = rgb
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        _r.font.color.rgb = RGBColor({r}, {g}, {b})
    """


def _gold_set_font_size(idx: int, pt: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        _r.font.size = Pt({pt})
    """


# Validation fix: title-only mutators for templates whose instruction reads
# "the title" but the legacy `_gold_set_font_*` loops over ALL shapes on the
# slide. Body shapes' font_size diff causes FALSE-FAIL when the agent (per
# instruction) only edits the title textbox. shape[0] is the title textbox
# placed first in `_slide_deck_body`.

def _gold_set_title_font_size(idx: int, pt: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _tb = _sl.shapes[0]
        if _tb.has_text_frame:
            for _r in _tb.text_frame.paragraphs[0].runs:
                _r.font.size = Pt({pt})
    """


def _gold_set_title_font_color(idx: int, rgb: tuple[int, int, int]) -> str:
    r, g, b = rgb
    return f"""\
        _sl = prs.slides[{idx}]
        _tb = _sl.shapes[0]
        if _tb.has_text_frame:
            for _r in _tb.text_frame.paragraphs[0].runs:
                _r.font.color.rgb = RGBColor({r}, {g}, {b})
    """


def _gold_set_font_name(idx: int, name: str) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        _r.font.name = {name!r}
    """


# Validation note (Fix 1): localize font-name mutation to title-only / body-
# only on the target slide. Previously `_gold_set_font_name` rewrote EVERY shape
# on the slide, which mismatches an agent that correctly follows a "title font"
# or "body font" instruction (the un-touched shape's run font.name diverges →
# eval=0). Helpers below mirror `_gold_set_title_font_color`'s shape-0 access
# pattern but for font.name, plus a body-only (shape-1) variant for compound /
# body-font specs. Use these in place of `_gold_set_font_name` whenever the
# instruction localizes the change.
def _gold_set_title_font_name(idx: int, name: str) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _tb = _sl.shapes[0]
        if _tb.has_text_frame:
            for _r in _tb.text_frame.paragraphs[0].runs:
                _r.font.name = {name!r}
    """


def _gold_set_body_font_name(idx: int, name: str) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _bb = _sl.shapes[1]
        if _bb.has_text_frame:
            for _p in _bb.text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.name = {name!r}
    """


def _gold_set_caption_font_name(idx: int, name: str) -> str:
    """Set font name on the caption textbox (shapes[2]) on slide `idx`.

    Hero-photo slides: shapes[0]=title, shapes[1]=photo, shapes[2]=caption.
    """
    return f"""\
        _sl = prs.slides[{idx}]
        if len(_sl.shapes) > 2 and _sl.shapes[2].has_text_frame:
            for _p in _sl.shapes[2].text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.name = {name!r}
    """


# Validation note: body-only helpers (Fix 2). When instruction
# says "the body text" but gold mutates all shapes including title, agent
# following instruction correctly still fails. This was fixed first for
# font_name, then extended to bold/color/alignment.
def _gold_set_body_bold(idx: int) -> str:
    # When instruction says "the body text" but gold mutates all
    # shapes including title, agent following instruction correctly still fails
    # (fixed first for font_name, then extended to bold/color/alignment).
    return f"""\
        _sl = prs.slides[{idx}]
        _bb = _sl.shapes[1]
        if _bb.has_text_frame:
            for _p in _bb.text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.bold = True
    """


def _gold_set_body_font_color(idx: int, rgb: tuple[int, int, int]) -> str:
    # When instruction says "the body text" but gold mutates all
    # shapes including title, agent following instruction correctly still fails
    # (fixed first for font_name, then extended to bold/color/alignment).
    r, g, b = rgb
    return f"""\
        _sl = prs.slides[{idx}]
        _bb = _sl.shapes[1]
        if _bb.has_text_frame:
            for _p in _bb.text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.color.rgb = RGBColor({r}, {g}, {b})
    """


def _gold_set_body_text_alignment(idx: int, align_name: str) -> str:
    # When instruction says "the body text" but gold mutates all
    # shapes including title, agent following instruction correctly still fails
    # (fixed first for font_name, then extended to bold/color/alignment).
    return f"""\
        _sl = prs.slides[{idx}]
        _bb = _sl.shapes[1]
        if _bb.has_text_frame:
            for _p in _bb.text_frame.paragraphs:
                _p.alignment = PP_ALIGN.{align_name}
    """


def _gold_set_text_alignment(idx: int, align_name: str) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    _p.alignment = PP_ALIGN.{align_name}
    """


def _gold_set_bold(idx: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        _r.font.bold = True
    """


def _gold_set_italic(idx: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        _r.font.italic = True
    """


def _gold_set_underline(idx: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.has_text_frame:
                for _p in _sh.text_frame.paragraphs:
                    for _r in _p.runs:
                        _r.font.underline = True
    """


def _gold_set_background(idx: int, rgb: tuple[int, int, int]) -> str:
    r, g, b = rgb
    return f"""\
        _sl = prs.slides[{idx}]
        _bg = _sl.background
        _bg.fill.solid()
        _bg.fill.fore_color.rgb = RGBColor({r}, {g}, {b})
    """


def _gold_set_speaker_note(idx: int, note: str) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _sl.notes_slide.notes_text_frame.text = {note!r}
    """


def _gold_swap_slides(idx_a: int, idx_b: int) -> str:
    """Swap two slide IDs in the slide list (operates on the XML sldIdLst).

    Mirrors perturb's `reorder_slides` snippet shape but does a true 2-way swap
    rather than a single move; LO round-trips slide-id ordering directly.
    """
    lo, hi = sorted((idx_a, idx_b))
    return f"""\
        _lst = prs.slides._sldIdLst
        _items = list(_lst)
        _a, _b = _items[{lo}], _items[{hi}]
        # Rebuild the list with the two entries swapped.
        _new = list(_items)
        _new[{lo}], _new[{hi}] = _b, _a
        for _x in _items:
            _lst.remove(_x)
        for _x in _new:
            _lst.append(_x)
    """


def _gold_title_to_bottom(idx: int) -> str:
    """Move the first shape (our title textbox) on slide `idx` down so its
    `top` is well past slide_height/2 — required by examine_title_bottom_position
    eval branch (`shape1.top > shape2.top` AND `shape1.top > 3600000` Emu).
    """
    return f"""\
        _sl = prs.slides[{idx}]
        _sw = prs.slide_width
        _shh = prs.slide_height
        _shape = _sl.shapes[0]
        _shape.top = _shh - _shape.height - Emu(457200)
    """


# --- Title-only / body-only attr setters (PD 1b: gold mutates only the
# textbox the instruction names — broad `_gold_set_bold` etc. would mutate
# both title and body, which a title-only agent can't reproduce). Title is
# `shapes[0]`, body/caption is `shapes[1]` (matches `_slide_deck_body`).

def _gold_set_title_bold(idx: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _sh = _sl.shapes[0]
        if _sh.has_text_frame:
            for _p in _sh.text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.bold = True
    """


def _gold_set_title_italic(idx: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _sh = _sl.shapes[0]
        if _sh.has_text_frame:
            for _p in _sh.text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.italic = True
    """


def _gold_set_title_underline(idx: int) -> str:
    return f"""\
        _sl = prs.slides[{idx}]
        _sh = _sl.shapes[0]
        if _sh.has_text_frame:
            for _p in _sh.text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.underline = True
    """


def _gold_set_caption_italic(idx: int) -> str:
    """Italicize only the caption textbox (shapes[2]) on slide `idx`.

    Hero-photo slides: shapes[0]=title, shapes[1]=photo, shapes[2]=caption.
    """
    return f"""\
        _sl = prs.slides[{idx}]
        if len(_sl.shapes) > 2 and _sl.shapes[2].has_text_frame:
            for _p in _sl.shapes[2].text_frame.paragraphs:
                for _r in _p.runs:
                    _r.font.italic = True
    """


def _gold_set_title_text_alignment(idx: int, align_name: str) -> str:
    """Set alignment of the title textbox (shapes[0]) only."""
    return f"""\
        _sl = prs.slides[{idx}]
        _sh = _sl.shapes[0]
        if _sh.has_text_frame:
            for _p in _sh.text_frame.paragraphs:
                _p.alignment = PP_ALIGN.{align_name}
    """


# --- D-IMP-57..67 new gold-mutate helpers (P2 image-size, P3 element-move,
# P4 multi-slide subset, P5 audio, P6 long-tail eval funcs) -----------------

def _gold_resize_picture_wh(idx: int, w_cm: float, h_cm: float) -> str:
    """Resize the FIRST picture on slide `idx` to (w_cm, h_cm). Differs from
    `_gold_resize_picture` (which only sets height) so that examine_image_size
    branch sees both width and height as the diff signal."""
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.shape_type == 13:
                _sh.width = Cm({w_cm})
                _sh.height = Cm({h_cm})
                break
    """


def _gold_move_picture(idx: int, left_cm: float, top_cm: float) -> str:
    """Move the FIRST picture on slide `idx` to (left_cm, top_cm). Width and
    height unchanged. Pairs with examine_shape=True so the position diff is
    the only visible signal between source and gold."""
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.shape_type == 13:
                _sh.left = Cm({left_cm})
                _sh.top = Cm({top_cm})
                break
    """


def _gold_move_title(idx: int, left_cm: float, top_cm: float) -> str:
    """Move the title textbox (shapes[0]) on slide `idx` to (left_cm, top_cm).
    Use with examine_shape=True so the move is detected."""
    return f"""\
        _sl = prs.slides[{idx}]
        _sh = _sl.shapes[0]
        _sh.left = Cm({left_cm})
        _sh.top = Cm({top_cm})
    """


def _gold_multi_slide_title_color(slide_idxs: list[int], rgb: tuple[int, int, int]) -> str:
    """Apply title-color mutation to a NON-CONTIGUOUS subset of slides. Echoes
    eval task `4e...` "Set color of titles in slides 2,3,5". Uses
    examine_color_rgb — the body shapes on the untouched slides also have
    runs but their color stays unchanged on both source and gold so it
    doesn't trip the per-shape color diff (the helper only mutates the
    title shape's run color)."""
    r, g, b = rgb
    idx_lit = ", ".join(str(i) for i in slide_idxs)
    return f"""\
        for _i in [{idx_lit}]:
            _sl = prs.slides[_i]
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.color.rgb = RGBColor({r}, {g}, {b})
    """


def _gold_multi_slide_title_bold_underline(slide_idxs: list[int]) -> str:
    """Bold AND underline the title on a non-contiguous subset of slides."""
    idx_lit = ", ".join(str(i) for i in slide_idxs)
    return f"""\
        for _i in [{idx_lit}]:
            _sl = prs.slides[_i]
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.bold = True
                    _r.font.underline = True
    """


def _gold_multi_slide_title_italic(slide_idxs: list[int]) -> str:
    """Italicize the title on a non-contiguous subset of slides."""
    idx_lit = ", ".join(str(i) for i in slide_idxs)
    return f"""\
        for _i in [{idx_lit}]:
            _sl = prs.slides[_i]
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.italic = True
    """


def _gold_all_slides_title_color(rgb: tuple[int, int, int]) -> str:
    """Doc-wide framing variant of `_gold_multi_slide_title_color`: apply the
    title color to EVERY slide in the deck (libreoffice_impress
    `slide_anchor.doc_wide` bridge). Instructions phrased "all titles" /
    "across the deck" rather than "on slide N"."""
    r, g, b = rgb
    return f"""\
        for _sl in prs.slides:
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.color.rgb = RGBColor({r}, {g}, {b})
    """


def _gold_all_slides_title_bold() -> str:
    """Bold the title text on every slide in the deck (doc-wide framing)."""
    return """\
        for _sl in prs.slides:
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.bold = True
    """


def _gold_all_slides_title_underline() -> str:
    """Underline the title text on every slide in the deck (doc-wide framing)."""
    return """\
        for _sl in prs.slides:
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.underline = True
    """


def _gold_all_slides_title_font_name(name: str) -> str:
    """Set the title-text font family on EVERY slide (doc-wide framing).
    Used by implicit-voice "Change the title font ..." instructions."""
    return f"""\
        for _sl in prs.slides:
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.name = {name!r}
    """


def _gold_all_slides_title_alignment(align_name: str) -> str:
    """Set the title-paragraph alignment on EVERY slide (doc-wide framing)."""
    return f"""\
        from pptx.enum.text import PP_ALIGN as _PP_ALIGN_M
        for _sl in prs.slides:
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                _tb.text_frame.paragraphs[0].alignment = getattr(_PP_ALIGN_M, {align_name!r})
    """


def _gold_all_slides_append_body_bullet(text: str) -> str:
    """Append a paragraph to the body textbox (shapes[1]) on EVERY slide
    that has a body. Used by implicit-voice "Add a bullet point ..."
    instructions on multi-slide decks (mirrors eval-row 46 / f23acfd2)."""
    return f"""\
        for _sl in prs.slides:
            if len(_sl.shapes) > 1 and _sl.shapes[1].has_text_frame:
                _tf = _sl.shapes[1].text_frame
                _new_p = _tf.add_paragraph()
                _new_p.text = {text!r}
    """


def _gold_all_slides_background(rgb: tuple[int, int, int]) -> str:
    """Set the slide background fill RGB on EVERY slide (doc-wide framing).
    Matches eval `examine_background_color` per-slide diff but the instruction
    targets all slides rather than one ordinal."""
    r, g, b = rgb
    return f"""\
        from pptx.util import Pt as _Pt  # noqa: F401
        from pptx.dml.color import RGBColor as _RGBColor
        for _sl in prs.slides:
            _bg = _sl.background
            _fill = _bg.fill
            _fill.solid()
            _fill.fore_color.rgb = _RGBColor({r}, {g}, {b})
    """


# --- validation add_slide / table / bullet / body-text helpers ---------------
# Echoes eval rows: 22 (Add "Page 1"), 28 (duplicate last two slides), 36
# (Summary Slide), 38 (six blank slides), 46 (add bullet point), 33 (move
# table), 18 (change first row of table). Each helper mutates `prs` in-place
# the same way the existing _gold_* helpers do, so it composes with the
# default factory body.replace(prs.save) machinery.

def _gold_add_blank_slide_at_end() -> str:
    """Append a single blank slide at the end of the deck. Used for
    add_slide tasks where gold has exactly one more slide than source."""
    return """\
        _blank = prs.slide_layouts[6]
        prs.slides.add_slide(_blank)
    """


def _gold_add_blank_slides_at_end(n: int) -> str:
    """Append `n` blank slides at the end of the deck (eval row 38 pattern)."""
    return f"""\
        _blank = prs.slide_layouts[6]
        for _ in range({n}):
            prs.slides.add_slide(_blank)
    """


def _gold_duplicate_last_slide() -> str:
    """Append a copy of the last slide's title+body text content as a new
    blank slide. Exact XML-level duplication is fragile in python-pptx, so we
    re-create the title/body shapes from the last slide's text. Matches the
    eval-row-28 pattern ("duplicate the last slide") well enough that a deck
    with N+1 slides ends up shape-wise indistinguishable from source's last
    slide's textual content."""
    return """\
        _blank = prs.slide_layouts[6]
        _last = prs.slides[-1]
        _texts = []
        for _sh in _last.shapes:
            if _sh.has_text_frame:
                _texts.append(_sh.text_frame.text)
        _new = prs.slides.add_slide(_blank)
        if len(_texts) >= 1:
            _tb = _new.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(22.0), Cm(2.5))
            _tb.text_frame.text = _texts[0]
        if len(_texts) >= 2:
            _bb = _new.shapes.add_textbox(Cm(1.0), Cm(4.0), Cm(22.0), Cm(8.0))
            _bb.text_frame.text = _texts[1]
    """


def _gold_add_summary_slide(summary_title: str = "Summary") -> str:
    """Append a 'Summary Slide' at the end whose body lists every previous
    slide's title (eval row 36 pattern). The agent's task is to add an
    Impress 'Summary Slide' at the end of the deck."""
    return f"""\
        _titles = []
        for _sl in prs.slides:
            if _sl.shapes and _sl.shapes[0].has_text_frame:
                _titles.append(_sl.shapes[0].text_frame.text)
        _blank = prs.slide_layouts[6]
        _new = prs.slides.add_slide(_blank)
        _tb = _new.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(22.0), Cm(2.5))
        _tb.text_frame.text = {summary_title!r}
        _bb = _new.shapes.add_textbox(Cm(1.0), Cm(4.0), Cm(22.0), Cm(8.0))
        _bb.text_frame.text = "\\n".join(_titles)
    """


def _gold_set_body_text(idx: int, text: str) -> str:
    """Overwrite the body textbox (shapes[1]) text on slide `idx`. Used for
    "add Page 1 to content" style tasks (eval row 22) where the gold body
    text is the source body text replaced with the target string. Reuses
    examine_text=True default to detect the change."""
    return f"""\
        _sl = prs.slides[{idx}]
        if len(_sl.shapes) > 1 and _sl.shapes[1].has_text_frame:
            _sl.shapes[1].text_frame.text = {text!r}
    """


def _gold_append_to_body(idx: int, text: str) -> str:
    """Append a paragraph to the body textbox (shapes[1]) on slide `idx`.
    Used for bullet/paragraph-add tasks (eval row 46). The new paragraph is
    appended on a new line so the existing body text is preserved."""
    return f"""\
        _sl = prs.slides[{idx}]
        if len(_sl.shapes) > 1 and _sl.shapes[1].has_text_frame:
            _tf = _sl.shapes[1].text_frame
            _new_p = _tf.add_paragraph()
            _new_p.text = {text!r}
    """


def _gold_insert_table_with_headers(idx: int, rows_n: int, cols_n: int,
                                     headers: list[str]) -> str:
    """Insert a table on slide `idx` with the first row pre-filled with
    `headers`. Used by table_change_row variants — gold's first row text
    diverges from source (which has no table at all)."""
    headers_lit = repr(headers)
    return f"""\
        _sl = prs.slides[{idx}]
        _tbl_shape = _sl.shapes.add_table({rows_n}, {cols_n},
            Cm(2.5), Cm(8.0), Cm(20.0), Cm(5.0))
        _hdrs = {headers_lit}
        for _r_i, _row in enumerate(_tbl_shape.table.rows):
            for _c_i, _cell in enumerate(_row.cells):
                if _r_i == 0 and _c_i < len(_hdrs):
                    _cell.text = _hdrs[_c_i]
                else:
                    _cell.text = ''
    """


def _gold_insert_table_at_pos(idx: int, rows_n: int, cols_n: int,
                               left_cm: float, top_cm: float) -> str:
    """Insert a table on slide `idx` at (left_cm, top_cm). Used by
    table_move variants — gold's table is at a non-default position (e.g.
    bottom-of-slide) so the "move the table to the bottom" instruction
    has a stable expected geometry."""
    return f"""\
        _sl = prs.slides[{idx}]
        _tbl_shape = _sl.shapes.add_table({rows_n}, {cols_n},
            Cm({left_cm}), Cm({top_cm}), Cm(20.0), Cm(4.0))
        for _row in _tbl_shape.table.rows:
            for _cell in _row.cells:
                _cell.text = ''
    """


# --- D-IMP-63 audio: pre_config step + oracle helpers ----------------------

def _ffmpeg_sine_mp3_step(dst_path: str, *, freq: int = 440, duration: int = 2) -> dict:
    """Generate a deterministic sine-wave .mp3 at `dst_path` via ffmpeg.

    Output bytes are reproducible across runs given the same (freq, duration).
    Used by D-IMP-63 audio_insert: the same .mp3 is referenced by both the
    eval (audio_in_slide vs vm_file) and the oracle (embed into pptx slide 1).
    """
    cmd = (
        f"rm -f '{dst_path}' && "
        f"ffmpeg -y -hide_banner -loglevel error "
        f"-f lavfi -i 'sine=frequency={freq}:duration={duration}' "
        f"-c:a libmp3lame -b:a 64k '{dst_path}'"
    )
    return _execute(cmd)


def _embed_audio_py(pptx_path: str, mp3_path: str, audio_name: str) -> str:
    """Build python that extracts pptx, drops the mp3 into ppt/media, adds
    the slide1 relationship + Content_Types entry, and re-zips. Mirrors the
    perturb-side helper for the same task."""
    return (
        "import zipfile, shutil, os\n"
        "import xml.etree.ElementTree as ET\n"
        f"pptx_path = {pptx_path!r}\n"
        f"mp3_path = {mp3_path!r}\n"
        f"audio_name = {audio_name!r}\n"
        "tmp_dir = '/tmp/_pptx_audio_embed_synth'\n"
        "if os.path.exists(tmp_dir):\n"
        "    shutil.rmtree(tmp_dir)\n"
        "os.makedirs(tmp_dir)\n"
        "with zipfile.ZipFile(pptx_path, 'r') as z:\n"
        "    z.extractall(tmp_dir)\n"
        "media_dir = os.path.join(tmp_dir, 'ppt/media')\n"
        "os.makedirs(media_dir, exist_ok=True)\n"
        "shutil.copy2(mp3_path, os.path.join(media_dir, audio_name))\n"
        "rels_path = os.path.join(tmp_dir, 'ppt/slides/_rels/slide1.xml.rels')\n"
        "if os.path.exists(rels_path):\n"
        "    tree = ET.parse(rels_path); root = tree.getroot()\n"
        "    max_id = 0\n"
        "    for r in root:\n"
        "        rid = r.get('Id', '')\n"
        "        if rid.startswith('rId'):\n"
        "            try: max_id = max(max_id, int(rid[3:]))\n"
        "            except ValueError: pass\n"
        "    new_rid = f'rId{max_id + 1}'\n"
        "    ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/relationships')\n"
        "    rel = ET.SubElement(root, 'Relationship')\n"
        "    rel.set('Id', new_rid)\n"
        "    rel.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio')\n"
        "    rel.set('Target', f'../media/{audio_name}')\n"
        "    tree.write(rels_path, xml_declaration=True, encoding='UTF-8')\n"
        "ct_path = os.path.join(tmp_dir, '[Content_Types].xml')\n"
        "ct_tree = ET.parse(ct_path); ct_root = ct_tree.getroot()\n"
        "ET.register_namespace('', 'http://schemas.openxmlformats.org/package/2006/content-types')\n"
        "if not any(e.get('Extension') == 'mp3' for e in ct_root):\n"
        "    e = ET.SubElement(ct_root, 'Default')\n"
        "    e.set('Extension', 'mp3'); e.set('ContentType', 'audio/mpeg')\n"
        "ct_tree.write(ct_path, xml_declaration=True, encoding='UTF-8')\n"
        "os.remove(pptx_path)\n"
        "with zipfile.ZipFile(pptx_path, 'w', zipfile.ZIP_DEFLATED) as zf:\n"
        "    for root_dir, _, files in os.walk(tmp_dir):\n"
        "        for f in files:\n"
        "            fp = os.path.join(root_dir, f)\n"
        "            zf.write(fp, os.path.relpath(fp, tmp_dir))\n"
        "shutil.rmtree(tmp_dir)\n"
    )


# --- D-IMP-64..67 long-tail eval helpers (xcu / pptx XML mutations) --------

_LO_REGISTRY_PATH = "/home/user/.config/libreoffice/4/user/registrymodifications.xcu"


def _xcu_template(items_xml: str) -> str:
    """Minimal valid registrymodifications.xcu wrapper. Mirrors perturb's
    `_registry_xcu_template` so the eval-side readers see the same shape."""
    return (
        '<?xml version="1.0" encoding="UTF-8"?>\n'
        '<oor:items xmlns:oor="http://openoffice.org/2001/registry" '
        'xmlns:xs="http://www.w3.org/2001/XMLSchema" '
        'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">\n'
        f'{items_xml}\n'
        '</oor:items>\n'
    )


def _xcu_presenter_disable_item() -> str:
    return (
        '  <item oor:path="/org.openoffice.Office.Impress/Misc/Start">'
        '<prop oor:name="EnablePresenterScreen" oor:op="fuse">'
        '<value>false</value></prop></item>'
    )


def _xcu_autosave_item(minutes: int) -> str:
    return (
        '  <item oor:path="/org.openoffice.Office.Common/Save/Document">'
        '<prop oor:name="AutoSave" oor:op="fuse"><value>true</value></prop>'
        '<prop oor:name="AutoSaveTimeIntervall" oor:op="fuse">'
        f'<value>{minutes}</value></prop></item>'
    )


def _write_xcu_step(items_xml: str) -> dict:
    """Emit an execute step that writes registrymodifications.xcu with `items_xml`."""
    xcu = _xcu_template(items_xml)
    py = (
        "import os, pathlib\n"
        f"p = pathlib.Path({_LO_REGISTRY_PATH!r})\n"
        "p.parent.mkdir(parents=True, exist_ok=True)\n"
        f"p.write_text({xcu!r}, encoding='utf-8')\n"
    )
    return _py_step(py)


def _portrait_swap_py(pptx_path: str) -> str:
    """Inline python that swaps slide_width/slide_height on `pptx_path`
    if currently landscape — so the resulting deck is portrait. Used as the
    D-IMP-66 oracle."""
    return (
        "from pptx import Presentation\n"
        f"prs = Presentation({pptx_path!r})\n"
        "if prs.slide_width >= prs.slide_height:\n"
        "    prs.slide_width, prs.slide_height = prs.slide_height, prs.slide_width\n"
        f"prs.save({pptx_path!r})\n"
    )


# Module-level template list. §I.h appends to it.
TEMPLATES: list[SynthTemplate] = []


# ===========================================================================
# §I. File-task templates (Batch, dataclass form, topic-rotating)
#
# Symmetric across all synth/*.py — same dataclass
# shape as synth/libreoffice_calc.py §I, with the **addition of a topic pool
# §I.b** that the parameterized source builders sample from per seed.
#
# Impress-specific: deck STRUCTURE is the FILE axis (slide count, layout,
# footer/notes/master, orientation). TOPIC is a *separate* random ingredient
# resolved INSIDE the parameterized source builders — free per-seed visual
# augmentation, not a template-level axis. Same FILE built with different
# seeds yields different slide titles / body content (food deck / wildlife
# deck / etc.) without affecting the task / eval / instruction triple.
#
# Layout (mirrors calc.py §I + topic pool):
#   §I.a  Caps
#   §I.b  Topic family pool + _pick_topic (impress-specific addition)
#   §I.c  Parameterized source builders (text_deck / hero_photo / gallery /
#         footer / notes / portrait)
#   §I.d  Dataclasses (File / Param / FileTask)
#   §I.e  File instances
#   §I.f  Factory  (_to_synth_template + _emit_templates)
#   §I.g  FILE_TASKS — flat list
#   §I.h  Emission
# ===========================================================================


# §I.a — caps.
SYNTH_CAP_TASKS_PER_FILE: int = 2
SYNTH_CAP_PARAMS_PER_TASK: int = 2


# §I.b — Topic themes (impress-specific).
#
# A `TopicTheme` is what a single pptx is *about*. The whole deck — titles,
# body lines, photos — is a coherent mini-narrative on ONE topic, not 6
# disconnected category-bucket items. `photo_dirs` lists which `photos/<dir>/`
# directories the deck's photos are sampled from at codegen time. A topic
# can span multiple dirs (e.g. "family weekend cooking" pulls from
# `photos/food/` AND `photos/portrait/`).
#
# `slide_titles` and `slide_bodies` each have ≥6 coherent strings forming a
# narrative arc (intro → details → reflection). Cycling helper extends them
# when n_slides > pool size.

@dataclass(frozen=True)
class TopicTheme:
    name: str
    slide_titles: list[str]
    slide_bodies: list[str]
    photo_dirs: list[str]   # subset of `lite/.../data/assets/synth/photos/<dir>/`


_TOPIC_FAMILIES: list[TopicTheme] = [
    TopicTheme(
        name="family_weekend_cooking",
        slide_titles=[
            "Weekend Family Cooking Plan",
            "Today's Main Dish: Hand-Rolled Pasta",
            "Helper Tasks for the Kids",
            "Table Setting in Warm Tones",
            "After-Dinner Board Game Time",
            "Next Weekend's Menu Ideas",
        ],
        slide_bodies=[
            "A full-family cooking night, from ingredients to table.",
            "Tomato base with fresh basil — simple but flavourful.",
            "Washing, kneading, plating — good tasks for school-age kids.",
            "Warm yellow cloth and clay dishes for a cosy mood.",
            "An hour of board games reinforces the weekly ritual.",
            "Try Spanish seafood paella next; pre-order the seafood.",
        ],
        photo_dirs=["food", "portrait"],
    ),
    TopicTheme(
        name="city_food_tour",
        slide_titles=[
            "City Food Tour — Downtown Loop",
            "Brunch at the Old Market",
            "Mid-Day Bakery Crawl",
            "Architecture Highlights Along the Way",
            "Sunset Drinks on a Rooftop",
            "Reservation Notes for Next Visit",
        ],
        slide_bodies=[
            "A six-stop walking tour, weather permitting.",
            "Egg dishes and pour-overs anchor the morning.",
            "Three small bakeries within four city blocks.",
            "Two heritage facades worth a quick photo stop.",
            "Roof terrace with skyline view; arrive by 17:30.",
            "Three places filled within the hour — book ahead.",
        ],
        photo_dirs=["city", "food", "architecture"],
    ),
    TopicTheme(
        name="wildlife_photography_journal",
        slide_titles=[
            "Wildlife Photography Journal — Field Notes",
            "Morning Light at the Wetland Edge",
            "Patience: Waiting for the Right Frame",
            "Lens Setup for Distant Subjects",
            "A Surprise Encounter at Dusk",
            "What I'd Pack Differently Next Time",
        ],
        slide_bodies=[
            "Three days in the field with two cameras and a tripod.",
            "Long shadows and pastel tones in the first thirty minutes.",
            "Most of the day is waiting; a few minutes is the shoot.",
            "Long-tele lens with monopod for steady mid-distance frames.",
            "Light rain brought wildlife out of cover near the trail.",
            "Lighter tripod, second battery, rain cover for the body.",
        ],
        photo_dirs=["wildlife", "nature", "landscape"],
    ),
    TopicTheme(
        name="office_productivity_review",
        slide_titles=[
            "Office Productivity Review — Q1",
            "Tools That Worked",
            "Workflow Bottlenecks We Hit",
            "Hardware Refresh Plan",
            "Meeting Cadence Adjustments",
            "Goals for the Next Quarter",
        ],
        slide_bodies=[
            "A retrospective on team workflows over the last three months.",
            "Shared note-app and async standups cut sync meetings by half.",
            "PR review queue and shared device backlog slowed shipping.",
            "Standing desks for three desks; second monitor for two seats.",
            "Trim weekly all-hands to thirty minutes; raise 1:1 frequency.",
            "Ship two features; reduce review-to-merge time by 40 percent.",
        ],
        photo_dirs=["office", "product"],
    ),
    TopicTheme(
        name="classroom_year_recap",
        slide_titles=[
            "Classroom Year in Review",
            "Learning Highlights from the Spring Term",
            "Group Project Showcase",
            "Reading-Corner Updates",
            "Parent-Teacher Night Recap",
            "Plans for the Summer Camp Weeks",
        ],
        slide_bodies=[
            "Twenty-two students across three subjects; mixed-ability cohort.",
            "Strong gains in science vocabulary; some gaps in long division.",
            "Three group posters on local-history themes earned wide praise.",
            "New shelf for graphic novels; rotate titles every two weeks.",
            "High attendance; questions focused on standardised testing.",
            "Two-week STEM camp; partner with the city library on visits.",
        ],
        photo_dirs=["classroom", "education", "portrait"],
    ),
    TopicTheme(
        name="healthy_eating_program",
        slide_titles=[
            "Healthy Eating Program — Week One",
            "Sample Menu for the First Three Days",
            "Portion Control Basics",
            "Common Pitfalls and How to Avoid Them",
            "Tracking Progress Without Obsessing",
            "Week Two Adjustments",
        ],
        slide_bodies=[
            "A four-week program for adults building sustainable habits.",
            "Whole grains, leafy greens, and one protein per main meal.",
            "Visual cues: a quarter plate of grains, a half plate of veg.",
            "Weekend slip-ups are normal; resume Monday without penance.",
            "Weekly weigh-in only; daily focus is on the plate, not the scale.",
            "Add a fourth weekly walk; reduce ultra-processed snacks further.",
        ],
        photo_dirs=["food", "medical"],
    ),
    TopicTheme(
        name="architecture_photo_essay",
        slide_titles=[
            "Architecture Photo Essay — Geometry in the City",
            "Repeating Patterns in Mid-Century Facades",
            "Light and Shadow at the Plaza",
            "Materials: Concrete, Glass, Tile",
            "Two Buildings Worth Revisiting",
            "Equipment Notes for Future Outings",
        ],
        slide_bodies=[
            "An afternoon spent looking up; one camera, one lens.",
            "Three blocks downtown where window grids dominate the frame.",
            "Low sun at 16:00 carves the plaza into bright and dark wedges.",
            "Texture contrasts make the strongest mid-distance compositions.",
            "The library colonnade and the rail-bridge underside.",
            "Wider lens next time; a polariser for glass-heavy scenes.",
        ],
        photo_dirs=["architecture", "city", "landscape"],
    ),
    TopicTheme(
        name="product_launch_recap",
        slide_titles=[
            "Product Launch Recap — Q1 Release",
            "What Shipped on Day One",
            "Day-Two Issues and Fixes",
            "Press and Community Coverage",
            "Early Customer Feedback Themes",
            "Roadmap Adjustments for Q2",
        ],
        slide_bodies=[
            "Combined hardware and software launch; team of twelve.",
            "Three core features and the integration with the partner SDK.",
            "Two hot-fixes deployed for a configuration edge case.",
            "Coverage in two trade outlets and one consumer publication.",
            "Onboarding flow praised; pricing tiers need clearer naming.",
            "Move pricing-page rework forward; defer the analytics module.",
        ],
        photo_dirs=["product", "office", "event"],
    ),
    TopicTheme(
        name="nature_hike_diary",
        slide_titles=[
            "Nature Hike Diary — Three Days on the Ridge Trail",
            "Day One: Forest Floor to Treeline",
            "Day Two: Alpine Meadow Crossing",
            "Wildlife Spotted Along the Way",
            "Camp Notes and Weather Patterns",
            "Gear That Worked and What I'd Change",
        ],
        slide_bodies=[
            "Forty kilometres total; two overnights, two trail companions.",
            "Steady climb through ferns and mossy boulders to first views.",
            "Wildflowers in bloom; light winds; soft afternoon clouds.",
            "Marmots, a distant eagle, and one mule deer at sunset.",
            "Cooler nights than expected; rain held off until the descent.",
            "Lighter pack, swap shells for the lighter rain jacket, more salt.",
        ],
        photo_dirs=["landscape", "nature", "wildlife"],
    ),
    TopicTheme(
        name="community_event_recap",
        slide_titles=[
            "Community Event Recap — Spring Block Party",
            "Setup and Volunteer Coordination",
            "Activities That Drew the Most Visitors",
            "Food and Drink Vendor Highlights",
            "Photos from the Afternoon Crowd",
            "Lessons Learned and Next Steps",
        ],
        slide_bodies=[
            "Three hundred attendees across the four-hour event.",
            "Seventeen volunteers; two-hour shifts kept everyone fresh.",
            "Kids' craft tables and the local-band stage drew steady crowds.",
            "Four food trucks and two coffee vendors; long lines past noon.",
            "Group portrait by the mural at the closing of the event.",
            "Move the food row away from the entrance; book live music earlier.",
        ],
        photo_dirs=["event", "portrait", "city"],
    ),
    TopicTheme(
        name="space_exploration_overview",
        slide_titles=[
            "Apollo Era — Engineering the Lunar Program",
            "Inner Solar System — Mars Robotic Exploration",
            "Outer Worlds — Jupiter's Moons in Focus",
            "Cassini at Saturn — Titan's Hazy Veil",
            "Looking Back — Earth from Lunar Orbit",
            "Next Decade — Crewed Return to the Moon",
        ],
        slide_bodies=[
            "Apollo 11 demonstrated that crewed lunar landings could be engineered to launch windows and trajectory budgets that left almost no margin. The mission's enduring lesson was the centrality of mission-control discipline.",
            "Curiosity and Perseverance have mapped Gale and Jezero in resolutions that ground geologists would envy. Sample caching at Jezero positions the next decade of Mars science around an eventual sample-return campaign.",
            "Galileo and Juno reframed Europa and Io as worlds in their own right — one likely hosting a subsurface ocean, the other the most volcanically active body in the solar system.",
            "Cassini's thirteen-year tour of the Saturnian system included a dedicated Titan radar campaign that revealed methane lakes and an Earth-like surface beneath the haze.",
            "The Apollo 8 Earthrise frame remains the most-reproduced photograph of the twentieth century. Its scientific value is modest; its cultural value is immense.",
            "Artemis aims to land crew near the lunar south pole within the decade. Open questions include cadence, sustainable surface infrastructure, and the role of commercial landers.",
        ],
        photo_dirs=["space", "portrait"],
    ),
    TopicTheme(
        name="renewable_energy_transition",
        slide_titles=[
            "Wind on the Plain — Capacity Factor at Scale",
            "Solar at Utility Scale — Field Deployments",
            "Hydroelectric Storage — Hoover Era to Present",
            "Nuclear Baseload — Cooling Tower Operations",
            "Grid Integration — Variable-Source Coordination",
            "Outlook — Adoption Trajectories Through 2040",
        ],
        slide_bodies=[
            "Onshore wind in plains regions now routinely clears a thirty-five percent capacity factor, with new turbines pushing past forty. Siting trade-offs remain the dominant constraint, not the underlying economics.",
            "Utility-scale photovoltaic farms have moved from boutique pilots to multi-gigawatt projects in under a decade. Levelised cost has fallen faster than almost any other generation source on record.",
            "Conventional hydro provides both energy and dispatchable storage; pumped-hydro retrofits at existing dams offer some of the cheapest grid-scale storage available today.",
            "Light-water reactors continue to anchor low-carbon baseload in several major grids. Cooling-tower thermal performance and tritium management remain the routine operational concerns.",
            "Integrating variable wind and solar at scale requires forecasting, fast-ramping reserves, and interregional transmission. None of these are technically novel; all of them are politically slow.",
            "Most credible 2040 scenarios converge on a grid roughly seventy percent renewable, with nuclear and gas-with-capture filling the remaining firm-power gap.",
        ],
        photo_dirs=["energy", "landscape", "industrial"],
    ),
    TopicTheme(
        name="manufacturing_floor_review",
        slide_titles=[
            "Q1 Production Throughput — Floor Metrics",
            "Loom Conversion Line — Defect Rate Trends",
            "Automotive Assembly — Cycle-Time Reduction",
            "Port Logistics — Container Throughput",
            "Quality Sampling — Method and Findings",
            "Forward Look — Capacity Plan H2",
        ],
        slide_bodies=[
            "Q1 throughput cleared the revised plan by four percent, driven mostly by the second-shift ramp on Line B. Unscheduled downtime fell to under two percent of scheduled hours for the first time in six quarters.",
            "The textile loom retrofit cut warp-break incidents by roughly a third. Remaining defects cluster around dye-lot transitions; a tighter purge protocol is being trialled this quarter.",
            "Assembly cycle time on the SUV line dropped from eighty-four to seventy-six seconds after the torque-station rework. Operator ergonomics scores improved in parallel.",
            "Container throughput at the export terminal hit a record in March, supported by the new rail-side gantry. Dwell times for transhipment containers remain above target.",
            "AQL sampling at receiving caught two supplier drift events early. The standing recommendation is to keep the tightened sampling plan through Q2 rather than relax it.",
            "Capacity headroom for H2 is comfortable on the body side and tight on the paint side. A second clearcoat booth is on the capex shortlist for the August review.",
        ],
        photo_dirs=["industrial", "office", "product"],
    ),
]


def _pick_topic(seed: int, salt: str = "") -> TopicTheme:
    """Pick one TopicTheme deterministically from seed + salt."""
    idx = (seed ^ (_stable_hash(salt) & 0xFFFFFFFF)) % len(_TOPIC_FAMILIES)
    return _TOPIC_FAMILIES[idx]


def _photos_in(category: str) -> list[str]:
    """Return rel paths (under assets/synth/) for all jpg/png photos in
    `photos/<category>/`. Empty list if dir missing."""
    dir_path = _ASSET_ROOT / "photos" / category
    if not dir_path.is_dir():
        return []
    out: list[str] = []
    for f in sorted(dir_path.iterdir()):
        if f.is_file() and f.suffix.lower() in (".jpg", ".jpeg", ".png"):
            out.append(f"photos/{category}/{f.name}")
    return out


def _sample_topic_photos(topic: TopicTheme, n: int, seed: int) -> list[str]:
    """Sample n photo rel-paths from topic.photo_dirs union (with replacement
    if pool < n). Deterministic in seed."""
    rng = random.Random(seed ^ (_stable_hash(f"photos_{topic.name}") & 0xFFFFFFFF))
    all_paths: list[str] = []
    for d in topic.photo_dirs:
        all_paths.extend(_photos_in(d))
    if not all_paths:
        raise FileNotFoundError(
            f"No photos under photos/{topic.photo_dirs} for topic {topic.name!r}"
        )
    if len(all_paths) >= n:
        return rng.sample(all_paths, n)
    return [rng.choice(all_paths) for _ in range(n)]


# §I.c — Parameterized source builders.
#
# Each takes (out_path, seed) + structural kwargs and returns a LIST of
# pre_config steps. Single-step builders return [build_step]; photo-using
# builders prepend N `_stage_asset` steps for real photos sampled from the
# topic's photo_dirs.

def _cycle_to(pool: list[str], n: int) -> list[str]:
    """Cycle `pool` (length k) until length n. If n ≤ k, just slice."""
    if n <= len(pool):
        return list(pool[:n])
    reps = (n // len(pool)) + 1
    return (pool * reps)[:n]


def _src_text_deck(
    out_path: str, seed: int, *,
    n_slides: int,
    layout: str = "title_body",        # ∈ {title_body, title_only, title_subtitle}
    source_font: str = "Liberation Sans",
) -> list[dict]:
    """Pure-text deck (no photos). N slides, each with title (+ body | subtitle).

    Titles and bodies come from the topic theme's coherent slide narrative.
    Cycles the pool if n_slides > len(pool).
    """
    topic = _pick_topic(seed, "text_deck")
    titles = _cycle_to(topic.slide_titles, n_slides)
    if layout == "title_only":
        bodies = ["" for _ in range(n_slides)]
    elif layout == "title_subtitle":
        bodies = [f"— {b.split('.')[0]}." for b in _cycle_to(topic.slide_bodies, n_slides)]
    else:
        bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = _slide_deck_body(titles, bodies, source_font=source_font)
    return [_build_pptx_step(out_path, deck_body)]


def _src_hero_photo_deck(
    out_path: str, seed: int, *,
    n_slides: int,
    photo_pos: str = "center",
    caption: bool = True,
) -> list[dict]:
    """Hero-photo deck: 1 real photo per slide + optional caption.

    Photos are sampled from topic.photo_dirs (e.g. food + portrait for a
    family-cooking topic) and staged via `_stage_asset`. The deck builder
    references the staged container paths in `add_picture(...)`.
    """
    topic = _pick_topic(seed, "hero_photo")
    titles = _cycle_to(topic.slide_titles, n_slides)
    captions = _cycle_to(topic.slide_bodies, n_slides) if caption else ["" for _ in range(n_slides)]
    pos_map = {
        # (left_cm, top_cm, width_cm)  — height auto from aspect ratio
        "center":  ("Cm(7.0)", "Cm(3.5)", "Cm(11.0)"),
        "top":     ("Cm(7.0)", "Cm(1.0)", "Cm(11.0)"),
        "bottom":  ("Cm(7.0)", "Cm(5.5)", "Cm(11.0)"),
        "corner":  ("Cm(1.0)", "Cm(1.0)", "Cm(8.0)"),
    }
    left, top, w = pos_map.get(photo_pos, pos_map["center"])

    # Sample + stage photos
    rel_paths = _sample_topic_photos(topic, n_slides, seed)
    basename = out_path.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    stage_steps: list[dict] = []
    container_paths: list[str] = []
    for i, rel in enumerate(rel_paths):
        ext = rel.rsplit(".", 1)[-1]
        dst = f"/tmp/{basename}_photo_{i}.{ext}"
        stage_steps.append(_stage_asset(rel, dst))
        container_paths.append(dst)

    py_lines = ["blank = prs.slide_layouts[6]"]
    for i, (title, cap, photo_path) in enumerate(zip(titles, captions, container_paths)):
        py_lines.append("slide = prs.slides.add_slide(blank)")
        # Title at top
        py_lines.append("tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.4), Cm(22.0), Cm(1.8))")
        py_lines.append("tf = tb.text_frame")
        py_lines.append("tf.text = " + repr(title))
        py_lines.append("tf.paragraphs[0].runs[0].font.size = Pt(24)")
        # Real photo via add_picture
        py_lines.append(
            f"slide.shapes.add_picture({photo_path!r}, {left}, {top}, width={w})"
        )
        if cap:
            py_lines.append("cb = slide.shapes.add_textbox(Cm(1.0), Cm(12.0), Cm(22.0), Cm(1.5))")
            py_lines.append("cf = cb.text_frame")
            py_lines.append("cf.text = " + repr(cap))
            py_lines.append("cf.paragraphs[0].runs[0].font.size = Pt(14)")
    build_step = _build_pptx_step(out_path, "\n".join(py_lines))
    return stage_steps + [build_step]


def _src_gallery_deck(
    out_path: str, seed: int, *,
    n_slides: int,
    grid: str = "2x2",
) -> list[dict]:
    """Photo gallery deck: each slide is a grid of REAL photos sampled
    from topic.photo_dirs."""
    topic = _pick_topic(seed, "gallery")
    titles = _cycle_to(topic.slide_titles, n_slides)
    grid_specs = {
        "2x2":   [(c, r) for r in range(2) for c in range(2)],
        "3x3":   [(c, r) for r in range(3) for c in range(3)],
        "1+3":   [(0, 0)] + [(1, r) for r in range(3)],
        "strip": [(c, 0) for c in range(4)],
    }
    cells = grid_specs.get(grid, grid_specs["2x2"])
    n_photos_total = n_slides * len(cells)

    # Sample + stage all photos (one per cell across all slides).
    rel_paths = _sample_topic_photos(topic, n_photos_total, seed)
    basename = out_path.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    stage_steps: list[dict] = []
    container_paths: list[str] = []
    for i, rel in enumerate(rel_paths):
        ext = rel.rsplit(".", 1)[-1]
        dst = f"/tmp/{basename}_photo_{i}.{ext}"
        stage_steps.append(_stage_asset(rel, dst))
        container_paths.append(dst)

    py_lines = ["blank = prs.slide_layouts[6]"]
    photo_idx = 0
    for i, title in enumerate(titles):
        py_lines.append("slide = prs.slides.add_slide(blank)")
        py_lines.append("tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.4), Cm(22.0), Cm(1.5))")
        py_lines.append("tb.text_frame.text = " + repr(title))
        py_lines.append("tb.text_frame.paragraphs[0].runs[0].font.size = Pt(22)")
        if grid == "3x3":
            cw, ch, x0, y0 = 6.0, 3.5, 3.5, 2.5
        elif grid == "1+3":
            cw, ch, x0, y0 = 5.0, 3.0, 4.5, 3.0
        elif grid == "strip":
            cw, ch, x0, y0 = 5.0, 5.0, 2.0, 4.0
        else:
            cw, ch, x0, y0 = 9.0, 5.0, 3.5, 3.0
        for j, (c, r) in enumerate(cells):
            if grid == "1+3" and j == 0:
                w_lit, px, py_ = "Cm(10.0)", "Cm(1.0)", "Cm(3.0)"
            else:
                w_lit = f"Cm({cw})"
                px = f"Cm({x0 + c * (cw + 0.3)})"
                py_ = f"Cm({y0 + r * (ch + 0.3)})"
            photo_path = container_paths[photo_idx]
            photo_idx += 1
            py_lines.append(
                f"slide.shapes.add_picture({photo_path!r}, {px}, {py_}, width={w_lit})"
            )
    build_step = _build_pptx_step(out_path, "\n".join(py_lines))
    return stage_steps + [build_step]


def _src_footer_deck(
    out_path: str, seed: int, *,
    n_slides: int,
    footer_type: str = "basic",
) -> list[dict]:
    """Text deck with footer / page number target (no photos)."""
    topic = _pick_topic(seed, "footer")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = _slide_deck_body(titles, bodies)
    footer_text = {
        "basic":            f"{topic.name.replace('_', ' ').title()} • Internal • Page",
        "with_page_num":    "Page",
        "with_logo_corner": f"{topic.name.replace('_', ' ').title()} 2026",
    }.get(footer_type, "Footer")
    footer_body = textwrap.dedent(f"""
        for _i, _sl in enumerate(prs.slides):
            _fb = _sl.shapes.add_textbox(Cm(1.0), Cm(13.5), Cm(22.0), Cm(0.6))
            _ft = _fb.text_frame
            _ft.text = {footer_text!r} + (f" {{_i + 1}}" if {footer_type!r} != "with_logo_corner" else "")
            if _ft.paragraphs[0].runs:
                _ft.paragraphs[0].runs[0].font.size = Pt(10)
                _ft.paragraphs[0].runs[0].font.italic = True
    """)
    return [_build_pptx_step(out_path, deck_body + "\n" + footer_body)]


def _src_notes_deck(
    out_path: str, seed: int, *,
    n_slides: int,
) -> list[dict]:
    """Text deck with speaker notes pre-populated (no photos)."""
    topic = _pick_topic(seed, "notes")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = _slide_deck_body(titles, bodies)
    notes_body = textwrap.dedent(f"""
        _initial_notes = {[f"Reminder: " + b for b in bodies]!r}
        for _i, _sl in enumerate(prs.slides):
            if _i < len(_initial_notes):
                _sl.notes_slide.notes_text_frame.text = _initial_notes[_i]
    """)
    return [_build_pptx_step(out_path, deck_body + "\n" + notes_body)]


def _src_portrait_deck(
    out_path: str, seed: int, *,
    n_slides: int,
) -> list[dict]:
    """Portrait-orientation text deck (no photos)."""
    topic = _pick_topic(seed, "portrait")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = (
        "prs.slide_width = Cm(14.29)\n"
        "prs.slide_height = Cm(25.4)\n"
        + _slide_deck_body(titles, bodies)
    )
    return [_build_pptx_step(out_path, deck_body)]


# --- D-IMP-51..56 source builders (eval-gap fillers, custom evaluators) ----
#
# These builders ship "raw" pre_config steps that the specialized factories
# below replay verbatim for the gold pptx (the gold factory swaps the path,
# then appends the variant-specific mutation). They do NOT go through the
# default `_to_synth_template`; their FileTasks set `make_template=` to the
# matching specialized factory.


# Inline PIL helper — builds a small RGB png on the container. Borrowed from
# the legacy `_PIL_PNG_BUILD` (kept here so the legacy block can be deleted).
_INLINE_PIL_PNG = textwrap.dedent("""\
    import os
    from PIL import Image, ImageDraw
    _img_path = '/tmp/_synth_impress_img.png'
    if not os.path.exists(_img_path):
        _img = Image.new('RGB', (320, 240), 'white')
        _draw = ImageDraw.Draw(_img)
        for _y in range(240):
            for _x in range(320):
                _img.putpixel((_x, _y), ((_x * 255 // 320), (_y * 255 // 240), 128))
        _draw.rectangle((40, 40, 280, 200), outline='black', width=4)
        _img.save(_img_path)
""")


def _build_pptx_with_pil_step(out_path: str, builder_body: str) -> dict:
    """python-pptx build step that first synthesises a PIL png at /tmp."""
    py = (
        _INLINE_PIL_PNG
        + _PPTX_PREAMBLE
        + textwrap.dedent(f"""
            prs = Presentation()
            prs.slide_width = Cm(25.4)
            prs.slide_height = Cm(14.29)
            path = {out_path!r}
        """)
        + "\n"
        + textwrap.dedent(builder_body)
        + "\nprs.save(path)\n"
    )
    return _py_step(py)


def _src_table_target_deck(
    out_path: str, seed: int, *,
    n_slides: int = 4,
) -> list[dict]:
    """Sparse text deck — slide K is intentionally light on content so a
    table fits below the title. Used by D-IMP-51 (insert_table) tasks."""
    topic = _pick_topic(seed, "table")
    titles = _cycle_to(topic.slide_titles, n_slides)
    # Shorter body lines so adding a 4-6 row table doesn't overflow.
    short = [b.split(".")[0] + "." for b in _cycle_to(topic.slide_bodies, n_slides)]
    deck_body = _slide_deck_body(titles, short)
    return [_build_pptx_step(out_path, deck_body)]


def _src_resize_hero_deck(
    out_path: str, seed: int, *,
    n_slides: int = 4,
    init_h_cm: float = 8.0,
) -> list[dict]:
    """Hero-photo deck where each slide carries ONE real photo at a known
    initial height — gold variants resize the picture on a target slide.

    Differs from `_src_hero_photo_deck` in that picture height is fixed (so
    the gold's `examine_modify_height` can detect a height delta cleanly)
    and there is no caption (so the slide has exactly 1 textbox + 1 picture).
    """
    topic = _pick_topic(seed, "resize_hero")
    titles = _cycle_to(topic.slide_titles, n_slides)
    rel_paths = _sample_topic_photos(topic, n_slides, seed)
    basename = out_path.rsplit("/", 1)[-1].rsplit(".", 1)[0]
    stage_steps: list[dict] = []
    container_paths: list[str] = []
    for i, rel in enumerate(rel_paths):
        ext = rel.rsplit(".", 1)[-1]
        dst = f"/tmp/{basename}_photo_{i}.{ext}"
        stage_steps.append(_stage_asset(rel, dst))
        container_paths.append(dst)

    py_lines = ["blank = prs.slide_layouts[6]"]
    for i, (title, photo_path) in enumerate(zip(titles, container_paths)):
        py_lines.append("slide = prs.slides.add_slide(blank)")
        py_lines.append("tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.4), Cm(22.0), Cm(1.8))")
        py_lines.append("tb.text_frame.text = " + repr(title))
        py_lines.append("tb.text_frame.paragraphs[0].runs[0].font.size = Pt(22)")
        py_lines.append(
            f"slide.shapes.add_picture({photo_path!r}, Cm(8.0), Cm(3.0), height=Cm({init_h_cm}))"
        )
    build_step = _build_pptx_step(out_path, "\n".join(py_lines))
    return stage_steps + [build_step]


def _src_master_bg_white_deck(
    out_path: str, seed: int, *,
    n_slides: int = 5,
) -> list[dict]:
    """Text deck with explicit WHITE slide backgrounds — gold mutate recolors
    every slide's background to the target RGB (for
    `evaluate_presentation_fill_to_rgb_distance`)."""
    topic = _pick_topic(seed, "master_bg")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = _slide_deck_body(titles, bodies)
    init_bg_py = textwrap.dedent("""\
        for _sl in prs.slides:
            _bg = _sl.background
            _bg.fill.solid()
            _bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    """)
    return [_build_pptx_step(out_path, deck_body + "\n" + init_bg_py)]


def _src_pagenum_master_deck(
    out_path: str, seed: int, *,
    n_slides: int = 5,
) -> list[dict]:
    """Text deck whose slide master carries a sldNum placeholder pre-coloured
    gray (#808080). The gold mutate (master XML patch) flips the placeholder
    colour to the target — used by D-IMP-55 (`check_page_number_colors`).

    Returns TWO steps: the python-pptx build, then a zip-rewrite that injects
    the sldNum placeholder into slideMaster1.xml. The specialized factory
    replays BOTH steps for the gold path so the master patch lands consistently.
    """
    topic = _pick_topic(seed, "pagenum_master")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = _slide_deck_body(titles, bodies)
    build_step = _build_pptx_step(out_path, deck_body)
    inject_py = textwrap.dedent(f"""\
        import zipfile, shutil, os
        src = {out_path!r}
        tmpd = '/tmp/_pn_init_' + os.path.basename(src).replace('.pptx', '')
        if os.path.exists(tmpd):
            shutil.rmtree(tmpd)
        os.makedirs(tmpd)
        with zipfile.ZipFile(src, 'r') as z:
            z.extractall(tmpd)
        mp = os.path.join(tmpd, 'ppt/slideMasters/slideMaster1.xml')
        with open(mp, 'r', encoding='utf-8') as f:
            xml = f.read()
        sldnum_sp = (
            '<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
            ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
            '<p:nvSpPr><p:cNvPr id="99" name="SlideNumberPlaceholder"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            '<p:nvPr><p:ph type="sldNum" sz="quarter" idx="2147483647"/></p:nvPr></p:nvSpPr>'
            '<p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/><a:p>'
            '<a:r><a:rPr lang="en-US"><a:solidFill><a:srgbClr val="808080"/></a:solidFill></a:rPr>'
            '<a:t>&lt;#&gt;</a:t></a:r></a:p></p:txBody></p:sp>'
        )
        xml = xml.replace('</p:spTree>', sldnum_sp + '</p:spTree>', 1)
        with open(mp, 'w', encoding='utf-8') as f:
            f.write(xml)
        os.remove(src)
        with zipfile.ZipFile(src, 'w', zipfile.ZIP_DEFLATED) as zf:
            for root, _, files in os.walk(tmpd):
                for fn in files:
                    fp = os.path.join(root, fn)
                    zf.write(fp, os.path.relpath(fp, tmpd))
        shutil.rmtree(tmpd)
    """)
    return [build_step, _py_step(inject_py)]


def _src_offcenter_pic_deck(
    out_path: str, seed: int, *,
    n_slides: int = 3,
) -> list[dict]:
    """Text deck with a small off-center PIL-built picture on slide 0. Used
    by D-IMP-56 (`check_image_stretch_and_center`) — gold leaves the deck
    structure unchanged and only swaps the picture's size + position."""
    topic = _pick_topic(seed, "offcenter_pic")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    deck_body = _slide_deck_body(titles, bodies)
    pic_extra = textwrap.dedent("""\
        _img_path = '/tmp/_synth_impress_img.png'
        _sl = prs.slides[0]
        _sl.shapes.add_picture(_img_path, Cm(1.0), Cm(1.0),
            width=Cm(5.0), height=Cm(4.0))
    """)
    return [_build_pptx_with_pil_step(out_path, deck_body + "\n" + pic_extra)]


# --- Master-XML colour-patch helper for D-IMP-55 ---------------------------

# --- Gold-mutate snippet helpers for D-IMP-51..56 --------------------------

def _gold_insert_table(idx: int, rows_n: int, cols_n: int) -> str:
    """Append `rows_n × cols_n` table to slide `idx`. Cells empty (eval cares
    only about shape count + table presence)."""
    return f"""\
        _sl = prs.slides[{idx}]
        _tbl_shape = _sl.shapes.add_table({rows_n}, {cols_n},
            Cm(2.5), Cm(8.0), Cm(20.0), Cm(5.0))
        for _row in _tbl_shape.table.rows:
            for _cell in _row.cells:
                _cell.text = ''
    """


def _gold_set_transition(idx: int, ttype: str) -> str:
    """Inject a `<p:transition><p:{ttype}/></p:transition>` element into
    slide `idx`. ttype must be one of fade/dissolve/wipe/push (LO-stable)."""
    NS_PPTX = "http://schemas.openxmlformats.org/presentationml/2006/main"
    xml = f'<p:transition xmlns:p="{NS_PPTX}"><p:{ttype}/></p:transition>'
    return f"""\
        _sld = prs.slides[{idx}]._element
        for _tr in list(_sld.findall(qn('p:transition'))):
            _sld.remove(_tr)
        _sld.append(etree.fromstring({xml!r}))
    """


def _gold_resize_picture(idx: int, target_h_cm: float) -> str:
    """Resize the FIRST picture on slide `idx` to `target_h_cm`. Uses the
    `_pic_*` index pattern: scans shapes for the first one with shape_type
    PICTURE (13). Source's picture was added with a known init height; the
    gold simply re-sets `.height` to a different cm value."""
    return f"""\
        _sl = prs.slides[{idx}]
        for _sh in _sl.shapes:
            if _sh.shape_type == 13:
                _sh.height = Cm({target_h_cm})
                break
    """


def _oracle_stretch_image(slide_idx: int = 0) -> str:
    """Inline-python ORACLE that stretches the picture on `slide_idx` of
    `src_path` to fill the slide and saves in place. Requires {src_path} to
    be substituted by the factory at call time — the factory passes
    `src_path` via the params dict but here we embed it via the same path
    template the source builder uses (D-IMP-56 sets src_basename d_imp_56.pptx)."""
    return textwrap.dedent(f"""\
        from pptx import Presentation
        prs = Presentation('{_DESKTOP}/d_imp_56.pptx')
        slide = prs.slides[{slide_idx}]
        sw = prs.slide_width; shh = prs.slide_height
        for shape in slide.shapes:
            if shape.shape_type == 13:
                shape.width = sw
                shape.height = shh
                shape.left = 0
                shape.top = 0
                break
        prs.save('{_DESKTOP}/d_imp_56.pptx')
    """)


def _master_color_patch_py(pptx_path: str, hex_color: str) -> str:
    """Generate inline python that rewrites every non-black srgbClr in
    slideMaster1.xml of `pptx_path` to `hex_color` (no leading '#'). Lifted
    from the legacy `_build_master_color_patch` so the legacy block can go."""
    return (
        "import zipfile, shutil, os, re\n"
        f"src = {pptx_path!r}\n"
        f"tmpd = '/tmp/_master_patch_{hex_color}'\n"
        "if os.path.exists(tmpd):\n"
        "    shutil.rmtree(tmpd)\n"
        "os.makedirs(tmpd)\n"
        "with zipfile.ZipFile(src, 'r') as z:\n"
        "    z.extractall(tmpd)\n"
        "mp = os.path.join(tmpd, 'ppt/slideMasters/slideMaster1.xml')\n"
        "if os.path.exists(mp):\n"
        "    with open(mp, 'r', encoding='utf-8') as f:\n"
        "        xml = f.read()\n"
        f"    xml = re.sub(r'<a:srgbClr val=\"[0-9A-Fa-f]{{6}}\"', '<a:srgbClr val=\"{hex_color}\"', xml)\n"
        "    with open(mp, 'w', encoding='utf-8') as f:\n"
        "        f.write(xml)\n"
        "os.remove(src)\n"
        "with zipfile.ZipFile(src, 'w', zipfile.ZIP_DEFLATED) as zf:\n"
        "    for root, _, files in os.walk(tmpd):\n"
        "        for fn in files:\n"
        "            fp = os.path.join(root, fn)\n"
        "            zf.write(fp, os.path.relpath(fp, tmpd))\n"
        "shutil.rmtree(tmpd)\n"
    )


# §I.d — Dataclasses (impress-shaped Param).

@dataclass(frozen=True)
class File:
    """One structurally distinct deck.

    `src(path, seed) -> list[dict]` returns the LIST of pre_config steps
    that build the source pptx. For text-only decks this is just
    `[_build_pptx_step(...)]`; for photo decks it is N `_stage_asset` steps
    followed by `_build_pptx_step(...)`. Topic theme is picked inside `src`
    per seed.
    """
    id: str
    setup_class: str
    basename: str
    src: Callable[[str, int], list[dict]]


@dataclass(frozen=True)
class Param:
    """One concrete parameterization of an impress task.

    `gold_mutate` is a python-pptx code snippet appended to the gold deck
    builder body (e.g. `_gold_set_title_font_color(0, (255, 0, 0))`).
    `examine_field` is the canonical `compare_pptx_files` `examine_*` key to
    enable (validated against `_VALID_EXAMINE_FIELDS` at build time). May be a
    single string or a tuple of keys for compound-mutation Params.
    `instr` is the rendered instruction. `extra_examine` opens additional
    `examine_*` flags (also whitelist-validated by `_to_synth_template`); the
    custom factories (transition / master_bg / pagenum_color / image_stretch)
    also use it as a generic rule carrier and bypass that validation path.
    """
    gold_mutate: str
    examine_field: str | tuple[str, ...]
    instr: str
    extra_examine: dict = field(default_factory=dict)


@dataclass(frozen=True)
class FileTask:
    """One (file, task) pair → one SynthTemplate at emit time.

    `make_template` overrides the default `_to_synth_template` factory for
    tasks that need a non-`compare_pptx_files` evaluator (insert_table /
    transition / image_resize / master_slide_bg / page_number_color /
    image_stretch_and_center). When set, the factory receives the FileTask
    and emits a fully custom SynthTemplate.
    """
    file: File
    task_id: str
    eval_class: str
    params: list[Param] = field(default_factory=list)
    make_template: Callable[["FileTask"], SynthTemplate] | None = None
    # When True, the default `_to_synth_template` factory replaces
    # `post_open_config_steps=_LO_POSTLAUNCH_SETTLE` with the richer
    # `_IMPRESS_BODY_FOCUS_STEPS` (extra Escape/F6 keystrokes that dodge
    # tip-of-the-day dialogs and ensure slide-canvas focus). Per
    # Axis A H-trigger: applied to 6 base title/body-style templates
    # (imp_04 / imp_05 / imp_12 / imp_25 / imp_30 / imp_32).
    body_focus: bool = False


# §I.e — File instances. Each is defined ONCE. Adding a structurally new
# deck = one new File entry. Topic variety is FREE per seed inside src.
# Use lambdas to bind structural kwargs (n_slides, layout, etc.).

D_IMP_01 = File("D-IMP-01", "text_deck_3s",  "d_imp_01.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_02 = File("D-IMP-02", "text_deck_5s",  "d_imp_02.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_03 = File("D-IMP-03", "text_deck_6s",  "d_imp_03.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=6, layout="title_body"))
D_IMP_04 = File("D-IMP-04", "title_only",    "d_imp_04.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_only"))
D_IMP_05 = File("D-IMP-05", "title_subtitle", "d_imp_05.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_subtitle"))
D_IMP_06 = File("D-IMP-06", "hero_photo_4s_center", "d_imp_06.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=4, photo_pos="center", caption=True))
D_IMP_07 = File("D-IMP-07", "hero_photo_6s_banner", "d_imp_07.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=6, photo_pos="top",    caption=False))
D_IMP_08 = File("D-IMP-08", "hero_photo_corner",    "d_imp_08.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=4, photo_pos="corner", caption=True))
D_IMP_09 = File("D-IMP-09", "gallery_2x2",   "d_imp_09.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=3, grid="2x2"))
D_IMP_10 = File("D-IMP-10", "gallery_3x3",   "d_imp_10.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=3, grid="3x3"))
D_IMP_11 = File("D-IMP-11", "gallery_1plus3", "d_imp_11.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=3, grid="1+3"))
D_IMP_12 = File("D-IMP-12", "gallery_strip", "d_imp_12.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=3, grid="strip"))
D_IMP_13 = File("D-IMP-13", "footer_basic_5s", "d_imp_13.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=5, footer_type="basic"))
D_IMP_14 = File("D-IMP-14", "footer_pagenum_5s", "d_imp_14.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=5, footer_type="with_page_num"))
D_IMP_15 = File("D-IMP-15", "footer_logo_5s", "d_imp_15.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=5, footer_type="with_logo_corner"))
D_IMP_16 = File("D-IMP-16", "notes_5s",      "d_imp_16.pptx",
                lambda o, s: _src_notes_deck(o, s, n_slides=5))
D_IMP_17 = File("D-IMP-17", "notes_6s",      "d_imp_17.pptx",
                lambda o, s: _src_notes_deck(o, s, n_slides=6))
D_IMP_18 = File("D-IMP-18", "portrait_5s",   "d_imp_18.pptx",
                lambda o, s: _src_portrait_deck(o, s, n_slides=5))
D_IMP_19 = File("D-IMP-19", "text_pos_4s",   "d_imp_19.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_body"))
D_IMP_20 = File("D-IMP-20", "text_swap_5s",  "d_imp_20.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))

# Additional structurally distinct decks (n_slides + layout +
# photo arrangement permutations). Total 40 files. Cap-2×2 → 80 emit rows.

D_IMP_21 = File("D-IMP-21", "text_deck_7s_long", "d_imp_21.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=7, layout="title_body"))
D_IMP_22 = File("D-IMP-22", "title_only_4s",     "d_imp_22.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_only"))
D_IMP_23 = File("D-IMP-23", "title_only_8s",     "d_imp_23.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=8, layout="title_only"))
D_IMP_24 = File("D-IMP-24", "title_subtitle_3s", "d_imp_24.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_subtitle"))
D_IMP_25 = File("D-IMP-25", "title_subtitle_6s", "d_imp_25.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=6, layout="title_subtitle"))
D_IMP_26 = File("D-IMP-26", "hero_photo_3s_nocap", "d_imp_26.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=3, photo_pos="center", caption=False))
D_IMP_27 = File("D-IMP-27", "hero_photo_5s_top",   "d_imp_27.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=5, photo_pos="top", caption=True))
D_IMP_28 = File("D-IMP-28", "hero_photo_5s_bottom","d_imp_28.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=5, photo_pos="bottom", caption=True))
D_IMP_29 = File("D-IMP-29", "hero_photo_8s_center","d_imp_29.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=8, photo_pos="center", caption=True))
D_IMP_30 = File("D-IMP-30", "gallery_2x2_4s",      "d_imp_30.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=4, grid="2x2"))
D_IMP_31 = File("D-IMP-31", "gallery_3x3_4s",      "d_imp_31.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=4, grid="3x3"))
D_IMP_32 = File("D-IMP-32", "gallery_1plus3_5s",   "d_imp_32.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=5, grid="1+3"))
D_IMP_33 = File("D-IMP-33", "gallery_strip_5s",    "d_imp_33.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=5, grid="strip"))
D_IMP_34 = File("D-IMP-34", "footer_basic_6s",     "d_imp_34.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=6, footer_type="basic"))
D_IMP_35 = File("D-IMP-35", "footer_pagenum_7s",   "d_imp_35.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=7, footer_type="with_page_num"))
D_IMP_36 = File("D-IMP-36", "footer_logo_8s",      "d_imp_36.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=8, footer_type="with_logo_corner"))
D_IMP_37 = File("D-IMP-37", "notes_7s",            "d_imp_37.pptx",
                lambda o, s: _src_notes_deck(o, s, n_slides=7))
D_IMP_38 = File("D-IMP-38", "notes_8s",            "d_imp_38.pptx",
                lambda o, s: _src_notes_deck(o, s, n_slides=8))
D_IMP_39 = File("D-IMP-39", "portrait_3s",         "d_imp_39.pptx",
                lambda o, s: _src_portrait_deck(o, s, n_slides=3))
D_IMP_40 = File("D-IMP-40", "portrait_7s",         "d_imp_40.pptx",
                lambda o, s: _src_portrait_deck(o, s, n_slides=7))

# Additional font variants + extra hero/gallery combos.

D_IMP_41 = File("D-IMP-41", "text_serif_5s",       "d_imp_41.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body",
                                            source_font="Liberation Serif"))
D_IMP_42 = File("D-IMP-42", "text_dejavu_sans_5s", "d_imp_42.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body",
                                            source_font="DejaVu Sans"))
D_IMP_43 = File("D-IMP-43", "text_dejavu_serif_4s","d_imp_43.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_body",
                                            source_font="DejaVu Serif"))
D_IMP_44 = File("D-IMP-44", "hero_photo_3s_top",   "d_imp_44.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=3, photo_pos="top",    caption=True))
D_IMP_45 = File("D-IMP-45", "hero_photo_5s_corner","d_imp_45.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=5, photo_pos="corner", caption=False))
D_IMP_46 = File("D-IMP-46", "hero_photo_7s_center","d_imp_46.pptx",
                lambda o, s: _src_hero_photo_deck(o, s, n_slides=7, photo_pos="center", caption=True))
D_IMP_47 = File("D-IMP-47", "gallery_2x2_5s",      "d_imp_47.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=5, grid="2x2"))
D_IMP_48 = File("D-IMP-48", "gallery_strip_4s",    "d_imp_48.pptx",
                lambda o, s: _src_gallery_deck(o, s, n_slides=4, grid="strip"))
D_IMP_49 = File("D-IMP-49", "footer_basic_4s",     "d_imp_49.pptx",
                lambda o, s: _src_footer_deck(o, s, n_slides=4, footer_type="basic"))
D_IMP_50 = File("D-IMP-50", "notes_4s",            "d_imp_50.pptx",
                lambda o, s: _src_notes_deck(o, s, n_slides=4))

# Eval-skill gap fillers (insert_table / transition /
# image_resize / master bg / page_number_color / image_stretch_and_center).
# These File instances pair with specialized factories via FileTask.make_template.

D_IMP_51 = File("D-IMP-51", "table_target_4s",   "d_imp_51.pptx",
                lambda o, s: _src_table_target_deck(o, s, n_slides=4))
D_IMP_52 = File("D-IMP-52", "transition_5s",     "d_imp_52.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_53 = File("D-IMP-53", "resize_hero_4s",    "d_imp_53.pptx",
                lambda o, s: _src_resize_hero_deck(o, s, n_slides=4, init_h_cm=8.0))
D_IMP_54 = File("D-IMP-54", "master_bg_white_5s","d_imp_54.pptx",
                lambda o, s: _src_master_bg_white_deck(o, s, n_slides=5))
D_IMP_55 = File("D-IMP-55", "pagenum_master_5s", "d_imp_55.pptx",
                lambda o, s: _src_pagenum_master_deck(o, s, n_slides=5))
D_IMP_56 = File("D-IMP-56", "offcenter_pic_3s",  "d_imp_56.pptx",
                lambda o, s: _src_offcenter_pic_deck(o, s, n_slides=3))

# Additional eval-gap fillers from instruction-shape validation.
# D-IMP-57..58 image add/resize at coord+size (echoes eval `c8...`).
# D-IMP-59..60 element-position move (echoes eval `2b.../15...`).
# D-IMP-61..62 multi-slide subset compound (echoes eval `4e...`).
# D-IMP-63 audio insert (echoes eval `c59742c0` compare_audios).
# D-IMP-64..67 long-tail eval funcs (presenter / autosave / orientation /
# left_panel).

D_IMP_57 = File("D-IMP-57", "resize_hero_5s_p2", "d_imp_57.pptx",
                lambda o, s: _src_resize_hero_deck(o, s, n_slides=5, init_h_cm=6.0))
D_IMP_58 = File("D-IMP-58", "resize_hero_4s_pos", "d_imp_58.pptx",
                lambda o, s: _src_resize_hero_deck(o, s, n_slides=4, init_h_cm=7.0))
D_IMP_59 = File("D-IMP-59", "text_title_move_4s", "d_imp_59.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_body"))
D_IMP_60 = File("D-IMP-60", "hero_photo_move_3s", "d_imp_60.pptx",
                lambda o, s: _src_resize_hero_deck(o, s, n_slides=3, init_h_cm=6.0))
D_IMP_61 = File("D-IMP-61", "text_subset_color_5s", "d_imp_61.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_62 = File("D-IMP-62", "text_subset_format_5s", "d_imp_62.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_63 = File("D-IMP-63", "audio_target_3s",  "d_imp_63.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_64 = File("D-IMP-64", "presenter_target_3s", "d_imp_64.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_65 = File("D-IMP-65", "autosave_target_3s", "d_imp_65.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_66 = File("D-IMP-66", "orientation_landscape_4s", "d_imp_66.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_body"))
D_IMP_67 = File("D-IMP-67", "left_panel_target_3s", "d_imp_67.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))

# Save-As / Export coverage Files (libreoffice_impress Axis E,
# mirror upstream eval `a097acff` Save-As .pptx pattern). Generic text
# decks; the per-task `_make_save_as_filetask_template` factory wires
# the agent's typed Save-As filename via Param.extra_examine['save_as_name'].
D_IMP_68 = File("D-IMP-68", "save_as_target_3s",   "d_imp_68.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_69 = File("D-IMP-69", "save_as_target_4s",   "d_imp_69.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_subtitle"))

# Doc-wide framing coverage Files (libreoffice_impress
# `slide_anchor.doc_wide` bridge). Instructions reference the entire
# deck rather than ordinal slide indices. Eval = `compare_pptx_files`
# default factory; gold helper applies the mutation to every slide.
D_IMP_70 = File("D-IMP-70", "doc_wide_title_5s",   "d_imp_70.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_71 = File("D-IMP-71", "doc_wide_title_4s",   "d_imp_71.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_body"))
D_IMP_72 = File("D-IMP-72", "doc_wide_bg_5s",      "d_imp_72.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_73 = File("D-IMP-73", "doc_wide_title_3s",   "d_imp_73.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_subtitle"))
D_IMP_74 = File("D-IMP-74", "doc_wide_title_6s",   "d_imp_74.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=6, layout="title_body"))
D_IMP_75 = File("D-IMP-75", "doc_wide_bg_4s",      "d_imp_75.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_body"))
D_IMP_76 = File("D-IMP-76", "doc_wide_title_4s_to","d_imp_76.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=4, layout="title_only"))
D_IMP_77 = File("D-IMP-77", "doc_wide_title_5s_sub","d_imp_77.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_subtitle"))
D_IMP_78 = File("D-IMP-78", "doc_wide_bg_3s",      "d_imp_78.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))

# Cycle-l7 ADDS — compound `compare_pptx_files`×2 / PNG-export /
# build-deck-from-scratch coverage. Closes the atom_2 0%→26% synth gap and
# the compare_images / "scratch deck" eval-row holes.

# Three-textbox decks for the per-textbox-color eval row (`04578141`):
# the target slide has THREE textboxes; the gold colors them top→bottom.
D_IMP_79 = File("D-IMP-79", "three_textbox_3s",    "d_imp_79.pptx",
                lambda o, s: _src_three_textbox_deck(o, s, n_slides=3, target_slide_idx=0))
D_IMP_80 = File("D-IMP-80", "three_textbox_5s",    "d_imp_80.pptx",
                lambda o, s: _src_three_textbox_deck(o, s, n_slides=5, target_slide_idx=0))

# Multi-slide alignment / title-style / per-line strikethrough decks for
# `05dd4c1d` / `08aced46` / `4ed5abd0` / `550ce7e7`.
D_IMP_81 = File("D-IMP-81", "multi_align_5s",      "d_imp_81.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
D_IMP_82 = File("D-IMP-82", "title_text_align_3s", "d_imp_82.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_83 = File("D-IMP-83", "multi_title_underline_5s", "d_imp_83.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))
def _src_todo_list_deck(out_path: str, seed: int, *, n_slides: int = 3) -> list[dict]:
    """Multi-paragraph body deck for `_gold_strikethrough_lines` (D_IMP_84).

    Validation fix: `_src_text_deck` writes the body via `text_frame.text = body`
    which produces a single paragraph. `_gold_strikethrough_lines` walks
    `text_frame.paragraphs[_li]` for indices 0/1/2 and silently no-ops past
    paragraph 0 (`if _li < len(paragraphs)` guard), so gold == source and the
    task evaluates as vacuous AND the agent cannot locate "first and third
    lines" → INFEASIBLE_CLAIM_TRAIN. Build each slide's body as a true
    multi-paragraph to-do list (4 items) so both the eval contract and the
    user-visible referent ("lines") are well-formed.
    """
    _TODO_ITEMS: list[list[str]] = [
        [
            "Finalise the Q3 sprint retrospective notes",
            "Push the new analytics dashboard to staging",
            "Schedule the customer onboarding follow-ups",
            "Renew the team's compliance training enrolment",
        ],
        [
            "Update the integration test fixtures for the new schema",
            "Triage the open production alerts from last week",
            "Coordinate the cross-team release calendar",
            "Capture the architecture decision in the design doc",
        ],
        [
            "Send the budget reconciliation memo to finance",
            "Review the procurement requests pending sign-off",
            "Confirm the venue booking for the off-site",
            "Distribute the meeting minutes from the leadership sync",
        ],
    ]
    titles = ["Last Week's To-Do List", "Sprint Backlog Carry-Over", "Operations Action Items"][:n_slides]
    todos = _TODO_ITEMS[:n_slides]
    py_lines = ["blank = prs.slide_layouts[6]"]
    for title, items in zip(titles, todos):
        py_lines.append("slide = prs.slides.add_slide(blank)")
        py_lines.append("tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(22.0), Cm(2.0))")
        py_lines.append("tb.text_frame.text = " + repr(title))
        py_lines.append("tb.text_frame.paragraphs[0].runs[0].font.size = Pt(28)")
        py_lines.append("bb = slide.shapes.add_textbox(Cm(1.0), Cm(3.5), Cm(22.0), Cm(10.0))")
        py_lines.append("bf = bb.text_frame")
        py_lines.append("bf.word_wrap = True")
        py_lines.append("bf.text = " + repr(items[0]))
        for item in items[1:]:
            py_lines.append("_p = bf.add_paragraph()")
            py_lines.append("_p.text = " + repr(item))
        py_lines.append("for _p in bf.paragraphs:")
        py_lines.append("    if _p.runs:")
        py_lines.append("        _p.runs[0].font.size = Pt(18)")
    return [_build_pptx_step(out_path, "\n".join(py_lines))]


D_IMP_84 = File("D-IMP-84", "strikethrough_lines_3s", "d_imp_84.pptx",
                lambda o, s: _src_todo_list_deck(o, s, n_slides=3))

# PNG-export target deck (`455d3c66`).
D_IMP_85 = File("D-IMP-85", "png_export_target_3s", "d_imp_85.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))

# Build-deck-from-scratch target (`bf4e9888`). Source is a 1-slide stub; the
# agent's job is to add K image-only blank slides.
D_IMP_86 = File("D-IMP-86", "build_from_images_stub", "d_imp_86.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=1, layout="title_body"))

# Extra compound coverage decks (compound bold+bg / underline+alignment
# already exist at D_IMP_05 / D_IMP_08 in non-compound 1-arm form; we now
# emit dedicated 2-arm compound rows on fresh files).
D_IMP_87 = File("D-IMP-87", "compound_color_bold_5s", "d_imp_87.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=5, layout="title_body"))

# --- Validation additions -----------------------------------------------
# 4-arm compound (emulates osworld_libreoffice_impress_f23acfd2 "Add a
# bullet point to the content of this slide" — single-slide source so
# "this slide" reads as unambiguous; gold mutator appends one bullet to
# the body textbox).
D_IMP_88 = File("D-IMP-88", "bullet_add_1s",      "d_imp_88.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=1, layout="title_body"))

# Deck-wide-implicit framing decks (libreoffice_impress
# `target_scope.deck_wide_or_implicit` bridge): instructions use NAKED
# voice — no "every slide" / "all slides" / "slide N" / ordinal anchor —
# so the regex in measure_gap.impress_target_scope falls through to
# `deck_wide_or_implicit`. Sources are single-slide decks so "the title"
# / "the background" reads unambiguously as the only slide; gold helpers
# still iterate `prs.slides` to keep behaviour identical if the seed bumps
# the slide count later.
D_IMP_89 = File("D-IMP-89", "implicit_title_3s",  "d_imp_89.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_90 = File("D-IMP-90", "implicit_bg_3s",     "d_imp_90.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_91 = File("D-IMP-91", "implicit_font_3s",   "d_imp_91.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_92 = File("D-IMP-92", "implicit_align_3s",  "d_imp_92.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_93 = File("D-IMP-93", "implicit_bullet_3s", "d_imp_93.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))
D_IMP_94 = File("D-IMP-94", "implicit_color_3s",  "d_imp_94.pptx",
                lambda o, s: _src_text_deck(o, s, n_slides=3, layout="title_body"))


# §I.f — Factory.

def _to_synth_template(ft: FileTask) -> SynthTemplate:
    """Turn ONE impress FileTask into ONE SynthTemplate.

    Per-seed: pick params[seed % len(params)]; build source via ft.file.src
    (topic random per seed); build gold by replaying the source builder for
    the gold path + appending variant.gold_mutate; build evaluator with
    examine_<field>=True + any extra_examine overrides.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        # ft.file.src returns a list of pre_config steps. The LAST one is the
        # python-pptx build step; any leading steps are `_stage_asset` host_push
        # steps for real photos referenced by the build step. The gold deck
        # replays the same build body (so it references the same staged photos)
        # then appends the variant's mutate snippet.
        src_steps = ft.file.src(src_path, seed)
        build_step = src_steps[-1]
        src_cmd = build_step["parameters"]["command"]
        body = src_cmd.split("python3 << 'PYEOF'\n", 1)[1].rsplit("\nPYEOF", 1)[0]
        gold_py = (
            body.replace(f"path = {src_path!r}", f"path = {gold_path!r}")
                .rsplit("prs.save(path)", 1)[0]
            + textwrap.dedent(variant.gold_mutate)
            + "\nprs.save(path)\n"
        )
        gold_step = _py_step(gold_py)
        options = _examine_options(variant.examine_field)
        # Validation landmine fix: validate compound-Param
        # extra_examine keys so eval_class strings can't sneak in as no-op
        # options. The custom factories (transition / master_bg / pagenum_color)
        # use extra_examine as a generic rule carrier and bypass this path.
        for _k in variant.extra_examine:
            assert _k in _VALID_EXAMINE_FIELDS, (
                f"Param.extra_examine: {_k!r} is not a valid compare_pptx_files "
                f"key. Valid keys: {sorted(_VALID_EXAMINE_FIELDS)}"
            )
        options.update(variant.extra_examine)
        post_steps = list(_IMPRESS_BODY_FOCUS_STEPS) if ft.body_focus else list(_LO_POSTLAUNCH_SETTLE)
        instr = variant.instr
        # Name-only color tasks give only a color NAME, but the gold is
        # an arbitrary RGB the agent cannot infer (judged at ±2 even via the
        # tolerant compare — e.g. "pale-blue"→#C8E6FF, "purple"→#800080). If the
        # instruction doesn't already name a hex, append the exact gold hex
        # (extracted from the gold_mutate's RGBColor) so the task is winnable.
        _ef = str(variant.examine_field)
        if ("examine_color_rgb" in _ef or "examine_background_color" in _ef) and "#" not in instr:
            _m = re.search(r"RGBColor\((\d+),\s*(\d+),\s*(\d+)\)", variant.gold_mutate)
            if _m:
                _hex = "#{:02X}{:02X}{:02X}".format(int(_m.group(1)), int(_m.group(2)), int(_m.group(3)))
                instr = f"{instr.rstrip()} (Custom Color {_hex})"
        return {
            "instruction":      instr,
            "pre_config_steps": [*src_steps, gold_step],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "_eval_options":    options,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": post_steps,
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        # Validation note: this was deferred from the title-color cluster
        # title_color cluster ≥15 variants): when the gold mutates RGB color,
        # route through the RGB-tolerant local override (eval/metrics.py) so
        # the Custom-Color picker's 1-3 byte per-channel drift no longer
        # false-fails the eval.
        eval_options = p["_eval_options"]
        if eval_options.get("examine_color_rgb"):
            func = "compare_pptx_files_color_tolerant"
        else:
            func = "compare_pptx_files"
        return {
            "func": func,
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": eval_options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="libreoffice_impress",
        instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=len(pool),
        setup_class=ft.file.setup_class,
        eval_class=ft.eval_class,
    )


def _replay_src_for_gold(
    src_steps: list[dict], src_path: str, gold_path: str, gold_mutate: str,
) -> list[dict]:
    """Replay the source builder steps but redirect the python-pptx build
    to `gold_path` and append `gold_mutate` to the python body. Multi-step
    builders (e.g. pagenum_master) keep all leading steps verbatim except
    that any reference to `src_path` flips to `gold_path`.

    `host_push` stage_asset steps (no `command` field, just `files`) are
    passed through verbatim — they stage real photos at fixed `/tmp/...`
    paths that BOTH the source and the gold builder reference, so we don't
    need a second copy of the same photos for the gold deck.
    """
    out: list[dict] = []
    for step in src_steps:
        if step["type"] == "host_push":
            # Photo-staging step — same photos serve source AND gold; pass through.
            continue
        cmd = step["parameters"]["command"]
        new_cmd = cmd.replace(repr(src_path), repr(gold_path))
        # Append `gold_mutate` to the LAST python heredoc only — that's the
        # one that owns the `prs.save(path)` line we're piggybacking on.
        if step is src_steps[-1] and "prs.save(path)" in new_cmd:
            new_cmd = new_cmd.replace(
                "prs.save(path)\n",
                textwrap.dedent(gold_mutate).rstrip() + "\nprs.save(path)\n",
                1,
            )
        out.append({"type": step["type"], "parameters": {**step["parameters"], "command": new_cmd}})
    return out


def _make_table_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-51 (insert_table). Eval =
    `compare_pptx_files` with `examine_shape=False`. The instruction says
    "On slide N, insert a table with R rows and C columns" — it does NOT
    commit to the table's pixel position, so the shape-presence check
    relies on the slide's shape-count delta (gold has table, source
    doesn't → len(shapes) differs → eval fails when agent didn't insert).
    `examine_shape=True` would strictly check (l,t,w,h) per-shape index
    against the gold's hardcoded `_gold_insert_table` position, false-
    failing agent placements at any other reasonable position.

    Validation fix: relaxed examine_shape=True → False (trigger O —
    output-format brittleness)."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    options = {
        "examine_shape": False,
        "examine_run_count": False,
        "examine_indent": False,
        "examine_strike_through": False,
        "examine_bullets": False,
        "examine_text": False,
    }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_pptx_files",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_add_slide_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for add_slide tasks (eval rows 28/36/38 family).
    Gold appends N≥1 slides; eval reads `examine_number_of_slides=True` with
    `examine_text=False` and `examine_shape=False` so the only diff signal is
    the slide count delta. compare_pptx_files always checks per-pair shape
    count at L302-305 (not gated by examine_shape) — to keep the per-pair
    comparison cheap, the gold helper appends slides at the END, so source's
    first-N slides match gold's first-N slides shape-for-shape and only the
    new tail slide(s) differ (and only fail the slide-count check, which is
    what the agent's task closes)."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    options = {
        "examine_number_of_slides": True,
        "examine_shape": False,
        "examine_run_count": False,
        "examine_indent": False,
        "examine_strike_through": False,
        "examine_bullets": False,
        "examine_text": False,
        "examine_font_name": False,
        "examine_font_size": False,
        "examine_font_bold": False,
        "examine_font_italic": False,
        "examine_color_rgb": False,
        "examine_font_underline": False,
        "examine_alignment": False,
        "examine_background_color": False,
        "examine_note": False,
    }

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_pptx_files",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_transition_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-52 (transition). Eval = `check_transition`
    with rule-based expected (slide_idx + transition_type). Param's
    `gold_mutate` field carries the transition-XML inject snippet; param's
    `examine_field` field is overloaded to carry the transition_type string
    for the eval rule (uses the field as a generic key/value carrier)."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        # extra_examine carries transition rule {"slide_idx": K, "transition_type": "fade"}
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "_rule":            dict(variant.extra_examine),
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "check_transition",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "rule", "rules": p["_rule"]},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_image_resize_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-53 (image_resize). Eval =
    `compare_pptx_files` with `examine_modify_height=True`. The source deck
    embeds a real photo at `init_h_cm`; the variant's gold_mutate rebuilds
    the deck and then resizes the picture on slide K to a different height."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    options = _examine_options("examine_modify_height")
    options["examine_shape"] = False  # required: modify_height branch needs shape OFF
    options["examine_text"] = False   # picture has no text run
    # Validation bug fix (audit §5g C5 + §2): widen upstream
    # `is_approximately_equal` from default 0.5% to 10% relative so LO
    # round-trip EMU drift on picture height (sub-mm at typical 4-12cm) no
    # longer false-fails the eval. Instructions are qualitative ("about Xcm
    # tall") so 10% relative is well within semantic tolerance.
    options["approximately_tolerance"] = 0.1

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_pptx_files",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_master_bg_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-54 (master slide bg). Eval =
    `evaluate_presentation_fill_to_rgb_distance` (rule-based: rgb +
    original_rgb). The variant's gold_mutate recolors all slide backgrounds
    to the target RGB; extra_examine carries the rule {rgb, original_rgb}."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "_rule":            dict(variant.extra_examine),
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "evaluate_presentation_fill_to_rgb_distance",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "rule", "rules": p["_rule"]},
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_pagenum_color_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-55 (page_number_color). Eval =
    `check_page_number_colors` (rule-based: color name). Source pre-injects
    a sldNum placeholder coloured gray; the gold mutate is a separate
    zip-rewrite step that flips the master's non-black srgbClr to the target
    colour. extra_examine carries {"color": "<name>"} for the eval rule.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        # Replay the WHOLE source pipeline against gold_path (build + sldNum
        # inject), then run the master colour-patch py against gold_path.
        # Skip host_push steps (no `command` field; staged photos serve both).
        gold_replay = [
            {"type": s["type"],
             "parameters": {**s["parameters"],
                            "command": s["parameters"]["command"].replace(repr(src_path), repr(gold_path))}}
            for s in src_steps if s["type"] != "host_push"
        ]
        # variant.gold_mutate is the hex colour (e.g. "FF0000").
        gold_patch = _py_step(_master_color_patch_py(gold_path, variant.gold_mutate))
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_replay, gold_patch],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "_rule":            dict(variant.extra_examine),
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "check_page_number_colors",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "rule", "rules": p["_rule"]},
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_image_stretch_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-56 (image_stretch_and_center). Eval =
    `check_image_stretch_and_center` (compares result vs an expected SNAPSHOT
    that has the SAME image blob but different size/position). Pre_config:
    (1) build source (small off-center pic); (2) snapshot src to expected path
    BEFORE oracle mutation. Oracle: mutate src so the picture fills the slide.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    expected_path = f"/tmp/expected_original_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        # Snapshot src to the expected_path BEFORE the agent mutates it.
        snapshot_step = _execute(f"cp '{src_path}' '{expected_path}'")
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, snapshot_step],
            "out_path":         src_path,
            "expected_path":    expected_path,
            "_oracle_py":       variant.gold_mutate,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "check_image_stretch_and_center",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"],
                         "dest": "expected_original.pptx"},
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: [_py_step(p["_oracle_py"])],
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_image_size_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-57/58 (P2 — image add at coord+size).

    Eval = `compare_pptx_files` with `examine_image_size=True` and
    `examine_shape=False`. Source pre-stages a real photo at the initial
    (w_init, h_init); gold replays + resizes the picture to the target
    (w, h). The agent's task is to size+position the image to those target
    dimensions. examine_image_size only checks picture shapes' width+height
    so other shape positions on the slide can diverge slightly without
    false-fail.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    options = _examine_options("examine_image_size")
    options["examine_shape"] = False
    options["examine_text"] = False
    # Validation bug fix (audit §5g C5 + §2): widen upstream
    # `is_approximately_equal` from default 0.5% to 10% relative so LO
    # round-trip EMU drift on picture (left, top, width, height) no longer
    # false-fails. Instructions are qualitative ("about XcmxYcm") so 10%
    # relative is well within semantic tolerance.
    options["approximately_tolerance"] = 0.1

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_pptx_files",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_position_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-59/60 (P3 — element-position move).

    Eval = `compare_pptx_files` with `examine_shape=True`. examine_shape
    compares (left, top, width, height) per shape index — so a gold that
    moves shape K from (l0,t0) to (l1,t1) detects exactly that move. The
    rest of the deck is identical between source and gold, so the only diff
    is the moved shape's position.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    options = {
        "examine_shape": True,
        "examine_run_count": False,
        "examine_indent": False,
        "examine_strike_through": False,
        "examine_bullets": False,
        "examine_text": False,
        # Validation bug fix (audit §5g C5 + §2): widen upstream
        # `is_approximately_equal` from default 0.5% to 10% relative on top
        # of the existing cm-grid snap. Position-compare instructions are
        # qualitative ("top-right corner", "bottom of the slide") so 10%
        # relative tolerance is well within semantic intent.
        "approximately_tolerance": 0.1,
    }

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            # validation fix (was deferred from validation
            # cluster escalation growth #1, position-compare cluster: d_imp_59 +
            # d_imp_60 x3 = 5 members): route through Cm-tolerant local
            # override (eval/metrics.py) so LO round-trip's sub-mm EMU drift
            # no longer false-fails the eval.
            "func": "compare_pptx_files_position_tolerant",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_audio_insert_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-63 (P5 — audio insert).

    Eval = `compare_audios` with `result.type=audio_in_slide` (reads the
    audio embedded in the pptx via the eval slides metric) and
    `expected.type=vm_file` pointing at the staged mp3 on Desktop.
    Pre_config: build source deck, ffmpeg-generate a sine-wave mp3 on
    Desktop. Oracle: zip-rewrite the pptx to embed the mp3 in slide 1.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        # variant.extra_examine carries {audio_name, freq, duration}
        audio_name = variant.extra_examine["audio_name"]
        mp3_path = f"{_DESKTOP}/{audio_name}"
        freq = variant.extra_examine["freq"]
        duration = variant.extra_examine["duration"]
        audio_step = _ffmpeg_sine_mp3_step(mp3_path, freq=freq, duration=duration)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, audio_step],
            "out_path":         src_path,
            "_mp3_path":        mp3_path,
            "_audio_name":      audio_name,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_audios",
            "result": {
                "type": "audio_in_slide",
                "ppt_file_path": p["out_path"],
                "slide_index": 0,
                "dest": p["_audio_name"],
            },
            "expected": {
                "type": "vm_file", "path": p["_mp3_path"],
                "dest": f"expected_{p['_audio_name']}",
            },
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    def _oracle(p: dict) -> list[dict]:
        # Kill any open soffice so the zip-rewrite isn't fighting an open lock,
        # then embed the mp3 into slide 1 via the shared helper.
        return [
            _execute("killall soffice.bin 2>/dev/null; sleep 1; true"),
            _py_step(_embed_audio_py(p["out_path"], p["_mp3_path"], p["_audio_name"])),
        ]

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=_oracle,
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_xcu_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-64/65 (P6 long-tail — xcu-based evals:
    check_presenter_console_disable / check_auto_saving_time).

    Pre_config: build source deck + write an "initial" xcu (without the
    target property, or with the opposite value) so eval reads 0 before the
    oracle runs. Oracle: rewrite the xcu to set the target property.
    `variant.gold_mutate` carries the FINAL xcu items_xml (string). `variant.
    extra_examine` carries the eval func name and (optional) `rules`.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        # Initial xcu: opposite-of-target so eval=0 before oracle.
        init_items = variant.extra_examine["init_items"]
        init_step = _write_xcu_step(init_items)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, init_step],
            "out_path":         src_path,
            "_target_items":    variant.gold_mutate,
            "_eval_func":       variant.extra_examine["func"],
            "_eval_rules":      variant.extra_examine.get("rules"),
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        # validation BUG-fix (D-IMP-64/65 family): force-quit soffice before eval
        # reads registrymodifications.xcu. LO buffers xcu writes; without a
        # flush, the eval reads the pre-mutation file and returns 0.
        evaluator: dict = {
            "func": p["_eval_func"],
            "result": {"type": "vm_file", "path": _LO_REGISTRY_PATH,
                       "dest": "registrymodifications.xcu"},
            "options": {},
            "postconfig": [
                {"type": "execute", "parameters": {
                    "command": ["bash", "-c",
                                "killall soffice.bin 2>/dev/null; sleep 2; true"]}},
            ],
        }
        if p["_eval_rules"] is not None:
            evaluator["expected"] = {"type": "rule", "rules": p["_eval_rules"]}
        else:
            evaluator["expected"] = {}
        return evaluator

    def _oracle(p: dict) -> list[dict]:
        return [_write_xcu_step(p["_target_items"])]

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=_oracle,
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_orientation_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-66 (P6 — check_slide_orientation_Portrait).

    Source deck is landscape. Oracle swaps slide_width/height to make it
    portrait. Eval reads the saved pptx and returns 1 iff width<height.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": list(src_steps),
            "out_path":         src_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "check_slide_orientation_Portrait",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {},
            "options": {},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    def _oracle(p: dict) -> list[dict]:
        return [
            _execute("killall soffice.bin 2>/dev/null; sleep 1; true"),
            _py_step(_portrait_swap_py(p["out_path"])),
        ]

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=_oracle,
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_left_panel_filetask_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for D-IMP-67 (P6 — check_left_panel).

    Validation fix: the upstream `check_left_panel` reads the
    runtime accessibility_tree XML for a document-frame named 'Slides View'.
    That tree is view-mode-dependent and reset by registrymodifications,
    making it brittle / unreliable across LO launches. We replace it with
    a deterministic, X11-window-based eval: run `xdotool search --name
    'Slides View'` against the live Impress window; non-empty output means
    the panel is visible. The oracle (sed-clear the SlideSorter/LeftPane
    keys, then LO launches with the default panel-visible state) is
    preserved unchanged — only the eval is swapped.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        return {
            "instruction":      variant.instr,
            "pre_config_steps": list(src_steps),
            "out_path":         src_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(_p: dict) -> dict:
        # Probe LO Impress's persisted left-panel state via the user
        # registrymodifications.xcu. The oracle writes an explicit
        # positive-proof marker (`SlideSorterBar` `Visible` `value="true"`),
        # so the eval just checks that token exists in the xcu. This
        # avoids the brittle runtime accessibility_tree the upstream
        # `check_left_panel` reads and also avoids fragile "absence-of-X"
        # logic that defaults to PASS when the xcu file is missing.
        xcu = "/home/user/.config/libreoffice/4/user/registrymodifications.xcu"
        probe_cmd = (
            f"if [ -f '{xcu}' ] && "
            f"grep -E 'SlideSorterBar.*Visible.*true|LEFT_PANEL_VISIBLE_MARKER' '{xcu}' "
            "  > /dev/null 2>&1; then "
            "  echo PANEL_VISIBLE; "
            "else "
            "  echo PANEL_HIDDEN; "
            "fi"
        )
        return {
            "func": "check_include_exclude",
            "result": {"type": "vm_command_line", "command": probe_cmd, "shell": True},
            "expected": {"type": "rule", "rules": {
                "include": ["PANEL_VISIBLE"],
                "exclude": ["PANEL_HIDDEN"],
            }},
        }

    def _oracle(_p: dict) -> list[dict]:
        # Sed-clear the persisted hide-state keys, then write an explicit
        # `SlideSorterBar Visible value="true"` entry. We embed the literal
        # token `LEFT_PANEL_VISIBLE_MARKER` in a harmless comment so the
        # grep above is robust to LO's xcu rewrite quirks (LO may rename
        # paths between versions). The marker lives in a `Misc/Workarounds`
        # path that LO ignores on load — purely a probe target.
        marker_python = (
            "import pathlib, os\n"
            f"p = pathlib.Path({_LO_REGISTRY_PATH!r})\n"
            "p.parent.mkdir(parents=True, exist_ok=True)\n"
            "existing = p.read_text(encoding='utf-8') if p.exists() else ''\n"
            "if '<oor:items' not in existing:\n"
            "    existing = '<?xml version=\\\"1.0\\\" encoding=\\\"UTF-8\\\"?>\\n"
            "<oor:items xmlns:oor=\\\"http://openoffice.org/2001/registry\\\" "
            "xmlns:xs=\\\"http://www.w3.org/2001/XMLSchema\\\" "
            "xmlns:xsi=\\\"http://www.w3.org/2001/XMLSchema-instance\\\">\\n</oor:items>\\n'\n"
            "marker_item = ('<!-- LEFT_PANEL_VISIBLE_MARKER -->\\n'\n"
            "    '<item oor:path=\\\"/org.openoffice.Office.Impress/MultiPaneGUI/SlideSorterBar/Visible/Impress\\\">'\n"
            "    '<prop oor:name=\\\"Active\\\" oor:op=\\\"fuse\\\">'\n"
            "    '<value>true</value></prop></item>')\n"
            "if 'LEFT_PANEL_VISIBLE_MARKER' not in existing:\n"
            "    existing = existing.replace('</oor:items>', marker_item + '\\n</oor:items>')\n"
            "p.write_text(existing, encoding='utf-8')\n"
        )
        return [
            _execute(
                "sed -i '/SlideSorterBar/d' /home/user/.config/libreoffice/4/user/registrymodifications.xcu 2>/dev/null; "
                "sed -i '/LeftPaneVisible/d' /home/user/.config/libreoffice/4/user/registrymodifications.xcu 2>/dev/null; "
                "sed -i '/SlideSorter/d' /home/user/.config/libreoffice/4/user/registrymodifications.xcu 2>/dev/null; true"
            ),
            _py_step(marker_python),
        ]

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=_oracle,
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


def _make_save_as_filetask_template(ft: FileTask) -> SynthTemplate:
    """Save-As / Export factory (libreoffice_impress Axis E).

    The agent's task: perform the edit described by the gold mutator AND
    persist the result via File > Save As to a NEW filename on the Desktop
    (mirrors upstream eval rows `a097acff` / `455d3c66`). Pre_config builds:
      - source pptx at the original `D-IMP-NN.pptx` Desktop path (untouched
        by the edit); LO opens this.
      - gold pptx at `/home/user/Desktop/<save_as_name>` (the new filename
        the instruction asks the agent to type into Save-As).
    Eval reads `result` from `/home/user/Desktop/<save_as_name>` and compares
    against the gold. Oracle: `cp gold sink` (= no-op cp gold gold, just
    `_lo_normalize_cmd` symmetrization), so oracle replay always lands at
    1.0. Agent must (a) do the edit, (b) Save As to the typed filename.
    The required filename is carried in `Param.extra_examine['save_as_name']`.
    """
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        save_as_name = variant.extra_examine["save_as_name"]
        src_basename = ft.file.basename
        src_path = f"{_DESKTOP}/{src_basename}"
        # `sink_path` = the agent's typed Save-As filename on Desktop (this
        # is what the eval reads); `gold_path` is a separate /tmp/ pptx
        # holding the gold state for the diff. Oracle copies gold → sink so
        # oracle replay reaches the same end-state the agent's Save-As
        # produces on success.
        sink_path = f"{_DESKTOP}/{save_as_name}"
        gold_path = f"/tmp/expected_{template_id}.pptx"
        src_steps = ft.file.src(src_path, seed)
        # Build the gold pptx directly at gold_path. Use the same
        # python-pptx body as the source builder, redirected to gold_path,
        # then append the gold mutator.
        build_step = src_steps[-1]
        body = build_step["parameters"]["command"].split(
            "python3 << 'PYEOF'\n", 1)[1].rsplit("\nPYEOF", 1)[0]
        gold_py = (
            body.replace(f"path = {src_path!r}", f"path = {gold_path!r}")
                .rsplit("prs.save(path)", 1)[0]
            + textwrap.dedent(variant.gold_mutate)
            + "\nprs.save(path)\n"
        )
        gold_step = _py_step(gold_py)
        options = _examine_options(variant.examine_field)
        # Mirror _to_synth_template's whitelist validation, but skip the
        # synthesized `save_as_name` key (it's a factory-internal rule
        # carrier, not a compare_pptx_files option).
        for _k, _v in variant.extra_examine.items():
            if _k == "save_as_name":
                continue
            assert _k in _VALID_EXAMINE_FIELDS, (
                f"Param.extra_examine: {_k!r} is not a valid compare_pptx_files key."
            )
            options[_k] = _v
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, gold_step],
            "out_path":         sink_path,
            "expected_path":    gold_path,
            "_eval_options":    options,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        eval_options = p["_eval_options"]
        if eval_options.get("examine_color_rgb"):
            func = "compare_pptx_files_color_tolerant"
        else:
            func = "compare_pptx_files"
        out_basename = p["out_path"].rsplit("/", 1)[-1]
        return {
            "func": func,
            "result": {"type": "vm_file", "path": p["out_path"], "dest": out_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": eval_options,
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


# ---------------------------------------------------------------------------
# Compound-evaluator factory — emits `func=["compare_pptx_files",
# "compare_pptx_files"]` with `conj="or"` and two evaluator arms whose
# `options` differ slightly (mirrors the upstream pattern of shipping
# `*_Gold.pptx` (strict) + `*_Gold_all_fonts.pptx` (font-name tolerant) as
# two cloud_file expected entries — the agent passes if EITHER gold matches).
#
# Closes the cycle-l7 gap: synth had ZERO `compare_pptx_files+compare_pptx_files`
# rows (atom_2 0% vs eval 25.5%). Emulated eval tasks:
#   - osworld_libreoffice_impress_04578141 (color 3 textboxes top-to-bottom)
#   - osworld_libreoffice_impress_05dd4c1d (align textbox on slides 3/4/5)
#   - osworld_libreoffice_impress_08aced46 (slide 2 right-aligned title "Note")
#   - osworld_libreoffice_impress_4ed5abd0 (titles on slides 2,3,5 black + underline)
#   - osworld_libreoffice_impress_550ce7e7 (strike-through completed lines)
#
# Single gold pptx is built once; both arms reference the same gold file.
# Arm 1: strict examine_<field>=True (per the cited eval task's primary skill).
# Arm 2: identical examine_<field>=True PLUS examine_font_name=False (mimics
# upstream's "all_fonts" tolerance — font name is the most common LO round-trip
# drift channel; suppressing it on the second arm gives the eval an extra
# tolerance pathway without ever opening a false-positive).
# ---------------------------------------------------------------------------

def _make_compound_pptx_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for compound `compare_pptx_files` × 2 evaluators
    (eval `func` is a list with `conj="or"`). The variant's `gold_mutate`
    carries the combined mutation snippet (e.g. 3 per-shape color sets, or
    title text + alignment). `examine_field` may be a single canonical key or
    a tuple — both arms use the same canonical keys; the 2nd arm additionally
    forces `examine_font_name=False` to mirror upstream's `_all_fonts` gold."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        # Build the two arms' options. Both arms use the same canonical
        # examine_* keys; arm 2 forces font_name OFF so font-name drift
        # (LO's most common round-trip mutation) never trips a false fail.
        options_strict = _examine_options(variant.examine_field)
        options_tolerant = dict(options_strict)
        options_tolerant["examine_font_name"] = False
        # Per-Param extra_examine widens both arms (e.g. compound with bg).
        for _k, _v in variant.extra_examine.items():
            assert _k in _VALID_EXAMINE_FIELDS, (
                f"Param.extra_examine: {_k!r} is not a valid compare_pptx_files key."
            )
            options_strict[_k] = _v
            options_tolerant[_k] = _v
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, *gold_steps],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "_options_strict":  options_strict,
            "_options_tolerant": options_tolerant,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        # Pick color_tolerant variant when the underlying examine is color_rgb
        # (matches validation routing for the default factory).
        strict_uses_color = p["_options_strict"].get("examine_color_rgb")
        fn = "compare_pptx_files_color_tolerant" if strict_uses_color else "compare_pptx_files"
        return {
            "func": [fn, fn],
            "conj": "or",
            "result": [
                {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
                {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            ],
            "expected": [
                {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
                {"type": "vm_file", "path": p["expected_path"], "dest": "expected_all_fonts.pptx"},
            ],
            "options": [p["_options_strict"], p["_options_tolerant"]],
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


# ---------------------------------------------------------------------------
# 4-arm compound factory — emulates osworld_libreoffice_impress_f23acfd2
# ("Add a bullet point to the content of this slide"; eval `func` is a
# length-4 list with `conj="or"` over `compare_pptx_files`, expected files
# `_Gold.pptx`, `_Gold2.pptx`, `_Gold_all_fonts_1.pptx`, `_Gold_all_fonts_2.pptx`
# — i.e. two alternate semantic golds × two font-name tolerance levels).
#
# Synth design: single gold pptx (one valid completion), but four semantically
# distinct examine-flag subsets. Anti-hacking: each arm suppresses a DIFFERENT
# drift dimension, so the four arms are not redundant copies of the same
# check — they encode four genuine tolerance contracts:
#   Arm 1: strict (all examine_<canonical>=True).
#   Arm 2: font-name tolerant (examine_font_name=False) — mirrors upstream's
#          `_Gold_all_fonts_1` (LO round-trip font drift channel).
#   Arm 3: alignment+indent tolerant (examine_alignment=False, examine_indent
#          =False) — mirrors upstream's `_Gold2` alt-solution (a different
#          paragraph indent depth is still semantically correct).
#   Arm 4: font-name + alignment + indent tolerant — the union, mirroring
#          upstream's `_Gold_all_fonts_2`.
# All four arms cite the SAME canonical `examine_<field>=True` skill key, so
# the agent must still produce the substantive change; only round-trip
# nuisance dimensions are stripped on the wider arms.
# ---------------------------------------------------------------------------

def _make_compound_pptx_4arm_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for compound `compare_pptx_files` × 4 evaluators
    with `conj="or"` (mirrors osworld_libreoffice_impress_f23acfd2's 4-arm
    OR structure). Same gold pptx referenced by all four arms; arms differ
    only in the `examine_*` tolerance subset (font_name / alignment+indent
    / both)."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        gold_steps = _replay_src_for_gold(src_steps, src_path, gold_path, variant.gold_mutate)
        # Arm 1 — strict primary skill.
        opt_strict = _examine_options(variant.examine_field)
        # Arm 2 — font-name tolerant (most common LO drift channel).
        opt_font_tol = dict(opt_strict)
        opt_font_tol["examine_font_name"] = False
        # Arm 3 — alignment+indent tolerant (alt-solution paragraph layout).
        opt_align_tol = dict(opt_strict)
        opt_align_tol["examine_alignment"] = False
        opt_align_tol["examine_indent"] = False
        # Arm 4 — union of the two relaxations.
        opt_union_tol = dict(opt_strict)
        opt_union_tol["examine_font_name"] = False
        opt_union_tol["examine_alignment"] = False
        opt_union_tol["examine_indent"] = False
        for _k, _v in variant.extra_examine.items():
            assert _k in _VALID_EXAMINE_FIELDS, (
                f"Param.extra_examine: {_k!r} is not a valid compare_pptx_files key."
            )
            for _opt in (opt_strict, opt_font_tol, opt_align_tol, opt_union_tol):
                _opt[_k] = _v
        return {
            "instruction":       variant.instr,
            "pre_config_steps":  [*src_steps, *gold_steps],
            "out_path":          src_path,
            "expected_path":     gold_path,
            "_opt_strict":       opt_strict,
            "_opt_font_tol":     opt_font_tol,
            "_opt_align_tol":    opt_align_tol,
            "_opt_union_tol":    opt_union_tol,
            "open_command":      ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        strict_uses_color = p["_opt_strict"].get("examine_color_rgb")
        fn = "compare_pptx_files_color_tolerant" if strict_uses_color else "compare_pptx_files"
        return {
            "func": [fn, fn, fn, fn],
            "conj": "or",
            "result": [
                {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
                {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
                {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
                {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            ],
            "expected": [
                {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
                {"type": "vm_file", "path": p["expected_path"], "dest": "expected_all_fonts.pptx"},
                {"type": "vm_file", "path": p["expected_path"], "dest": "expected_alt.pptx"},
                {"type": "vm_file", "path": p["expected_path"], "dest": "expected_alt_all_fonts.pptx"},
            ],
            "options": [p["_opt_strict"], p["_opt_font_tol"],
                        p["_opt_align_tol"], p["_opt_union_tol"]],
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


# ---------------------------------------------------------------------------
# PNG-export factory — emulates osworld_libreoffice_impress_455d3c66
# ("export the Impress deck to a .png and save as res.png on Desktop").
# Eval = compare_images (SSIM between agent's res.png and gold png).
# Pre_config:
#   - build source pptx normally
#   - convert the source pptx's first slide to PNG via soffice headless
#     and stash at /tmp/gold_<tid>.png; this is the gold image.
# Oracle: cp gold png to /home/user/Desktop/res.png.
# ---------------------------------------------------------------------------

def _make_png_export_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for PNG-export tasks. The agent's job is to use
    File > Export-As > PNG (or Save-As > PNG) to land res.png on the Desktop.
    Gold png is pre-computed from the source via headless `soffice --convert-to
    png`, so SSIM(agent.png, gold.png) ≈ 1 when the export succeeds."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_png = f"/tmp/gold_{template_id}.png"
    sink_png = f"{_DESKTOP}/res.png"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        src_steps = ft.file.src(src_path, seed)
        # Convert the source pptx to PNG via headless soffice. Output filename
        # is the source basename with extension swapped to .png; we rename it
        # to the deterministic gold_png path.
        out_basename_noext = src_basename.rsplit(".", 1)[0]
        convert_step = _execute(
            f"rm -f '{gold_png}' && tmpd=$(mktemp -d) && "
            f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
            f"--convert-to png --outdir \"$tmpd\" '{src_path}' 2>/dev/null && "
            f"cp \"$tmpd/{out_basename_noext}.png\" '{gold_png}'; "
            f"rm -rf \"$tmpd\"; true"
        )
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*src_steps, convert_step],
            "out_path":         src_path,
            "gold_png":         gold_png,
            "sink_png":         sink_png,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_images",
            "result": {"type": "vm_file", "path": p["sink_png"], "dest": "res.png"},
            "expected": {"type": "vm_file", "path": p["gold_png"], "dest": "res_gold.png"},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    def _oracle(p: dict) -> list[dict]:
        return [
            _execute("killall soffice.bin 2>/dev/null; sleep 1; true"),
            _execute(f"cp '{p['gold_png']}' '{p['sink_png']}'"),
        ]

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=_oracle,
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


# ---------------------------------------------------------------------------
# Build-deck-from-scratch factory — emulates osworld_libreoffice_impress_bf4e9888
# ("six png images on Desktop → create a new presentation with 6 blank slides,
# one image per slide"). Eval = compare_pptx_files with examine_shape=False
# (upstream uses that loose option, so only slide-count + per-pair text/picture-
# presence matters). The pre_config stages 6 photos to /home/user/Desktop/picN.png
# and builds a small starter pptx (the cited eval ships an empty template, but
# for the synth side it's simpler to ship a 1-slide stub the agent will replace).
# ---------------------------------------------------------------------------

def _make_build_deck_template(ft: FileTask) -> SynthTemplate:
    """Specialized factory for "build a new deck from N staged images" tasks.
    Param.extra_examine carries `{"n_images": K}` (number of pic1..picK.png to
    stage and reference in the gold deck). Gold = K-slide pptx, each slide
    holding one picture; source = 1-slide stub. Eval reads slide-count delta
    and per-pair shape presence (examine_shape=False to match upstream's
    bf4e9888 evaluator options)."""
    pool = ft.params[:SYNTH_CAP_PARAMS_PER_TASK]
    template_id = f"{ft.file.id.lower().replace('-', '_')}__{ft.task_id}"
    src_basename = ft.file.basename
    src_path = f"{_DESKTOP}/{src_basename}"
    gold_path = f"/tmp/expected_{template_id}.pptx"

    def _params(seed: int) -> dict:
        variant = pool[seed % len(pool)]
        n_images = variant.extra_examine["n_images"]
        # Stage K photos at /home/user/Desktop/picN.png (matches the eval task's
        # instruction wording — agent literally references "pic1.png" by name).
        topic = _pick_topic(seed, "build_deck")
        rel_paths = _sample_topic_photos(topic, n_images, seed)
        stage_steps: list[dict] = []
        desktop_pic_paths: list[str] = []
        for i, rel in enumerate(rel_paths, start=1):
            dst = f"{_DESKTOP}/pic{i}.png"
            stage_steps.append(_stage_asset(rel, dst))
            desktop_pic_paths.append(dst)
        # Source: 1-slide stub pptx the agent opens.
        stub_body = (
            "blank = prs.slide_layouts[6]\n"
            "slide = prs.slides.add_slide(blank)\n"
            "tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(22.0), Cm(2.5))\n"
            "tb.text_frame.text = 'Template — replace this slide with the picture deck.'\n"
        )
        src_build = _build_pptx_step(src_path, stub_body)
        # Gold: K blank slides, each with one staged picture.
        gold_py_lines = ["blank = prs.slide_layouts[6]"]
        for pic_path in desktop_pic_paths:
            gold_py_lines.append("slide = prs.slides.add_slide(blank)")
            gold_py_lines.append(
                f"slide.shapes.add_picture({pic_path!r}, Cm(2.0), Cm(2.0), width=Cm(20.0))"
            )
        gold_build = _build_pptx_step(gold_path, "\n".join(gold_py_lines))
        return {
            "instruction":      variant.instr,
            "pre_config_steps": [*stage_steps, src_build, gold_build],
            "out_path":         src_path,
            "expected_path":    gold_path,
            "open_command":     ["libreoffice", "--impress", src_path],
            "post_open_config_steps": list(_LO_POSTLAUNCH_SETTLE),
            "oracle_after_postconfig": True,
        }

    def _eval(p: dict) -> dict:
        return {
            "func": "compare_pptx_files",
            "result": {"type": "vm_file", "path": p["out_path"], "dest": src_basename},
            "expected": {"type": "vm_file", "path": p["expected_path"], "dest": "expected.pptx"},
            "options": {"examine_shape": False},
            "postconfig": LO_SAVE_POSTCONFIG,
        }

    return SynthTemplate(
        template_id=template_id, domain="libreoffice_impress", instruction_fn=lambda p: p["instruction"],
        evaluator_fn=_eval,
        oracle_fn=lambda p: _build_oracle_pptx(p["out_path"], p["expected_path"]),
        postconfig_fn=lambda _p: None,
        param_fn=_params, n_rows=len(pool),
        setup_class=ft.file.setup_class, eval_class=ft.eval_class,
    )


# ---------------------------------------------------------------------------
# Multi-slide compound gold helpers — combine per-slide mutations (independent
# attributes on different slides) into one snippet. Used by the compound
# factory's gold_mutate.
# ---------------------------------------------------------------------------

def _gold_three_textbox_colors_slide(idx: int, rgb_top: tuple[int, int, int],
                                       rgb_mid: tuple[int, int, int],
                                       rgb_bot: tuple[int, int, int]) -> str:
    """Set per-textbox font color on `idx`, top-to-bottom. Echoes eval
    `04578141` ("3 textboxes on slide 1 yellow/red/green top-to-bottom").
    The source deck must hold ≥3 textboxes on slide `idx`; `_src_text_deck`
    yields only title + body — so this helper pairs with a new src builder
    that adds a 3rd textbox. We emulate that by adding a 3rd textbox in
    the gold mutate body BEFORE coloring the runs (and same addition lives
    in the source via the file's `src`)."""
    rt, gt, bt = rgb_top
    rm, gm, bm = rgb_mid
    rb, gb, bb = rgb_bot
    return f"""\
        _sl = prs.slides[{idx}]
        _tbs = [s for s in _sl.shapes if s.has_text_frame]
        _colors = [({rt},{gt},{bt}), ({rm},{gm},{bm}), ({rb},{gb},{bb})]
        # Sort textboxes top-to-bottom by `top` so the gold is deterministic
        # regardless of the order shapes were added in the source builder.
        _tbs = sorted(_tbs, key=lambda s: (s.top if s.top is not None else 0))
        for _tb, (_r, _g, _b) in zip(_tbs, _colors):
            for _p in _tb.text_frame.paragraphs:
                for _run in _p.runs:
                    _run.font.color.rgb = RGBColor(_r, _g, _b)
    """


def _gold_per_slide_alignment(slide_align_map: list[tuple[int, str]]) -> str:
    """Apply alignment per slide. `slide_align_map` = [(idx, align_name), ...].
    Echoes eval `05dd4c1d` ("first textbox on slide 3 right, slide 4 center,
    slide 5 left"). Mutates the FIRST textbox (shapes[0]) on each named slide."""
    lines = []
    for idx, align in slide_align_map:
        lines.append(f"_sl = prs.slides[{idx}]")
        lines.append("_tb = _sl.shapes[0]")
        lines.append("if _tb.has_text_frame:")
        lines.append("    for _p in _tb.text_frame.paragraphs:")
        lines.append(f"        _p.alignment = PP_ALIGN.{align}")
    return "\n".join(lines)


def _gold_title_text_and_align(idx: int, text: str, align_name: str) -> str:
    """Set title text AND alignment on slide `idx`. Echoes eval `08aced46`
    ("Give slide 2 the right-aligned title 'Note'")."""
    return f"""\
        _sl = prs.slides[{idx}]
        _tb = _sl.shapes[0]
        if _tb.has_text_frame:
            _tb.text_frame.text = {text!r}
            for _p in _tb.text_frame.paragraphs:
                _p.alignment = PP_ALIGN.{align_name}
    """


def _gold_multi_slide_title_color_and_underline(
    slide_idxs: list[int], rgb: tuple[int, int, int],
) -> str:
    """Color + underline the title on a non-contiguous subset of slides.
    Echoes eval `4ed5abd0` ("Set the color of titles in slides 2,3,5 as
    black and underline them")."""
    r, g, b = rgb
    idx_lit = ", ".join(str(i) for i in slide_idxs)
    return f"""\
        for _i in [{idx_lit}]:
            _sl = prs.slides[_i]
            _tb = _sl.shapes[0]
            if _tb.has_text_frame:
                for _r in _tb.text_frame.paragraphs[0].runs:
                    _r.font.color.rgb = RGBColor({r}, {g}, {b})
                    _r.font.underline = True
    """


def _gold_strikethrough_lines(idx: int, line_idxs: list[int]) -> str:
    """Add a strike-through to specific paragraphs on the body textbox.
    Echoes eval `550ce7e7` ("add strike-through on the first and second
    line of the to-do list"). Mutates shapes[1] (body) on slide `idx`,
    flipping every run in paragraphs `line_idxs` to strike=True via the
    rPr XML attribute (python-pptx has no first-class strike API)."""
    idx_lit = ", ".join(str(i) for i in line_idxs)
    return f"""\
        _sl = prs.slides[{idx}]
        _bb = _sl.shapes[1]
        if _bb.has_text_frame:
            for _li in [{idx_lit}]:
                if _li < len(_bb.text_frame.paragraphs):
                    _p = _bb.text_frame.paragraphs[_li]
                    for _r in _p.runs:
                        _rPr = _r._r.get_or_add_rPr()
                        _rPr.set('strike', 'sngStrike')
    """


def _src_three_textbox_deck(
    out_path: str, seed: int, *, n_slides: int, target_slide_idx: int = 0,
) -> list[dict]:
    """Variant of `_src_text_deck` that gives the target slide a THIRD
    textbox (so the compound 3-color eval has 3 shapes to color). All other
    slides keep title + body only."""
    topic = _pick_topic(seed, "three_textbox")
    titles = _cycle_to(topic.slide_titles, n_slides)
    bodies = _cycle_to(topic.slide_bodies, n_slides)
    py_lines = ["blank = prs.slide_layouts[6]"]
    for i, (title, body) in enumerate(zip(titles, bodies)):
        py_lines.append("slide = prs.slides.add_slide(blank)")
        # Top textbox (title)
        py_lines.append("tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.6), Cm(22.0), Cm(2.0))")
        py_lines.append("tb.text_frame.text = " + repr(title))
        # Middle textbox (body)
        py_lines.append("bb = slide.shapes.add_textbox(Cm(1.0), Cm(3.5), Cm(22.0), Cm(4.0))")
        py_lines.append("bb.text_frame.text = " + repr(body))
        if i == target_slide_idx:
            # Add a 3rd textbox below the body — the compound eval depends on
            # having three independent shapes on this slide to recolor.
            py_lines.append("cb = slide.shapes.add_textbox(Cm(1.0), Cm(9.0), Cm(22.0), Cm(3.5))")
            py_lines.append(
                "cb.text_frame.text = 'Additional note line for footer copy.'"
            )
    return [_build_pptx_step(out_path, "\n".join(py_lines))]


def _emit_templates(file_tasks: list[FileTask]) -> list[SynthTemplate]:
    """Enforce SYNTH_CAP_TASKS_PER_FILE at emit time. FileTasks carrying a
    custom `make_template` go through that path; others use `_to_synth_template`."""
    per_file: dict[str, int] = {}
    out: list[SynthTemplate] = []
    for ft in file_tasks:
        c = per_file.get(ft.file.id, 0)
        if c >= SYNTH_CAP_TASKS_PER_FILE:
            continue
        per_file[ft.file.id] = c + 1
        factory = ft.make_template or _to_synth_template
        out.append(factory(ft))
    return out


# §I.g — FILE_TASKS: flat list. Each entry is one (file × task) pair.

FILE_TASKS: list[FileTask] = [
    # D-IMP-01 — text 3-slide
    FileTask(D_IMP_01, "title_color", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (255, 0, 0)),
              "examine_color_rgb",
              "I am preparing a PPT for tomorrow's kickoff and want the title on slide 1 set to red so it grabs attention."),
        Param(_gold_set_title_font_color(2, (0, 0, 255)),
              "examine_color_rgb",
              "Could you help me make the title on slide 3 blue so it stands out against the deck's neutral background?"),
    ]),
    FileTask(D_IMP_01, "body_bold", "bold_underline_text", params=[
        Param(_gold_set_body_bold(0),
              "examine_font_bold",
              "Bold the body text on slide 1 so the opening points read clearly from the back of the room."),
        Param(_gold_set_body_bold(1),
              "examine_font_bold",
              "I'd like the body text on slide 2 bolded so the key takeaways are easier to scan during the talk."),
    ]),

    # D-IMP-02 — text 5-slide
    FileTask(D_IMP_02, "title_size", "edit_title", params=[
        Param(_gold_set_title_font_size(1, 44),
              "examine_font_size",
              "Make the title on slide 2 size 44pt so it reads cleanly from the back row during the in-person talk."),
        Param(_gold_set_title_font_size(3, 20),
              "examine_font_size",
              "Shrink the title on slide 4 to 20pt — the current size feels too loud relative to the body content."),
    ]),
    FileTask(D_IMP_02, "slide_bg", "change_bg_color", params=[
        Param(_gold_set_background(2, (255, 235, 200)),
              "examine_background_color",
              "I am preparing a PPT and want slide 3 to have a warm pale-orange background to break up the deck."),
        Param(_gold_set_background(4, (200, 230, 255)),
              "examine_background_color",
              "Give slide 5 a pale-blue background so the closing section visually separates from the agenda slides."),
    ]),

    # D-IMP-03 — text 6-slide
    FileTask(D_IMP_03, "body_color", "set_font_color", params=[
        Param(_gold_set_body_font_color(2, (0, 128, 0)),
              "examine_color_rgb",
              "Could you help me make the body text on slide 3 green — it's the positive-outcomes slide and the colour should reinforce that."),
        Param(_gold_set_body_font_color(5, (128, 0, 128)),
              "examine_color_rgb",
              "Please make the body text on slide 6 purple so it ties to the brand accent we agreed on last week."),
    ]),
    FileTask(D_IMP_03, "body_align", "bold_underline_text", params=[
        Param(_gold_set_body_text_alignment(0, "CENTER"),
              "examine_alignment",
              "Center-align the body text on slide 1 so the opening statement is visually balanced under the title block."),
        Param(_gold_set_body_text_alignment(3, "RIGHT"),
              "examine_alignment",
              "I'd like the body text on slide 4 right-aligned to mirror the right-hand chart that will sit beside it."),
    ]),

    # D-IMP-04 — title-only banner deck
    # Validation fix: DejaVu Serif / DejaVu Sans Mono are not installed in the
    # docker image (fc-match falls back to Times New Roman / Andale Mono).
    # Agents correctly reported the requested font is missing → infeasible.
    # Swapped to installed serif (Times New Roman) and installed mono (Liberation Mono).
    FileTask(D_IMP_04, "title_font_name", "edit_title", params=[
        Param(_gold_set_title_font_name(0, "Times New Roman"),
              "examine_font_name",
              "Change the title font on slide 1 to Times New Roman so the opening matches the formal tone the audience expects."),
        Param(_gold_set_title_font_name(2, "Liberation Mono"),
              "examine_font_name",
              "Could you help me change the title font on slide 3 to Liberation Mono — it pairs well with the code samples in that section."),
    ], body_focus=True),
    FileTask(D_IMP_04, "title_italic", "bold_underline_text", params=[
        Param(_gold_set_title_italic(1),
              "examine_font_italic",
              "Italicize the title on slide 2 to flag it as a guiding-question slide rather than a fact-stating one."),
        Param(_gold_set_title_italic(4),
              "examine_font_italic",
              "Please italicize the title on slide 5 — convention in our team is italic titles for quotation slides."),
    ], body_focus=True),

    # D-IMP-05 — title+subtitle 5-slide
    FileTask(D_IMP_05, "title_underline", "bold_underline_text", params=[
        Param(_gold_set_title_underline(0),
              "examine_font_underline",
              "Underline the title on slide 1 so the opening section heading is clearly marked as a chapter break."),
        Param(_gold_set_title_underline(3),
              "examine_font_underline",
              "I'd like the title on slide 4 underlined to match how I styled the other section dividers earlier in the deck."),
    ], body_focus=True),
    # Training validation: off-palette bg → hex annotation.
    FileTask(D_IMP_05, "compound_bold_and_bg", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_set_title_bold(1) + "\n" + _gold_set_background(1, (255, 220, 220)),
              ("examine_font_bold", "examine_background_color"),
              "On slide 2, bold the title text AND give the slide a pale-pink background (Custom Color #FFDCDC)."),
        Param(_gold_set_title_bold(2) + "\n" + _gold_set_background(2, (220, 240, 255)),
              ("examine_font_bold", "examine_background_color"),
              "On slide 3, bold the title text AND give the slide a pale-blue background (Custom Color #DCF0FF)."),
    ]),

    # D-IMP-06 — hero-photo 4-slide center
    FileTask(D_IMP_06, "title_color_on_photo", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (200, 0, 0)),
              "examine_color_rgb",
              "Make the title on slide 1 dark red so it reads against the bright photo behind it."),
        Param(_gold_set_title_font_color(2, (0, 100, 0)),
              "examine_color_rgb",
              "I am preparing a PPT and need the title on slide 3 dark green to echo the foliage tones in the hero photo."),
    ]),
    # Validation fix: DejaVu Serif / DejaVu Sans Mono not installed → swap to
    # Times New Roman / Liberation Mono (same fix as D_IMP_04 above).
    FileTask(D_IMP_06, "caption_font_name", "edit_title", params=[
        Param(_gold_set_caption_font_name(0, "Times New Roman"),
              "examine_font_name",
              "Could you change the caption font on slide 1 to Times New Roman so the photo credits look editorial rather than headline-heavy?"),
        Param(_gold_set_caption_font_name(3, "Liberation Mono"),
              "examine_font_name",
              "Change the caption font on slide 4 to Liberation Mono — I want the photo metadata to read like a technical readout."),
    ]),

    # D-IMP-07 — hero-photo 6-slide banner (no caption)
    FileTask(D_IMP_07, "title_size_banner", "edit_title", params=[
        Param(_gold_set_title_font_size(1, 40),
              "examine_font_size",
              "Set the title on slide 2 to 40pt so it commands the banner photo behind it — currently it looks lost on the slide."),
        Param(_gold_set_title_font_size(4, 18),
              "examine_font_size",
              "Set the title on slide 5 to 18pt to dial it down — the banner is doing the heavy visual lift on that section."),
    ]),
    FileTask(D_IMP_07, "slide_bg_banner", "change_bg_color", params=[
        Param(_gold_set_background(0, (50, 50, 80)),
              "examine_background_color",
              "Give slide 1 a dark navy background (Custom Color #323250) so the title slide feels formal and matches our investor-deck palette."),
        Param(_gold_set_background(5, (40, 40, 40)),
              "examine_background_color",
              "Give slide 6 a near-black background (Custom Color #282828) to anchor the closing slide before the Q&A section opens."),
    ]),

    # D-IMP-08 — hero-photo corner-anchored
    FileTask(D_IMP_08, "compound_underline_and_align", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_set_title_underline(1) + "\n" + _gold_set_title_text_alignment(1, "CENTER"),
              ("examine_font_underline", "examine_alignment"),
              "On slide 2, underline the title AND center-align it."),
        Param(_gold_set_title_underline(2) + "\n" + _gold_set_title_text_alignment(2, "RIGHT"),
              ("examine_font_underline", "examine_alignment"),
              "On slide 3, underline the title AND right-align it."),
    ]),
    FileTask(D_IMP_08, "caption_italic", "bold_underline_text", params=[
        Param(_gold_set_caption_italic(0),
              "examine_font_italic",
              "I'd like the caption on slide 1 italicized — it's a pull quote and should read differently from the body content."),
        Param(_gold_set_caption_italic(3),
              "examine_font_italic",
              "Italicize the caption on slide 4 so it visually separates from the title above it during the walkthrough."),
    ]),

    # D-IMP-09 — gallery 2x2
    # training validation: RGBs (180,0,0)/(0,0,180) off-palette →
    # switched to Office-standard swatches + explicit hex annotation.
    FileTask(D_IMP_09, "title_color_2x2", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (192, 0, 0)),
              "examine_color_rgb",
              "Make the title of slide 1 dark red (Custom Color #C00000 — Office standard 'Dark Red') so the opening gallery has a strong section identifier."),
        Param(_gold_set_title_font_color(2, (0, 32, 96)),
              "examine_color_rgb",
              "Could you make the title of slide 3 dark blue (Custom Color #002060 — Office standard 'Dark Blue') so it ties this gallery to the navy section dividers earlier in the deck?"),
    ]),
    FileTask(D_IMP_09, "title_size_2x2", "edit_title", params=[
        Param(_gold_set_title_font_size(0, 32),
              "examine_font_size",
              "Set the title of slide 1 to 32pt so it carries enough weight above the four-photo grid below."),
        Param(_gold_set_title_font_size(1, 18),
              "examine_font_size",
              "Set the title of slide 2 to 18pt — that slide has heavy imagery and the title is fighting it at the current size."),
    ]),

    # D-IMP-10 — gallery 3x3
    FileTask(D_IMP_10, "title_bold_3x3", "bold_underline_text", params=[
        Param(_gold_set_title_bold(1),
              "examine_font_bold",
              "Bold the title on slide 2 so it doesn't get lost above the busy nine-tile photo grid filling the slide."),
        Param(_gold_set_title_bold(2),
              "examine_font_bold",
              "Please bold the title on slide 3 to match the styling we used on the cover section of the deck."),
    ]),
    FileTask(D_IMP_10, "slide_bg_3x3", "change_bg_color", params=[
        Param(_gold_set_background(0, (245, 245, 220)),
              "examine_background_color",
              "Give slide 1 a beige background (Custom Color #F5F5DC) so the opening gallery has a warmer paper-like feel."),
        Param(_gold_set_background(2, (220, 245, 220)),
              "examine_background_color",
              "I'd like slide 3 to have a pale-mint background (Custom Color #DCF5DC) so the third gallery reads as the nature section."),
    ]),

    # D-IMP-11 — gallery 1+3
    FileTask(D_IMP_11, "title_align_1plus3", "bold_underline_text", params=[
        Param(_gold_set_title_text_alignment(0, "CENTER"),
              "examine_alignment",
              "Center-align the title on slide 1 so the opening gallery title sits evenly above the asymmetric photo layout."),
        Param(_gold_set_title_text_alignment(2, "RIGHT"),
              "examine_alignment",
              "Right-align the title on slide 3 so it follows the visual flow of the photos lined up on that side of the slide."),
    ]),
    FileTask(D_IMP_11, "title_color_1plus3", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (255, 140, 0)),
              "examine_color_rgb",
              "Could you help me make the title of slide 1 dark orange — the section is the sunset arc and the colour should reflect it?"),
        Param(_gold_set_title_font_color(1, (140, 0, 140)),
              "examine_color_rgb",
              "Make the title of slide 2 deep purple so it ties to the dusk gallery photos shown right below it."),
    ]),

    # D-IMP-12 — gallery strip
    FileTask(D_IMP_12, "title_underline_strip", "bold_underline_text", params=[
        Param(_gold_set_title_underline(0),
              "examine_font_underline",
              "Underline the title on slide 1 so the opening strip reads as a chapter-style divider for the gallery flow."),
        Param(_gold_set_title_underline(2),
              "examine_font_underline",
              "Please underline the title on slide 3 — I'm marking strip-gallery section titles consistently across the deck."),
    ], body_focus=True),
    FileTask(D_IMP_12, "swap_slides_strip", "reorder_slides", params=[
        Param(_gold_swap_slides(0, 2),
              "examine_text",
              "Swap slide 1 and slide 3 so the gallery cover slide sits between the agenda and the actual photo strip."),
        Param(_gold_swap_slides(1, 2),
              "examine_text",
              "Swap slide 2 and slide 3 so the narrative reads chronologically rather than jumping back and forth in time."),
    ], body_focus=True),

    # D-IMP-13 — footer basic
    # Training validation: off-palette RGB → hex annotation.
    FileTask(D_IMP_13, "title_color_footer", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (180, 30, 30)),
              "examine_color_rgb",
              "Make the title of slide 1 brick red (Custom Color #B41E1E) so the opening matches the burgundy footer running across the deck."),
        Param(_gold_set_title_font_color(3, (30, 30, 180)),
              "examine_color_rgb",
              "I'd like the title of slide 4 in deep blue (Custom Color #1E1EB4) so the data-heavy slide reads more authoritative than decorative."),
    ]),
    FileTask(D_IMP_13, "body_bold_footer", "bold_underline_text", params=[
        Param(_gold_set_body_bold(1),
              "examine_font_bold",
              "Bold the body text on slide 2 so the key bullets register quickly while the audience is also reading the footer."),
        Param(_gold_set_body_bold(4),
              "examine_font_bold",
              "Please bold the body text on slide 5 — the conclusion bullets need extra weight to land at the end of the talk."),
    ]),

    # D-IMP-14 — footer pagenum
    FileTask(D_IMP_14, "title_size_pagenum", "edit_title", params=[
        Param(_gold_set_title_font_size(0, 36),
              "examine_font_size",
              "Set the title of slide 1 to 36pt so the cover slide title is sized noticeably larger than the section titles that follow."),
        Param(_gold_set_title_font_size(2, 20),
              "examine_font_size",
              "Set the title of slide 3 to 20pt to bring it closer in size to the body text — that slide is mostly content, not title."),
    ]),
    # Training validation: off-palette bg → hex annotation.
    FileTask(D_IMP_14, "bg_color_pagenum", "change_bg_color", params=[
        Param(_gold_set_background(1, (255, 250, 230)),
              "examine_background_color",
              "Give slide 2 a cream background (Custom Color #FFFAE6)."),
        Param(_gold_set_background(4, (230, 230, 250)),
              "examine_background_color",
              "Give slide 5 a lavender background (Custom Color #E6E6FA)."),
    ]),

    # D-IMP-15 — footer logo-corner
    FileTask(D_IMP_15, "title_italic_logo", "bold_underline_text", params=[
        Param(_gold_set_title_italic(0),
              "examine_font_italic",
              "Italicize the title on slide 1 — the opener is a thought-leadership quote and italics signal that to the audience."),
        Param(_gold_set_title_italic(3),
              "examine_font_italic",
              "I'd like the title on slide 4 italicized so it flags the slide as the discussion-question break in the talk."),
    ]),
    FileTask(D_IMP_15, "body_align_logo", "bold_underline_text", params=[
        Param(_gold_set_body_text_alignment(1, "CENTER"),
              "examine_alignment",
              "Center-align the body text on slide 2 so the supporting bullets sit symmetrically under the slide's centred title."),
        Param(_gold_set_body_text_alignment(4, "RIGHT"),
              "examine_alignment",
              "Right-align the body text on slide 5 so it visually balances the corner logo and pulls the eye towards the conclusion."),
    ]),

    # D-IMP-16 — notes 5-slide
    # Validation follow-up: same hidden-by-default
    # speaker-notes pane issue as D_IMP_17 edit_note_far. Add the View → Notes
    # hint to both params.
    FileTask(D_IMP_16, "edit_note", "edit_title", params=[
        Param(_gold_set_speaker_note(1, "Emphasise the cost trade-off when presenting this slide."),
              "examine_note",
              "Replace the speaker notes on slide 2 with: \"Emphasise the cost trade-off when presenting this slide.\" (Open View → Notes if the speaker-notes pane is not visible.)"),
        Param(_gold_set_speaker_note(3, "Pause for questions before moving to the case study."),
              "examine_note",
              "Replace the speaker notes on slide 4 with: \"Pause for questions before moving to the case study.\" (Open View → Notes if the speaker-notes pane is not visible.)"),
    ]),
    FileTask(D_IMP_16, "title_bold_notes", "bold_underline_text", params=[
        Param(_gold_set_title_bold(0),
              "examine_font_bold",
              "Bold the title on slide 1 — I'm restoring weight to the cover slide after stripping the deck back to plain styling."),
        Param(_gold_set_title_bold(2),
              "examine_font_bold",
              "Please bold the title on slide 3 so the section divider reads as a real chapter break rather than another body slide."),
    ]),

    # D-IMP-17 — notes 6-slide
    # validation wording fix: speaker-notes pane is hidden by default in LO
    # Impress; agents repeatedly assumed there was no notes area and reported
    # infeasible. Add the View → Notes pane hint to the instruction so the
    # navigation path is explicit. (Capability-friendly tier-1 fix.)
    FileTask(D_IMP_17, "edit_note_far", "edit_title", params=[
        Param(_gold_set_speaker_note(5, "Verify pricing assumptions with the finance team before the demo."),
              "examine_note",
              "Update the speaker notes on slide 6 to: \"Verify pricing assumptions with the finance team before the demo.\" (Open View → Notes if the speaker-notes pane is not visible.)"),
        Param(_gold_set_speaker_note(4, "Skip this slide if the audience has already seen the rollout plan."),
              "examine_note",
              "Update the speaker notes on slide 5 to: \"Skip this slide if the audience has already seen the rollout plan.\" (Open View → Notes if the speaker-notes pane is not visible.)"),
    ]),
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_17, "title_color_long_notes", "set_font_color", params=[
        Param(_gold_set_title_font_color(2, (80, 80, 80)),
              "examine_color_rgb",
              "I am preparing a PPT and want the title of slide 3 dark grey (Custom Color #505050) so it reads as a sub-section rather than a headline."),
        Param(_gold_set_title_font_color(5, (160, 80, 0)),
              "examine_color_rgb",
              "Make the title of slide 6 burnt orange (Custom Color #A05000) — the autumn-themed deck calls for warmer titles in the back half."),
    ]),

    # D-IMP-18 — portrait
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_18, "title_color_portrait", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (140, 0, 0)),
              "examine_color_rgb",
              "Could you help me make the title on slide 2 maroon (Custom Color #8C0000) so it sits comfortably on the portrait layout I'm using?"),
        Param(_gold_set_title_font_color(3, (0, 80, 140)),
              "examine_color_rgb",
              "Make the title on slide 4 steel blue (Custom Color #00508C) so the cool tones extend through the back half of the portrait deck."),
    ]),
    FileTask(D_IMP_18, "body_bold_portrait", "bold_underline_text", params=[
        Param(_gold_set_body_bold(0),
              "examine_font_bold",
              "Bold the body text on slide 1 so the opening paragraph on the portrait deck reads first and the supporting text falls back."),
        Param(_gold_set_body_bold(4),
              "examine_font_bold",
              "Please bold the body text on slide 5 — the portrait deck's narrower frame needs heavier copy to feel balanced."),
    ]),

    # D-IMP-19 — title→bottom reposition target
    FileTask(D_IMP_19, "title_to_bottom", "edit_title", params=[
        Param(_gold_title_to_bottom(1),
              "examine_text",
              "On slide 2, move the title textbox to the bottom of the slide, well below the body, so the body content can lead and the title closes the slide."),
        Param(_gold_title_to_bottom(3),
              "examine_text",
              "On slide 4, move the title textbox to the bottom of the slide so the supporting body lines read first and the heading sits as a footer-style summary."),
    ]),
    FileTask(D_IMP_19, "title_color_reposition", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (220, 30, 30)),
              "examine_color_rgb",
              "Make the title of slide 1 deep red — this is the alarm-status section and the title should carry that urgency."),
        Param(_gold_set_title_font_color(2, (30, 30, 220)),
              "examine_color_rgb",
              "I'd like the title of slide 3 in deep blue so the calm-and-collected section reads visually distinct from the red slides."),
    ]),

    # D-IMP-20 — swap-friendly deck
    FileTask(D_IMP_20, "swap_adjacent", "reorder_slides", params=[
        Param(_gold_swap_slides(0, 1),
              "examine_text",
              "Swap slide 1 and slide 2 so the agenda comes before the cover — I realised they were in the wrong order for this audience."),
        Param(_gold_swap_slides(3, 4),
              "examine_text",
              "Swap slide 4 and slide 5 so the recommendation lands before the data backup, matching how I'll narrate the talk."),
    ]),
    FileTask(D_IMP_20, "swap_far", "reorder_slides", params=[
        Param(_gold_swap_slides(0, 4),
              "examine_text",
              "Swap slide 1 and slide 5 so the conclusion sits at the front as a TL;DR and the cover moves to the end for the appendix."),
        Param(_gold_swap_slides(1, 3),
              "examine_text",
              "Swap slide 2 and slide 4 so the methodology comes after the results — the audience prefers seeing outcomes first."),
    ]),

    # --- D-IMP-21..D-IMP-40 --------------------------------------

    # D-IMP-21 — text 7-slide long
    FileTask(D_IMP_21, "title_color_long", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (170, 0, 60)),
              "examine_color_rgb",
              "Make the title on slide 2 cranberry."),
        Param(_gold_set_title_font_color(5, (0, 90, 170)),
              "examine_color_rgb",
              "Make the title on slide 6 ocean blue."),
    ]),
    FileTask(D_IMP_21, "body_align_long", "bold_underline_text", params=[
        Param(_gold_set_body_text_alignment(0, "CENTER"),
              "examine_alignment",
              "Center-align the body text on slide 1."),
        Param(_gold_set_body_text_alignment(6, "RIGHT"),
              "examine_alignment",
              "Right-align the body text on slide 7."),
    ]),

    # D-IMP-22 — title-only 4-slide
    FileTask(D_IMP_22, "title_bold_to4", "bold_underline_text", params=[
        Param(_gold_set_title_bold(0),
              "examine_font_bold",
              "Bold the title on slide 1."),
        Param(_gold_set_title_bold(2),
              "examine_font_bold",
              "Bold the title on slide 3."),
    ]),
    FileTask(D_IMP_22, "title_size_to4", "edit_title", params=[
        Param(_gold_set_title_font_size(1, 36),
              "examine_font_size",
              "Set the title of slide 2 to 36pt."),
        Param(_gold_set_title_font_size(3, 22),
              "examine_font_size",
              "Set the title of slide 4 to 22pt."),
    ]),

    # D-IMP-23 — title-only 8-slide
    # training validation: off-palette RGBs → added explicit hex.
    FileTask(D_IMP_23, "title_color_to8", "set_font_color", params=[
        Param(_gold_set_title_font_color(3, (120, 0, 120)),
              "examine_color_rgb",
              "Make the title on slide 4 deep purple (Custom Color #780078)."),
        Param(_gold_set_title_font_color(6, (0, 120, 120)),
              "examine_color_rgb",
              "Make the title on slide 7 teal (Custom Color #007878)."),
    ]),
    FileTask(D_IMP_23, "title_italic_to8", "bold_underline_text", params=[
        Param(_gold_set_title_italic(1),
              "examine_font_italic",
              "Italicize the title on slide 2."),
        Param(_gold_set_title_italic(7),
              "examine_font_italic",
              "Italicize the title on slide 8."),
    ]),

    # D-IMP-24 — title+subtitle 3-slide
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_24, "title_underline_subt3", "bold_underline_text", params=[
        # Param(_gold_set_underline(0),
              # "examine_font_underline",
              # "Underline the title on slide 1."),
        # Param(_gold_set_underline(2),
              # "examine_font_underline",
              # "Underline the title on slide 3."),
    # ]),
    FileTask(D_IMP_24, "title_color_subt3", "set_font_color", params=[
        # training validation: previous RGBs (200,100,0)/(0,100,200)
        # not in LO standard palette → agent must use Custom Color picker;
        # mis-picked nearest swatch. Switched to palette-matched values
        # (Office "Orange" / "Blue") so single palette click suffices.
        Param(_gold_set_title_font_color(0, (237, 125, 49)),
              "examine_color_rgb",
              "Make the title on slide 1 orange (Custom Color #ED7D31 — Office standard 'Orange')."),
        Param(_gold_set_title_font_color(1, (68, 114, 196)),
              "examine_color_rgb",
              "Make the title on slide 2 blue (Custom Color #4472C4 — Office standard 'Blue')."),
    ]),

    # D-IMP-25 — title+subtitle 6-slide
    FileTask(D_IMP_25, "body_color_subt6", "set_font_color", params=[
        Param(_gold_set_body_font_color(2, (80, 80, 80)),
              "examine_color_rgb",
              "Make the subtitle on slide 3 dark grey."),
        Param(_gold_set_body_font_color(4, (180, 80, 0)),
              "examine_color_rgb",
              "Make the subtitle on slide 5 burnt orange."),
    ], body_focus=True),
    FileTask(D_IMP_25, "title_size_subt6", "edit_title", params=[
        Param(_gold_set_title_font_size(0, 40),
              "examine_font_size",
              "Set the title of slide 1 to 40pt."),
        Param(_gold_set_title_font_size(5, 20),
              "examine_font_size",
              "Set the title of slide 6 to 20pt."),
    ], body_focus=True),

    # D-IMP-26 — hero 3-slide no caption
    FileTask(D_IMP_26, "title_color_h3nc", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (180, 30, 30)),
              "examine_color_rgb",
              "Make the title on slide 1 brick red so the opening hero shot has a strong typographic anchor."),
        Param(_gold_set_title_font_color(2, (30, 130, 30)),
              "examine_color_rgb",
              "I'd like the title on slide 3 in forest green to mirror the foliage tones in the photo behind the title."),
    ]),
    FileTask(D_IMP_26, "title_size_h3nc", "edit_title", params=[
        Param(_gold_set_title_font_size(0, 32),
              "examine_font_size",
              "Set the title of slide 1 to 32pt."),
        Param(_gold_set_title_font_size(1, 18),
              "examine_font_size",
              "Set the title of slide 2 to 18pt."),
    ]),

    # D-IMP-27 — hero 5-slide top
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_27, "title_bold_h5top", "bold_underline_text", params=[
        # Param(_gold_set_bold(1),
              # "examine_font_bold",
              # "Bold the title on slide 2."),
        # Param(_gold_set_bold(3),
              # "examine_font_bold",
              # "Bold the title on slide 4."),
    # ]),
    # Training validation: off-palette bg → hex annotation.
    FileTask(D_IMP_27, "slide_bg_h5top", "change_bg_color", params=[
        Param(_gold_set_background(0, (255, 245, 230)),
              "examine_background_color",
              "Give slide 1 a cream background (Custom Color #FFF5E6)."),
        Param(_gold_set_background(4, (230, 245, 255)),
              "examine_background_color",
              "Give slide 5 a frost-blue background (Custom Color #E6F5FF)."),
    ]),

    # D-IMP-28 — hero 5-slide bottom
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_28, "caption_italic_h5b", "bold_underline_text", params=[
        # Param(_gold_set_italic(0),
              # "examine_font_italic",
              # "Italicize the caption on slide 1."),
        # Param(_gold_set_italic(2),
              # "examine_font_italic",
              # "Italicize the caption on slide 3."),
    # ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_28, "title_align_h5b", "bold_underline_text", params=[
        # Param(_gold_set_text_alignment(1, "CENTER"),
              # "examine_alignment",
              # "Center-align the title on slide 2."),
        # Param(_gold_set_text_alignment(4, "RIGHT"),
              # "examine_alignment",
              # "Right-align the title on slide 5."),
    # ]),

    # D-IMP-29 — hero 8-slide center
    FileTask(D_IMP_29, "title_color_h8", "set_font_color", params=[
        Param(_gold_set_title_font_color(2, (120, 60, 0)),
              "examine_color_rgb",
              "Make the title on slide 3 saddle brown."),
        Param(_gold_set_title_font_color(6, (0, 60, 120)),
              "examine_color_rgb",
              "Make the title on slide 7 steel blue."),
    ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_29, "title_underline_h8", "bold_underline_text", params=[
        # Param(_gold_set_underline(0),
              # "examine_font_underline",
              # "Underline the title on slide 1."),
        # Param(_gold_set_underline(7),
              # "examine_font_underline",
              # "Underline the title on slide 8."),
    # ]),

    # D-IMP-30 — gallery 2x2 4-slide
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_30, "title_bold_g2x2_4", "bold_underline_text", params=[
        # Param(_gold_set_bold(0),
              # "examine_font_bold",
              # "Bold the title on slide 1."),
        # Param(_gold_set_bold(2),
              # "examine_font_bold",
              # "Bold the title on slide 3."),
    # ]),
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_30, "title_color_g2x2_4", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (160, 30, 110)),
              "examine_color_rgb",
              "Make the title of slide 2 magenta (Custom Color #A01E6E)."),
        Param(_gold_set_title_font_color(3, (30, 110, 160)),
              "examine_color_rgb",
              "Make the title of slide 4 cyan-blue (Custom Color #1E6EA0)."),
    ], body_focus=True),

    # D-IMP-31 — gallery 3x3 4-slide
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_31, "title_size_g3x3_4", "edit_title", params=[
        # Param(_gold_set_title_font_size(1, 28),
              # "examine_font_size",
              # "Set the title of slide 2 to 28pt."),
        # Param(_gold_set_title_font_size(3, 20),
              # "examine_font_size",
              # "Set the title of slide 4 to 20pt."),
    # ]),
    FileTask(D_IMP_31, "slide_bg_g3x3_4", "change_bg_color", params=[
        # Validation note: "sand"/"pale blue" don't map to a specific RGB,
        # and the agent's LibreOffice named swatches differ from the gold's
        # custom RGB. Bake the hex into the instruction (mirrors F-CALC-31's
        # "Custom Color #FF6347" pattern).
        Param(_gold_set_background(0, (245, 235, 220)),
              "examine_background_color",
              "Give slide 1 a sand background (Custom Color #F5EBDC)."),
        Param(_gold_set_background(2, (220, 235, 245)),
              "examine_background_color",
              "Give slide 3 a pale blue background (Custom Color #DCEBF5)."),
    ]),

    # D-IMP-32 — gallery 1+3 5-slide
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_32, "title_italic_g1p3_5", "bold_underline_text", params=[
        # Param(_gold_set_italic(0),
              # "examine_font_italic",
              # "Italicize the title on slide 1."),
        # Param(_gold_set_italic(3),
              # "examine_font_italic",
              # "Italicize the title on slide 4."),
    # ]),
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_32, "title_color_g1p3_5", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (210, 90, 0)),
              "examine_color_rgb",
              "Make the title of slide 2 vivid orange (Custom Color #D25A00)."),
        Param(_gold_set_title_font_color(4, (0, 90, 210)),
              "examine_color_rgb",
              "Make the title of slide 5 vivid blue (Custom Color #005AD2)."),
    ], body_focus=True),

    # D-IMP-33 — gallery strip 5-slide
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_33, "title_align_strip5", "bold_underline_text", params=[
        # Param(_gold_set_text_alignment(1, "CENTER"),
              # "examine_alignment",
              # "Center-align the title on slide 2."),
        # Param(_gold_set_text_alignment(4, "RIGHT"),
              # "examine_alignment",
              # "Right-align the title on slide 5."),
    # ]),
    FileTask(D_IMP_33, "swap_slides_strip5", "reorder_slides", params=[
        Param(_gold_swap_slides(0, 3),
              "examine_text",
              "Swap slide 1 and slide 4."),
        Param(_gold_swap_slides(1, 4),
              "examine_text",
              "Swap slide 2 and slide 5."),
    ]),

    # D-IMP-34 — footer basic 6-slide
    FileTask(D_IMP_34, "title_color_fb6", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (160, 0, 90)),
              "examine_color_rgb",
              "Make the title of slide 1 raspberry."),
        Param(_gold_set_title_font_color(5, (0, 90, 160)),
              "examine_color_rgb",
              "Make the title of slide 6 royal blue."),
    ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_34, "body_bold_fb6", "bold_underline_text", params=[
        # Param(_gold_set_body_bold(2),
              # "examine_font_bold",
              # "Bold the body text on slide 3."),
        # Param(_gold_set_body_bold(4),
              # "examine_font_bold",
              # "Bold the body text on slide 5."),
    # ]),

    # D-IMP-35 — footer pagenum 7-slide
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_35, "title_size_fp7", "edit_title", params=[
        # Param(_gold_set_title_font_size(1, 32),
              # "examine_font_size",
              # "Set the title of slide 2 to 32pt."),
        # Param(_gold_set_title_font_size(5, 18),
              # "examine_font_size",
              # "Set the title of slide 6 to 18pt."),
    # ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_35, "title_underline_fp7", "bold_underline_text", params=[
        # Param(_gold_set_underline(0),
              # "examine_font_underline",
              # "Underline the title on slide 1."),
        # Param(_gold_set_underline(6),
              # "examine_font_underline",
              # "Underline the title on slide 7."),
    # ]),

    # D-IMP-36 — footer logo 8-slide
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_36, "slide_bg_fl8", "change_bg_color", params=[
        Param(_gold_set_background(2, (240, 240, 220)),
              "examine_background_color",
              "Give slide 3 a parchment background (Custom Color #F0F0DC)."),
        Param(_gold_set_background(6, (220, 240, 240)),
              "examine_background_color",
              "Give slide 7 a pale cyan background (Custom Color #DCF0F0)."),
    ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_36, "body_align_fl8", "bold_underline_text", params=[
        # Param(_gold_set_body_text_alignment(0, "CENTER"),
              # "examine_alignment",
              # "Center-align the body text on slide 1."),
        # Param(_gold_set_body_text_alignment(7, "RIGHT"),
              # "examine_alignment",
              # "Right-align the body text on slide 8."),
    # ]),

    # D-IMP-37 — notes 7-slide
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_37, "edit_note_n7", "edit_title", params=[
        # Param(_gold_set_speaker_note(2, "Highlight the regulatory implications before moving on."),
              # "examine_note",
              # "Replace the speaker notes on slide 3 with: \"Highlight the regulatory implications before moving on.\""),
        # Param(_gold_set_speaker_note(5, "Keep this slide brief — refer to the appendix for details."),
              # "examine_note",
              # "Replace the speaker notes on slide 6 with: \"Keep this slide brief — refer to the appendix for details.\""),
    # ]),
    FileTask(D_IMP_37, "title_color_n7", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (90, 90, 90)),
              "examine_color_rgb",
              "Make the title of slide 2 medium grey."),
        Param(_gold_set_title_font_color(6, (190, 0, 0)),
              "examine_color_rgb",
              "Make the title of slide 7 vermilion."),
    ]),

    # D-IMP-38 — notes 8-slide
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_38, "edit_note_n8", "edit_title", params=[
        # Param(_gold_set_speaker_note(3, "Confirm the rollout timeline with engineering before this slide."),
              # "examine_note",
              # "Update the speaker notes on slide 4 to: \"Confirm the rollout timeline with engineering before this slide.\""),
        # Param(_gold_set_speaker_note(6, "Wrap up with a clear call to action and contact info."),
              # "examine_note",
              # "Update the speaker notes on slide 7 to: \"Wrap up with a clear call to action and contact info.\""),
    # ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_38, "title_bold_n8", "bold_underline_text", params=[
        # Param(_gold_set_bold(0),
              # "examine_font_bold",
              # "Bold the title on slide 1."),
        # Param(_gold_set_bold(7),
              # "examine_font_bold",
              # "Bold the title on slide 8."),
    # ]),

    # D-IMP-39 — portrait 3-slide
    FileTask(D_IMP_39, "title_color_p3", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (160, 60, 0)),
              "examine_color_rgb",
              "Make the title on slide 1 rust orange."),
        Param(_gold_set_title_font_color(2, (0, 60, 160)),
              "examine_color_rgb",
              "Make the title on slide 3 cobalt."),
    ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_39, "body_bold_p3", "bold_underline_text", params=[
        # Param(_gold_set_body_bold(1),
              # "examine_font_bold",
              # "Bold the body text on slide 2."),
        # Param(_gold_set_body_bold(2),
              # "examine_font_bold",
              # "Bold the body text on slide 3."),
    # ]),

    # D-IMP-40 — portrait 7-slide
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_40, "title_size_p7", "edit_title", params=[
        # Param(_gold_set_title_font_size(2, 36),
              # "examine_font_size",
              # "Set the title of slide 3 to 36pt."),
        # Param(_gold_set_title_font_size(5, 22),
              # "examine_font_size",
              # "Set the title of slide 6 to 22pt."),
    # ]),
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_40, "slide_bg_p7", "change_bg_color", params=[
        Param(_gold_set_background(0, (245, 240, 230)),
              "examine_background_color",
              "Give slide 1 a warm off-white background (Custom Color #F5F0E6)."),
        Param(_gold_set_background(6, (230, 240, 245)),
              "examine_background_color",
              "Give slide 7 a cool off-white background (Custom Color #E6F0F5)."),
    ]),

    # --- D-IMP-41..D-IMP-50 (font + extra layouts) ----------------

    # D-IMP-41 — text serif 5-slide. Headline skill: font-name swap (back to a
    # sans family) — gold mutates only the title font on the named slide.
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_41, "title_font_to_sans", "edit_title", params=[
        # Param(_gold_set_title_font_name(0, "DejaVu Sans"),
              # "examine_font_name",
              # "Change the title font on slide 1 to DejaVu Sans."),
        # Param(_gold_set_title_font_name(2, "DejaVu Sans Mono"),
              # "examine_font_name",
              # "Change the title font on slide 3 to DejaVu Sans Mono."),
    # ]),
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_41, "body_color_serif", "set_font_color", params=[
        Param(_gold_set_body_font_color(1, (60, 60, 60)),
              "examine_color_rgb",
              "Make the body text on slide 2 charcoal grey (Custom Color #3C3C3C)."),
        Param(_gold_set_body_font_color(3, (110, 0, 80)),
              "examine_color_rgb",
              "Make the body text on slide 4 plum (Custom Color #6E0050)."),
    ]),

    # D-IMP-42 — text DejaVu Sans 5-slide.
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_42, "title_bold_dejavu", "bold_underline_text", params=[
        # Param(_gold_set_title_bold(1),
              # "examine_font_bold",
              # "Bold the title on slide 2."),
        # Param(_gold_set_title_bold(4),
              # "examine_font_bold",
              # "Bold the title on slide 5."),
    # ]),
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_42, "title_size_dejavu", "edit_title", params=[
        # Param(_gold_set_title_font_size(0, 36),
              # "examine_font_size",
              # "Set the title of slide 1 to 36pt."),
        # Param(_gold_set_title_font_size(3, 22),
              # "examine_font_size",
              # "Set the title of slide 4 to 22pt."),
    # ]),

    # D-IMP-43 — text DejaVu Serif 4-slide.
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_43, "title_underline_dejavu_serif", "bold_underline_text", params=[
        # Param(_gold_set_title_underline(0),
              # "examine_font_underline",
              # "Underline the title on slide 1."),
        # Param(_gold_set_title_underline(3),
              # "examine_font_underline",
              # "Underline the title on slide 4."),
    # ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_43, "body_align_dejavu_serif", "bold_underline_text", params=[
        # Param(_gold_set_body_text_alignment(1, "JUSTIFY"),
              # "examine_alignment",
              # "Justify the body text on slide 2."),
        # Param(_gold_set_body_text_alignment(2, "CENTER"),
              # "examine_alignment",
              # "Center-align the body text on slide 3."),
    # ]),

    # D-IMP-44 — hero photo 3-slide top-anchored.
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_44, "title_color_h3top", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (160, 0, 90)),
              "examine_color_rgb",
              "Make the title on slide 1 raspberry (Custom Color #A0005A)."),
        Param(_gold_set_title_font_color(2, (0, 90, 160)),
              "examine_color_rgb",
              "Make the title on slide 3 royal blue (Custom Color #005AA0)."),
    ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_44, "caption_italic_h3top", "bold_underline_text", params=[
        # Param(_gold_set_caption_italic(0),
              # "examine_font_italic",
              # "Italicize the caption on slide 1."),
        # Param(_gold_set_caption_italic(2),
              # "examine_font_italic",
              # "Italicize the caption on slide 3."),
    # ]),

    # D-IMP-45 — hero photo 5-slide corner anchored (no caption).
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_45, "title_align_h5corner", "bold_underline_text", params=[
        # Param(_gold_set_title_text_alignment(1, "CENTER"),
              # "examine_alignment",
              # "Center-align the title on slide 2."),
        # Param(_gold_set_title_text_alignment(3, "RIGHT"),
              # "examine_alignment",
              # "Right-align the title on slide 4."),
    # ]),
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_45, "title_size_h5corner", "edit_title", params=[
        # Param(_gold_set_title_font_size(0, 30),
              # "examine_font_size",
              # "Set the title on slide 1 to 30pt."),
        # Param(_gold_set_title_font_size(4, 20),
              # "examine_font_size",
              # "Set the title on slide 5 to 20pt."),
    # ]),

    # D-IMP-46 — hero photo 7-slide center.
    FileTask(D_IMP_46, "swap_slides_h7", "reorder_slides", params=[
        Param(_gold_swap_slides(1, 5),
              "examine_text",
              "Swap slide 2 and slide 6 so the supporting hero photo lands at the start and the cover-style slide closes the section."),
        Param(_gold_swap_slides(2, 6),
              "examine_text",
              "Could you swap slide 3 and slide 7 — I'd like the mid-section hero photo to close the deck and the current last slide to move to position three?"),
    ]),
    FileTask(D_IMP_46, "title_color_h7", "set_font_color", params=[
        Param(_gold_set_title_font_color(2, (200, 100, 0)),
              "examine_color_rgb",
              "Make the title on slide 3 burnt orange so it picks up the warm tones in the autumn hero photo."),
        Param(_gold_set_title_font_color(5, (0, 100, 200)),
              "examine_color_rgb",
              "I'd like the title on slide 6 in cobalt blue so the closing section feels distinct from the warmer slides earlier."),
    ]),

    # D-IMP-47 — gallery 2x2 5-slide.
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_47, "title_italic_g2x2_5", "bold_underline_text", params=[
        # Param(_gold_set_title_italic(0),
              # "examine_font_italic",
              # "Italicize the title on slide 1."),
        # Param(_gold_set_title_italic(3),
              # "examine_font_italic",
              # "Italicize the title on slide 4."),
    # ]),
    # Training validation: off-palette bg → hex annotation (mirrors D-IMP-75/90).
    FileTask(D_IMP_47, "slide_bg_g2x2_5", "change_bg_color", params=[
        Param(_gold_set_background(1, (250, 240, 230)),
              "examine_background_color",
              "Give slide 2 a warm cream background (Custom Color #FAF0E6)."),
        Param(_gold_set_background(4, (230, 240, 250)),
              "examine_background_color",
              "Give slide 5 a cool ice background (Custom Color #E6F0FA)."),
    ]),

    # D-IMP-48 — gallery strip 4-slide.
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_48, "title_bold_strip4", "bold_underline_text", params=[
        # Param(_gold_set_title_bold(0),
              # "examine_font_bold",
              # "Bold the title on slide 1."),
        # Param(_gold_set_title_bold(2),
              # "examine_font_bold",
              # "Bold the title on slide 3."),
    # ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_48, "title_underline_strip4", "bold_underline_text", params=[
        # Param(_gold_set_title_underline(1),
              # "examine_font_underline",
              # "Underline the title on slide 2."),
        # Param(_gold_set_title_underline(3),
              # "examine_font_underline",
              # "Underline the title on slide 4."),
    # ]),

    # D-IMP-49 — footer basic 4-slide.
    FileTask(D_IMP_49, "body_color_fb4", "set_font_color", params=[
        Param(_gold_set_body_font_color(0, (90, 60, 0)),
              "examine_color_rgb",
              "Make the body text on slide 1 dark amber."),
        Param(_gold_set_body_font_color(2, (0, 60, 90)),
              "examine_color_rgb",
              "Make the body text on slide 3 deep teal."),
    ]),
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_49, "title_size_fb4", "edit_title", params=[
        # Param(_gold_set_title_font_size(1, 32),
              # "examine_font_size",
              # "Set the title on slide 2 to 32pt."),
        # Param(_gold_set_title_font_size(3, 20),
              # "examine_font_size",
              # "Set the title on slide 4 to 20pt."),
    # ]),

    # D-IMP-50 — notes 4-slide.
    # Pruned (impress rebalance, eval_class=edit_title OVER):
    # FileTask(D_IMP_50, "edit_note_n4", "edit_title", params=[
        # Param(_gold_set_speaker_note(1, "Cite the August field-trial numbers when challenged on this slide."),
              # "examine_note",
              # "Replace the speaker notes on slide 2 with: \"Cite the August field-trial numbers when challenged on this slide.\""),
        # Param(_gold_set_speaker_note(3, "Yield the floor for questions before introducing the case study."),
              # "examine_note",
              # "Replace the speaker notes on slide 4 with: \"Yield the floor for questions before introducing the case study.\""),
    # ]),
    # Pruned (impress rebalance, eval_class=bold_underline_text OVER):
    # FileTask(D_IMP_50, "title_bold_n4", "bold_underline_text", params=[
        # Param(_gold_set_title_bold(0),
              # "examine_font_bold",
              # "Bold the title on slide 1."),
        # Param(_gold_set_title_bold(2),
              # "examine_font_bold",
              # "Bold the title on slide 3."),
    # ]),

    # --- D-IMP-51..D-IMP-56 (eval-skill gap fillers) --------------
    #
    # These FileTasks set `make_template=` to a specialized factory because
    # their evaluators (insert_table / transition / image_resize / master bg /
    # page_number_color / image_stretch_and_center) don't fit the default
    # `compare_pptx_files` shape used by D-IMP-01..40.

    # D-IMP-51 — insert_table (eval = compare_pptx_files + examine_shape).
    # Param.gold_mutate carries the table-insert snippet; the specialized
    # factory wraps it in a shape-only options bundle so untouched runs don't
    # false-fail.
    FileTask(D_IMP_51, "insert_table_4x3_slide2", "add_slide",
             make_template=_make_table_template, params=[
        Param(_gold_insert_table(1, 4, 3),
              "examine_shape",
              "On slide 2, insert a table with 4 rows and 3 columns."),
        Param(_gold_insert_table(1, 5, 4),
              "examine_shape",
              "On slide 2, insert a table with 5 rows and 4 columns."),
    ]),
    FileTask(D_IMP_51, "insert_table_3x2_slide3", "add_slide",
             make_template=_make_table_template, params=[
        Param(_gold_insert_table(2, 3, 2),
              "examine_shape",
              "On slide 3, insert a table with 3 rows and 2 columns."),
        Param(_gold_insert_table(2, 6, 3),
              "examine_shape",
              "On slide 3, insert a table with 6 rows and 3 columns."),
    ]),

    # D-IMP-52 — slide transition (eval = check_transition, rule-based).
    # Param.extra_examine doubles as the eval rule {slide_idx, transition_type}.
    FileTask(D_IMP_52, "transition_fade_slide2", "add_transition",
             make_template=_make_transition_filetask_template, params=[
        Param(_gold_set_transition(1, "fade"),
              "examine_text",  # unused (custom eval)
              "Apply a fade transition to slide 2.",
              extra_examine={"slide_idx": 1, "transition_type": "fade"}),
        Param(_gold_set_transition(3, "fade"),
              "examine_text",
              "Apply a fade transition to slide 4.",
              extra_examine={"slide_idx": 3, "transition_type": "fade"}),
    ]),
    FileTask(D_IMP_52, "transition_wipe_slide3", "add_transition",
             make_template=_make_transition_filetask_template, params=[
        Param(_gold_set_transition(2, "wipe"),
              "examine_text",
              "Apply a wipe transition to slide 3.",
              extra_examine={"slide_idx": 2, "transition_type": "wipe"}),
        Param(_gold_set_transition(0, "push"),
              "examine_text",
              "Apply a push transition to slide 1.",
              extra_examine={"slide_idx": 0, "transition_type": "push"}),
    ]),

    # D-IMP-53 — image_resize (eval = compare_pptx_files + examine_modify_height).
    # Source picture on each slide is at init_h_cm=8.0. Variants resize the
    # picture on slide K to a clearly-different height so the agent's "shrink
    # the photo" op produces a visible delta.
    FileTask(D_IMP_53, "resize_picture_slide1", "image_stretch",
             make_template=_make_image_resize_filetask_template, params=[
        Param(_gold_resize_picture(0, 4.0),
              "examine_modify_height",
              "Resize the photo on slide 1 to about 4cm tall."),
        Param(_gold_resize_picture(0, 12.0),
              "examine_modify_height",
              "Resize the photo on slide 1 to about 12cm tall — make it noticeably bigger."),
    ]),
    FileTask(D_IMP_53, "resize_picture_slide3", "image_stretch",
             make_template=_make_image_resize_filetask_template, params=[
        Param(_gold_resize_picture(2, 5.0),
              "examine_modify_height",
              "Resize the photo on slide 3 down to about 5cm tall."),
        Param(_gold_resize_picture(2, 11.0),
              "examine_modify_height",
              "Resize the photo on slide 3 to about 11cm tall."),
    ]),

    # D-IMP-54 — master slide bg color (eval =
    # evaluate_presentation_fill_to_rgb_distance, rule-based with original_rgb).
    # Source ships every slide bg = white; gold recolors every slide bg to
    # the target RGB. extra_examine carries the {rgb, original_rgb} rule.
    FileTask(D_IMP_54, "master_bg_blue", "evaluate_presentation_fill_to_rgb_distance",
             make_template=_make_master_bg_filetask_template, params=[
        Param("\n".join(_gold_set_background(i, (0, 0, 255)) for i in range(5)),
              "examine_text",
              "Set the background of every slide to blue.",
              extra_examine={"rgb": [0, 0, 255], "original_rgb": [255, 255, 255]}),
        Param("\n".join(_gold_set_background(i, (255, 255, 0)) for i in range(5)),
              "examine_text",
              "Set the background of every slide to yellow.",
              extra_examine={"rgb": [255, 255, 0], "original_rgb": [255, 255, 255]}),
    ]),
    FileTask(D_IMP_54, "master_bg_green", "evaluate_presentation_fill_to_rgb_distance",
             make_template=_make_master_bg_filetask_template, params=[
        Param("\n".join(_gold_set_background(i, (0, 128, 0)) for i in range(5)),
              "examine_text",
              "Make every slide's background green.",
              extra_examine={"rgb": [0, 128, 0], "original_rgb": [255, 255, 255]}),
        Param("\n".join(_gold_set_background(i, (200, 200, 255)) for i in range(5)),
              "examine_text",
              "Make every slide's background pale lavender.",
              extra_examine={"rgb": [200, 200, 255], "original_rgb": [255, 255, 255]}),
    ]),

    # D-IMP-55 — page_number_color via slide-master patch (eval =
    # check_page_number_colors, rule-based color-name). Param.gold_mutate
    # overloaded: carries the HEX color (no '#') for the master patch.
    # extra_examine carries the {"color": "<name>"} eval rule.
    FileTask(D_IMP_55, "pagenum_red", "check_page_number_colors",
             make_template=_make_pagenum_color_filetask_template, params=[
        Param("FF0000",
              "examine_text",
              "The slide number is hard to read — change its color to red across every slide. "
              "Open Insert > Header & Footer to enable slide numbers, then go to View > Master Slide "
              "and change the slide-number placeholder color to red via Format > Character > Font Color.",
              extra_examine={"color": "red"}),
        Param("0000FF",
              "examine_text",
              "Recolor the page number on every slide to blue. "
              "Open Insert > Header & Footer to enable slide numbers, then go to View > Master Slide "
              "and change the slide-number placeholder color to blue via Format > Character > Font Color.",
              extra_examine={"color": "blue"}),
    ]),
    FileTask(D_IMP_55, "pagenum_green", "check_page_number_colors",
             make_template=_make_pagenum_color_filetask_template, params=[
        Param("00C800",
              "examine_text",
              "Change the slide-number color to green so it stands out on every slide. "
              "Open Insert > Header & Footer to enable slide numbers, then go to View > Master Slide "
              "and change the slide-number placeholder color to green via Format > Character > Font Color.",
              extra_examine={"color": "green"}),
        Param("FF0000",
              "examine_text",
              "Make the page number on every slide red — they're invisible against the deck right now. "
              "Open Insert > Header & Footer to enable slide numbers, then go to View > Master Slide "
              "and change the slide-number placeholder color to red via Format > Character > Font Color.",
              extra_examine={"color": "red"}),
    ]),

    # D-IMP-56 — image stretch and center (eval =
    # check_image_stretch_and_center). Source has small off-center pic on
    # slide 1. Param.gold_mutate carries the ORACLE python (which mutates
    # `out_path` to the stretched layout). Eval matches the image blob across
    # source vs the pre-mutation snapshot at expected_path.
    # Validation PARAM_REDUCIBLE: dropped duplicate paraphrase
    # Param. Original 2 Params shared identical gold+eval and only differed
    # in instruction wording; pixel-tolerance eval was hostile, so keep one
    # Param (skill diversity preserved by sibling stretch tasks).
    FileTask(D_IMP_56, "stretch_to_full_slide", "check_image_stretch_and_center",
             make_template=_make_image_stretch_filetask_template, params=[
        Param(_oracle_stretch_image(0),
              "examine_text",
              "Stretch the image on slide 1 to fill the entire slide while keeping it centered."),
    ]),

    # --- D-IMP-57..D-IMP-67 (instruction-shape validation) --------

    # D-IMP-57 — P2 image add at coord+size (echoes eval `c8...`). Source has
    # a hero photo at init height 6cm. Gold resizes it to the target (w,h) cm.
    # The agent's task is to set the image to that size on the named slide.
    FileTask(D_IMP_57, "add_image_at_size", "image_stretch",
             make_template=_make_image_size_filetask_template, params=[
        Param(_gold_resize_picture_wh(1, 5.0, 5.0),
              "examine_image_size",
              "I am preparing a PPT and want the hero photo on slide 2 resized to 5cm wide by 5cm tall so it sits neatly beside the body copy."),
        Param(_gold_resize_picture_wh(2, 8.0, 4.0),
              "examine_image_size",
              "Could you resize the photo on slide 3 to 8cm wide by 4cm tall so it reads as a banner under the slide title?"),
    ]),

    # D-IMP-58 — P2 image add at position (echoes eval `c8...`). Same shape
    # as D-IMP-57 but exercises different target sizes (corner-anchored).
    FileTask(D_IMP_58, "add_image_position", "image_stretch",
             make_template=_make_image_size_filetask_template, params=[
        Param(_gold_resize_picture_wh(2, 6.0, 6.0),
              "examine_image_size",
              "Resize the photo on slide 3 to 6cm wide by 6cm tall so it tucks into the top-right corner as a square thumbnail."),
        Param(_gold_resize_picture_wh(0, 4.0, 3.0),
              "examine_image_size",
              "I'd like the photo on slide 1 resized to 4cm wide by 3cm tall — it's the bottom-left avatar that introduces the speaker."),
    ]),

    # D-IMP-59 — P3 title-position move (echoes eval `15...`). The gold moves
    # the title textbox to a target (left, top). examine_shape=True so the
    # move is the diff signal; the body textbox stays put on both sides.
    FileTask(D_IMP_59, "title_move_position", "image_stretch",
             make_template=_make_position_filetask_template, params=[
        Param(_gold_move_title(1, 1.0, 12.0),
              "examine_shape",
              "On slide 2, move the title textbox to the bottom of the slide so the supporting body lines lead and the heading closes the slide."),
        Param(_gold_move_title(2, 15.0, 0.5),
              "examine_shape",
              "Could you move the title on slide 3 to the top-right corner so the title sits beside the photo rather than above it?"),
    ]),

    # D-IMP-60 — P3 image-position move (echoes eval `2b...`). Source has a
    # hero photo at the default center; gold moves it to right / left. The
    # agent's task is to reposition the existing image.
    FileTask(D_IMP_60, "image_move_position", "image_stretch",
             make_template=_make_position_filetask_template, params=[
        Param(_gold_move_picture(1, 16.0, 3.0),
              "examine_shape",
              "Move the photo on slide 2 to the right side of the slide so the body text on the left has room to breathe."),
        Param(_gold_move_picture(2, 0.5, 3.0),
              "examine_shape",
              "I'd like the photo on slide 3 moved to the left side of the slide — the caption I'm adding fits better on the right."),
    ]),

    # D-IMP-61 — P4 multi-slide subset color (echoes eval `4ed5abd0` family
    # — title color on slides 1,3,5). Compound 2-arm so the row lands as
    # atom_2 compare_pptx_files+compare_pptx_files (eval signature).
    FileTask(D_IMP_61, "multi_slide_subset_color", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_multi_slide_title_color([0, 2, 4], (200, 30, 30)),
              "examine_color_rgb",
              "Set the title color on slides 1, 3, and 5 to deep red so the three case-study slides share a strong red header."),
        Param(_gold_multi_slide_title_color([1, 3], (30, 30, 200)),
              "examine_color_rgb",
              "Recolor the titles on slides 2 and 4 to deep blue so the two methodology slides stand apart from the result slides."),
    ]),

    # D-IMP-62 — P4 multi-slide subset format (echoes eval `4ed5abd0`). Gold
    # bold+underline titles on slides 2,3,5 OR italic on 1,4. Compound 2-arm.
    FileTask(D_IMP_62, "multi_slide_subset_format", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_multi_slide_title_bold_underline([1, 2, 4]),
              ("examine_font_bold", "examine_font_underline"),
              "Could you bold and underline the titles on slides 2, 3, and 5 so the three milestone slides have heavier-styled headings?"),
        Param(_gold_multi_slide_title_italic([0, 3]),
              "examine_font_italic",
              "I'd like the titles on slides 1 and 4 italicized — those are the framing-question slides and italics signal that to the audience."),
    ]),

    # D-IMP-63 — P5 audio insert (echoes eval `c59742c0`). Pre_config builds
    # the source pptx and ffmpeg-generates a sine-wave mp3 on Desktop.
    # Oracle embeds the mp3 into slide 1 via zip rewrite. Eval =
    # compare_audios (audio_in_slide vs vm_file).
    FileTask(D_IMP_63, "audio_insert", "compare_audios",
             make_template=_make_audio_insert_filetask_template, params=[
        Param("",  # unused — oracle is built from extra_examine
              "examine_text",
              "I am preparing a PPT and want to insert the audio file 'intro_tone.mp3' from the Desktop into slide 1 so it plays when the slideshow starts.",
              extra_examine={"audio_name": "intro_tone.mp3", "freq": 440, "duration": 2}),
        Param("",
              "examine_text",
              "Could you help me embed 'opening_chime.mp3' from the Desktop into the first slide of this Impress deck so the talk opens with the chime?",
              extra_examine={"audio_name": "opening_chime.mp3", "freq": 660, "duration": 2}),
    ]),

    # D-IMP-64 — P6 check_presenter_console_disable. Source ships an xcu with
    # EnablePresenterScreen=true (init=on); oracle rewrites it to false.
    FileTask(D_IMP_64, "presenter_console_disable", "check_presenter_console_disable",
             make_template=_make_xcu_filetask_template, params=[
        Param(_xcu_presenter_disable_item(),
              "examine_text",
              "I have two screens but only want the slideshow on one — disable the LibreOffice Impress presenter console so the second monitor stays free for my notes app.",
              extra_examine={
                  "func": "check_presenter_console_disable",
                  "init_items": (
                      '  <item oor:path="/org.openoffice.Office.Impress/Misc/Start">'
                      '<prop oor:name="EnablePresenterScreen" oor:op="fuse">'
                      '<value>true</value></prop></item>'
                  ),
              }),
        Param(_xcu_presenter_disable_item(),
              "examine_text",
              "Could you turn off the presenter screen feature in LibreOffice Impress so running the slideshow uses only one display rather than both monitors?",
              extra_examine={
                  "func": "check_presenter_console_disable",
                  "init_items": (
                      '  <item oor:path="/org.openoffice.Office.Impress/Misc/Start">'
                      '<prop oor:name="EnablePresenterScreen" oor:op="fuse">'
                      '<value>true</value></prop></item>'
                  ),
              }),
    ]),

    # D-IMP-65 — P6 check_auto_saving_time. Source ships xcu with autosave
    # interval=1 (init), oracle rewrites to target interval (3 or 5).
    FileTask(D_IMP_65, "auto_saving_time", "check_auto_saving_time",
             make_template=_make_xcu_filetask_template, params=[
        Param(_xcu_autosave_item(3),
              "examine_text",
              "Enable auto-save every 3 minutes in LibreOffice so I don't have to keep hitting Ctrl+S while I work on this deck.",
              extra_examine={
                  "func": "check_auto_saving_time",
                  "rules": {"minutes": 3},
                  "init_items": _xcu_autosave_item(1),
              }),
        Param(_xcu_autosave_item(5),
              "examine_text",
              "Configure LibreOffice to auto-save documents every 5 minutes — currently it barely saves and I'm worried about losing work.",
              extra_examine={
                  "func": "check_auto_saving_time",
                  "rules": {"minutes": 5},
                  "init_items": _xcu_autosave_item(1),
              }),
    ]),

    # D-IMP-66 — P6 check_slide_orientation_Portrait. Source ships a landscape
    # deck. Oracle swaps slide_width/height so the saved pptx is portrait.
    FileTask(D_IMP_66, "slide_orientation_portrait", "check_slide_orientation_Portrait",
             make_template=_make_orientation_filetask_template, params=[
        Param("",  # oracle is built into the factory; param body unused
              "examine_text",
              "Please set my slides to portrait orientation instead of landscape — I want them upright so they print well on letter paper."),
        Param("",
              "examine_text",
              "Switch the page setup for this deck from landscape to portrait orientation so the slides are taller than they are wide."),
    ]),

    # D-IMP-67 — P6 left-panel-visible. validation INFEASIBLE_BUG DROP:
    # upstream `check_left_panel` reads the runtime accessibility_tree which
    # is view-mode-dependent and brittle across LO launches. We tried two
    # replacement evals: (a) xdotool window-name probe for 'Slides View' —
    # that string is an a11y document-frame name not an X11 window name,
    # so xdotool can't find it; (b) registrymodifications.xcu file-state
    # probe — the LO panel state isn't reliably persisted to that file
    # (default-launch may not write any explicit entry, or LO rewrites
    # paths between versions). Neither yields a deterministic positive
    # proof of "panel visible right now without trivial-passing on a
    # fresh install". DROP the FileTask rather than ship a flaky row.
    # `_make_left_panel_filetask_template` is preserved (other code may
    # reference it) but no longer wired into FILE_TASKS.
    # FileTask(D_IMP_67, "left_panel", "check_left_panel",
    #          make_template=_make_left_panel_filetask_template, params=[
    #     Param("",
    #           "examine_text",
    #           "I closed the left slide-panel by accident and can't figure out how to bring it back — please restore the Slides View panel on the left."),
    #     Param("",
    #           "examine_text",
    #           "Could you help me re-enable the left slide-sorter panel in Impress so I can see the slide thumbnails while I edit the deck?"),
    # ]),

    # --- validation eval-anchored ADDS (post-cut rebalance) -----------------
    # Fill open slots on Files where main agent pruned over-amped FileTasks.
    # Priorities: add_slide (eval rows 28/36/38/46/22), table ops (eval rows
    # 11/18/33), set_font_color (eval color rows), image_op (rows 7/22/40),
    # reorder_slides, compound_pptx.

    # add_slide: blank-slide appends (eval rows 28/38) -----------------------
    FileTask(D_IMP_24, "add_blank_slide_subt3", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Could you add one more blank slide at the end of the deck so I have room for the closing remarks I'm planning to draft?"),
        Param(_gold_add_blank_slides_at_end(2),
              "examine_number_of_slides",
              "Please add two more blank slides at the end of this deck so I have placeholders for the appendix slides I still need to write."),
    ]),

    FileTask(D_IMP_29, "duplicate_last_slide_h8", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_duplicate_last_slide(),
              "examine_number_of_slides",
              "Could you duplicate the last slide of this deck and append the copy at the end so I can use it as a template for the next section?"),
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Add one new blank slide at the end of this hero-photo deck — I want a placeholder for the closing summary."),
    ]),

    FileTask(D_IMP_31, "add_summary_slide_g3x3", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_summary_slide("Summary"),
              "examine_number_of_slides",
              "Add a summary slide at the end of this gallery deck that lists all the slide titles so the audience has a recap to refer to."),
        Param(_gold_add_summary_slide("Recap"),
              "examine_number_of_slides",
              "Please append a recap slide at the end of the deck — title it 'Recap' and include the headings of every other slide so it acts as a table of contents."),
    ]),

    FileTask(D_IMP_34, "add_blank_slide_fb6", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "I am preparing a PPT and need one extra blank slide at the end of this footer-styled deck so I can sketch a Q&A placeholder."),
        Param(_gold_add_blank_slides_at_end(3),
              "examine_number_of_slides",
              "Add three blank slides at the end of the deck so I have room for the closing exercises section I haven't drafted yet."),
    ]),

    FileTask(D_IMP_37, "duplicate_last_slide_n7", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_duplicate_last_slide(),
              "examine_number_of_slides",
              "Could you duplicate the last slide of this notes deck and append the copy at the end — I want to reuse its layout for the follow-up slide."),
        Param(_gold_add_summary_slide("Summary"),
              "examine_number_of_slides",
              "Add a summary slide titled 'Summary' at the end of this deck that lists every previous slide title so the audience has a recap to refer to."),
    ]),

    FileTask(D_IMP_44, "add_blank_slide_h3top", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Add one blank slide at the end of this top-anchored hero-photo deck so I can sketch a Q&A placeholder before the meeting."),
        Param(_gold_duplicate_last_slide(),
              "examine_number_of_slides",
              "Duplicate the last slide of this deck and place the copy at the end so I can adapt it for the next photo section."),
    ]),

    FileTask(D_IMP_63, "add_blank_slide_audio_target", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Please append one blank slide at the end of this deck — I'm going to drop the audio-narrated closing section onto it."),
        Param(_gold_add_blank_slides_at_end(2),
              "examine_number_of_slides",
              "Could you add two blank slides at the end of the deck so I have placeholders for the two audio-narrated outro sections?"),
    ]),

    FileTask(D_IMP_66, "add_blank_slide_landscape", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Add one more blank slide to the end of this landscape deck — I need a spare for the discussion-question section."),
        Param(_gold_add_summary_slide("Agenda Recap"),
              "examine_number_of_slides",
              "Append a recap slide at the end titled 'Agenda Recap' that lists every previous slide's title so the audience sees a clear summary."),
    ]),

    # table ops (eval rows 11/18/33) — insert / change-row / move-table ------
    # D-IMP-41 has open slot; uses _src_text_deck so adding a table on slide 1
    # works (the slide has title + body but room remains for a small table).
    FileTask(D_IMP_41, "insert_table_features_5x2", "add_slide",
             make_template=_make_table_template, params=[
        Param(_gold_insert_table(0, 5, 2),
              "examine_shape",
              "On the first slide of this serif-styled deck, insert a 5-row, 2-column table beneath the title so we have somewhere to drop the comparison data."),
        Param(_gold_insert_table(2, 4, 3),
              "examine_shape",
              "On slide 3, insert a table with 4 rows and 3 columns so the features-vs-benefits split has somewhere structured to live."),
    ]),

    FileTask(D_IMP_49, "insert_table_footer_4x3", "add_slide",
             make_template=_make_table_template, params=[
        Param(_gold_insert_table(1, 4, 3),
              "examine_shape",
              "On slide 2 of this footer-styled deck, insert a 4-row, 3-column table — it'll hold the cost breakdown I'm about to add."),
        Param(_gold_insert_table_with_headers(0, 5, 4, ["T1", "T2", "T3", "T4"]),
              "examine_shape",
              "On slide 1, insert a 5-row, 4-column table and label the first row cells T1, T2, T3, T4 left-to-right."),
    ]),

    FileTask(D_IMP_65, "insert_table_with_headers", "add_slide",
             make_template=_make_table_template, params=[
        Param(_gold_insert_table_with_headers(0, 4, 4, ["T1", "T2", "T3", "T4"]),
              "examine_shape",
              "On slide 1, insert a 4-row, 4-column table and set the first row to T1, T2, T3, T4 in left-to-right order."),
        Param(_gold_insert_table_at_pos(1, 3, 3, 2.5, 11.0),
              "examine_shape",
              "On slide 2, insert a 3-row, 3-column table near the bottom of the slide (around 11cm from the top) so it sits below the body content."),
    ]),

    # set_font_color (eval has 14 color rows — add 5 more variants) ---------
    FileTask(D_IMP_27, "title_color_h5top", "set_font_color", params=[
        Param(_gold_set_title_font_color(2, (180, 60, 0)),
              "examine_color_rgb",
              "Make the title on slide 3 burnt orange so it reflects the autumn-foliage tones in the top-anchored hero photo."),
        Param(_gold_set_title_font_color(4, (0, 60, 180)),
              "examine_color_rgb",
              "I'd like the title on slide 5 in steel blue so the closing hero-photo slide reads cool relative to the warmer earlier slides."),
    ]),

    # Validation PARAM_REDUCIBLE: dropped the body-color Param.
    # Kept title-only — sibling set_font_color FileTasks are title-only; the
    # body-color variant was a noisy skill outlier.
    FileTask(D_IMP_30, "title_color_g2x2_4_b", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (220, 0, 0)),
              "examine_color_rgb",
              "Make the title of slide 1 a strong deep red so the gallery opener has a bold typographic anchor."),
    ], body_focus=True),

    FileTask(D_IMP_47, "title_color_g2x2_5", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (140, 30, 110)),
              "examine_color_rgb",
              "Could you make the title of slide 1 magenta so the opening gallery has a distinctive section identifier?"),
        Param(_gold_set_title_font_color(3, (30, 110, 140)),
              "examine_color_rgb",
              "Make the title of slide 4 teal blue so it ties to the cool palette used in the closing section."),
    ]),

    FileTask(D_IMP_62, "title_color_subset_format", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (192, 0, 0)),
              "examine_color_rgb",
              "Make the title on slide 2 dark red (Custom Color #C00000 — Office standard 'Dark Red') so the milestone slide stands apart from the framing-question slides on either side."),
        Param(_gold_set_title_font_color(4, (0, 0, 128)),
              "examine_color_rgb",
              "Recolor the title on slide 5 to navy (Custom Color #000080 — Office standard 'Dark Blue') to mirror the navy used on the cover slide of the deck."),
    ]),

    FileTask(D_IMP_67, "title_color_left_panel", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (160, 0, 0)),
              "examine_color_rgb",
              "Make the title on slide 1 dark red so the cover slide has a strong typographic identifier."),
        Param(_gold_set_title_font_color(2, (0, 80, 160)),
              "examine_color_rgb",
              "I'd like the title on slide 3 in royal blue so the closing slide reads cooler than the opening."),
    ]),

    # reorder_slides (3 more variants) ---------------------------------------
    FileTask(D_IMP_32, "swap_slides_g1p3", "reorder_slides", params=[
        Param(_gold_swap_slides(0, 2),
              "examine_text",
              "Swap slide 1 and slide 3 of this 1+3 gallery deck so the standout shot leads and the original cover moves to the middle."),
        Param(_gold_swap_slides(1, 4),
              "examine_text",
              "Swap slide 2 and slide 5 so the closing gallery slide moves up to the start of the picture sequence and the original second slide closes things out."),
    ], body_focus=True),

    FileTask(D_IMP_33, "swap_slides_strip5_b", "reorder_slides", params=[
        Param(_gold_swap_slides(0, 4),
              "examine_text",
              "Swap slide 1 and slide 5 of this strip-gallery deck so the closing strip becomes the new opener and the cover moves to the end as a wrap-up."),
        Param(_gold_swap_slides(2, 3),
              "examine_text",
              "Swap slide 3 and slide 4 of the strip-gallery deck so the middle two photo strips appear in the correct chronological order."),
    ]),

    FileTask(D_IMP_40, "swap_slides_p7", "reorder_slides", params=[
        Param(_gold_swap_slides(1, 5),
              "examine_text",
              "Swap slide 2 and slide 6 of this portrait deck so the supporting detail slide moves to the start and the cover-style slide moves to the closing section."),
        # NB: avoid (0,6) — D_IMP_40 is a 7-slide deck with period-6 title
        # cycle, so slide 0 and slide 6 share the same title (no-op swap →
        # trivial_pass). Use (2,4) instead — both indices are unique.
        Param(_gold_swap_slides(2, 4),
              "examine_text",
              "Swap slide 3 and slide 5 of the portrait deck so the mid-section preview lands earlier and the rooftop sunset slide moves to its more chronological position."),
    ]),

    # image_op (image_move / image_resize variants — eval rows 7/22) --------
    FileTask(D_IMP_56, "image_resize_offcenter", "image_stretch",
             make_template=_make_image_resize_filetask_template, params=[
        Param(_gold_resize_picture(0, 10.0),
              "examine_modify_height",
              "Resize the image on slide 1 of this deck to about 10cm tall so it reads as a larger feature image rather than a small accent."),
        Param(_gold_resize_picture(0, 3.0),
              "examine_modify_height",
              "Shrink the image on slide 1 to about 3cm tall so it tucks into the corner as a small reference thumbnail."),
    ]),

    FileTask(D_IMP_57, "resize_picture_extra", "image_stretch",
             make_template=_make_image_resize_filetask_template, params=[
        Param(_gold_resize_picture(1, 9.0),
              "examine_modify_height",
              "Resize the hero photo on slide 2 to about 9cm tall so it occupies a more dominant share of the slide."),
        Param(_gold_resize_picture(3, 4.0),
              "examine_modify_height",
              "Resize the photo on slide 4 to about 4cm tall — it's a supporting visual and shouldn't compete with the body copy."),
    ]),

    FileTask(D_IMP_58, "image_resize_corner", "image_stretch",
             make_template=_make_image_resize_filetask_template, params=[
        Param(_gold_resize_picture(0, 5.0),
              "examine_modify_height",
              "Shrink the picture on slide 1 to about 5cm tall so it fits neatly in the corner alongside the speaker bio text."),
        Param(_gold_resize_picture(2, 12.0),
              "examine_modify_height",
              "Resize the picture on slide 3 to about 12cm tall so the gallery-style image dominates the slide."),
    ]),

    FileTask(D_IMP_60, "image_move_extra", "image_stretch",
             make_template=_make_position_filetask_template, params=[
        Param(_gold_move_picture(0, 14.0, 5.0),
              "examine_shape",
              "Move the photo on slide 1 to the right side of the slide so the title and any body copy on the left have room to breathe."),
        Param(_gold_move_picture(2, 1.0, 1.0),
              "examine_shape",
              "Move the photo on slide 3 to the top-left corner of the slide so it sits as a small reference image above the body content."),
    ]),

    # compound_pptx (paired-mutation tasks) ---------------------------------
    FileTask(D_IMP_59, "compound_color_and_bold", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_set_title_font_color(0, (200, 30, 30)) + "\n" + _gold_set_title_bold(0),
              ("examine_color_rgb", "examine_font_bold"),
              "On slide 1, make the title red AND bold so it commands attention as the section opener."),
        Param(_gold_set_title_font_color(2, (30, 30, 200)) + "\n" + _gold_set_title_underline(2),
              ("examine_color_rgb", "examine_font_underline"),
              "On slide 3, make the title blue AND underlined so it reads as a clearly marked section divider."),
    ]),

    # Training validation: off-palette compound color → hex annotation.
    FileTask(D_IMP_61, "compound_color_and_align", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_set_body_font_color(1, (60, 60, 60)) + "\n" + _gold_set_body_text_alignment(1, "CENTER"),
              ("examine_color_rgb", "examine_alignment"),
              "On slide 2, make the body text dark grey (Custom Color #3C3C3C) AND center-align it so the supporting copy reads as a balanced quote block."),
        Param(_gold_set_body_font_color(3, (0, 90, 0)) + "\n" + _gold_set_body_text_alignment(3, "RIGHT"),
              ("examine_color_rgb", "examine_alignment"),
              "On slide 4, make the body text forest green (Custom Color #005A00) AND right-align it so the supporting copy mirrors the right-hand visual."),
    ]),

    # Two more set_font_color adds on empty Files (D-IMP-28, D-IMP-50) to
    # close the color skill gap while filling otherwise-unused Files.
    # Training validation: off-palette RGB → hex annotation.
    FileTask(D_IMP_28, "title_color_h5bottom", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (140, 70, 0)),
              "examine_color_rgb",
              "Make the title on slide 1 of this bottom-anchored hero deck saddle brown (Custom Color #8C4600) so it echoes the warm tones in the photos."),
        Param(_gold_set_title_font_color(3, (0, 70, 140)),
              "examine_color_rgb",
              "Make the title on slide 4 deep blue (Custom Color #00468C) so the cool closing slide visually separates from the warmer slides earlier."),
    ]),

    # training validation: off-palette RGBs → hex annotation
    # (mirrors validation D-IMP-09/23/28/41/47/64 pattern).
    FileTask(D_IMP_50, "title_color_n4", "set_font_color", params=[
        Param(_gold_set_title_font_color(1, (160, 0, 80)),
              "examine_color_rgb",
              "Make the title on slide 2 of this notes deck raspberry (Custom Color #A00050) so the speaker can tell at a glance which slide carries the key Q&A."),
        Param(_gold_set_title_font_color(3, (0, 80, 160)),
              "examine_color_rgb",
              "Make the title on slide 4 royal blue (Custom Color #0050A0) so the closing notes slide reads cooler than the warm-toned earlier section."),
    ]),

    # Additional eval-alignment tasks — finish filling remaining 1-slot Files so the distribution
    # is symmetric across the catalog (D-IMP-28, 36, 39, 50, 64 each had one
    # open slot; pair each with an eval-anchored skill).

    FileTask(D_IMP_28, "duplicate_last_h5bottom", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_duplicate_last_slide(),
              "examine_number_of_slides",
              "Could you duplicate the last slide of this bottom-anchored hero-photo deck and append the copy at the end so I can adapt it for the next chapter?"),
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Append one blank slide at the end of this bottom-anchored hero deck so I have a placeholder for the closing slide I still need to draft."),
    ]),

    FileTask(D_IMP_36, "add_blank_slide_fl8", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Add one more blank slide at the end of this footer-with-logo deck so I have a slot for the appendix slide I haven't drafted yet."),
        Param(_gold_add_summary_slide("Summary"),
              "examine_number_of_slides",
              "Please append a summary slide at the end of this deck — title it 'Summary' and list every previous slide's title so the audience has a recap."),
    ]),

    FileTask(D_IMP_39, "add_blank_slide_p3", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Add one blank slide at the end of this portrait-orientation deck so I can sketch the closing summary panel I'm planning."),
        Param(_gold_duplicate_last_slide(),
              "examine_number_of_slides",
              "Duplicate the last slide of this portrait deck and append the copy at the end — I want to reuse its layout for the next section."),
    ]),

    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_64, "title_color_presenter", "set_font_color", params=[
        Param(_gold_set_title_font_color(0, (170, 30, 30)),
              "examine_color_rgb",
              "Make the title on slide 1 brick red (Custom Color #AA1E1E) so the cover slide of this presenter-target deck has a strong opening statement."),
        Param(_gold_set_title_font_color(2, (30, 30, 170)),
              "examine_color_rgb",
              "I'd like the title on slide 3 in deep blue (Custom Color #1E1EAA) so the closing slide reads with a clear, calm tone."),
    ]),

    FileTask(D_IMP_50, "add_blank_slide_n4", "add_slide",
             make_template=_make_add_slide_template, params=[
        Param(_gold_add_blank_slide_at_end(),
              "examine_number_of_slides",
              "Add one blank slide at the end of this notes-styled deck so I have a placeholder for the closing Q&A slide I still need to draft."),
        Param(_gold_duplicate_last_slide(),
              "examine_number_of_slides",
              "Duplicate the last slide of this notes deck and append the copy at the end so I can adapt the layout for the next section."),
    ]),

    # Save-As / Export coverage (libreoffice_impress Axis E). Two
    # templates mirror upstream eval pattern `a097acff` (Save-As .pptx) —
    # agent edits and persists the result via File > Save As to a NEW
    # filename on the Desktop. Custom factory `_make_save_as_filetask_template`
    # routes the eval to the typed Save-As path, oracle plants gold there.
    # Save-As / Export coverage — pure Save-As (no content edit). Gold = source
    # (the gold mutator is a no-op `pass`); the agent's only task is File >
    # Save As to the typed filename on the Desktop, keeping .pptx format.
    # Eval reads the new path via `compare_pptx_files`. examine_number_of_slides
    # gives the diff signal (a missing-file or wrong-content sink will trip the
    # slide-count check). Instructions deliberately avoid title/bold/font/etc.
    # keywords so the op_family classifier routes them to `save_or_export`.
    FileTask(D_IMP_68, "save_as_pptx_filename", "edit_title",
             make_template=_make_save_as_filetask_template, params=[
        Param("pass",
              "examine_number_of_slides",
              "Save a copy of this presentation as `presentation_final.pptx` on the Desktop using File > Save As — keep the .pptx format when the dialog asks.",
              extra_examine={"save_as_name": "presentation_final.pptx"}),
        Param("pass",
              "examine_number_of_slides",
              "Use File > Save As to save the open deck as `report_v2.pptx` on the Desktop; keep .pptx format in the format-confirmation dialog.",
              extra_examine={"save_as_name": "report_v2.pptx"}),
    ]),
    FileTask(D_IMP_69, "save_as_pptx_export", "edit_title",
             make_template=_make_save_as_filetask_template, params=[
        Param("pass",
              "examine_number_of_slides",
              "Export this deck as `deck_archive.pptx` on the Desktop via File > Save As, keeping the .pptx format in the format prompt that appears.",
              extra_examine={"save_as_name": "deck_archive.pptx"}),
        Param("pass",
              "examine_number_of_slides",
              "I need to keep an extra copy — Save As `backup_copy.pptx` on the Desktop, .pptx format, so the original file stays untouched.",
              extra_examine={"save_as_name": "backup_copy.pptx"}),
    ]),

    # Doc-wide framing coverage (libreoffice_impress
    # `slide_anchor.doc_wide` bridge): instructions reference "all slides" /
    # "throughout the deck" rather than "on slide N". Gold helpers apply
    # the mutation deck-wide via a `for _sl in prs.slides:` loop.
    # Training validation: off-palette → switch to Office swatches + hex.
    FileTask(D_IMP_70, "doc_wide_title_color", "set_font_color", params=[
        Param(_gold_all_slides_title_color((0, 32, 96)),
              "examine_color_rgb",
              "Recolor every title throughout this deck to deep navy blue (Custom Color #002060 — Office 'Dark Blue') so the section headings share a single consistent brand color across the whole presentation."),
        Param(_gold_all_slides_title_color((192, 0, 0)),
              "examine_color_rgb",
              "Apply a brick-red color (Custom Color #C00000 — Office 'Dark Red') to all of the slide titles across the entire deck — I want a unified header treatment from cover to close."),
    ]),
    FileTask(D_IMP_71, "doc_wide_title_bold", "bold_underline_text", params=[
        Param(_gold_all_slides_title_bold(),
              "examine_font_bold",
              "Bold every title across the deck so each section heading reads with the same weight from start to finish."),
        Param(_gold_all_slides_title_underline(),
              "examine_font_underline",
              "Underline all of the slide titles throughout this presentation for a uniform section-divider look across the whole deck."),
    ]),
    # Training validation: doc-wide off-palette bg → hex annotation
    # (mirrors D-IMP-75 validation pattern; missed in earlier sweep).
    FileTask(D_IMP_72, "doc_wide_bg", "change_bg_color", params=[
        Param(_gold_all_slides_background((245, 245, 230)),
              "examine_background_color",
              "Give every slide in this deck a cream-paper background (Custom Color #F5F5E6) — I want the whole presentation to feel warm and editorial rather than stark white."),
        Param(_gold_all_slides_background((230, 240, 250)),
              "examine_background_color",
              "Apply a pale-blue background fill (Custom Color #E6F0FA) to all slides across the deck so the whole presentation reads with a calm, cool wash from cover to close."),
    ]),
    FileTask(D_IMP_73, "doc_wide_title_italic", "bold_underline_text", params=[
        Param(_gold_multi_slide_title_italic([0, 1, 2]),
              "examine_font_italic",
              "Italicize every title throughout this 3-slide deck so the headings carry a consistent editorial tone across the whole presentation."),
        Param(_gold_all_slides_title_bold(),
              "examine_font_bold",
              "Bold all of the slide titles across the entire deck — the whole presentation should share a heavy header treatment from start to end."),
    ]),
    # Training validation: off-palette doc-wide title color → hex annotation.
    FileTask(D_IMP_74, "doc_wide_title_color_b", "set_font_color", params=[
        Param(_gold_all_slides_title_color((200, 60, 0)),
              "examine_color_rgb",
              "Apply a warm orange color (Custom Color #C83C00) to every title throughout this deck so the section headings share one unified brand accent across all slides."),
        Param(_gold_all_slides_title_color((60, 60, 60)),
              "examine_color_rgb",
              "Recolor all of the slide titles across the deck to a neutral charcoal grey (Custom Color #3C3C3C) so the headings feel restrained and editorial throughout."),
    ]),
    FileTask(D_IMP_75, "doc_wide_bg_b", "change_bg_color", params=[
        # training validation: target RGBs not in LO standard palette
        # → agent forced into Custom Color Hex field, mistyped or rounded.
        # Added explicit hex annotation mirroring D-IMP-90 validation pattern.
        Param(_gold_all_slides_background((250, 248, 240)),
              "examine_background_color",
              "Set the slide background fill across the entire deck to a soft ivory (Custom Color #FAF8F0) so every slide reads as part of one warm-paper presentation."),
        Param(_gold_all_slides_background((232, 244, 232)),
              "examine_background_color",
              "Apply a light mint background (Custom Color #E8F4E8) to all slides in this presentation so the whole deck reads with a fresh, unified palette from start to finish."),
    ]),
    FileTask(D_IMP_76, "doc_wide_title_bold_to", "bold_underline_text", params=[
        Param(_gold_all_slides_title_bold(),
              "examine_font_bold",
              "Bold every slide title throughout this title-only deck so each banner heading reads with the same weight from the cover to the closing."),
        Param(_gold_all_slides_title_underline(),
              "examine_font_underline",
              "Apply underline formatting to every slide title across this banner-style deck."),
    ]),
    # Training validation: off-palette → hex annotation.
    FileTask(D_IMP_77, "doc_wide_title_color_sub", "set_font_color", params=[
        Param(_gold_all_slides_title_color((30, 100, 30)),
              "examine_color_rgb",
              "Recolor every title across this subtitle-style deck to forest green (Custom Color #1E641E) so all the section headings share one consistent brand color throughout."),
        Param(_gold_all_slides_title_color((100, 30, 100)),
              "examine_color_rgb",
              "Apply a deep plum color (Custom Color #641E64) to all of the slide titles across the entire deck for a unified, editorial header treatment from start to finish."),
    ]),
    # training validation: off-palette doc-wide bg → hex annotation
    # (mirrors validation D-IMP-75 + validation D-IMP-72 pattern).
    FileTask(D_IMP_78, "doc_wide_bg_c", "change_bg_color", params=[
        Param(_gold_all_slides_background((248, 240, 232)),
              "examine_background_color",
              "Apply a warm sand background fill (Custom Color #F8F0E8) to every slide of this deck so the whole presentation reads with one calm, earth-tone palette throughout."),
        Param(_gold_all_slides_background((240, 240, 248)),
              "examine_background_color",
              "Set a pale lavender background (Custom Color #F0F0F8) across all slides of this presentation so the entire deck shares one soft, unified colour from cover to close."),
    ]),

    # -----------------------------------------------------------------------
    # Cycle-l7 ADDS — compound `compare_pptx_files`×2 (atom_2 closes 0%→target)
    # + PNG-export + build-deck-from-scratch. Every entry below cites the
    # emulated osworld_libreoffice_impress_XXXX eval task in its comment.
    # -----------------------------------------------------------------------

    # D-IMP-79 / D-IMP-80 — emulates osworld_libreoffice_impress_04578141
    # ("3 textboxes on slide 1 yellow/red/green top-to-bottom"). The target
    # slide has THREE textboxes (built by `_src_three_textbox_deck`); gold
    # colors them top→bottom in 3 distinct RGB values. Compound 2-arm eval.
    FileTask(D_IMP_79, "three_textbox_colors_yrg", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_three_textbox_colors_slide(0, (255, 255, 0), (255, 0, 0), (0, 255, 0)),
              "examine_color_rgb",
              "Change the text color in the textboxes on slide 1 to yellow, red, and green respectively, in top-to-bottom order. Use exactly these colors — no variations."),
        Param(_gold_three_textbox_colors_slide(0, (0, 0, 255), (255, 165, 0), (128, 0, 128)),
              "examine_color_rgb",
              "On slide 1 set the three textbox text colors top-to-bottom: the top one blue, the middle one orange, and the bottom one purple — exact colors, no shades."),
    ]),
    FileTask(D_IMP_80, "three_textbox_colors_rbg", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_three_textbox_colors_slide(0, (255, 0, 0), (0, 0, 255), (0, 128, 0)),
              "examine_color_rgb",
              "On slide 1 of this deck, set the textbox text colors to red, blue, and green respectively from top to bottom. Use the exact named colors."),
        Param(_gold_three_textbox_colors_slide(0, (0, 0, 0), (128, 128, 128), (255, 255, 255)),
              "examine_color_rgb",
              "On slide 1 set the text color (font/character color, NOT the box fill) of the three textboxes to black, grey, and white in top-to-bottom order — three distinct values, no in-between shades."),
    ]),

    # D-IMP-81 — emulates osworld_libreoffice_impress_05dd4c1d
    # ("Align first textbox on slide 3 right, slide 4 center, slide 5 left").
    # 3-slide independent alignment assertions — quintessential compound row.
    FileTask(D_IMP_81, "per_slide_alignment_rcl", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_per_slide_alignment([(2, "RIGHT"), (3, "CENTER"), (4, "LEFT")]),
              "examine_alignment",
              "Align the first textbox on slide 3 to the right, on slide 4 to the center, and on slide 5 to the left. Ensure the alignment is applied correctly on each slide."),
        Param(_gold_per_slide_alignment([(0, "CENTER"), (1, "RIGHT"), (2, "JUSTIFY")]),
              "examine_alignment",
              "Set the first textbox alignment per slide: slide 1 centered, slide 2 right-aligned, slide 3 justified. Each slide gets its own independent alignment."),
    ]),

    # D-IMP-82 — emulates osworld_libreoffice_impress_08aced46
    # ("Give the slide 2 the right-aligned title 'Note'"). Compound = TEXT
    # (examine_text) AND ALIGNMENT (examine_alignment) on the same shape.
    FileTask(D_IMP_82, "title_text_and_align", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_title_text_and_align(1, "Note", "RIGHT"),
              ("examine_text", "examine_alignment"),
              "Give slide 2 the right-aligned title 'Note'. The title text and its alignment both need to be applied together."),
        Param(_gold_title_text_and_align(0, "Summary", "CENTER"),
              ("examine_text", "examine_alignment"),
              "On slide 1, set the title text to 'Summary' and center-align it. Both the text and the alignment need to be in place."),
    ]),

    # D-IMP-83 — emulates osworld_libreoffice_impress_4ed5abd0
    # ("Set color of titles in slides 2,3,5 as black and underline them").
    # Compound = COLOR_RGB AND UNDERLINE on a non-contiguous slide subset.
    # Training validation: navy blue (30,60,160) off-palette → switch to
    # Office "Dark Blue" (0,32,96) + explicit hex annotation.
    FileTask(D_IMP_83, "multi_title_color_and_underline", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_multi_slide_title_color_and_underline([1, 2, 4], (0, 0, 0)),
              ("examine_color_rgb", "examine_font_underline"),
              "Set the color of titles in slides 2, 3, and 5 as black and underline them. Apply both formatting changes to each of those three slides."),
        Param(_gold_multi_slide_title_color_and_underline([0, 3], (0, 32, 96)),
              ("examine_color_rgb", "examine_font_underline"),
              "Make the titles on slides 1 and 4 navy blue (Custom Color #002060 — Office standard 'Dark Blue') and underline them — both changes on both slides so the section headers share a uniform navy-underlined style."),
    ]),

    # D-IMP-84 — emulates osworld_libreoffice_impress_550ce7e7
    # ("Strike-through completed to-do list lines"). Compound: per-paragraph
    # strikethrough on the body textbox (independent of other lines).
    FileTask(D_IMP_84, "strikethrough_first_two_lines", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_strikethrough_lines(0, [0, 1]),
              "examine_strike_through",
              "I am checking the to-do list for last week and adding strike-through to the lines we already accomplished. Could you add a strike-through on the first and second lines of the body on slide 1?"),
        Param(_gold_strikethrough_lines(1, [0, 2]),
              "examine_strike_through",
              "On slide 2, add a strike-through to the first and third lines of the body content so the completed items in the list are crossed out."),
    ]),

    # D-IMP-85 — emulates osworld_libreoffice_impress_455d3c66
    # ("Export the Impress file to .png as res.png on the Desktop"). Eval =
    # compare_images (SSIM).
    FileTask(D_IMP_85, "png_export_res", "compare_images",
             make_template=_make_png_export_template, params=[
        Param("",  # unused — oracle/eval don't read gold_mutate
              "examine_text",
              "Could you help me export this Impress file to a .png image and save it as `res.png` on the Desktop? Default export settings are fine."),
        Param("",
              "examine_text",
              "Please export the open presentation to a PNG image on the Desktop named `res.png` — keep the default export options when the dialog appears."),
    ]),

    # D-IMP-86 — emulates osworld_libreoffice_impress_bf4e9888
    # ("6 png images on Desktop → create a new presentation with 6 blank slides,
    # one image per slide in numerical order"). Gold = K-slide image-only deck.
    FileTask(D_IMP_86, "build_six_image_deck", "add_slide",
             make_template=_make_build_deck_template, params=[
        # Training Validation note: 6-image variant DROPPED — agent needs
        # ~6-8 turns per image (new slide → Insert Image → dialog → resize),
        # 6 images blows past the 30-turn budget (4/4 truncated in S1+S2).
        # Keeping 4-image variant only (still exercises the add_slide skill).
        Param("",
              "examine_text",
              "I have pic1.png through pic4.png on the Desktop. Please build a brand-new four-slide deck where each slide is blank except for one of those pictures, inserted in numerical order — pic1 on slide 1, pic2 on slide 2, and so on.",
              extra_examine={"n_images": 4}),
    ]),

    # D-IMP-87 — extra compound coverage: title bold + body color on the same
    # slide (echoes the broader `compound_pptx` family). Two independent
    # canonical-field checks via the 2-arm factory.
    FileTask(D_IMP_87, "compound_title_bold_body_color", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_set_title_bold(1) + "\n" + _gold_set_body_font_color(1, (30, 30, 200)),
              ("examine_font_bold", "examine_color_rgb"),
              "On slide 2, bold the title text and recolor the body paragraph to deep blue — apply both formatting changes together."),
        Param(_gold_set_title_underline(2) + "\n" + _gold_set_body_font_color(2, (30, 130, 30)),
              ("examine_font_underline", "examine_color_rgb"),
              "On slide 3, underline the title and color the body text forest green — both edits need to land on the same slide."),
    ]),

    # -----------------------------------------------------------------------
    # validation ADDS — closes (a) atom_3plus 0%→target via 4-arm
    # `compare_pptx_files` compound (emulates f23acfd2) and (b) target_scope
    # .deck_wide_or_implicit -32pp gap via 6 NAKED-voice templates. The
    # `deck_wide_or_implicit` bucket fires when the instruction has NO
    # ordinal anchor AND NO explicit deck-wide phrase ("every slide" / "all
    # slides" / "across the deck" / "throughout" / "whole presentation" /
    # "entire deck" / "each slide") — see measure_gap.impress_target_scope.
    # These templates carefully avoid all those phrases yet gold iterates
    # the full deck, so the natural-language reading matches the action.
    # -----------------------------------------------------------------------

    # D-IMP-88 — emulates osworld_libreoffice_impress_f23acfd2
    # ("Add a bullet point to the content of this slide"; 4-arm OR over
    # compare_pptx_files with two alt golds × two font-tolerance levels).
    # Single-slide source so "this slide" is unambiguous; gold appends one
    # paragraph to the body textbox on slide 0 (eval-row 46 pattern).
    FileTask(D_IMP_88, "add_bullet_implicit", "compound_pptx",
             make_template=_make_compound_pptx_4arm_template, params=[
        # emulates osworld_libreoffice_impress_f23acfd2
        Param(_gold_append_to_body(0, "Follow up with the team next week"),
              "examine_bullets",
              "Add a bullet point reading 'Follow up with the team next week' to the content of this slide."),
        # emulates osworld_libreoffice_impress_f23acfd2
        Param(_gold_append_to_body(0, "Open question: budget for Q4"),
              "examine_bullets",
              "I need one more bullet under the body content here — append the follow-up line 'Open question: budget for Q4' so the list closes with an action item."),
    ]),

    # D-IMP-89 — implicit-voice title bold. Naked "Bold the slide titles"
    # / "Make the section titles bold" — no per-slide anchor, no explicit
    # "every slide" phrase. Gold iterates every slide title.
    # deck_wide_implicit: matches eval phrasing pattern (no specific eval task)
    FileTask(D_IMP_89, "implicit_title_bold", "bold_underline_text", params=[
        Param(_gold_all_slides_title_bold(),
              "examine_font_bold",
              "Bold the slide titles so the section headers stand out from the body copy when this deck is projected."),
        Param(_gold_all_slides_title_underline(),
              "examine_font_underline",
              "Underline the section titles to give the headers a clearer visual divider against the body text."),
    ]),

    # D-IMP-90 — implicit-voice background recolor. Naked "Change the
    # background to <color>" — no anchor, no explicit deck-wide phrase.
    # deck_wide_implicit: matches eval phrasing pattern (no specific eval task)
    FileTask(D_IMP_90, "implicit_bg_recolor", "change_bg_color", params=[
        # training validation: previous RGBs (230,240,250)/(245,245,230)
        # not in LO standard palette → agent forced into Custom Color picker
        # and consistently mis-picked. Switched to palette-matched values +
        # explicit hex annotation so Custom Color Hex field entry is reliable.
        Param(_gold_all_slides_background((218, 238, 243)),
              "examine_background_color",
              "Change the slide background to pale blue (Custom Color #DAEEF3) so the slides feel cooler under a projector."),
        Param(_gold_all_slides_background((253, 233, 217)),
              "examine_background_color",
              "Set the slide background to warm peach (Custom Color #FDE9D9) — the current stark white looks too clinical for the audience."),
    ]),

    # D-IMP-91 — implicit-voice title font-family change. Naked "Change the
    # title font to <name>" — single naked imperative, gold rewrites every
    # title font-name. New helper `_gold_all_slides_title_font_name`.
    # deck_wide_implicit: matches eval phrasing pattern (no specific eval task)
    FileTask(D_IMP_91, "implicit_title_font_name", "edit_title", params=[
        Param(_gold_all_slides_title_font_name("Liberation Serif"),
              "examine_font_name",
              "Change the title font to Liberation Serif so the headings have a more editorial, book-like feel than the default sans face."),
        # Validation fix: DejaVu Sans not installed → swap to Carlito (a Calibri
        # metric-compatible sans face, installed via crosextra-carlito).
        Param(_gold_all_slides_title_font_name("Carlito"),
              "examine_font_name",
              "Switch the title typeface to Carlito for a cleaner, modern look on the section headers."),
    ]),

    # D-IMP-92 — implicit-voice title alignment. Naked "Center the titles"
    # — no per-slide anchor, no "every slide" phrase. Gold sets paragraph
    # alignment on every slide's title.
    # deck_wide_implicit: matches eval phrasing pattern (no specific eval task)
    FileTask(D_IMP_92, "implicit_title_align", "bold_underline_text", params=[
        Param(_gold_all_slides_title_alignment("CENTER"),
              "examine_alignment",
              "Center the slide titles so the section headings sit on the visual midline of each page."),
        Param(_gold_all_slides_title_alignment("RIGHT"),
              "examine_alignment",
              "Right-align the titles to match the rest of the brand template's flush-right header treatment."),
    ]),

    # D-IMP-93 — implicit-voice bullet append. Naked "Add a follow-up
    # bullet to the body" — gold appends one paragraph to body on every
    # slide (echoes the f23acfd2 pattern but in multi-slide implicit form).
    # deck_wide_implicit: matches eval phrasing pattern (no specific eval task)
    FileTask(D_IMP_93, "implicit_body_bullet", "compound_pptx",
             make_template=_make_compound_pptx_template, params=[
        Param(_gold_all_slides_append_body_bullet("Action item: review next week"),
              "examine_bullets",
              "Add a follow-up bullet to the body content — I want the action-item line 'Action item: review next week' under the existing copy so the audience knows the next step."),
        Param(_gold_all_slides_append_body_bullet("Open question: confirm budget"),
              "examine_bullets",
              "Append the open-question bullet 'Open question: confirm budget' under the body text so the slides surface the unresolved budget point."),
    ]),

    # D-IMP-94 — implicit-voice title color. Naked "Make the titles
    # <color>" — no anchor, no explicit deck-wide phrase. Gold sets every
    # title text color.
    # deck_wide_implicit: matches eval phrasing pattern (no specific eval task)
    FileTask(D_IMP_94, "implicit_title_color", "set_font_color", params=[
        # training validation: previous RGBs (30,60,160)/(140,30,30)
        # not in LO standard palette → agent picked palette nearest swatch
        # (pure red / pure blue) which didn't match exact RGB. Switched to
        # palette-matched values + explicit Hex annotation.
        Param(_gold_all_slides_title_color((0, 32, 96)),
              "examine_color_rgb",
              "Make the titles dark blue (Custom Color #002060 — Office standard 'Dark Blue') — the section headers need a consistent corporate accent instead of plain black."),
        Param(_gold_all_slides_title_color((192, 0, 0)),
              "examine_color_rgb",
              "Recolor the section titles to dark red (Custom Color #C00000 — Office standard 'Dark Red') so the headings carry the campaign palette through the slides."),
    ]),
]


# §I.h — Emission.
TEMPLATES.extend(_emit_templates(FILE_TASKS))
