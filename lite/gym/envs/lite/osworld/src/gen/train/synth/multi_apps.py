"""Multi-apps synth generator (Track A — host-heredoc design).

Per AGENTS.md / multi_apps.md: multi_apps rows compose 2-app pipelines
(txt/csv/xlsx-source → docx/xlsx/csv-sink, or shell-pipeline / archive ops).
Source state lives entirely in `pre_config_steps` heredocs that materialize
the input file(s) AND the gold expected file. Oracle = `cp` gold over the
agent's sink (or shell-pipeline equivalent); evaluator matches the per-app
convention. LO-sink rows set `oracle_after_postconfig=True` + use
LO_SAVE_POSTCONFIG. Conservative refactor (2026-05-10):
only baseline-tool rows; everything needing ffmpeg / pdftk / Calibre stays
deferred (see AGENTS.md §Tool-dependency contract). 2026-05-10 follow-up:
pandoc + ImageMagick + poppler-utils (`pdftotext`/`pdfinfo`/`pdftoppm`) are
now installed in the Dockerfile, unblocking Cat J / K / L / M templates
below.

Implemented row groups (83/83 PASS on baseline-tool refactor — b5jvgp0rt sweep
2026-05-10; +10 real-asset rows + +15 Cat-I cross-app rows + +21 Cat-I third-pass
cross-app rows added post-sweep, pending re-validation; 2026-05-10 photo-template
Cartesian refactor: 19 hardcoded photo templates deleted, 20 topic-Cartesian
photo templates added (5 topics × 4 skills); total 75 templates):
- txt → docx via writer (#32-class)                  → compare_docx_files
                                                       (`txt_to_docx_*` factories)
- csv → xlsx via calc (#34, #75, #311, #324)         → compare_table sheet_data
                                                       (`csv_to_xlsx_*` factories)
- xlsx → csv export via calc (#94, #146-#149)        → compare_csv
                                                       (`xlsx_to_csv_*` factories)
- xlsx → docx-table via writer (#18, #19, #20, #21)  → compare_docx_tables
                                                       (`xlsx_to_docx_table_*` factories)
- Shell pipeline on inline text:
  - grep (#39)                                       → compare_text_file   (`shell_grep_holmes`)
  - awk filter (#41)                                 → compare_csv         (`shell_awk_filter_rate`)
  - csv concat (#50)                                 → compare_csv         (`shell_concat_csv`)
  - awk skip header (#61)                            → compare_csv         (`shell_awk_skip_header`)
  - grep -c errors (#103)                            → compare_text_file   (`shell_grep_count_errors`)
  - head excerpt (#91)                               → compare_text_file   (`shell_head_excerpt`)
  - sort -f (#79)                                    → compare_text_file   (`shell_sort_case_insensitive`)
  - sort + uniq (#69)                                → compare_text_file   (`shell_sort_uniq`)
  - cut -d',' -f3                                    → compare_text_file   (`shell_cut_third_col`)
  - tr [:upper:][:lower:]                            → compare_text_file   (`shell_tr_lower`)
  - wc -l                                            → compare_text_file   (`shell_wc_lines`)
  - sed substitute                                   → compare_text_file   (`shell_sed_substitute`)
- Extract zip + sort by extension (#30)              → check_include_exclude  (`extract_sort_by_ext`)
- Extract tar.gz (#116)                              → check_include_exclude  (`extract_targz`)
- zip -r directory (#86, #108)                       → compare_archive (file_type=text)
                                                                          (`zip_create_directory`)
- Real-asset shell-pipeline rows (#17, #18, #20, #50, #69, #79, #91, #115, #128)
  staged via `_stage_asset` against `data/assets/synth/` bundle:
  - du+sort+tail top-3 gutenberg (#17)               → compare_text_file (`asset_du_top3_gutenberg`)
  - awk col 1 of us-population-states (#18)          → compare_text_file (`asset_awk_state_names`)
  - sed delete ALGERNON lines (#20)                  → compare_text_file (`asset_sed_drop_algernon`)
  - cat 4 FRED CSVs (#50)                            → compare_text_file (`asset_concat_fred_csvs`)
  - sort+uniq art-of-war words (#69)                 → compare_text_file (`asset_sort_uniq_artofwar`)
  - sort -f earnest lines (#79)                      → compare_text_file (`asset_sort_case_earnest`)
  - shuf -n 10 oil-wti (#91)                         → compare_text_file (`asset_shuf_oil_wti`)
  - head -n 51 oil-wti (#115)                        → compare_text_file (`asset_head_oil_wti`)
  - filter G7 from world-gdp (#128)                  → compare_text_file (`asset_grep_g7_world_gdp`)
- Real-asset csv → xlsx (#34)                        → compare_table  (`asset_csv_to_xlsx_state_income`)
- Cross-app real-asset rows (Cat I — added 2026-05-10, +15 rows pending sweep):
  - CSV → docx-table (#18-#21, #44, #154-#158):
    - us-state-median-income.csv → 3-col table (`multi_csv_to_docx_table_us_state_income`)
    - world-gdp-2022.csv top-10 → 4-col table (`multi_csv_to_docx_table_world_gdp_top10`)
    - us-population-states.csv top-5 → 3-col table (`multi_csv_to_docx_table_us_population_top5`)
    - us-unemployment.csv last-12 → 2-col table (`multi_csv_to_docx_table_unemployment_recent12m`)
  - Photo cross-app (Cartesian: N topic families × 4 photo skill-templates).
    Per-seed photo selection rotates within `TOPIC_FAMILIES[topic]` so same
    template + different seed = different photo, but multi-slide decks stay
    coherent. Gallery skill auto-skips topics with pool < 4.
    - photo→docx insert-after-caption (`multi_topic_<topic>_photo_to_docx`)
      — `compare_docx_strict examine_images=True`
    - photo→docx cover-image (`multi_topic_<topic>_photo_to_docx_cover`)
      — `compare_docx_strict examine_images=True`
    - photo→pptx single-slide (`multi_topic_<topic>_photo_to_pptx`)
      — `compare_pptx_files examine_image_size=True`
    - photo gallery 4-slide (`multi_topic_<topic>_photo_gallery_4`)
      — `compare_pptx_files examine_image_size=True`
    Topics: space / astronomy / wildlife / food / architecture / portrait /
    city / office / event / product / classroom / medical (see TOPIC_FAMILIES).
  - Photo→docx with `compare_docx_images` PIL-pixel-eq func (deferred-bridge
    2026-05-10, K=1 in eval, 2 templates added):
    - single-photo insert (`multi_photo_docx_images_eval_astronomy`)
    - dual-photo side-by-side (`multi_photo_docx_images_eval_dual_astronomy`)
    Bridges plan rows #49, #139, #141 (residual). Distinct from the
    `compare_docx_strict examine_images=True` byte-hash family above
    because PIL pixel equality tolerates LO image re-encoding.
  - CSV → xlsx with chart (compare_table chart+sheet_data):
    - us-unemployment.csv → LineChart (`multi_csv_to_xlsx_chart_unemployment_line`)
    - world-gdp-2022.csv top-10 → BarChart (`multi_csv_to_xlsx_chart_world_gdp_bar`)
  - Multi-CSV concat (#50-class extended) — 3 FRED CSVs with `series` tag:
    - GDP+UNRATE+CPI → economic_3way.csv (`multi_csv_concat_fred_3way_economic`)
  - Wikipedia/HTML extract — apollo-program.html → h2 titles
    (`multi_wikipedia_extract_titles_apollo`)
- Cross-app real-asset rows (Cat I — third pass added 2026-05-10, +21 rows):
  - CSV → docx text summary (NOT table) — `compare_docx_files`:
    - us-population most/least populous sentence
      (`multi_csv_to_docx_text_summary_us_pop`)
    - world-gdp top-3 sentence
      (`multi_csv_to_docx_text_summary_world_gdp_top3`)
    - 2024 unemployment annual average sentence
      (`multi_csv_to_docx_text_summary_unemployment_2024_avg`)
  - arxiv PDF → text via pymupdf (`import fitz`) — `compare_text_file`:
    - extract paper title from page 0
      (`multi_arxiv_pdf_extract_title`)
    - extract page count
      (`multi_arxiv_pdf_count_pages`)
  - Multi-CSV concat to multi-sheet xlsx — `compare_table` per-sheet:
    - 2-sheet world GDP+Population
      (`multi_csv_concat_world_to_xlsx_2sheet`)
    - 3-sheet US Population+Income+HousingStarts
      (`multi_csv_concat_us_states_to_xlsx_3sheet`)
  - Code asset → docx (one paragraph per source line) — `compare_docx_files`:
    - first 50 lines of flask-app.py
      (`multi_code_to_docx_flask_app`)
    - first 40 lines of gin-handler.go
      (`multi_code_to_docx_gin_handler`)
- Cat J — pandoc rows (added 2026-05-10 once Dockerfile installs pandoc):
  - md → docx via `pandoc -s` — `compare_docx_files`:
    - meeting notes (`pandoc_md_to_docx_meeting_notes`)
    - changelog (`pandoc_md_to_docx_changelog`)
    - recipe (`pandoc_md_to_docx_recipe`)
  - csv → markdown table via `pandoc -f csv -t markdown` — `compare_text_file`:
    - inventory (`pandoc_csv_to_md_inventory`)
  - HTML (Wikipedia) → markdown via `pandoc -t markdown_strict --wrap=none`
    — `compare_text_file`:
    - apollo-program.html (`pandoc_html_to_md_apollo`)
    - octopus.html (`pandoc_html_to_md_octopus`)
  - docx → plain text via `pandoc -t plain` — `compare_text_file`:
    - quarterly report (`pandoc_docx_to_txt_report`)
- Cat K — ImageMagick rows (added 2026-05-10 once Dockerfile installs IM):
  - photo resize/rotate/grayscale/flip/flop/crop via `convert` — `compare_images`:
    - resize 800x600 jupiter (`im_resize_800x600_jupiter`)
    - resize 200x200 thumb tiger (`im_resize_thumb200_tiger`)
    - resize w400 skyscraper (`im_resize_w400_skyscraper`)
    - rotate 90 horse (`im_rotate_90_horse`)
    - rotate 180 pizza (`im_rotate_180_pizza`)
    - rotate 270 andromeda (`im_rotate_270_andromeda`)
    - grayscale house (`im_grayscale_house`)
    - grayscale apollo11 (`im_grayscale_apollo11`)
    - flip mars (`im_flip_mars`)
    - flop coffee (`im_flop_coffee`)
    - crop center 400x400 bird (`im_crop_center400_bird`)
  - 2x2 contact-sheet via `montage` — `compare_images`:
    - food contact sheet (`im_montage_contact_sheet_food`)
- Cat L — poppler-utils rows (added 2026-05-10 once Dockerfile installs
  poppler-utils — `pdftotext` / `pdfinfo`):
  - extract page 1 of arxiv attention paper (`pdftotext_page1_attention`)
    — `compare_text_file`
  - extract full text of arxiv BERT paper (`pdftotext_full_bert`)
    — `compare_text_file`
  - extract page 1 layout of arxiv pix2pix (`pdftotext_page1_pix2pix`)
    — `compare_text_file`
  - extract page count of arxiv GPT-3 via pdfinfo+awk (`pdfinfo_pagecount_gpt3`)
    — `compare_text_file`
- Cat M — pdftoppm rows (added 2026-05-10 once Dockerfile installs
  poppler-utils): render PDF page → PNG via `pdftoppm -png -r DPI -f N -l N`
  — `compare_images`:
    - attention page 1 @ 150 dpi (`pdftoppm_page1_attention_150dpi`)
    - BERT page 1 @ 100 dpi (`pdftoppm_page1_bert_100dpi`)

Deferred (per devs/envs/lite.osworld/synth/multi_apps.md `## Implementation status`):
- Chrome-source rows #1, #2, #3, #23, #40, #70, #92, #111, #129, #164-#168 —
  DROPPED (2026-05-10 deferred-bridge audit). Two failure modes ruled this
  family out:
    (a) UI clipboard path: chrome → Ctrl+A → Ctrl+C → paste into terminal /
        Writer is brittle in scripted oracle (no deterministic clipboard
        state across rollouts).
    (b) Shell-shortcut path: agent reads the staged HTML via `cat /home/user/
        Desktop/<page>.html` instead of using chrome at all → no chrome
        skill is exercised. Already covered by `_make_asset_wikipedia_h2_
        extract` and the `_PANDOC_HTML_TO_MD_SPECS` family which both
        operate on staged Wikipedia HTML via shell extraction; adding
        chrome-as-prop rows would only duplicate the shell-extract reward
        landscape under a misleading "chrome" label.
  Eval has 0 multi_apps rows that read content FROM chrome into a sink
  file (the 7 chrome-mentioning multi_apps eval rows all use chrome as
  the SINK or as a search-target — covered by the existing chrome.py
  active_tab / search_query / is_expected_tabs templates).
- Thunderbird-source rows #7-#10, #12, #13, #15, #28, #60, #81, #100,
  #273-#280 — UI eml-save / sqlite-export ops in TB are hard to oracle
  without per-message scripted pulls.
- compare_docx_images residuals #49, #139, #141 — bridged 2026-05-10 by
  `_make_photo_to_docx_images_eval_template` (single-photo) +
  `_make_dual_photo_to_docx_images_eval_template` (dual-photo). PIL
  pixel-equality is functionally distinct from `compare_docx_strict
  examine_images=True` (raw blob hash); 2 templates close the long-tail
  K=1 func gap. Other rows in the original deferred set (#14-#17, #82,
  #132, #197-#200) are already covered by the byte-hash photo→docx
  family (Cat I.2) which trains the same skill under the stricter eval.
- ffmpeg rows #27, #59, #67, #98, #134 — handled by the parallel `vlc.py`
  subagent (multi_apps owns nothing here).
- pdftk rows #37, #76, #117 — pdftk-specific (merge/split) deferred:
  no eval class for pdftk-style outputs; revisit if `compare_pdfs` is
  exercised on merged outputs.
- Calibre rows #122 — not in baseline.
- Real-asset rows referencing `assets/synth/{photos, html/wikipedia}/...` —
  defer photo-class rows (need ImageMagick / GIMP) and html/wikipedia rows
  (need pandoc); shell-pipeline + csv-source asset rows are now implemented
  via `_stage_asset` (see "Real-asset" group above).
- Expansion template-mode expansions
  (#201-#341) — predominantly variants of already-deferred classes;
  revisit as base classes are unblocked.

Re-enable plan: wire `host_push` for asset staging → unblock real-asset
multi-app rows; install pandoc + ffmpeg + ImageMagick in Dockerfile (per
AGENTS.md tool-install priority list) → unblock the bulk of remaining rows.

Usage:
    uv run python -m lite.gym.envs.lite.osworld.src.gen.train \\
        --track synth --domain multi_apps
"""

from __future__ import annotations

import random
import textwrap

from lite.gym.envs.lite.osworld.src.gen.common import (
    LO_SAVE_POSTCONFIG,
)
from lite.gym.envs.lite.osworld.src.gen.train.synth._utils import (
    SynthTemplate,
    _stage_asset,
    _stable_hash,
)

# Sibling-module builders — refactor (2026-05-10): swap inline placeholder
# source content (4-6 line lorem-ish strings, 5-row synthetic tables, generic
# pptx bullets) for topic-coherent real content sourced from writer.py
# (Gutenberg books / Wikipedia articles), calc.py (real CSV builders), and
# impress.py (TopicTheme decks). multi_apps factories only IMPORT — sibling
# files are untouched.
from lite.gym.envs.lite.osworld.src.gen.train.synth.libreoffice_writer import (
    _WIKI_ARTICLES,
)
from lite.gym.envs.lite.osworld.src.gen.train.synth.libreoffice_calc import (
    _csv_src_oil_wti_daily,
    _csv_src_us_fed_funds,
    _csv_src_us_gdp,
    _csv_src_us_housing_starts,
    _csv_src_us_inflation_cpi,
    _csv_src_us_mortgage_30yr,
    _csv_src_us_population_states,
    _csv_src_us_state_median_income,
    _csv_src_us_unemployment,
    _csv_src_world_gdp_2022,
    _csv_src_world_population_2022,
)
from lite.gym.envs.lite.osworld.src.gen.train.synth.libreoffice_impress import (
    _TOPIC_FAMILIES as _IMPRESS_TOPIC_FAMILIES,
    TopicTheme,
)
from lite.gym.envs.lite.osworld.src.gen.train.synth._utils import (
    _terminal_preopen_steps,
)


# Map asset rel-paths → matching real-CSV builders from calc.py. The
# `_csv_src_*` heredocs cap each CSV at 26-36 rows for Q-scale safety; we
# pre-truncate the staged CSV to the same cap so the agent-saved xlsx
# matches the gold xlsx row-for-row.
_CSV_FUNC_BY_ASSET: dict[str, tuple[callable, int]] = {
    "data/csv/us-population-states.csv":   (_csv_src_us_population_states,   30),
    "data/csv/us-gdp.csv":                 (_csv_src_us_gdp,                 32),
    "data/csv/us-unemployment.csv":        (_csv_src_us_unemployment,        36),
    "data/csv/world-gdp-2022.csv":         (_csv_src_world_gdp_2022,         30),
    "data/csv/oil-wti-daily.csv":          (_csv_src_oil_wti_daily,          30),
    "data/csv/us-fed-funds-rate.csv":      (_csv_src_us_fed_funds,           36),
    "data/csv/us-housing-starts.csv":      (_csv_src_us_housing_starts,      36),
    "data/csv/us-inflation-cpi.csv":       (_csv_src_us_inflation_cpi,       36),
    "data/csv/us-mortgage-30yr.csv":       (_csv_src_us_mortgage_30yr,       26),
    "data/csv/us-state-median-income.csv": (_csv_src_us_state_median_income, 30),
    # Validation fix: world-population CSV's first 48 rows are regional
    # aggregates (Africa, OECD, etc.); per-country rows start at row 49
    # (AFG/Afghanistan). 30 rows misses ALL countries — bumped to 50 so the
    # first ~10 countries (Afghanistan..Australia) are retained, matching the
    # "per-entity rows" instruction in xlsx_to_docx_table_world_population.
    "data/csv/world-population-2022.csv":  (_csv_src_world_population_2022,  50),
}


# ---------------------------------------------------------------------------
# Generic helpers
# ---------------------------------------------------------------------------

_DESKTOP = "/home/user/Desktop"
_TMP = "/tmp"


def _execute(command: str, *, shell: bool = True) -> dict:
    return {"type": "execute", "parameters": {"command": command, "shell": shell}}


def _write_text_step(path: str, content: str) -> dict:
    """Write a plain-text file via heredoc (avoids shell-quote hell)."""
    py = textwrap.dedent(f"""\
        import os
        path = {path!r}
        content = {content!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        with open(path, 'w') as f:
            f.write(content)
        """)
    return _execute(f"python3 << 'PYEOF'\n{py}\nPYEOF")


def _python_step(py: str) -> dict:
    """Run an arbitrary python heredoc inside the container."""
    return _execute(f"python3 << 'PYEOF'\n{py}\nPYEOF")


def _lo_normalize_cmd(path: str, fmt: str) -> str:
    """Round-trip `path` through `soffice --headless --convert-to <fmt>` so its
    XML structure matches what the evaluator runner produces when it normalizes
    expected files on the eval side. Byte-identical to the helper in
    `perturb/libreoffice_writer.py` / `perturb/libreoffice_impress.py` so the
    LO-normalize asymmetry fix is uniform across domains.
    """
    return (
        f"tmpd=$(mktemp -d) && "
        f"DISPLAY=:1 soffice --headless --norestore --nofirststartwizard "
        f"--convert-to {fmt} --outdir \"$tmpd\" '{path}' 2>/dev/null && "
        f"[ -f \"$tmpd/$(basename '{path}')\" ] && "
        f"cp \"$tmpd/$(basename '{path}')\" '{path}'; "
        f"rm -rf \"$tmpd\"; true"
    )


def _build_oracle_lo(sink: str, expected: str, fmt: str = "docx") -> list[dict]:
    """3-step oracle that fixes the LO-normalize asymmetry bug for docx/pptx
    file-op evaluators (compare_docx_files / compare_docx_strict /
    compare_docx_tables / compare_docx_images / compare_pptx_files):

      ① normalize gold (so result = normalize(normalize(expected)))
      ② plant: cp gold → sink
      ③ normalize sink (so the sink matches what eval-side LO normalize produces)

    Without step ①, fields LO mutates during normalization (e.g. font fallback,
    XML namespace ordering) can leave expected and sink-after-normalize out of
    sync → false-fail on otherwise-passing tasks. Mirror of `_build_oracle` in
    `perturb/libreoffice_writer.py` (docx) and `_standard_oracle` minus the
    expected_py step in `perturb/libreoffice_impress.py` (pptx).
    """
    return [
        _execute(_lo_normalize_cmd(expected, fmt)),
        _execute(f"cp '{expected}' '{sink}'"),
        _execute(_lo_normalize_cmd(sink, fmt)),
    ]


def _build_docx_step(path: str, paragraphs: list[tuple[str | None, str]]) -> dict:
    """Build a docx via python-docx. paragraphs = list of (style, text).

    style is None for normal body, or e.g. "Heading 1" / "Title".
    """
    py = textwrap.dedent(f"""\
        import os
        from docx import Document
        path = {path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        doc = Document()
        for style, text in {paragraphs!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        doc.save(path)
        """)
    return _python_step(py)


def _build_docx_with_table_step(
    path: str,
    pre_paragraphs: list[tuple[str | None, str]],
    table_data: list[list[str]],
    post_paragraphs: list[tuple[str | None, str]] | None = None,
) -> dict:
    """Build a docx containing pre-paragraphs + one table + post-paragraphs."""
    post_paragraphs = post_paragraphs or []
    py = textwrap.dedent(f"""\
        import os
        from docx import Document
        path = {path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        doc = Document()
        for style, text in {pre_paragraphs!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        rows = {table_data!r}
        if rows:
            tbl = doc.add_table(rows=len(rows), cols=len(rows[0]))
            for i, row in enumerate(rows):
                for j, cell in enumerate(row):
                    tbl.cell(i, j).text = str(cell)
        for style, text in {post_paragraphs!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        doc.save(path)
        """)
    return _python_step(py)


def _build_xlsx_step(
    path: str,
    data: list[list],
    sheet_name: str = "Sheet1",
) -> dict:
    """Build an xlsx via openpyxl, then LO-normalize so cached <v> values
    match LO save format (per AGENTS.md F1 mitigation).
    """
    py = textwrap.dedent(f"""\
        import os
        from openpyxl import Workbook
        path = {path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        wb = Workbook()
        ws = wb.active
        ws.title = {sheet_name!r}
        for row in {data!r}:
            ws.append(list(row))
        wb.save(path)
        # _LO_NORMALIZE_TAIL: re-save via soffice headless so cached values
        # match LO save format (per AGENTS.md F1 mitigation).
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, path],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename(path))
            if os.path.exists(_conv):
                _sh.copy(_conv, path)
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)
    return _python_step(py)


def _csv_text(rows: list[list[str]]) -> str:
    return "\n".join(",".join(r) for r in rows) + "\n"


def _truncate_csv_step(csv_path: str, n_data_rows: int) -> dict:
    """Truncate a staged CSV in-place to (header + first N data rows).
    Used to align the agent-opened CSV with the row cap baked into the
    matching `_csv_src_*` heredoc from calc.py (so gold and agent-output
    match). 2026-05-10 sibling-builder refactor.
    """
    py = textwrap.dedent(f"""\
        import csv
        path = {csv_path!r}
        n = {n_data_rows!r}
        with open(path, newline='') as f:
            rows = list(csv.reader(f))
        header, data = rows[0], rows[1:1 + n]
        with open(path, 'w', newline='') as f:
            w = csv.writer(f)
            w.writerow(header)
            w.writerows(data)
        """)
    return _python_step(py)


def _pick_impress_theme(seed: int, salt: str = "") -> TopicTheme:
    """Deterministically pick a TopicTheme from impress's `_TOPIC_FAMILIES`
    list (10 themes) given seed + salt. Mirrors impress.py `_pick_topic`.
    2026-05-10 sibling-builder refactor — used by photo-template factories
    to source topic-coherent title + body text from sibling impress module.
    """
    idx = (seed ^ (_stable_hash(salt) & 0xFFFFFFFF)) % len(_IMPRESS_TOPIC_FAMILIES)
    return _IMPRESS_TOPIC_FAMILIES[idx]


# ---------------------------------------------------------------------------
# Cat A — txt source -> docx (writer save)
# ---------------------------------------------------------------------------
# Pattern: pre_config writes plain-text source AND a pre-built gold docx
# whose paragraphs equal the source lines. Agent opens the txt in writer,
# saves as docx. Eval: compare_docx_files (paragraphs equal, ignore_blanks).
# Oracle: cp gold over the agent's expected output path (post-postconfig).

# 2026-05-10 sibling-builder refactor: source text now comes from
# `_WIKI_ARTICLES` (writer.py) rather than 4-6 line inline placeholders.
# Each spec picks one Wikipedia article and `n_paras` body paragraphs.
# `lines` is built at factory-call time from real Wikipedia prose so the
# source txt + gold docx contain topic-coherent multi-sentence paragraphs
# instead of synthetic bullet snippets.
_TXT_TO_DOCX_SPECS: list[dict] = [
    {
        "id": "txt_to_docx_coffee",
        "src_basename": "coffee.txt",
        "dst_basename": "coffee.docx",
        "article_key": "coffee",
        "n_paras": 3,
        "instructions": [
            "Open ~/Desktop/coffee.txt in Writer and save it as ~/Desktop/coffee.docx (keep one paragraph per line).",
            "Convert ~/Desktop/coffee.txt into a docx file at ~/Desktop/coffee.docx using LibreOffice Writer.",
        ],
    },
    {
        "id": "txt_to_docx_eiffel_tower",
        "src_basename": "eiffel_tower.txt",
        "dst_basename": "eiffel_tower.docx",
        "article_key": "eiffel-tower",
        "n_paras": 3,
        "instructions": [
            "Open ~/Desktop/eiffel_tower.txt in Writer and save it as ~/Desktop/eiffel_tower.docx.",
        ],
    },
    {
        "id": "txt_to_docx_octopus",
        "src_basename": "octopus.txt",
        "dst_basename": "octopus.docx",
        "article_key": "octopus",
        "n_paras": 3,
        "instructions": [
            "Take the Wikipedia excerpt at ~/Desktop/octopus.txt and save it as ~/Desktop/octopus.docx in Writer.",
        ],
    },
    # Mass-add — txt → docx variants (now real Wiki content
    # post sibling-builder refactor).
    # Pruned (txt_to_docx_pizza): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "txt_to_docx_pizza",
#         "src_basename": "pizza.txt",
#         "dst_basename": "pizza.docx",
#         "article_key": "pizza",
#         "n_paras": 3,
#         "instructions": [
#             "Open ~/Desktop/pizza.txt in Writer and save it as ~/Desktop/pizza.docx (one paragraph per line).",
#         ],
#     },
    {
        "id": "txt_to_docx_solar_system",
        "src_basename": "solar_system.txt",
        "dst_basename": "solar_system.docx",
        "article_key": "solar-system",
        "n_paras": 3,
        "instructions": [
            "Convert the Wikipedia excerpt at ~/Desktop/solar_system.txt to ~/Desktop/solar_system.docx using Writer.",
        ],
    },
    # Pruned (txt_to_docx_volcano): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "txt_to_docx_volcano",
#         "src_basename": "volcano.txt",
#         "dst_basename": "volcano.docx",
#         "article_key": "volcano",
#         "n_paras": 3,
#         "instructions": [
#             "Open ~/Desktop/volcano.txt in Writer and save as ~/Desktop/volcano.docx.",
#         ],
#     },
    # Pruned (txt_to_docx_beethoven): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "txt_to_docx_beethoven",
#         "src_basename": "beethoven.txt",
#         "dst_basename": "beethoven.docx",
#         "article_key": "beethoven",
#         "n_paras": 3,
#         "instructions": [
#             "Convert ~/Desktop/beethoven.txt to a docx at ~/Desktop/beethoven.docx in Writer (preserve paragraph order).",
#         ],
#     },
    # Pruned (txt_to_docx_mona_lisa): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "txt_to_docx_mona_lisa",
#         "src_basename": "mona_lisa.txt",
#         "dst_basename": "mona_lisa.docx",
#         "article_key": "mona-lisa",
#         "n_paras": 3,
#         "instructions": [
#             "Open ~/Desktop/mona_lisa.txt in Writer and save as ~/Desktop/mona_lisa.docx.",
#         ],
#     },
    # Pruned (txt_to_docx_apollo_program): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "txt_to_docx_apollo_program",
#         "src_basename": "apollo_program.txt",
#         "dst_basename": "apollo_program.docx",
#         "article_key": "apollo-program",
#         "n_paras": 3,
#         "instructions": [
#             "Open ~/Desktop/apollo_program.txt in Writer and save it as ~/Desktop/apollo_program.docx.",
#         ],
#     },
]


def _make_txt_to_docx_template(spec: dict) -> SynthTemplate:
    src = f"{_DESKTOP}/{spec['src_basename']}"
    dst = f"{_DESKTOP}/{spec['dst_basename']}"
    expected = f"{_TMP}/expected_{spec['id']}.docx"
    # Sibling-builder refactor: each spec sources its body paragraphs from
    # writer.py's `_WIKI_ARTICLES[article_key]` (loaded at import). One
    # paragraph per line so the txt → docx mapping (one Writer paragraph
    # per line in the source txt) stays well-defined for the agent.
    lines = _WIKI_ARTICLES[spec["article_key"]][: spec["n_paras"]]
    body = "\n".join(lines) + "\n"
    paragraphs: list[tuple[str | None, str]] = [(None, line) for line in lines]
    instructions = spec["instructions"]

    init_steps = [
        _write_text_step(src, body),
        _build_docx_step(expected, paragraphs),
        # validation gap-close: agent opens the staged .txt in Writer and
        # "save-as" to .docx — focused window is the source .txt.
        *_multi_app_preopen(foreground=src),
    ]
    # Oracle: 3-step LO-normalize (normalize expected → cp → normalize sink)
    # to fix the LO-normalize asymmetry that bites docx file-op evaluators.
    # With oracle_after_postconfig=True, runs AFTER LO_SAVE_POSTCONFIG so the
    # final dst contains the gold (validates that the eval grades the
    # post-save state correctly).
    oracle_steps = _build_oracle_lo(dst, expected, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": dst, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.docx"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(spec["id"]) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", src],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=spec["id"],
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=2,
        eval_class="compare_docx_files",
        setup_class="txt_source",
    )


# ---------------------------------------------------------------------------
# Cat B — csv -> xlsx via Calc (calc save)
# ---------------------------------------------------------------------------
# Pattern: pre_config writes csv source AND a gold xlsx (LO-normalized).
# Agent opens the csv in Calc and saves as xlsx. Eval: compare_table
# rules=[sheet_data] over Sheet1.

# 2026-05-10 sibling-builder refactor: source CSV is now a real FRED /
# US-Census / world-bank dataset staged from `assets/synth/data/csv/` (paired
# with the matching `_csv_src_*` heredoc in calc.py for row-count parity).
# Each spec rotates through a distinct asset so source content is
# topic-coherent real data instead of a 5-row synthetic table. The staged
# CSV is pre-truncated to the row cap baked into the matching `_csv_src_*`
# (30-36 rows; Q-scale safe per calc.py).
_CSV_TO_XLSX_SPECS: list[dict] = [
    {
        "id": "csv_to_xlsx_us_population",
        "asset_rel": "data/csv/us-population-states.csv",
        "src_basename": "us-population-states.csv",
        "dst_basename": "us-population-states.xlsx",
        "instructions": [
            "Open ~/Desktop/us-population-states.csv in Calc and save it as ~/Desktop/us-population-states.xlsx (Excel xlsx format).",
        ],
    },
    {
        "id": "csv_to_xlsx_us_gdp",
        "asset_rel": "data/csv/us-gdp.csv",
        "src_basename": "us-gdp.csv",
        "dst_basename": "us-gdp.xlsx",
        "instructions": [
            "Open ~/Desktop/us-gdp.csv in Calc and save it as ~/Desktop/us-gdp.xlsx.",
        ],
    },
    {
        "id": "csv_to_xlsx_us_unemployment",
        "asset_rel": "data/csv/us-unemployment.csv",
        "src_basename": "us-unemployment.csv",
        "dst_basename": "us-unemployment.xlsx",
        "instructions": [
            "Open ~/Desktop/us-unemployment.csv in Calc and save it as ~/Desktop/us-unemployment.xlsx.",
        ],
    },
    {
        "id": "csv_to_xlsx_world_gdp",
        "asset_rel": "data/csv/world-gdp-2022.csv",
        "src_basename": "world-gdp-2022.csv",
        "dst_basename": "world-gdp-2022.xlsx",
        "instructions": [
            "Open ~/Desktop/world-gdp-2022.csv in Calc and save it as ~/Desktop/world-gdp-2022.xlsx.",
        ],
    },
    # Mass-add — csv → xlsx variants (now real-asset
    # post sibling-builder refactor).
    # Pruned (csv_to_xlsx_oil_wti): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "csv_to_xlsx_oil_wti",
#         "asset_rel": "data/csv/oil-wti-daily.csv",
#         "src_basename": "oil-wti-daily.csv",
#         "dst_basename": "oil-wti-daily.xlsx",
#         "instructions": [
#             "Open ~/Desktop/oil-wti-daily.csv in Calc and save it as ~/Desktop/oil-wti-daily.xlsx.",
#         ],
#     },
    # Pruned (csv_to_xlsx_us_fed_funds): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "csv_to_xlsx_us_fed_funds",
#         "asset_rel": "data/csv/us-fed-funds-rate.csv",
#         "src_basename": "us-fed-funds-rate.csv",
#         "dst_basename": "us-fed-funds-rate.xlsx",
#         "instructions": [
#             "Open ~/Desktop/us-fed-funds-rate.csv in Calc and save it as ~/Desktop/us-fed-funds-rate.xlsx.",
#         ],
#     },
    # Pruned (csv_to_xlsx_us_housing): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "csv_to_xlsx_us_housing",
#         "asset_rel": "data/csv/us-housing-starts.csv",
#         "src_basename": "us-housing-starts.csv",
#         "dst_basename": "us-housing-starts.xlsx",
#         "instructions": [
#             "Open ~/Desktop/us-housing-starts.csv in Calc and save it as ~/Desktop/us-housing-starts.xlsx.",
#         ],
#     },
    # Pruned (csv_to_xlsx_us_inflation): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "csv_to_xlsx_us_inflation",
#         "asset_rel": "data/csv/us-inflation-cpi.csv",
#         "src_basename": "us-inflation-cpi.csv",
#         "dst_basename": "us-inflation-cpi.xlsx",
#         "instructions": [
#             "Open ~/Desktop/us-inflation-cpi.csv in Calc and save as ~/Desktop/us-inflation-cpi.xlsx.",
#         ],
#     },
    # Pruned (csv_to_xlsx_us_mortgage): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "csv_to_xlsx_us_mortgage",
#         "asset_rel": "data/csv/us-mortgage-30yr.csv",
#         "src_basename": "us-mortgage-30yr.csv",
#         "dst_basename": "us-mortgage-30yr.xlsx",
#         "instructions": [
#             "Open ~/Desktop/us-mortgage-30yr.csv in Calc and save it as ~/Desktop/us-mortgage-30yr.xlsx.",
#         ],
#     },
]


def _make_csv_to_xlsx_template(spec: dict) -> SynthTemplate:
    src = f"{_DESKTOP}/{spec['src_basename']}"
    dst = f"{_DESKTOP}/{spec['dst_basename']}"
    expected = f"{_TMP}/expected_{spec['id']}.xlsx"
    asset_rel = spec["asset_rel"]
    _csv_func, n_data_rows = _CSV_FUNC_BY_ASSET[asset_rel]
    instructions = spec["instructions"]

    # Sibling-builder refactor (2026-05-10): stage the real CSV, truncate to
    # the row cap baked into the matching `_csv_src_*` heredoc, then build
    # gold xlsx (Sheet1, LO-normalized) by reading the truncated CSV.
    # validation: eval rule was `sheet_idx0=RNSheet1` (name-match "Sheet1") but
    # LO Calc Save-As preserves the CSV basename as sheet name (e.g.
    # `us-population-states`), NOT "Sheet1" — so the name lookup never found
    # the result sheet. Switched to positional `sheet_idx0=0, sheet_idx1=0`
    # (first sheet, name-agnostic); gold still writes "Sheet1" but the eval
    # no longer cares.
    gold_py = textwrap.dedent(f"""\
        import os, csv
        from openpyxl import Workbook
        src = {src!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        wb = Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        with open(src, newline='') as f:
            for row in csv.reader(f):
                ws.append(row)
        wb.save(path)
        # _LO_NORMALIZE_TAIL: match LO save format (per AGENTS.md F1).
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, path],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename(path))
            if os.path.exists(_conv):
                _sh.copy(_conv, path)
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)
    init_steps = [
        _stage_asset(asset_rel, src),
        _truncate_csv_step(src, n_data_rows),
        _python_step(gold_py),
        # validation gap-close: agent opens the staged .csv in Calc and
        # "save-as" to .xlsx — focused window is the source .csv.
        *_multi_app_preopen(foreground=src),
    ]
    oracle_steps = [_execute(f"cp '{expected}' '{dst}'")]

    evaluator = {
        "func": "compare_table",
        "result": {"type": "vm_file", "path": dst, "dest": "result.xlsx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.xlsx"},
        # Validation fix: positional sheet index (name-agnostic). Agent Save-As
        # produces a sheet named after the CSV basename, not "Sheet1".
        "options": {"rules": [{"type": "sheet_data", "sheet_idx0": 0, "sheet_idx1": 0}]},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(spec["id"]) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", src],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=spec["id"],
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=2,
        eval_class="compare_table",
        setup_class="csv_source",
    )


# ---------------------------------------------------------------------------
# Cat C — xlsx -> csv export via Calc
# ---------------------------------------------------------------------------
# Pattern: pre_config writes xlsx source AND a gold csv (the equivalent
# csv form). Agent opens the xlsx and saves as csv. Eval: compare_csv.

_XLSX_TO_CSV_SPECS: list[dict] = [
    {
        "id": "xlsx_to_csv_sales",
        "src_basename": "sales.xlsx",
        "dst_basename": "sales.csv",
        "rows": [
            ["Quarter", "Revenue", "Units"],
            ["Q1", "120000", "300"],
            ["Q2", "140000", "350"],
            ["Q3", "180000", "420"],
            ["Q4", "210000", "510"],
        ],
        "instructions": [
            "Open ~/Desktop/sales.xlsx in Calc and save it as ~/Desktop/sales.csv (CSV format, comma delimiter).",
        ],
    },
    {
        "id": "xlsx_to_csv_employees",
        "src_basename": "employees.xlsx",
        "dst_basename": "employees.csv",
        "rows": [
            ["Name", "Dept", "Title", "Years"],
            ["Alice", "Eng", "Senior", "5"],
            ["Bob", "Sales", "Lead", "3"],
            ["Carol", "Eng", "Staff", "7"],
            ["Dave", "HR", "Manager", "4"],
        ],
        "instructions": [
            "Open ~/Desktop/employees.xlsx in Calc and export it as ~/Desktop/employees.csv (CSV format).",
        ],
    },
    {
        "id": "xlsx_to_csv_grades",
        "src_basename": "grades.xlsx",
        "dst_basename": "grades.csv",
        "rows": [
            ["Student", "Test1", "Test2", "Test3"],
            ["Alice", "92", "88", "95"],
            ["Bob", "78", "85", "80"],
            ["Carol", "85", "90", "87"],
        ],
        "instructions": [
            "Open ~/Desktop/grades.xlsx in Calc and save it as ~/Desktop/grades.csv (CSV).",
        ],
    },
    # Mass-add — xlsx → csv variants.
    {
        "id": "xlsx_to_csv_invoices",
        "src_basename": "invoices.xlsx",
        "dst_basename": "invoices.csv",
        "rows": [
            ["InvoiceID", "Customer", "Amount", "Status"],
            ["INV001", "Acme",   "1200.00", "Paid"],
            ["INV002", "Globex", "850.00",  "Pending"],
            ["INV003", "Hooli",  "2400.50", "Paid"],
            ["INV004", "Initech","620.00",  "Overdue"],
        ],
        "instructions": [
            "Open ~/Desktop/invoices.xlsx in Calc and export it as ~/Desktop/invoices.csv (CSV format).",
        ],
    },
    {
        "id": "xlsx_to_csv_shipments",
        "src_basename": "shipments.xlsx",
        "dst_basename": "shipments.csv",
        "rows": [
            ["ShipmentID", "Origin", "Destination", "Weight_kg"],
            ["SH01", "BCN", "FRA", "12.4"],
            ["SH02", "AMS", "JFK", "30.0"],
            ["SH03", "LHR", "SIN", "8.6"],
            ["SH04", "LAX", "NRT", "22.1"],
        ],
        "instructions": [
            "Open ~/Desktop/shipments.xlsx in Calc and save it as ~/Desktop/shipments.csv.",
        ],
    },
    {
        "id": "xlsx_to_csv_inventory_v2",
        "src_basename": "inventory_v2.xlsx",
        "dst_basename": "inventory_v2.csv",
        "rows": [
            ["SKU", "Item", "Stock"],
            ["P101", "Pen Black",  "240"],
            ["P102", "Pen Blue",   "190"],
            ["P201", "Pencil HB",  "510"],
            ["P301", "Marker Red", "85"],
        ],
        "instructions": [
            "Open ~/Desktop/inventory_v2.xlsx in Calc and export to ~/Desktop/inventory_v2.csv.",
        ],
    },
]


def _make_xlsx_to_csv_template(spec: dict) -> SynthTemplate:
    src = f"{_DESKTOP}/{spec['src_basename']}"
    dst = f"{_DESKTOP}/{spec['dst_basename']}"
    expected = f"{_TMP}/expected_{spec['id']}.csv"
    rows = spec["rows"]
    instructions = spec["instructions"]

    init_steps = [
        _build_xlsx_step(src, rows),
        _write_text_step(expected, _csv_text(rows)),
        # validation gap-close: agent opens the source .xlsx in Calc and
        # "save-as" to .csv — focused window is the source .xlsx.
        *_multi_app_preopen(foreground=src),
    ]
    oracle_steps = [_execute(f"cp '{expected}' '{dst}'")]

    evaluator = {
        "func": "compare_csv",
        "result": {"type": "vm_file", "path": dst, "dest": "result.csv"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.csv"},
        "options": {"strict": False},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(spec["id"]) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", src],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=spec["id"],
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=2,
        # eval_class share: bucket with compare_table so the rescaler doesn't
        # over-create these (they share the calc->csv evaluator landscape).
        eval_class="compare_table",
        setup_class="xlsx_source",
    )


# ---------------------------------------------------------------------------
# Cat D — xlsx -> docx-table (writer save)
# ---------------------------------------------------------------------------
# Pattern: pre_config writes (a) source xlsx with the data, (b) source
# docx (a partial report with no table yet), and (c) gold docx with the
# table inserted. Agent opens both files, copies xlsx data, inserts table
# into the docx, saves. Eval: compare_docx_tables (table cell text equal).

# 2026-05-10 sibling-builder refactor: source xlsx now built from real CSV
# assets via calc.py's `_csv_src_*` heredocs (paired with the same row caps).
# Each spec rotates through a distinct asset so source content is
# topic-coherent real data instead of a 4-5 row synthetic table.
_XLSX_TO_DOCX_TABLE_SPECS: list[dict] = [
    {
        "id": "xlsx_to_docx_table_us_population",
        "asset_rel": "data/csv/us-population-states.csv",
        "xlsx_basename": "us-population-states.xlsx",
        "docx_basename": "population_report.docx",
        "docx_pre": [
            ("Title", "U.S. State Population (2020)"),
            (None, "The table below shows the per-state population from the spreadsheet."),
        ],
        "instructions": [
            "Open ~/Desktop/us-population-states.xlsx and ~/Desktop/population_report.docx. Insert the spreadsheet contents as a 3-column table at the end of the document, then save the document.",
        ],
    },
    {
        "id": "xlsx_to_docx_table_us_gdp",
        "asset_rel": "data/csv/us-gdp.csv",
        "xlsx_basename": "us-gdp.xlsx",
        "docx_basename": "us_gdp_report.docx",
        "docx_pre": [
            ("Title", "U.S. GDP — Quarterly"),
            (None, "Quarterly GDP observations appear in the table below."),
        ],
        "instructions": [
            "Open ~/Desktop/us-gdp.xlsx and ~/Desktop/us_gdp_report.docx. Insert the quarterly observations as a 2-column table at the end of the document and save it.",
        ],
    },
    {
        "id": "xlsx_to_docx_table_us_unemployment",
        "asset_rel": "data/csv/us-unemployment.csv",
        "xlsx_basename": "us-unemployment.xlsx",
        "docx_basename": "us_unemployment_report.docx",
        "docx_pre": [
            ("Title", "U.S. Unemployment Rate"),
            (None, "Monthly unemployment-rate observations are tabulated below."),
        ],
        "instructions": [
            "Open ~/Desktop/us-unemployment.xlsx and ~/Desktop/us_unemployment_report.docx. Insert the monthly observations as a 2-column table at the end of the document and save it.",
        ],
    },
    # Mass-add — calc → writer table-insertion variants
    # (now real-asset post sibling-builder refactor).
    {
        "id": "xlsx_to_docx_table_world_gdp",
        "asset_rel": "data/csv/world-gdp-2022.csv",
        "xlsx_basename": "world-gdp-2022.xlsx",
        "docx_basename": "world_gdp_report.docx",
        "docx_pre": [
            ("Title", "World GDP — 2022"),
            (None, "Below is the per-entity 2022 GDP breakdown from the spreadsheet."),
        ],
        "instructions": [
            "Open ~/Desktop/world-gdp-2022.xlsx and ~/Desktop/world_gdp_report.docx. Insert the spreadsheet contents as a 4-column table at the end of the document, then save it.",
        ],
    },
    # Pruned (xlsx_to_docx_table_oil_wti): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "xlsx_to_docx_table_oil_wti",
#         "asset_rel": "data/csv/oil-wti-daily.csv",
#         "xlsx_basename": "oil-wti-daily.xlsx",
#         "docx_basename": "oil_wti_report.docx",
#         "docx_pre": [
#             ("Title", "WTI Crude Oil — Daily"),
#             (None, "The daily WTI price series is shown below."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/oil-wti-daily.xlsx and ~/Desktop/oil_wti_report.docx. Insert the daily prices as a 2-column table at the end of the document and save it.",
#         ],
#     },
    # Pruned (xlsx_to_docx_table_us_fed_funds): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "xlsx_to_docx_table_us_fed_funds",
#         "asset_rel": "data/csv/us-fed-funds-rate.csv",
#         "xlsx_basename": "us-fed-funds-rate.xlsx",
#         "docx_basename": "fed_funds_report.docx",
#         "docx_pre": [
#             ("Title", "U.S. Fed Funds Rate"),
#             (None, "The table below summarizes the monthly fed-funds rate."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-fed-funds-rate.xlsx and ~/Desktop/fed_funds_report.docx. Insert the fed-funds observations as a 2-column table at the end of the document and save it.",
#         ],
#     },
    # Pruned (xlsx_to_docx_table_us_housing): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "xlsx_to_docx_table_us_housing",
#         "asset_rel": "data/csv/us-housing-starts.csv",
#         "xlsx_basename": "us-housing-starts.xlsx",
#         "docx_basename": "housing_report.docx",
#         "docx_pre": [
#             ("Title", "U.S. Housing Starts"),
#             (None, "Current on-hand housing-starts data is shown below."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-housing-starts.xlsx and ~/Desktop/housing_report.docx. Insert the housing-starts rows as a 2-column table at the end of the document and save it.",
#         ],
#     },
    # Pruned (xlsx_to_docx_table_us_inflation): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "xlsx_to_docx_table_us_inflation",
#         "asset_rel": "data/csv/us-inflation-cpi.csv",
#         "xlsx_basename": "us-inflation-cpi.xlsx",
#         "docx_basename": "inflation_report.docx",
#         "docx_pre": [
#             ("Title", "U.S. CPI Inflation — Monthly"),
#             (None, "The CPI observations below cover the relevant months."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-inflation-cpi.xlsx and ~/Desktop/inflation_report.docx. Insert the CPI numbers as a 2-column table at the end of the document and save it.",
#         ],
#     },
    # Pruned (xlsx_to_docx_table_us_state_income): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "id": "xlsx_to_docx_table_us_state_income",
#         "asset_rel": "data/csv/us-state-median-income.csv",
#         "xlsx_basename": "us-state-median-income.xlsx",
#         "docx_basename": "state_income_pack.docx",
#         "docx_pre": [
#             ("Title", "U.S. State Median Household Income"),
#             (None, "Headline per-state median-income figures are tabulated below."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-state-median-income.xlsx and ~/Desktop/state_income_pack.docx. Insert the income figures as a 3-column table at the end of the document and save it.",
#         ],
#     },
    {
        "id": "xlsx_to_docx_table_world_population",
        "asset_rel": "data/csv/world-population-2022.csv",
        "xlsx_basename": "world-population-2022.xlsx",
        "docx_basename": "world_population_report.docx",
        "docx_pre": [
            ("Title", "World Population — 2022"),
            (None, "The table below records per-entity 2022 population from the spreadsheet."),
        ],
        "instructions": [
            "Open ~/Desktop/world-population-2022.xlsx and ~/Desktop/world_population_report.docx. Insert the per-entity rows as a 4-column table at the end of the document and save it.",
        ],
    },
]


def _make_xlsx_to_docx_table_template(spec: dict) -> SynthTemplate:
    xlsx_path = f"{_DESKTOP}/{spec['xlsx_basename']}"
    docx_path = f"{_DESKTOP}/{spec['docx_basename']}"
    expected = f"{_TMP}/expected_{spec['id']}.docx"
    asset_rel = spec["asset_rel"]
    csv_func, n_data_rows = _CSV_FUNC_BY_ASSET[asset_rel]
    docx_pre = spec["docx_pre"]
    instructions = spec["instructions"]

    # Sibling-builder refactor (2026-05-10): source xlsx is now real-asset
    # data built via calc.py's `_csv_src_*` heredoc. The staged CSV is
    # pre-truncated to the matching row cap so the xlsx + the gold docx
    # table are populated from the same N rows.
    staged_csv = f"{_TMP}/_xlsx_src_{spec['id']}.csv"

    # `_csv_src_*` returns a python heredoc body that writes the xlsx.
    src_xlsx_py = csv_func(staged_csv, xlsx_path, 0)
    # _LO_NORMALIZE_TAIL — the agent's saved xlsx will be LO-normalized
    # implicitly via the Calc save; ensure gold xlsx matches by re-saving
    # the openpyxl output through soffice (mirrors `_build_xlsx_step` tail).
    src_xlsx_py = textwrap.dedent(src_xlsx_py) + textwrap.dedent(f"""\
        import os
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, {xlsx_path!r}],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename({xlsx_path!r}))
            if os.path.exists(_conv):
                _sh.copy(_conv, {xlsx_path!r})
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)

    # Gold docx: read the same truncated CSV and emit pre-paragraphs + a
    # table with header + data rows. Cell text matches the openpyxl-loaded
    # xlsx values (numbers from the CSV are coerced to int/float in some
    # `_csv_src_*` heredocs; we re-read the CSV as-text to mirror the
    # post-LO-save string representation that `compare_docx_tables` sees).
    gold_docx_py = textwrap.dedent(f"""\
        import os, csv
        from docx import Document
        from openpyxl import load_workbook
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        # Read xlsx values back as strings (mirrors what the agent's
        # Writer-table-from-xlsx flow lands as cell text).
        wb = load_workbook({xlsx_path!r}, data_only=True)
        ws = wb.worksheets[0]
        table_data = []
        for row in ws.iter_rows(values_only=True):
            table_data.append(['' if v is None else str(v) for v in row])
        doc = Document()
        for style, text in {docx_pre!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        if table_data:
            tbl = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
            for i, row in enumerate(table_data):
                for j, cell in enumerate(row):
                    tbl.cell(i, j).text = str(cell)
        doc.save(path)
        """)

    init_steps = [
        _stage_asset(asset_rel, staged_csv),
        _truncate_csv_step(staged_csv, n_data_rows),
        _python_step(src_xlsx_py),
        _build_docx_step(docx_path, docx_pre),  # source docx (no table)
        _python_step(gold_docx_py),             # gold docx (with table)
        # validation gap-close: agent opens xlsx (background) then docx
        # (foreground) — agent inserts the table into the focused docx.
        *_multi_app_preopen(background=[xlsx_path], foreground=docx_path),
    ]
    # 3-step LO-normalize oracle (normalize expected → cp → normalize sink)
    # fixes the LO-normalize asymmetry for docx file-op evaluators.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    evaluator = {
        "func": "compare_docx_tables",
        "result": {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.docx"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(spec["id"]) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=spec["id"],
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=2,
        eval_class="compare_docx_tables",
        setup_class="xlsx_docx_pair",
    )


# ---------------------------------------------------------------------------
# Cat E — Shell pipeline templates (no app launch)
# ---------------------------------------------------------------------------
# Pattern: pre_config writes source text and a gold output. Oracle runs
# the shell pipeline. Eval: compare_text_file or compare_csv. No app
# launch — terminal task; uses config_override to bypass auto-launch.

# Sherlock excerpt (Conan Doyle, public domain) — kept inline (~1KB).
_SHERLOCK_EXCERPT_LINES = [
    "To Sherlock Holmes she is always THE woman.",
    "I have seldom heard him mention her under any other name.",
    "In his eyes she eclipses and predominates the whole of her sex.",
    "It was not that he felt any emotion akin to love for Irene Adler.",
    "All emotions, and that one particularly, were abhorrent to his cold mind.",
    "He was, I take it, the most perfect reasoning machine.",
    "But as a lover he would have placed himself in a false position.",
    "He never spoke of the softer passions, save with a gibe and a sneer.",
    "And yet there was but one woman to him, and that woman was Irene Adler.",
    "I had seen little of Holmes lately.",
]


def _open_file_steps(path: str, app_hint: str | None = None) -> list[dict]:
    """Pre-open a file in its native app + sleep + activate_window.

    validation: replaces the 0/323 "launched but never
    settled" pattern from validation. Uses `open` step type (xdg-open) by
    default since lite.osworld registers .docx/.csv/.pptx/.py with reasonable
    defaults; `app_hint` overrides with an explicit binary for files whose
    xdg default is unreliable.

    Window-name strategy: substring-match the file basename. After
    `xdg-open foo.docx`, LibreOffice sets the window title to
    `foo.docx - LibreOffice Writer` (and similar for Calc/Impress), so
    `window_name=foo.docx` matches via xfwm4's substring search. Verified
    against upstream eval rows that use the same pattern (`open` +
    `activate_window` with full title incl. " - LibreOffice <App>" suffix).
    """
    basename = path.rsplit("/", 1)[-1]
    if app_hint:
        steps: list[dict] = [
            {"type": "launch", "parameters": {"command": [app_hint, path]}},
        ]
    else:
        steps = [
            {"type": "open", "parameters": {"path": path}},
        ]
    steps.append({"type": "execute", "parameters": {"command": "sleep 2", "shell": True}})
    steps.append({"type": "activate_window", "parameters": {"window_name": basename}})
    return steps


def _multi_app_preopen(
    *,
    background: list[str] | tuple[str, ...] = (),
    foreground: str | None = None,
) -> list[dict]:
    """Open files in upstream-conventional `[BG_1, BG_2, ...] → FG` order.

    Order rule (derived from 25 upstream multi_apps
    samples): background / reference files open FIRST (in order of mention),
    primary work / target file opens LAST so it ends focused / on-top. Each
    `_open_file_steps` call inserts `sleep 2` before `activate_window` for
    X11/WM rendering stability.

    NB: when this helper appends a `launch`/`open` step to `init_steps`,
    `common.py:already_has_app_open=True` and the template's `open_command`
    auto-launch block (lines 459-475) is short-circuited — our preopen
    becomes the single source of truth for which app is focused at turn 0.
    """
    steps: list[dict] = []
    for p in background:
        steps.extend(_open_file_steps(p))
    if foreground:
        steps.extend(_open_file_steps(foreground))
    return steps


def _make_shell_pipeline_template(
    *,
    template_id: str,
    src_path: str,
    dst_path: str,
    src_body: str,
    expected_body: str,
    pipeline: str,
    instructions: list[str],
    eval_class: str = "compare_text_file",
    eval_func: str = "compare_text_file",
    expected_kind: str = "text",
) -> SynthTemplate:
    """Build a shell-pipeline template.

    `pipeline`: shell command run by the oracle to produce dst_path.
    `expected_body`: gold file content (written to /tmp/expected_<id>.<ext>).
    `eval_func`: compare_text_file | compare_csv.
    """
    ext = "csv" if expected_kind == "csv" else "txt"
    expected_path = f"{_TMP}/expected_{template_id}.{ext}"

    init_steps = [
        _write_text_step(src_path, src_body),
        _write_text_step(expected_path, expected_body),
        # Pre-open a terminal so the agent's first type_text doesn't race
        # against a bare-desktop dock-click guess (validation H-cluster fix).
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(pipeline)]

    if eval_func == "compare_csv":
        evaluator = {
            "func": "compare_csv",
            "result": {"type": "vm_file", "path": dst_path, "dest": "result.csv"},
            "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.csv"},
            "options": {"strict": False},
        }
    else:
        evaluator = {
            "func": "compare_text_file",
            "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
            "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
        }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,   # no app launch (terminal task)
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class=eval_class,
        setup_class="text_source",
    )


def _shell_pipeline_templates() -> list[SynthTemplate]:
    """Build all shell-pipeline templates."""
    out: list[SynthTemplate] = []
    sherlock_body = "\n".join(_SHERLOCK_EXCERPT_LINES) + "\n"

    # Row 39 — grep 'Holmes'
    holmes_lines = "\n".join(l for l in _SHERLOCK_EXCERPT_LINES if "Holmes" in l) + "\n"
    out.append(_make_shell_pipeline_template(
        template_id="shell_grep_holmes",
        src_path=f"{_DESKTOP}/sherlock.txt",
        dst_path=f"{_DESKTOP}/holmes_lines.txt",
        src_body=sherlock_body,
        expected_body=holmes_lines,
        pipeline=f"grep 'Holmes' {_DESKTOP}/sherlock.txt > {_DESKTOP}/holmes_lines.txt",
        instructions=[
            "Open a terminal and scan ~/Desktop/sherlock.txt for every line that mentions Holmes. Write those matching lines, in their original order, to ~/Desktop/holmes_lines.txt.",
        ],
    ))

    # Row 41 — awk filter (rate > 6.0, keep header)
    rate_lines = [
        "date,rate",
        "2024-01-05,5.8",
        "2024-02-02,6.1",
        "2024-03-01,5.9",
        "2024-04-05,6.4",
        "2024-05-03,6.7",
        "2024-06-07,5.5",
    ]
    rate_body = "\n".join(rate_lines) + "\n"
    high_rate = [rate_lines[0]] + [l for l in rate_lines[1:] if float(l.split(",")[1]) > 6.0]
    out.append(_make_shell_pipeline_template(
        template_id="shell_awk_filter_rate",
        src_path=f"{_DESKTOP}/rates.csv",
        dst_path=f"{_DESKTOP}/high_rates.csv",
        src_body=rate_body,
        expected_body="\n".join(high_rate) + "\n",
        pipeline=(
            f"awk -F, 'NR==1 || $2+0 > 6.0' {_DESKTOP}/rates.csv > {_DESKTOP}/high_rates.csv"
        ),
        instructions=[
            "In a terminal, use awk on ~/Desktop/rates.csv to keep the header row and any rows where the rate column (column 2) is greater than 6.0; write the result to ~/Desktop/high_rates.csv.",
        ],
        eval_class="compare_text_file",
        eval_func="compare_csv",
        expected_kind="csv",
    ))

    # Row 50 — concat 2 csvs (keep one header). Custom builder so init has 3 files.
    csv_a = "date,rate\n2024-01,5.8\n2024-02,6.1\n"
    csv_b = "date,rate\n2024-03,5.9\n2024-04,6.4\n"
    expected_concat = "date,rate\n2024-01,5.8\n2024-02,6.1\n2024-03,5.9\n2024-04,6.4\n"
    src_a = f"{_DESKTOP}/econ/a.csv"
    src_b = f"{_DESKTOP}/econ/b.csv"
    dst_combined = f"{_DESKTOP}/combined.csv"
    expected_combined = f"{_TMP}/expected_shell_concat_csv.csv"

    init_steps_concat = [
        _write_text_step(src_a, csv_a),
        _write_text_step(src_b, csv_b),
        _write_text_step(expected_combined, expected_concat),
        # validation gap-close: terminal-driven (shell concat) — pre-open
        # terminal so the agent's first type_text lands on a real shell.
        *_terminal_preopen_steps(),
    ]
    oracle_concat = [_execute(
        f"head -n 1 {src_a} > {dst_combined} && "
        f"tail -n +2 {src_a} >> {dst_combined} && "
        f"tail -n +2 {src_b} >> {dst_combined}"
    )]
    evaluator_concat = {
        "func": "compare_csv",
        "result": {"type": "vm_file", "path": dst_combined, "dest": "result.csv"},
        "expected": {"type": "vm_file", "path": expected_combined, "dest": "expected.csv"},
        "options": {"strict": False},
    }
    instr_concat = [
        "In a terminal, concatenate ~/Desktop/econ/a.csv and ~/Desktop/econ/b.csv into ~/Desktop/combined.csv. Keep the header row only once: take the header from a.csv and skip the header row of b.csv.",
    ]

    def _params_concat(seed: int) -> dict:
        return {"instr": instr_concat[0], "config_override": init_steps_concat}

    out.append(SynthTemplate(
        template_id="shell_concat_csv",
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator_concat,
        oracle_fn=lambda _p: oracle_concat,
        postconfig_fn=lambda _p: None,
        param_fn=_params_concat,
        n_rows=2,
        eval_class="compare_text_file",
        setup_class="text_source",
    ))

    # Row 61 — awk skip-header + filter col 3 > 100
    dataset_lines = [
        "id,name,score",
        "1,alpha,80",
        "2,beta,150",
        "3,gamma,99",
        "4,delta,200",
        "5,epsilon,40",
    ]
    dataset_body = "\n".join(dataset_lines) + "\n"
    clean_lines = [l for l in dataset_lines[1:] if int(l.split(",")[2]) > 100]
    out.append(_make_shell_pipeline_template(
        template_id="shell_awk_skip_header",
        src_path=f"{_DESKTOP}/dataset.csv",
        dst_path=f"{_DESKTOP}/clean.csv",
        src_body=dataset_body,
        expected_body="\n".join(clean_lines) + "\n",
        pipeline=(
            f"awk -F, 'NR>1 && $3+0 > 100' {_DESKTOP}/dataset.csv > {_DESKTOP}/clean.csv"
        ),
        instructions=[
            "In a terminal, use awk on ~/Desktop/dataset.csv to skip the header and keep only rows whose third column (score) is greater than 100; write the result (no header) to ~/Desktop/clean.csv.",
        ],
        eval_class="compare_text_file",
        eval_func="compare_csv",
        expected_kind="csv",
    ))

    # Row 103 — grep -c 'ERROR'
    log_lines = [
        "2026-05-09 12:01:00 INFO startup",
        "2026-05-09 12:02:00 ERROR connection refused",
        "2026-05-09 12:03:00 INFO retry",
        "2026-05-09 12:04:00 ERROR timeout",
        "2026-05-09 12:05:00 ERROR disk full",
        "2026-05-09 12:06:00 INFO shutdown",
    ]
    log_body = "\n".join(log_lines) + "\n"
    err_count = sum(1 for l in log_lines if "ERROR" in l)
    out.append(_make_shell_pipeline_template(
        template_id="shell_grep_count_errors",
        src_path=f"{_DESKTOP}/raw.log",
        dst_path=f"{_DESKTOP}/error_count.txt",
        src_body=log_body,
        expected_body=f"{err_count}\n",
        pipeline=f"grep -c 'ERROR' {_DESKTOP}/raw.log > {_DESKTOP}/error_count.txt",
        instructions=[
            "From a terminal, count how many lines in ~/Desktop/raw.log contain the substring ERROR. Save just that integer (no extra text) followed by a newline to ~/Desktop/error_count.txt.",
        ],
    ))

    # Row 91 — head
    head_n = 5
    expected_head = "\n".join(_SHERLOCK_EXCERPT_LINES[:head_n]) + "\n"
    out.append(_make_shell_pipeline_template(
        template_id="shell_head_excerpt",
        src_path=f"{_DESKTOP}/full_excerpt.txt",
        dst_path=f"{_DESKTOP}/excerpt_head.txt",
        src_body=sherlock_body,
        expected_body=expected_head,
        pipeline=f"head -n {head_n} {_DESKTOP}/full_excerpt.txt > {_DESKTOP}/excerpt_head.txt",
        instructions=[
            f"In a terminal, save the first {head_n} lines of ~/Desktop/full_excerpt.txt to ~/Desktop/excerpt_head.txt (use `head -n {head_n}`).",
        ],
    ))

    # Row 79 — sort -f case-insensitive
    mixed_lines = ["Banana", "apple", "Cherry", "blueberry", "Apricot"]
    mixed_body = "\n".join(mixed_lines) + "\n"
    sorted_body = "\n".join(sorted(mixed_lines, key=str.casefold)) + "\n"
    out.append(_make_shell_pipeline_template(
        template_id="shell_sort_case_insensitive",
        src_path=f"{_DESKTOP}/words.txt",
        dst_path=f"{_DESKTOP}/words_sorted.txt",
        src_body=mixed_body,
        expected_body=sorted_body,
        pipeline=f"sort -f {_DESKTOP}/words.txt > {_DESKTOP}/words_sorted.txt",
        instructions=[
            "In a terminal, sort ~/Desktop/words.txt case-insensitively (use `sort -f`) and save the sorted lines to ~/Desktop/words_sorted.txt.",
        ],
    ))

    # sort -u (sort + uniq)
    dup_lines = ["alpha", "beta", "alpha", "gamma", "beta", "delta"]
    dup_body = "\n".join(dup_lines) + "\n"
    uniq_sorted = "\n".join(sorted(set(dup_lines))) + "\n"
    out.append(_make_shell_pipeline_template(
        template_id="shell_sort_uniq",
        src_path=f"{_DESKTOP}/with_dups.txt",
        dst_path=f"{_DESKTOP}/unique.txt",
        src_body=dup_body,
        expected_body=uniq_sorted,
        pipeline=f"sort -u {_DESKTOP}/with_dups.txt > {_DESKTOP}/unique.txt",
        instructions=[
            "In a terminal, deduplicate ~/Desktop/with_dups.txt by sorting and removing duplicates (use `sort -u`); write the unique lines to ~/Desktop/unique.txt.",
        ],
    ))

    # cut -d',' -f3
    csv_for_cut_lines = ["id,name,score", "1,alpha,80", "2,beta,150", "3,gamma,99"]
    csv_for_cut = "\n".join(csv_for_cut_lines) + "\n"
    cut_expected = "\n".join(l.split(",")[2] for l in csv_for_cut_lines) + "\n"
    out.append(_make_shell_pipeline_template(
        template_id="shell_cut_third_col",
        src_path=f"{_DESKTOP}/scores.csv",
        dst_path=f"{_DESKTOP}/scores_only.txt",
        src_body=csv_for_cut,
        expected_body=cut_expected,
        pipeline=f"cut -d',' -f3 {_DESKTOP}/scores.csv > {_DESKTOP}/scores_only.txt",
        instructions=[
            "In a terminal, use `cut -d',' -f3` to extract the third column from ~/Desktop/scores.csv and write it to ~/Desktop/scores_only.txt.",
        ],
    ))

    # tr lowercase
    upper_body = "HELLO WORLD\nFOO BAR\nBaz QUX\n"
    lower_body = upper_body.lower()
    out.append(_make_shell_pipeline_template(
        template_id="shell_tr_lower",
        src_path=f"{_DESKTOP}/upper.txt",
        dst_path=f"{_DESKTOP}/lower.txt",
        src_body=upper_body,
        expected_body=lower_body,
        pipeline=(
            f"tr '[:upper:]' '[:lower:]' < {_DESKTOP}/upper.txt > {_DESKTOP}/lower.txt"
        ),
        instructions=[
            "In a terminal, lowercase the contents of ~/Desktop/upper.txt with `tr` and write the result to ~/Desktop/lower.txt.",
        ],
    ))

    # wc -l (count only)
    wc_count = len(_SHERLOCK_EXCERPT_LINES)
    out.append(_make_shell_pipeline_template(
        template_id="shell_wc_lines",
        src_path=f"{_DESKTOP}/excerpt.txt",
        dst_path=f"{_DESKTOP}/line_count.txt",
        src_body=sherlock_body,
        expected_body=f"{wc_count}\n",
        # `wc -l < file` writes only the count, no filename suffix.
        pipeline=f"wc -l < {_DESKTOP}/excerpt.txt > {_DESKTOP}/line_count.txt",
        instructions=[
            "In a terminal, count the lines in ~/Desktop/excerpt.txt and save just the count (no filename) to ~/Desktop/line_count.txt. Hint: `wc -l < file > out`.",
        ],
    ))

    # sed substitute
    sed_src = (
        "Sun Tzu said: The art of war is of vital importance to the State.\n"
        "It is a matter of life and death.\n"
        "Hence it is a subject of inquiry which can on no account be neglected.\n"
    )
    sed_expected = sed_src.replace("Sun Tzu", "Sunzi")
    out.append(_make_shell_pipeline_template(
        template_id="shell_sed_substitute",
        src_path=f"{_DESKTOP}/quote.txt",
        dst_path=f"{_DESKTOP}/quote_fixed.txt",
        src_body=sed_src,
        expected_body=sed_expected,
        pipeline=(
            f"sed 's/Sun Tzu/Sunzi/g' {_DESKTOP}/quote.txt > {_DESKTOP}/quote_fixed.txt"
        ),
        instructions=[
            "In a terminal, use sed to replace every occurrence of 'Sun Tzu' with 'Sunzi' in ~/Desktop/quote.txt and write the result to ~/Desktop/quote_fixed.txt.",
        ],
    ))

    return out


# ---------------------------------------------------------------------------
# Cat F — File org via shell (extract + sort by ext)
# ---------------------------------------------------------------------------

def _make_extract_sort_template() -> SynthTemplate:
    """Row 30 — extract a zip, sort files into csv/png/docx subdirs."""
    template_id = "extract_sort_by_ext"
    archive = f"{_DESKTOP}/mixed_files.zip"
    extract_dir = f"{_DESKTOP}/extracted"

    build_zip_py = textwrap.dedent(f"""\
        import os, zipfile, shutil
        archive = {archive!r}
        os.makedirs(os.path.dirname(archive), exist_ok=True)
        if os.path.exists(archive):
            os.remove(archive)
        text_files = {{
            'sales-q1.csv': 'q,total\\nQ1,12000\\n',
            'sales-q2.csv': 'q,total\\nQ2,15000\\n',
            'notes.docx': 'PK\\x03\\x04' + '\\x00' * 64,
            'summary.docx': 'PK\\x03\\x04' + '\\x00' * 64,
        }}
        binary_files = {{
            'logo.png': b'\\x89PNG\\r\\n\\x1a\\n' + b'\\x00' * 32,
            'banner.png': b'\\x89PNG\\r\\n\\x1a\\n' + b'\\x00' * 32,
        }}
        with zipfile.ZipFile(archive, 'w') as zf:
            for name, body in text_files.items():
                zf.writestr(name, body)
            for name, body in binary_files.items():
                zf.writestr(name, body)
        ed = {extract_dir!r}
        shutil.rmtree(ed, ignore_errors=True)
        os.makedirs(ed, exist_ok=True)
        """)
    init_step = _python_step(build_zip_py)

    oracle_steps = [_execute(
        f"cd {extract_dir} && unzip -o {archive} && "
        f"mkdir -p csv png docx && "
        f"mv -f *.csv csv/ 2>/dev/null; "
        f"mv -f *.png png/ 2>/dev/null; "
        f"mv -f *.docx docx/ 2>/dev/null; true"
    )]
    probe_cmd = f"ls -R {extract_dir}"
    evaluator = {
        "func": "check_include_exclude",
        "result": {"type": "vm_command_line", "command": probe_cmd, "shell": True},
        "expected": {"type": "rule", "rules": {
            "include": [
                "csv:", "png:", "docx:",
                "sales-q1.csv", "sales-q2.csv",
                "logo.png", "banner.png",
                "notes.docx", "summary.docx",
            ],
            "exclude": [],
        }},
    }
    instructions = [
        "In a terminal, extract the archive ~/Desktop/mixed_files.zip into ~/Desktop/extracted/, then create subdirectories csv/, png/, and docx/ inside extracted/, and move each extracted file into the subdir matching its extension.",
        "Please help me unpack ~/Desktop/mixed_files.zip into ~/Desktop/extracted/, then make csv/, png/, and docx/ subfolders there and sort each extracted file into the subfolder matching its extension.",
        "I'd like ~/Desktop/mixed_files.zip extracted into ~/Desktop/extracted/ from a terminal — afterwards, create csv/, png/, and docx/ subdirectories under extracted/ and move every extracted file into the one matching its extension.",
        "From a terminal, please extract ~/Desktop/mixed_files.zip into ~/Desktop/extracted/. Then add three subdirectories named csv, png, and docx inside extracted/, and relocate each unpacked file into the subdirectory whose name matches the file's extension.",
        "Help me tidy up ~/Desktop/mixed_files.zip: in a terminal, extract it into ~/Desktop/extracted/, then create csv/, png/, and docx/ subdirs under extracted/ and move each extracted file into the subdir that matches its extension.",
    ]

    def _params(seed: int) -> dict:
        # validation gap-close: terminal preopen for shell-driven extract task.
        init_steps_full = [init_step, *_terminal_preopen_steps()]
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps_full}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="check_include_exclude",
        setup_class="archive_source",
    )


def _make_extract_targz_template() -> SynthTemplate:
    """Row 116 — extract a tar.gz; verify files appear in extracted/."""
    template_id = "extract_targz"
    archive = f"{_DESKTOP}/source.tar.gz"
    extract_dir = f"{_DESKTOP}/unpacked"

    build_targz_py = textwrap.dedent(f"""\
        import os, tarfile, io, shutil
        archive = {archive!r}
        os.makedirs(os.path.dirname(archive), exist_ok=True)
        if os.path.exists(archive):
            os.remove(archive)
        files = {{'a.txt': 'alpha\\n', 'b.txt': 'beta\\n', 'c.txt': 'gamma\\n'}}
        with tarfile.open(archive, 'w:gz') as tf:
            for name, body in files.items():
                data = body.encode()
                info = tarfile.TarInfo(name=name)
                info.size = len(data)
                tf.addfile(info, io.BytesIO(data))
        ed = {extract_dir!r}
        shutil.rmtree(ed, ignore_errors=True)
        os.makedirs(ed, exist_ok=True)
        """)
    init_step = _python_step(build_targz_py)
    oracle_steps = [_execute(f"tar -xzf {archive} -C {extract_dir}")]
    probe_cmd = f"ls {extract_dir}"
    evaluator = {
        "func": "check_include_exclude",
        "result": {"type": "vm_command_line", "command": probe_cmd, "shell": True},
        "expected": {"type": "rule", "rules": {
            "include": ["a.txt", "b.txt", "c.txt"],
            "exclude": [],
        }},
    }
    instructions = [
        "Open a terminal and extract everything inside the gzipped tarball ~/Desktop/source.tar.gz into the directory ~/Desktop/unpacked/ (the destination already exists).",
        "I have a tarball at ~/Desktop/source.tar.gz and need its files unpacked into ~/Desktop/unpacked/. Please do it from a terminal.",
        "From a terminal, please decompress ~/Desktop/source.tar.gz so that every archived file ends up inside ~/Desktop/unpacked/.",
        "Could you unpack the gzipped tar archive at ~/Desktop/source.tar.gz into the existing folder ~/Desktop/unpacked/? Use a terminal session, please.",
        "Use a terminal to expand the contents of ~/Desktop/source.tar.gz into the ~/Desktop/unpacked/ directory.",
    ]

    def _params(seed: int) -> dict:
        # validation gap-close: terminal preopen for shell-driven tar -xzf task.
        init_steps_full = [init_step, *_terminal_preopen_steps()]
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps_full}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_include_exclude",
        setup_class="archive_source",
    )


# ---------------------------------------------------------------------------
# Cat G — Archive create
# ---------------------------------------------------------------------------

def _make_zip_create_template() -> SynthTemplate:
    """Row 86 / 108 — bundle a directory of .txt as a zip. compare_archive."""
    template_id = "zip_create_directory"
    src_dir = f"{_DESKTOP}/articles"
    out_zip = f"{_DESKTOP}/articles.zip"
    expected_zip = f"{_TMP}/expected_articles.zip"

    files = {
        "a.txt": "alpha contents\n",
        "b.txt": "beta contents\n",
        "c.txt": "gamma contents\n",
    }
    build_files_py = textwrap.dedent(f"""\
        import os, shutil, zipfile
        src_dir = {src_dir!r}
        out_zip = {out_zip!r}
        expected_zip = {expected_zip!r}
        shutil.rmtree(src_dir, ignore_errors=True)
        os.makedirs(src_dir, exist_ok=True)
        files = {files!r}
        for name, body in files.items():
            with open(os.path.join(src_dir, name), 'w') as f:
                f.write(body)
        for p in (out_zip, expected_zip):
            if os.path.exists(p):
                os.remove(p)
        with zipfile.ZipFile(expected_zip, 'w', zipfile.ZIP_DEFLATED) as zf:
            for name in sorted(files.keys()):
                zf.write(os.path.join(src_dir, name), arcname=name)
        """)
    init_step = _python_step(build_files_py)

    # Oracle: zip from inside the dir so arcnames are bare basenames (matches gold)
    oracle_steps = [_execute(f"cd {src_dir} && zip -r {out_zip} . -i '*.txt'")]

    evaluator = {
        "func": "compare_archive",
        "result": {"type": "vm_file", "path": out_zip, "dest": "result.zip"},
        "expected": {"type": "vm_file", "path": expected_zip, "dest": "expected.zip"},
        "options": {"file_type": "text"},
    }
    instructions = [
        "From a terminal, bundle every .txt file inside ~/Desktop/articles/ into a single zip archive at ~/Desktop/articles.zip. The entries inside the zip should be bare filenames (a.txt, b.txt, ...) with no leading directory.",
        "I have a folder of articles at ~/Desktop/articles/ — please zip just its .txt files into ~/Desktop/articles.zip, making sure each archive entry is the bare filename (no 'articles/' prefix inside the zip).",
        "Could you produce a zip archive at ~/Desktop/articles.zip that contains all .txt files from ~/Desktop/articles/? Use a terminal, and ensure the archive entries are stored as plain filenames (no parent directory in the path).",
        "From a terminal, please pack every .txt file under ~/Desktop/articles/ into ~/Desktop/articles.zip so that the entries are stored as flat filenames at the root of the archive.",
        "I need a zip at ~/Desktop/articles.zip with all of the .txt files from ~/Desktop/articles/. Each entry inside the zip should be the bare basename (no folder prefix).",
    ]

    def _params(seed: int) -> dict:
        # validation gap-close: terminal preopen for shell-driven `zip -r` task.
        init_steps_full = [init_step, *_terminal_preopen_steps()]
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps_full}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_archive",
        setup_class="text_dir",
    )


# ---------------------------------------------------------------------------
# Cat H — Real-asset shell-pipeline templates (staged via _stage_asset)
# ---------------------------------------------------------------------------
# Pattern (mirrors gimp.py `_make_asset_filter_row`):
# 1. `_stage_asset(rel, dst)` host-pushes a real bundle asset onto the VM.
# 2. A python heredoc reads the staged source and writes the gold file (or,
#    for fixed-output rows, writes the gold directly).
# 3. Oracle runs the shell pipeline; eval = compare_text_file (or compare_csv).
# This unblocks the multi_apps.md "Real-asset" deferred rows by leaning on
# the assets/synth/ bundle now that `host_push` is wired in.


def _make_asset_shell_row(spec: dict) -> SynthTemplate:
    """Build a real-asset shell-pipeline row.

    Spec keys:
      template_id        : str
      asset_rels         : list[str] — paths under data/assets/synth/
      src_paths          : list[str] — corresponding container dst paths
      gold_builder_body  : str | None — python body that opens src(s) and writes
                                 `gold_path` (a variable bound for the body).
      gold_shell_cmd     : str | None — shell command to build the gold on
                                 the VM. Use `{expected_path}` placeholder
                                 for the gold output path. Mutually exclusive
                                 with `gold_builder_body`. Use this when the
                                 Python gold builder cannot precisely match
                                 the shell pipeline (locale-dependent sort
                                 order, CRLF preservation, etc.) — running
                                 the same shell pipeline as the oracle (with
                                 different output path) is byte-equivalent
                                 by construction.
      pipeline           : str — shell command run by oracle to produce dst.
      dst_path           : str — container path the agent writes to.
      instructions       : list[str]
      eval_func          : "compare_text_file" | "compare_csv" (default text)
      expected_kind      : "text" | "csv" (controls /tmp/expected ext + dest)
    """
    template_id   = spec["template_id"]
    asset_rels    = spec["asset_rels"]
    src_paths     = spec["src_paths"]
    gold_body     = spec.get("gold_builder_body")
    gold_shell_cmd = spec.get("gold_shell_cmd")
    pipeline      = spec["pipeline"]
    dst_path      = spec["dst_path"]
    instructions  = spec["instructions"]
    eval_func     = spec.get("eval_func", "compare_text_file")
    expected_kind = spec.get("expected_kind", "text")
    assert len(asset_rels) == len(src_paths)
    assert (gold_body is None) != (gold_shell_cmd is None), (
        f"{template_id}: must set exactly one of gold_builder_body / gold_shell_cmd"
    )

    ext = "csv" if expected_kind == "csv" else "txt"
    expected_path = f"{_TMP}/expected_{template_id}.{ext}"

    # Step 1: host_push each asset to its container path.
    init_steps: list[dict] = [
        _stage_asset(rel, src) for rel, src in zip(asset_rels, src_paths)
    ]
    # Step 1b: pre-open a terminal (validation H-cluster fix) — these rows all
    # instruct "In a terminal, ..." but lite.osworld desktop has no visible
    # terminal dock icon, so the agent's first type_text would race against a
    # blind dock-click guess. Pre-opening gnome-terminal removes the race.
    init_steps.extend(_terminal_preopen_steps())
    # Step 2: build the gold from the staged source(s).
    if gold_shell_cmd is not None:
        # Run the oracle-equivalent shell pipeline at config time so the
        # gold is byte-identical to what the agent's pipeline produces
        # (locale + CRLF + collation parity by construction).
        init_steps.append(_execute(gold_shell_cmd.format(expected_path=expected_path)))
    else:
        # Build at column 0 — no textwrap.dedent. The previous
        # `dedent(f"...{indent(body, '        ').lstrip()}...")` produced
        # broken Python: dedent's common-prefix was governed by the
        # least-indented body line, so outer scaffolding ended up at col 4
        # while body lines straddled cols 0 / 4 / 8, breaking `with` blocks.
        # gold_body already comes pre-indented (4-space `with` bodies); the
        # outer lines must stay at col 0 to keep the indent shape consistent.
        gold_py = (
            "import os\n"
            f"gold_path = {expected_path!r}\n"
            f"src_paths = {src_paths!r}\n"
            "os.makedirs(os.path.dirname(gold_path) or '.', exist_ok=True)\n"
            + gold_body
        )
        init_steps.append(_python_step(gold_py))

    # Oracle: run the shell pipeline (writes dst_path).
    oracle_steps = [_execute(pipeline)]

    if eval_func == "compare_csv":
        evaluator = {
            "func": "compare_csv",
            "result": {"type": "vm_file", "path": dst_path, "dest": "result.csv"},
            "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.csv"},
            "options": {"strict": False},
        }
    else:
        evaluator = {
            "func": "compare_text_file",
            "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
            "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
        }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,   # no app launch (terminal task)
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file" if eval_func == "compare_text_file" else "compare_table",
        setup_class="real_asset_text",
    )


_ASSET_SHELL_SPECS: list[dict] = [
    # Row #17 — top-3 largest gutenberg .txt by byte size (du -b | sort -n | tail -3).
    # Oracle uses GNU coreutils `du`/`sort`/`tail`; gold built by reading the
    # staged sizes via os.path.getsize and emitting the same `<size>\t<path>`
    # tab-separated lines that `du -b` produces.
    {
        "template_id": "asset_du_top3_gutenberg",
        "asset_rels": [
            "docs/gutenberg/alice-in-wonderland.txt",
            "docs/gutenberg/art-of-war.txt",
            "docs/gutenberg/frankenstein.txt",
            "docs/gutenberg/importance-of-being-earnest.txt",
            "docs/gutenberg/metamorphosis.txt",
            "docs/gutenberg/moby-dick.txt",
            "docs/gutenberg/pride-and-prejudice.txt",
            "docs/gutenberg/sherlock-adventures.txt",
            "docs/gutenberg/tale-of-two-cities.txt",
            "docs/gutenberg/tom-sawyer.txt",
            "docs/gutenberg/treasure-island.txt",
        ],
        "src_paths": [
            f"{_DESKTOP}/gutenberg/alice-in-wonderland.txt",
            f"{_DESKTOP}/gutenberg/art-of-war.txt",
            f"{_DESKTOP}/gutenberg/frankenstein.txt",
            f"{_DESKTOP}/gutenberg/importance-of-being-earnest.txt",
            f"{_DESKTOP}/gutenberg/metamorphosis.txt",
            f"{_DESKTOP}/gutenberg/moby-dick.txt",
            f"{_DESKTOP}/gutenberg/pride-and-prejudice.txt",
            f"{_DESKTOP}/gutenberg/sherlock-adventures.txt",
            f"{_DESKTOP}/gutenberg/tale-of-two-cities.txt",
            f"{_DESKTOP}/gutenberg/tom-sawyer.txt",
            f"{_DESKTOP}/gutenberg/treasure-island.txt",
        ],
        "gold_builder_body": (
            "sized = sorted((os.path.getsize(p), p) for p in src_paths)\n"
            "with open(gold_path, 'w') as f:\n"
            "    for sz, p in sized[-3:]:\n"
            "        f.write(f'{sz}\\t{p}\\n')\n"
        ),
        "pipeline": (
            f"du -b {_DESKTOP}/gutenberg/*.txt | sort -n | tail -3 "
            f"> {_DESKTOP}/top3_books.txt"
        ),
        "dst_path": f"{_DESKTOP}/top3_books.txt",
        "instructions": [
            "In a terminal, find the three largest .txt files inside ~/Desktop/gutenberg/ by byte size (smallest first, largest last) and save the result to ~/Desktop/top3_books.txt. Each output line should be the byte count, a tab, and the file path (the same format GNU du with byte counts produces).",
        ],
    },
    # Row #18 — awk first column (state name) of us-population-states.csv,
    # skip header, save first 5 names.
    {
        "template_id": "asset_awk_state_names",
        "asset_rels": ["data/csv/us-population-states.csv"],
        "src_paths": [f"{_DESKTOP}/us-population-states.csv"],
        "gold_builder_body": (
            "with open(src_paths[0]) as f:\n"
            "    lines = f.read().splitlines()\n"
            "names = [ln.split(',')[0] for ln in lines[1:6]]\n"
            "with open(gold_path, 'w') as f:\n"
            "    for n in names:\n"
            "        f.write(n + '\\n')\n"
        ),
        "pipeline": (
            f"awk -F, 'NR>1 && NR<=6 {{print $1}}' "
            f"{_DESKTOP}/us-population-states.csv "
            f"> {_DESKTOP}/first5_states.txt"
        ),
        "dst_path": f"{_DESKTOP}/first5_states.txt",
        "instructions": [
            "From a terminal, pull the state name (first column) from the first 5 data rows of ~/Desktop/us-population-states.csv, skipping the header row. Save the five names, one per line, to ~/Desktop/first5_states.txt.",
        ],
    },
    # Row #20 — sed delete every line containing 'ALGERNON' from the staged
    # play. Mirrors the inline os.py row but uses the real Wilde text so the
    # eval surface matches the bundle.
    {
        "template_id": "asset_sed_drop_algernon",
        "asset_rels": ["docs/gutenberg/importance-of-being-earnest.txt"],
        "src_paths": [f"{_DESKTOP}/importance-of-being-earnest.txt"],
        "gold_builder_body": (
            "with open(src_paths[0]) as f:\n"
            "    lines = f.readlines()\n"
            "with open(gold_path, 'w') as f:\n"
            "    for ln in lines:\n"
            "        if 'ALGERNON' not in ln:\n"
            "            f.write(ln)\n"
        ),
        "pipeline": (
            f"sed '/ALGERNON/d' {_DESKTOP}/importance-of-being-earnest.txt "
            f"> {_DESKTOP}/earnest_no_algernon.txt"
        ),
        "dst_path": f"{_DESKTOP}/earnest_no_algernon.txt",
        "instructions": [
            "From a terminal, drop every line that contains the word ALGERNON from ~/Desktop/importance-of-being-earnest.txt (leave all other lines untouched) and write the filtered text to ~/Desktop/earnest_no_algernon.txt.",
        ],
    },
    # Row #50 — concat 4 FRED CSVs (us-fed-funds-rate / us-mortgage-30yr /
    # us-housing-starts / oil-wti-daily) into one all_indicators.csv.
    # The four files have DIFFERENT schemas (different col-2 names), so a
    # plain concat with a single canonical header would mis-label rows.
    # Instead the gold prepends a 3-col header (date,series,value) and emits
    # `<date>,<series>,<value>` rows from each source. The oracle pipeline
    # mirrors this with awk so eval matches.
    {
        "template_id": "asset_concat_fred_csvs",
        "asset_rels": [
            "data/csv/us-fed-funds-rate.csv",
            "data/csv/us-mortgage-30yr.csv",
            "data/csv/us-housing-starts.csv",
            "data/csv/oil-wti-daily.csv",
        ],
        "src_paths": [
            f"{_DESKTOP}/fred/us-fed-funds-rate.csv",
            f"{_DESKTOP}/fred/us-mortgage-30yr.csv",
            f"{_DESKTOP}/fred/us-housing-starts.csv",
            f"{_DESKTOP}/fred/oil-wti-daily.csv",
        ],
        "gold_builder_body": (
            "series_names = ['FEDFUNDS', 'MORTGAGE30US', 'HOUST', 'DCOILWTICO']\n"
            "with open(gold_path, 'w') as out:\n"
            "    out.write('date,series,value\\n')\n"
            "    for src, name in zip(src_paths, series_names):\n"
            "        with open(src) as f:\n"
            "            lines = f.read().splitlines()\n"
            "        for ln in lines[1:]:\n"
            "            date, val = ln.split(',', 1)\n"
            "            out.write(f'{date},{name},{val}\\n')\n"
        ),
        "pipeline": (
            f"{{ echo 'date,series,value'; "
            f"awk -F, 'NR>1 {{print $1\",FEDFUNDS,\"$2}}' "
            f"{_DESKTOP}/fred/us-fed-funds-rate.csv; "
            f"awk -F, 'NR>1 {{print $1\",MORTGAGE30US,\"$2}}' "
            f"{_DESKTOP}/fred/us-mortgage-30yr.csv; "
            f"awk -F, 'NR>1 {{print $1\",HOUST,\"$2}}' "
            f"{_DESKTOP}/fred/us-housing-starts.csv; "
            f"awk -F, 'NR>1 {{print $1\",DCOILWTICO,\"$2}}' "
            f"{_DESKTOP}/fred/oil-wti-daily.csv; "
            f"}} > {_DESKTOP}/all_indicators.csv"
        ),
        "dst_path": f"{_DESKTOP}/all_indicators.csv",
        "instructions": [
            "In a terminal, combine the four FRED CSVs in ~/Desktop/fred/ (us-fed-funds-rate.csv, us-mortgage-30yr.csv, us-housing-starts.csv, oil-wti-daily.csv) into a single ~/Desktop/all_indicators.csv with header `date,series,value`. For each input file, drop its header and emit `<date>,<series_name>,<value>` rows where series_name is FEDFUNDS, MORTGAGE30US, HOUST, and DCOILWTICO respectively (in that order).",
        ],
        "eval_func": "compare_csv",
        "expected_kind": "csv",
    },
    # Row #69 — sort + uniq on the ART OF WAR text, line-level.
    # validation (2026-05-10): the original `gold_builder_body` rebuilt the gold
    # in Python via `sorted(set(lines))` after `ln.rstrip('\\n')`, but the
    # source is CRLF and Python's universal-newline mode strips '\\r' on
    # read while shell `sort -u` operates on raw bytes (preserves CRLF).
    # The two outputs differed at byte 1 (CRLF vs LF). Fix: build the gold
    # by running the same shell pipeline as the oracle on the VM. Pipeline
    # also pinned to LC_ALL=C for byte-deterministic collation.
    {
        "template_id": "asset_sort_uniq_artofwar",
        "asset_rels": ["docs/gutenberg/art-of-war.txt"],
        "src_paths": [f"{_DESKTOP}/art-of-war.txt"],
        "gold_shell_cmd": (
            f"LC_ALL=C sort -u {_DESKTOP}/art-of-war.txt "
            f"> {{expected_path}}"
        ),
        "pipeline": (
            f"LC_ALL=C sort -u {_DESKTOP}/art-of-war.txt "
            f"> {_DESKTOP}/art-of-war-uniq.txt"
        ),
        "dst_path": f"{_DESKTOP}/art-of-war-uniq.txt",
        "instructions": [
            "In a terminal, sort the lines of ~/Desktop/art-of-war.txt and remove duplicates with `LC_ALL=C sort -u`. Save the deduplicated, sorted result to ~/Desktop/art-of-war-uniq.txt.",
        ],
    },
    # Row #79 — sort -f case-insensitive on Wilde play.
    # validation (2026-05-10): the original Python `key=str.casefold` differs
    # from shell `sort -f` on bytes like '[' (0x5b) vs 'I' (0x49) — `sort -f`
    # uppercases for comparison so 'I' < '[', while `str.casefold` lowercases
    # so '[' < 'i'. Mismatch at line 977. Fix: build the gold by running the
    # same shell pipeline (LC_ALL=C-pinned for determinism) as the oracle.
    {
        "template_id": "asset_sort_case_earnest",
        "asset_rels": ["docs/gutenberg/importance-of-being-earnest.txt"],
        "src_paths": [f"{_DESKTOP}/importance-of-being-earnest.txt"],
        "gold_shell_cmd": (
            f"LC_ALL=C sort -f {_DESKTOP}/importance-of-being-earnest.txt "
            f"> {{expected_path}}"
        ),
        "pipeline": (
            f"LC_ALL=C sort -f {_DESKTOP}/importance-of-being-earnest.txt "
            f"> {_DESKTOP}/earnest-sorted.txt"
        ),
        "dst_path": f"{_DESKTOP}/earnest-sorted.txt",
        "instructions": [
            "In a terminal, sort ~/Desktop/importance-of-being-earnest.txt case-insensitively (use `LC_ALL=C sort -f`) and save the sorted lines to ~/Desktop/earnest-sorted.txt.",
        ],
    },
    # Row #91 — `shuf -n 10` on oil-wti-daily.csv. NOTE: shuf is randomized,
    # so the gold can't replicate the agent's output line-by-line. Instead we
    # use a DETERMINISTIC selection (first 10 data rows) for both the gold
    # AND the oracle. The agent's instruction directs `shuf -n 10`-style
    # selection, but eval grades the deterministic baseline so we keep the
    # contract; rewrite the pipeline to use `head -n 11 | tail -n 10`.
    {
        "template_id": "asset_shuf_oil_wti",
        "asset_rels": ["data/csv/oil-wti-daily.csv"],
        "src_paths": [f"{_DESKTOP}/oil-wti-daily.csv"],
        "gold_builder_body": (
            "with open(src_paths[0]) as f:\n"
            "    lines = f.read().splitlines()\n"
            "selected = lines[1:11]\n"
            "with open(gold_path, 'w') as f:\n"
            "    for ln in selected:\n"
            "        f.write(ln + '\\n')\n"
        ),
        "pipeline": (
            f"tail -n +2 {_DESKTOP}/oil-wti-daily.csv | head -n 10 "
            f"> {_DESKTOP}/oil-wti-sample.csv"
        ),
        "dst_path": f"{_DESKTOP}/oil-wti-sample.csv",
        "instructions": [
            "From a terminal, pull the first 10 data rows out of ~/Desktop/oil-wti-daily.csv (skip the header row entirely) and save just those 10 rows to ~/Desktop/oil-wti-sample.csv.",
        ],
    },
    # Row #115 — head -n 51 of oil-wti-daily.csv (first 50 data rows + header).
    {
        "template_id": "asset_head_oil_wti",
        "asset_rels": ["data/csv/oil-wti-daily.csv"],
        "src_paths": [f"{_DESKTOP}/oil-wti-daily.csv"],
        "gold_builder_body": (
            "with open(src_paths[0]) as f:\n"
            "    lines = f.read().splitlines()\n"
            "with open(gold_path, 'w') as f:\n"
            "    for ln in lines[:51]:\n"
            "        f.write(ln + '\\n')\n"
        ),
        "pipeline": (
            f"head -n 51 {_DESKTOP}/oil-wti-daily.csv "
            f"> {_DESKTOP}/oil-wti-head.csv"
        ),
        "dst_path": f"{_DESKTOP}/oil-wti-head.csv",
        "instructions": [
            "In a terminal, save the first 51 lines of ~/Desktop/oil-wti-daily.csv (header + first 50 data rows) to ~/Desktop/oil-wti-head.csv. Use `head -n 51`.",
        ],
    },
    # Row #128 — filter G7 countries from world-gdp-2022.csv. Gold is the
    # 7 lines whose CountryCode is in {CAN,FRA,DEU,ITA,JPN,GBR,USA}.
    {
        "template_id": "asset_grep_g7_world_gdp",
        "asset_rels": ["data/csv/world-gdp-2022.csv"],
        "src_paths": [f"{_DESKTOP}/world-gdp-2022.csv"],
        "gold_builder_body": (
            "g7 = {'CAN', 'FRA', 'DEU', 'ITA', 'JPN', 'GBR', 'USA'}\n"
            "with open(src_paths[0]) as f:\n"
            "    lines = f.read().splitlines()\n"
            "matches = [ln for ln in lines[1:] if ln.split(',')[0] in g7]\n"
            "with open(gold_path, 'w') as f:\n"
            "    for ln in matches:\n"
            "        f.write(ln + '\\n')\n"
        ),
        "pipeline": (
            f"grep -E '^(CAN|FRA|DEU|ITA|JPN|GBR|USA),' "
            f"{_DESKTOP}/world-gdp-2022.csv "
            f"> {_DESKTOP}/g7-gdp.csv"
        ),
        "dst_path": f"{_DESKTOP}/g7-gdp.csv",
        "instructions": [
            "From a terminal, filter ~/Desktop/world-gdp-2022.csv down to only the G7 country rows — those whose CountryCode (first column) is one of CAN, FRA, DEU, ITA, JPN, GBR, or USA. Save the matching rows in their original order to ~/Desktop/g7-gdp.csv (no header).",
        ],
    },
]


def _make_asset_csv_to_xlsx_template() -> SynthTemplate:
    """Row #34 — stage us-state-median-income.csv, agent opens in Calc and
    saves as xlsx. Gold xlsx built by reading the staged csv via openpyxl
    (then LO-normalized so cached <v> values match LO's save format).
    Eval: compare_table sheet_data over Sheet1.
    """
    template_id = "asset_csv_to_xlsx_state_income"
    asset_rel = "data/csv/us-state-median-income.csv"
    src = f"{_DESKTOP}/us-state-median-income.csv"
    dst = f"{_DESKTOP}/us-state-median-income.xlsx"
    expected = f"{_TMP}/expected_{template_id}.xlsx"

    init_src = _stage_asset(asset_rel, src)
    gold_py = textwrap.dedent(f"""\
        import os, csv
        from openpyxl import Workbook
        src = {src!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        wb = Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        with open(src, newline='') as f:
            reader = csv.reader(f)
            for row in reader:
                ws.append(row)
        wb.save(path)
        # _LO_NORMALIZE_TAIL: re-save via soffice headless so cached values
        # match LO save format (per AGENTS.md F1 mitigation).
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, path],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename(path))
            if os.path.exists(_conv):
                _sh.copy(_conv, path)
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)
    init_gold = _python_step(gold_py)
    init_steps = [
        init_src,
        init_gold,
        # validation gap-close: agent opens the staged .csv in Calc and
        # "save-as" to .xlsx — focused window is the source .csv.
        *_multi_app_preopen(foreground=src),
    ]
    oracle_steps = [_execute(f"cp '{expected}' '{dst}'")]

    evaluator = {
        "func": "compare_table",
        "result": {"type": "vm_file", "path": dst, "dest": "result.xlsx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.xlsx"},
        # Validation fix: positional sheet index (LO Calc Save-As preserves CSV
        # basename as sheet name, not "Sheet1"). See _make_csv_to_xlsx_template.
        "options": {"rules": [{"type": "sheet_data", "sheet_idx0": 0, "sheet_idx1": 0}]},
    }

    instructions = [
        "Open ~/Desktop/us-state-median-income.csv in Calc and save it as ~/Desktop/us-state-median-income.xlsx (Excel xlsx format).",
        "Convert ~/Desktop/us-state-median-income.csv into an xlsx workbook at ~/Desktop/us-state-median-income.xlsx using LibreOffice Calc.",
    ]

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", src],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_table",
        setup_class="real_asset_csv",
    )


# ---------------------------------------------------------------------------
# Cat I — Real-asset CROSS-APP rows (CSV → docx-table, photo → docx,
#         photo → pptx, CSV → xlsx-with-chart, multi-CSV concat,
#         Wikipedia HTML extract). Mirrors `_make_asset_shell_row`.
# ---------------------------------------------------------------------------
# Pattern: stage real asset(s) via `_stage_asset`, then a python heredoc
# reads the staged source(s) and builds source + gold artefacts. LO-sink
# rows (docx/xlsx/pptx) set `oracle_after_postconfig=True` + use
# `LO_SAVE_POSTCONFIG`; non-LO rows (compare_text_file / compare_csv on shell
# output) keep the terminal-task config_override pattern.
# ---------------------------------------------------------------------------


# ---- Cat I.1 — CSV → docx-table -------------------------------------------
# Source docx has a "TODO: insert table here" placeholder paragraph; gold
# docx replaces that placeholder with a python-docx table built from a
# top-N projection of the staged CSV. Eval: compare_docx_tables.

_CSV_TO_DOCX_TABLE_SPECS: list[dict] = [
    {
        "template_id": "multi_csv_to_docx_table_us_state_income",
        "asset_rel": "data/csv/us-state-median-income.csv",
        "csv_basename": "us-state-median-income.csv",
        "docx_basename": "state_income_report.docx",
        # Source CSV header is State,MedianHouseholdIncome_USD,StateFIPS (3 cols).
        # Take all 51 rows; emit a 3-col table (header + 51 data rows).
        "select_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "header, data = rows[0], rows[1:]\n"
            "table_data = [header] + data\n"
        ),
        "docx_pre": [
            ("Title", "U.S. State Median Household Income"),
            (None, "The table below summarizes the per-state median household income."),
        ],
        # Validation fix: instruction said "header + 51 data rows" but the
        # actual asset has 52 data rows (53 lines incl. header). Use the
        # safer ambiguity-free phrasing "all data rows".
        "instructions": [
            "Transfer the data from ~/Desktop/us-state-median-income.csv (open in Calc) into the docx table at ~/Desktop/state_income_report.docx (you're currently editing in Writer). Insert the CSV header row plus every data row as a 3-column table at the end of the document, then save the document.",
        ],
    },
    {
        "template_id": "multi_csv_to_docx_table_world_gdp_top10",
        "asset_rel": "data/csv/world-gdp-2022.csv",
        "csv_basename": "world-gdp-2022.csv",
        "docx_basename": "gdp_summary.docx",
        # Top 10 countries by GDP_USD (descending). Note: world-gdp-2022.csv
        # contains both regional aggregates AND individual country rows; we
        # take the top 10 by raw GDP across the whole file (deterministic).
        "select_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "header, data = rows[0], rows[1:]\n"
            "data.sort(key=lambda r: float(r[3]), reverse=True)\n"
            "table_data = [header] + data[:10]\n"
        ),
        "docx_pre": [
            ("Title", "GDP Summary"),
            (None, "The table below lists the top 10 entities by 2022 GDP (USD)."),
        ],
        "instructions": [
            "Transfer the data from ~/Desktop/world-gdp-2022.csv (open in Calc) into the docx you're currently editing at ~/Desktop/gdp_summary.docx (in Writer). Sort the CSV by GDP_USD descending and insert the header + top 10 rows as a 4-column table at the end of the document. Save the document.",
        ],
    },
    {
        "template_id": "multi_csv_to_docx_table_us_population_top5",
        "asset_rel": "data/csv/us-population-states.csv",
        "csv_basename": "us-population-states.csv",
        "docx_basename": "population_brief.docx",
        # Top 5 states by Population2020 (descending).
        "select_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "header, data = rows[0], rows[1:]\n"
            "data.sort(key=lambda r: int(r[1]), reverse=True)\n"
            "table_data = [header] + data[:5]\n"
        ),
        "docx_pre": [
            ("Title", "Top 5 U.S. States by Population (2020)"),
            (None, "The table below shows the five most populous states from the 2020 census."),
        ],
        "instructions": [
            "Transfer the data from ~/Desktop/us-population-states.csv (open in Calc) into the docx you're currently editing at ~/Desktop/population_brief.docx (in Writer). Sort the CSV by Population2020 descending and insert the header + top 5 rows as a 3-column table at the end of the document. Save the document.",
        ],
    },
    {
        "template_id": "multi_csv_to_docx_table_unemployment_recent12m",
        "asset_rel": "data/csv/us-unemployment.csv",
        "csv_basename": "us-unemployment.csv",
        "docx_basename": "unemployment_recent.docx",
        # Last 12 data rows (most recent 12 months).
        "select_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "header, data = rows[0], rows[1:]\n"
            "table_data = [header] + data[-12:]\n"
        ),
        "docx_pre": [
            ("Title", "U.S. Unemployment Rate — Last 12 Months"),
            (None, "The table below shows the unemployment rate for the most recent 12 monthly observations."),
        ],
        "instructions": [
            "Transfer the data from ~/Desktop/us-unemployment.csv (open in Calc) into the docx you're currently editing at ~/Desktop/unemployment_recent.docx (in Writer). Insert the header + the LAST 12 data rows of the CSV as a 2-column table at the end of the document. Save the document.",
        ],
    },
    # Mass-add — additional CSV→docx-table real-asset variants.
    {
        "template_id": "multi_csv_to_docx_table_oil_wti_recent10",
        "asset_rel": "data/csv/oil-wti-daily.csv",
        "csv_basename": "oil-wti-daily.csv",
        "docx_basename": "oil_wti_recent.docx",
        "select_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "header, data = rows[0], rows[1:]\n"
            "table_data = [header] + data[-10:]\n"
        ),
        "docx_pre": [
            ("Title", "WTI Crude Oil — 10 Most Recent Days"),
            (None, "The table below lists the 10 latest daily WTI crude observations."),
        ],
        "instructions": [
            "Transfer the data from ~/Desktop/oil-wti-daily.csv (open in Calc) into the docx you're currently editing at ~/Desktop/oil_wti_recent.docx (in Writer). Insert the header + the LAST 10 data rows as a 2-column table at the end of the document. Save it.",
        ],
    },
    {
        "template_id": "multi_csv_to_docx_table_us_gdp_recent8",
        "asset_rel": "data/csv/us-gdp.csv",
        "csv_basename": "us-gdp.csv",
        "docx_basename": "us_gdp_recent.docx",
        "select_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "header, data = rows[0], rows[1:]\n"
            "table_data = [header] + data[-8:]\n"
        ),
        "docx_pre": [
            ("Title", "U.S. GDP — 8 Most Recent Quarters"),
            (None, "The table below shows the latest eight quarterly GDP observations."),
        ],
        "instructions": [
            "Transfer the data from ~/Desktop/us-gdp.csv (open in Calc) into the docx you're currently editing at ~/Desktop/us_gdp_recent.docx (in Writer). Insert the header + the LAST 8 data rows as a 2-column table at the end of the document. Save it.",
        ],
    },
    # Pruned (multi_csv_to_docx_table_world_population_top10): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_table_world_population_top10",
#         "asset_rel": "data/csv/world-population-2022.csv",
#         "csv_basename": "world-population-2022.csv",
#         "docx_basename": "world_population_top10.docx",
#         # Header: CountryCode,CountryName,Year,Population (col 3 numeric).
#         "select_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "header, data = rows[0], rows[1:]\n"
#             "data.sort(key=lambda r: float(r[3]), reverse=True)\n"
#             "table_data = [header] + data[:10]\n"
#         ),
#         "docx_pre": [
#             ("Title", "World Population — Top 10 (2022)"),
#             (None, "The table below ranks the top 10 entities by 2022 population (mixes countries and regional aggregates)."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/world-population-2022.csv and ~/Desktop/world_population_top10.docx. Sort the CSV by Population descending and insert the header + top 10 rows as a 4-column table at the end. Save it.",
#         ],
#     },
    # Pruned (multi_csv_to_docx_table_us_inflation_last24m): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_table_us_inflation_last24m",
#         "asset_rel": "data/csv/us-inflation-cpi.csv",
#         "csv_basename": "us-inflation-cpi.csv",
#         "docx_basename": "us_inflation_recent.docx",
#         "select_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "header, data = rows[0], rows[1:]\n"
#             "table_data = [header] + data[-24:]\n"
#         ),
#         "docx_pre": [
#             ("Title", "U.S. Inflation (CPI) — Last 24 Months"),
#             (None, "The table below shows the most recent 24 monthly CPI observations."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-inflation-cpi.csv and ~/Desktop/us_inflation_recent.docx. Insert the header + LAST 24 data rows as a 2-column table at the end of the document. Save it.",
#         ],
#     },
    # Pruned (multi_csv_to_docx_table_us_mortgage_recent12): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_table_us_mortgage_recent12",
#         "asset_rel": "data/csv/us-mortgage-30yr.csv",
#         "csv_basename": "us-mortgage-30yr.csv",
#         "docx_basename": "us_mortgage_recent.docx",
#         "select_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "header, data = rows[0], rows[1:]\n"
#             "table_data = [header] + data[-12:]\n"
#         ),
#         "docx_pre": [
#             ("Title", "U.S. 30-Year Mortgage Rate — Last 12 Weeks"),
#             (None, "The table below shows the 12 most recent weekly mortgage-rate observations."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-mortgage-30yr.csv and ~/Desktop/us_mortgage_recent.docx. Insert the header + LAST 12 data rows as a 2-column table at the end of the document. Save it.",
#         ],
#     },
    # Pruned (multi_csv_to_docx_table_us_fedfunds_recent24m): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_table_us_fedfunds_recent24m",
#         "asset_rel": "data/csv/us-fed-funds-rate.csv",
#         "csv_basename": "us-fed-funds-rate.csv",
#         "docx_basename": "us_fedfunds_recent.docx",
#         "select_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "header, data = rows[0], rows[1:]\n"
#             "table_data = [header] + data[-24:]\n"
#         ),
#         "docx_pre": [
#             ("Title", "U.S. Fed Funds Rate — Last 24 Months"),
#             (None, "The table below shows the 24 most recent monthly fed-funds observations."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-fed-funds-rate.csv and ~/Desktop/us_fedfunds_recent.docx. Insert the header + LAST 24 data rows as a 2-column table at the end of the document. Save it.",
#         ],
#     },
    # Pruned (multi_csv_to_docx_table_us_housing_recent24m): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_table_us_housing_recent24m",
#         "asset_rel": "data/csv/us-housing-starts.csv",
#         "csv_basename": "us-housing-starts.csv",
#         "docx_basename": "us_housing_recent.docx",
#         "select_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "header, data = rows[0], rows[1:]\n"
#             "table_data = [header] + data[-24:]\n"
#         ),
#         "docx_pre": [
#             ("Title", "U.S. Housing Starts — Last 24 Months"),
#             (None, "The table below shows the 24 most recent monthly housing-starts observations."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-housing-starts.csv and ~/Desktop/us_housing_recent.docx. Insert the header + LAST 24 data rows as a 2-column table at the end of the document. Save it.",
#         ],
#     },
]


def _make_csv_to_docx_table_template(spec: dict) -> SynthTemplate:
    template_id  = spec["template_id"]
    asset_rel    = spec["asset_rel"]
    csv_path     = f"{_DESKTOP}/{spec['csv_basename']}"
    docx_path    = f"{_DESKTOP}/{spec['docx_basename']}"
    expected     = f"{_TMP}/expected_{template_id}.docx"
    select_py    = spec["select_py"]
    docx_pre     = spec["docx_pre"]
    instructions = spec["instructions"]

    init_csv = _stage_asset(asset_rel, csv_path)
    init_docx = _build_docx_step(docx_path, docx_pre)

    # Build gold by reading staged CSV, projecting via `select_py`, then
    # writing a docx with `docx_pre` paragraphs followed by the table.
    gold_py = textwrap.dedent(f"""\
        import os
        from docx import Document
        src = {csv_path!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        {textwrap.indent(select_py, '        ').lstrip()}
        doc = Document()
        for style, text in {docx_pre!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        if table_data:
            tbl = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
            for i, row in enumerate(table_data):
                for j, cell in enumerate(row):
                    tbl.cell(i, j).text = str(cell)
        doc.save(path)
        """)
    init_steps = [
        init_csv, init_docx, _python_step(gold_py),
        # validation: [csv (Calc bg)] → docx (Writer fg).
        # Bypasses `open_command` auto-launch (no sleep+activate for multi_apps).
        *_multi_app_preopen(background=[csv_path], foreground=docx_path),
    ]
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    evaluator = {
        "func": "compare_docx_tables",
        "result": {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.docx"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_tables",
        setup_class="real_asset_csv_docx",
    )


# ---- Cat I.2 — Photo → docx with inline image -----------------------------
# Mirror of writer's `_make_image_insert_row`: stage a photo, build a source
# docx with a caption paragraph (and a placeholder body), build gold by
# inserting the staged image after the caption. Eval: compare_docx_strict
# with examine_images=True (image-byte-hash multiset equality).

# ---------------------------------------------------------------------------
# Topic-family × skill-template Cartesian-product infrastructure.
#
# Same shape as `libreoffice_impress.TOPIC_FAMILIES` (kept in this module
# so multi_apps stays self-contained — duplication is intentional for
# locality; downstream skill factories pick photos via `_pick_photo_for_seed`
# / `_pick_photos_for_seed`). Each pool MUST have ≥4 distinct entries so
# 4-slide galleries don't dup. Asset paths verified at codegen time via
# `_stage_asset` which asserts existence.
# ---------------------------------------------------------------------------

# (asset_rel, slide_title, caption_text)
_PhotoSpec = tuple[str, str, str]

TOPIC_FAMILIES: dict[str, list[_PhotoSpec]] = {
    "space": [
        ("photos/space/jupiter-full-disk.jpg",         "Jupiter — Full Disk",        "Jupiter, Voyager full-disk view"),
        ("photos/space/mars-curiosity-panorama.jpg",   "Mars — Curiosity Panorama",  "Mars Curiosity panorama, Gale Crater"),
        ("photos/space/saturn-moon-titan.jpg",         "Saturn's Moon — Titan",      "Titan, Saturn's largest moon"),
        ("photos/space/europa-ice-surface.jpg",        "Europa — Ice Surface",       "Europa's icy surface, Galileo"),
        ("photos/space/earth-blue-marble-apollo17.jpg","Earth — Blue Marble",        "Earth, Apollo 17 'Blue Marble'"),
        ("photos/space/io-volcanic-eruption.jpg",      "Io — Volcanic Eruption",     "Io, volcanic eruption captured by Galileo"),
        ("photos/space/earthrise-apollo8.jpg",         "Earthrise",                  "Earthrise over lunar horizon, Apollo 8 (1968)"),
        ("photos/space/mars-rover-vista.jpg",          "Mars — Rover Vista",         "Mars rover vista, NASA"),
    ],
    "astronomy": [
        ("photos/nature/galaxy-andromeda.jpg",             "Andromeda Galaxy",           "M31 Andromeda, Hubble composite"),
        ("photos/nature/galaxy-hubble.jpg",                "Hubble Deep Field",          "Hubble deep-field galaxies"),
        ("photos/education/nebula-pillars-of-creation.jpg","Pillars of Creation",        "Eagle Nebula M16, Hubble"),
        ("photos/education/sun-solar-flare.jpg",           "Solar Flare",                "Sun, X-class solar flare"),
    ],
    "wildlife": [
        ("photos/wildlife/tiger-closeup.jpg",   "Tiger Close-up",  "Bengal tiger close-up portrait"),
        ("photos/wildlife/horse-meadow.jpg",    "Horse in Meadow", "Adult horse grazing in a meadow"),
        ("photos/wildlife/bird-perch.jpg",      "Perched Bird",    "Bird perched, observed at sunrise"),
        ("photos/wildlife/fox-portrait.jpg",    "Red Fox Portrait", "Red fox portrait, naturalist series"),
    ],
    "food": [
        ("photos/food/pizza-dish.jpg",      "Wood-Fired Pizza",    "Margherita pizza, Naples"),
        ("photos/food/coffee-latte.jpg",    "Latte Coffee",        "Cafe latte with rosetta art"),
        ("photos/food/restaurant-meal.jpg", "Restaurant Breakfast","Restaurant breakfast plate"),
        ("photos/food/salad-bowl.jpg",      "Garden Salad Bowl",   "Mixed garden salad, vinaigrette"),
    ],
    "architecture": [
        ("photos/architecture/house-modern.jpg",      "Modern House",       "Modern house, exterior late-afternoon"),
        ("photos/architecture/skyscraper-glass.jpg",  "Glass Skyscraper",   "Glass-clad tower, downtown"),
        ("photos/architecture/nasa-facility.jpg",     "NASA Facility",      "NASA mission-support facility"),
        ("photos/architecture/office-building.jpg",   "Office Building",    "Mid-rise commercial office building"),
        ("photos/architecture/living-room.jpg",       "Modern Living Room", "Contemporary living-room interior"),
    ],
    "portrait": [
        ("photos/portrait/astronaut-apollo11-crew.jpg", "Apollo 11 Crew",   "Apollo 11 prime crew portrait, NASA"),
        ("photos/portrait/person-headshot-1.jpg",       "Headshot — One",   "Professional portrait headshot"),
        ("photos/portrait/person-headshot-2.jpg",       "Headshot — Two",   "Professional portrait headshot"),
        ("photos/portrait/person-headshot-3.jpg",       "Headshot — Three", "Professional portrait headshot"),
    ],
    # Additional iteration — wire previously-unused photo categories.
    # Mirrors impress.TOPIC_FAMILIES additions; multi_apps photo→docx /
    # photo→pptx Cartesian factories pick up these new families automatically.
    "city": [
        ("photos/city/asian-skyline.jpg",     "Asian Skyline",        "Night skyline of an Asian metropolis"),
        ("photos/city/european-street.jpg",   "European Street",      "Cobblestone street in a European old town"),
        ("photos/city/street-night.jpg",      "Night Street Scene",   "Illuminated city street after sunset"),
    ],
    "office": [
        ("photos/office/laptop-desk.jpg",     "Workspace Setup",      "Tidy desk workspace with laptop"),
        ("photos/office/server-room.jpg",     "Server Room",          "Datacenter server room aisle"),
        ("photos/office/team-meeting.jpg",    "Team Meeting",         "Professional team meeting in a conference room"),
    ],
    "event": [
        ("photos/event/birthday-table.jpg",   "Birthday Table",       "Birthday party table with cake and candles"),
        ("photos/event/concert-stage.jpg",    "Concert Stage",        "Concert stage under colored lighting"),
        ("photos/event/wedding-decor.jpg",    "Wedding Decor",        "Wedding reception with floral table decor"),
    ],
    "product": [
        ("photos/product/headphones.jpg",     "Headphones",           "Studio headphones product shot"),
        ("photos/product/sneakers.jpg",       "Sneakers",             "Athletic sneakers product shot"),
        ("photos/product/wristwatch.jpg",     "Wristwatch",           "Mechanical wristwatch product shot"),
    ],
    "classroom": [
        ("photos/classroom/desk-and-board.jpg", "Classroom Desk",     "Classroom desk with chalkboard backdrop"),
        ("photos/classroom/library-shelves.jpg","Library Shelves",    "Library reading room with stacked shelves"),
        ("photos/classroom/open-books.jpg",     "Open Books",         "Open books on a study table"),
    ],
    "medical": [
        ("photos/medical/lab-bench.jpg",      "Laboratory Bench",     "Laboratory bench with research equipment"),
        ("photos/medical/stethoscope.jpg",    "Stethoscope",          "Stethoscope on a clinician's clipboard"),
    ],
    "education": [
        ("photos/education/chalkboard-classroom.jpg",       "Math Class on the Board", "Teacher writing equations on the blackboard"),
        ("photos/education/graduation-ceremony.jpg",        "Graduation Day",          "Graduates tossing caps at commencement"),
        ("photos/education/students-studying.jpg",          "Reading Hour",            "Schoolchildren reading together in study time"),
        ("photos/education/nebula-pillars-of-creation.jpg", "Pillars of Creation",     "Hubble Eagle Nebula M16 — astronomy lesson"),
        ("photos/education/sun-solar-flare.jpg",            "Solar Flare",             "SDO mid-level solar flare — astronomy lesson"),
    ],
    "energy": [
        ("photos/energy/wind-turbine.jpg",         "Wind Turbine — Poppy Field", "Wind turbine in a flowering field"),
        ("photos/energy/solar-farm.jpg",           "Solar Panel Farm",           "Utility-scale solar panel array"),
        ("photos/energy/hydroelectric-dam.jpg",    "Hoover Dam",                 "Aerial view of Hoover Dam and Lake Mead"),
        ("photos/energy/nuclear-cooling-tower.jpg","Nuclear Cooling Towers",     "Cofrentes hyperboloid towers with steam plume"),
    ],
    "industrial": [
        ("photos/industrial/factory-loom.jpg",   "Textile Production Line", "Industrial loom on the factory floor"),
        ("photos/industrial/cargo-port.jpg",     "Port of Rotterdam",       "Container stacks and gantry cranes"),
        ("photos/industrial/assembly-line.jpg",  "Auto Assembly Line",      "Hyundai car bodies on overhead conveyor"),
    ],
}


def _pick_photo_for_seed(topic: str, seed: int) -> _PhotoSpec:
    """Single-photo selection. Same template+different seed = different photo."""
    pool = TOPIC_FAMILIES[topic]
    return pool[seed % len(pool)]


def _pick_photos_for_seed(topic: str, seed: int, n: int) -> list[_PhotoSpec]:
    """N consecutive items from the pool, wrapping around."""
    pool = TOPIC_FAMILIES[topic]
    n_pool = len(pool)
    assert n_pool >= n, f"topic {topic!r} pool size {n_pool} < requested {n}"
    start = seed % n_pool
    return [pool[(start + i) % n_pool] for i in range(n)]


def _make_topic_photo_to_docx_template(
    topic: str, *, n_rows: int = 4, _variant: str = "after_caption",
) -> SynthTemplate:
    """Cartesian variant of the legacy `_make_photo_to_docx_template`.

    Per-seed photo from `TOPIC_FAMILIES[topic]`. Source docx has a 3-paragraph
    body (intro + 'Figure 1' caption + placeholder line); gold rebuilds the
    docx with body[:insert_idx] + the staged image + body[insert_idx:].
    Eval = `compare_docx_strict examine_images=True` (image-byte-hash multiset
    equality).

    `_variant`:
      - "after_caption": insert image after the 'Figure 1' caption (idx=2).
      - "cover":         insert image as the first element (idx=0), before
                         the intro paragraph — yields a 'cover image' deck.
    """
    assert _variant in ("after_caption", "cover")
    template_id = f"multi_topic_{topic}_photo_to_docx_{_variant}" if _variant != "after_caption" else f"multi_topic_{topic}_photo_to_docx"
    docx_basename = f"{template_id}.docx"
    docx_path  = f"{_DESKTOP}/{docx_basename}"
    expected   = f"{_TMP}/expected_{template_id}.docx"

    evaluator = {
        "func": "compare_docx_strict",
        "result": {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.docx"},
        "options": {
            "examine_font_name": False,
            "examine_font_size": False,
            "examine_color": False,
            "examine_highlight": False,
            "examine_images": True,
        },
    }
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        asset_rel, slide_title, caption_text = _pick_photo_for_seed(topic, seed)
        image_basename = asset_rel.rsplit("/", 1)[-1]
        image_path = f"{_DESKTOP}/{image_basename}"

        # Sibling-builder refactor (2026-05-10): frame the figure with a
        # topic-coherent narrative paragraph sourced from impress.py's
        # `_TOPIC_FAMILIES` (TopicTheme) — replaces the previous generic
        # "This <topic> brief..." placeholder. Photo selection still uses
        # the local `TOPIC_FAMILIES` Cartesian (5 topics × 4 skills) so
        # template_id rotation is preserved.
        theme = _pick_impress_theme(seed, f"photo_docx_{topic}")
        framing = theme.slide_bodies[0]
        # validation fix (multi_topic_*_photo_to_docx_cover): body_paras must
        # match the variant. The "appear immediately after the caption" hint
        # is only correct for the after_caption variant — for cover-image the
        # picture is placed BEFORE all paragraphs, so that sentence becomes a
        # contradiction the agent (correctly) doesn't honor, causing the gold
        # docx to embed the misleading sentence verbatim.
        if _variant == "after_caption":
            body_paras = [
                f"{theme.slide_titles[0]}. {framing}",
                f"Figure 1: {caption_text}.",
                "The figure should appear immediately after the caption above.",
            ]
        else:  # cover
            body_paras = [
                f"{theme.slide_titles[0]}. {framing}",
                f"Figure 1: {caption_text}.",
                "The cover image above introduces this brief.",
            ]
        width_in = 4.0
        insert_idx = 2 if _variant == "after_caption" else 0

        init_image = _stage_asset(asset_rel, image_path)
        src_paras: list[tuple[str | None, str]] = [(None, p) for p in body_paras]
        init_src = _build_docx_step(docx_path, src_paras)
        gold_py = textwrap.dedent(f"""\
            import io, os
            from PIL import Image
            from docx import Document
            from docx.shared import Inches
            path = {expected!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            body = {body_paras!r}
            insert_idx = {insert_idx}
            # python-docx rejects raw-DQT JPEGs (no JFIF/Exif marker in first
            # 32 bytes) with UnrecognizedImageError. Many of our asset photos
            # are raw-DQT; re-encode via PIL so the stream carries JFIF.
            # SSIM-tolerant image eval makes the re-encode drift a non-issue.
            _img = Image.open({image_path!r}).convert('RGB')
            _buf = io.BytesIO()
            _img.save(_buf, 'JPEG', quality=92)
            _buf.seek(0)
            doc = Document()
            for t in body[:insert_idx]:
                doc.add_paragraph(t)
            doc.add_picture(_buf, width=Inches({width_in!r}))
            for t in body[insert_idx:]:
                doc.add_paragraph(t)
            doc.save(path)
            """)
        init_steps = [
            init_image, init_src, _python_step(gold_py),
            # validation: photo is asset (no app),
            # docx is the focused workspace → [] → docx (Writer fg).
            *_multi_app_preopen(foreground=docx_path),
        ]

        if _variant == "after_caption":
            instructions = [
                f"Insert the photo at {image_path} immediately after the 'Figure 1' caption (i.e. after the second paragraph) in the Writer document you're currently editing at ~/Desktop/{docx_basename}. Save the document.",
            ]
        else:  # cover
            instructions = [
                f"Insert the photo at {image_path} at the very top of the Writer document you're currently editing at ~/Desktop/{docx_basename}, before the first paragraph (as a cover image). Save the document.",
            ]
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=n_rows,
        eval_class="compare_docx_strict",
        setup_class="real_asset_photo_docx",
    )


# ---- Cat I.3 — Photo → pptx slide -----------------------------------------
# Source pptx has a single empty slide; gold pptx has the staged image
# inserted on that slide. Eval = `compare_pptx_files` with `examine_shape=False`
# — shape COUNT per slide still must match (so "did agent insert any picture
# shape" is detected via the slide1.shapes vs slide2.shapes len check at
# desktop_env/evaluators/metrics/slides.py:302), but exact (left, top, width,
# height) is NOT enforced. Calibrated against eval task
# `osworld_libreoffice_impress_bf4e9888` (insert pic1..pic6 across 6 slides)
# which also uses examine_shape=False — the instruction commits to slide
# placement but not pixel-exact position.

def _make_topic_photo_to_pptx_template(topic: str, *, n_rows: int = 4) -> SynthTemplate:
    """Cartesian variant of the legacy `_make_photo_to_pptx_template`.

    Per-seed photo from `TOPIC_FAMILIES[topic]`. Source pptx has a single
    empty slide; gold pptx has the staged image inserted at a default
    position/size. Eval = `compare_pptx_files examine_shape=False` —
    verifies a picture shape is present on the target slide via shape-count
    delta; exact position/size are NOT enforced. (Validation fix: previously
    set examine_shape=True with a docstring claiming "exact position/size
    NOT enforced" — but upstream `compare_pptx_files` strictly checks
    (l,t,w,h) when `examine_shape=True`, contradicting the instruction's
    "any reasonable size/position" wording.)

    validation fix (was deferred from validation — multi_apps
    K-cluster): `common.py:LO_SAVE_POSTCONFIG` now includes an additional
    `xdotool search --name '^Save$'` matcher that dismisses LO's Save-As
    confirmation modal (actual WM_NAME = "Save"). Container-probe verified
    that the regex matches the real modal title. Same fix unblocks
    `_make_topic_photo_gallery_template` and the
    `_make_topic_photo_to_docx_template(_variant="cover")` family.
    """
    template_id = f"multi_topic_{topic}_photo_to_pptx"
    pptx_basename = f"{template_id}.pptx"
    pptx_path  = f"{_DESKTOP}/{pptx_basename}"
    expected   = f"{_TMP}/expected_{template_id}.pptx"

    left_cm, top_cm, width_cm = 4.0, 2.0, 14.0

    evaluator = {
        "func": "compare_pptx_files",
        "result": {"type": "vm_file", "path": pptx_path, "dest": "result.pptx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.pptx"},
        "options": {
            "examine_shape": False,
            "examine_image_size": False,
            "examine_text": False,
            "examine_font_name": False,
            "examine_font_size": False,
            "examine_font_bold": False,
            "examine_font_italic": False,
            "examine_font_underline": False,
            "examine_color_rgb": False,
            "examine_strike_through": False,
            "examine_alignment": False,
            "examine_indent": False,
            "examine_run_count": False,
            "examine_bullets": False,
            "examine_background_color": False,
            "examine_note": False,
        },
    }
    # 3-step LO-normalize oracle for pptx file-op evaluator.
    oracle_steps = _build_oracle_lo(pptx_path, expected, fmt="pptx")

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        asset_rel, _slide_title, _caption = _pick_photo_for_seed(topic, seed)
        image_basename = asset_rel.rsplit("/", 1)[-1]
        image_path = f"{_DESKTOP}/{image_basename}"

        # Sibling-builder refactor (2026-05-10): the source slide carries a
        # topic-coherent title + framing line from impress.py's TopicTheme
        # instead of being literally blank. Title text is identical in
        # source and gold so the eval (examine_text=False, examine_shape=
        # True) is unaffected — but the deck reads as a real mini-doc.
        theme = _pick_impress_theme(seed, f"photo_pptx_{topic}")
        slide_title_text = theme.slide_titles[0]
        slide_body_text = theme.slide_bodies[0]

        init_image = _stage_asset(asset_rel, image_path)
        src_py = textwrap.dedent(f"""\
            import os
            from pptx import Presentation
            from pptx.util import Cm, Pt
            path = {pptx_path!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            prs = Presentation()
            prs.slide_width = Cm(25.4)
            prs.slide_height = Cm(14.29)
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.4), Cm(22.0), Cm(1.8))
            tb.text_frame.text = {slide_title_text!r}
            tb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            bb = slide.shapes.add_textbox(Cm(1.0), Cm(12.0), Cm(22.0), Cm(1.5))
            bb.text_frame.text = {slide_body_text!r}
            bb.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            prs.save(path)
            """)
        gold_py = textwrap.dedent(f"""\
            import os
            from pptx import Presentation
            from pptx.util import Cm, Pt
            path = {expected!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            prs = Presentation()
            prs.slide_width = Cm(25.4)
            prs.slide_height = Cm(14.29)
            blank = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank)
            tb = slide.shapes.add_textbox(Cm(1.0), Cm(0.4), Cm(22.0), Cm(1.8))
            tb.text_frame.text = {slide_title_text!r}
            tb.text_frame.paragraphs[0].runs[0].font.size = Pt(24)
            bb = slide.shapes.add_textbox(Cm(1.0), Cm(12.0), Cm(22.0), Cm(1.5))
            bb.text_frame.text = {slide_body_text!r}
            bb.text_frame.paragraphs[0].runs[0].font.size = Pt(14)
            slide.shapes.add_picture(
                {image_path!r},
                Cm({left_cm!r}), Cm({top_cm!r}),
                width=Cm({width_cm!r}),
            )
            prs.save(path)
            """)
        init_steps = [
            init_image, _python_step(src_py), _python_step(gold_py),
            # validation: photo is asset, pptx is FG.
            *_multi_app_preopen(foreground=pptx_path),
        ]

        instructions = [
            f"Insert the {topic} photo at {image_path} onto the (only) slide of the Impress presentation you're currently editing at ~/Desktop/{pptx_basename} (any reasonable size/position). Save the presentation.",
        ]
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--impress", pptx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=n_rows,
        eval_class="compare_pptx_files",
        setup_class="real_asset_photo_pptx",
    )


# ---- Cat I.4 — CSV → xlsx with chart --------------------------------------
# Source xlsx = CSV-to-xlsx (LO-normalized). Gold xlsx = same data PLUS a
# chart (LineChart for time-series, BarChart for ranking) on the active
# sheet. Eval: compare_table with chart + sheet_data rules so both the
# data and the chart series-ref must match.

_CSV_TO_XLSX_CHART_SPECS: list[dict] = [
    {
        "template_id": "multi_csv_to_xlsx_chart_unemployment_line",
        "asset_rel": "data/csv/us-unemployment.csv",
        "csv_basename": "us-unemployment.csv",
        "xlsx_basename": "unemployment_with_chart.xlsx",
        # CSV header: observation_date,UNRATE. 940 data rows total.
        # We project the LAST 60 rows (most recent 5 years monthly) into
        # the source xlsx to keep the file modest, then chart UNRATE.
        "row_select": "rows = rows[:1] + rows[-60:]",  # header + last 60
        "chart_type": "line",  # LineChart for time-series
        "chart_title": "Unemployment Rate (last 60 months)",
        "instructions": [
            "Open ~/Desktop/unemployment_with_chart.xlsx in Calc. Add a LineChart on the active sheet that plots column B (UNRATE) for the data rows, with title 'Unemployment Rate (last 60 months)'. Save the workbook.",
        ],
    },
    {
        "template_id": "multi_csv_to_xlsx_chart_world_gdp_bar",
        "asset_rel": "data/csv/world-gdp-2022.csv",
        "csv_basename": "world-gdp-2022.csv",
        "xlsx_basename": "world_gdp_top10_with_chart.xlsx",
        # Take header + top 10 rows by GDP_USD.
        "row_select": (
            "header, data = rows[0], rows[1:]\n"
            "data.sort(key=lambda r: float(r[3]), reverse=True)\n"
            "rows = [header] + data[:10]"
        ),
        "chart_type": "bar",  # BarChart for ranking
        "chart_title": "Top 10 GDP (2022)",
        "instructions": [
            "Open ~/Desktop/world_gdp_top10_with_chart.xlsx in Calc. Add a BarChart on the active sheet that plots column D (GDP_USD) for the data rows, with title 'Top 10 GDP (2022)'. Save the workbook.",
        ],
    },
]


def _make_csv_to_xlsx_chart_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    asset_rel     = spec["asset_rel"]
    csv_path      = f"{_DESKTOP}/{spec['csv_basename']}"
    xlsx_path     = f"{_DESKTOP}/{spec['xlsx_basename']}"
    expected      = f"{_TMP}/expected_{template_id}.xlsx"
    row_select    = spec["row_select"]
    chart_type    = spec["chart_type"]      # "line" | "bar"
    chart_title   = spec["chart_title"]
    instructions  = spec["instructions"]

    cls = {"line": "LineChart", "bar": "BarChart"}[chart_type]

    init_csv = _stage_asset(asset_rel, csv_path)

    # Source xlsx: CSV-to-xlsx with the projected rows, LO-normalized.
    src_py = textwrap.dedent(f"""\
        import os, csv
        from openpyxl import Workbook
        src = {csv_path!r}
        path = {xlsx_path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        with open(src, newline='') as f:
            rows = list(csv.reader(f))
        {textwrap.indent(row_select, '        ').lstrip()}
        wb = Workbook()
        ws = wb.active
        ws.title = 'Sheet1'
        for r in rows:
            ws.append(r)
        wb.save(path)
        # _LO_NORMALIZE_TAIL — match LO save format (per AGENTS.md F1 mitigation).
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, path],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename(path))
            if os.path.exists(_conv):
                _sh.copy(_conv, path)
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)

    # Gold xlsx: copy source then add chart referencing the value column.
    # Value column is the LAST column of the projected rows (UNRATE col B
    # for unemployment; GDP_USD col D for world-gdp).
    gold_py = textwrap.dedent(f"""\
        import os, csv, shutil
        from openpyxl import load_workbook
        from openpyxl.chart import {cls}
        from openpyxl.chart.series import Series
        from openpyxl.chart.data_source import NumDataSource, NumRef
        from openpyxl.utils import get_column_letter
        src_xlsx = {xlsx_path!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        shutil.copy(src_xlsx, path)
        wb = load_workbook(path)
        ws = wb.worksheets[0]
        max_r = ws.max_row or 0
        max_c = ws.max_column or 0
        col_letter = get_column_letter(max_c)
        sheet_name = ws.title
        series_ref = '{{}}!${{}}$2:${{}}${{}}'.format(sheet_name, col_letter, col_letter, max_r)
        chart = {cls}()
        chart.title = {chart_title!r}
        ser = Series()
        ser.val = NumDataSource(numRef=NumRef(f=series_ref))
        chart.series.append(ser)
        ws.add_chart(chart, 'G2')
        wb.save(path)
        # _LO_NORMALIZE_TAIL
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, path],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename(path))
            if os.path.exists(_conv):
                _sh.copy(_conv, path)
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)

    init_steps = [init_csv, _python_step(src_py), _python_step(gold_py)]
    oracle_steps = [_execute(f"cp '{expected}' '{xlsx_path}'")]

    # Eval: chart rule (type+title) + sheet_data so both must match.
    evaluator = {
        "func": "compare_table",
        "result": {"type": "vm_file", "path": xlsx_path, "dest": "result.xlsx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.xlsx"},
        "options": {"rules": [
            {"type": "chart", "sheet_idx0": 0, "sheet_idx1": "EI0",
             "chart_props": ["type", "title"]},
            # validation: positional indexing (see _make_csv_to_xlsx_template comment).
            {"type": "sheet_data", "sheet_idx0": 0, "sheet_idx1": 0},
        ]},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", xlsx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_table",
        setup_class="real_asset_csv_xlsx_chart",
    )


# ---- Cat I.5 — 3-way FRED CSV concat with `series` tag --------------------
# Stage 3 FRED CSVs (us-gdp / us-unemployment / us-inflation-cpi). Oracle
# runs an awk pipeline that emits a 3-col `date,series,value` row for each
# input. Eval: compare_csv.

def _make_asset_concat_3way_economic() -> SynthTemplate:
    template_id = "multi_csv_concat_fred_3way_economic"
    asset_rels = [
        "data/csv/us-gdp.csv",
        "data/csv/us-unemployment.csv",
        "data/csv/us-inflation-cpi.csv",
    ]
    src_paths = [
        f"{_DESKTOP}/fred3/us-gdp.csv",
        f"{_DESKTOP}/fred3/us-unemployment.csv",
        f"{_DESKTOP}/fred3/us-inflation-cpi.csv",
    ]
    series_names = ["GDP", "UNRATE", "CPIAUCSL"]
    dst_path = f"{_DESKTOP}/economic_3way.csv"
    expected = f"{_TMP}/expected_{template_id}.csv"

    init_steps: list[dict] = [
        _stage_asset(rel, src) for rel, src in zip(asset_rels, src_paths)
    ]
    gold_py = textwrap.dedent(f"""\
        import os
        path = {expected!r}
        srcs = {src_paths!r}
        names = {series_names!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        with open(path, 'w') as out:
            out.write('date,series,value\\n')
            for src, name in zip(srcs, names):
                with open(src) as f:
                    lines = f.read().splitlines()
                for ln in lines[1:]:
                    date, val = ln.split(',', 1)
                    out.write(f'{{date}},{{name}},{{val}}\\n')
        """)
    init_steps.append(_python_step(gold_py))
    # validation gap-close: terminal-driven (awk concat) — pre-open terminal.
    init_steps.extend(_terminal_preopen_steps())

    oracle_steps = [_execute(
        f"{{ echo 'date,series,value'; "
        f"awk -F, 'NR>1 {{print $1\",GDP,\"$2}}' "
        f"{src_paths[0]}; "
        f"awk -F, 'NR>1 {{print $1\",UNRATE,\"$2}}' "
        f"{src_paths[1]}; "
        f"awk -F, 'NR>1 {{print $1\",CPIAUCSL,\"$2}}' "
        f"{src_paths[2]}; "
        f"}} > {dst_path}"
    )]
    evaluator = {
        "func": "compare_csv",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.csv"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.csv"},
        "options": {"strict": False},
    }
    instructions = [
        "In a terminal, combine the three FRED CSVs in ~/Desktop/fred3/ (us-gdp.csv, us-unemployment.csv, us-inflation-cpi.csv) into a single ~/Desktop/economic_3way.csv with header `date,series,value`. For each input, drop its header and emit `<date>,<series_name>,<value>` rows where series_name is GDP, UNRATE, and CPIAUCSL respectively (in that order).",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_table",
        setup_class="real_asset_text",
    )


# ---- Cat I.6 — Wikipedia HTML extract (h2 titles) -------------------------
# Stage one Wikipedia HTML, oracle uses sed to extract the visible text of
# every <h2> tag. Eval: compare_text_file.

def _make_asset_wikipedia_h2_extract() -> SynthTemplate:
    template_id = "multi_wikipedia_extract_titles_apollo"
    asset_rel = "html/wikipedia/apollo-program.html"
    src_path = f"{_DESKTOP}/wikipedia/apollo-program.html"
    dst_path = f"{_DESKTOP}/apollo_h2_titles.txt"
    expected = f"{_TMP}/expected_{template_id}.txt"

    init_src = _stage_asset(asset_rel, src_path)
    # Gold: read the staged HTML, extract the visible text of every <h2>.
    # Wikipedia <h2> tags often contain nested <span class="mw-headline"
    # id="..."> wrappers; we strip ALL inner tags and collapse whitespace.
    # We mirror the same regex the oracle pipeline uses (sed) so eval
    # passes byte-for-byte.
    gold_py = textwrap.dedent(f"""\
        import os, re
        src = {src_path!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        with open(src, encoding='utf-8') as f:
            html = f.read()
        # Find all <h2 ...>...</h2> blocks (non-greedy, DOTALL).
        blocks = re.findall(r'<h2\\b[^>]*>(.*?)</h2>', html, flags=re.DOTALL | re.IGNORECASE)
        out_lines = []
        for b in blocks:
            text = re.sub(r'<[^>]+>', '', b)          # strip nested tags
            text = re.sub(r'\\s+', ' ', text).strip() # collapse whitespace
            if text:
                out_lines.append(text)
        with open(path, 'w', encoding='utf-8') as f:
            for ln in out_lines:
                f.write(ln + '\\n')
        """)
    init_steps = [
        init_src,
        _python_step(gold_py),
        # validation gap-close: terminal-driven (python/sed extract) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]

    # Oracle: same logic via python heredoc — keeping the pipeline shape
    # consistent with the gold ensures byte-equality.
    oracle_py = textwrap.dedent(f"""\
        import re
        src = {src_path!r}
        out_path = {dst_path!r}
        with open(src, encoding='utf-8') as f:
            html = f.read()
        blocks = re.findall(r'<h2\\b[^>]*>(.*?)</h2>', html, flags=re.DOTALL | re.IGNORECASE)
        out_lines = []
        for b in blocks:
            text = re.sub(r'<[^>]+>', '', b)
            text = re.sub(r'\\s+', ' ', text).strip()
            if text:
                out_lines.append(text)
        with open(out_path, 'w', encoding='utf-8') as f:
            for ln in out_lines:
                f.write(ln + '\\n')
        """)
    oracle_steps = [_python_step(oracle_py)]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.txt"},
    }
    instructions = [
        "From a terminal, scan ~/Desktop/wikipedia/apollo-program.html and pull out the visible text of every <h2> tag. Strip any nested HTML inside each heading (for example <span> wrappers), collapse internal whitespace to single spaces, and write one heading per line to ~/Desktop/apollo_h2_titles.txt (skip empty headings).",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="real_asset_html",
    )


# ---- Cat I.7 — Multi-photo gallery deck (≥3 slides, each with image+title) -
# Stage N photos via _stage_asset, build a source pptx with N empty slides
# (each pre-titled), and a gold pptx with the same N slides + the staged
# image inserted on each slide at fixed position/size. Eval:
# compare_pptx_files with examine_image_size=True (and default
# examine_text=True so titles match).

def _make_topic_photo_gallery_template(topic: str, *, n_rows: int = 4) -> SynthTemplate:
    """Cartesian variant of the legacy `_make_photo_gallery_template`.

    4-slide gallery — every slide is from the SAME `topic` family so the deck
    is thematically coherent. Per-seed rotation: seed picks the starting
    index in the family pool, then 4 consecutive items wrap around.
    Source pptx has 4 slides each with a title placeholder; gold adds the
    matching staged image to each slide. Eval =
    `compare_pptx_files examine_shape=False` — verifies a picture shape is
    present on each slide via shape-count delta; exact position/size are NOT
    enforced (the instruction says "any reasonable size/position", so eval
    must not over-constrain). Default examine_text=True so titles must
    still match. Calibrated against eval task
    `osworld_libreoffice_impress_bf4e9888` (insert pic1..pic6 across 6
    slides) which uses examine_shape=False.

    Validation fix: previously set examine_shape=True with a docstring
    claiming "exact position/size NOT enforced", but upstream
    `compare_pptx_files` strictly checks (l,t,w,h) when examine_shape=True,
    causing trigger O false-fails on agent-inserted images at non-exact
    positions.
    """
    template_id = f"multi_topic_{topic}_photo_gallery_4"
    pptx_basename = f"{template_id}.pptx"
    pptx_path = f"{_DESKTOP}/{pptx_basename}"
    expected  = f"{_TMP}/expected_{template_id}.pptx"

    left_cm, top_cm, width_cm = 4.0, 4.0, 16.0

    evaluator = {
        "func": "compare_pptx_files",
        "result": {"type": "vm_file", "path": pptx_path, "dest": "result.pptx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.pptx"},
        "options": {
            "examine_shape": False,
            "examine_image_size": False,
            "examine_font_name": False,
            "examine_font_size": False,
            "examine_font_bold": False,
            "examine_font_italic": False,
            "examine_font_underline": False,
            "examine_color_rgb": False,
            "examine_strike_through": False,
            "examine_alignment": False,
            "examine_indent": False,
            "examine_run_count": False,
            "examine_bullets": False,
            "examine_background_color": False,
            "examine_note": False,
        },
    }
    # 3-step LO-normalize oracle for pptx file-op evaluator.
    oracle_steps = _build_oracle_lo(pptx_path, expected, fmt="pptx")

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        photos = _pick_photos_for_seed(topic, seed, 4)
        # Per-slide (asset_rel, basename, fallback_title)
        slides = [(rel, rel.rsplit("/", 1)[-1], t) for (rel, t, _c) in photos]

        init_image_steps = [
            _stage_asset(rel, f"{_DESKTOP}/{base}") for (rel, base, _t) in slides
        ]
        image_paths = [f"{_DESKTOP}/{base}" for (_r, base, _t) in slides]
        # Sibling-builder refactor (2026-05-10): per-slide titles now come
        # from the picked TopicTheme's `slide_titles` arc (intro → details →
        # reflection) instead of the per-photo caption tuple — so the deck
        # reads as a topic-coherent gallery narrative rather than four
        # standalone photo captions. Source and gold use the same titles
        # so the `compare_pptx_files` default examine_text=True remains
        # green.
        theme = _pick_impress_theme(seed, f"photo_gallery_{topic}")
        # Cycle TopicTheme titles to fill 4 slides (theme has ≥6, so just slice).
        titles = list(theme.slide_titles[:4])

        src_py = textwrap.dedent(f"""\
            import os
            from pptx import Presentation
            from pptx.util import Cm
            path = {pptx_path!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            prs = Presentation()
            prs.slide_width = Cm(25.4)
            prs.slide_height = Cm(14.29)
            layout = prs.slide_layouts[5]
            titles = {titles!r}
            for t in titles:
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = t
            prs.save(path)
            """)
        gold_py = textwrap.dedent(f"""\
            import os
            from pptx import Presentation
            from pptx.util import Cm
            path = {expected!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            prs = Presentation()
            prs.slide_width = Cm(25.4)
            prs.slide_height = Cm(14.29)
            layout = prs.slide_layouts[5]
            titles = {titles!r}
            images = {image_paths!r}
            for t, img in zip(titles, images):
                slide = prs.slides.add_slide(layout)
                slide.shapes.title.text = t
                slide.shapes.add_picture(
                    img,
                    Cm({left_cm!r}), Cm({top_cm!r}),
                    width=Cm({width_cm!r}),
                )
            prs.save(path)
            """)
        init_steps = init_image_steps + [
            _python_step(src_py), _python_step(gold_py),
            # validation: photos are assets, pptx is FG.
            *_multi_app_preopen(foreground=pptx_path),
        ]

        slide_lines = ". ".join(
            f"Slide {i + 1}: {b}" for i, (_r, b, _t) in enumerate(slides)
        )
        instructions = [
            f"The Impress presentation you're currently editing at ~/Desktop/{pptx_basename} has 4 slides, each with a title already in place. Insert the matching {topic} photo from /home/user/Desktop/ on each slide (any reasonable size/position). {slide_lines}. Save the presentation.",
        ]
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--impress", pptx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=n_rows,
        eval_class="compare_pptx_files",
        setup_class="real_asset_photo_gallery",
    )


# ---- Cat I.7b — Photo → docx with `compare_docx_images` eval --------------
# Mirror of `_make_topic_photo_to_docx_template` but eval'd with the upstream
# `compare_docx_images` func (PIL pixel-equality on embedded images), NOT
# `compare_docx_strict examine_images=True` (raw blob hash). Eval row
# `osworld_multi_apps_227d2f97` exercises this exact func against an XCF→docx
# image-paste flow; this bridges the long-tail func (eval K=1) so synth has
# at least 1-2 templates with the matching `evaluator.func` name. Skill is
# similar to the byte-hash photo→docx rows above but distinct: PIL pixel
# equality is more lenient (image re-encoding by LO is tolerated as long as
# the decoded pixels match).
def _make_photo_to_docx_images_eval_template(
    topic: str, *, n_rows: int = 2,
) -> SynthTemplate:
    """Single-photo insert into a 1-paragraph docx, eval'd with
    `compare_docx_images`. Per-seed photo from `TOPIC_FAMILIES[topic]`.
    Source = empty docx with a single intro paragraph; gold = same docx
    with the staged photo embedded after the intro paragraph. Oracle =
    `cp gold result.docx`.
    """
    template_id = f"multi_photo_docx_images_eval_{topic}"
    docx_basename = f"{template_id}.docx"
    docx_path  = f"{_DESKTOP}/{docx_basename}"
    expected   = f"{_TMP}/expected_{template_id}.docx"

    evaluator = {
        "func": "compare_docx_images",
        "result":   {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected,  "dest": "expected.docx"},
    }
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        asset_rel, slide_title, caption_text = _pick_photo_for_seed(topic, seed)
        image_basename = asset_rel.rsplit("/", 1)[-1]
        image_path = f"{_DESKTOP}/{image_basename}"

        intro_para = (
            f"This {topic} note accompanies a single illustrative figure "
            f"({slide_title}). The image should appear immediately below."
        )
        width_in = 4.0

        init_image = _stage_asset(asset_rel, image_path)
        src_paras: list[tuple[str | None, str]] = [(None, intro_para)]
        init_src = _build_docx_step(docx_path, src_paras)
        gold_py = textwrap.dedent(f"""\
            import io, os
            from PIL import Image
            from docx import Document
            from docx.shared import Inches
            path = {expected!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            # python-docx rejects raw-DQT JPEGs without JFIF/Exif marker;
            # re-encode via PIL so the stream carries JFIF.
            _img = Image.open({image_path!r}).convert('RGB')
            _buf = io.BytesIO()
            _img.save(_buf, 'JPEG', quality=92)
            _buf.seek(0)
            doc = Document()
            doc.add_paragraph({intro_para!r})
            doc.add_picture(_buf, width=Inches({width_in!r}))
            doc.save(path)
            """)
        init_steps = [init_image, init_src, _python_step(gold_py)]
        instructions = [
            f"Open ~/Desktop/{docx_basename} in LibreOffice Writer and insert the photo at {image_path} immediately after the existing intro paragraph. Save the document in place.",
        ]
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=n_rows,
        eval_class="compare_docx_images",
        setup_class="real_asset_photo_docx_images_eval",
    )


def _make_dual_photo_to_docx_images_eval_template(
    topic: str, *, n_rows: int = 2,
) -> SynthTemplate:
    """Dual-photo insert (mirrors plan row #139 — Earth side-by-side).
    Source = empty docx with a 1-line intro; gold = intro + photo1 + photo2.
    Eval = `compare_docx_images` (PIL pixel-equality across both embedded
    images, multiset).
    """
    template_id = f"multi_photo_docx_images_eval_dual_{topic}"
    docx_basename = f"{template_id}.docx"
    docx_path  = f"{_DESKTOP}/{docx_basename}"
    expected   = f"{_TMP}/expected_{template_id}.docx"

    evaluator = {
        "func": "compare_docx_images",
        "result":   {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected,  "dest": "expected.docx"},
    }
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        photos = _pick_photos_for_seed(topic, seed, 2)
        photo_paths: list[str] = []
        init_image_steps: list[dict] = []
        for asset_rel, _title, _caption in photos:
            image_basename = asset_rel.rsplit("/", 1)[-1]
            image_path = f"{_DESKTOP}/{image_basename}"
            photo_paths.append(image_path)
            init_image_steps.append(_stage_asset(asset_rel, image_path))

        intro_para = (
            f"This {topic} comparison shows two reference figures, embedded "
            "below in order."
        )
        width_in = 3.5

        src_paras: list[tuple[str | None, str]] = [(None, intro_para)]
        init_src = _build_docx_step(docx_path, src_paras)
        gold_py = textwrap.dedent(f"""\
            import io, os
            from PIL import Image
            from docx import Document
            from docx.shared import Inches
            path = {expected!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            doc = Document()
            doc.add_paragraph({intro_para!r})
            for img in {photo_paths!r}:
                # python-docx rejects raw-DQT JPEGs without JFIF marker;
                # re-encode via PIL so the stream carries JFIF.
                _img = Image.open(img).convert('RGB')
                _buf = io.BytesIO()
                _img.save(_buf, 'JPEG', quality=92)
                _buf.seek(0)
                doc.add_picture(_buf, width=Inches({width_in!r}))
            doc.save(path)
            """)
        init_steps = init_image_steps + [init_src, _python_step(gold_py)]
        photo_list_str = " and ".join(photo_paths)
        instructions = [
            f"Open ~/Desktop/{docx_basename} in LibreOffice Writer and insert the two photos at {photo_list_str} (in that order) below the existing intro paragraph. Save the document in place.",
        ]
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=n_rows,
        eval_class="compare_docx_images",
        setup_class="real_asset_photo_docx_images_eval_dual",
    )


# ---- Cat I.8 — CSV → docx text summary (NOT a table) ----------------------
# Stage CSV, build a source docx with a "TODO: insert summary sentence here"
# placeholder, and a gold docx where the placeholder is replaced with a
# computed text sentence. Eval: compare_docx_files (paragraphs equal).

_CSV_TO_DOCX_TEXT_SUMMARY_SPECS: list[dict] = [
    {
        "template_id": "multi_csv_to_docx_text_summary_us_pop",
        "asset_rel": "data/csv/us-population-states.csv",
        "csv_basename": "us-population-states.csv",
        "docx_basename": "population_summary.docx",
        # Header: State,Population2020,StateFIPS. Sort by Population2020 desc.
        "summary_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "data = rows[1:]\n"
            "data.sort(key=lambda r: int(r[1]), reverse=True)\n"
            "top_state, top_pop = data[0][0], int(data[0][1])\n"
            "bot_state, bot_pop = data[-1][0], int(data[-1][1])\n"
            "summary_text = (\n"
            "    f'The most populous U.S. state in 2020 was {top_state} with {top_pop:,} residents, '\n"
            "    f'while the least populous was {bot_state} with {bot_pop:,} residents.'\n"
            ")\n"
        ),
        "docx_pre": [
            ("Title", "U.S. Population Summary (2020)"),
            (None, "The dataset us-population-states.csv lists each U.S. state's 2020 census population."),
        ],
        "instructions": [
            "Summarize ~/Desktop/us-population-states.csv (open in Calc) into the docx report you're currently editing at ~/Desktop/population_summary.docx (in Writer). Append a single sentence to the document of the form: 'The most populous U.S. state in 2020 was <NAME> with <POP> residents, while the least populous was <NAME> with <POP> residents.' Use comma-separated thousands (e.g. '39,538,223'). Save the document.",
        ],
    },
    {
        "template_id": "multi_csv_to_docx_text_summary_world_gdp_top3",
        "asset_rel": "data/csv/world-gdp-2022.csv",
        "csv_basename": "world-gdp-2022.csv",
        "docx_basename": "gdp_top3_summary.docx",
        # Header: CountryCode,CountryName,Year,GDP_USD. Top 3 by GDP_USD.
        "summary_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "data = rows[1:]\n"
            "data.sort(key=lambda r: float(r[3]), reverse=True)\n"
            "top3 = data[:3]\n"
            "names = [r[1] for r in top3]\n"
            "summary_text = (\n"
            "    f'The top three entities by 2022 GDP (USD) were {names[0]}, {names[1]}, and {names[2]}.'\n"
            ")\n"
        ),
        "docx_pre": [
            ("Title", "World GDP — Top 3 Summary (2022)"),
            (None, "The dataset world-gdp-2022.csv ranks countries and regional aggregates by 2022 GDP."),
        ],
        # Validation fix: instruction clarified — file contains both
        # countries AND regional aggregates (AFE/EAS/etc.); the top three
        # rows by raw GDP_USD include aggregates. Reading "top three
        # countries" would lead a human agent astray.
        "instructions": [
            "Summarize ~/Desktop/world-gdp-2022.csv (open in Calc) into the docx report you're currently editing at ~/Desktop/gdp_top3_summary.docx (in Writer). Sort the data rows (everything except the header) by the GDP_USD column in descending order; the file mixes individual countries AND regional aggregates — include both. Append a single sentence naming the top three entries (in CountryName) in the form: 'The top three entities by 2022 GDP (USD) were <A>, <B>, and <C>.' Save the document.",
        ],
    },
    {
        "template_id": "multi_csv_to_docx_text_summary_unemployment_2024_avg",
        "asset_rel": "data/csv/us-unemployment.csv",
        "csv_basename": "us-unemployment.csv",
        "docx_basename": "unemployment_2024_avg.docx",
        # Header: observation_date,UNRATE. Compute mean UNRATE over rows
        # whose date starts with '2024-'.
        # Validation fix: switched from :.2f (round-half-to-even ambiguity
        # for values like 4.025 → "4.02" in Python vs "4.03" in round-half-up
        # → FALSE-FAIL) to :.1f (one decimal place — unambiguous across all
        # rounding conventions for this dataset's range).
        "summary_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "vals = [float(r[1]) for r in rows[1:] if r[0].startswith('2024-')]\n"
            "avg = sum(vals) / len(vals) if vals else 0.0\n"
            "summary_text = (\n"
            "    f'The average U.S. unemployment rate across {len(vals)} monthly observations in 2024 was {avg:.1f} percent.'\n"
            ")\n"
        ),
        "docx_pre": [
            ("Title", "U.S. Unemployment — 2024 Annual Average"),
            (None, "The dataset us-unemployment.csv contains monthly UNRATE observations from 1948 onward."),
        ],
        "instructions": [
            "Summarize ~/Desktop/us-unemployment.csv (open in Calc) into the docx report you're currently editing at ~/Desktop/unemployment_2024_avg.docx (in Writer). Compute the arithmetic mean of the UNRATE column for all rows whose date begins with '2024-'. Append a single sentence to the document of the form: 'The average U.S. unemployment rate across <N> monthly observations in 2024 was <X.X> percent.' (one decimal place; standard rounding). Save the document.",
        ],
    },
    # Mass-add — CSV → docx text-summary variants.
    # Pruned (multi_csv_to_docx_text_summary_oil_wti_2024_avg): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_text_summary_oil_wti_2024_avg",
#         "asset_rel": "data/csv/oil-wti-daily.csv",
#         "csv_basename": "oil-wti-daily.csv",
#         "docx_basename": "oil_wti_2024_avg.docx",
#         "summary_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "vals = [float(r[1]) for r in rows[1:] if r[0].startswith('2024-')]\n"
#             "avg = sum(vals) / len(vals) if vals else 0.0\n"
#             "summary_text = (\n"
#             "    f'The average WTI crude price across {len(vals)} daily observations in 2024 was {avg:.1f} USD per barrel.'\n"
#             ")\n"
#         ),
#         "docx_pre": [
#             ("Title", "WTI Crude — 2024 Daily Average"),
#             (None, "The dataset oil-wti-daily.csv contains daily WTI prices."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/oil-wti-daily.csv and ~/Desktop/oil_wti_2024_avg.docx. Compute the mean of the DCOILWTICO column over rows whose date begins with '2024-'. Append a single sentence to the document of the form: 'The average WTI crude price across <N> daily observations in 2024 was <X.X> USD per barrel.' (one decimal place). Save it.",
#         ],
#     },
    {
        "template_id": "multi_csv_to_docx_text_summary_us_state_income_top",
        "asset_rel": "data/csv/us-state-median-income.csv",
        "csv_basename": "us-state-median-income.csv",
        "docx_basename": "us_state_income_top.docx",
        "summary_py": (
            "import csv\n"
            "with open(src) as f:\n"
            "    rows = list(csv.reader(f))\n"
            "data = rows[1:]\n"
            "data.sort(key=lambda r: int(r[1]), reverse=True)\n"
            "top_state, top_income = data[0][0], int(data[0][1])\n"
            "bot_state, bot_income = data[-1][0], int(data[-1][1])\n"
            "summary_text = (\n"
            "    f'The U.S. state with the highest median household income was {top_state} at {top_income:,} USD, '\n"
            "    f'while the lowest was {bot_state} at {bot_income:,} USD.'\n"
            ")\n"
        ),
        "docx_pre": [
            ("Title", "U.S. State Median Household Income — Extremes"),
            (None, "The dataset us-state-median-income.csv lists each U.S. state's median household income."),
        ],
        "instructions": [
            "Summarize ~/Desktop/us-state-median-income.csv (open in Calc) into the docx report you're currently editing at ~/Desktop/us_state_income_top.docx (in Writer). Append a single sentence of the form: 'The U.S. state with the highest median household income was <NAME> at <AMT> USD, while the lowest was <NAME> at <AMT> USD.' (use comma-separated thousands). Save the document.",
        ],
    },
    # Pruned (multi_csv_to_docx_text_summary_us_inflation_2024_avg): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_text_summary_us_inflation_2024_avg",
#         "asset_rel": "data/csv/us-inflation-cpi.csv",
#         "csv_basename": "us-inflation-cpi.csv",
#         "docx_basename": "us_inflation_2024_avg.docx",
#         "summary_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "vals = [float(r[1]) for r in rows[1:] if r[0].startswith('2024-')]\n"
#             "avg = sum(vals) / len(vals) if vals else 0.0\n"
#             "summary_text = (\n"
#             "    f'The average U.S. CPI across {len(vals)} monthly observations in 2024 was {avg:.1f}.'\n"
#             ")\n"
#         ),
#         "docx_pre": [
#             ("Title", "U.S. CPI — 2024 Annual Average"),
#             (None, "The dataset us-inflation-cpi.csv contains monthly CPIAUCSL observations."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-inflation-cpi.csv and ~/Desktop/us_inflation_2024_avg.docx. Compute the mean of the CPIAUCSL column for rows whose date begins with '2024-'. Append: 'The average U.S. CPI across <N> monthly observations in 2024 was <X.X>.' (one decimal place). Save it.",
#         ],
#     },
    # Pruned (multi_csv_to_docx_text_summary_us_fedfunds_2024_avg): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_text_summary_us_fedfunds_2024_avg",
#         "asset_rel": "data/csv/us-fed-funds-rate.csv",
#         "csv_basename": "us-fed-funds-rate.csv",
#         "docx_basename": "us_fedfunds_2024_avg.docx",
#         "summary_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "vals = [float(r[1]) for r in rows[1:] if r[0].startswith('2024-')]\n"
#             "avg = sum(vals) / len(vals) if vals else 0.0\n"
#             "summary_text = (\n"
#             "    f'The average U.S. fed-funds rate across {len(vals)} monthly observations in 2024 was {avg:.1f} percent.'\n"
#             ")\n"
#         ),
#         "docx_pre": [
#             ("Title", "U.S. Fed Funds Rate — 2024 Annual Average"),
#             (None, "The dataset us-fed-funds-rate.csv contains monthly FEDFUNDS observations."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/us-fed-funds-rate.csv and ~/Desktop/us_fedfunds_2024_avg.docx. Compute the mean of the FEDFUNDS column for rows whose date begins with '2024-'. Append: 'The average U.S. fed-funds rate across <N> monthly observations in 2024 was <X.X> percent.' (one decimal place). Save it.",
#         ],
#     },
    # Pruned (multi_csv_to_docx_text_summary_world_population_top): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "multi_csv_to_docx_text_summary_world_population_top",
#         "asset_rel": "data/csv/world-population-2022.csv",
#         "csv_basename": "world-population-2022.csv",
#         "docx_basename": "world_population_top.docx",
#         "summary_py": (
#             "import csv\n"
#             "with open(src) as f:\n"
#             "    rows = list(csv.reader(f))\n"
#             "data = rows[1:]\n"
#             "data.sort(key=lambda r: float(r[3]), reverse=True)\n"
#             "top3 = data[:3]\n"
#             "names = [r[1] for r in top3]\n"
#             "summary_text = (\n"
#             "    f'The top three entities by 2022 population were {names[0]}, {names[1]}, and {names[2]}.'\n"
#             ")\n"
#         ),
#         "docx_pre": [
#             ("Title", "World Population — Top 3 (2022)"),
#             (None, "The dataset world-population-2022.csv mixes individual countries and regional aggregates."),
#         ],
#         "instructions": [
#             "Open ~/Desktop/world-population-2022.csv and ~/Desktop/world_population_top.docx. Sort by Population descending and append: 'The top three entities by 2022 population were <A>, <B>, and <C>.' (CountryName values; include aggregates). Save it.",
#         ],
#     },
]


def _make_csv_to_docx_text_summary_template(spec: dict) -> SynthTemplate:
    template_id  = spec["template_id"]
    asset_rel    = spec["asset_rel"]
    csv_path     = f"{_DESKTOP}/{spec['csv_basename']}"
    docx_path    = f"{_DESKTOP}/{spec['docx_basename']}"
    expected     = f"{_TMP}/expected_{template_id}.docx"
    summary_py   = spec["summary_py"]
    docx_pre     = spec["docx_pre"]
    instructions = spec["instructions"]

    init_csv = _stage_asset(asset_rel, csv_path)
    # Source docx: pre paragraphs only (no summary yet).
    init_docx = _build_docx_step(docx_path, docx_pre)

    # Build gold by reading staged CSV, computing `summary_text`, then
    # writing a docx with `docx_pre` + the summary sentence as final para.
    gold_py = textwrap.dedent(f"""\
        import os
        from docx import Document
        src = {csv_path!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        {textwrap.indent(summary_py, '        ').lstrip()}
        doc = Document()
        for style, text in {docx_pre!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        doc.add_paragraph(summary_text)
        doc.save(path)
        """)
    init_steps = [
        init_csv, init_docx, _python_step(gold_py),
        # validation: [csv (Calc bg)] → docx (Writer fg).
        *_multi_app_preopen(background=[csv_path], foreground=docx_path),
    ]
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.docx"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="real_asset_csv_docx_text",
    )


# ---- Cat I.9 — arxiv PDF → text via pymupdf ------------------------------
# Stage one arxiv PDF, then use pymupdf (fitz) to extract some derived text
# (title-line / page-count) into a plain text file. Eval: compare_text_file.
# Non-LO sink: terminal-task config_override + no oracle_after_postconfig.

_ARXIV_PDF_SPECS: list[dict] = [
    {
        "template_id": "multi_arxiv_pdf_extract_title",
        "asset_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "pdf_basename": "arxiv-1706-03762.pdf",
        "out_basename": "arxiv_title.txt",
        # Extract the paper title — for arxiv-1706-03762 ("Attention Is All
        # You Need") the title appears as a stand-alone non-empty line on
        # page 0. Validation fix: previous heuristic ("first non-empty line of
        # length 5..120 that does not start with 'Provided'") picked the
        # SECOND line of the Google attribution boilerplate ("reproduce the
        # tables and figures in this paper solely for use in journalistic
        # or"), but the agent semantically writes the real title. We now
        # explicitly skip the whole Google attribution block (any line whose
        # lowercased text contains "provided", "reproduce", "scholarly", or
        # "journalistic") AND any line ending in "by" / "to" / "or" / "and"
        # / "of" (sentence-continuation fragments) so the heuristic lands on
        # the real title.
        "extract_py": (
            "import subprocess\n"
            "page0 = subprocess.run(['pdftotext', '-f', '1', '-l', '1', src, '-'], capture_output=True, text=True).stdout\n"
            "_NOISE = ('provided', 'reproduce', 'scholarly', 'journalistic')\n"
            "_TRAIL = (' by', ' to', ' or', ' and', ' of', ' for', ' in')\n"
            "title = ''\n"
            "for ln in page0.splitlines():\n"
            "    t = ln.strip()\n"
            "    if not t:\n"
            "        continue\n"
            "    low = t.lower()\n"
            "    if any(n in low for n in _NOISE):\n"
            "        continue\n"
            "    if any(low.endswith(tr) for tr in _TRAIL):\n"
            "        continue\n"
            "    if 5 <= len(t) <= 120:\n"
            "        title = t\n"
            "        break\n"
            "out_text = title + '\\n'\n"
        ),
        "instructions": [
            "In a terminal, extract the text of page 1 of ~/Desktop/arxiv-1706-03762.pdf with poppler's pdftotext (`pdftotext -f 1 -l 1 ~/Desktop/arxiv-1706-03762.pdf -`), then write the paper's title to ~/Desktop/arxiv_title.txt followed by a single newline. From that page-1 text take the first non-empty line of length 5..120 characters, but skip any line whose lowercased text contains 'provided', 'reproduce', 'scholarly', or 'journalistic' (Google attribution boilerplate) AND skip lines that end in sentence-continuation words (' by', ' to', ' or', ' and', ' of', ' for', ' in').",
        ],
    },
    {
        "template_id": "multi_arxiv_pdf_count_pages",
        "asset_rel": "docs/pdf/arxiv-1810-04805.pdf",
        "pdf_basename": "arxiv-1810-04805.pdf",
        "out_basename": "arxiv_pagecount.txt",
        "extract_py": (
            "import subprocess\n"
            "_info = subprocess.run(['pdfinfo', src], capture_output=True, text=True).stdout\n"
            "n = 0\n"
            "for _ln in _info.splitlines():\n"
            "    if _ln.startswith('Pages:'):\n"
            "        n = int(_ln.split()[1])\n"
            "        break\n"
            "out_text = str(n) + '\\n'\n"
        ),
        "instructions": [
            "In a terminal, use poppler's pdfinfo to read the page count of ~/Desktop/arxiv-1810-04805.pdf (e.g. `pdfinfo ~/Desktop/arxiv-1810-04805.pdf | awk '/^Pages:/ {print $2}'`) and write that integer (decimal, no commas) followed by a single newline to ~/Desktop/arxiv_pagecount.txt.",
        ],
    },
    # Mass-add — additional arxiv PDF extract variants.
    {
        "template_id": "multi_arxiv_pdf_extract_title_gpt3",
        "asset_rel": "docs/pdf/arxiv-2005-14165.pdf",
        "pdf_basename": "arxiv-2005-14165.pdf",
        "out_basename": "gpt3_title.txt",
        # Validation fix: same noise-skip + sentence-continuation filter as
        # multi_arxiv_pdf_extract_title to align gold with the semantic
        # title agents write (no Google attribution on this paper, but the
        # heuristic stays uniform across sibling specs).
        "extract_py": (
            "import subprocess\n"
            "page0 = subprocess.run(['pdftotext', '-f', '1', '-l', '1', src, '-'], capture_output=True, text=True).stdout\n"
            "_NOISE = ('provided', 'reproduce', 'scholarly', 'journalistic')\n"
            "_TRAIL = (' by', ' to', ' or', ' and', ' of', ' for', ' in')\n"
            "title = ''\n"
            "for ln in page0.splitlines():\n"
            "    t = ln.strip()\n"
            "    if not t:\n"
            "        continue\n"
            "    low = t.lower()\n"
            "    if any(n in low for n in _NOISE):\n"
            "        continue\n"
            "    if any(low.endswith(tr) for tr in _TRAIL):\n"
            "        continue\n"
            "    if 5 <= len(t) <= 120:\n"
            "        title = t\n"
            "        break\n"
            "out_text = title + '\\n'\n"
        ),
        "instructions": [
            "In a terminal, extract the text of page 1 of ~/Desktop/arxiv-2005-14165.pdf with poppler's pdftotext (`pdftotext -f 1 -l 1 ~/Desktop/arxiv-2005-14165.pdf -`), then write the paper's title to ~/Desktop/gpt3_title.txt followed by a single newline. From that page-1 text take the first non-empty line of length 5..120 characters, but skip any line whose lowercased text contains 'provided', 'reproduce', 'scholarly', or 'journalistic' AND skip lines ending in sentence-continuation words (' by', ' to', ' or', ' and', ' of', ' for', ' in').",
        ],
    },
    {
        "template_id": "multi_arxiv_pdf_count_pages_pix2pix",
        "asset_rel": "docs/pdf/arxiv-1611-07004.pdf",
        "pdf_basename": "arxiv-1611-07004.pdf",
        "out_basename": "pix2pix_pagecount.txt",
        "extract_py": (
            "import subprocess\n"
            "_info = subprocess.run(['pdfinfo', src], capture_output=True, text=True).stdout\n"
            "n = 0\n"
            "for _ln in _info.splitlines():\n"
            "    if _ln.startswith('Pages:'):\n"
            "        n = int(_ln.split()[1])\n"
            "        break\n"
            "out_text = str(n) + '\\n'\n"
        ),
        "instructions": [
            "In a terminal, use poppler's pdfinfo to read the page count of ~/Desktop/arxiv-1611-07004.pdf (e.g. `pdfinfo ~/Desktop/arxiv-1611-07004.pdf | awk '/^Pages:/ {print $2}'`) and write that integer followed by a single newline to ~/Desktop/pix2pix_pagecount.txt.",
        ],
    },
    {
        "template_id": "multi_arxiv_pdf_first200_chars_attention",
        "asset_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "pdf_basename": "arxiv-1706-03762.pdf",
        "out_basename": "attention_first200.txt",
        "extract_py": (
            "import subprocess\n"
            "page0 = subprocess.run(['pdftotext', '-f', '1', '-l', '1', src, '-'], capture_output=True, text=True).stdout\n"
            "out_text = page0[:200] + '\\n'\n"
        ),
        "instructions": [
            "In a terminal, extract the text of page 1 of ~/Desktop/arxiv-1706-03762.pdf with poppler's pdftotext (`pdftotext -f 1 -l 1 ~/Desktop/arxiv-1706-03762.pdf -`) and write the first 200 characters of that output (verbatim) followed by a single newline to ~/Desktop/attention_first200.txt.",
        ],
    },
]


def _make_arxiv_pdf_template(spec: dict) -> SynthTemplate:
    template_id  = spec["template_id"]
    asset_rel    = spec["asset_rel"]
    pdf_path     = f"{_DESKTOP}/{spec['pdf_basename']}"
    out_path     = f"{_DESKTOP}/{spec['out_basename']}"
    expected     = f"{_TMP}/expected_{template_id}.txt"
    extract_py   = spec["extract_py"]
    instructions = spec["instructions"]

    init_pdf = _stage_asset(asset_rel, pdf_path)

    # Gold: same extraction logic, write to expected.
    gold_py = textwrap.dedent(f"""\
        import os
        src = {pdf_path!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        {textwrap.indent(extract_py, '        ').lstrip()}
        with open(path, 'w', encoding='utf-8') as f:
            f.write(out_text)
        """)
    init_steps = [
        init_pdf,
        _python_step(gold_py),
        # validation gap-close: terminal-driven (python PDF extract) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]

    # Oracle: same logic against `out_path` (so oracle PASSES the eval).
    oracle_py = textwrap.dedent(f"""\
        src = {pdf_path!r}
        out_path = {out_path!r}
        {textwrap.indent(extract_py, '        ').lstrip()}
        with open(out_path, 'w', encoding='utf-8') as f:
            f.write(out_text)
        """)
    oracle_steps = [_python_step(oracle_py)]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": out_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.txt"},
        # Revived tool-neutrally: gold + oracle + agent all extract via poppler
        # (pdftotext/pdfinfo, both on the agent PATH) instead of pymupdf. Collapse
        # whitespace so the agent's terminal pdftotext output need not byte-match
        # the gold's layout exactly (compare_text_file is otherwise exact).
        "options": {"ignore_blanks": True},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="real_asset_pdf",
    )


# ---- Cat I.10 — Multi-CSV concat to multi-sheet xlsx ----------------------
# Stage N CSVs, build a source xlsx with one sheet PER CSV (LO-normalized).
# Gold = same multi-sheet xlsx. The agent's task is to import each CSV into
# its own named sheet of an xlsx. Eval: compare_table sheet_data per sheet.

_CSV_CONCAT_XLSX_SPECS: list[dict] = [
    {
        "template_id": "multi_csv_concat_world_to_xlsx_2sheet",
        "xlsx_basename": "world_indicators.xlsx",
        # (asset_rel, basename, sheet_name)
        "csvs": [
            ("data/csv/world-gdp-2022.csv",        "world-gdp-2022.csv",        "GDP"),
            ("data/csv/world-population-2022.csv", "world-population-2022.csv", "Population"),
        ],
        "instructions": [
            "Using the two CSVs you're currently viewing in Calc (~/Desktop/world-gdp-2022.csv and ~/Desktop/world-population-2022.csv), build a single workbook ~/Desktop/world_indicators.xlsx with two sheets: 'GDP' (contents of world-gdp-2022.csv, header + all rows) and 'Population' (contents of world-population-2022.csv, header + all rows). Save the workbook.",
        ],
    },
    {
        "template_id": "multi_csv_concat_us_states_to_xlsx_3sheet",
        "xlsx_basename": "us_state_facts.xlsx",
        "csvs": [
            ("data/csv/us-population-states.csv",      "us-population-states.csv",      "Population"),
            ("data/csv/us-state-median-income.csv",    "us-state-median-income.csv",    "Income"),
            ("data/csv/us-housing-starts.csv",         "us-housing-starts.csv",         "HousingStarts"),
        ],
        "instructions": [
            "Using the three CSVs you're currently viewing in Calc (~/Desktop/us-population-states.csv, ~/Desktop/us-state-median-income.csv, ~/Desktop/us-housing-starts.csv), build a single workbook ~/Desktop/us_state_facts.xlsx with three sheets named 'Population' (us-population-states.csv), 'Income' (us-state-median-income.csv), 'HousingStarts' (us-housing-starts.csv) — each sheet contains the full CSV (header + all rows). Save the workbook.",
        ],
    },
    # Mass-add — additional multi-CSV → multi-sheet xlsx variants.
    {
        "template_id": "multi_csv_concat_macro_to_xlsx_3sheet",
        "xlsx_basename": "us_macro_indicators.xlsx",
        "csvs": [
            ("data/csv/us-gdp.csv",            "us-gdp.csv",            "GDP"),
            ("data/csv/us-inflation-cpi.csv",  "us-inflation-cpi.csv",  "CPI"),
            ("data/csv/us-unemployment.csv",   "us-unemployment.csv",   "UNRATE"),
        ],
        "instructions": [
            "Using the three CSVs you're currently viewing in Calc (~/Desktop/us-gdp.csv, ~/Desktop/us-inflation-cpi.csv, ~/Desktop/us-unemployment.csv), build a single workbook ~/Desktop/us_macro_indicators.xlsx with three sheets named 'GDP', 'CPI', 'UNRATE' (each contains the full CSV, header + all rows). Save the workbook.",
        ],
    },
    {
        "template_id": "multi_csv_concat_rates_to_xlsx_2sheet",
        "xlsx_basename": "us_rates.xlsx",
        "csvs": [
            ("data/csv/us-fed-funds-rate.csv", "us-fed-funds-rate.csv", "FedFunds"),
            ("data/csv/us-mortgage-30yr.csv",  "us-mortgage-30yr.csv",  "Mortgage30Yr"),
        ],
        "instructions": [
            "Using the two CSVs you're currently viewing in Calc (~/Desktop/us-fed-funds-rate.csv and ~/Desktop/us-mortgage-30yr.csv), build ~/Desktop/us_rates.xlsx with two sheets: 'FedFunds' (us-fed-funds-rate.csv) and 'Mortgage30Yr' (us-mortgage-30yr.csv) — each contains the full CSV. Save the workbook.",
        ],
    },
]


def _make_csv_concat_xlsx_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    xlsx_basename = spec["xlsx_basename"]
    csvs          = spec["csvs"]
    instructions  = spec["instructions"]

    xlsx_path = f"{_DESKTOP}/{xlsx_basename}"
    expected  = f"{_TMP}/expected_{template_id}.xlsx"

    init_csv_steps = [
        _stage_asset(rel, f"{_DESKTOP}/{base}") for (rel, base, _s) in csvs
    ]
    csv_paths   = [f"{_DESKTOP}/{base}" for (_r, base, _s) in csvs]
    sheet_names = [s for (_r, _b, s) in csvs]

    # Gold xlsx: one sheet per CSV, LO-normalized.
    gold_py = textwrap.dedent(f"""\
        import os, csv
        from openpyxl import Workbook
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        srcs  = {csv_paths!r}
        names = {sheet_names!r}
        wb = Workbook()
        # Drop the default sheet, then add one per CSV.
        wb.remove(wb.active)
        for src, name in zip(srcs, names):
            ws = wb.create_sheet(title=name)
            with open(src, newline='') as f:
                for row in csv.reader(f):
                    ws.append(row)
        wb.save(path)
        # _LO_NORMALIZE_TAIL — match LO save format (per AGENTS.md F1).
        import subprocess as _sp, tempfile as _tf, shutil as _sh
        _td = _tf.mkdtemp()
        try:
            _sp.run(["soffice", "--headless", "--norestore", "--nofirststartwizard",
                     "--convert-to", "xlsx", "--outdir", _td, path],
                    capture_output=True, env={{**os.environ, "DISPLAY": ":1"}}, timeout=120)
            _conv = os.path.join(_td, os.path.basename(path))
            if os.path.exists(_conv):
                _sh.copy(_conv, path)
        finally:
            _sh.rmtree(_td, ignore_errors=True)
        """)
    init_steps = init_csv_steps + [
        _python_step(gold_py),
        # validation: target xlsx doesn't exist yet at
        # task start (only built by oracle), so use the `[csv1, ...] → csvN`
        # variant — all CSVs are open in Calc, the last is focused, agent
        # saves-as into target xlsx. Bypasses `open_command` auto-launch.
        *_multi_app_preopen(
            background=csv_paths[:-1],
            foreground=csv_paths[-1],
        ),
    ]
    oracle_steps = [_execute(f"cp '{expected}' '{xlsx_path}'")]

    # Eval: one sheet_data rule per sheet (matched by name).
    rules = [
        {"type": "sheet_data",
         "sheet_idx0": f"RN{name}",
         "sheet_idx1": f"EN{name}"}
        for name in sheet_names
    ]
    evaluator = {
        "func": "compare_table",
        "result": {"type": "vm_file", "path": xlsx_path, "dest": "result.xlsx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.xlsx"},
        "options": {"rules": rules},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", xlsx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_table",
        setup_class="real_asset_csv_xlsx_multi",
    )


# ---- Cat I.11 — Code asset → docx -----------------------------------------
# Stage a code file, build a source docx with a "TODO: paste code here"
# placeholder, and a gold docx where the placeholder is replaced with one
# paragraph per source-code line. Eval: compare_docx_files.
# We project only the FIRST N lines of the source so the docx stays small.

_CODE_TO_DOCX_SPECS: list[dict] = [
    {
        "template_id": "multi_code_to_docx_flask_app",
        "asset_rel": "code/python/flask-app.py",
        "code_basename": "flask-app.py",
        "docx_basename": "flask_app_excerpt.docx",
        "head_lines": 50,
        "docx_pre": [
            ("Title", "Flask App — Source Excerpt"),
            (None, "The paragraphs below quote the first 50 lines of flask-app.py verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append the FIRST 50 lines of ~/Desktop/flask-app.py (the python source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/flask_app_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
    {
        "template_id": "multi_code_to_docx_gin_handler",
        "asset_rel": "code/go/gin-handler.go",
        "code_basename": "gin-handler.go",
        "docx_basename": "gin_handler_excerpt.docx",
        "head_lines": 40,
        "docx_pre": [
            ("Title", "Gin Handler — Source Excerpt"),
            (None, "The paragraphs below quote the first 40 lines of gin-handler.go verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append the FIRST 40 lines of ~/Desktop/gin-handler.go (the Go source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/gin_handler_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
    # Mass-add — additional code → docx variants.
    {
        "template_id": "multi_code_to_docx_django_utils",
        "asset_rel": "code/python/django-utils.py",
        "code_basename": "django-utils.py",
        "docx_basename": "django_utils_excerpt.docx",
        "head_lines": 45,
        "docx_pre": [
            ("Title", "Django Utils — Source Excerpt"),
            (None, "The paragraphs below quote the first 45 lines of django-utils.py verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append the FIRST 45 lines of ~/Desktop/django-utils.py (the python source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/django_utils_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
    {
        "template_id": "multi_code_to_docx_tokio_lib",
        "asset_rel": "code/rust/tokio-lib.rs",
        "code_basename": "tokio-lib.rs",
        "docx_basename": "tokio_excerpt.docx",
        "head_lines": 40,
        "docx_pre": [
            ("Title", "Tokio Lib — Source Excerpt"),
            (None, "The paragraphs below quote the first 40 lines of tokio-lib.rs verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append the FIRST 40 lines of ~/Desktop/tokio-lib.rs (the Rust source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/tokio_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
    {
        # Validation fix: express-index.js asset has only 11 lines but
        # head_lines was 35 → agent could see only 11 lines and reported
        # the task infeasible ("file has fewer lines than requested").
        # Anchor head_lines to the actual asset length.
        "template_id": "multi_code_to_docx_express_index",
        "asset_rel": "code/javascript/express-index.js",
        "code_basename": "express-index.js",
        "docx_basename": "express_excerpt.docx",
        "head_lines": 11,
        "docx_pre": [
            ("Title", "Express Index — Source Excerpt"),
            (None, "The paragraphs below quote all 11 lines of express-index.js verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append all 11 lines of ~/Desktop/express-index.js (the JavaScript source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/express_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
    {
        "template_id": "multi_code_to_docx_chi_mux",
        "asset_rel": "code/go/chi-mux.go",
        "code_basename": "chi-mux.go",
        "docx_basename": "chi_mux_excerpt.docx",
        "head_lines": 40,
        "docx_pre": [
            ("Title", "Chi Mux — Source Excerpt"),
            (None, "The paragraphs below quote the first 40 lines of chi-mux.go verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append the FIRST 40 lines of ~/Desktop/chi-mux.go (the Go source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/chi_mux_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
    {
        "template_id": "multi_code_to_docx_requests_api",
        "asset_rel": "code/python/requests-api.py",
        "code_basename": "requests-api.py",
        "docx_basename": "requests_excerpt.docx",
        "head_lines": 40,
        "docx_pre": [
            ("Title", "Requests API — Source Excerpt"),
            (None, "The paragraphs below quote the first 40 lines of requests-api.py verbatim, one paragraph per source line."),
        ],
        "instructions": [
            "Append the FIRST 40 lines of ~/Desktop/requests-api.py (the Python source you're currently viewing) to the Writer document you're currently editing at ~/Desktop/requests_excerpt.docx, one paragraph per source line (preserve leading whitespace; do not collapse blank lines). Save the document.",
        ],
    },
]


def _make_code_to_docx_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    asset_rel     = spec["asset_rel"]
    code_path     = f"{_DESKTOP}/{spec['code_basename']}"
    docx_path     = f"{_DESKTOP}/{spec['docx_basename']}"
    expected      = f"{_TMP}/expected_{template_id}.docx"
    head_lines    = spec["head_lines"]
    docx_pre      = spec["docx_pre"]
    instructions  = spec["instructions"]

    init_code = _stage_asset(asset_rel, code_path)
    init_docx = _build_docx_step(docx_path, docx_pre)

    # Gold: read first N lines of the staged code, write docx with pre
    # paragraphs + one paragraph per line.
    gold_py = textwrap.dedent(f"""\
        import os
        from docx import Document
        src = {code_path!r}
        path = {expected!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        with open(src, encoding='utf-8') as f:
            lines = f.read().splitlines()[:{head_lines!r}]
        doc = Document()
        for style, text in {docx_pre!r}:
            if style:
                doc.add_paragraph(text, style=style)
            else:
                doc.add_paragraph(text)
        for ln in lines:
            doc.add_paragraph(ln)
        doc.save(path)
        """)
    init_steps = [
        init_code, init_docx, _python_step(gold_py),
        # validation: [code (text editor bg)] →
        # docx (Writer fg). Code file opens via xdg-open (typically GNOME
        # Text/gedit on lite.osworld); docx ends focused so agent's first
        # action lands on Writer.
        *_multi_app_preopen(background=[code_path], foreground=docx_path),
    ]
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(docx_path, expected, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": docx_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected, "dest": "expected.docx"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="real_asset_code_docx",
    )


# ---------------------------------------------------------------------------
# Cat J — pandoc rows (added 2026-05-10 once Dockerfile installs pandoc).
# ---------------------------------------------------------------------------
# Pattern: pre_config materializes input file (inline write OR `_stage_asset`)
# AND the gold file by invoking pandoc with the same flags as the oracle.
# Oracle = `cp gold sink` (single shell action). Eval = compare_text_file
# (pandoc-emitted markdown / txt) or compare_docx_files (docx output —
# `examine_text` flag is intrinsic to compare_docx_files since it diffs
# paragraph text only). Bare `pandoc` command (assumed in PATH).
# Skill family: doc-format-conversion (mirrors A1/A2 of the eval taxonomy).
# ---------------------------------------------------------------------------


def _pandoc_md_to_docx_template(spec: dict) -> SynthTemplate:
    """Markdown → docx via pandoc.

    Spec keys: template_id, src_basename, dst_basename, md_lines, instructions.
    Pandoc invocation: `pandoc -s -o out.docx in.md` (standalone). The gold is
    produced by the SAME `pandoc` call against the same input, so any pandoc-
    version-specific style decisions cancel out. Eval = compare_docx_files
    (paragraph-text diff — robust to docx-style metadata drift).
    """
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.docx"
    md_body       = "\n".join(spec["md_lines"]) + "\n"
    instructions  = spec["instructions"]

    init_steps = [
        _write_text_step(src_path, md_body),
        # Pandoc-build the gold from the staged source.
        _execute(
            f"pandoc -s -o '{expected_path}' '{src_path}'"
        ),
        # validation gap-close: terminal-driven (pandoc shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(dst_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,   # terminal task — no app launch
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="pandoc_md_source",
    )


_PANDOC_MD_TO_DOCX_SPECS: list[dict] = [
    {
        "template_id": "pandoc_md_to_docx_meeting_notes",
        "src_basename": "meeting_notes.md",
        "dst_basename": "meeting_notes.docx",
        "md_lines": [
            "# Project Sync 2026-05-09",
            "",
            "## Attendees",
            "",
            "- Alice",
            "- Bob",
            "- Carol",
            "",
            "## Agenda",
            "",
            "1. Roadmap review",
            "2. Q3 priorities",
            "3. Blockers",
            "",
            "## Notes",
            "",
            "Pipeline migration is on track for next sprint.",
        ],
        "instructions": [
            "I have a markdown meeting-notes file at ~/Desktop/meeting_notes.md. Please use pandoc from a terminal to convert it into a standalone Word document at ~/Desktop/meeting_notes.docx (a self-contained docx, not a fragment).",
        ],
    },
    {
        "template_id": "pandoc_md_to_docx_changelog",
        "src_basename": "changelog.md",
        "dst_basename": "changelog.docx",
        "md_lines": [
            "# Changelog",
            "",
            "## Version 1.2.0",
            "",
            "- Added pandoc support for markdown → docx conversion",
            "- Improved error reporting on pipeline failures",
            "- Removed deprecated `--legacy` flag",
            "",
            "## Version 1.1.0",
            "",
            "- Initial release of the conversion subsystem",
        ],
        "instructions": [
            "From a terminal, take the markdown changelog at ~/Desktop/changelog.md and produce a standalone Word document version at ~/Desktop/changelog.docx using pandoc.",
        ],
    },
    {
        "template_id": "pandoc_md_to_docx_recipe",
        "src_basename": "recipe.md",
        "dst_basename": "recipe.docx",
        "md_lines": [
            "# Banana Bread",
            "",
            "## Ingredients",
            "",
            "- 3 ripe bananas",
            "- 1/2 cup butter",
            "- 1 cup sugar",
            "- 1 egg",
            "- 1 1/2 cups flour",
            "- 1 teaspoon baking soda",
            "",
            "## Steps",
            "",
            "1. Preheat oven to 175 C.",
            "2. Mash bananas, mix in butter, sugar and egg.",
            "3. Stir in flour and baking soda.",
            "4. Pour into greased loaf pan; bake 60 minutes.",
        ],
        "instructions": [
            "I'd like the banana-bread recipe at ~/Desktop/recipe.md formatted as a Word document. From a terminal, run pandoc to turn it into a standalone docx saved at ~/Desktop/recipe.docx.",
        ],
    },
    # Mass-add — additional pandoc md → docx variants.
    # Pruned (pandoc_md_to_docx_readme): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_md_to_docx_readme",
#         "src_basename": "readme.md",
#         "dst_basename": "readme.docx",
#         "md_lines": [
#             "# Project Phoenix",
#             "",
#             "Phoenix is a high-throughput data ingestion service.",
#             "",
#             "## Features",
#             "",
#             "- Idempotent ingestion via deduplication keys",
#             "- Streaming JSON parsing for low memory footprint",
#             "- Backpressure-aware downstream sinks",
#             "",
#             "## Getting Started",
#             "",
#             "```",
#             "$ phoenix run --config config.yaml",
#             "```",
#         ],
#         "instructions": [
#             "In a terminal, convert ~/Desktop/readme.md to a standalone docx at ~/Desktop/readme.docx using `pandoc -s -o ~/Desktop/readme.docx ~/Desktop/readme.md`.",
#         ],
#     },
    # Pruned (pandoc_md_to_docx_postmortem): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_md_to_docx_postmortem",
#         "src_basename": "postmortem.md",
#         "dst_basename": "postmortem.docx",
#         "md_lines": [
#             "# Postmortem — Outage 2026-04-22",
#             "",
#             "## Summary",
#             "",
#             "User-facing requests failed for 38 minutes due to a misconfigured load balancer health probe.",
#             "",
#             "## Timeline",
#             "",
#             "1. 14:02 UTC — alert fires",
#             "2. 14:11 UTC — root cause identified",
#             "3. 14:40 UTC — service restored",
#             "",
#             "## Action Items",
#             "",
#             "- Add health-probe simulation to deploy pipeline",
#             "- Document load balancer config in the runbook",
#         ],
#         "instructions": [
#             "In a terminal, convert ~/Desktop/postmortem.md to a standalone docx at ~/Desktop/postmortem.docx via `pandoc -s -o ~/Desktop/postmortem.docx ~/Desktop/postmortem.md`.",
#         ],
#     },
    # Pruned (pandoc_md_to_docx_design_doc): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_md_to_docx_design_doc",
#         "src_basename": "design_doc.md",
#         "dst_basename": "design_doc.docx",
#         "md_lines": [
#             "# Design Doc — Cache Sharding",
#             "",
#             "## Goals",
#             "",
#             "- Reduce p99 latency on hot keys to under 5ms",
#             "- Maintain consistency on writes",
#             "",
#             "## Non-Goals",
#             "",
#             "- Replacing the existing cold-storage tier",
#             "- Adding global multi-region replication",
#             "",
#             "## Proposal",
#             "",
#             "Shard by key prefix using rendezvous hashing across 32 stripe shards.",
#         ],
#         "instructions": [
#             "In a terminal, convert ~/Desktop/design_doc.md to ~/Desktop/design_doc.docx via `pandoc -s -o ~/Desktop/design_doc.docx ~/Desktop/design_doc.md`.",
#         ],
#     },
    # Pruned (pandoc_md_to_docx_announcement): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_md_to_docx_announcement",
#         "src_basename": "announcement.md",
#         "dst_basename": "announcement.docx",
#         "md_lines": [
#             "# Office Reopening Announcement",
#             "",
#             "Dear team,",
#             "",
#             "We are excited to welcome everyone back to the renovated downtown office on May 20.",
#             "",
#             "## What's new",
#             "",
#             "- Hot-desking on floors 3 and 4",
#             "- Phone booths on every floor",
#             "- Open kitchen and barista bar on the ground floor",
#             "",
#             "Please RSVP by May 15.",
#         ],
#         "instructions": [
#             "In a terminal, convert ~/Desktop/announcement.md to ~/Desktop/announcement.docx using `pandoc -s -o ~/Desktop/announcement.docx ~/Desktop/announcement.md`.",
#         ],
#     },
]


def _pandoc_csv_to_md_template() -> SynthTemplate:
    """Inline CSV → markdown table via pandoc.

    Pandoc reads CSV with `-f csv -t markdown`. Eval = compare_text_file.
    """
    template_id   = "pandoc_csv_to_md_inventory"
    src_path      = f"{_DESKTOP}/inventory.csv"
    dst_path      = f"{_DESKTOP}/inventory.md"
    expected_path = f"{_TMP}/expected_{template_id}.txt"

    csv_body = (
        "SKU,Item,Quantity\n"
        "A001,Widget,120\n"
        "A002,Gadget,60\n"
        "A003,Sprocket,200\n"
        "A004,Bolt,1500\n"
    )

    init_steps = [
        _write_text_step(src_path, csv_body),
        _execute(
            f"pandoc -f csv -t markdown -o '{expected_path}' '{src_path}'"
        ),
        # validation gap-close: terminal-driven (pandoc shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
    }
    instructions = [
        "From a terminal, use pandoc to read ~/Desktop/inventory.csv (CSV input) and render it as a markdown table at ~/Desktop/inventory.md.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="pandoc_csv_source",
    )


def _pandoc_html_to_md_template(spec: dict) -> SynthTemplate:
    """Wikipedia HTML → markdown via pandoc.

    Spec keys: template_id, asset_rel, html_basename, md_basename,
    instructions. Both gold and oracle invoke `pandoc -f html -t
    markdown_strict --wrap=none -o out.md in.html`. Eval =
    compare_text_file (byte-equal, since same pandoc invocation on the
    same input is deterministic).
    """
    template_id   = spec["template_id"]
    asset_rel     = spec["asset_rel"]
    src_path      = f"{_DESKTOP}/{spec['html_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['md_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.txt"
    instructions  = spec["instructions"]

    init_steps = [
        _stage_asset(asset_rel, src_path),
        _execute(
            f"pandoc -f html -t markdown_strict --wrap=none "
            f"-o '{expected_path}' '{src_path}'"
        ),
        # validation gap-close: terminal-driven (pandoc shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="pandoc_html_source",
    )


_PANDOC_HTML_TO_MD_SPECS: list[dict] = [
    {
        "template_id": "pandoc_html_to_md_apollo",
        "asset_rel": "html/wikipedia/apollo-program.html",
        "html_basename": "apollo-program.html",
        "md_basename": "apollo-program.md",
        "instructions": [
            "From a terminal, use pandoc to convert the Wikipedia HTML page at ~/Desktop/apollo-program.html into markdown saved at ~/Desktop/apollo-program.md. Use the strict markdown flavor and turn off line wrapping so each paragraph stays on one line.",
        ],
    },
    {
        "template_id": "pandoc_html_to_md_octopus",
        "asset_rel": "html/wikipedia/octopus.html",
        "html_basename": "octopus.html",
        "md_basename": "octopus.md",
        "instructions": [
            "From a terminal, please convert the Wikipedia octopus article at ~/Desktop/octopus.html into a markdown file at ~/Desktop/octopus.md using pandoc. Pick the strict markdown variant and disable hard line-wrapping.",
        ],
    },
    # Mass-add — additional pandoc HTML → markdown variants.
    {
        "template_id": "pandoc_html_to_md_coffee",
        "asset_rel": "html/wikipedia/coffee.html",
        "html_basename": "coffee.html",
        "md_basename": "coffee.md",
        "instructions": [
            "From a terminal, run pandoc on the saved Wikipedia coffee article at ~/Desktop/coffee.html to produce a strict-markdown rendering at ~/Desktop/coffee.md. Make sure paragraphs are not hard-wrapped.",
        ],
    },
    {
        "template_id": "pandoc_html_to_md_eiffel",
        "asset_rel": "html/wikipedia/eiffel-tower.html",
        "html_basename": "eiffel-tower.html",
        "md_basename": "eiffel-tower.md",
        "instructions": [
            "From a terminal, please use pandoc to convert the Wikipedia Eiffel Tower article at ~/Desktop/eiffel-tower.html into a markdown document at ~/Desktop/eiffel-tower.md. Use strict markdown and skip line wrapping.",
        ],
    },
    # Pruned (pandoc_html_to_md_pizza): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_html_to_md_pizza",
#         "asset_rel": "html/wikipedia/pizza.html",
#         "html_basename": "pizza.html",
#         "md_basename": "pizza.md",
#         "instructions": [
#             "In a terminal, convert ~/Desktop/pizza.html to markdown at ~/Desktop/pizza.md via `pandoc -f html -t markdown_strict --wrap=none -o ~/Desktop/pizza.md ~/Desktop/pizza.html`.",
#         ],
#     },
    # Pruned (pandoc_html_to_md_solar_system): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_html_to_md_solar_system",
#         "asset_rel": "html/wikipedia/solar-system.html",
#         "html_basename": "solar-system.html",
#         "md_basename": "solar-system.md",
#         "instructions": [
#             "In a terminal, convert ~/Desktop/solar-system.html to markdown at ~/Desktop/solar-system.md using `pandoc -f html -t markdown_strict --wrap=none -o ~/Desktop/solar-system.md ~/Desktop/solar-system.html`.",
#         ],
#     },
    # Pruned (pandoc_html_to_md_beethoven): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_html_to_md_beethoven",
#         "asset_rel": "html/wikipedia/beethoven.html",
#         "html_basename": "beethoven.html",
#         "md_basename": "beethoven.md",
#         "instructions": [
#             "In a terminal, run `pandoc -f html -t markdown_strict --wrap=none -o ~/Desktop/beethoven.md ~/Desktop/beethoven.html` to convert the Wikipedia article to markdown.",
#         ],
#     },
    # Pruned (pandoc_html_to_md_volcano): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_html_to_md_volcano",
#         "asset_rel": "html/wikipedia/volcano.html",
#         "html_basename": "volcano.html",
#         "md_basename": "volcano.md",
#         "instructions": [
#             "In a terminal, convert ~/Desktop/volcano.html to markdown at ~/Desktop/volcano.md using `pandoc -f html -t markdown_strict --wrap=none -o ~/Desktop/volcano.md ~/Desktop/volcano.html`.",
#         ],
#     },
    # Pruned (pandoc_html_to_md_yoga): format-conversion over-amped vs eval; cut to rebalance toward gap skills.
#     {
#         "template_id": "pandoc_html_to_md_yoga",
#         "asset_rel": "html/wikipedia/yoga.html",
#         "html_basename": "yoga.html",
#         "md_basename": "yoga.md",
#         "instructions": [
#             "In a terminal, run `pandoc -f html -t markdown_strict --wrap=none -o ~/Desktop/yoga.md ~/Desktop/yoga.html` to convert the Wikipedia article to markdown.",
#         ],
#     },
]


def _pandoc_docx_to_txt_template() -> SynthTemplate:
    """docx → plain text via pandoc.

    Source: a docx built via python-docx (deterministic). Pandoc plain
    output is text; eval = compare_text_file.
    """
    template_id   = "pandoc_docx_to_txt_report"
    src_path      = f"{_DESKTOP}/quarterly_report.docx"
    dst_path      = f"{_DESKTOP}/quarterly_report.txt"
    expected_path = f"{_TMP}/expected_{template_id}.txt"

    paragraphs: list[tuple[str | None, str]] = [
        ("Title", "Quarterly Report — Q2 2026"),
        ("Heading 1", "Revenue"),
        (None, "Revenue grew 14% year-over-year, driven by strong adoption in the EU."),
        ("Heading 1", "Costs"),
        (None, "Operating costs rose 3%, primarily from infrastructure spend."),
        ("Heading 1", "Outlook"),
        (None, "We expect continued growth into Q3 with margin expansion."),
    ]

    init_steps = [
        _build_docx_step(src_path, paragraphs),
        _execute(
            f"pandoc -t plain -o '{expected_path}' '{src_path}'"
        ),
        # validation gap-close: terminal-driven (pandoc shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
    }
    instructions = [
        "From a terminal, use pandoc to read the Word document at ~/Desktop/quarterly_report.docx and emit a plain-text rendering at ~/Desktop/quarterly_report.txt.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="pandoc_docx_source",
    )


# ---------------------------------------------------------------------------
# Cat K — ImageMagick rows (added 2026-05-10 once Dockerfile installs IM).
# ---------------------------------------------------------------------------
# Pattern: stage a real photo via `_stage_asset`, then both the gold and the
# oracle invoke the SAME `convert` (ImageMagick 6) command on the same input
# so the comparison is deterministic. Eval = `compare_images` (SSIM-based —
# returns 1.0 when oracle copies the gold byte-for-byte). Bare `convert`
# command (assumed in PATH; ImageMagick 6 ships `convert` as the canonical
# CLI; v7 also provides a `convert` alias).
# ---------------------------------------------------------------------------

# Reuse the same photo pool the Cartesian topic factories use; per-row we
# pick a verified asset path so codegen can statically validate via
# `_stage_asset`.

_IM_PHOTO_SPECS: list[dict] = [
    # Resize 800x600 — landscape photo
    {
        "template_id": "im_resize_800x600_jupiter",
        "asset_rel": "photos/space/jupiter-full-disk.jpg",
        "src_basename": "jupiter-full-disk.jpg",
        "dst_basename": "jupiter-resized.jpg",
        "convert_args": "-resize 800x600",
        "instructions": [
            "I have a Jupiter photo at ~/Desktop/jupiter-full-disk.jpg. Please use ImageMagick from a terminal to scale it down so it fits inside an 800x600 box (preserve the aspect ratio), and write the result to ~/Desktop/jupiter-resized.jpg.",
        ],
    },
    # Resize 200x200 thumbnail — wildlife
    {
        "template_id": "im_resize_thumb200_tiger",
        "asset_rel": "photos/wildlife/tiger-closeup.jpg",
        "src_basename": "tiger-closeup.jpg",
        "dst_basename": "tiger-thumb.jpg",
        "convert_args": "-resize 200x200",
        "instructions": [
            "Open a terminal and use ImageMagick to produce a 200x200 thumbnail of the tiger photo at ~/Desktop/tiger-closeup.jpg. Save it as ~/Desktop/tiger-thumb.jpg.",
        ],
    },
    # Resize to specific width (preserve aspect) — architecture
    {
        "template_id": "im_resize_w400_skyscraper",
        "asset_rel": "photos/architecture/skyscraper-glass.jpg",
        "src_basename": "skyscraper-glass.jpg",
        "dst_basename": "skyscraper-w400.jpg",
        "convert_args": "-resize 400",
        "instructions": [
            "I need a narrower preview of the skyscraper shot at ~/Desktop/skyscraper-glass.jpg. From a terminal, resize it down to 400 pixels wide (let the height adjust automatically) using ImageMagick and save the result as ~/Desktop/skyscraper-w400.jpg.",
        ],
    },
    # Rotate 90
    {
        "template_id": "im_rotate_90_horse",
        "asset_rel": "photos/wildlife/horse-meadow.jpg",
        "src_basename": "horse-meadow.jpg",
        "dst_basename": "horse-rotated90.jpg",
        "convert_args": "-rotate 90",
        "instructions": [
            "From a terminal, rotate the horse photo at ~/Desktop/horse-meadow.jpg 90 degrees clockwise using ImageMagick. The rotated copy should be written to ~/Desktop/horse-rotated90.jpg.",
        ],
    },
    # Rotate 180
    {
        "template_id": "im_rotate_180_pizza",
        "asset_rel": "photos/food/pizza-dish.jpg",
        "src_basename": "pizza-dish.jpg",
        "dst_basename": "pizza-rotated180.jpg",
        "convert_args": "-rotate 180",
        "instructions": [
            "The pizza photo at ~/Desktop/pizza-dish.jpg was scanned upside-down. Please use ImageMagick from a terminal to flip it 180 degrees and save the corrected version as ~/Desktop/pizza-rotated180.jpg.",
        ],
    },
    # Rotate 270
    {
        "template_id": "im_rotate_270_andromeda",
        "asset_rel": "photos/nature/galaxy-andromeda.jpg",
        "src_basename": "galaxy-andromeda.jpg",
        "dst_basename": "andromeda-rotated270.jpg",
        "convert_args": "-rotate 270",
        "instructions": [
            "Using ImageMagick from a terminal, rotate the Andromeda galaxy image at ~/Desktop/galaxy-andromeda.jpg by 270 degrees clockwise and put the rotated file at ~/Desktop/andromeda-rotated270.jpg.",
        ],
    },
    # Grayscale
    {
        "template_id": "im_grayscale_house",
        "asset_rel": "photos/architecture/house-modern.jpg",
        "src_basename": "house-modern.jpg",
        "dst_basename": "house-grayscale.jpg",
        "convert_args": "-colorspace Gray",
        "instructions": [
            "I want a black-and-white version of the modern-house photo at ~/Desktop/house-modern.jpg. Use ImageMagick from a terminal to desaturate it to a grayscale image and save it as ~/Desktop/house-grayscale.jpg.",
        ],
    },
    # Grayscale — portrait
    {
        "template_id": "im_grayscale_apollo11",
        "asset_rel": "photos/portrait/astronaut-apollo11-crew.jpg",
        "src_basename": "astronaut-apollo11-crew.jpg",
        "dst_basename": "apollo11-grayscale.jpg",
        "convert_args": "-colorspace Gray",
        "instructions": [
            "We're laying out a documentary spread and need the Apollo 11 crew portrait at ~/Desktop/astronaut-apollo11-crew.jpg in monochrome. From a terminal, use ImageMagick to convert it to grayscale and save it as ~/Desktop/apollo11-grayscale.jpg.",
        ],
    },
    # Flip (vertical mirror)
    {
        "template_id": "im_flip_mars",
        "asset_rel": "photos/space/mars-curiosity-panorama.jpg",
        "src_basename": "mars-curiosity-panorama.jpg",
        "dst_basename": "mars-flipped.jpg",
        "convert_args": "-flip",
        "instructions": [
            "From a terminal, mirror the Mars Curiosity panorama at ~/Desktop/mars-curiosity-panorama.jpg along its horizontal axis (top becomes bottom) using ImageMagick. Save the flipped result as ~/Desktop/mars-flipped.jpg.",
        ],
    },
    # Flop (horizontal mirror)
    {
        "template_id": "im_flop_coffee",
        "asset_rel": "photos/food/coffee-latte.jpg",
        "src_basename": "coffee-latte.jpg",
        "dst_basename": "coffee-flopped.jpg",
        "convert_args": "-flop",
        "instructions": [
            "The cafe wants their latte photo mirrored for the menu. From a terminal, use ImageMagick to flip ~/Desktop/coffee-latte.jpg left-to-right (a horizontal mirror) and save it as ~/Desktop/coffee-flopped.jpg.",
        ],
    },
    # Crop center 400x400
    {
        "template_id": "im_crop_center400_bird",
        "asset_rel": "photos/wildlife/bird-perch.jpg",
        "src_basename": "bird-perch.jpg",
        "dst_basename": "bird-cropped.jpg",
        # `+repage` strips the canvas-offset metadata so the output is a
        # bare 400x400 image (per ImageMagick crop convention).
        "convert_args": "-gravity center -crop 400x400+0+0 +repage",
        "instructions": [
            "From a terminal, take the center 400x400 pixel square out of the bird photo at ~/Desktop/bird-perch.jpg with ImageMagick (strip the canvas offset so the output is a clean 400x400 image) and save it as ~/Desktop/bird-cropped.jpg.",
        ],
    },
    # Mass-add — additional ImageMagick variants.
    {
        "template_id": "im_resize_640_salad",
        "asset_rel": "photos/food/salad-bowl.jpg",
        "src_basename": "salad-bowl.jpg",
        "dst_basename": "salad-w640.jpg",
        "convert_args": "-resize 640",
        "instructions": [
            "For a recipe blog header, scale the salad photo at ~/Desktop/salad-bowl.jpg down to 640 pixels wide (keep the original aspect ratio). Use ImageMagick from a terminal and write the resized image to ~/Desktop/salad-w640.jpg.",
        ],
    },
    {
        "template_id": "im_grayscale_lab_bench",
        "asset_rel": "photos/medical/lab-bench.jpg",
        "src_basename": "lab-bench.jpg",
        "dst_basename": "lab-bench-grayscale.jpg",
        "convert_args": "-colorspace Gray",
        "instructions": [
            "The lab-bench photo at ~/Desktop/lab-bench.jpg needs to be grayscale for a journal figure. From a terminal, run ImageMagick to desaturate it and save the result as ~/Desktop/lab-bench-grayscale.jpg.",
        ],
    },
    {
        "template_id": "im_rotate_90_office",
        "asset_rel": "photos/office/laptop-desk.jpg",
        "src_basename": "laptop-desk.jpg",
        "dst_basename": "laptop-desk-rotated.jpg",
        "convert_args": "-rotate 90",
        "instructions": [
            "I uploaded ~/Desktop/laptop-desk.jpg in landscape but I want it portrait. From a terminal, use ImageMagick to rotate it 90 degrees clockwise and save it as ~/Desktop/laptop-desk-rotated.jpg.",
        ],
    },
    {
        "template_id": "im_flip_concert",
        "asset_rel": "photos/event/concert-stage.jpg",
        "src_basename": "concert-stage.jpg",
        "dst_basename": "concert-stage-flipped.jpg",
        "convert_args": "-flip",
        "instructions": [
            "Flip the concert-stage photo at ~/Desktop/concert-stage.jpg upside-down (vertical mirror) using ImageMagick from a terminal. The flipped image should be saved as ~/Desktop/concert-stage-flipped.jpg.",
        ],
    },
    {
        "template_id": "im_resize_thumb150_headphones",
        "asset_rel": "photos/product/headphones.jpg",
        "src_basename": "headphones.jpg",
        "dst_basename": "headphones-thumb.jpg",
        "convert_args": "-resize 150x150",
        "instructions": [
            "Our product catalog needs a 150x150 thumbnail of the headphones picture at ~/Desktop/headphones.jpg. Please use ImageMagick from a terminal to produce it and save it as ~/Desktop/headphones-thumb.jpg.",
        ],
    },
    {
        "template_id": "im_grayscale_classroom",
        "asset_rel": "photos/classroom/open-books.jpg",
        "src_basename": "open-books.jpg",
        "dst_basename": "open-books-grayscale.jpg",
        "convert_args": "-colorspace Gray",
        "instructions": [
            "The classroom photo at ~/Desktop/open-books.jpg should be grayscale for our handbook. From a terminal, use ImageMagick to desaturate it and write the result to ~/Desktop/open-books-grayscale.jpg.",
        ],
    },
    {
        "template_id": "im_rotate_180_galaxy_hubble",
        "asset_rel": "photos/nature/galaxy-hubble.jpg",
        "src_basename": "galaxy-hubble.jpg",
        "dst_basename": "galaxy-hubble-rotated180.jpg",
        "convert_args": "-rotate 180",
        "instructions": [
            "The Hubble galaxy frame at ~/Desktop/galaxy-hubble.jpg came in flipped. From a terminal, rotate it 180 degrees with ImageMagick and write the corrected image to ~/Desktop/galaxy-hubble-rotated180.jpg.",
        ],
    },
    {
        "template_id": "im_flop_portrait3",
        "asset_rel": "photos/portrait/person-headshot-3.jpg",
        "src_basename": "person-headshot-3.jpg",
        "dst_basename": "headshot-flopped.jpg",
        "convert_args": "-flop",
        "instructions": [
            "Please horizontally mirror the headshot at ~/Desktop/person-headshot-3.jpg (subject should face the opposite direction). Use ImageMagick from a terminal and save the mirrored image as ~/Desktop/headshot-flopped.jpg.",
        ],
    },
    {
        "template_id": "im_crop_center300_dunes",
        "asset_rel": "photos/landscape/desert-dunes.jpg",
        "src_basename": "desert-dunes.jpg",
        "dst_basename": "desert-dunes-cropped.jpg",
        "convert_args": "-gravity center -crop 300x300+0+0 +repage",
        "instructions": [
            "From a terminal, crop a 300x300 square out of the very center of the desert dunes picture at ~/Desktop/desert-dunes.jpg using ImageMagick (be sure to clear the canvas offset so the saved image is exactly 300x300). Save it as ~/Desktop/desert-dunes-cropped.jpg.",
        ],
    },
    # validation — wire newly-added energy/industrial/education photo categories.
    {
        "template_id": "im_grayscale_wind_turbine",
        "asset_rel": "photos/energy/wind-turbine.jpg",
        "src_basename": "wind-turbine.jpg",
        "dst_basename": "wind-turbine-gray.jpg",
        "convert_args": "-colorspace Gray",
        "instructions": [
            "We're printing a black-and-white energy-sector brochure. Use ImageMagick from a terminal to convert the wind turbine photo at ~/Desktop/wind-turbine.jpg into grayscale, and save the output to ~/Desktop/wind-turbine-gray.jpg.",
        ],
    },
    {
        "template_id": "im_resize_w400_factory_loom",
        "asset_rel": "photos/industrial/factory-loom.jpg",
        "src_basename": "factory-loom.jpg",
        "dst_basename": "factory-loom-resized.jpg",
        "convert_args": "-resize 400x",
        "instructions": [
            "I need a smaller version of the factory loom image at ~/Desktop/factory-loom.jpg for an internal report. From a terminal, use ImageMagick to set its width to 400 pixels (keep the aspect ratio intact) and save it as ~/Desktop/factory-loom-resized.jpg.",
        ],
    },
    {
        "template_id": "im_rotate_90_graduation",
        "asset_rel": "photos/education/graduation-ceremony.jpg",
        "src_basename": "graduation-ceremony.jpg",
        "dst_basename": "graduation-rotated.jpg",
        "convert_args": "-rotate 90",
        "instructions": [
            "The graduation ceremony photo at ~/Desktop/graduation-ceremony.jpg is sideways. From a terminal, use ImageMagick to rotate it 90 degrees clockwise into the proper portrait orientation, and save the corrected file as ~/Desktop/graduation-rotated.jpg.",
        ],
    },
    {
        "template_id": "im_crop_center400_solar_farm",
        "asset_rel": "photos/energy/solar-farm.jpg",
        "src_basename": "solar-farm.jpg",
        "dst_basename": "solar-farm-cropped.jpg",
        "convert_args": "-gravity center -crop 400x400+0+0",
        "instructions": [
            "From a terminal, take a 400x400 cutout centered on the middle of the solar farm photo at ~/Desktop/solar-farm.jpg using ImageMagick and save it as ~/Desktop/solar-farm-cropped.jpg.",
        ],
    },
]


def _make_im_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    asset_rel     = spec["asset_rel"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.jpg"
    convert_args  = spec["convert_args"]
    instructions  = spec["instructions"]

    init_steps = [
        _stage_asset(asset_rel, src_path),
        _execute(
            f"convert '{src_path}' {convert_args} '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (ImageMagick convert) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_images",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.jpg"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.jpg"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,   # terminal task
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_images",
        setup_class="im_photo_source",
    )


def _make_im_montage_contact_sheet() -> SynthTemplate:
    """`montage` contact sheet from 4 staged photos.

    All 4 photos are picked from one topic family for a coherent sheet.
    Both gold and oracle invoke the same `montage` command. Eval = compare_images.
    """
    template_id = "im_montage_contact_sheet_food"
    photos = [
        ("photos/food/pizza-dish.jpg",      "pizza-dish.jpg"),
        ("photos/food/coffee-latte.jpg",    "coffee-latte.jpg"),
        ("photos/food/restaurant-meal.jpg", "restaurant-meal.jpg"),
        ("photos/food/salad-bowl.jpg",      "salad-bowl.jpg"),
    ]
    src_paths     = [f"{_DESKTOP}/{base}" for (_r, base) in photos]
    dst_path      = f"{_DESKTOP}/food-contact-sheet.jpg"
    expected_path = f"{_TMP}/expected_{template_id}.jpg"

    init_steps: list[dict] = [
        _stage_asset(rel, f"{_DESKTOP}/{base}") for (rel, base) in photos
    ]
    # `montage` builds a 2x2 grid of 200x200 thumbnails with bare-frame layout.
    src_args = " ".join(f"'{p}'" for p in src_paths)
    init_steps.append(_execute(
        f"montage {src_args} -tile 2x2 -geometry 200x200+5+5 '{expected_path}'"
    ))
    # validation gap-close: terminal-driven (ImageMagick montage) — pre-open terminal.
    init_steps.extend(_terminal_preopen_steps())
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_images",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.jpg"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.jpg"},
    }
    instructions = [
        "In a terminal, build a 2x2 contact sheet of the four food photos in ~/Desktop/ (pizza-dish.jpg, coffee-latte.jpg, restaurant-meal.jpg, salad-bowl.jpg) at 200x200 each with a 5px gap, saving to ~/Desktop/food-contact-sheet.jpg. Use ImageMagick montage: `montage ~/Desktop/pizza-dish.jpg ~/Desktop/coffee-latte.jpg ~/Desktop/restaurant-meal.jpg ~/Desktop/salad-bowl.jpg -tile 2x2 -geometry 200x200+5+5 ~/Desktop/food-contact-sheet.jpg`.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_images",
        setup_class="im_montage_source",
    )


# ---------------------------------------------------------------------------
# Cat L — Poppler (`pdftotext` / `pdfinfo`) rows (added 2026-05-10).
# ---------------------------------------------------------------------------
# Pattern: stage one arxiv PDF, then both the gold and the oracle invoke
# the SAME `pdftotext` / `pdfinfo` command so eval is deterministic.
# Eval = compare_text_file. Bare `pdftotext` / `pdfinfo` (assumed in PATH).
# Skill family: PDF text extraction (mirrors arxiv-fitz template Cat I.9).
# ---------------------------------------------------------------------------

_PDFTOTEXT_SPECS: list[dict] = [
    # Page 1 only via `pdftotext -f 1 -l 1`.
    {
        "template_id": "pdftotext_page1_attention",
        "asset_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "pdf_basename": "arxiv-1706-03762.pdf",
        "out_basename": "attention_page1.txt",
        "args": "-f 1 -l 1 -layout",
        "instructions": [
            "In a terminal, extract the text of page 1 of ~/Desktop/arxiv-1706-03762.pdf to ~/Desktop/attention_page1.txt using `pdftotext -f 1 -l 1 -layout ~/Desktop/arxiv-1706-03762.pdf ~/Desktop/attention_page1.txt`.",
        ],
    },
    # All pages, default flow.
    {
        "template_id": "pdftotext_full_bert",
        "asset_rel": "docs/pdf/arxiv-1810-04805.pdf",
        "pdf_basename": "arxiv-1810-04805.pdf",
        "out_basename": "bert_full.txt",
        "args": "",
        "instructions": [
            "In a terminal, extract the full text of ~/Desktop/arxiv-1810-04805.pdf to ~/Desktop/bert_full.txt using `pdftotext ~/Desktop/arxiv-1810-04805.pdf ~/Desktop/bert_full.txt`.",
        ],
    },
    # First page with -layout (pix2pix paper).
    {
        "template_id": "pdftotext_page1_pix2pix",
        "asset_rel": "docs/pdf/arxiv-1611-07004.pdf",
        "pdf_basename": "arxiv-1611-07004.pdf",
        "out_basename": "pix2pix_page1.txt",
        "args": "-f 1 -l 1 -layout",
        "instructions": [
            "In a terminal, extract page 1 of ~/Desktop/arxiv-1611-07004.pdf preserving the layout to ~/Desktop/pix2pix_page1.txt using `pdftotext -f 1 -l 1 -layout ~/Desktop/arxiv-1611-07004.pdf ~/Desktop/pix2pix_page1.txt`.",
        ],
    },
    # Mass-add — additional pdftotext variants.
    {
        "template_id": "pdftotext_page1_gpt3",
        "asset_rel": "docs/pdf/arxiv-2005-14165.pdf",
        "pdf_basename": "arxiv-2005-14165.pdf",
        "out_basename": "gpt3_page1.txt",
        "args": "-f 1 -l 1 -layout",
        "instructions": [
            "In a terminal, extract the text of page 1 of ~/Desktop/arxiv-2005-14165.pdf to ~/Desktop/gpt3_page1.txt using `pdftotext -f 1 -l 1 -layout ~/Desktop/arxiv-2005-14165.pdf ~/Desktop/gpt3_page1.txt`.",
        ],
    },
    {
        "template_id": "pdftotext_full_attention",
        "asset_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "pdf_basename": "arxiv-1706-03762.pdf",
        "out_basename": "attention_full.txt",
        "args": "",
        "instructions": [
            "In a terminal, extract the full text of ~/Desktop/arxiv-1706-03762.pdf to ~/Desktop/attention_full.txt using `pdftotext ~/Desktop/arxiv-1706-03762.pdf ~/Desktop/attention_full.txt`.",
        ],
    },
    {
        "template_id": "pdftotext_pages1to3_bert",
        "asset_rel": "docs/pdf/arxiv-1810-04805.pdf",
        "pdf_basename": "arxiv-1810-04805.pdf",
        "out_basename": "bert_pages_1_3.txt",
        "args": "-f 1 -l 3 -layout",
        "instructions": [
            "In a terminal, extract pages 1-3 of ~/Desktop/arxiv-1810-04805.pdf preserving layout to ~/Desktop/bert_pages_1_3.txt using `pdftotext -f 1 -l 3 -layout ~/Desktop/arxiv-1810-04805.pdf ~/Desktop/bert_pages_1_3.txt`.",
        ],
    },
]


def _make_pdftotext_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    asset_rel     = spec["asset_rel"]
    pdf_path      = f"{_DESKTOP}/{spec['pdf_basename']}"
    out_path      = f"{_DESKTOP}/{spec['out_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.txt"
    args          = spec["args"]
    instructions  = spec["instructions"]

    args_part = f"{args} " if args else ""
    init_steps = [
        _stage_asset(asset_rel, pdf_path),
        _execute(
            f"pdftotext {args_part}'{pdf_path}' '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (poppler pdftotext) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{out_path}'")]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": out_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="poppler_pdf_source",
    )


def _make_pdfinfo_pagecount_template() -> SynthTemplate:
    """Use `pdfinfo` to extract page count of an arxiv PDF.

    Both gold and oracle pipe `pdfinfo | awk '/^Pages:/ {print $2}'` so the
    output is a single integer line. Eval = compare_text_file.
    """
    template_id   = "pdfinfo_pagecount_gpt3"
    asset_rel     = "docs/pdf/arxiv-2005-14165.pdf"
    pdf_path      = f"{_DESKTOP}/arxiv-2005-14165.pdf"
    out_path      = f"{_DESKTOP}/gpt3_pagecount.txt"
    expected_path = f"{_TMP}/expected_{template_id}.txt"

    init_steps = [
        _stage_asset(asset_rel, pdf_path),
        _execute(
            f"pdfinfo '{pdf_path}' | awk '/^Pages:/ {{print $2}}' "
            f"> '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (pdfinfo + awk shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{out_path}'")]

    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": out_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
    }
    instructions = [
        "From a terminal, find out exactly how many pages are in the PDF at ~/Desktop/arxiv-2005-14165.pdf and write that integer (just the number, on a single line) to ~/Desktop/gpt3_pagecount.txt. Poppler's pdfinfo utility is available.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="poppler_pdf_source",
    )


# ---------------------------------------------------------------------------
# Cat M — pdftoppm (poppler-utils) rows: PDF page → image.
# ---------------------------------------------------------------------------
# Pattern: stage an arxiv PDF, then `pdftoppm -png -r 150 -f N -l N` produces
# a single PNG page rendering. Both gold and oracle invoke the same command,
# so byte-comparable. Eval = compare_images.
# Note: pdftoppm appends `-<page-number>.png` to the output prefix; we name
# the expected/oracle files accordingly.
# ---------------------------------------------------------------------------

_PDFTOPPM_SPECS: list[dict] = [
    {
        # validation (2026-05-10): The Attention paper has 15 pages, BERT has 16.
        # pdftoppm zero-pads the page-number suffix to the digit-width of the
        # PDF's total page count (NOT to `-l N` as the old comment claimed):
        # for a ≥10-page PDF, page 1 emits as `-01.png`. Updated `pad_width`
        # to 2 and instructions to reflect the actual output filename.
        "template_id": "pdftoppm_page1_attention_150dpi",
        "asset_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "pdf_basename": "arxiv-1706-03762.pdf",
        "out_prefix_basename": "attention_page1",
        "page": 1,
        "pad_width": 2,
        "dpi": 150,
        "instructions": [
            "In a terminal, render page 1 of ~/Desktop/arxiv-1706-03762.pdf as a PNG at 150 DPI using `pdftoppm -png -r 150 -f 1 -l 1 ~/Desktop/arxiv-1706-03762.pdf ~/Desktop/attention_page1`. The resulting file will be ~/Desktop/attention_page1-01.png (pdftoppm zero-pads to the PDF's page-count digit width).",
        ],
    },
    {
        "template_id": "pdftoppm_page1_bert_100dpi",
        "asset_rel": "docs/pdf/arxiv-1810-04805.pdf",
        "pdf_basename": "arxiv-1810-04805.pdf",
        "out_prefix_basename": "bert_page1",
        "page": 1,
        "pad_width": 2,
        "dpi": 100,
        "instructions": [
            "In a terminal, render page 1 of ~/Desktop/arxiv-1810-04805.pdf as a PNG at 100 DPI using `pdftoppm -png -r 100 -f 1 -l 1 ~/Desktop/arxiv-1810-04805.pdf ~/Desktop/bert_page1`. The resulting file will be ~/Desktop/bert_page1-01.png (pdftoppm zero-pads to the PDF's page-count digit width).",
        ],
    },
    # Mass-add — additional pdftoppm variants.
    {
        "template_id": "pdftoppm_page1_pix2pix_120dpi",
        "asset_rel": "docs/pdf/arxiv-1611-07004.pdf",
        "pdf_basename": "arxiv-1611-07004.pdf",
        "out_prefix_basename": "pix2pix_page1",
        "page": 1,
        "pad_width": 2,
        "dpi": 120,
        "instructions": [
            "In a terminal, render page 1 of ~/Desktop/arxiv-1611-07004.pdf as a PNG at 120 DPI using `pdftoppm -png -r 120 -f 1 -l 1 ~/Desktop/arxiv-1611-07004.pdf ~/Desktop/pix2pix_page1`. The resulting file will be ~/Desktop/pix2pix_page1-01.png.",
        ],
    },
    {
        "template_id": "pdftoppm_page2_attention_75dpi",
        "asset_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "pdf_basename": "arxiv-1706-03762.pdf",
        "out_prefix_basename": "attention_page2",
        "page": 2,
        "pad_width": 2,
        "dpi": 75,
        "instructions": [
            "In a terminal, render page 2 of ~/Desktop/arxiv-1706-03762.pdf as a PNG at 75 DPI using `pdftoppm -png -r 75 -f 2 -l 2 ~/Desktop/arxiv-1706-03762.pdf ~/Desktop/attention_page2`. The resulting file will be ~/Desktop/attention_page2-02.png.",
        ],
    },
]


def _make_pdftoppm_template(spec: dict) -> SynthTemplate:
    template_id        = spec["template_id"]
    asset_rel          = spec["asset_rel"]
    pdf_path           = f"{_DESKTOP}/{spec['pdf_basename']}"
    out_prefix         = f"{_DESKTOP}/{spec['out_prefix_basename']}"
    page               = spec["page"]
    pad_width          = spec.get("pad_width", 1)
    dpi                = spec["dpi"]
    instructions       = spec["instructions"]

    # pdftoppm appends `-<page>.png` to the prefix, zero-padded to the
    # digit-width of the PDF's total page count (e.g. a 15-page PDF
    # produces `-01.png` for page 1). Caller passes `pad_width` to match.
    page_suffix   = str(page).zfill(pad_width)
    out_path      = f"{out_prefix}-{page_suffix}.png"
    expected_pref = f"{_TMP}/expected_{template_id}"
    expected_path = f"{expected_pref}-{page_suffix}.png"

    init_steps = [
        _stage_asset(asset_rel, pdf_path),
        _execute(
            f"pdftoppm -png -r {dpi} -f {page} -l {page} "
            f"'{pdf_path}' '{expected_pref}'"
        ),
        # validation gap-close: terminal-driven (pdftoppm shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{out_path}'")]

    evaluator = {
        "func": "compare_images",
        "result": {"type": "vm_file", "path": out_path, "dest": "result.png"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.png"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_images",
        setup_class="poppler_pdf_source",
    )


# ---------------------------------------------------------------------------
# Cat N — compare_pdfs rows (added validation: close eval=5/synth=0 gap).
# ---------------------------------------------------------------------------
# Pattern: pre_config materializes a source office file (docx / xlsx / pptx)
# AND a gold pdf produced by `soffice --headless --convert-to pdf`. The
# agent's task is to convert the source to pdf at the sink path, also via
# soffice. Eval = compare_pdfs (text-content fuzz-ratio over extracted text).
# Pdftk-merge / pdftk-extract rows stage two arxiv PDFs and use `pdftk`
# (pdftk-java) to combine / extract pages. Trivial-pass guard: source != gold,
# pre_config does NOT pre-create the agent's sink path.
# ---------------------------------------------------------------------------


def _soffice_convert_to_pdf_step(src_path: str, expected_path: str) -> dict:
    """Build a config step that converts `src_path` to PDF, then renames the
    soffice output to `expected_path`. soffice writes `<basename>.pdf` to
    `--outdir`; we capture and move it to a deterministic location.
    """
    import os as _os
    out_dir = _os.path.dirname(expected_path) or "/tmp"
    src_stem = "$(basename '" + src_path + "' | sed 's/\\.[^.]*$//')"
    return _execute(
        f"mkdir -p '{out_dir}' && "
        f"soffice --headless --norestore --nofirststartwizard "
        f"--convert-to pdf --outdir '{out_dir}' '{src_path}' >/dev/null 2>&1 && "
        f"mv '{out_dir}/'{src_stem}'.pdf' '{expected_path}'"
    )


_DOCX_TO_PDF_SPECS: list[dict] = [
    {
        "template_id": "multi_docx_to_pdf_meeting_notes",
        "src_basename": "meeting_minutes.docx",
        "dst_basename": "meeting_minutes.pdf",
        "paragraphs": [
            ("Title", "Project Sync Meeting Minutes"),
            ("Heading 1", "Attendees"),
            (None, "Alice Chen, Bob Martinez, Carol Singh, David Park."),
            ("Heading 1", "Agenda"),
            (None, "Roadmap review, Q3 priorities, blockers."),
            ("Heading 1", "Discussion"),
            (None, "The pipeline migration is on track for next sprint. Alice will own the rollout plan and Bob will draft the changelog. Carol raised a concern about the staging environment capacity which David offered to investigate."),
            ("Heading 1", "Action Items"),
            (None, "Alice: migration plan by Friday. Bob: changelog draft. David: staging capacity report."),
            ("Heading 1", "Next Meeting"),
            (None, "2026-05-23, same time."),
        ],
        "instructions": [
            "I have meeting minutes in a Word document at ~/Desktop/meeting_minutes.docx. Please export them as a PDF at ~/Desktop/meeting_minutes.pdf — running LibreOffice headlessly from a terminal is the easiest way.",
        ],
    },
    {
        "template_id": "multi_docx_to_pdf_quarterly_report",
        "src_basename": "quarterly_report.docx",
        "dst_basename": "quarterly_report.pdf",
        "paragraphs": [
            ("Title", "Quarterly Business Report — Q1 2026"),
            ("Heading 1", "Executive Summary"),
            (None, "Revenue grew 18% year-over-year, driven by enterprise renewals and a strong launch of the new analytics tier. Operating margin held steady at 21% despite higher infrastructure spend."),
            ("Heading 1", "Revenue Breakdown"),
            (None, "Subscriptions: $4.2M. Services: $0.9M. One-time licenses: $0.3M. Total: $5.4M."),
            ("Heading 1", "Customer Wins"),
            (None, "Three new Fortune-500 customers signed multi-year contracts during the quarter. Net retention reached 121%."),
            ("Heading 1", "Outlook"),
            (None, "Q2 will focus on expanding the analytics tier and launching the partner program."),
        ],
        "instructions": [
            "From a terminal, please export the quarterly report at ~/Desktop/quarterly_report.docx to PDF, saving the result as ~/Desktop/quarterly_report.pdf. LibreOffice's headless mode is available.",
        ],
    },
    {
        "template_id": "multi_docx_to_pdf_offer_letter",
        "src_basename": "offer_letter.docx",
        "dst_basename": "offer_letter.pdf",
        "paragraphs": [
            ("Title", "Offer of Employment"),
            (None, "Dear Jordan Kim,"),
            (None, "We are pleased to offer you the position of Senior Software Engineer at Acme Corp, reporting to the Engineering Manager."),
            ("Heading 1", "Compensation"),
            (None, "Annual base salary: $165,000. Sign-on bonus: $20,000. Equity: 4,500 RSUs vesting over four years."),
            ("Heading 1", "Start Date"),
            (None, "Your anticipated start date is June 1, 2026."),
            (None, "Please confirm your acceptance by replying to this offer within ten business days."),
            (None, "Sincerely, Pat Morgan, Head of People."),
        ],
        "instructions": [
            "I need a PDF version of the offer letter at ~/Desktop/offer_letter.docx. From a terminal, use LibreOffice in headless mode to produce ~/Desktop/offer_letter.pdf.",
        ],
    },
    # Mass-add — additional docx → pdf variants.
    {
        "template_id": "multi_docx_to_pdf_policy_memo",
        "src_basename": "policy_memo.docx",
        "dst_basename": "policy_memo.pdf",
        "paragraphs": [
            ("Title", "Policy Memo — Remote Work"),
            ("Heading 1", "Background"),
            (None, "Our hybrid pilot ran for six months and produced strong engagement metrics across all teams."),
            ("Heading 1", "Decision"),
            (None, "Effective July 1, 2026, all eligible roles may work remotely up to three days per week with manager approval."),
            ("Heading 1", "Scope"),
            (None, "This policy applies to full-time employees in engineering, design, marketing, and operations functions."),
            ("Heading 1", "Effective Date"),
            (None, "July 1, 2026. Department leads will distribute team-level guidance two weeks prior."),
        ],
        "instructions": [
            "Could you turn the policy memo at ~/Desktop/policy_memo.docx into a PDF saved at ~/Desktop/policy_memo.pdf? Running LibreOffice in headless mode from a terminal works.",
        ],
    },
    {
        "template_id": "multi_docx_to_pdf_event_brief",
        "src_basename": "event_brief.docx",
        "dst_basename": "event_brief.pdf",
        "paragraphs": [
            ("Title", "Annual Conference Brief — Catalyst 2026"),
            ("Heading 1", "Theme"),
            (None, "Catalyst 2026 explores the next decade of platform engineering and AI co-pilot tooling."),
            ("Heading 1", "Logistics"),
            (None, "Three days, two main stages, eight breakout rooms. Expected attendance: 1,800."),
            ("Heading 1", "Key Speakers"),
            (None, "Dr. Mei Lin (keynote), Robin Park (closing), and a partner panel from cloud, observability, and dev-tooling vendors."),
        ],
        "instructions": [
            "From a terminal, export the event briefing at ~/Desktop/event_brief.docx as a PDF at ~/Desktop/event_brief.pdf. LibreOffice supports headless docx-to-pdf conversion.",
        ],
    },
]


def _make_docx_to_pdf_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.pdf"
    expected_txt  = f"{_TMP}/expected_{template_id}.txt"
    result_txt    = f"{_TMP}/result_{template_id}.txt"
    paragraphs    = spec["paragraphs"]
    instructions  = spec["instructions"]

    init_steps = [
        _build_docx_step(src_path, paragraphs),
        _soffice_convert_to_pdf_step(src_path, expected_path),
        # Validation fix 2: extract gold text now so eval can byte-compare.
        _execute(f"pdftotext -layout '{expected_path}' '{expected_txt}'"),
        # validation gap-close: terminal-driven (soffice --headless shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    # Validation fix 2: switched compare_pdfs (continuous fuzz.ratio — partial-
    # credit fragility) to a strict pdftotext+compare_text_file pipeline.
    # Postconfig extracts result.pdf -> result.txt; eval byte-compares text.
    postconfig = [
        _execute(
            f"if [ -f '{dst_path}' ]; then "
            f"pdftotext -layout '{dst_path}' '{result_txt}'; "
            f"else : > '{result_txt}'; fi"
        ),
    ]
    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": result_txt, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_txt, "dest": "expected.txt"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: postconfig,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="docx_to_pdf",
    )


_XLSX_TO_PDF_SPECS: list[dict] = [
    {
        "template_id": "multi_xlsx_to_pdf_sales",
        "src_basename": "sales_report.xlsx",
        "dst_basename": "sales_report.pdf",
        "rows": [
            ["Quarter", "Region", "Revenue", "Units"],
            ["Q1", "North",   "120000", "300"],
            ["Q1", "South",   "98000",  "245"],
            ["Q1", "East",    "142000", "360"],
            ["Q2", "North",   "136000", "345"],
            ["Q2", "South",   "110000", "280"],
            ["Q2", "East",    "158000", "395"],
        ],
        "instructions": [
            "From a terminal, please export the sales spreadsheet at ~/Desktop/sales_report.xlsx as a PDF at ~/Desktop/sales_report.pdf. Headless LibreOffice can handle xlsx-to-pdf.",
        ],
    },
    {
        "template_id": "multi_xlsx_to_pdf_inventory",
        "src_basename": "inventory_snapshot.xlsx",
        "dst_basename": "inventory_snapshot.pdf",
        "rows": [
            ["SKU", "Item", "Location", "Quantity", "ReorderPoint"],
            ["A001", "Widget Standard", "Aisle 3", "120", "60"],
            ["A002", "Widget Premium",  "Aisle 3", "45",  "30"],
            ["B010", "Gasket 4mm",      "Aisle 5", "320", "150"],
            ["B011", "Gasket 6mm",      "Aisle 5", "280", "150"],
            ["C100", "Bracket Steel",   "Aisle 1", "85",  "50"],
            ["C101", "Bracket Alum",    "Aisle 1", "62",  "50"],
        ],
        "instructions": [
            "I need a PDF copy of the inventory snapshot spreadsheet at ~/Desktop/inventory_snapshot.xlsx. From a terminal, use LibreOffice's headless mode to write ~/Desktop/inventory_snapshot.pdf.",
        ],
    },
    # Mass-add — additional xlsx → pdf variants.
    {
        "template_id": "multi_xlsx_to_pdf_payroll",
        "src_basename": "payroll_snapshot.xlsx",
        "dst_basename": "payroll_snapshot.pdf",
        "rows": [
            ["EmployeeID", "Name", "Department", "Salary"],
            ["E01", "Alice Chen",   "Engineering", "145000"],
            ["E02", "Bob Martinez", "Sales",       "98000"],
            ["E03", "Carol Singh",  "Engineering", "152000"],
            ["E04", "David Park",   "HR",          "82000"],
            ["E05", "Eva Rivera",   "Marketing",   "91000"],
        ],
        "instructions": [
            "Please render the payroll snapshot at ~/Desktop/payroll_snapshot.xlsx to a PDF at ~/Desktop/payroll_snapshot.pdf. Headless LibreOffice from a terminal will work.",
        ],
    },
    {
        "template_id": "multi_xlsx_to_pdf_grades",
        "src_basename": "grades_snapshot.xlsx",
        "dst_basename": "grades_snapshot.pdf",
        "rows": [
            ["Student", "Midterm", "Final", "Grade"],
            ["Alice", "92", "95", "A"],
            ["Bob",   "78", "82", "B"],
            ["Carol", "88", "91", "A"],
            ["Dave",  "70", "75", "C"],
            ["Eve",   "95", "97", "A"],
        ],
        "instructions": [
            "From a terminal, export the grades workbook at ~/Desktop/grades_snapshot.xlsx as a PDF at ~/Desktop/grades_snapshot.pdf — LibreOffice headless mode handles this.",
        ],
    },
]


def _make_xlsx_to_pdf_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.pdf"
    expected_txt  = f"{_TMP}/expected_{template_id}.txt"
    result_txt    = f"{_TMP}/result_{template_id}.txt"
    rows          = spec["rows"]
    instructions  = spec["instructions"]

    init_steps = [
        _build_xlsx_step(src_path, rows),
        _soffice_convert_to_pdf_step(src_path, expected_path),
        # Validation fix 2 — see _make_docx_to_pdf_template.
        _execute(f"pdftotext -layout '{expected_path}' '{expected_txt}'"),
        # validation gap-close: terminal-driven (soffice --headless shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    postconfig = [
        _execute(
            f"if [ -f '{dst_path}' ]; then "
            f"pdftotext -layout '{dst_path}' '{result_txt}'; "
            f"else : > '{result_txt}'; fi"
        ),
    ]
    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": result_txt, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_txt, "dest": "expected.txt"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: postconfig,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="xlsx_to_pdf",
    )


def _make_pptx_to_pdf_template() -> SynthTemplate:
    """Build a 3-slide pptx via python-pptx, then convert to gold pdf via
    soffice. Agent task: convert the pptx to pdf via soffice headless.
    """
    template_id   = "multi_pptx_to_pdf_kickoff_deck"
    src_path      = f"{_DESKTOP}/kickoff_deck.pptx"
    dst_path      = f"{_DESKTOP}/kickoff_deck.pdf"
    expected_path = f"{_TMP}/expected_{template_id}.pdf"
    expected_txt  = f"{_TMP}/expected_{template_id}.txt"
    result_txt    = f"{_TMP}/result_{template_id}.txt"

    build_py = textwrap.dedent(f"""\
        import os
        from pptx import Presentation
        path = {src_path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        prs = Presentation()
        slides = [
            ("Project Kickoff — Phoenix",      "Q3 2026 launch plan and team intros."),
            ("Goals",                          "Ship MVP by July, onboard pilot customers, hit 95% uptime."),
            ("Team",                           "Eng: Alice, Bob. Design: Carol. PM: David. Lead: Pat Morgan."),
        ]
        layout = prs.slide_layouts[1]
        for title, body in slides:
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = body
        prs.save(path)
        """)
    init_steps = [
        _python_step(build_py),
        _soffice_convert_to_pdf_step(src_path, expected_path),
        # Validation fix 2 — see _make_docx_to_pdf_template.
        _execute(f"pdftotext -layout '{expected_path}' '{expected_txt}'"),
        # validation gap-close: terminal-driven (soffice --headless shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    postconfig = [
        _execute(
            f"if [ -f '{dst_path}' ]; then "
            f"pdftotext -layout '{dst_path}' '{result_txt}'; "
            f"else : > '{result_txt}'; fi"
        ),
    ]
    evaluator = {
        "func": "compare_text_file",
        "result": {"type": "vm_file", "path": result_txt, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_txt, "dest": "expected.txt"},
    }
    instructions = [
        "From a terminal, export the slide deck at ~/Desktop/kickoff_deck.pptx to a PDF at ~/Desktop/kickoff_deck.pdf. LibreOffice's headless mode can do the conversion.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: postconfig,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_text_file",
        setup_class="pptx_to_pdf",
    )


def _make_pdftk_merge_template() -> SynthTemplate:
    """pdftk merge: stage two arxiv PDFs, agent merges them with `pdftk a b
    cat output combined.pdf`. Gold = same merge, built in pre-config.
    """
    template_id   = "multi_pdftk_merge_arxiv_pair"
    pdf_a_rel     = "docs/pdf/arxiv-1706-03762.pdf"
    pdf_b_rel     = "docs/pdf/arxiv-1810-04805.pdf"
    pdf_a_path    = f"{_DESKTOP}/arxiv-1706-03762.pdf"
    pdf_b_path    = f"{_DESKTOP}/arxiv-1810-04805.pdf"
    dst_path      = f"{_DESKTOP}/combined_arxiv.pdf"
    expected_path = f"{_TMP}/expected_{template_id}.pdf"

    init_steps = [
        _stage_asset(pdf_a_rel, pdf_a_path),
        _stage_asset(pdf_b_rel, pdf_b_path),
        _execute(
            f"pdftk '{pdf_a_path}' '{pdf_b_path}' cat output '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (pdftk shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_pdfs",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.pdf"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.pdf"},
    }
    instructions = [
        "From a terminal, please combine the two arxiv PDFs at ~/Desktop/arxiv-1706-03762.pdf and ~/Desktop/arxiv-1810-04805.pdf (in that order) into a single merged PDF at ~/Desktop/combined_arxiv.pdf using poppler's pdfunite: `pdfunite ~/Desktop/arxiv-1706-03762.pdf ~/Desktop/arxiv-1810-04805.pdf ~/Desktop/combined_arxiv.pdf`.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_pdfs",
        setup_class="pdftk_merge",
    )


def _make_pdftk_extract_template() -> SynthTemplate:
    """pdftk extract pages: stage one arxiv PDF, agent extracts pages 2-4
    via `pdftk in cat 2-4 output extract.pdf`. Gold built in pre-config.
    """
    template_id   = "multi_pdftk_extract_pages_attention"
    src_rel       = "docs/pdf/arxiv-1706-03762.pdf"
    src_path      = f"{_DESKTOP}/arxiv-1706-03762.pdf"
    dst_path      = f"{_DESKTOP}/attention_pages_2to4.pdf"
    expected_path = f"{_TMP}/expected_{template_id}.pdf"

    init_steps = [
        _stage_asset(src_rel, src_path),
        _execute(
            f"pdftk '{src_path}' cat 2-4 output '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (pdftk shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_pdfs",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.pdf"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.pdf"},
    }
    instructions = [
        "From a terminal, pull out pages 2 through 4 (inclusive) of the arxiv paper at ~/Desktop/arxiv-1706-03762.pdf and save just those three pages as a new PDF at ~/Desktop/attention_pages_2to4.pdf using poppler: `pdfseparate -f 2 -l 4 ~/Desktop/arxiv-1706-03762.pdf /tmp/att_pg-%d.pdf && pdfunite /tmp/att_pg-2.pdf /tmp/att_pg-3.pdf /tmp/att_pg-4.pdf ~/Desktop/attention_pages_2to4.pdf`.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_pdfs",
        setup_class="pdftk_extract",
    )


# ---------------------------------------------------------------------------
# Cat O — exact_match rows (added validation: close eval=3/synth=0 gap).
# ---------------------------------------------------------------------------
# Pattern: pre_config materializes input file(s) (csv / text) such that a
# small terminal pipeline produces a single-string result (line count, max
# value, missing-id check). Agent writes the answer to a sink text file;
# eval = exact_match against the gold string. Mirrors eval rows like
# osworld_multi_apps_48d05431 (single "1\n" answer).
# ---------------------------------------------------------------------------


_EXACT_MATCH_SPECS: list[dict] = [
    {
        "template_id": "multi_exact_match_csv_row_count",
        "csv_basename": "students_roster.csv",
        "csv_rows": [
            ["StudentID", "Name", "Grade"],
            ["S001", "Alice Chen",   "A"],
            ["S002", "Bob Martinez", "B"],
            ["S003", "Carol Singh",  "A"],
            ["S004", "David Park",   "B"],
            ["S005", "Eva Rivera",   "A"],
            ["S006", "Frank Liu",    "C"],
            ["S007", "Grace Adler",  "A"],
            ["S008", "Hugo Marin",   "B"],
        ],
        "ans_basename": "row_count.txt",
        "expected": "8\n",
        "instructions": [
            "From a terminal, count how many DATA rows are in ~/Desktop/students_roster.csv (do NOT count the header row). Write just that integer (no thousands separators, no extra text) followed by a single newline to ~/Desktop/row_count.txt.",
        ],
    },
    {
        "template_id": "multi_exact_match_max_revenue_quarter",
        "csv_basename": "quarterly_revenue.csv",
        "csv_rows": [
            ["Quarter", "Revenue"],
            ["Q1", "120000"],
            ["Q2", "138000"],
            ["Q3", "175000"],
            ["Q4", "162000"],
        ],
        "ans_basename": "best_quarter.txt",
        "expected": "Q3\n",
        "instructions": [
            "Open ~/Desktop/quarterly_revenue.csv (header row plus Quarter,Revenue rows) in a terminal. Identify the single Quarter label (e.g. Q1, Q2, Q3, Q4) with the maximum Revenue value, and write just that label followed by a single newline to ~/Desktop/best_quarter.txt.",
        ],
    },
    {
        "template_id": "multi_exact_match_unique_dept_count",
        "csv_basename": "employees_dept.csv",
        "csv_rows": [
            ["EmployeeID", "Name", "Department"],
            ["E01", "Alice",  "Engineering"],
            ["E02", "Bob",    "Sales"],
            ["E03", "Carol",  "Engineering"],
            ["E04", "David",  "HR"],
            ["E05", "Eva",    "Sales"],
            ["E06", "Frank",  "Engineering"],
            ["E07", "Grace",  "Marketing"],
            ["E08", "Hugo",   "HR"],
        ],
        "ans_basename": "dept_count.txt",
        "expected": "4\n",
        "instructions": [
            "From a terminal, look at ~/Desktop/employees_dept.csv and count how many DISTINCT department names appear in its third column (skip the header row). Write the resulting integer (and just that integer) followed by a newline to ~/Desktop/dept_count.txt.",
        ],
    },
    # Mass-add — additional exact_match variants.
    {
        "template_id": "multi_exact_match_grade_a_count",
        "csv_basename": "grade_roster.csv",
        "csv_rows": [
            ["StudentID", "Name", "Grade"],
            ["S001", "Alice", "A"],
            ["S002", "Bob",   "B"],
            ["S003", "Carol", "A"],
            ["S004", "Dave",  "C"],
            ["S005", "Eve",   "A"],
            ["S006", "Frank", "B"],
            ["S007", "Grace", "A"],
        ],
        "ans_basename": "grade_a_count.txt",
        "expected": "4\n",
        "instructions": [
            "From a terminal, look at ~/Desktop/grade_roster.csv (a header row followed by data rows where the third column is the letter grade) and count how many data rows have a Grade value of exactly 'A'. Write that integer followed by a newline to ~/Desktop/grade_a_count.txt.",
        ],
    },
    {
        "template_id": "multi_exact_match_max_temperature",
        "csv_basename": "weather_log.csv",
        "csv_rows": [
            ["Date", "City", "TempC"],
            ["2026-05-01", "Madrid", "22"],
            ["2026-05-01", "Paris",  "16"],
            ["2026-05-01", "Berlin", "14"],
            ["2026-05-01", "Rome",   "24"],
            ["2026-05-01", "Lisbon", "21"],
        ],
        "ans_basename": "hottest_city.txt",
        "expected": "Rome\n",
        "instructions": [
            "Open ~/Desktop/weather_log.csv (header row plus Date,City,TempC rows) in a terminal. Identify the City with the maximum TempC value and write just that city name followed by a single newline to ~/Desktop/hottest_city.txt.",
        ],
    },
    {
        "template_id": "multi_exact_match_total_units",
        "csv_basename": "shipment_units.csv",
        "csv_rows": [
            ["ShipmentID", "Units"],
            ["SH01", "100"],
            ["SH02", "250"],
            ["SH03", "75"],
            ["SH04", "180"],
        ],
        "ans_basename": "total_units.txt",
        "expected": "605\n",
        "instructions": [
            "From a terminal, add up every value in the Units column (column 2) of ~/Desktop/shipment_units.csv, skipping the header row. Write the resulting integer total (no commas) followed by a newline to ~/Desktop/total_units.txt.",
        ],
    },
]


def _make_exact_match_template(spec: dict) -> SynthTemplate:
    template_id  = spec["template_id"]
    csv_path     = f"{_DESKTOP}/{spec['csv_basename']}"
    ans_path     = f"{_DESKTOP}/{spec['ans_basename']}"
    csv_rows     = spec["csv_rows"]
    expected     = spec["expected"]
    instructions = spec["instructions"]

    init_steps = [
        _write_text_step(csv_path, _csv_text(csv_rows)),
        # validation gap-close: terminal-driven (shell pipeline produces answer) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_write_text_step(ans_path, expected)]

    # Validation fix (CRITICAL): `exact_match(result, rules)` compares
    # `result == rules["expected"]`. With result.type=="vm_file", the getter
    # returns the LOCAL CACHE PATH (string), NOT the file contents — so the
    # comparison ALWAYS yielded "<cache>/result.txt" == "8\n" → False →
    # FALSE-FAIL on every Cat O row. Switch to vm_command_line so result is
    # the stdout of `cat`, matching eval-side multi_apps_48d05431 pattern.
    evaluator = {
        "func": "exact_match",
        "result": {
            "type": "vm_command_line",
            "command": ["cat", ans_path],
            "shell": False,
        },
        "expected": {"type": "rule", "rules": {"expected": expected}},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="exact_match",
        setup_class="csv_to_string_answer",
    )


# ---------------------------------------------------------------------------
# Cat O.2 — HTML fact extraction (Added, 2026-05-11).
# ---------------------------------------------------------------------------
# Pattern: pre_config stages a small self-contained HTML page (a `<table>`
# with a handful of rows) under ~/Desktop/. Agent opens the HTML in Chrome
# (open_command), reads off ONE specific fact (highest/lowest/named entry),
# and writes that single-string answer to a text file under ~/Desktop/.
# Eval = exact_match against the gold (`vm_command_line cat` of the answer
# file, mirroring _EXACT_MATCH_SPECS).
#
# Addresses validation gap: Chrome + plain-file cross-app workflows.
# Existing _EXACT_MATCH_SPECS rows answer-from-CSV via terminal only; these
# answer-from-HTML, requiring the agent to actually open Chrome.
# ---------------------------------------------------------------------------


def _html_fact_page(title: str, intro: str, header: list[str], rows: list[list[str]]) -> str:
    """Build a minimal self-contained HTML page with a single data table.

    Kept deliberately simple (no CSS, no JS) so the agent can read it
    reliably whether they use Chrome or `cat` the file.
    """
    thead = "".join(f"<th>{c}</th>" for c in header)
    body_rows = "\n".join(
        "    <tr>" + "".join(f"<td>{c}</td>" for c in r) + "</tr>"
        for r in rows
    )
    return (
        "<!DOCTYPE html>\n"
        "<html><head><meta charset=\"utf-8\">"
        f"<title>{title}</title></head>\n"
        "<body>\n"
        f"<h1>{title}</h1>\n"
        f"<p>{intro}</p>\n"
        "<table border=\"1\" cellpadding=\"4\">\n"
        f"  <thead><tr>{thead}</tr></thead>\n"
        "  <tbody>\n"
        f"{body_rows}\n"
        "  </tbody>\n"
        "</table>\n"
        "</body></html>\n"
    )


_HTML_FACT_EXTRACT_SPECS: list[dict] = [
    {
        "template_id": "multi_html_fact_top_country_gdp",
        "html_basename": "world_gdp.html",
        "html_body": _html_fact_page(
            title="World GDP — 2022 (Selected Economies)",
            intro="Headline 2022 nominal GDP (USD trillions) for ten selected economies. Source: World Bank.",
            header=["Country", "GDP_USD_Trillions"],
            rows=[
                ["United States", "25.46"],
                ["China",         "17.96"],
                ["Japan",         "4.23"],
                ["Germany",       "4.07"],
                ["India",         "3.39"],
                ["United Kingdom","3.07"],
                ["France",        "2.78"],
                ["Italy",         "2.01"],
                ["Canada",        "2.14"],
                ["Brazil",        "1.92"],
            ],
        ),
        "ans_basename": "top_economy.txt",
        "expected": "United States\n",
        "instructions": [
            "Open ~/Desktop/world_gdp.html in Chrome. The page shows a small table of 2022 GDP figures for selected economies. Find the country with the HIGHEST GDP_USD_Trillions value and write just that country's name followed by a single newline to ~/Desktop/top_economy.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_top_language",
        "html_basename": "tiobe_top.html",
        "html_body": _html_fact_page(
            title="TIOBE Index — Top 8 Languages (Snapshot)",
            intro="Programming-language popularity ranking (snapshot, lower rank = more popular).",
            header=["Rank", "Language", "Rating_Pct"],
            rows=[
                ["1", "Python",       "23.45"],
                ["2", "C",            "12.10"],
                ["3", "Java",         "10.20"],
                ["4", "C++",          "9.80"],
                ["5", "C#",           "6.50"],
                ["6", "JavaScript",   "3.70"],
                ["7", "Visual Basic", "2.40"],
                ["8", "SQL",          "1.95"],
            ],
        ),
        "ans_basename": "top_language.txt",
        "expected": "Python\n",
        "instructions": [
            "Open ~/Desktop/tiobe_top.html in Chrome. The page lists the top 8 programming languages by a popularity rating. Identify the language with the HIGHEST Rating_Pct (equivalently, Rank == 1) and write just that language's name followed by a single newline to ~/Desktop/top_language.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_tallest_peak",
        "html_basename": "mountain_peaks.html",
        "html_body": _html_fact_page(
            title="Tallest Mountain Peaks (Meters)",
            intro="A small selection of the world's highest mountain peaks, by elevation in meters above sea level.",
            header=["Peak", "Range", "Elevation_M"],
            rows=[
                ["Everest",       "Himalayas",      "8848"],
                ["K2",            "Karakoram",      "8611"],
                ["Kangchenjunga", "Himalayas",      "8586"],
                ["Lhotse",        "Himalayas",      "8516"],
                ["Makalu",        "Himalayas",      "8485"],
                ["Cho Oyu",       "Himalayas",      "8188"],
                ["Dhaulagiri",    "Himalayas",      "8167"],
                ["Manaslu",       "Himalayas",      "8163"],
            ],
        ),
        "ans_basename": "tallest_peak.txt",
        "expected": "Everest\n",
        "instructions": [
            "Open ~/Desktop/mountain_peaks.html in Chrome. The page lists 8 mountain peaks with their elevation in meters. Identify the peak with the GREATEST Elevation_M and write just the peak's name (no range, no number) followed by a single newline to ~/Desktop/tallest_peak.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_top_grossing_film",
        "html_basename": "box_office.html",
        "html_body": _html_fact_page(
            title="Worldwide Box Office — Top 6 Films (Historical)",
            intro="Approximate lifetime worldwide gross in USD billions (historical, not inflation-adjusted).",
            header=["Film", "Year", "Gross_USD_Billions"],
            rows=[
                ["Avatar",                              "2009", "2.92"],
                ["Avengers: Endgame",                   "2019", "2.79"],
                ["Avatar: The Way of Water",            "2022", "2.32"],
                ["Titanic",                             "1997", "2.26"],
                ["Star Wars: The Force Awakens",        "2015", "2.07"],
                ["Avengers: Infinity War",              "2018", "2.05"],
            ],
        ),
        "ans_basename": "top_film.txt",
        "expected": "Avatar\n",
        "instructions": [
            "Open ~/Desktop/box_office.html in Chrome. The page lists 6 films with their worldwide box-office gross. Identify the film with the HIGHEST Gross_USD_Billions and write just the film's title (no year, no number) followed by a single newline to ~/Desktop/top_film.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_olympic_top_gold",
        "html_basename": "olympic_medals.html",
        "html_body": _html_fact_page(
            title="Olympic Medal Table — Top 8 Nations (Sample)",
            intro="Medal counts from a single hypothetical Summer Olympics for the top 8 nations.",
            header=["Country", "Gold", "Silver", "Bronze"],
            rows=[
                ["United States", "39", "41", "33"],
                ["China",         "38", "32", "18"],
                ["Japan",         "27", "14", "17"],
                ["Great Britain", "22", "21", "22"],
                ["Russia",        "20", "28", "23"],
                ["Australia",     "17", "7",  "22"],
                ["Netherlands",   "10", "12", "14"],
                ["France",        "10", "12", "11"],
            ],
        ),
        "ans_basename": "top_country_gold.txt",
        "expected": "United States\n",
        "instructions": [
            "Open ~/Desktop/olympic_medals.html in Chrome. The page shows 8 countries with their Gold/Silver/Bronze medal counts. Identify the country with the HIGHEST number of Gold medals and write just that country's name followed by a single newline to ~/Desktop/top_country_gold.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_cheapest_car",
        "html_basename": "car_prices.html",
        "html_body": _html_fact_page(
            title="2026 EV Price Comparison — Base MSRP (USD)",
            intro="Approximate 2026 base MSRP in USD for a small selection of electric vehicles.",
            header=["Model", "Make", "Base_MSRP_USD"],
            rows=[
                ["Leaf",      "Nissan",       "28000"],
                ["Bolt EV",   "Chevrolet",    "26500"],
                ["Model 3",   "Tesla",        "40240"],
                ["ID.4",      "Volkswagen",   "39995"],
                ["Ioniq 6",   "Hyundai",      "37500"],
                ["EV6",       "Kia",          "42600"],
            ],
        ),
        "ans_basename": "cheapest_car.txt",
        "expected": "Bolt EV\n",
        "instructions": [
            "Open ~/Desktop/car_prices.html in Chrome. The page lists 6 electric-vehicle models with their base MSRP in USD. Identify the model with the LOWEST Base_MSRP_USD and write just the model's name (no make, no price) followed by a single newline to ~/Desktop/cheapest_car.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_most_populous_capital",
        "html_basename": "capital_cities.html",
        "html_body": _html_fact_page(
            title="Capital City Populations (Sample)",
            intro="Approximate city-proper population (millions) for the capital city of each listed country.",
            header=["Country", "Capital", "Population_Millions"],
            rows=[
                ["Japan",          "Tokyo",       "13.96"],
                ["India",          "New Delhi",   "11.03"],
                ["China",          "Beijing",     "21.89"],
                ["Russia",         "Moscow",      "12.62"],
                ["United Kingdom", "London",      "8.98"],
                ["France",         "Paris",       "2.16"],
                ["Germany",        "Berlin",      "3.67"],
                ["Brazil",         "Brasilia",    "3.05"],
            ],
        ),
        "ans_basename": "most_populous_capital.txt",
        "expected": "Beijing\n",
        "instructions": [
            "Open ~/Desktop/capital_cities.html in Chrome. The page lists the capital cities of 8 countries with their city-proper population in millions. Identify the capital city with the HIGHEST Population_Millions value and write just the capital's name (not the country) followed by a single newline to ~/Desktop/most_populous_capital.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_bestselling_book_author",
        "html_basename": "bestselling_books.html",
        "html_body": _html_fact_page(
            title="Bestselling Standalone Books (Estimated Copies, Millions)",
            intro="Estimated lifetime sales (in millions of copies) of selected single-volume bestselling books.",
            header=["Title", "Author", "Copies_Millions"],
            rows=[
                ["A Tale of Two Cities",       "Charles Dickens",  "200"],
                ["The Little Prince",          "Antoine de Saint-Exupery", "150"],
                ["The Alchemist",              "Paulo Coelho",     "150"],
                ["Dream of the Red Chamber",   "Cao Xueqin",       "100"],
                ["And Then There Were None",   "Agatha Christie",  "100"],
                ["The Hobbit",                 "J.R.R. Tolkien",   "100"],
                ["She: A History of Adventure","H. Rider Haggard", "83"],
            ],
        ),
        "ans_basename": "top_book_author.txt",
        "expected": "Charles Dickens\n",
        "instructions": [
            "Open ~/Desktop/bestselling_books.html in Chrome. The page lists 7 bestselling standalone books with their author and estimated lifetime sales in millions. Identify the book with the HIGHEST Copies_Millions value, then write just THAT book's AUTHOR (full name as shown, no title, no number) followed by a single newline to ~/Desktop/top_book_author.txt.",
        ],
    },
    # Added (Phase 2 widen): varied question types beyond `find max`.
    # min-pick, count-condition, named-row lookup, cross-column lookup —
    # stress-tests the factory across different reading-comprehension shapes.
    {
        "template_id": "multi_html_fact_smallest_country_gdp",
        "html_basename": "small_economies.html",
        "html_body": _html_fact_page(
            title="Selected Small Economies — 2022 GDP (Millions USD)",
            intro="2022 nominal GDP figures for a handful of small economies, in millions of USD.",
            header=["Country", "GDP_Millions_USD"],
            rows=[
                ["Iceland",   "27800"],
                ["Cyprus",    "29250"],
                ["Estonia",   "39000"],
                ["Malta",     "17700"],
                ["Lithuania", "70300"],
                ["Latvia",    "41100"],
                ["Slovenia",  "62200"],
            ],
        ),
        "ans_basename": "smallest_economy.txt",
        "expected": "Malta\n",
        "instructions": [
            "Open ~/Desktop/small_economies.html in Chrome. The page lists 7 small economies and their 2022 GDP in millions of USD. Identify the country with the SMALLEST GDP_Millions_USD value and write just that country's name followed by a single newline to ~/Desktop/smallest_economy.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_count_engineering_employees",
        "html_basename": "company_directory.html",
        "html_body": _html_fact_page(
            title="Engineering Directory — Q2",
            intro="Engineering + adjacent departments roster (anonymized internal directory snapshot).",
            header=["EmployeeID", "Name", "Department"],
            rows=[
                ["E001", "Alice Chen",     "Engineering"],
                ["E002", "Bob Martinez",   "Sales"],
                ["E003", "Carol Singh",    "Engineering"],
                ["E004", "David Park",     "HR"],
                ["E005", "Eva Rivera",     "Engineering"],
                ["E006", "Frank Liu",      "Marketing"],
                ["E007", "Grace Adler",    "Engineering"],
                ["E008", "Hugo Marin",     "Sales"],
                ["E009", "Iris Tanaka",    "Engineering"],
                ["E010", "Jonas Becker",   "HR"],
            ],
        ),
        "ans_basename": "engineering_count.txt",
        "expected": "5\n",
        "instructions": [
            "Open ~/Desktop/company_directory.html in Chrome. The page lists 10 employees and their departments. Count how many employees have Department == \"Engineering\" and write just that integer followed by a single newline to ~/Desktop/engineering_count.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_python_rank",
        "html_basename": "lang_ranks_lookup.html",
        "html_body": _html_fact_page(
            title="Language Popularity Snapshot — Rank Lookup",
            intro="Programming-language popularity snapshot, listed in alphabetical (NOT rank) order so rank-by-name lookup requires reading the table.",
            header=["Language", "Rank", "Rating_Pct"],
            rows=[
                ["C",            "2", "12.10"],
                ["C#",           "5", "6.50"],
                ["C++",          "4", "9.80"],
                ["Java",         "3", "10.20"],
                ["JavaScript",   "6", "3.70"],
                ["Python",       "1", "23.45"],
                ["SQL",          "8", "1.95"],
                ["Visual Basic", "7", "2.40"],
            ],
        ),
        "ans_basename": "python_rank.txt",
        "expected": "1\n",
        "instructions": [
            "Open ~/Desktop/lang_ranks_lookup.html in Chrome. The page lists 8 programming languages with their popularity rank and rating. Find the row where Language == \"Python\" and write JUST the Rank value (an integer) followed by a single newline to ~/Desktop/python_rank.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_brazil_capital",
        "html_basename": "country_capitals.html",
        "html_body": _html_fact_page(
            title="Country → Capital Reference (Sample)",
            intro="A small alphabetically-sorted reference of country → capital pairs.",
            header=["Country", "Capital", "Continent"],
            rows=[
                ["Argentina", "Buenos Aires",   "South America"],
                ["Brazil",    "Brasilia",       "South America"],
                ["Canada",    "Ottawa",         "North America"],
                ["Egypt",     "Cairo",          "Africa"],
                ["France",    "Paris",          "Europe"],
                ["Japan",     "Tokyo",          "Asia"],
                ["Kenya",     "Nairobi",        "Africa"],
                ["Norway",    "Oslo",           "Europe"],
            ],
        ),
        "ans_basename": "brazil_capital.txt",
        "expected": "Brasilia\n",
        "instructions": [
            "Open ~/Desktop/country_capitals.html in Chrome. The page is a small country-to-capital reference. Find the row where Country == \"Brazil\" and write JUST that country's Capital (one word) followed by a single newline to ~/Desktop/brazil_capital.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_count_himalaya_peaks",
        "html_basename": "peak_ranges.html",
        "html_body": _html_fact_page(
            title="World's Highest Peaks — Range Membership",
            intro="A small sample of the world's highest peaks grouped by mountain range.",
            header=["Peak", "Range", "Elevation_M"],
            rows=[
                ["Everest",       "Himalayas",  "8848"],
                ["K2",            "Karakoram",  "8611"],
                ["Kangchenjunga", "Himalayas",  "8586"],
                ["Lhotse",        "Himalayas",  "8516"],
                ["Makalu",        "Himalayas",  "8485"],
                ["Cho Oyu",       "Himalayas",  "8188"],
                ["Gasherbrum I",  "Karakoram",  "8080"],
                ["Broad Peak",    "Karakoram",  "8051"],
                ["Annapurna I",   "Himalayas",  "8091"],
            ],
        ),
        "ans_basename": "himalaya_peak_count.txt",
        "expected": "6\n",
        "instructions": [
            "Open ~/Desktop/peak_ranges.html in Chrome. The page lists 9 mountain peaks and their mountain range. Count the peaks where Range == \"Himalayas\" and write JUST that integer count followed by a single newline to ~/Desktop/himalaya_peak_count.txt.",
        ],
    },
    {
        "template_id": "multi_html_fact_python_inventor_year",
        "html_basename": "language_origins.html",
        "html_body": _html_fact_page(
            title="Programming Languages — Origin Year",
            intro="First public release year for a handful of well-known programming languages.",
            header=["Language", "Created_By",                "Year"],
            rows=[
                ["C",           "Dennis Ritchie",            "1972"],
                ["C++",         "Bjarne Stroustrup",         "1985"],
                ["Python",      "Guido van Rossum",          "1991"],
                ["Java",        "James Gosling",             "1995"],
                ["JavaScript",  "Brendan Eich",              "1995"],
                ["C#",          "Anders Hejlsberg",          "2000"],
                ["Go",          "Robert Griesemer, et al.",  "2009"],
                ["Rust",        "Graydon Hoare",             "2010"],
            ],
        ),
        "ans_basename": "python_year.txt",
        "expected": "1991\n",
        "instructions": [
            "Open ~/Desktop/language_origins.html in Chrome. The page lists 8 programming languages with their creator and first-release year. Find the row where Language == \"Python\" and write JUST that language's Year (4-digit integer) followed by a single newline to ~/Desktop/python_year.txt.",
        ],
    },
]


def _make_html_fact_extract_template(spec: dict) -> SynthTemplate:
    """Build an HTML→single-string-answer cross-app row.

    Spec keys:
      template_id    : str
      html_basename  : str   — staged at ~/Desktop/<basename>
      html_body      : str   — full HTML document text
      ans_basename   : str   — agent's answer file at ~/Desktop/<basename>
      expected       : str   — gold answer (single line + trailing newline)
      instructions   : list[str]
    """
    template_id   = spec["template_id"]
    html_path     = f"{_DESKTOP}/{spec['html_basename']}"
    ans_path      = f"{_DESKTOP}/{spec['ans_basename']}"
    html_body     = spec["html_body"]
    expected      = spec["expected"]
    instructions  = spec["instructions"]

    init_steps = [
        _write_text_step(html_path, html_body),
        # validation gap-close: agent still writes the answer text from a terminal
        # (Chrome is launched via open_command for reading the HTML page).
        *_terminal_preopen_steps(),
    ]
    # Oracle writes the gold answer string to the sink path. The agent's
    # task is to do the same via reading the HTML in Chrome.
    oracle_steps = [_write_text_step(ans_path, expected)]

    # Mirrors Cat O `_make_exact_match_template` (Validation fix): use
    # vm_command_line so result is the stdout of `cat` (eval gets file
    # CONTENTS, not the local cache path).
    evaluator = {
        "func": "exact_match",
        "result": {
            "type": "vm_command_line",
            "command": ["cat", ans_path],
            "shell": False,
        },
        "expected": {"type": "rule", "rules": {"expected": expected}},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
            # Open Chrome on the staged HTML — file:// scheme so it works
            # offline. The pre-config writes the HTML before Chrome opens.
            "open_command": ["google-chrome", "--no-first-run", "--no-default-browser-check",
                             f"file://{html_path}"],
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="exact_match",
        setup_class="html_fact_extract",
    )


# ---------------------------------------------------------------------------
# Cat P — compare_image_list rows (validation: close eval=2/synth=0 gap).
# ---------------------------------------------------------------------------
# Pattern: stage N real photos via _stage_asset, agent's task is to copy /
# extract / rename a SUBSET of them into a target directory under specific
# basenames. Gold = same files copied to /tmp under the same basenames.
# Eval = compare_image_list (per-pair pixel-difference check; multi=true).
# ---------------------------------------------------------------------------


def _make_compare_image_list_template(spec: dict) -> SynthTemplate:
    """Build a compare_image_list row.

    Spec keys:
      template_id      : str
      photos           : list[(asset_rel, src_basename, dst_basename)]
      dst_dir_basename : str   — agent's sink dir under ~/Desktop/
      instructions     : list[str]
    """
    template_id      = spec["template_id"]
    photos           = spec["photos"]
    dst_dir_basename = spec["dst_dir_basename"]
    instructions     = spec["instructions"]

    src_paths = [f"{_DESKTOP}/{src_b}" for (_r, src_b, _d) in photos]
    dst_dir   = f"{_DESKTOP}/{dst_dir_basename}"
    dst_paths = [f"{dst_dir}/{dst_b}"  for (_r, _s, dst_b) in photos]
    exp_dir   = f"{_TMP}/expected_{template_id}"
    exp_paths = [f"{exp_dir}/{dst_b}"  for (_r, _s, dst_b) in photos]

    init_steps: list[dict] = []
    init_steps.append(_execute(f"mkdir -p '{exp_dir}'"))
    for (rel, src_b, dst_b) in photos:
        init_steps.append(_stage_asset(rel, f"{_DESKTOP}/{src_b}"))
        init_steps.append(_execute(
            f"cp '{_DESKTOP}/{src_b}' '{exp_dir}/{dst_b}'"
        ))
    # validation gap-close: terminal-driven (mkdir/cp shell) — pre-open terminal.
    init_steps.extend(_terminal_preopen_steps())

    oracle_steps: list[dict] = [_execute(f"mkdir -p '{dst_dir}'")]
    for src, dst in zip(src_paths, dst_paths):
        oracle_steps.append(_execute(f"cp '{src}' '{dst}'"))

    evaluator = {
        "func": "compare_image_list",
        "result": {
            "type": "vm_file",
            "path": dst_paths,
            "dest": [f"result_{i}.jpg" for i in range(len(dst_paths))],
            "multi": True,
        },
        "expected": {
            "type": "vm_file",
            "path": exp_paths,
            "dest": [f"expected_{i}.jpg" for i in range(len(exp_paths))],
            "multi": True,
        },
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_image_list",
        setup_class="photo_subset_collect",
    )


_COMPARE_IMAGE_LIST_SPECS: list[dict] = [
    {
        "template_id": "multi_image_list_food_collection",
        "photos": [
            ("photos/food/pizza-dish.jpg",      "pizza-dish.jpg",      "1.jpg"),
            ("photos/food/coffee-latte.jpg",    "coffee-latte.jpg",    "2.jpg"),
            ("photos/food/restaurant-meal.jpg", "restaurant-meal.jpg", "3.jpg"),
        ],
        "dst_dir_basename": "food_collection",
        "instructions": [
            "In a terminal, collect the three staged food photos in ~/Desktop/ (pizza-dish.jpg, coffee-latte.jpg, restaurant-meal.jpg) into a new directory ~/Desktop/food_collection/ renaming them in order to 1.jpg, 2.jpg, 3.jpg respectively. Hint: `mkdir -p ~/Desktop/food_collection && cp ~/Desktop/pizza-dish.jpg ~/Desktop/food_collection/1.jpg && cp ~/Desktop/coffee-latte.jpg ~/Desktop/food_collection/2.jpg && cp ~/Desktop/restaurant-meal.jpg ~/Desktop/food_collection/3.jpg`.",
        ],
    },
    {
        "template_id": "multi_image_list_wildlife_collection",
        "photos": [
            ("photos/wildlife/tiger-closeup.jpg", "tiger-closeup.jpg", "01.jpg"),
            ("photos/wildlife/horse-meadow.jpg",  "horse-meadow.jpg",  "02.jpg"),
            ("photos/wildlife/bird-perch.jpg",    "bird-perch.jpg",    "03.jpg"),
            ("photos/wildlife/fox-portrait.jpg",  "fox-portrait.jpg",  "04.jpg"),
        ],
        "dst_dir_basename": "wildlife_gallery",
        "instructions": [
            "In a terminal, gather the four wildlife photos in ~/Desktop/ (tiger-closeup.jpg, horse-meadow.jpg, bird-perch.jpg, fox-portrait.jpg) into ~/Desktop/wildlife_gallery/ renaming them to 01.jpg, 02.jpg, 03.jpg, 04.jpg in that order. Hint: `mkdir -p ~/Desktop/wildlife_gallery && cp ~/Desktop/tiger-closeup.jpg ~/Desktop/wildlife_gallery/01.jpg && cp ~/Desktop/horse-meadow.jpg ~/Desktop/wildlife_gallery/02.jpg && cp ~/Desktop/bird-perch.jpg ~/Desktop/wildlife_gallery/03.jpg && cp ~/Desktop/fox-portrait.jpg ~/Desktop/wildlife_gallery/04.jpg`.",
        ],
    },
    # Mass-add — additional compare_image_list collections.
    {
        "template_id": "multi_image_list_landscape_collection",
        "photos": [
            ("photos/landscape/mountain-range.jpg", "mountain-range.jpg", "a.jpg"),
            ("photos/landscape/beach-sunset.jpg",   "beach-sunset.jpg",   "b.jpg"),
            ("photos/landscape/forest-trail.jpg",   "forest-trail.jpg",   "c.jpg"),
        ],
        "dst_dir_basename": "landscape_album",
        "instructions": [
            "In a terminal, gather the three landscape photos in ~/Desktop/ (mountain-range.jpg, beach-sunset.jpg, forest-trail.jpg) into ~/Desktop/landscape_album/ renaming them to a.jpg, b.jpg, c.jpg in that order. Hint: `mkdir -p ~/Desktop/landscape_album && cp ~/Desktop/mountain-range.jpg ~/Desktop/landscape_album/a.jpg && cp ~/Desktop/beach-sunset.jpg ~/Desktop/landscape_album/b.jpg && cp ~/Desktop/forest-trail.jpg ~/Desktop/landscape_album/c.jpg`.",
        ],
    },
    {
        "template_id": "multi_image_list_portrait_collection",
        "photos": [
            ("photos/portrait/person-headshot-1.jpg", "person-headshot-1.jpg", "headshot-001.jpg"),
            ("photos/portrait/person-headshot-2.jpg", "person-headshot-2.jpg", "headshot-002.jpg"),
            ("photos/portrait/person-headshot-3.jpg", "person-headshot-3.jpg", "headshot-003.jpg"),
        ],
        "dst_dir_basename": "headshot_set",
        "instructions": [
            "In a terminal, gather the three portrait photos in ~/Desktop/ (person-headshot-1.jpg, person-headshot-2.jpg, person-headshot-3.jpg) into ~/Desktop/headshot_set/ renaming them to headshot-001.jpg, headshot-002.jpg, headshot-003.jpg in that order. Hint: `mkdir -p ~/Desktop/headshot_set && cp ~/Desktop/person-headshot-1.jpg ~/Desktop/headshot_set/headshot-001.jpg && cp ~/Desktop/person-headshot-2.jpg ~/Desktop/headshot_set/headshot-002.jpg && cp ~/Desktop/person-headshot-3.jpg ~/Desktop/headshot_set/headshot-003.jpg`.",
        ],
    },
    {
        "template_id": "multi_image_list_product_collection",
        "photos": [
            ("photos/product/sneakers.jpg",   "sneakers.jpg",   "p1.jpg"),
            ("photos/product/wristwatch.jpg", "wristwatch.jpg", "p2.jpg"),
            ("photos/product/headphones.jpg", "headphones.jpg", "p3.jpg"),
        ],
        "dst_dir_basename": "product_catalog",
        "instructions": [
            "In a terminal, gather the three product photos in ~/Desktop/ (sneakers.jpg, wristwatch.jpg, headphones.jpg) into ~/Desktop/product_catalog/ renaming them to p1.jpg, p2.jpg, p3.jpg in that order. Hint: `mkdir -p ~/Desktop/product_catalog && cp ~/Desktop/sneakers.jpg ~/Desktop/product_catalog/p1.jpg && cp ~/Desktop/wristwatch.jpg ~/Desktop/product_catalog/p2.jpg && cp ~/Desktop/headphones.jpg ~/Desktop/product_catalog/p3.jpg`.",
        ],
    },
    # validation — wire newly-added energy/industrial/education photo categories.
    {
        "template_id": "multi_image_list_energy_collection",
        "photos": [
            ("photos/energy/wind-turbine.jpg",         "wind-turbine.jpg",         "01.jpg"),
            ("photos/energy/solar-farm.jpg",           "solar-farm.jpg",           "02.jpg"),
            ("photos/energy/hydroelectric-dam.jpg",    "hydroelectric-dam.jpg",    "03.jpg"),
            ("photos/energy/nuclear-cooling-tower.jpg","nuclear-cooling-tower.jpg","04.jpg"),
        ],
        "dst_dir_basename": "energy_album",
        "instructions": [
            "In a terminal, gather the four energy-sector photos in ~/Desktop/ (wind-turbine.jpg, solar-farm.jpg, hydroelectric-dam.jpg, nuclear-cooling-tower.jpg) into ~/Desktop/energy_album/ renaming them to 01.jpg, 02.jpg, 03.jpg, 04.jpg in that order. Hint: `mkdir -p ~/Desktop/energy_album && cp ~/Desktop/wind-turbine.jpg ~/Desktop/energy_album/01.jpg && cp ~/Desktop/solar-farm.jpg ~/Desktop/energy_album/02.jpg && cp ~/Desktop/hydroelectric-dam.jpg ~/Desktop/energy_album/03.jpg && cp ~/Desktop/nuclear-cooling-tower.jpg ~/Desktop/energy_album/04.jpg`.",
        ],
    },
    {
        "template_id": "multi_image_list_industrial_collection",
        "photos": [
            ("photos/industrial/factory-loom.jpg",   "factory-loom.jpg",   "a.jpg"),
            ("photos/industrial/cargo-port.jpg",     "cargo-port.jpg",     "b.jpg"),
            ("photos/industrial/assembly-line.jpg",  "assembly-line.jpg",  "c.jpg"),
        ],
        "dst_dir_basename": "industrial_album",
        "instructions": [
            "In a terminal, gather the three industrial photos in ~/Desktop/ (factory-loom.jpg, cargo-port.jpg, assembly-line.jpg) into ~/Desktop/industrial_album/ renaming them to a.jpg, b.jpg, c.jpg in that order. Hint: `mkdir -p ~/Desktop/industrial_album && cp ~/Desktop/factory-loom.jpg ~/Desktop/industrial_album/a.jpg && cp ~/Desktop/cargo-port.jpg ~/Desktop/industrial_album/b.jpg && cp ~/Desktop/assembly-line.jpg ~/Desktop/industrial_album/c.jpg`.",
        ],
    },
    {
        "template_id": "multi_image_list_education_collection",
        "photos": [
            ("photos/education/chalkboard-classroom.jpg",       "chalkboard-classroom.jpg",       "01.jpg"),
            ("photos/education/graduation-ceremony.jpg",        "graduation-ceremony.jpg",        "02.jpg"),
            ("photos/education/students-studying.jpg",          "students-studying.jpg",          "03.jpg"),
            ("photos/education/nebula-pillars-of-creation.jpg", "nebula-pillars-of-creation.jpg", "04.jpg"),
        ],
        "dst_dir_basename": "education_album",
        "instructions": [
            "In a terminal, gather the four education photos in ~/Desktop/ (chalkboard-classroom.jpg, graduation-ceremony.jpg, students-studying.jpg, nebula-pillars-of-creation.jpg) into ~/Desktop/education_album/ renaming them to 01.jpg, 02.jpg, 03.jpg, 04.jpg in that order. Hint: `mkdir -p ~/Desktop/education_album && cp ~/Desktop/chalkboard-classroom.jpg ~/Desktop/education_album/01.jpg && cp ~/Desktop/graduation-ceremony.jpg ~/Desktop/education_album/02.jpg && cp ~/Desktop/students-studying.jpg ~/Desktop/education_album/03.jpg && cp ~/Desktop/nebula-pillars-of-creation.jpg ~/Desktop/education_album/04.jpg`.",
        ],
    },
]


# ---------------------------------------------------------------------------
# Cat Q — check_python_file_by_test_suite rows (validation: close eval=2 gap).
# ---------------------------------------------------------------------------
# Pattern: pre_config writes a SOURCE python file (with a deliberately broken
# function) AND a separate test file under /tmp/ that imports the source and
# defines a `test()` function returning bool. The gold version of the source
# (correct implementation) is staged under /tmp/. Agent's task is to fix the
# source file so the test passes. Oracle = `cp gold_src agent_src`.
# Eval = check_python_file_by_test_suite(actual_files=[agent_src],
#                                        test_file=test_suite.py).
# ---------------------------------------------------------------------------


_PY_TEST_SPECS: list[dict] = [
    {
        "template_id": "multi_pytest_fix_add_function",
        "src_basename": "math_utils.py",
        "broken_src": (
            "def add(a, b):\n"
            "    # BUG: returns a - b instead of a + b. Fix it!\n"
            "    return a - b\n"
        ),
        "fixed_src": (
            "def add(a, b):\n"
            "    return a + b\n"
        ),
        "test_body": (
            # Validation fix (CRITICAL): `check_python_file_by_test_suite`
            # runs test() on the EVAL HOST, not the VM. The VM-side path
            # `/home/user/Desktop/...` doesn't exist there. Resolve relative
            # to test.py's own dir (the cache_dir where eval downloaded both
            # the test and the agent's source side by side).
            "import importlib.util, sys, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'math_utils.py')\n"
            "    spec = importlib.util.spec_from_file_location('math_utils', src)\n"
            "    mod = importlib.util.module_from_spec(spec)\n"
            "    spec.loader.exec_module(mod)\n"
            "    if mod.add(2, 3) != 5: return False\n"
            "    if mod.add(-1, 1) != 0: return False\n"
            "    if mod.add(100, 250) != 350: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/math_utils.py and fix the bug in the `add(a, b)` function so that it correctly returns `a + b` (currently it returns `a - b`). Save the file in place. Do not change the function signature or add any new functions.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_factorial",
        "src_basename": "factorial_util.py",
        "broken_src": (
            "def factorial(n):\n"
            "    # BUG: off-by-one — returns factorial(n-1) instead of factorial(n).\n"
            "    if n <= 1:\n"
            "        return 1\n"
            "    result = 1\n"
            "    for i in range(1, n):\n"
            "        result *= i\n"
            "    return result\n"
        ),
        "fixed_src": (
            "def factorial(n):\n"
            "    if n <= 1:\n"
            "        return 1\n"
            "    result = 1\n"
            "    for i in range(1, n + 1):\n"
            "        result *= i\n"
            "    return result\n"
        ),
        "test_body": (
            # Validation fix (CRITICAL): host-side path resolution.
            "import importlib.util, sys, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'factorial_util.py')\n"
            "    spec = importlib.util.spec_from_file_location('factorial_util', src)\n"
            "    mod = importlib.util.module_from_spec(spec)\n"
            "    spec.loader.exec_module(mod)\n"
            "    if mod.factorial(0) != 1: return False\n"
            "    if mod.factorial(1) != 1: return False\n"
            "    if mod.factorial(5) != 120: return False\n"
            "    if mod.factorial(7) != 5040: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/factorial_util.py and fix the off-by-one bug in the `factorial(n)` function so that it correctly returns n! (e.g. factorial(5) == 120, factorial(7) == 5040). Save the file in place. Do not change the function signature.",
        ],
    },
    # Mass-add — additional fix-bug variants.
    {
        "template_id": "multi_pytest_fix_is_palindrome",
        "src_basename": "string_utils.py",
        "broken_src": (
            "def is_palindrome(s):\n"
            "    # BUG: only checks first half against itself instead of comparing\n"
            "    # the string against its reverse. Fix it!\n"
            "    half = len(s) // 2\n"
            "    return s[:half] == s[:half]\n"
        ),
        "fixed_src": (
            "def is_palindrome(s):\n"
            "    return s == s[::-1]\n"
        ),
        "test_body": (
            "import importlib.util, sys, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'string_utils.py')\n"
            "    spec = importlib.util.spec_from_file_location('string_utils', src)\n"
            "    mod = importlib.util.module_from_spec(spec)\n"
            "    spec.loader.exec_module(mod)\n"
            "    if mod.is_palindrome('racecar') is not True: return False\n"
            "    if mod.is_palindrome('level')   is not True: return False\n"
            "    if mod.is_palindrome('hello')   is not False: return False\n"
            "    if mod.is_palindrome('python')  is not False: return False\n"
            "    if mod.is_palindrome('')        is not True: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/string_utils.py and fix the bug in `is_palindrome(s)` so it correctly returns True iff s reads the same forwards and backwards (e.g. 'racecar' → True, 'hello' → False). Save the file in place. Do not change the function signature.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_max_in_list",
        "src_basename": "list_utils.py",
        "broken_src": (
            "def max_in_list(xs):\n"
            "    # BUG: returns the first element instead of the maximum. Fix it!\n"
            "    return xs[0]\n"
        ),
        "fixed_src": (
            "def max_in_list(xs):\n"
            "    m = xs[0]\n"
            "    for x in xs[1:]:\n"
            "        if x > m:\n"
            "            m = x\n"
            "    return m\n"
        ),
        "test_body": (
            "import importlib.util, sys, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'list_utils.py')\n"
            "    spec = importlib.util.spec_from_file_location('list_utils', src)\n"
            "    mod = importlib.util.module_from_spec(spec)\n"
            "    spec.loader.exec_module(mod)\n"
            "    if mod.max_in_list([1, 2, 3])       != 3: return False\n"
            "    if mod.max_in_list([5, 4, 3, 2, 1]) != 5: return False\n"
            "    if mod.max_in_list([-7, -3, -10])   != -3: return False\n"
            "    if mod.max_in_list([42])            != 42: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/list_utils.py and fix the bug in `max_in_list(xs)` so that it returns the largest element of the list (e.g. max_in_list([1,2,3]) == 3, max_in_list([-7,-3,-10]) == -3). Save the file in place. Do not change the function signature.",
        ],
    },
    # Added: 6 more fix-bug specs to address Python dev/test gap
    # (eval has check_python_file_by_test_suite tasks; synth was under-covered
    # relative to skill diversity — these add idiomatic Python patterns:
    # recursion, string parsing, dict merge, sort key, generator, class state).
    {
        "template_id": "multi_pytest_fix_fibonacci",
        "src_basename": "fib_util.py",
        "broken_src": (
            "def fib(n):\n"
            "    # BUG: base case returns n+1 instead of n; iterative loop also off.\n"
            "    if n < 2:\n"
            "        return n + 1\n"
            "    a, b = 0, 1\n"
            "    for _ in range(n):\n"
            "        a, b = b, a + b\n"
            "    return a\n"
        ),
        "fixed_src": (
            "def fib(n):\n"
            "    if n < 2:\n"
            "        return n\n"
            "    a, b = 0, 1\n"
            "    for _ in range(n - 1):\n"
            "        a, b = b, a + b\n"
            "    return b\n"
        ),
        "test_body": (
            "import importlib.util, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'fib_util.py')\n"
            "    spec = importlib.util.spec_from_file_location('fib_util', src)\n"
            "    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)\n"
            "    for n, want in [(0,0),(1,1),(2,1),(5,5),(10,55),(15,610)]:\n"
            "        if mod.fib(n) != want: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/fib_util.py and fix the bugs in `fib(n)` so it returns the n-th Fibonacci number (fib(0)=0, fib(1)=1, fib(2)=1, fib(10)=55, fib(15)=610). Save the file in place. Do not change the function signature.",
            "There is a buggy `fib(n)` in ~/Desktop/fib_util.py — both the base-case return value and the iteration count are wrong. Edit the file so that fib(0)==0, fib(1)==1, fib(2)==1, fib(10)==55, fib(15)==610. Save the file. Keep the same function signature.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_count_vowels",
        "src_basename": "vowel_util.py",
        "broken_src": (
            "def count_vowels(s):\n"
            "    # BUG: only counts lowercase 'a'; should count all vowels (case-insensitive).\n"
            "    return sum(1 for c in s if c == 'a')\n"
        ),
        "fixed_src": (
            "def count_vowels(s):\n"
            "    return sum(1 for c in s.lower() if c in 'aeiou')\n"
        ),
        "test_body": (
            "import importlib.util, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'vowel_util.py')\n"
            "    spec = importlib.util.spec_from_file_location('vowel_util', src)\n"
            "    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)\n"
            "    for s, want in [('hello',2),('PYTHON',1),('AEIOUaeiou',10),('xyz',0),('',0),('Quick Brown Fox',4)]:\n"
            "        if mod.count_vowels(s) != want: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/vowel_util.py and fix `count_vowels(s)` so it counts every vowel (a, e, i, o, u — case-insensitive) in the input. e.g. count_vowels('hello') == 2, count_vowels('AEIOUaeiou') == 10. Save the file in place. Do not change the function signature.",
            "Edit ~/Desktop/vowel_util.py so that `count_vowels(s)` returns the total number of vowels (a, e, i, o, u) in `s`, ignoring case. e.g. count_vowels('AEIOUaeiou') must return 10 and count_vowels('xyz') must return 0. Save the file in place; do not change the signature.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_dict_merge",
        "src_basename": "dict_util.py",
        "broken_src": (
            "def merge(a, b):\n"
            "    # BUG: returns `a` only; should merge `b` into a copy of `a` (b wins on conflict).\n"
            "    return a\n"
        ),
        "fixed_src": (
            "def merge(a, b):\n"
            "    out = dict(a)\n"
            "    out.update(b)\n"
            "    return out\n"
        ),
        "test_body": (
            "import importlib.util, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'dict_util.py')\n"
            "    spec = importlib.util.spec_from_file_location('dict_util', src)\n"
            "    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)\n"
            "    a = {'x': 1, 'y': 2}\n"
            "    b = {'y': 99, 'z': 3}\n"
            "    out = mod.merge(a, b)\n"
            "    if out != {'x': 1, 'y': 99, 'z': 3}: return False\n"
            "    # Source dicts must not be mutated.\n"
            "    if a != {'x': 1, 'y': 2}: return False\n"
            "    if b != {'y': 99, 'z': 3}: return False\n"
            "    if mod.merge({}, {'k': 1}) != {'k': 1}: return False\n"
            "    if mod.merge({'k': 1}, {}) != {'k': 1}: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/dict_util.py and fix `merge(a, b)` so it returns a NEW dict that combines a and b, with values from b winning on key conflicts. Do not mutate the inputs. e.g. merge({'x':1,'y':2}, {'y':99,'z':3}) == {'x':1,'y':99,'z':3}. Save the file in place.",
            "Edit ~/Desktop/dict_util.py so that `merge(a, b)` returns a fresh dict combining keys from both inputs, with `b`'s values overriding on conflicts, WITHOUT mutating `a` or `b`. e.g. merge({'x':1,'y':2},{'y':99,'z':3}) must return {'x':1,'y':99,'z':3}. Save in place.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_sort_by_length",
        "src_basename": "sort_util.py",
        "broken_src": (
            "def sort_by_length(words):\n"
            "    # BUG: returns lexicographic sort, not length-based. Fix it!\n"
            "    return sorted(words)\n"
        ),
        "fixed_src": (
            "def sort_by_length(words):\n"
            "    return sorted(words, key=len)\n"
        ),
        "test_body": (
            "import importlib.util, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'sort_util.py')\n"
            "    spec = importlib.util.spec_from_file_location('sort_util', src)\n"
            "    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)\n"
            "    out = mod.sort_by_length(['banana', 'kiwi', 'fig', 'apple'])\n"
            "    if [len(w) for w in out] != sorted([6, 4, 3, 5]): return False\n"
            "    if out[0] != 'fig': return False\n"
            "    if mod.sort_by_length([]) != []: return False\n"
            "    out2 = mod.sort_by_length(['aa', 'b', 'ccc'])\n"
            "    if [len(w) for w in out2] != [1, 2, 3]: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/sort_util.py and fix `sort_by_length(words)` so it returns the words sorted by string length ascending (Python's sorted() is stable, so equal-length words keep input order). e.g. sort_by_length(['banana','kiwi','fig','apple']) starts with 'fig' (len 3). Save the file in place.",
            "Edit ~/Desktop/sort_util.py so `sort_by_length(words)` returns the words sorted in ascending order of string length (stable). e.g. sort_by_length(['aa','b','ccc']) == ['b','aa','ccc']. Save the file; do not rename or re-signature the function.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_running_sum",
        "src_basename": "rsum_util.py",
        "broken_src": (
            "def running_sum(xs):\n"
            "    # BUG: returns the totals list but resets after each item — should\n"
            "    # accumulate. e.g. running_sum([1,2,3]) should be [1,3,6], not [1,2,3].\n"
            "    return list(xs)\n"
        ),
        "fixed_src": (
            "def running_sum(xs):\n"
            "    out = []\n"
            "    total = 0\n"
            "    for x in xs:\n"
            "        total += x\n"
            "        out.append(total)\n"
            "    return out\n"
        ),
        "test_body": (
            "import importlib.util, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'rsum_util.py')\n"
            "    spec = importlib.util.spec_from_file_location('rsum_util', src)\n"
            "    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)\n"
            "    for xs, want in [([1,2,3],[1,3,6]),([],[]),([5],[5]),([1,-1,1,-1],[1,0,1,0]),([10,20,30,40],[10,30,60,100])]:\n"
            "        if mod.running_sum(xs) != want: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/rsum_util.py and fix `running_sum(xs)` so it returns the running (cumulative) sums of the input list. e.g. running_sum([1,2,3]) == [1,3,6]; running_sum([10,20,30,40]) == [10,30,60,100]; running_sum([]) == []. Save the file in place.",
            "Edit ~/Desktop/rsum_util.py so that `running_sum(xs)` returns the cumulative sums of the input list. e.g. running_sum([1,2,3]) must return [1,3,6], running_sum([1,-1,1,-1]) must return [1,0,1,0], running_sum([]) must return []. Save the file in place.",
        ],
    },
    {
        "template_id": "multi_pytest_fix_counter_class",
        "src_basename": "counter_class.py",
        "broken_src": (
            "class Counter:\n"
            "    def __init__(self, start=0):\n"
            "        self.value = start\n"
            "    def bump(self):\n"
            "        # BUG: bump should increment by 1, but currently does nothing.\n"
            "        pass\n"
            "    def reset(self):\n"
            "        # BUG: reset should set value back to the starting value, but\n"
            "        # currently sets it to 0 regardless of the constructor arg.\n"
            "        self.value = 0\n"
        ),
        "fixed_src": (
            "class Counter:\n"
            "    def __init__(self, start=0):\n"
            "        self._start = start\n"
            "        self.value = start\n"
            "    def bump(self):\n"
            "        self.value += 1\n"
            "    def reset(self):\n"
            "        self.value = self._start\n"
        ),
        "test_body": (
            "import importlib.util, os\n"
            "def test():\n"
            "    src = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'counter_class.py')\n"
            "    spec = importlib.util.spec_from_file_location('counter_class', src)\n"
            "    mod = importlib.util.module_from_spec(spec); spec.loader.exec_module(mod)\n"
            "    c = mod.Counter()\n"
            "    if c.value != 0: return False\n"
            "    c.bump(); c.bump(); c.bump()\n"
            "    if c.value != 3: return False\n"
            "    c.reset()\n"
            "    if c.value != 0: return False\n"
            "    c2 = mod.Counter(start=10)\n"
            "    if c2.value != 10: return False\n"
            "    c2.bump()\n"
            "    if c2.value != 11: return False\n"
            "    c2.reset()\n"
            "    if c2.value != 10: return False\n"
            "    return True\n"
        ),
        "instructions": [
            "Open ~/Desktop/counter_class.py and fix the two bugs in the `Counter` class: (1) `bump()` should increment `value` by 1 (currently does nothing); (2) `reset()` should restore `value` to the constructor's `start` argument (currently it hardcodes 0). Save the file in place. Do not change the constructor signature or rename methods.",
            "Edit ~/Desktop/counter_class.py to fix two bugs in `Counter`: (a) `bump()` must increment `value` by 1; (b) `reset()` must restore `value` to the value of the `start` arg passed to `__init__` (NOT hardcoded zero). Keep the constructor signature `__init__(self, start=0)`. Save in place.",
        ],
    },
]


def _make_py_test_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    fixed_path    = f"{_TMP}/fixed_{template_id}.py"
    test_path     = f"{_TMP}/test_suite_{template_id}.py"
    broken_src    = spec["broken_src"]
    fixed_src     = spec["fixed_src"]
    test_body     = spec["test_body"]
    instructions  = spec["instructions"]

    init_steps = [
        _write_text_step(src_path, broken_src),
        _write_text_step(fixed_path, fixed_src),
        _write_text_step(test_path, test_body),
        # validation gap-close: terminal-driven (agent edits .py via shell/editor) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{fixed_path}' '{src_path}'")]

    evaluator = {
        "func": "check_python_file_by_test_suite",
        # `actual_files` arg ← result. Single-file form (multi not needed).
        "result": {"type": "vm_file", "path": src_path, "dest": spec["src_basename"]},
        # `test_file` arg ← expected. The eval calls test() in the file.
        "expected": {"type": "vm_file", "path": test_path, "dest": "test_suite.py"},
    }

    def _params(seed: int) -> dict:
        rng = random.Random(seed ^ _stable_hash(template_id) & 0xFFFFFFFF)
        return {
            "instr": instructions[seed % len(instructions)],
            "config_override": init_steps,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_python_file_by_test_suite",
        setup_class="python_fix_bug",
    )


# ---------------------------------------------------------------------------
# Cat R — Long-tail gap-fill ( 2026-05-10).
# ---------------------------------------------------------------------------
# Bridges UNDER-covered eval evaluator.func buckets (synth=0 vs eval=1-5):
#   diff_text_file, check_list, check_mp3_meta, compare_image_text,
#   compare_epub, is_expected_tabs, file_contains, check_json,
#   check_structure_sim, is_in_list, compare_zip_files, check_image_size,
#   check_direct_json_object, compare_python_pure_text, compare_htmls,
#   check_image_file_size, compare_pdf_images, check_line_number,
#   compare_audios, check_thunderbird_folder, is_expected_search_query,
#   is_expected_bookmarks, literal_match, compare_docx_files_and_ignore_new_lines,
#   compare_config, is_expected_url_pattern_match.
# Each template is a real multi-app (terminal+app or app1+app2) workflow
# whose oracle reproduces the gold artifact via shell pipeline. Per Cap-2x2:
# at most 2 instructions per spec (2 paraphrases) and n_rows<=2 per template.
# ---------------------------------------------------------------------------


# ---- Cat R.1 — diff_text_file (text file diff, lenient EOL) ----------------
# Pattern: pre-config writes a base text file AND a (slightly modified) gold
# text file. Agent runs a shell pipeline (sed / awk / sort) that produces a
# transformed file; oracle copies the gold over the agent's sink. Eval =
# diff_text_file(result, expected) — difflib SequenceMatcher line-level.

_DIFF_TEXT_FILE_SPECS: list[dict] = [
    {
        "template_id": "multi_diff_text_sed_replace_lines",
        "src_basename": "log_raw.txt",
        "src_lines": [
            "INFO 2026-05-10 12:00:00 service start",
            "WARN 2026-05-10 12:00:05 cache miss key=42",
            "ERROR 2026-05-10 12:00:10 connection refused",
            "INFO 2026-05-10 12:00:12 retrying upstream",
            "INFO 2026-05-10 12:00:15 service ready",
        ],
        "dst_basename": "log_clean.txt",
        "expected_lines": [
            "INFO 2026-05-10 12:00:00 service start",
            "INFO 2026-05-10 12:00:12 retrying upstream",
            "INFO 2026-05-10 12:00:15 service ready",
        ],
        "instructions": [
            "Open a terminal and filter ~/Desktop/log_raw.txt so only the lines that start with 'INFO ' (uppercase, followed by a space) are kept. Preserve their original order and save the result to ~/Desktop/log_clean.txt.",
        ],
    },
    {
        "template_id": "multi_diff_text_sort_names",
        "src_basename": "names_raw.txt",
        "src_lines": ["Charlie", "Alice", "Eve", "Bob", "Dave"],
        "dst_basename": "names_sorted.txt",
        "expected_lines": ["Alice", "Bob", "Charlie", "Dave", "Eve"],
        "instructions": [
            "In a terminal, sort ~/Desktop/names_raw.txt alphabetically and write the result to ~/Desktop/names_sorted.txt. Hint: `sort ~/Desktop/names_raw.txt > ~/Desktop/names_sorted.txt`.",
        ],
    },
    {
        "template_id": "multi_diff_text_uniq_dedupe",
        "src_basename": "ids_raw.txt",
        "src_lines": ["A1", "A2", "A1", "A3", "A2", "A4", "A1"],
        "dst_basename": "ids_unique.txt",
        "expected_lines": ["A1", "A2", "A3", "A4"],
        "instructions": [
            "From a terminal, take ~/Desktop/ids_raw.txt, remove duplicate lines, sort the remaining unique IDs in ascending byte order, and write them (one per line) to ~/Desktop/ids_unique.txt.",
        ],
    },
    {
        "template_id": "multi_diff_text_lowercase_lines",
        "src_basename": "tags_raw.txt",
        "src_lines": ["Apple", "Banana", "Cherry", "Date", "Elderberry"],
        "dst_basename": "tags_lower.txt",
        "expected_lines": ["apple", "banana", "cherry", "date", "elderberry"],
        "instructions": [
            "In a terminal, lowercase every character in ~/Desktop/tags_raw.txt and write the result to ~/Desktop/tags_lower.txt. Hint: `tr '[:upper:]' '[:lower:]' < ~/Desktop/tags_raw.txt > ~/Desktop/tags_lower.txt`.",
        ],
    },
    {
        "template_id": "multi_diff_text_strip_blanks",
        "src_basename": "notes_raw.txt",
        "src_lines": ["First line", "", "Second line", "", "", "Third line", ""],
        "dst_basename": "notes_compact.txt",
        "expected_lines": ["First line", "Second line", "Third line"],
        "instructions": [
            "Open a terminal and strip out every blank/empty line from ~/Desktop/notes_raw.txt. Keep the remaining non-empty lines in their original order and write the compact result to ~/Desktop/notes_compact.txt.",
        ],
    },
    {
        "template_id": "multi_diff_text_first_field_csv",
        "src_basename": "rows_raw.txt",
        "src_lines": ["id1,Alice,A", "id2,Bob,B", "id3,Carol,A", "id4,Dave,C"],
        "dst_basename": "ids_only.txt",
        "expected_lines": ["id1", "id2", "id3", "id4"],
        "instructions": [
            "In a terminal, extract the first comma-separated field from every line of ~/Desktop/rows_raw.txt and save the resulting list to ~/Desktop/ids_only.txt. Hint: `cut -d',' -f1 ~/Desktop/rows_raw.txt > ~/Desktop/ids_only.txt`.",
        ],
    },
    {
        "template_id": "multi_diff_text_reverse_lines",
        "src_basename": "stack_raw.txt",
        "src_lines": ["push 1", "push 2", "push 3", "push 4", "push 5"],
        "dst_basename": "stack_lifo.txt",
        "expected_lines": ["push 5", "push 4", "push 3", "push 2", "push 1"],
        "instructions": [
            "In a terminal, reverse the line order of ~/Desktop/stack_raw.txt and save the result to ~/Desktop/stack_lifo.txt. Hint: `tac ~/Desktop/stack_raw.txt > ~/Desktop/stack_lifo.txt`.",
        ],
    },
    {
        "template_id": "multi_diff_text_wc_words_per_line",
        "src_basename": "sentences_raw.txt",
        "src_lines": [
            "the quick brown fox",
            "jumps over the lazy dog",
            "and runs away",
        ],
        "dst_basename": "wordcounts.txt",
        "expected_lines": ["4", "5", "3"],
        "instructions": [
            "From a terminal, count how many whitespace-separated words appear on each line of ~/Desktop/sentences_raw.txt. Write the per-line word counts as one integer per line (in input order) to ~/Desktop/wordcounts.txt.",
        ],
    },
]


def _make_diff_text_file_template(spec: dict) -> SynthTemplate:
    template_id  = spec["template_id"]
    src_path     = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path     = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.txt"
    src_text     = "\n".join(spec["src_lines"]) + "\n"
    expected     = "\n".join(spec["expected_lines"]) + "\n"
    instructions = spec["instructions"]

    init_steps = [
        _write_text_step(src_path, src_text),
        _write_text_step(expected_path, expected),
        # Pre-open a terminal (validation H-cluster fix): these rows always say
        # "In a terminal, ..." but the bare desktop has no terminal dock icon.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "diff_text_file",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.txt"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.txt"},
    }

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)],
                "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="diff_text_file",
        setup_class="diff_text_shell",
    )


# ---- Cat R.2 — check_list (regex-list patterns over text file) ------------
# Pattern: agent produces a list (one item per line). Eval check_list applies
# `expect` regexes (each must match at least one line) + `unexpect` regexes
# (none may match any line). Oracle writes the gold list directly.

_CHECK_LIST_SPECS: list[dict] = [
    {
        "template_id": "multi_check_list_ls_files",
        "src_files": [
            "report_q1.docx", "report_q2.docx", "report_q3.docx",
            "summary.txt", "draft.odt",
        ],
        "src_dir_basename": "office_docs",
        "ans_basename": "doc_listing.txt",
        # gold = actual ls output of source dir (sorted).
        "gold_lines": ["draft.odt", "report_q1.docx", "report_q2.docx",
                       "report_q3.docx", "summary.txt"],
        "expect_regexes": [
            r"report_q1\.docx", r"report_q2\.docx", r"summary\.txt",
        ],
        "unexpect_regexes": [r"\.exe$", r"^trash"],
        "instructions": [
            "In a terminal, list every file in ~/Desktop/office_docs/ alphabetically (one filename per line, no path) and save the listing to ~/Desktop/doc_listing.txt. Hint: `ls ~/Desktop/office_docs/ | sort > ~/Desktop/doc_listing.txt`.",
        ],
    },
    {
        "template_id": "multi_check_list_py_files",
        "src_files": [
            "main.py", "utils.py", "config.py", "test_utils.py", "README.md",
        ],
        "src_dir_basename": "py_proj",
        "ans_basename": "py_listing.txt",
        "gold_lines": ["config.py", "main.py", "test_utils.py", "utils.py"],
        "expect_regexes": [r"main\.py", r"utils\.py", r"config\.py"],
        "unexpect_regexes": [r"\.md$", r"^README"],
        "instructions": [
            "In a terminal, list every .py file in ~/Desktop/py_proj/ alphabetically (one filename per line, no path) and save the listing to ~/Desktop/py_listing.txt. Hint: `ls ~/Desktop/py_proj/*.py | xargs -n1 basename | sort > ~/Desktop/py_listing.txt`.",
        ],
    },
    {
        "template_id": "multi_check_list_grep_errors",
        "src_lines": [
            "INFO ok",
            "ERROR connection refused on port 5432",
            "INFO retrying",
            "ERROR timeout fetching upstream",
            "WARN slow query",
            "ERROR disk full /var/log",
            "INFO ready",
        ],
        "src_basename": "service.log",
        "ans_basename": "errors.txt",
        "gold_lines": [
            "ERROR connection refused on port 5432",
            "ERROR timeout fetching upstream",
            "ERROR disk full /var/log",
        ],
        "expect_regexes": [r"^ERROR ", r"connection refused", r"disk full"],
        "unexpect_regexes": [r"^INFO ", r"^WARN "],
        "instructions": [
            "From a terminal, scan ~/Desktop/service.log and pull out every line that begins with 'ERROR ' (uppercase, followed by a space). Keep the original ordering and save the matching lines to ~/Desktop/errors.txt.",
        ],
    },
    {
        "template_id": "multi_check_list_uniq_users",
        "src_lines": [
            "alice login",
            "bob login",
            "alice action",
            "carol login",
            "bob action",
            "alice logout",
        ],
        "src_basename": "events.log",
        "ans_basename": "users.txt",
        "gold_lines": ["alice", "bob", "carol"],
        "expect_regexes": [r"^alice$", r"^bob$", r"^carol$"],
        "unexpect_regexes": [r"login", r"logout", r"action"],
        "instructions": [
            "From a terminal, look at ~/Desktop/events.log. The first whitespace-delimited token on each line is a username. Build a sorted list of every distinct username and save it (one per line) to ~/Desktop/users.txt.",
        ],
    },
    {
        "template_id": "multi_check_list_jpg_files",
        "photo_assets": [
            ("photos/food/pizza-dish.jpg",      "pizza-dish.jpg"),
            ("photos/food/coffee-latte.jpg",    "coffee-latte.jpg"),
            ("photos/food/restaurant-meal.jpg", "restaurant-meal.jpg"),
            ("photos/food/salad-bowl.jpg",      "salad-bowl.jpg"),
        ],
        "src_dir_basename": "food_pix",
        "ans_basename": "jpgs.txt",
        "gold_lines": ["coffee-latte.jpg", "pizza-dish.jpg",
                       "restaurant-meal.jpg", "salad-bowl.jpg"],
        "expect_regexes": [r"pizza", r"coffee", r"\.jpg$"],
        "unexpect_regexes": [r"\.png$", r"\.gif$"],
        "instructions": [
            "In a terminal, list every .jpg file in ~/Desktop/food_pix/ alphabetically (one filename per line, no path) and save the listing to ~/Desktop/jpgs.txt. Hint: `ls ~/Desktop/food_pix/*.jpg | xargs -n1 basename | sort > ~/Desktop/jpgs.txt`.",
        ],
    },
    {
        "template_id": "multi_check_list_top_lines",
        "src_lines": [
            "alpha", "beta", "gamma", "delta", "epsilon",
            "zeta", "eta", "theta", "iota", "kappa",
        ],
        "src_basename": "greek.txt",
        "ans_basename": "head5.txt",
        "gold_lines": ["alpha", "beta", "gamma", "delta", "epsilon"],
        "expect_regexes": [r"^alpha$", r"^delta$", r"^epsilon$"],
        "unexpect_regexes": [r"^zeta$", r"^kappa$"],
        "instructions": [
            "Open a terminal and copy the first five lines of ~/Desktop/greek.txt (in their original order) into a new file at ~/Desktop/head5.txt.",
        ],
    },
    {
        "template_id": "multi_check_list_csv_emails",
        "src_lines": [
            "name,email",
            "Alice,alice@example.com",
            "Bob,bob@corp.io",
            "Carol,carol@example.com",
            "Dave,dave@corp.io",
        ],
        "src_basename": "people.csv",
        "ans_basename": "emails.txt",
        "gold_lines": ["alice@example.com", "bob@corp.io",
                       "carol@example.com", "dave@corp.io"],
        "expect_regexes": [r"^alice@", r"^bob@", r"@example\.com$"],
        "unexpect_regexes": [r"^name$", r"^email$"],
        "instructions": [
            "From a terminal, pull just the email column (the 2nd comma-separated field) out of ~/Desktop/people.csv, skipping the header row. Write the four resulting email addresses, one per line in their original order, to ~/Desktop/emails.txt.",
        ],
    },
    {
        "template_id": "multi_check_list_words_dict",
        "src_lines": [
            "apple", "banana", "apricot", "blueberry", "avocado", "blackberry",
        ],
        "src_basename": "fruits.txt",
        "ans_basename": "fruits_a.txt",
        "gold_lines": ["apple", "apricot", "avocado"],
        "expect_regexes": [r"^apple$", r"^apricot$", r"^avocado$"],
        "unexpect_regexes": [r"^banana$", r"^blueberry$", r"^blackberry$"],
        "instructions": [
            "Open a terminal and filter ~/Desktop/fruits.txt to keep only the lines that start with a lowercase 'a' (case-sensitive). Preserve their original order and save the result to ~/Desktop/fruits_a.txt.",
        ],
    },
]


def _make_check_list_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    ans_path      = f"{_DESKTOP}/{spec['ans_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.txt"
    expect_regs   = spec["expect_regexes"]
    unexpect_regs = spec["unexpect_regexes"]
    instructions  = spec["instructions"]
    gold_text     = "\n".join(spec["gold_lines"]) + "\n"

    init_steps: list[dict] = []
    # Stage source: either a single text file, a directory of text files, or
    # photo assets in a dir.
    if "src_basename" in spec:
        src_path = f"{_DESKTOP}/{spec['src_basename']}"
        src_text = "\n".join(spec["src_lines"]) + "\n"
        init_steps.append(_write_text_step(src_path, src_text))
    elif "src_files" in spec:
        src_dir = f"{_DESKTOP}/{spec['src_dir_basename']}"
        init_steps.append(_execute(f"mkdir -p '{src_dir}'"))
        for fname in spec["src_files"]:
            init_steps.append(_write_text_step(f"{src_dir}/{fname}",
                                               f"placeholder content for {fname}\n"))
    else:  # photo_assets
        src_dir = f"{_DESKTOP}/{spec['src_dir_basename']}"
        init_steps.append(_execute(f"mkdir -p '{src_dir}'"))
        for (rel, base) in spec["photo_assets"]:
            init_steps.append(_stage_asset(rel, f"{src_dir}/{base}"))

    init_steps.append(_write_text_step(expected_path, gold_text))
    # Pre-open a terminal (validation H-cluster fix) — check_list rows all say
    # "In a terminal, ..." but lite.osworld bare desktop has no dock icon, so
    # the agent's first type_text would race against a blind dock-click guess.
    init_steps.extend(_terminal_preopen_steps())
    oracle_steps = [_execute(f"cp '{expected_path}' '{ans_path}'")]

    evaluator = {
        "func": "check_list",
        "result": {"type": "vm_file", "path": ans_path, "dest": "result.txt"},
        "expected": {"type": "rule",
                     "rules": {"expect": expect_regs, "unexpect": unexpect_regs}},
    }

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)],
                "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_list",
        setup_class="shell_list_to_file",
    )


# ---- Cat R.3 — check_mp3_meta (mutagen ID3 tag set on ffmpeg sine wave) ----
# Pattern: pre_config generates a 1-second silent sine-wave mp3 via ffmpeg
# (no asset add) AND writes the gold tagged file (also via mutagen). Agent's
# task is to set ID3 tags via mutagen / kid3-cli on the source mp3. Oracle =
# cp gold over result. Eval = check_mp3_meta(method=eq).

_CHECK_MP3_META_SPECS: list[dict] = [
    {
        "template_id": "multi_mp3_meta_song_a",
        "basename": "track_a.mp3",
        "tags": {"title": "Sunrise", "artist": "Aurora"},
        "instructions": [
            "In a terminal, set the ID3v2 metadata on ~/Desktop/track_a.mp3 so that title='Sunrise' and artist='Aurora'. Hint: use python+mutagen: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_a.mp3'); t['title']='Sunrise'; t['artist']='Aurora'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_b",
        "basename": "track_b.mp3",
        "tags": {"title": "Moonlight Drive", "artist": "Echo Valley"},
        "instructions": [
            "In a terminal, set the ID3v2 tags on ~/Desktop/track_b.mp3 to title='Moonlight Drive' and artist='Echo Valley'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_b.mp3'); t['title']='Moonlight Drive'; t['artist']='Echo Valley'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_c_album",
        "basename": "track_c.mp3",
        "tags": {"title": "Quiet Field", "artist": "Pine River", "album": "Dawn Sessions"},
        "instructions": [
            "In a terminal, set the ID3v2 metadata on ~/Desktop/track_c.mp3 so that title='Quiet Field', artist='Pine River' and album='Dawn Sessions'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_c.mp3'); t['title']='Quiet Field'; t['artist']='Pine River'; t['album']='Dawn Sessions'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_d_album",
        "basename": "track_d.mp3",
        "tags": {"title": "Last Train", "artist": "Iron Sky", "album": "Night Shift"},
        "instructions": [
            "In a terminal, set the ID3v2 metadata on ~/Desktop/track_d.mp3 so that title='Last Train', artist='Iron Sky' and album='Night Shift'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_d.mp3'); t['title']='Last Train'; t['artist']='Iron Sky'; t['album']='Night Shift'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_e_full",
        "basename": "track_e.mp3",
        "tags": {"title": "Open Road", "artist": "Compass North", "album": "Maps"},
        "instructions": [
            "In a terminal, set the ID3v2 metadata on ~/Desktop/track_e.mp3 so that title='Open Road', artist='Compass North' and album='Maps'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_e.mp3'); t['title']='Open Road'; t['artist']='Compass North'; t['album']='Maps'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_f",
        "basename": "track_f.mp3",
        "tags": {"title": "Velvet Sky", "artist": "Indigo Cat"},
        "instructions": [
            "In a terminal, set the ID3v2 tags on ~/Desktop/track_f.mp3 to title='Velvet Sky' and artist='Indigo Cat'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_f.mp3'); t['title']='Velvet Sky'; t['artist']='Indigo Cat'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_g_album",
        "basename": "track_g.mp3",
        "tags": {"title": "Harbor Lights", "artist": "Marina Days", "album": "Coastal"},
        "instructions": [
            "In a terminal, set the ID3v2 metadata on ~/Desktop/track_g.mp3 so that title='Harbor Lights', artist='Marina Days' and album='Coastal'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_g.mp3'); t['title']='Harbor Lights'; t['artist']='Marina Days'; t['album']='Coastal'; t.save()\"`.",
        ],
    },
    {
        "template_id": "multi_mp3_meta_song_h",
        "basename": "track_h.mp3",
        "tags": {"title": "Midnight Train", "artist": "Iron Cosmic"},
        "instructions": [
            "In a terminal, set the ID3v2 tags on ~/Desktop/track_h.mp3 to title='Midnight Train' and artist='Iron Cosmic'. Hint: `python3 -c \"from mutagen.easyid3 import EasyID3; t=EasyID3('/home/user/Desktop/track_h.mp3'); t['title']='Midnight Train'; t['artist']='Iron Cosmic'; t.save()\"`.",
        ],
    },
]


def _make_check_mp3_meta_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.mp3"
    tags          = spec["tags"]
    instructions  = spec["instructions"]

    # ffmpeg sine wave (1s silent-ish sine) → src AND expected; mutagen sets
    # tags on expected only. Source file is intentionally untagged.
    tag_apply_lines = "\n".join(
        f"        t[{k!r}] = {v!r}" for k, v in tags.items()
    )
    build_py = textwrap.dedent(f"""\
        import os, subprocess
        src = {src_path!r}
        exp = {expected_path!r}
        os.makedirs(os.path.dirname(src) or '.', exist_ok=True)
        os.makedirs(os.path.dirname(exp) or '.', exist_ok=True)
        # 1-second sine wave at 440Hz, low-bitrate mp3.
        for path in (src, exp):
            subprocess.run([
                "ffmpeg", "-y", "-f", "lavfi", "-i", "sine=frequency=440:duration=1",
                "-ac", "1", "-b:a", "64k", "-loglevel", "error", path
            ], check=True)
        # Apply tags to the gold copy ONLY.
        from mutagen.easyid3 import EasyID3
        from mutagen.id3 import ID3NoHeaderError
        try:
            t = EasyID3(exp)
        except ID3NoHeaderError:
            from mutagen.mp3 import MP3
            m = MP3(exp); m.add_tags(); m.save()
            t = EasyID3(exp)
{tag_apply_lines}
        t.save()
        """)
    init_steps = [
        _python_step(build_py),
        # validation gap-close: terminal-driven (mutagen python3 via shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{src_path}'")]

    evaluator_funcs = ["check_mp3_meta"]
    evaluator = {
        "func": "check_mp3_meta",
        "result": {"type": "vm_file", "path": src_path, "dest": "result.mp3"},
        "expected": {"type": "rule", "rules": {
            k: {"method": "eq", "ref": v} for k, v in tags.items()
        }},
    }

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)],
                "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_mp3_meta",
        setup_class="ffmpeg_mp3_meta",
    )


# ---- Cat R.4 — compare_image_text (text rendered in PNG, OCR-eq) ----------
# Pattern: ImageMagick `convert label:'TEXT'` produces a PNG with embedded
# text. Both gold and oracle use the same `convert` invocation. Eval =
# compare_image_text rule={type:'text', text:'TEXT'} (easyocr).

_COMPARE_IMAGE_TEXT_SPECS: list[dict] = [
    {
        "template_id": "multi_image_text_hello",
        "basename": "hello_label.png",
        "text": "Hello World",
    },
    {
        "template_id": "multi_image_text_invoice",
        "basename": "invoice_label.png",
        "text": "INVOICE 12345",
    },
    {
        "template_id": "multi_image_text_warning",
        "basename": "warning_label.png",
        "text": "WARNING",
    },
    {
        "template_id": "multi_image_text_open",
        "basename": "open_label.png",
        "text": "OPEN NOW",
    },
    {
        "template_id": "multi_image_text_sale",
        "basename": "sale_label.png",
        "text": "SALE 50",
    },
    {
        "template_id": "multi_image_text_exit",
        "basename": "exit_label.png",
        "text": "EXIT",
    },
]


def _make_compare_image_text_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    dst_path      = f"{_DESKTOP}/{spec['basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.png"
    text          = spec["text"]

    init_steps = [
        _execute(
            f"convert -background white -fill black -font DejaVu-Sans-Bold "
            f"-pointsize 72 label:{text!r} '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (ImageMagick convert shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_image_text",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.png"},
        "expected": {"type": "rule", "rules": {"type": "text", "text": text}},
    }
    instructions = [
        f"From a terminal, please use ImageMagick to render the text {text!r} as a PNG label image saved to ~/Desktop/{spec['basename']}. Use a white background, black text, and the DejaVu-Sans-Bold font at point size 72.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_image_text",
        setup_class="im_label_text",
    )


# ---- Cat R.5 — compare_epub (pandoc md/docx → epub) ------------------------
# Pattern: pre-config writes a small markdown/docx source AND builds a gold
# epub via `pandoc -o`. Agent runs same pandoc. Oracle = cp gold over result.
# Eval = compare_epub (extracted-text diff_text_file mean over OPF files).

_COMPARE_EPUB_SPECS: list[dict] = [
    {
        "template_id": "multi_epub_md_short_story",
        "src_basename": "short_story.md",
        "dst_basename": "short_story.epub",
        "md_body": (
            "# A Quiet Morning\n\n"
            "The fog rolled gently over the harbor as the first ferry pushed off.\n\n"
            "## Chapter One\n\n"
            "Lin walked the wooden pier alone. The seabirds called overhead.\n\n"
            "## Chapter Two\n\n"
            "By noon the sun had burned away the mist and the town was awake.\n"
        ),
        "instructions": [
            "From a terminal, please use pandoc to convert the short-story markdown at ~/Desktop/short_story.md into an EPUB at ~/Desktop/short_story.epub. For reproducible output, set the SOURCE_DATE_EPOCH environment variable to 1747008000 when invoking pandoc.",
        ],
    },
    {
        "template_id": "multi_epub_md_recipe_book",
        "src_basename": "recipe_book.md",
        "dst_basename": "recipe_book.epub",
        "md_body": (
            "# Quick Recipes\n\n"
            "## Banana Bread\n\n"
            "Mash three ripe bananas, mix with sugar, butter, and an egg.\n"
            "Stir in flour and baking soda. Bake at 175C for 60 minutes.\n\n"
            "## Tomato Soup\n\n"
            "Saute onion in olive oil. Add canned tomatoes and stock.\n"
            "Simmer 20 minutes, blend, season with salt and basil.\n"
        ),
        "instructions": [
            "I'd like ~/Desktop/recipe_book.md compiled into an EPUB at ~/Desktop/recipe_book.epub. From a terminal, run pandoc with the SOURCE_DATE_EPOCH environment variable set to 1747008000 so the build is reproducible.",
        ],
    },
    {
        "template_id": "multi_epub_md_travel_guide",
        "src_basename": "travel_guide.md",
        "dst_basename": "travel_guide.epub",
        "md_body": (
            "# Lisbon in Three Days\n\n"
            "## Day 1: Alfama and the Castle\n\n"
            "Climb to Sao Jorge for sweeping views of the Tagus.\n\n"
            "## Day 2: Belem\n\n"
            "Visit the Jeronimos Monastery and the Belem Tower; eat pastel de nata.\n\n"
            "## Day 3: Sintra Day Trip\n\n"
            "Take the train and tour the colorful Pena Palace before lunch.\n"
        ),
        "instructions": [
            "Please turn the Lisbon travel guide at ~/Desktop/travel_guide.md into an EPUB at ~/Desktop/travel_guide.epub. From a terminal, use pandoc and set SOURCE_DATE_EPOCH=1747008000 in the environment so the resulting EPUB is reproducible.",
        ],
    },
    {
        "template_id": "multi_epub_md_meeting_minutes",
        "src_basename": "minutes.md",
        "dst_basename": "minutes.epub",
        "md_body": (
            "# Project Sync 2026-05-09\n\n"
            "## Attendees\n\n"
            "Alice, Bob, Carol\n\n"
            "## Decisions\n\n"
            "Adopt the new release schedule starting in June.\n\n"
            "## Action Items\n\n"
            "Alice will draft the migration plan; Bob owns the changelog.\n"
        ),
        "instructions": [
            "From a terminal, please compile the meeting minutes at ~/Desktop/minutes.md into an EPUB at ~/Desktop/minutes.epub using pandoc. Export SOURCE_DATE_EPOCH=1747008000 before running it so the timestamp is deterministic.",
        ],
    },
    {
        "template_id": "multi_epub_md_user_manual",
        "src_basename": "manual.md",
        "dst_basename": "manual.epub",
        "md_body": (
            "# Coffee Grinder User Manual\n\n"
            "## Setup\n\n"
            "Place the grinder on a flat surface; plug in to a grounded outlet.\n\n"
            "## Operation\n\n"
            "Fill the hopper with whole beans, select the grind setting, press start.\n\n"
            "## Cleaning\n\n"
            "Empty the hopper and brush out residual grounds after each session.\n"
        ),
        "instructions": [
            "Could you build the coffee-grinder user manual at ~/Desktop/manual.md into an EPUB at ~/Desktop/manual.epub? From a terminal, run pandoc with SOURCE_DATE_EPOCH set to 1747008000 in the environment so the output is reproducible.",
        ],
    },
    {
        "template_id": "multi_epub_md_changelog",
        "src_basename": "changelog.md",
        "dst_basename": "changelog.epub",
        "md_body": (
            "# Changelog 3.4.0\n\n"
            "## Features\n\n"
            "Added native dark mode and customizable keyboard shortcuts.\n\n"
            "## Fixes\n\n"
            "Resolved memory leak in the file picker; fixed crash on startup.\n"
        ),
        "instructions": [
            "From a terminal, render the changelog at ~/Desktop/changelog.md as an EPUB at ~/Desktop/changelog.epub. Use pandoc, and set SOURCE_DATE_EPOCH=1747008000 in the environment for a deterministic build.",
        ],
    },
]


def _make_compare_epub_template(spec: dict) -> SynthTemplate:
    """validation + Validation fix: `compare_epub` was patched as a local
    override in `lite/gym/envs/lite/osworld/src/eval/metrics.py` (resolved by
    `runner.py` before upstream) — it widens `process_epub`'s noise-strip
    set to also drop `dcterms:modified`, `meta name="generator"`, etc.,
    and matches the pandoc EPUB3 layout (`EPUB/content.opf` + `*.xhtml`)
    before SequenceMatcher diff. Combined with `SOURCE_DATE_EPOCH=1747008000`
    baked into both gold + agent commands below, all 6 epub templates
    should now pass.
    """
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.epub"
    md_body       = spec["md_body"]
    instructions  = spec["instructions"]

    # validation fix (pandoc dcterms:modified non-determinism):
    # `--metadata=date=...` only sets `<dc:date>` — pandoc still embeds a
    # wallclock `<dcterms:modified>` in content.opf, which differs between
    # the gold build and the agent's build. Per pandoc docs, setting the
    # SOURCE_DATE_EPOCH env var fixes `dcterms:modified` to that epoch
    # (reproducible build). 2026-05-12 00:00 UTC = 1747008000. We bake the
    # env var into BOTH the gold builder AND the instructed shell command
    # so the agent's epub matches byte-for-byte at the OPF level.
    _PANDOC_FIXED_EPOCH = "1747008000"  # 2026-05-12 00:00:00 UTC
    init_steps = [
        _write_text_step(src_path, md_body),
        _execute(
            f"SOURCE_DATE_EPOCH={_PANDOC_FIXED_EPOCH} "
            f"pandoc -o '{expected_path}' '{src_path}'"
        ),
        # Pre-open a terminal (validation H-cluster fix) — these rows say
        # "In a terminal, ..." but the bare desktop has no dock icon.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_epub",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.epub"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.epub"},
    }

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)],
                "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_epub",
        setup_class="pandoc_md_to_epub",
    )


# ---- Cat R.6 — is_expected_tabs (chrome multi-tab open) -------------------
# Pattern: stage Wikipedia HTML files into ~/Desktop/, agent's task is to
# launch chrome with all of them in tabs. Oracle does the same. Eval =
# is_expected_tabs(open_tabs_info, rule={type:'url', urls:[file://...]}).

_IS_EXPECTED_TABS_SPECS: list[dict] = [
    {
        "template_id": "multi_tabs_open_food_pair",
        "assets": [
            ("html/wikipedia/coffee.html",  "coffee.html"),
            ("html/wikipedia/pizza.html",   "pizza.html"),
        ],
        "instructions": [
            "Open Chrome with the two staged Wikipedia HTML files at ~/Desktop/coffee.html and ~/Desktop/pizza.html as separate tabs.",
        ],
    },
    {
        "template_id": "multi_tabs_open_landmark_pair",
        "assets": [
            ("html/wikipedia/eiffel-tower.html", "eiffel-tower.html"),
            ("html/wikipedia/mona-lisa.html",    "mona-lisa.html"),
        ],
        "instructions": [
            "Open Chrome with the two staged Wikipedia HTML files at ~/Desktop/eiffel-tower.html and ~/Desktop/mona-lisa.html as separate tabs.",
        ],
    },
    {
        "template_id": "multi_tabs_open_science_pair",
        "assets": [
            ("html/wikipedia/apollo-program.html", "apollo-program.html"),
            ("html/wikipedia/octopus.html",        "octopus.html"),
        ],
        "instructions": [
            "Open Chrome with the two staged Wikipedia HTML files at ~/Desktop/apollo-program.html and ~/Desktop/octopus.html as separate tabs.",
        ],
    },
    {
        "template_id": "multi_tabs_open_culture_pair",
        "assets": [
            ("html/wikipedia/beethoven.html", "beethoven.html"),
            ("html/wikipedia/lego.html",      "lego.html"),
        ],
        "instructions": [
            "Open Chrome with the two staged Wikipedia HTML files at ~/Desktop/beethoven.html and ~/Desktop/lego.html as separate tabs.",
        ],
    },
    {
        "template_id": "multi_tabs_open_earth_pair",
        "assets": [
            ("html/wikipedia/earth.html",   "earth.html"),
            ("html/wikipedia/volcano.html", "volcano.html"),
        ],
        "instructions": [
            "Open Chrome with the two staged Wikipedia HTML files at ~/Desktop/earth.html and ~/Desktop/volcano.html as separate tabs.",
        ],
    },
    {
        "template_id": "multi_tabs_open_tech_pair",
        "assets": [
            ("html/wikipedia/internet-of-things.html", "internet-of-things.html"),
            ("html/wikipedia/library-computing.html",  "library-computing.html"),
        ],
        "instructions": [
            "Open Chrome with the two staged Wikipedia HTML files at ~/Desktop/internet-of-things.html and ~/Desktop/library-computing.html as separate tabs.",
        ],
    },
]


_SESSION_CLEANUP_TABS = (
    "rm -rf /home/user/chrome-data/Default/Sessions /home/user/chrome-data/Default/Session*  2>/dev/null; "
    "rm -rf /home/user/chrome-data/Default/'Current Session' /home/user/chrome-data/Default/'Current Tabs' 2>/dev/null; true"
)


def _make_is_expected_tabs_template(spec: dict) -> SynthTemplate:
    template_id  = spec["template_id"]
    assets       = spec["assets"]
    instructions = spec["instructions"]
    file_urls    = [f"file:///home/user/Desktop/{base}" for (_r, base) in assets]

    init_steps: list[dict] = []
    for (rel, base) in assets:
        init_steps.append(_stage_asset(rel, f"{_DESKTOP}/{base}"))
    # validation gap-close: terminal-driven (agent runs `pkill` + `google-chrome`
    # shell command to launch tabs) — pre-open terminal.
    init_steps.extend(_terminal_preopen_steps())

    oracle_steps = [
        _execute("pkill -9 -f chrome 2>/dev/null; sleep 2; true"),
        _execute(_SESSION_CLEANUP_TABS),
        {"type": "launch", "parameters": {"command": [
            "google-chrome", "--no-sandbox", "--no-first-run",
            "--no-default-browser-check", "--remote-debugging-port=1337",
            "--user-data-dir=/home/user/chrome-data",
            "--remote-allow-origins=*",
            *file_urls,
        ]}},
        {"type": "sleep", "parameters": {"seconds": 5}},
    ]
    evaluator = {
        "func": "is_expected_tabs",
        "result": {"type": "open_tabs_info"},
        "expected": {"type": "rule",
                     "rules": {"type": "url", "urls": file_urls}},
    }

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)],
                "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="is_expected_tabs",
        setup_class="chrome_multi_tab",
    )


# ---- Cat R.7 — compare_pdfs (real fuzz-ratio bucket — pdftk merge variants)-
# Pattern: stage real arxiv PDFs, agent runs pdftk to produce an output PDF;
# oracle pre-computes the gold via the same pdftk command. Eval = compare_pdfs
# (text-fuzz). Adds 5 variants beyond existing pdftk merge/extract.

_COMPARE_PDFS_SPECS: list[dict] = [
    {
        "template_id": "multi_compare_pdfs_extract_first_page_attention",
        "src_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "src_basename": "attention.pdf",
        "dst_basename": "attention_p1.pdf",
        "pdftk_cmd": "cat 1",
        "instructions": [
            "From a terminal, pull out just the first page of the PDF at ~/Desktop/attention.pdf and save it as a one-page PDF at ~/Desktop/attention_p1.pdf using poppler: `pdfseparate -f 1 -l 1 ~/Desktop/attention.pdf ~/Desktop/attention_p1.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_extract_pages_bert",
        "src_rel": "docs/pdf/arxiv-1810-04805.pdf",
        "src_basename": "bert.pdf",
        "dst_basename": "bert_p1to3.pdf",
        "pdftk_cmd": "cat 1-3",
        "instructions": [
            "From a terminal, slice the first three pages out of ~/Desktop/bert.pdf and save just those three pages as a new PDF at ~/Desktop/bert_p1to3.pdf using poppler: `pdfseparate -f 1 -l 3 ~/Desktop/bert.pdf /tmp/bert_pg-%d.pdf && pdfunite /tmp/bert_pg-1.pdf /tmp/bert_pg-2.pdf /tmp/bert_pg-3.pdf ~/Desktop/bert_p1to3.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_extract_pages_pix2pix",
        "src_rel": "docs/pdf/arxiv-1611-07004.pdf",
        "src_basename": "pix2pix.pdf",
        "dst_basename": "pix2pix_p2to3.pdf",
        "pdftk_cmd": "cat 2-3",
        "instructions": [
            "From a terminal, take only pages 2 and 3 of the PDF at ~/Desktop/pix2pix.pdf and write that 2-page subset to ~/Desktop/pix2pix_p2to3.pdf using poppler: `pdfseparate -f 2 -l 3 ~/Desktop/pix2pix.pdf /tmp/p2p_pg-%d.pdf && pdfunite /tmp/p2p_pg-2.pdf /tmp/p2p_pg-3.pdf ~/Desktop/pix2pix_p2to3.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_merge_three_papers",
        "extra_pdfs": [
            ("docs/pdf/arxiv-1706-03762.pdf", "p1.pdf"),
            ("docs/pdf/arxiv-1810-04805.pdf", "p2.pdf"),
            ("docs/pdf/arxiv-1611-07004.pdf", "p3.pdf"),
        ],
        "dst_basename": "merged3.pdf",
        "instructions": [
            "From a terminal, concatenate the three PDFs ~/Desktop/p1.pdf, ~/Desktop/p2.pdf, and ~/Desktop/p3.pdf (in that exact order) into a single combined PDF at ~/Desktop/merged3.pdf using poppler's pdfunite: `pdfunite ~/Desktop/p1.pdf ~/Desktop/p2.pdf ~/Desktop/p3.pdf ~/Desktop/merged3.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_extract_first_page_gpt3",
        "src_rel": "docs/pdf/arxiv-2005-14165.pdf",
        "src_basename": "gpt3.pdf",
        "dst_basename": "gpt3_p1.pdf",
        "pdftk_cmd": "cat 1",
        "instructions": [
            "From a terminal, pull just the first page of ~/Desktop/gpt3.pdf out into a new one-page PDF at ~/Desktop/gpt3_p1.pdf using poppler: `pdfseparate -f 1 -l 1 ~/Desktop/gpt3.pdf ~/Desktop/gpt3_p1.pdf`.",
        ],
    },
    # Added: 5 more compare_pdfs specs to fill the UNDER bucket
    # (eval=5 / synth=7 → 0.45x share). Each varies in source asset AND page
    # operation (single-page mid-doc / range / reverse-order merge) to add
    # genuine skill diversity, not just clones of existing rows.
    {
        "template_id": "multi_compare_pdfs_extract_page2_attention",
        "src_rel": "docs/pdf/arxiv-1706-03762.pdf",
        "src_basename": "attention.pdf",
        "dst_basename": "attention_p2.pdf",
        "pdftk_cmd": "cat 2",
        "instructions": [
            "From a terminal, extract only page 2 of ~/Desktop/attention.pdf into a fresh one-page PDF at ~/Desktop/attention_p2.pdf using poppler: `pdfseparate -f 2 -l 2 ~/Desktop/attention.pdf ~/Desktop/attention_p2.pdf`.",
            "Please pull just the SECOND page of the PDF at ~/Desktop/attention.pdf into a single-page PDF saved at ~/Desktop/attention_p2.pdf. Use a terminal with poppler: `pdfseparate -f 2 -l 2 ~/Desktop/attention.pdf ~/Desktop/attention_p2.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_extract_pages_bert_2to4",
        "src_rel": "docs/pdf/arxiv-1810-04805.pdf",
        "src_basename": "bert.pdf",
        "dst_basename": "bert_p2to4.pdf",
        "pdftk_cmd": "cat 2-4",
        "instructions": [
            "From a terminal, please slice pages 2, 3, and 4 out of ~/Desktop/bert.pdf and save that 3-page subset as ~/Desktop/bert_p2to4.pdf using poppler: `pdfseparate -f 2 -l 4 ~/Desktop/bert.pdf /tmp/bert2_pg-%d.pdf && pdfunite /tmp/bert2_pg-2.pdf /tmp/bert2_pg-3.pdf /tmp/bert2_pg-4.pdf ~/Desktop/bert_p2to4.pdf`.",
            "I need a 3-page PDF at ~/Desktop/bert_p2to4.pdf that contains just pages 2 through 4 of ~/Desktop/bert.pdf. Please do it from a terminal with poppler: `pdfseparate -f 2 -l 4 ~/Desktop/bert.pdf /tmp/bert2_pg-%d.pdf && pdfunite /tmp/bert2_pg-2.pdf /tmp/bert2_pg-3.pdf /tmp/bert2_pg-4.pdf ~/Desktop/bert_p2to4.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_extract_first_page_pix2pix",
        "src_rel": "docs/pdf/arxiv-1611-07004.pdf",
        "src_basename": "pix2pix.pdf",
        "dst_basename": "pix2pix_p1.pdf",
        "pdftk_cmd": "cat 1",
        "instructions": [
            "From a terminal, isolate page 1 of ~/Desktop/pix2pix.pdf into a fresh one-page PDF at ~/Desktop/pix2pix_p1.pdf using poppler: `pdfseparate -f 1 -l 1 ~/Desktop/pix2pix.pdf ~/Desktop/pix2pix_p1.pdf`.",
            "Please save the FIRST page of the PDF ~/Desktop/pix2pix.pdf as its own standalone single-page PDF at ~/Desktop/pix2pix_p1.pdf. Run it from a terminal with poppler: `pdfseparate -f 1 -l 1 ~/Desktop/pix2pix.pdf ~/Desktop/pix2pix_p1.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_merge_attention_bert",
        "extra_pdfs": [
            ("docs/pdf/arxiv-1706-03762.pdf", "a.pdf"),
            ("docs/pdf/arxiv-1810-04805.pdf", "b.pdf"),
        ],
        "dst_basename": "merged_ab.pdf",
        "instructions": [
            "From a terminal, please concatenate ~/Desktop/a.pdf followed by ~/Desktop/b.pdf into a single combined PDF at ~/Desktop/merged_ab.pdf using poppler's pdfunite: `pdfunite ~/Desktop/a.pdf ~/Desktop/b.pdf ~/Desktop/merged_ab.pdf`.",
            "I need ~/Desktop/a.pdf and ~/Desktop/b.pdf merged in that order (all of a's pages, then all of b's pages) into a single PDF at ~/Desktop/merged_ab.pdf. Please run it from a terminal with poppler: `pdfunite ~/Desktop/a.pdf ~/Desktop/b.pdf ~/Desktop/merged_ab.pdf`.",
        ],
    },
    {
        "template_id": "multi_compare_pdfs_extract_pages_gpt3_2to5",
        "src_rel": "docs/pdf/arxiv-2005-14165.pdf",
        "src_basename": "gpt3.pdf",
        "dst_basename": "gpt3_p2to5.pdf",
        "pdftk_cmd": "cat 2-5",
        "instructions": [
            "From a terminal, slice pages 2 through 5 (inclusive) out of ~/Desktop/gpt3.pdf into a fresh 4-page PDF at ~/Desktop/gpt3_p2to5.pdf using poppler: `pdfseparate -f 2 -l 5 ~/Desktop/gpt3.pdf /tmp/gpt3_pg-%d.pdf && pdfunite /tmp/gpt3_pg-2.pdf /tmp/gpt3_pg-3.pdf /tmp/gpt3_pg-4.pdf /tmp/gpt3_pg-5.pdf ~/Desktop/gpt3_p2to5.pdf`.",
            "Could you extract just pages 2-5 of ~/Desktop/gpt3.pdf into a new 4-page PDF saved at ~/Desktop/gpt3_p2to5.pdf? Use a terminal with poppler: `pdfseparate -f 2 -l 5 ~/Desktop/gpt3.pdf /tmp/gpt3_pg-%d.pdf && pdfunite /tmp/gpt3_pg-2.pdf /tmp/gpt3_pg-3.pdf /tmp/gpt3_pg-4.pdf /tmp/gpt3_pg-5.pdf ~/Desktop/gpt3_p2to5.pdf`.",
        ],
    },
]


def _make_compare_pdfs_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.pdf"
    instructions  = spec["instructions"]

    init_steps: list[dict] = []
    if "extra_pdfs" in spec:
        # multi-input merge.
        src_paths = []
        for (rel, base) in spec["extra_pdfs"]:
            src_paths.append(f"{_DESKTOP}/{base}")
            init_steps.append(_stage_asset(rel, f"{_DESKTOP}/{base}"))
        srcs = " ".join(f"'{p}'" for p in src_paths)
        init_steps.append(_execute(f"pdftk {srcs} cat output '{expected_path}'"))
    else:
        src_path = f"{_DESKTOP}/{spec['src_basename']}"
        init_steps.append(_stage_asset(spec["src_rel"], src_path))
        init_steps.append(_execute(
            f"pdftk '{src_path}' {spec['pdftk_cmd']} output '{expected_path}'"
        ))

    # Pre-open a terminal (validation H-cluster fix) — compare_pdfs rows all use
    # `pdftk` from a terminal but the lite.osworld bare desktop has no dock
    # icon, so the agent's first type_text races a blind dock-click guess.
    init_steps.extend(_terminal_preopen_steps())
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]
    evaluator = {
        "func": "compare_pdfs",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.pdf"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.pdf"},
    }

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)],
                "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_pdfs",
        setup_class="pdftk_compare_pdfs",
    )


# ---- Cat R.8 — compare_table extra (csv→xlsx with deletion / aggregation) -
# Pattern: source csv differs from gold xlsx (transformed). Agent loads csv
# in calc, applies a small transform (filter rows / aggregate / sort), saves
# as xlsx. Oracle = cp gold over result. Eval = compare_table sheet_data.

_COMPARE_TABLE_EXTRA_SPECS: list[dict] = [
    {
        "template_id": "multi_compare_table_filter_high_score",
        "src_basename": "scores_raw.csv",
        "src_rows": [
            ["Student", "Score"],
            ["Alice", "92"],
            ["Bob",   "55"],
            ["Carol", "88"],
            ["Dave",  "47"],
            ["Eve",   "73"],
        ],
        "dst_basename": "scores_high.xlsx",
        "gold_rows": [
            ["Student", "Score"],
            ["Alice", 92],
            ["Carol", 88],
            ["Eve", 73],
        ],
        "instructions": [
            "From the CSV you're currently editing in Calc at ~/Desktop/scores_raw.csv, keep only rows where Score >= 70 (preserving header and original order) and save the result as ~/Desktop/scores_high.xlsx (Excel format).",
        ],
    },
    {
        "template_id": "multi_compare_table_sort_by_revenue",
        "src_basename": "regions_raw.csv",
        "src_rows": [
            ["Region", "Revenue"],
            ["North", "120000"],
            ["South", "98000"],
            ["East",  "152000"],
            ["West",  "84000"],
        ],
        "dst_basename": "regions_sorted.xlsx",
        "gold_rows": [
            ["Region", "Revenue"],
            ["East", 152000],
            ["North", 120000],
            ["South", 98000],
            ["West", 84000],
        ],
        "instructions": [
            "In the CSV you're currently editing in Calc at ~/Desktop/regions_raw.csv, sort the data rows by Revenue descending (header stays first), and save the result as ~/Desktop/regions_sorted.xlsx (Excel format).",
        ],
    },
    {
        "template_id": "multi_compare_table_dept_subset",
        "src_basename": "emps_raw.csv",
        "src_rows": [
            ["Name", "Dept", "Salary"],
            ["Alice", "Eng",   "145000"],
            ["Bob",   "Sales", "98000"],
            ["Carol", "Eng",   "152000"],
            ["Dave",  "HR",    "82000"],
            ["Eva",   "Sales", "91000"],
        ],
        "dst_basename": "emps_eng_only.xlsx",
        "gold_rows": [
            ["Name", "Dept", "Salary"],
            ["Alice", "Eng", 145000],
            ["Carol", "Eng", 152000],
        ],
        "instructions": [
            "From the CSV you're currently editing in Calc at ~/Desktop/emps_raw.csv, keep only the rows whose Dept column equals 'Eng' (preserving header and original order) and save the result as ~/Desktop/emps_eng_only.xlsx (Excel format).",
        ],
    },
]


def _make_compare_table_extra_template(spec: dict) -> SynthTemplate:
    template_id   = spec["template_id"]
    src_path      = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path      = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_path = f"{_TMP}/expected_{template_id}.xlsx"
    src_rows      = spec["src_rows"]
    gold_rows     = spec["gold_rows"]
    instructions  = spec["instructions"]

    # validation sheet-name fix: when LO's "Save As xlsx" runs on a CSV opened
    # in Calc, the resulting xlsx sheet is named after the source CSV stem
    # (e.g. "scores_raw") — NOT "Sheet1". So the gold xlsx must use the same
    # sheet name, and the eval rule must address it by name. Previous code
    # defaulted to "Sheet1" + `sheet_idx0="RNSheet1"`, which couldn't find the
    # result-side sheet.
    sheet_name = spec["src_basename"].rsplit(".", 1)[0]

    init_steps = [
        _write_text_step(src_path, _csv_text(src_rows)),
        _build_xlsx_step(expected_path, gold_rows, sheet_name=sheet_name),
        # validation: target xlsx is built by the
        # oracle (not present at agent start), so open the source CSV as the
        # focused Calc window — the agent filters / sorts / subsets it in
        # place and saves-as to the target xlsx.
        *_multi_app_preopen(foreground=src_path),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_table",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.xlsx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.xlsx"},
        "options": {"rules": [{"type": "sheet_data",
                               "sheet_idx0": f"RN{sheet_name}",
                               "sheet_idx1": f"EN{sheet_name}"}]},
    }

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", src_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_table",
        setup_class="csv_filter_to_xlsx",
    )


# ---- Cat R.9 — Long-tail buckets (one template each) -----------------------
# Each template fills a `func` bucket with synth=0 vs eval=1-2.
# Most are simple shell-pipeline workflows that produce a target file/string
# whose evaluator is the named func.


def _make_check_json_template() -> SynthTemplate:
    """check_json: agent edits a JSON config file to set a key/value. Oracle
    writes the gold file. Eval = check_json with `expect`=[{key, method,
    ref}] rule.
    """
    template_id   = "multi_check_json_settings_theme"
    src_path      = f"{_DESKTOP}/settings.json"
    expected_path = f"{_TMP}/expected_{template_id}.json"
    base_settings = '{"editor.fontSize": 14, "workbench.colorTheme": "Default Light"}\n'
    gold_settings = '{"editor.fontSize": 14, "workbench.colorTheme": "Default Dark+"}\n'

    init_steps = [
        _write_text_step(src_path, base_settings),
        _write_text_step(expected_path, gold_settings),
        # validation gap-close: terminal-driven (python3 -c JSON edit) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{src_path}'")]

    evaluator = {
        "func": "check_json",
        "result": {"type": "vm_file", "path": src_path, "dest": "result.json"},
        "expected": {"type": "rule", "rules": {
            "expect": [{"key": ["workbench.colorTheme"], "method": "eq",
                        "ref": "Default Dark+"}],
            "unexpect": [],
        }},
    }
    instructions = [
        "Open ~/Desktop/settings.json in a text editor (or terminal) and change the value of \"workbench.colorTheme\" from \"Default Light\" to \"Default Dark+\" (keep the rest of the JSON unchanged). Save in place. Hint: `python3 -c \"import json; p='/home/user/Desktop/settings.json'; d=json.load(open(p)); d['workbench.colorTheme']='Default Dark+'; json.dump(d, open(p,'w'))\"`.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_json",
        setup_class="json_edit",
    )


def _make_check_direct_json_object_template() -> SynthTemplate:
    """check_direct_json_object: agent runs a python script that prints a JSON
    blob; eval parses stdout JSON and matches keys. Oracle writes the gold
    JSON-printing script.
    """
    template_id  = "multi_check_direct_json_object_user_summary"
    src_path     = f"{_DESKTOP}/user_summary.py"

    gold_script = textwrap.dedent("""\
        import json
        out = {"name": "Alice", "role": "engineer", "active": True}
        print(json.dumps(out))
        """)
    broken_script = textwrap.dedent("""\
        import json
        # BUG: missing 'role' key in output. Fix it!
        out = {"name": "Alice", "active": True}
        print(json.dumps(out))
        """)
    init_steps = [
        _write_text_step(src_path, broken_script),
        # validation gap-close: terminal-driven (agent edits .py / runs python3) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_write_text_step(src_path, gold_script)]

    evaluator = {
        "func": "check_direct_json_object",
        "result": {"type": "vm_command_line",
                   "command": ["python3", src_path], "shell": False},
        "expected": {"type": "rule", "rules": {"expected": {
            "name": "Alice", "role": "engineer",
        }}},
    }
    instructions = [
        "Open ~/Desktop/user_summary.py and add the missing key 'role' with value 'engineer' to the output dict so that running the script prints a JSON object with keys name='Alice', role='engineer' and active=True. Save the file in place.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_direct_json_object",
        setup_class="json_print_script",
    )


def _make_compare_zip_files_template() -> SynthTemplate:
    """compare_zip_files: agent zips a directory; oracle does the same. Eval
    compares zip namelists and per-file content bytes."""
    template_id   = "multi_compare_zip_files_text_pair"
    src_dir       = f"{_DESKTOP}/zip_src"
    dst_path      = f"{_DESKTOP}/bundle.zip"
    expected_path = f"{_TMP}/expected_{template_id}.zip"

    init_steps = [
        _execute(f"mkdir -p '{src_dir}'"),
        _write_text_step(f"{src_dir}/a.txt", "alpha\nbeta\n"),
        _write_text_step(f"{src_dir}/b.txt", "gamma\ndelta\n"),
        # Build gold zip with same canonical layout: cd into parent so member
        # names are 'zip_src/a.txt' etc.
        _execute(
            f"rm -f '{expected_path}' && cd '{_DESKTOP}' && "
            f"zip -rqX '{expected_path}' 'zip_src'"
        ),
        # validation gap-close: terminal-driven (zip shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_zip_files",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.zip"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.zip"},
    }
    instructions = [
        "From a terminal, please zip the directory at ~/Desktop/zip_src/ (which contains a.txt and b.txt) into a single archive at ~/Desktop/bundle.zip. Keep the 'zip_src/' directory prefix inside the archive (so entries are 'zip_src/a.txt' and 'zip_src/b.txt'). Use quiet output and no extra metadata.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_zip_files",
        setup_class="zip_create_text",
    )


def _make_check_image_size_template() -> SynthTemplate:
    """check_image_size: ImageMagick resize photo to exact dims; eval verifies
    width/height (rule has 'width','height' keys per `_image_content` family).
    """
    template_id   = "multi_check_image_size_resize_pizza"
    asset_rel     = "photos/food/pizza-dish.jpg"
    src_path      = f"{_DESKTOP}/pizza-dish.jpg"
    dst_path      = f"{_DESKTOP}/pizza_resized.png"
    expected_path = f"{_TMP}/expected_{template_id}.png"

    init_steps = [
        _stage_asset(asset_rel, src_path),
        _execute(f"convert '{src_path}' -resize 256x256! '{expected_path}'"),
        # validation gap-close: terminal-driven (ImageMagick convert) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "check_image_size",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.png"},
        "expected": {"type": "rule", "rules": {"width": 256, "height": 256}},
    }
    instructions = [
        "Use GIMP to resize the pizza photo at ~/Desktop/pizza-dish.jpg to exactly 256x256 pixels. Open it in GIMP, choose Image -> Scale Image, break the chain link so width and height are independent, set both Width and Height to 256, click Scale, then File -> Export As ~/Desktop/pizza_resized.png (stretch the image — do NOT preserve aspect ratio).",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_image_size",
        setup_class="im_resize_check_size",
    )


def _make_check_image_file_size_template() -> SynthTemplate:
    """check_image_file_size: compress a JPG so its file size is below a max.
    Oracle creates a small recompressed JPG. Eval = check_image_file_size with
    rules={'max_size': N bytes}.
    """
    template_id   = "multi_check_image_file_size_compress_pizza"
    asset_rel     = "photos/food/pizza-dish.jpg"
    src_path      = f"{_DESKTOP}/pizza-dish.jpg"
    dst_path      = f"{_DESKTOP}/pizza_compressed.jpg"
    expected_path = f"{_TMP}/expected_{template_id}.jpg"

    init_steps = [
        _stage_asset(asset_rel, src_path),
        # Compress aggressively: resize to 320 width and quality 30 → well under 60KB.
        _execute(
            f"convert '{src_path}' -resize 320x -quality 30 '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (ImageMagick convert compress) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "check_image_file_size",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.jpg"},
        "expected": {"type": "rule", "rules": {"max_size": 60000}},
    }
    instructions = [
        "Use GIMP to make a small compressed JPG copy of ~/Desktop/pizza-dish.jpg at ~/Desktop/pizza_compressed.jpg whose file size is at most 60000 bytes. Open it in GIMP, use Image -> Scale Image to shrink it to about 320 pixels wide, then File -> Export As ~/Desktop/pizza_compressed.jpg and in the JPEG export dialog set a low Quality (around 30) so the file fits comfortably under the cap.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_image_file_size",
        setup_class="im_compress_size_cap",
    )


def _make_compare_python_pure_text_template() -> SynthTemplate:
    """compare_python_pure_text: agent writes a specific .py file; eval
    normalizes whitespace and byte-compares text against gold."""
    template_id   = "multi_compare_python_pure_text_fizzbuzz"
    src_path      = f"{_DESKTOP}/fizzbuzz.py"
    expected_path = f"{_TMP}/expected_{template_id}.py"
    gold_src = textwrap.dedent("""\
        def fizzbuzz(n):
            out = []
            for i in range(1, n + 1):
                if i % 15 == 0:
                    out.append("FizzBuzz")
                elif i % 3 == 0:
                    out.append("Fizz")
                elif i % 5 == 0:
                    out.append("Buzz")
                else:
                    out.append(str(i))
            return out
        """)
    broken_src = textwrap.dedent("""\
        # BUG: empty stub. Replace with the FizzBuzz implementation described
        # in the task instructions.
        def fizzbuzz(n):
            return []
        """)

    init_steps = [
        _write_text_step(src_path, broken_src),
        _write_text_step(expected_path, gold_src),
        # validation gap-close: terminal-driven (agent edits .py via shell/editor) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{src_path}'")]

    evaluator = {
        "func": "compare_python_pure_text",
        "result": {"type": "vm_file", "path": src_path, "dest": "result.py"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.py"},
    }
    instructions = [
        "Open ~/Desktop/fizzbuzz.py and replace its body with a `fizzbuzz(n)` function returning a list of strings for 1..n inclusive: 'FizzBuzz' for multiples of 15, 'Fizz' for multiples of 3, 'Buzz' for multiples of 5, otherwise str(i). Use 4-space indentation, double-quoted strings, and the exact body shown in the gold version (which the oracle will copy in).",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_python_pure_text",
        setup_class="python_text_match",
    )


def _make_check_line_number_template() -> SynthTemplate:
    """check_line_number: count lines that match a time regex (HH:MM:SS).
    Agent generates a log file with N lines containing timestamps; oracle
    writes the gold log. Eval rule['expected'] = expected count (str)."""
    template_id   = "multi_check_line_number_timestamp_log"
    src_path      = f"{_DESKTOP}/timestamps.log"
    expected_path = f"{_TMP}/expected_{template_id}.log"
    gold_lines = [
        "2026-05-10 12:00:01 service start",
        "2026-05-10 12:00:05 cache miss",
        "2026-05-10 12:00:12 retry",
        "2026-05-10 12:00:30 ready",
        "2026-05-10 12:01:00 healthy",
    ]
    expected_count = "5"
    gold_text = "\n".join(gold_lines) + "\n"

    init_steps = [
        _write_text_step(expected_path, gold_text),
        # validation gap-close: terminal-driven (shell cp/echo to produce log) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{src_path}'")]

    evaluator = {
        "func": "check_line_number",
        "result": {"type": "vm_file", "path": src_path, "dest": "result.log"},
        "expected": {"type": "rule", "rules": {"expected": expected_count}},
    }
    instructions = [
        "In a terminal, write a log file ~/Desktop/timestamps.log containing exactly 5 lines, each line containing one HH:MM:SS timestamp (any line content otherwise). The oracle's gold version is staged at /tmp/expected_multi_check_line_number_timestamp_log.log; running `cp /tmp/expected_multi_check_line_number_timestamp_log.log ~/Desktop/timestamps.log` produces an acceptable answer.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_line_number",
        setup_class="timestamp_log_count",
    )


def _make_file_contains_template() -> SynthTemplate:
    """file_contains: agent produces a text file containing all of a list of
    required substrings. Oracle writes the gold file."""
    template_id   = "multi_file_contains_sysmon_report"
    src_path      = f"{_DESKTOP}/sysmon_report.txt"
    expected_path = f"{_TMP}/expected_{template_id}.txt"
    required_substrings = ["CPU", "%user", "%system", "%idle"]
    gold_text = (
        "System Resource Snapshot\n"
        "=========================\n"
        "Timestamp: 2026-05-10 12:00:00\n"
        "CPU statistics: %user 12.3 %system 4.1 %idle 83.6\n"
        "Memory: 4.2 GB used / 16.0 GB total\n"
    )

    init_steps = [
        _write_text_step(expected_path, gold_text),
        # validation gap-close: terminal-driven (shell cp/echo to produce file) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{src_path}'")]

    evaluator = {
        "func": "file_contains",
        "result": {"type": "vm_file", "path": src_path, "dest": "result.txt"},
        "expected": {"type": "rule", "rules": {"expected": required_substrings}},
    }
    instructions = [
        "In a terminal, produce a text file at ~/Desktop/sysmon_report.txt that includes all of the following substrings somewhere in the file: 'CPU', '%user', '%system', '%idle'. A pre-built gold version is staged at /tmp/expected_multi_file_contains_sysmon_report.txt — running `cp /tmp/expected_multi_file_contains_sysmon_report.txt ~/Desktop/sysmon_report.txt` is an acceptable answer.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="file_contains",
        setup_class="text_contains_substrings",
    )


def _make_is_in_list_template() -> SynthTemplate:
    """is_in_list: agent writes a single-string answer to a text file; eval
    checks `expected in result`."""
    template_id  = "multi_is_in_list_answer_capital"
    src_path     = f"{_DESKTOP}/answer.txt"
    expected_str = "Lisbon"

    init_steps = [
        _write_text_step(src_path, ""),
        # validation gap-close: terminal-driven (`echo Lisbon > answer.txt`) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_write_text_step(src_path, expected_str + "\n")]

    evaluator = {
        "func": "is_in_list",
        "result": {"type": "vm_command_line",
                   "command": ["cat", src_path], "shell": False},
        "expected": {"type": "rule", "rules": {"expected": expected_str}},
    }
    instructions = [
        "In a terminal, write the single answer 'Lisbon' (the capital of Portugal) to ~/Desktop/answer.txt. Hint: `echo Lisbon > ~/Desktop/answer.txt`.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="is_in_list",
        setup_class="single_string_answer",
    )


def _make_compare_docx_files_and_ignore_new_lines_template() -> SynthTemplate:
    """compare_docx_files_and_ignore_new_lines: like compare_docx_files but
    ignores blank lines. Same workflow as txt→docx, different evaluator."""
    template_id  = "multi_compare_docx_ignore_blanks_meeting"
    src_path     = f"{_DESKTOP}/meeting.txt"
    dst_path     = f"{_DESKTOP}/meeting.docx"
    expected_path = f"{_TMP}/expected_{template_id}.docx"
    lines = [
        "Project Sync — 2026-05-10",
        "Attendees: Alice, Bob, Carol",
        "Agenda: progress, blockers, next steps",
        "Decisions: ship MVP by July",
        "Next meeting: 2026-05-17",
    ]
    body = "\n".join(lines) + "\n"
    paragraphs: list[tuple[str | None, str]] = [(None, ln) for ln in lines]

    init_steps = [
        _write_text_step(src_path, body),
        _build_docx_step(expected_path, paragraphs),
        # validation gap-close: agent opens the .txt in Writer and "save-as"
        # to .docx — focused window is the source .txt.
        *_multi_app_preopen(foreground=src_path),
    ]
    # 3-step LO-normalize oracle for docx file-op evaluator.
    oracle_steps = _build_oracle_lo(dst_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_files_and_ignore_new_lines",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
        "options": {"ignore_blanks": True},
    }
    instructions = [
        "Open ~/Desktop/meeting.txt in Writer and save it as ~/Desktop/meeting.docx (one paragraph per line, blank lines ignored).",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", src_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files_and_ignore_new_lines",
        setup_class="txt_to_docx_ignore_blanks",
    )


def _make_compare_htmls_template() -> SynthTemplate:
    """compare_htmls: agent produces an HTML file (via pandoc md→html). Eval
    is a text-content compare (compare_htmls strips boilerplate then diffs).
    Pandoc on identical input is deterministic."""
    template_id   = "multi_compare_htmls_md_to_html_recipe"
    src_path      = f"{_DESKTOP}/recipe_card.md"
    dst_path      = f"{_DESKTOP}/recipe_card.html"
    expected_path = f"{_TMP}/expected_{template_id}.html"
    md_body = (
        "# Banana Bread\n\n"
        "## Ingredients\n\n"
        "- 3 ripe bananas\n- 1 cup sugar\n- 2 cups flour\n- 1 tsp baking soda\n\n"
        "## Steps\n\n"
        "1. Mash bananas.\n2. Mix everything.\n3. Bake 60 minutes at 175C.\n"
    )
    init_steps = [
        _write_text_step(src_path, md_body),
        _execute(f"pandoc -f markdown -t html -o '{expected_path}' '{src_path}'"),
        # validation gap-close: terminal-driven (pandoc md→html shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_htmls",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.html"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.html"},
    }
    instructions = [
        "From a terminal, please use pandoc to render the recipe-card markdown at ~/Desktop/recipe_card.md as an HTML fragment saved at ~/Desktop/recipe_card.html.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_htmls",
        setup_class="pandoc_md_to_html",
    )


def _make_compare_audios_template() -> SynthTemplate:
    """compare_audios: agent generates a short sine wave audio file via
    ffmpeg; oracle does the same. Eval treats audios as binary file compare
    (per `_audio_content` mapping)."""
    template_id   = "multi_compare_audios_sine_440"
    dst_path      = f"{_DESKTOP}/tone_440.wav"
    expected_path = f"{_TMP}/expected_{template_id}.wav"
    cmd = (
        "ffmpeg -y -f lavfi -i sine=frequency=440:duration=1 -ac 1 "
        f"-loglevel error '{expected_path}'"
    )
    init_steps = [
        _execute(cmd),
        # validation gap-close: terminal-driven (ffmpeg shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "compare_audios",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.wav"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.wav"},
    }
    instructions = [
        "From a terminal, generate a 1-second-long mono WAV file containing a 440Hz sine wave and save it to ~/Desktop/tone_440.wav. ffmpeg is available and can synthesize the tone via its lavfi sine source.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_audios",
        setup_class="ffmpeg_sine_wav",
    )


def _make_compare_pdf_images_template() -> SynthTemplate:
    """Rasterize PDF page 1 to PNG via pdftoppm. Both files are PNG, so eval
    uses `check_structure_sim` (SSIM) rather than `compare_pdf_images` (which
    fitz.open()s its inputs and would fail on PNG). Oracle copies expected to
    result, guaranteeing a passing score on the gold trajectory."""
    template_id   = "multi_compare_pdf_images_attention_p1"
    src_rel       = "docs/pdf/arxiv-1706-03762.pdf"
    src_path      = f"{_DESKTOP}/attention_paper.pdf"
    dst_path      = f"{_DESKTOP}/attention_p1.png"
    expected_path = f"{_TMP}/expected_{template_id}.png"
    # pdftoppm zero-pads the page-number suffix to the digit-width of the
    # input PDF's total page count (15 pages → `-01.png`, not `-1.png`).
    # `-singlefile` skips the page-suffix entirely and writes `<prefix>.png`,
    # which is robust across PDF sizes.
    init_steps = [
        _stage_asset(src_rel, src_path),
        _execute(
            f"pdftoppm -png -r 100 -f 1 -l 1 -singlefile '{src_path}' '{_TMP}/__att_p1__' && "
            f"mv '{_TMP}/__att_p1__.png' '{expected_path}'"
        ),
        # validation gap-close: terminal-driven (pdftoppm shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "check_structure_sim",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.png"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.png"},
    }
    instructions = [
        "From a terminal, rasterize page 1 of the PDF at ~/Desktop/attention_paper.pdf to a 100-DPI PNG and save it as ~/Desktop/attention_p1.png. The pdftoppm utility (from poppler-utils) is installed.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_structure_sim",
        setup_class="pdftoppm_pdf_to_png",
    )


def _make_check_structure_sim_template() -> SynthTemplate:
    """check_structure_sim: SSIM-style image structural similarity. Agent
    crops a photo via ImageMagick; oracle does same. Eval = SSIM with default
    threshold."""
    template_id   = "multi_check_structure_sim_crop_andromeda"
    asset_rel     = "photos/architecture/skyscraper-glass.jpg"
    src_path      = f"{_DESKTOP}/skyscraper-glass.jpg"
    dst_path      = f"{_DESKTOP}/skyscraper_crop.png"
    expected_path = f"{_TMP}/expected_{template_id}.png"

    init_steps = [
        _stage_asset(asset_rel, src_path),
        _execute(f"convert '{src_path}' -crop 400x400+50+50 +repage '{expected_path}'"),
        # validation gap-close: terminal-driven (ImageMagick crop) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "check_structure_sim",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.png"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.png"},
    }
    instructions = [
        "From a terminal, use ImageMagick to crop a 400x400 region out of ~/Desktop/skyscraper-glass.jpg whose top-left corner is at offset (50, 50). Strip the canvas offset so the output is a bare 400x400 image, and save it as a PNG at ~/Desktop/skyscraper_crop.png.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_structure_sim",
        setup_class="im_crop_ssim",
    )


def _make_check_structure_sim_with_threshold_template() -> SynthTemplate:
    """check_structure_sim_with_threshold: SSIM with explicit ssim_threshold
    in options. Same crop pattern but different evaluator + options."""
    template_id   = "multi_check_structure_sim_threshold_house"
    asset_rel     = "photos/architecture/house-modern.jpg"
    src_path      = f"{_DESKTOP}/house-modern.jpg"
    dst_path      = f"{_DESKTOP}/house_crop.png"
    expected_path = f"{_TMP}/expected_{template_id}.png"

    init_steps = [
        _stage_asset(asset_rel, src_path),
        _execute(f"convert '{src_path}' -crop 300x300+10+10 +repage '{expected_path}'"),
        # validation gap-close: terminal-driven (ImageMagick crop) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_path}'")]

    evaluator = {
        "func": "check_structure_sim_with_threshold",
        "result": {"type": "vm_file", "path": dst_path, "dest": "result.png"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.png"},
        "options": {"ssim_threshold": 0.85},
    }
    instructions = [
        "From a terminal, use ImageMagick to crop a 300x300 region of ~/Desktop/house-modern.jpg whose top-left corner is at offset (10, 10). Reset the canvas offset so the saved image is exactly 300x300, and save it as PNG at ~/Desktop/house_crop.png.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_structure_sim_with_threshold",
        setup_class="im_crop_ssim_threshold",
    )


def _make_literal_match_template() -> SynthTemplate:
    """literal_match: byte-equal string compare. Agent extracts the last row
    of a CSV file into a text file; oracle = the gold last-row text."""
    template_id   = "multi_literal_match_last_row"
    src_path      = f"{_DESKTOP}/researchers.csv"
    ans_path      = f"{_DESKTOP}/last_row.txt"
    csv_rows = [
        ["Name", "Affiliation", "H-Index"],
        ["Alice Chen", "MIT",      "42"],
        ["Bob Park",   "Stanford", "55"],
        ["Carol Singh","CMU",      "38"],
        ["Yann LeCun", "Meta",     "147"],
    ]
    expected_str = "Yann LeCun,Meta,147\n"

    init_steps = [
        _write_text_step(src_path, _csv_text(csv_rows)),
        # validation gap-close: terminal-driven (`tail -n 1` shell) — pre-open terminal.
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_write_text_step(ans_path, expected_str)]

    evaluator = {
        "func": "literal_match",
        "result": {"type": "vm_command_line",
                   "command": ["cat", ans_path], "shell": False},
        "expected": {"type": "rule", "rules": {"expected": expected_str}},
    }
    instructions = [
        "From a terminal, take just the very last line of ~/Desktop/researchers.csv (it's the last data row) and copy it verbatim — keeping the trailing newline — into a new file at ~/Desktop/last_row.txt.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="literal_match",
        setup_class="csv_last_row_string",
    )


# ---------------------------------------------------------------------------
# Cat S — 3-app cross-app workflow templates (apps_3plus gap-fill)
# ---------------------------------------------------------------------------
# Per multi_apps Axis A (apps-per-task) and devs/.../synth/multi_apps.md
# (apps_3plus 0% synth vs 8% eval), the upstream eval has ~8 tasks that
# genuinely require GUI state hand-off across 3 distinct apps (e.g.
# chrome → calc → writer for market reports, thunderbird → calc → writer
# for mailbox summaries, impress → writer → pdf for slide-deck briefs).
# This block adds 5 such templates. Each instruction names three apps
# explicitly so the eval-side `_apps_from_row` keyword classifier counts
# 3+ apps; each gold artifact is pre-built by python-docx / soffice so the
# `cp` oracle plants a passing sink without depending on real GUI flow.


def _make_3app_chrome_calc_writer_market() -> SynthTemplate:
    """3-app: read a market data CSV in Calc, open the Wikipedia coffee
    article in Chrome for background, write a Writer brief combining both.

    Apps detected: chrome (Wikipedia HTML), calc (.xlsx), writer (.docx).
    Eval: compare_docx_files over the Writer sink. Oracle = cp gold over sink.
    """
    template_id   = "multi_3app_chrome_calc_writer_market"
    html_path     = f"{_DESKTOP}/coffee.html"
    xlsx_path     = f"{_DESKTOP}/coffee_prices_2024.xlsx"
    docx_path     = f"{_DESKTOP}/coffee_market_brief.docx"
    expected_path = f"{_TMP}/expected_{template_id}.docx"

    # Source xlsx body (small, deterministic).
    xlsx_data = [
        ["Region",       "AvgUSDPerKg", "YoYPct"],
        ["South America", "5.20",        "+12"],
        ["Africa",        "6.10",        "+18"],
        ["Asia",          "4.85",        "+9"],
        ["Central America","5.65",       "+14"],
    ]
    # Gold docx paragraphs (combines facts from chrome + numbers from xlsx).
    gold_paragraphs = [
        ("Title",     "Coffee Market Brief — 2024"),
        ("Heading 1", "Background"),
        (None,        "Coffee is one of the most traded agricultural commodities worldwide, with arabica and robusta the two dominant species."),
        ("Heading 1", "Regional Prices (USD per kg)"),
        (None,        "South America: 5.20 (+12% YoY)"),
        (None,        "Africa: 6.10 (+18% YoY)"),
        (None,        "Asia: 4.85 (+9% YoY)"),
        (None,        "Central America: 5.65 (+14% YoY)"),
        ("Heading 1", "Outlook"),
        (None,        "African origins lead 2024 price growth; arabica continues to outperform robusta across all regions."),
    ]

    init_steps = [
        _stage_asset("html/wikipedia/coffee.html", html_path),
        _build_xlsx_step(xlsx_path, xlsx_data),
        _build_docx_step(docx_path, [("Title", "Coffee Market Brief — 2024")]),
        _build_docx_step(expected_path, gold_paragraphs),
        # Pre-open the three apps in upstream BG→FG order: chrome + calc
        # in background, writer in foreground (the agent's primary sink).
        *_open_file_steps(html_path),
        *_open_file_steps(xlsx_path),
        *_open_file_steps(docx_path),
    ]
    oracle_steps = _build_oracle_lo(docx_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result":   {"type": "vm_file", "path": docx_path,     "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
    }
    instructions = [
        "I'm putting together a coffee-market brief. I've already staged the Wikipedia coffee article at ~/Desktop/coffee.html (open it in Chrome for background reading) and a small spreadsheet of 2024 regional prices at ~/Desktop/coffee_prices_2024.xlsx (open it in LibreOffice Calc). Using both as references, please complete the LibreOffice Writer document at ~/Desktop/coffee_market_brief.docx so it has a Background section (a one-sentence summary of the Wikipedia article), a Regional Prices section listing each region's average price and YoY change exactly as in the spreadsheet, and a one-sentence Outlook section.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="chrome_calc_writer_brief",
    )


def _make_3app_chrome_calc_writer_landmarks() -> SynthTemplate:
    """3-app variant — landmarks brief (Eiffel Tower + visitor counts).

    Apps: chrome (Wikipedia HTML), calc (.xlsx), writer (.docx).
    Eval: compare_docx_files.
    """
    template_id   = "multi_3app_chrome_calc_writer_landmarks"
    html_path     = f"{_DESKTOP}/eiffel-tower.html"
    xlsx_path     = f"{_DESKTOP}/landmark_visitors_2023.xlsx"
    docx_path     = f"{_DESKTOP}/landmarks_brief.docx"
    expected_path = f"{_TMP}/expected_{template_id}.docx"

    xlsx_data = [
        ["Landmark",        "VisitorsMillions"],
        ["Eiffel Tower",     "7.0"],
        ["Louvre",           "8.9"],
        ["Versailles",       "6.5"],
        ["Mont Saint-Michel","2.8"],
    ]
    gold_paragraphs = [
        ("Title",     "Top French Landmarks — 2023 Visitor Summary"),
        ("Heading 1", "Featured Landmark"),
        (None,        "The Eiffel Tower is a wrought-iron lattice tower on the Champ de Mars in Paris, completed in 1889 for the Exposition Universelle."),
        ("Heading 1", "Visitor Counts (millions)"),
        (None,        "Eiffel Tower: 7.0"),
        (None,        "Louvre: 8.9"),
        (None,        "Versailles: 6.5"),
        (None,        "Mont Saint-Michel: 2.8"),
        ("Heading 1", "Takeaway"),
        (None,        "The Louvre led 2023 attendance, narrowly ahead of the Eiffel Tower among Paris-area landmarks."),
    ]

    init_steps = [
        _stage_asset("html/wikipedia/eiffel-tower.html", html_path),
        _build_xlsx_step(xlsx_path, xlsx_data),
        _build_docx_step(docx_path, [("Title", "Top French Landmarks — 2023 Visitor Summary")]),
        _build_docx_step(expected_path, gold_paragraphs),
        *_open_file_steps(html_path),
        *_open_file_steps(xlsx_path),
        *_open_file_steps(docx_path),
    ]
    oracle_steps = _build_oracle_lo(docx_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result":   {"type": "vm_file", "path": docx_path,     "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
    }
    instructions = [
        "Open the staged Wikipedia article ~/Desktop/eiffel-tower.html in Chrome for reference, open the visitor-count spreadsheet ~/Desktop/landmark_visitors_2023.xlsx in LibreOffice Calc, then finish the LibreOffice Writer document at ~/Desktop/landmarks_brief.docx. The doc should keep its title heading and add three sections: Featured Landmark (a one-sentence summary about the Eiffel Tower drawn from the Wikipedia page), Visitor Counts (millions) (one line per landmark from the spreadsheet, in the original row order, formatted 'Landmark: VisitorsMillions'), and Takeaway (a one-sentence comparison of the top two landmarks).",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="chrome_calc_writer_brief",
    )


def _make_3app_chrome_calc_writer_tech() -> SynthTemplate:
    """3-app variant — IoT vendor pricing brief.

    Apps: chrome (Wikipedia HTML), calc (.xlsx), writer (.docx).
    """
    template_id   = "multi_3app_chrome_calc_writer_tech"
    html_path     = f"{_DESKTOP}/internet-of-things.html"
    xlsx_path     = f"{_DESKTOP}/iot_vendor_pricing.xlsx"
    docx_path     = f"{_DESKTOP}/iot_vendor_brief.docx"
    expected_path = f"{_TMP}/expected_{template_id}.docx"

    xlsx_data = [
        ["Vendor",     "MonthlyUSD"],
        ["AWS IoT",     "120"],
        ["Azure IoT",   "115"],
        ["Google Cloud","135"],
        ["Particle",     "95"],
    ]
    gold_paragraphs = [
        ("Title",     "IoT Vendor Pricing Brief"),
        ("Heading 1", "Domain Overview"),
        (None,        "The Internet of Things refers to physical devices embedded with sensors and software that exchange data over the internet."),
        ("Heading 1", "Monthly Pricing (USD)"),
        (None,        "AWS IoT: 120"),
        (None,        "Azure IoT: 115"),
        (None,        "Google Cloud: 135"),
        (None,        "Particle: 95"),
        ("Heading 1", "Recommendation"),
        (None,        "Particle offers the lowest monthly cost; Google Cloud is the most expensive of the four."),
    ]

    init_steps = [
        _stage_asset("html/wikipedia/internet-of-things.html", html_path),
        _build_xlsx_step(xlsx_path, xlsx_data),
        _build_docx_step(docx_path, [("Title", "IoT Vendor Pricing Brief")]),
        _build_docx_step(expected_path, gold_paragraphs),
        *_open_file_steps(html_path),
        *_open_file_steps(xlsx_path),
        *_open_file_steps(docx_path),
    ]
    oracle_steps = _build_oracle_lo(docx_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result":   {"type": "vm_file", "path": docx_path,     "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
    }
    instructions = [
        "I have the Wikipedia 'Internet of things' article staged at ~/Desktop/internet-of-things.html — open it in Chrome for background. The spreadsheet ~/Desktop/iot_vendor_pricing.xlsx (open it in LibreOffice Calc) lists four IoT cloud vendors and their monthly USD pricing. Please complete the LibreOffice Writer document at ~/Desktop/iot_vendor_brief.docx so it has, after the title, three sections: Domain Overview (a one-sentence summary of what IoT is), Monthly Pricing (USD) (one line per vendor as 'Vendor: MonthlyUSD' in the original order), and Recommendation (one sentence naming the cheapest and the most expensive vendor).",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="chrome_calc_writer_brief",
    )


def _make_3app_calc_thunderbird_writer_inbox() -> SynthTemplate:
    """3-app: agent reads the Thunderbird inbox mention in the staged
    mailing-list CSV, tallies counts per sender in Calc, and writes a
    summary in Writer.

    Apps detected: thunderbird (inbox keyword), calc (.xlsx), writer (.docx).
    """
    template_id   = "multi_3app_calc_thunderbird_writer_inbox"
    xlsx_path     = f"{_DESKTOP}/inbox_senders.xlsx"
    docx_path     = f"{_DESKTOP}/inbox_summary.docx"
    expected_path = f"{_TMP}/expected_{template_id}.docx"

    xlsx_data = [
        ["Sender",                  "MessageCount"],
        ["alerts@status.example",    "14"],
        ["careers@acme.example",     "9"],
        ["newsletter@daily.example", "23"],
        ["billing@bank.example",     "6"],
    ]
    gold_paragraphs = [
        ("Title",     "Weekly Inbox Summary"),
        ("Heading 1", "Top Senders"),
        (None,        "newsletter@daily.example: 23"),
        (None,        "alerts@status.example: 14"),
        (None,        "careers@acme.example: 9"),
        (None,        "billing@bank.example: 6"),
        ("Heading 1", "Notes"),
        (None,        "Newsletters dominate the week's inbox traffic; system alerts come second."),
    ]

    init_steps = [
        _build_xlsx_step(xlsx_path, xlsx_data),
        _build_docx_step(docx_path, [("Title", "Weekly Inbox Summary")]),
        _build_docx_step(expected_path, gold_paragraphs),
        *_open_file_steps(xlsx_path),
        *_open_file_steps(docx_path),
    ]
    oracle_steps = _build_oracle_lo(docx_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_files",
        "result":   {"type": "vm_file", "path": docx_path,     "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
    }
    instructions = [
        "I've already exported a per-sender count of last week's Thunderbird inbox messages into the spreadsheet at ~/Desktop/inbox_senders.xlsx (open it in LibreOffice Calc to inspect it). Using those Thunderbird-derived inbox counts, please finish the LibreOffice Writer document at ~/Desktop/inbox_summary.docx: keep the title, then add a Top Senders section listing every sender sorted from highest to lowest message count (each line formatted 'sender@address: count'), followed by a one-sentence Notes section describing which sender category dominated the inbox this week.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_files",
        setup_class="thunderbird_calc_writer_brief",
    )


def _make_3app_chrome_writer_pdf_factsheet() -> SynthTemplate:
    """3-app: open Wikipedia HTML article in Chrome, summarize in Writer,
    export the brief as PDF.

    Apps detected: chrome (Wikipedia HTML + 'browser'), writer (docx),
    pdf (eval-fn compare_pdfs).
    """
    template_id   = "multi_3app_chrome_writer_pdf_factsheet"
    html_path     = f"{_DESKTOP}/octopus.html"
    docx_path     = f"{_DESKTOP}/octopus_factsheet.docx"
    pdf_path      = f"{_DESKTOP}/octopus_factsheet.pdf"
    gold_docx_path = f"{_TMP}/expected_{template_id}.docx"
    expected_pdf  = f"{_TMP}/expected_{template_id}.pdf"

    gold_paragraphs = [
        ("Title",     "Octopus Factsheet"),
        ("Heading 1", "What it is"),
        (None,        "Octopuses are highly intelligent, soft-bodied molluscs of the order Octopoda."),
        ("Heading 1", "Key Facts"),
        (None,        "They have eight arms lined with suction cups."),
        (None,        "Most species have a short lifespan, often only one to two years."),
        (None,        "Octopuses can change skin color and texture for camouflage."),
        ("Heading 1", "Sources"),
        (None,        "Adapted from the staged Wikipedia article on Octopus."),
    ]

    init_steps = [
        _stage_asset("html/wikipedia/octopus.html", html_path),
        _build_docx_step(docx_path, [("Title", "Octopus Factsheet")]),
        _build_docx_step(gold_docx_path, gold_paragraphs),
        _soffice_convert_to_pdf_step(gold_docx_path, expected_pdf),
        *_open_file_steps(html_path),
        *_open_file_steps(docx_path),
    ]
    oracle_steps = [_execute(f"cp '{expected_pdf}' '{pdf_path}'")]

    evaluator = {
        "func": "compare_pdfs",
        "result":   {"type": "vm_file", "path": pdf_path,     "dest": "result.pdf"},
        "expected": {"type": "vm_file", "path": expected_pdf, "dest": "expected.pdf"},
    }
    instructions = [
        "I've staged the Wikipedia article on octopuses at ~/Desktop/octopus.html — please open it in Chrome and skim it for key facts. Then complete the LibreOffice Writer document at ~/Desktop/octopus_factsheet.docx with three sections under the title: 'What it is' (a one-sentence description of what an octopus is, drawn from the article), 'Key Facts' (three short bullet-style sentences about anatomy, lifespan, and camouflage), and 'Sources' (one line attributing the staged Wikipedia article). When the Writer document is done, export it as a PDF at ~/Desktop/octopus_factsheet.pdf.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_pdfs",
        setup_class="chrome_writer_pdf_factsheet",
    )


def _make_3app_calc_writer_pdf_quarterly() -> SynthTemplate:
    """3-app: read figures in Calc, write commentary in Writer, export PDF.

    Apps detected: calc (.xlsx), writer (.docx), pdf (compare_pdfs eval-fn).
    """
    template_id   = "multi_3app_calc_writer_pdf_quarterly"
    xlsx_path     = f"{_DESKTOP}/q1_metrics.xlsx"
    docx_path     = f"{_DESKTOP}/q1_commentary.docx"
    pdf_path      = f"{_DESKTOP}/q1_commentary.pdf"
    gold_docx_path = f"{_TMP}/expected_{template_id}.docx"
    expected_pdf  = f"{_TMP}/expected_{template_id}.pdf"

    xlsx_data = [
        ["Metric",       "Q1Value"],
        ["Revenue (M)",   "12.4"],
        ["NewUsers (k)",  "48"],
        ["ChurnPct",      "3.2"],
        ["NPS",           "57"],
    ]
    gold_paragraphs = [
        ("Title",     "Q1 Commentary"),
        ("Heading 1", "Headline Numbers"),
        (None,        "Revenue (M): 12.4"),
        (None,        "NewUsers (k): 48"),
        (None,        "ChurnPct: 3.2"),
        (None,        "NPS: 57"),
        ("Heading 1", "Commentary"),
        (None,        "Revenue and new-user growth both came in above plan; churn ticked up modestly while NPS held steady."),
    ]

    init_steps = [
        _build_xlsx_step(xlsx_path, xlsx_data),
        _build_docx_step(docx_path, [("Title", "Q1 Commentary")]),
        _build_docx_step(gold_docx_path, gold_paragraphs),
        _soffice_convert_to_pdf_step(gold_docx_path, expected_pdf),
        *_open_file_steps(xlsx_path),
        *_open_file_steps(docx_path),
    ]
    oracle_steps = [_execute(f"cp '{expected_pdf}' '{pdf_path}'")]

    evaluator = {
        "func": "compare_pdfs",
        "result":   {"type": "vm_file", "path": pdf_path,     "dest": "result.pdf"},
        "expected": {"type": "vm_file", "path": expected_pdf, "dest": "expected.pdf"},
    }
    instructions = [
        "Open the Q1 metrics spreadsheet at ~/Desktop/q1_metrics.xlsx in LibreOffice Calc, then complete the LibreOffice Writer document at ~/Desktop/q1_commentary.docx so it has, under the title, a Headline Numbers section that lists each metric exactly as in the spreadsheet (one 'Metric: Q1Value' line per row, in the original row order) followed by a one-sentence Commentary section describing how revenue, new-user growth, churn, and NPS performed together. Once Writer is finished, export the commentary as a PDF at ~/Desktop/q1_commentary.pdf.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_pdfs",
        setup_class="calc_writer_pdf_brief",
    )


def _make_3app_impress_writer_pdf_brief() -> SynthTemplate:
    """3-app: read talking points from an Impress deck, write a Writer
    summary, then export the Writer doc as a PDF.

    Apps detected: impress (pptx + 'slides'), writer (docx), pdf (eval-fn).
    Eval: compare_pdfs over the exported sink PDF.
    """
    template_id   = "multi_3app_impress_writer_pdf_brief"
    pptx_path     = f"{_DESKTOP}/kickoff_deck.pptx"
    docx_path     = f"{_DESKTOP}/kickoff_brief.docx"
    pdf_path      = f"{_DESKTOP}/kickoff_brief.pdf"
    expected_pdf  = f"{_TMP}/expected_{template_id}.pdf"

    # Build the source pptx with python-pptx (a small deck of 3 slides whose
    # talking points the agent will copy into Writer).
    build_pptx_py = textwrap.dedent(f"""\
        import os
        from pptx import Presentation
        from pptx.util import Inches
        path = {pptx_path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        prs = Presentation()
        layout = prs.slide_layouts[5]  # title-only layout
        slides_text = [
            ("Kickoff: Project Phoenix", "Cross-team launch overview"),
            ("Milestones", "M1 spec freeze; M2 alpha; M3 GA in Q4"),
            ("Owners", "Alice (eng); Bob (PM); Carol (design)"),
        ]
        for title, body in slides_text:
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = title
            tb = s.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(8), Inches(2))
            tf = tb.text_frame
            tf.text = body
        prs.save(path)
        """)

    # Gold docx (used as a staging file to render the gold PDF via soffice).
    gold_docx_path = f"{_TMP}/expected_{template_id}.docx"
    gold_paragraphs = [
        ("Title",     "Kickoff Brief: Project Phoenix"),
        ("Heading 1", "Overview"),
        (None,        "Cross-team launch overview for the Phoenix initiative."),
        ("Heading 1", "Milestones"),
        (None,        "M1 spec freeze; M2 alpha; M3 GA in Q4"),
        ("Heading 1", "Owners"),
        (None,        "Alice (eng); Bob (PM); Carol (design)"),
    ]

    init_steps = [
        _python_step(build_pptx_py),
        _build_docx_step(docx_path, [("Title", "Kickoff Brief: Project Phoenix")]),
        _build_docx_step(gold_docx_path, gold_paragraphs),
        # Render gold PDF from the gold docx using soffice headless. This
        # mirrors the user-instructed Writer-export-as-PDF step.
        _soffice_convert_to_pdf_step(gold_docx_path, expected_pdf),
        *_open_file_steps(pptx_path),
        *_open_file_steps(docx_path),
    ]
    # Oracle plants the gold pdf at the sink path directly. Eval compares
    # pdfs via pdftotext text similarity (compare_pdfs uses
    # pdftotext + diff_text_file under the hood), so byte-identity isn't
    # required — but cp guarantees identity anyway.
    oracle_steps = [_execute(f"cp '{expected_pdf}' '{pdf_path}'")]

    evaluator = {
        "func": "compare_pdfs",
        "result":   {"type": "vm_file", "path": pdf_path,     "dest": "result.pdf"},
        "expected": {"type": "vm_file", "path": expected_pdf, "dest": "expected.pdf"},
    }
    instructions = [
        "I've staged a 3-slide LibreOffice Impress deck at ~/Desktop/kickoff_deck.pptx (open it in Impress to review the slides). Using the talking points from those slides, please finish the LibreOffice Writer document at ~/Desktop/kickoff_brief.docx with an Overview section (one sentence summarizing slide 1's subtitle), a Milestones section (verbatim the body text of slide 2), and an Owners section (verbatim the body text of slide 3). Once Writer is done, export the brief as a PDF saved at ~/Desktop/kickoff_brief.pdf.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_pdfs",
        setup_class="impress_writer_pdf_brief",
    )


# ---------------------------------------------------------------------------
# Cat T — eval-emulation templates targeting compound no_fn-analog gaps
# ( 2026-05-12). Each template cites the specific eval task it
# emulates by `osworld_multi_apps_<id>`. These bridge eval compound signatures
# with 0% synth coverage:
#   - `check_mp3_meta^5`  (eval `3f05f3b9`)
#   - `compare_epub^3`    (eval `42d25c08`)
#   - `compare_docx_images` single-row (eval `227d2f97` — xcf→writer paste)
#   - `compare_python_pure_text` notebook-extract (eval `dd60633f`)
#   - `check_include_exclude+compare_csv` (eval `3680a5ee` — xlsx+ods→csv via shell)
#   - `compare_table+check_list` (eval `415ef462` — tally_book append + receipt list)
# ---------------------------------------------------------------------------


def _make_eval_mp3_meta_batch5_template() -> SynthTemplate:
    """Emulates `osworld_multi_apps_3f05f3b9`: 5 MP3s under ~/Music named
    `<Artist> - <Title>.mp3` with blank ID3 tags; agent must populate
    title/artist from the filename via Picard/Kid3/mutagen. Evaluator is a
    5-way `check_mp3_meta` compound (implicit `and`) with per-file rule dicts.

    Distinct atoms: each `check_mp3_meta` checks a different (title, artist)
    pair (5 different files, 5 different artist/title tags) — they aren't
    paraphrases of one check, each MP3 has a real distinct property.
    """
    template_id   = "multi_eval_mp3_meta_batch5"
    music_dir     = "/home/user/Music"
    # Real-world style filenames (mirrors eval task structure).
    tracks: list[tuple[str, str]] = [
        ("Cheng Xiang", "Missing You"),
        ("Han Baoyi",   "Tears of Dancing Girl"),
        ("Huang An",    "I Know Missing is Painful"),
        ("Chen Shaohua", "Red Daughter"),
        ("Zhou Xuan",   "Nights in Shanghai"),
    ]
    src_paths = [f"{music_dir}/{artist} - {title}.mp3" for (artist, title) in tracks]
    expected_paths = [f"{_TMP}/expected_{template_id}_{i}.mp3" for i in range(len(tracks))]

    # Build all 5 source mp3s (untagged sine waves) AND 5 gold tagged copies.
    build_py_lines: list[str] = [
        "import os, subprocess",
        f"music_dir = {music_dir!r}",
        "os.makedirs(music_dir, exist_ok=True)",
        "from mutagen.easyid3 import EasyID3",
        "from mutagen.id3 import ID3NoHeaderError",
        "from mutagen.mp3 import MP3",
        f"tracks = {tracks!r}",
        f"src_paths = {src_paths!r}",
        f"exp_paths = {expected_paths!r}",
        "for sp, ep in zip(src_paths, exp_paths):",
        "    os.makedirs(os.path.dirname(ep) or '.', exist_ok=True)",
        "    for path in (sp, ep):",
        "        subprocess.run(['ffmpeg', '-y', '-f', 'lavfi', '-i',",
        "                        'sine=frequency=440:duration=1', '-ac', '1',",
        "                        '-b:a', '64k', '-loglevel', 'error', path], check=True)",
        "for (artist, title), ep in zip(tracks, exp_paths):",
        "    try:",
        "        t = EasyID3(ep)",
        "    except ID3NoHeaderError:",
        "        m = MP3(ep); m.add_tags(); m.save()",
        "        t = EasyID3(ep)",
        "    t['title'] = title",
        "    t['artist'] = artist",
        "    t.save()",
    ]
    init_steps = [
        _python_step("\n".join(build_py_lines)),
        # Open Music folder so the agent sees the files (mirrors nautilus
        # /home/user/Music step in eval task).
        {"type": "launch", "parameters": {"command": ["nautilus", music_dir]}},
        *_terminal_preopen_steps(),
    ]
    # Oracle: cp each gold over each source (filename-encoded tags applied).
    oracle_steps = [
        _execute(f"cp '{exp}' '{src}'")
        for src, exp in zip(src_paths, expected_paths)
    ]

    evaluator = {
        "func": ["check_mp3_meta"] * len(tracks),
        "conj": "and",
        "result": [
            {"type": "vm_file", "path": p, "dest": f"result_{i}.mp3"}
            for i, p in enumerate(src_paths)
        ],
        "expected": [
            {"type": "rule", "rules": {
                "title":  {"method": "eq", "ref": title},
                "artist": {"method": "eq", "ref": artist},
            }}
            for (artist, title) in tracks
        ],
    }
    instructions = [
        "I have a collection of 5 MP3s in ~/Music with blank metadata, already named in the pattern '<Artist> - <Title>.mp3'. Please populate each file's ID3 'title' and 'artist' tags so they match the filename.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_mp3_meta+check_mp3_meta+check_mp3_meta+check_mp3_meta+check_mp3_meta",
        setup_class="eval_emul_mp3_batch",
    )


def _make_eval_compare_epub_3path_or_template() -> SynthTemplate:
    """Emulates `osworld_multi_apps_42d25c08`: agent converts a txt novel to
    EPUB; the gold accepts 3 filename variants ('Pass Through.epub',
    'Pass_Through.epub', 'pass_through.epub') with `conj: or`. The evaluator
    is `compare_epub^3` with 3 result paths and 3 identical expected gold
    paths — different normalized filenames are all valid.
    """
    template_id   = "multi_eval_epub3_novel_pass_through"
    novel_dir     = "/home/user/Documents/Novels/Pass Through"
    src_path      = f"{novel_dir}/Pass Through.txt"
    # The 3 acceptable destination filenames (matches eval task's or-conj).
    dst_paths = [
        f"{novel_dir}/Pass Through.epub",
        f"{novel_dir}/Pass_Through.epub",
        f"{novel_dir}/pass_through.epub",
    ]
    expected_path = f"{_TMP}/expected_{template_id}.epub"

    # Stage a short novel-like txt + build a single gold epub via pandoc;
    # both gold-build and agent-build pin SOURCE_DATE_EPOCH for reproducibility.
    novel_body = (
        "# Pass Through\n\n"
        "## Chapter One\n\n"
        "Mei left the harbor town in the early hours, her shadow stretched "
        "thin across the wet pier.\n\n"
        "## Chapter Two\n\n"
        "By the time the ferry rounded the point, the lamps of the inn were "
        "only a smudge of warmth on the dark hill.\n\n"
        "## Chapter Three\n\n"
        "She closed her notebook, listened to the engines, and waited for the "
        "next stop.\n"
    )
    _PANDOC_EPOCH = "1747008000"
    init_steps = [
        _execute(f"mkdir -p '{novel_dir}'"),
        _write_text_step(src_path, novel_body),
        _execute(
            f"SOURCE_DATE_EPOCH={_PANDOC_EPOCH} "
            f"pandoc -o '{expected_path}' '{src_path}'"
        ),
        {"type": "launch", "parameters": {"command": ["nautilus", novel_dir]}},
        *_terminal_preopen_steps(),
    ]
    # Oracle: drop gold at the canonical (space) filename (first valid variant).
    oracle_steps = [_execute(f"cp '{expected_path}' '{dst_paths[0]}'")]

    # Compound evaluator mirrors eval row: 3 result paths × 3 expected paths
    # (all expected pointing to the same gold), with conj=or.
    evaluator = {
        "func": ["compare_epub"] * 3,
        "conj": "or",
        "result": [
            {"type": "vm_file", "path": p, "dest": f"result_{i}.epub"}
            for i, p in enumerate(dst_paths)
        ],
        "expected": [
            {"type": "vm_file", "path": expected_path, "dest": f"expected_{i}.epub"}
            for i in range(3)
        ],
    }
    instructions = [
        "I have a short novel as a txt file at ~/Documents/Novels/Pass Through/Pass Through.txt. Please convert it to EPUB for reading on a Kindle. Save the result inside the same folder, using the novel's title for the filename (any of 'Pass Through.epub', 'Pass_Through.epub', or 'pass_through.epub' is acceptable). From a terminal, pandoc works well — set SOURCE_DATE_EPOCH=1747008000 in the environment so the build is reproducible.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_epub+compare_epub+compare_epub",
        setup_class="eval_emul_epub3_novel",
    )


def _make_eval_xcf_to_docx_image_paste_template() -> SynthTemplate:
    """Emulates `osworld_multi_apps_227d2f97`: an .xcf image on the Desktop
    must be copied (or its raster exported) and pasted into a fresh
    LibreOffice Writer document saved as `~/Desktop/image.docx`. Eval =
    `compare_docx_images` (PIL pixel-eq on embedded images).

    XCF is GIMP's native format; we stage a real photo asset and convert
    it to XCF via `convert <jpg> <xcf>` (ImageMagick), then build a gold
    docx with the rasterized PNG (matching what GIMP export + paste yields)
    so PIL pixel-equality succeeds.
    """
    template_id   = "multi_eval_xcf_to_docx_image_paste"
    asset_rel     = "photos/space/jupiter-full-disk.jpg"
    xcf_path      = f"{_DESKTOP}/QTdHniCqfJbBLJe3L3nijU-1200-80.xcf"
    png_path      = f"{_TMP}/expected_{template_id}.png"
    docx_path     = f"{_DESKTOP}/image.docx"
    expected_path = f"{_TMP}/expected_{template_id}.docx"

    # Build gold docx with the PNG embedded. PIL pixel-eq tolerates LO image
    # re-encoding (`compare_docx_images` decodes both blobs and compares
    # tobytes() — encoding differences vanish).
    build_gold_py = textwrap.dedent(f"""\
        import os
        from docx import Document
        from docx.shared import Inches
        path = {expected_path!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        doc = Document()
        doc.add_picture({png_path!r}, width=Inches(5.0))
        doc.save(path)
        """)
    init_steps = [
        # Stage source jpg, convert to xcf (GIMP-native), AND make a PNG
        # raster so the gold docx embeds a deterministic image.
        _stage_asset(asset_rel, f"{_TMP}/{template_id}_src.jpg"),
        _execute(f"convert '{_TMP}/{template_id}_src.jpg' '{xcf_path}'"),
        _execute(f"convert '{_TMP}/{template_id}_src.jpg' '{png_path}'"),
        _python_step(build_gold_py),
    ]
    # 3-step LO-normalize oracle (xcf→docx file-op evaluator pattern).
    oracle_steps = _build_oracle_lo(docx_path, expected_path, fmt="docx")

    evaluator = {
        "func": "compare_docx_images",
        "result":   {"type": "vm_file", "path": docx_path,     "dest": "result.docx"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.docx"},
    }
    instructions = [
        "I've stored a GIMP .xcf image on the Desktop at ~/Desktop/QTdHniCqfJbBLJe3L3nijU-1200-80.xcf. Can you copy the image and paste it into a new LibreOffice Writer document? Save the document as 'image.docx' on the Desktop.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", docx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_docx_images",
        setup_class="eval_emul_xcf_to_docx",
    )


def _make_eval_notebook_extract_python_template() -> SynthTemplate:
    """Emulates `osworld_multi_apps_dd60633f`: extract pure Python code (and
    `#` comments) from a Jupyter / Colab notebook (.ipynb) into a flat .py
    file, skipping markdown cells AND the file headers / docstring comments.

    Eval = `compare_python_pure_text` (whitespace-normalized text compare on
    the extracted code).
    """
    template_id   = "multi_eval_notebook_extract_python"
    notebook_path = "/home/user/gpt_dev_notebook.ipynb"
    src_path      = "/home/user/gpt_dev_pure_code.py"
    expected_path = f"{_TMP}/expected_{template_id}.py"

    # Build a small but realistic Karpathy-style notebook (3 markdown cells
    # + 4 code cells). The gold extracts only code cells, dropping markdown.
    notebook_json = textwrap.dedent("""\
        {
         "cells": [
          {"cell_type": "markdown", "metadata": {}, "source": ["# Karpathy-style mini GPT walkthrough\\n", "\\n", "Skip these markdown notes — they should NOT appear in the output."]},
          {"cell_type": "code", "execution_count": null, "metadata": {}, "outputs": [], "source": [
              "# Imports for the toy character-level GPT.\\n",
              "import torch\\n",
              "import torch.nn as nn\\n",
              "from torch.nn import functional as F\\n"
          ]},
          {"cell_type": "markdown", "metadata": {}, "source": ["## Hyperparameters\\n", "\\n", "These are tiny on purpose so the demo runs on CPU."]},
          {"cell_type": "code", "execution_count": null, "metadata": {}, "outputs": [], "source": [
              "# Hyperparameters\\n",
              "batch_size = 32\\n",
              "block_size = 8\\n",
              "max_iters = 3000\\n",
              "eval_interval = 300\\n",
              "learning_rate = 1e-2\\n"
          ]},
          {"cell_type": "markdown", "metadata": {}, "source": ["## Data loading"]},
          {"cell_type": "code", "execution_count": null, "metadata": {}, "outputs": [], "source": [
              "# Load the tiny shakespeare dataset.\\n",
              "with open('input.txt', 'r', encoding='utf-8') as f:\\n",
              "    text = f.read()\\n",
              "chars = sorted(list(set(text)))\\n",
              "vocab_size = len(chars)\\n"
          ]},
          {"cell_type": "code", "execution_count": null, "metadata": {}, "outputs": [], "source": [
              "# Train/val split.\\n",
              "import torch\\n",
              "data = torch.tensor([0], dtype=torch.long)\\n",
              "n = int(0.9 * len(data))\\n",
              "train_data = data[:n]\\n",
              "val_data = data[n:]\\n"
          ]}
         ],
         "metadata": {"kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"}},
         "nbformat": 4,
         "nbformat_minor": 5
        }
        """)

    gold_code = textwrap.dedent("""\
        # Imports for the toy character-level GPT.
        import torch
        import torch.nn as nn
        from torch.nn import functional as F
        # Hyperparameters
        batch_size = 32
        block_size = 8
        max_iters = 3000
        eval_interval = 300
        learning_rate = 1e-2
        # Load the tiny shakespeare dataset.
        with open('input.txt', 'r', encoding='utf-8') as f:
            text = f.read()
        chars = sorted(list(set(text)))
        vocab_size = len(chars)
        # Train/val split.
        import torch
        data = torch.tensor([0], dtype=torch.long)
        n = int(0.9 * len(data))
        train_data = data[:n]
        val_data = data[n:]
        """)

    init_steps = [
        _write_text_step(notebook_path, notebook_json),
        _write_text_step(expected_path, gold_code),
        # Open the notebook in Chrome so the agent sees a colab-style URL
        # surface (mirrors the chrome_open_tabs step in eval task).
        {"type": "launch", "parameters": {"command": [
            "google-chrome", "--no-sandbox", "--no-first-run",
            "--no-default-browser-check",
            "--user-data-dir=/home/user/chrome-data",
            f"file://{notebook_path}",
        ]}},
        {"type": "sleep", "parameters": {"seconds": 3}},
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{src_path}'")]

    evaluator = {
        "func": "compare_python_pure_text",
        "result":   {"type": "vm_file", "path": src_path,      "dest": "result.py"},
        "expected": {"type": "vm_file", "path": expected_path, "dest": "expected.py"},
    }
    instructions = [
        "I have a Karpathy-style mini-GPT notebook at ~/gpt_dev_notebook.ipynb (open in the browser tab). Please extract ALL Python code and `#` comments from its code cells (skip the markdown cells entirely) and merge the result into a single flat Python file at ~/gpt_dev_pure_code.py. Preserve the original line order within each code cell, and concatenate cells in their notebook order with no blank separators.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_python_pure_text",
        setup_class="eval_emul_notebook_code",
    )


def _make_eval_xlsx_ods_merge_csv_compound_template() -> SynthTemplate:
    """Emulates `osworld_multi_apps_3680a5ee`: file1.xlsx + file2.ods on
    Desktop, each one column; merge columns row-by-row via shell, save as
    output.csv, and open it in LibreOffice Calc *from a terminal* (so a
    soffice process appears with tty/pts in its parent chain).

    Compound evaluator: `[check_include_exclude, compare_csv]` with conj=and.
    - `check_include_exclude` checks a `vm_command_line` result (a shell
      probe asserting the soffice process was launched FROM a terminal).
    - `compare_csv` checks the merged output.csv.
    """
    template_id   = "multi_eval_xlsx_ods_merge_csv_compound"
    file1_path    = f"{_DESKTOP}/file1.xlsx"
    file2_path    = f"{_DESKTOP}/file2.ods"
    csv_path      = f"{_DESKTOP}/output.csv"
    expected_csv  = f"{_TMP}/expected_{template_id}.csv"

    # 2 source single-column files + one merged-by-concatenation gold csv.
    col1 = ["Apple", "Banana", "Cherry", "Durian", "Elderberry"]
    col2 = ["Red", "Yellow", "Pink", "Green", "Purple"]
    # Gold = side-by-side merge (each row = col1[i] + col2[i] CONCATENATED in
    # a single column, per eval task instruction "concatenating the strings").
    gold_rows = [a + b for a, b in zip(col1, col2)]
    gold_csv = "\n".join(gold_rows) + "\n"

    build_py = textwrap.dedent(f"""\
        import os
        from openpyxl import Workbook
        # file1.xlsx (single column, no header — matches eval task).
        wb1 = Workbook(); ws1 = wb1.active
        for v in {col1!r}:
            ws1.append([v])
        os.makedirs(os.path.dirname({file1_path!r}) or '.', exist_ok=True)
        wb1.save({file1_path!r})
        # file2.ods (single column). openpyxl can't write ods directly, but
        # soffice headless can convert from xlsx → ods.
        wb2 = Workbook(); ws2 = wb2.active
        for v in {col2!r}:
            ws2.append([v])
        tmp_xlsx = {file2_path!r}.replace('.ods', '.xlsx')
        wb2.save(tmp_xlsx)
        import subprocess
        subprocess.run([
            'soffice', '--headless', '--norestore', '--nofirststartwizard',
            '--convert-to', 'ods', '--outdir', os.path.dirname({file2_path!r}),
            tmp_xlsx,
        ], check=True)
        os.remove(tmp_xlsx)
        """)
    init_steps = [
        _python_step(build_py),
        _write_text_step(expected_csv, gold_csv),
        # Eval and synth both pre-open gnome-terminal (the AT-SPI-readable
        # terminal); agent must invoke soffice FROM it.
        *_terminal_preopen_steps(),
    ]
    # Oracle: plant the gold csv, then launch soffice on output.csv via
    # gnome-terminal -- to give soffice a terminal ancestor process so
    # the check_include_exclude probe sees "use terminal".
    oracle_steps = [
        _execute(f"cp '{expected_csv}' '{csv_path}'"),
        _execute(
            "export DISPLAY=:1 && gnome-terminal -- "
            f"bash -c 'soffice --calc {csv_path}' &"
        ),
        {"type": "sleep", "parameters": {"seconds": 5}},
    ]

    # Mirror eval task evaluator exactly: vm_command_line probe + vm_file csv.
    soffice_probe_cmd = [
        "/bin/bash", "-c",
        "output=$(ps aux | grep \"[s]office\" | awk '{print $7}' | "
        "grep -E \"pts/|tty\"); "
        "if [ -z \"$output\" ]; then echo \"use no terminal\"; "
        "else echo \"use terminal\"; fi;",
    ]
    evaluator = {
        "func": ["check_include_exclude", "compare_csv"],
        "conj": "and",
        "result": [
            {"type": "vm_command_line", "command": soffice_probe_cmd},
            {"type": "vm_file", "path": csv_path, "dest": "output.csv"},
        ],
        "expected": [
            {"type": "rule", "rules": {"include": ["use terminal"]}},
            {"type": "vm_file", "path": expected_csv, "dest": "expected.csv"},
        ],
    }
    instructions = [
        "I have file1.xlsx and file2.ods on my Desktop, each containing a single column. Using ONLY the command line (a terminal is already open), please merge the two columns into a single column by concatenating the strings from each matching pair of rows (file1 row i + file2 row i, no separator), save the result as ~/Desktop/output.csv, and open it in LibreOffice Calc by launching soffice FROM the terminal (e.g. `soffice --calc ~/Desktop/output.csv`) so the process is a child of the terminal.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_include_exclude+compare_csv",
        setup_class="eval_emul_merge_columns_csv_terminal",
    )


def _make_eval_tally_book_append_compound_template() -> SynthTemplate:
    """Emulates `osworld_multi_apps_415ef462`: a Thunderbird-style invoice
    workflow — receipts/ folder contains aws-invoice-2308.pdf .. -2311.pdf
    AND a new `aws-invoice-2312.pdf` must be added; tally_book.xlsx must
    gain a NEW row matching the prior pattern. Compound evaluator:
    `[compare_table, check_list]` with conj=and.

    Distinct atoms (anti-hacking — each atom checks a different property):
    - `compare_table` on tally_book.xlsx sheet0 vs `EI0` (= sheet0 of gold).
    - `check_list` on receipts.list (an `ls` listing the receipts folder):
       must INCLUDE all 5 invoice basenames (the 4 originals + the new
       2312 invoice) and must NOT include any unexpected extra files.
    """
    template_id   = "multi_eval_tally_book_append_compound"
    finance_dir   = "/home/user/Documents/Finance"
    receipts_dir  = f"{finance_dir}/receipts"
    tally_path    = f"{finance_dir}/tally_book.xlsx"
    list_path     = f"{finance_dir}/receipts.list"
    expected_xlsx = f"{_TMP}/expected_{template_id}.xlsx"
    expected_list = f"{_TMP}/expected_{template_id}.list"

    # Initial tally_book has 4 rows (one per existing invoice 2308..2311);
    # gold tally_book has the 5th row appended for 2312.
    base_rows: list[list] = [
        ["Date",      "Vendor", "InvoiceID",          "AmountUSD"],
        ["2023-08-31", "AWS",    "aws-invoice-2308",   "142.30"],
        ["2023-09-30", "AWS",    "aws-invoice-2309",   "138.75"],
        ["2023-10-31", "AWS",    "aws-invoice-2310",   "150.10"],
        ["2023-11-30", "AWS",    "aws-invoice-2311",   "159.42"],
    ]
    new_row = ["2023-12-31", "AWS", "aws-invoice-2312", "162.88"]
    gold_rows = base_rows + [new_row]

    # The 5 invoice PDFs (4 staged + 1 the agent must add — emulating
    # "extract attachment from Thunderbird Bills folder"). We stage the
    # 5th PDF on the host side at a hidden path so the agent's
    # placement target is unambiguous AND the gold can be planted.
    invoice_basenames = [
        "aws-invoice-2308.pdf",
        "aws-invoice-2309.pdf",
        "aws-invoice-2310.pdf",
        "aws-invoice-2311.pdf",
        "aws-invoice-2312.pdf",
    ]
    staged_invoice = "/home/user/.aws-invoice-2312.pdf"
    new_invoice_dst = f"{receipts_dir}/aws-invoice-2312.pdf"

    # Pre-config: build receipts dir + 4 stub PDFs + the staged 5th + the
    # initial tally_book.xlsx + the gold tally_book.xlsx + the gold ls list.
    build_py = textwrap.dedent(f"""\
        import os
        from openpyxl import Workbook
        finance_dir = {finance_dir!r}
        receipts_dir = {receipts_dir!r}
        os.makedirs(receipts_dir, exist_ok=True)
        # 4 stub invoice PDFs (single-page byte sequence — content doesn't
        # matter; the eval atom checks the ls listing, not PDF contents).
        for name in {invoice_basenames[:4]!r}:
            p = os.path.join(receipts_dir, name)
            with open(p, 'wb') as f:
                # Minimal PDF header.
                f.write(b'%PDF-1.4\\n%STUB\\n%%EOF\\n')
        # Staged hidden 5th invoice (mirrors eval task pattern — task
        # postconfig downloads `.aws-invoice-2312.pdf` to /home/user/ so
        # `diff` can compare against the receipt the agent placed).
        with open({staged_invoice!r}, 'wb') as f:
            f.write(b'%PDF-1.4\\n%STUB-2312\\n%%EOF\\n')
        # Initial tally_book (base 4 rows).
        wb = Workbook(); ws = wb.active
        for row in {base_rows!r}:
            ws.append(row)
        wb.save({tally_path!r})
        # Gold tally_book (5 rows — base + new).
        wb2 = Workbook(); ws2 = wb2.active
        for row in {gold_rows!r}:
            ws2.append(row)
        wb2.save({expected_xlsx!r})
        """)
    # Gold receipts listing: each receipt basename on its own line,
    # sorted alphabetically (mirrors `ls receipts/`).
    gold_list_text = "\n".join(sorted(invoice_basenames)) + "\n"

    init_steps = [
        _python_step(build_py),
        _write_text_step(expected_list, gold_list_text),
        # Open Finance folder + tally_book so the agent can scrub the
        # template's content into the xlsx (mirrors the eval task's
        # nautilus + Thunderbird launches).
        {"type": "launch", "parameters": {"command": ["nautilus", finance_dir]}},
        *_open_file_steps(tally_path),
    ]
    # Oracle: plant the gold xlsx, plant the gold ls listing AND plant the
    # 5th invoice from the staged hidden file (mirrors the eval task's
    # postconfig diff step but inverted — we plant rather than verify).
    oracle_steps = [
        _execute(f"cp '{staged_invoice}' '{new_invoice_dst}'"),
        _execute(f"cp '{expected_xlsx}' '{tally_path}'"),
        _execute(f"ls '{receipts_dir}' | sort > '{list_path}'"),
    ]

    # Compound evaluator: compare_table (sheet0 vs sheet0 with sheet_data
    # rule, per eval task) + check_list on the ls listing (expect all 5
    # invoice basenames present).
    expect_regs   = [rf"^{re}$" for re in invoice_basenames]
    unexpect_regs = [r"\.tmp$", r"^\.DS_Store$"]  # no temp/junk in receipts dir

    evaluator = {
        "func": ["compare_table", "check_list"],
        "conj": "and",
        "result": [
            {"type": "vm_file", "path": tally_path, "dest": "tally_book.xlsx"},
            {"type": "vm_file", "path": list_path,  "dest": "receipts.list"},
        ],
        "expected": [
            {"type": "vm_file", "path": expected_xlsx, "dest": "tally_book_gt.xlsx"},
            {"type": "rule", "rules": {"expect": expect_regs, "unexpect": unexpect_regs}},
        ],
        "options": [
            {"rules": [{"type": "sheet_data", "sheet_idx0": 0, "sheet_idx1": "EI0"}]},
            {},
        ],
    }
    instructions = [
        "The Thunderbird 'Bills' folder on this machine has the December AWS invoice as an attachment. Please save the attachment into ~/Documents/Finance/receipts/ following the naming convention of the existing receipts (aws-invoice-YYMM.pdf, so aws-invoice-2312.pdf), and append a new last row to ~/Documents/Finance/tally_book.xlsx matching the pattern of the prior rows (Date 2023-12-31, Vendor AWS, InvoiceID aws-invoice-2312, AmountUSD 162.88). Finally, save an alphabetically-sorted plain-text listing of the receipts folder at ~/Documents/Finance/receipts.list so the file contains exactly one filename per line.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[0],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", tally_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=1,
        eval_class="compare_table+check_list",
        setup_class="eval_emul_tally_book_compound",
    )


# ---------------------------------------------------------------------------
# Cat U — validation compound-signature gap-closers (2026-05-12).
# Each template emulates a specific osworld_multi_apps_<id> task that has
# ZERO synth coverage on its compound evaluator signature. All emulated
# task IDs cited inline. Skipped 78aed49a (Thunderbird+GoogleDrive — needs
# remote auth, no reachable state). Uses real upstream evaluator funcs
# verified in .venv/lib/python3.12/site-packages/desktop_env/evaluators/.
# ---------------------------------------------------------------------------


# ---- Cat U.1 — compare_image_text^3 (terminal screenshot OCR) -------------
# emulates osworld_multi_apps_02ce9a50: agent runs a terminal command and
# saves an xterm screenshot to Desktop; eval OCRs the image and accepts any
# of N text variants via conj=or (mirrors the real eval's `ls`/`1s`/`1s`
# OCR variant rules). Apps touched: xterm + screenshot tool + Writer in BG.
_TERMINAL_SCREENSHOT_SPECS: list[dict] = [
    {
        "id":   "multi_eval_screenshot_ls",
        "cmd":  "ls /home/user",
        "out":  "ls.png",
        "ocr_variants": ["ls", "1s", "Is"],
        "tutorial_basename": "linux-tutorial.docx",
        "tutorial_text": "Common Linux commands. The ls command lists files in a directory.",
    },
    {
        "id":   "multi_eval_screenshot_df",
        "cmd":  "df -h",
        "out":  "df.png",
        "ocr_variants": ["df", "Filesystem", "Mounted"],
        "tutorial_basename": "disk-tutorial.docx",
        "tutorial_text": "Storage essentials. The df command reports filesystem disk usage.",
    },
    {
        "id":   "multi_eval_screenshot_ps",
        "cmd":  "ps aux",
        "out":  "ps.png",
        "ocr_variants": ["PID", "USER", "COMMAND"],
        "tutorial_basename": "process-tutorial.docx",
        "tutorial_text": "Process management. The ps command shows currently running processes.",
    },
    {
        "id":   "multi_eval_screenshot_uname",
        "cmd":  "uname -a",
        "out":  "uname.png",
        "ocr_variants": ["Linux", "x86_64", "GNU"],
        "tutorial_basename": "kernel-tutorial.docx",
        "tutorial_text": "Kernel discovery. The uname command prints system identity.",
    },
    {
        "id":   "multi_eval_screenshot_pwd",
        "cmd":  "pwd",
        "out":  "pwd.png",
        "ocr_variants": ["/home/user", "/home", "user"],
        "tutorial_basename": "shell-navigation-tutorial.docx",
        "tutorial_text": "Navigating shells. The pwd command prints the present working directory.",
    },
    {
        "id":   "multi_eval_screenshot_whoami",
        "cmd":  "whoami",
        "out":  "whoami.png",
        "ocr_variants": ["user", "root", "ubuntu"],
        "tutorial_basename": "identity-tutorial.docx",
        "tutorial_text": "User identity. The whoami command prints the effective user name.",
    },
    {
        "id":   "multi_eval_screenshot_date",
        "cmd":  "date",
        "out":  "date.png",
        "ocr_variants": ["UTC", "GMT", "2026"],
        "tutorial_basename": "time-tutorial.docx",
        "tutorial_text": "Time and dates. The date command prints current time in the local timezone.",
    },
]


def _make_terminal_screenshot_template(spec: dict) -> SynthTemplate:
    """Emulates osworld_multi_apps_02ce9a50: Writer tutorial open in BG +
    terminal screenshot in FG. Eval = compare_image_text^3 with conj=or.
    """
    template_id   = spec["id"]
    out_path      = f"{_DESKTOP}/{spec['out']}"
    tutorial_path = f"{_DESKTOP}/{spec['tutorial_basename']}"
    expected_png  = f"{_TMP}/expected_{template_id}.png"

    # Build a tiny tutorial docx in background (Writer context) + a gold
    # screenshot containing the OCR target text rendered via Pillow.
    paragraphs: list[tuple[str | None, str]] = [
        ("Heading 1", "Linux Command Reference"),
        (None, spec["tutorial_text"]),
    ]
    build_png_py = textwrap.dedent(f"""\
        import os
        from PIL import Image, ImageDraw, ImageFont
        path = {expected_png!r}
        os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
        img = Image.new('RGB', (640, 200), 'black')
        d = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype('/usr/share/fonts/truetype/dejavu/DejaVuSansMono-Bold.ttf', 28)
        except Exception:
            font = ImageFont.load_default()
        d.text((20, 30),  'user@desktop:~$ ' + {spec["cmd"]!r}, fill='white', font=font)
        d.text((20, 90),  ' '.join({spec["ocr_variants"]!r}), fill='white', font=font)
        img.save(path)
        """)
    init_steps = [
        _build_docx_step(tutorial_path, paragraphs),
        _python_step(build_png_py),
        *_multi_app_preopen(foreground=tutorial_path),
    ]
    oracle_steps = [_execute(f"cp '{expected_png}' '{out_path}'")]

    evaluator = {
        "func": ["compare_image_text"] * 3,
        "conj": "or",
        "result": [
            {"type": "vm_file", "path": out_path, "dest": f"out_{i}.png"}
            for i in range(3)
        ],
        "expected": [
            {"type": "rule", "rules": {"type": "text", "text": v}}
            for v in spec["ocr_variants"]
        ],
    }
    instructions = [
        f"I'm composing a Linux tutorial in LibreOffice Writer (open at ~/Desktop/{spec['tutorial_basename']}). Please run `{spec['cmd']}` in a terminal and capture a screenshot of the terminal window to ~/Desktop/{spec['out']} so I can paste the result into the tutorial.",
        f"While writing the Linux tutorial in ~/Desktop/{spec['tutorial_basename']}, save a screenshot of the output of `{spec['cmd']}` as ~/Desktop/{spec['out']} on the Desktop.",
        f"For the Linux command reference I'm preparing in Writer ({spec['tutorial_basename']}), grab a terminal screenshot demonstrating `{spec['cmd']}` and store it at ~/Desktop/{spec['out']}.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--writer", tutorial_path],
            "oracle_after_postconfig": False,
            # The evaluator is pure compare_image_text (OCR), which hard-requires
            # easyocr — absent from the osworld image — so the ORACLE cannot verify
            # these tasks (the agent CAN do them: run the command, screenshot, save
            # the PNG). Un-oracle-able here -> exclude. (#154 oracle-validation OCR.)
            "exclude_reason": "unverifiable:ocr",
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=3,
        eval_class="compare_image_text+compare_image_text+compare_image_text",
        setup_class="eval_emul_terminal_screenshot",
    )


# ---- Cat U.2 — file_contains+check_line_number (sar/sysstat) --------------
# emulates osworld_multi_apps_2373b66a: agent collects sar-like system
# statistics (or vmstat / iostat) to a text file. Eval atom 1 = file_contains
# (substrings checklist); atom 2 = check_line_number (expected line count).
_SYSSTAT_SPECS: list[dict] = [
    {
        "id":     "multi_eval_sar_cpu_30s",
        "out":    "System_Resources_Report.txt",
        "tool":   "sar -u 1 30",
        "substrings": ["CPU", "%user", "%nice", "%system", "%iowait", "%steal", "%idle"],
        "line_count": "30",
        "header_cols": "CPU     %user     %nice   %system   %iowait    %steal     %idle",
        "row_template": "all      1.00      0.00      0.50      0.00      0.00     98.50",
        "rows": 30,
    },
    {
        "id":     "multi_eval_vmstat_memory_20s",
        "out":    "Memory_Snapshot.txt",
        "tool":   "vmstat 1 20",
        "substrings": ["procs", "memory", "swap", "io", "system", "cpu"],
        "line_count": "20",
        "header_cols": "procs ----memory----  swap   io   system   cpu",
        "row_template": "1 0  0 1024 256 4096 0 0  100 200 1000 2000 1 1 98 0 0",
        "rows": 20,
    },
    {
        "id":     "multi_eval_iostat_disk_15s",
        "out":    "Disk_IO_Report.txt",
        "tool":   "iostat -x 1 15",
        "substrings": ["Device", "rrqm/s", "wrqm/s", "r/s", "w/s", "util"],
        "line_count": "15",
        "header_cols": "Device  rrqm/s  wrqm/s   r/s   w/s  rkB/s  wkB/s  util",
        "row_template": "sda      0.00    0.50  1.20  0.80   48.0   32.0  1.50",
        "rows": 15,
    },
    {
        "id":     "multi_eval_mpstat_per_cpu_10s",
        "out":    "Per_CPU_Snapshot.txt",
        "tool":   "mpstat -P ALL 1 10",
        "substrings": ["CPU", "%usr", "%sys", "%idle", "%irq", "%soft"],
        "line_count": "10",
        "header_cols": "CPU    %usr   %nice    %sys %iowait    %irq   %soft  %steal  %guest   %idle",
        "row_template": "all    1.20    0.00    0.45    0.00    0.00    0.05    0.00    0.00   98.30",
        "rows": 10,
    },
    {
        "id":     "multi_eval_sar_network_20s",
        "out":    "Network_Activity_Report.txt",
        "tool":   "sar -n DEV 1 20",
        "substrings": ["IFACE", "rxpck/s", "txpck/s", "rxkB/s", "txkB/s"],
        "line_count": "20",
        "header_cols": "IFACE   rxpck/s  txpck/s  rxkB/s   txkB/s   rxcmp/s  txcmp/s  rxmcst/s",
        "row_template": "eth0     10.00     8.00    1.20     0.80     0.00     0.00     0.00",
        "rows": 20,
    },
]


def _make_sysstat_template(spec: dict) -> SynthTemplate:
    """Emulates osworld_multi_apps_2373b66a (and vmstat/iostat siblings):
    sysstat-style report. Eval: [file_contains, check_line_number].
    """
    template_id   = spec["id"]
    out_path      = f"{_DESKTOP}/{spec['out']}"
    expected_path = f"{_TMP}/expected_{template_id}.txt"

    # upstream check_line_number counts lines matching HH:MM:SS regex, not total lines.
    # Prepend a synthetic timestamp to each data row so line_count == spec["rows"].
    header = f"Linux 5.15.0 (cua-vm)    05/12/2026     _x86_64_    (4 CPU)\n\nTime     {spec['header_cols']}\n"
    body_rows = "\n".join(f"12:{(i // 60) % 60:02d}:{i % 60:02d} {spec['row_template']}" for i in range(spec["rows"]))
    gold_text = header + body_rows + "\n"

    init_steps = [
        _write_text_step(expected_path, gold_text),
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [_execute(f"cp '{expected_path}' '{out_path}'")]

    # check_line_number DROPPED: it counts HH:MM:SS-matching lines and requires that to
    # EXACTLY equal spec["line_count"], but a real sysstat run can never match — sar emits
    # a timestamped column header (rows+1 matches), mpstat -P ALL emits one row per CPU,
    # and vmstat/iostat emit no HH:MM:SS lines at all. The synthetic gold only matched
    # because its header is the literal word "Time". file_contains (the column-header
    # tokens) still verifies the agent ran the right tool; the exact sample count is
    # unverifiable across hosts/tools, so requiring it made the whole family unsolvable.
    evaluator = {
        "func": ["file_contains"],
        "result": [
            {"type": "vm_file", "path": out_path, "dest": "report.txt"},
        ],
        "expected": [
            {"type": "rule", "rules": {"expected": spec["substrings"]}},
        ],
    }
    instructions = [
        f"Monitor Ubuntu system resource usage with `{spec['tool']}` from the sysstat toolkit (or its analog). Save the collected output to '~/Desktop/{spec['out']}' so the file contains the standard column headers and one data line per sample.",
        f"From a terminal, capture the output of `{spec['tool']}` to ~/Desktop/{spec['out']}. The resulting file must include all the expected column header tokens and have the right number of total lines for the sampling interval.",
        f"Run the sysstat sampling tool `{spec['tool']}` and redirect its output into a report at ~/Desktop/{spec['out']} (the report should retain the header line and one entry per sampling tick).",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=3,
        eval_class="file_contains",
        setup_class="eval_emul_sysstat_report",
    )


# ---- Cat U.3 — compare_table^3 (multi-sheet timetable) --------------------
# emulates osworld_multi_apps_3a93cae4: agent edits a multi-sheet xlsx
# (e.g. weekly course timetable) to add/edit cells; eval = compare_table^3
# (three sheets, each checked with sheet_data rules). Apps: nautilus + Calc.
_TIMETABLE_SPECS: list[dict] = [
    {
        "id": "multi_eval_timetable_wed_lecture",
        "sheets": ["Mon", "Tue", "Wed"],
        "init_data": {
            "Mon": [["Time","Slot"], ["09:00","Algebra"], ["11:00","Physics"]],
            "Tue": [["Time","Slot"], ["10:00","Chemistry"], ["13:00","English"]],
            "Wed": [["Time","Slot"], ["09:00","Biology"], ["14:00","History"]],
        },
        "gold_data": {
            "Mon": [["Time","Slot"], ["09:00","Algebra"], ["11:00","Physics"]],
            "Tue": [["Time","Slot"], ["10:00","Chemistry"], ["13:00","English"]],
            "Wed": [["Time","Slot"], ["09:00","Biology"], ["12:00","Lecture"], ["14:00","History"]],
        },
    },
    {
        "id": "multi_eval_timetable_fri_seminar",
        "sheets": ["Thu", "Fri", "Sat"],
        "init_data": {
            "Thu": [["Time","Slot"], ["10:00","Workshop"], ["15:00","Lab"]],
            "Fri": [["Time","Slot"], ["09:00","Research"], ["13:00","Meeting"]],
            "Sat": [["Time","Slot"], ["10:00","Office hours"]],
        },
        "gold_data": {
            "Thu": [["Time","Slot"], ["10:00","Workshop"], ["15:00","Lab"]],
            "Fri": [["Time","Slot"], ["09:00","Research"], ["11:00","Seminar"], ["13:00","Meeting"]],
            "Sat": [["Time","Slot"], ["10:00","Office hours"]],
        },
    },
    {
        "id": "multi_eval_timetable_thu_tutoring",
        "sheets": ["Mon", "Wed", "Thu"],
        "init_data": {
            "Mon": [["Time","Slot"], ["09:00","Geometry"]],
            "Wed": [["Time","Slot"], ["10:00","Lab"], ["14:00","Reading"]],
            "Thu": [["Time","Slot"], ["09:30","Calculus"], ["15:00","Statistics"]],
        },
        "gold_data": {
            "Mon": [["Time","Slot"], ["09:00","Geometry"]],
            "Wed": [["Time","Slot"], ["10:00","Lab"], ["14:00","Reading"]],
            "Thu": [["Time","Slot"], ["09:30","Calculus"], ["13:00","Tutoring"], ["15:00","Statistics"]],
        },
    },
]


def _make_timetable_template(spec: dict) -> SynthTemplate:
    template_id   = spec["id"]
    xlsx_path     = f"{_DESKTOP}/timetable.xlsx"
    expected_xlsx = f"{_TMP}/expected_{template_id}.xlsx"

    def _build_xlsx_multisheet_step(path: str, data_by_sheet: dict[str, list[list]]) -> dict:
        py = textwrap.dedent(f"""\
            import os
            from openpyxl import Workbook
            path = {path!r}
            data = {data_by_sheet!r}
            sheets = {spec["sheets"]!r}
            os.makedirs(os.path.dirname(path) or '.', exist_ok=True)
            wb = Workbook(); wb.remove(wb.active)
            for s in sheets:
                ws = wb.create_sheet(title=s)
                for row in data[s]:
                    ws.append(row)
            wb.save(path)
            import subprocess as _sp, tempfile as _tf, shutil as _sh
            _td = _tf.mkdtemp()
            try:
                _sp.run(['soffice','--headless','--norestore','--nofirststartwizard',
                         '--convert-to','xlsx','--outdir',_td,path],
                        capture_output=True, env={{**os.environ, 'DISPLAY':':1'}}, timeout=120)
                c = os.path.join(_td, os.path.basename(path))
                if os.path.exists(c): _sh.copy(c, path)
            finally:
                _sh.rmtree(_td, ignore_errors=True)
            """)
        return _python_step(py)

    init_steps = [
        _build_xlsx_multisheet_step(xlsx_path, spec["init_data"]),
        _build_xlsx_multisheet_step(expected_xlsx, spec["gold_data"]),
        *_multi_app_preopen(foreground=xlsx_path),
    ]
    oracle_steps = [_execute(f"cp '{expected_xlsx}' '{xlsx_path}'")]

    evaluator = {
        "func": ["compare_table"] * 3,
        "conj": "and",
        "result": [
            {"type": "vm_file", "path": xlsx_path, "dest": f"result_{i}.xlsx"}
            for i in range(3)
        ],
        "expected": [
            {"type": "vm_file", "path": expected_xlsx, "dest": f"expected_{i}.xlsx"}
            for i in range(3)
        ],
        "options": [
            {"rules": [{"type": "sheet_data", "sheet_idx0": i, "sheet_idx1": f"EI{i}"}]}
            for i in range(3)
        ],
    }
    diff_day = [s for s in spec["sheets"]
                if spec["init_data"][s] != spec["gold_data"][s]][0]
    diff_rows = [r for r in spec["gold_data"][diff_day]
                 if r not in spec["init_data"][diff_day]]
    diff_summary = ", ".join(" ".join(r) for r in diff_rows)
    instructions = [
        f"Open ~/Desktop/timetable.xlsx in LibreOffice Calc. The workbook has sheets {spec['sheets']!r}; please add the missing entry to the '{diff_day}' sheet so it includes: {diff_summary}. Leave the other sheets untouched.",
        f"My weekly timetable at ~/Desktop/timetable.xlsx is missing one slot on the '{diff_day}' tab — could you append the row [{diff_summary}] in time order on that sheet only? Don't change anything on the other sheets.",
        f"Could you add a missing class slot to my course timetable at ~/Desktop/timetable.xlsx? It belongs on the '{diff_day}' sheet — please insert [{diff_summary}] in chronological order without touching the other day tabs.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", xlsx_path],
            "oracle_after_postconfig": True,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: LO_SAVE_POSTCONFIG,
        param_fn=_params,
        n_rows=3,
        eval_class="compare_table+compare_table+compare_table",
        setup_class="eval_emul_timetable_multisheet",
    )


# ---- Cat U.4 — is_extension_installed+check_image_size --------------------
# emulates osworld_multi_apps_42f4d1c7: install a VS Code extension and
# resize an image via a script-runner. Apps: VS Code + GIMP/Python.
_VSCODE_EXT_IMG_SPECS: list[dict] = [
    {
        "id":   "multi_eval_vscode_lisp_resize_128",
        "ext":  "mattn.lisp",
        "src_basename": "character.png",
        "dst_basename": "resized.png",
        "width":  128, "height": 128,
        "asset_rel": "photos/wildlife/tiger-closeup.jpg",
    },
    {
        "id":   "multi_eval_vscode_yaml_resize_256",
        "ext":  "redhat.vscode-yaml",
        "src_basename": "logo_src.png",
        "dst_basename": "logo_thumb.png",
        "width":  256, "height": 256,
        "asset_rel": "photos/city/asian-skyline.jpg",
    },
]


def _make_vscode_ext_img_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    src_path    = f"{_DESKTOP}/{spec['src_basename']}"
    dst_path    = f"{_DESKTOP}/{spec['dst_basename']}"
    expected_png = f"{_TMP}/expected_{template_id}.png"

    init_steps = [
        _stage_asset(spec["asset_rel"], f"{_TMP}/{template_id}_src.jpg"),
        _execute(f"convert '{_TMP}/{template_id}_src.jpg' '{src_path}'"),
        _python_step(textwrap.dedent(f"""\
            from PIL import Image
            img = Image.open({src_path!r})
            img.resize(({spec['width']}, {spec['height']}), Image.LANCZOS).save({expected_png!r})
            """)),
    ]
    oracle_steps = [
        _execute(f"code --no-sandbox --install-extension {spec['ext']} --force 2>/dev/null || true"),
        _execute(f"cp '{expected_png}' '{dst_path}'"),
    ]

    evaluator = {
        "func": ["is_extension_installed", "check_image_size"],
        "result": [
            {"type": "vm_command_line", "command": ["code", "--list-extensions", "|", "grep", spec["ext"]]},
            {"type": "vm_file", "path": dst_path, "dest": spec["dst_basename"]},
        ],
        "expected": [
            {"type": "rule", "rules": {"type": "contain", "expected": spec["ext"]}},
            {"type": "rule", "rules": {"height": spec["height"], "width": spec["width"]}},
        ],
    }
    instructions = [
        f"Install the VS Code extension `{spec['ext']}` so its commands are available, then write a short script that resizes ~/Desktop/{spec['src_basename']} to {spec['width']}x{spec['height']} pixels and saves the output as ~/Desktop/{spec['dst_basename']}.",
        f"Set up the `{spec['ext']}` extension in VS Code, then use any script-runner to scale the image at ~/Desktop/{spec['src_basename']} down to a {spec['width']}x{spec['height']}-px PNG saved as ~/Desktop/{spec['dst_basename']}.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="is_extension_installed+check_image_size",
        setup_class="eval_emul_vscode_ext_image_resize",
    )


# ---- Cat U.5 — check_list+is_expected_tabs (workspace setup) --------------
# emulates osworld_multi_apps_48c46dc7: agent opens a project folder in
# terminal AND file manager, AND opens documentation URLs in Chrome.
# Eval: check_list (running processes contain terminal+nautilus) +
# is_expected_tabs (Chrome has the right URLs open). Apps: 3+ (chrome,
# terminal, nautilus).
_WORKSPACE_SETUP_SPECS: list[dict] = [
    {
        "id":   "multi_eval_workspace_osworld_python",
        "proj_dir": "/home/user/Documents/Projects/OSWorld",
        "tabs": ["https://github.com/", "https://docs.python.org/3/"],
        "tab_keywords": ["github", "python"],
    },
    # kernel.org root 200 directly; wiki.archlinux.org/ 301-redirects to
    # /title/Main_page — use the post-redirect URL so chrome's address-bar
    # tab URL matches the eval's `is_expected_tabs` strict path compare
    # (compare_urls normalizes domain but keeps path).
    {
        "id":   "multi_eval_workspace_kernel_arch",
        "proj_dir": "/home/user/Documents/Projects/Kernel",
        "tabs": ["https://www.kernel.org/", "https://wiki.archlinux.org/title/Main_page"],
        "tab_keywords": ["kernel", "archlinux"],
    },
]


def _make_workspace_setup_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    proj_dir    = spec["proj_dir"]
    list_path   = f"{_TMP}/process_list_{template_id}.txt"
    expected_list = f"{_TMP}/expected_{template_id}.list"

    init_steps = [
        _execute(f"mkdir -p '{proj_dir}'"),
        _write_text_step(expected_list, "\n".join(["bash", "nautilus"]) + "\n"),
        # socat for chrome devtools (mirrors eval).
        {"type": "launch", "parameters": {"command": ["socat", "tcp-listen:9222,fork", "tcp:localhost:1337"]}},
    ]
    # Oracle: launch terminal in proj_dir, nautilus on proj_dir, chrome with
    # the documentation tabs, then dump running-process list for `check_list`.
    chrome_cmd = ["google-chrome", "--remote-debugging-port=1337", "--no-first-run",
                  "--no-default-browser-check"] + list(spec["tabs"])
    oracle_steps = [
        _execute(f"cd '{proj_dir}' && gnome-terminal --working-directory='{proj_dir}' &"),
        {"type": "sleep", "parameters": {"seconds": 2}},
        {"type": "launch", "parameters": {"command": ["nautilus", proj_dir]}},
        {"type": "sleep", "parameters": {"seconds": 2}},
        {"type": "launch", "parameters": {"command": chrome_cmd}},
        {"type": "sleep", "parameters": {"seconds": 4}},
        _execute(f"ps -A -o comm | sort -u > '{list_path}'"),
    ]

    evaluator = {
        "func": ["check_list", "is_expected_tabs"],
        "result": [
            {"type": "vm_file", "path": list_path, "dest": "process_list.txt"},
            {"type": "open_tabs_info"},
        ],
        "expected": [
            {"type": "rule", "rules": {"expect": ["bash", "nautilus"]}},
            {"type": "rule", "rules": {"type": "url", "urls": spec["tabs"]}},
        ],
    }
    tab_summary = " and ".join(spec["tabs"])
    instructions = [
        f"Help me set up my workspace automatically: open the project directory {proj_dir} in both a terminal and the file manager, and open {tab_summary} in two separate Chrome tabs.",
        f"Bootstrap my dev environment: launch a terminal and a file-manager window both pointed at {proj_dir}, then start Chrome with the reference tabs {tab_summary}.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="check_list+is_expected_tabs",
        setup_class="eval_emul_workspace_setup",
    )


# ---- Cat U.6 — check_include_exclude+compare_config (VS Code from term) ---
# emulates osworld_multi_apps_510f64c8: start VS Code in a folder FROM the
# terminal. check_include_exclude probes ~/.bash_history for `code <dir>`;
# compare_config reads the recorded extension state.
_VSCODE_OPEN_FROM_TERM_SPECS: list[dict] = [
    {
        "id":   "multi_eval_vscode_open_project_from_term",
        "proj_dir": "/home/user/Desktop/project",
        "history_token": "code ",
        "config_value": "project",
    },
    {
        "id":   "multi_eval_vscode_open_demo_from_term",
        "proj_dir": "/home/user/Desktop/demo_app",
        "history_token": "code ",
        "config_value": "demo_app",
    },
]


def _make_vscode_open_from_term_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    proj_dir    = spec["proj_dir"]
    hist_path   = "/home/user/.bash_history"
    probe_out   = f"{_TMP}/probe_{template_id}.txt"
    config_out  = "/home/user/OpenProject.txt"

    init_steps = [
        _execute(f"mkdir -p '{proj_dir}/.vscode'"),
        _write_text_step(f"{proj_dir}/main.py", "print('hello')\n"),
        _write_text_step(f"{proj_dir}/README.md", f"# {spec['config_value']}\n"),
        _execute("history -c 2>/dev/null || true; echo > ~/.bash_history"),
        *_terminal_preopen_steps(),
    ]
    # Oracle: simulate the user running `code <proj_dir>` from a terminal,
    # plant the bash_history entry, plant the recorded OpenProject path.
    oracle_steps = [
        _execute(f"echo 'code {proj_dir}' >> '{hist_path}'"),
        _execute(
            "if cat ~/.bash_history | grep -q '[c]ode '; "
            f"then echo true > '{probe_out}'; else echo false > '{probe_out}'; fi"
        ),
        _write_text_step(config_out, spec["config_value"] + "\n"),
    ]

    evaluator = {
        "func": ["check_include_exclude", "compare_config"],
        "result": [
            {"type": "vm_command_line", "command": f"cat {probe_out}", "shell": True},
            {"type": "vm_file", "path": config_out, "dest": "OpenProject.txt"},
        ],
        "expected": [
            {"type": "rule", "rules": {"include": ["true"]}},
            {"type": "rule", "rules": {"expected": spec["config_value"]}},
        ],
    }
    instructions = [
        f"Could you start VS Code on the folder {proj_dir} FROM a terminal (so the launch is recorded in shell history)? The opened workspace name must match the folder basename.",
        f"Open the folder {proj_dir} as a VS Code workspace by running `code {proj_dir}` from an interactive shell — the invocation needs to be visible in ~/.bash_history afterwards.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="check_include_exclude+compare_config",
        setup_class="eval_emul_vscode_from_terminal",
    )


# ---- Cat U.7 — compare_pdfs+compare_docx_files (paper + summary) ----------
# emulates osworld_multi_apps_68a25bd4: agent extracts the first paper URL
# from an xlsx, downloads its PDF, AND writes a citing-paper title to a docx.
# Apps: Calc (open xlsx) + Chrome (download PDF) + Writer (save docx). 3-app.
_PAPER_PDF_DOCX_SPECS: list[dict] = [
    {
        "id":   "multi_eval_paper_pdf_citation_docx",
        "pdf_basename":  "paper01.pdf",
        "docx_basename": "ans.docx",
        "xlsx_basename": "papers.xlsx",
        "citing_title": "BERT: Pre-training of Deep Bidirectional Transformers for Language Understanding",
        "rows": [
            ["Title", "URL"],
            ["Attention Is All You Need", "https://arxiv.org/pdf/1706.03762"],
            ["BERT: Pre-training of Deep Bidirectional Transformers", "https://arxiv.org/pdf/1810.04805"],
        ],
    },
    {
        "id":   "multi_eval_paper_pdf_citation_docx_glue",
        "pdf_basename":  "paper01.pdf",
        "docx_basename": "ans.docx",
        "xlsx_basename": "ml-papers.xlsx",
        "citing_title": "GLUE: A Multi-Task Benchmark and Analysis Platform for Natural Language Understanding",
        "rows": [
            ["Title", "URL"],
            ["ELMo: Deep contextualized word representations", "https://arxiv.org/pdf/1802.05365"],
            ["GLUE: A Multi-Task Benchmark and Analysis Platform", "https://arxiv.org/pdf/1804.07461"],
        ],
    },
]


def _make_paper_pdf_docx_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    user_home   = "/home/user"
    pdf_path    = f"{user_home}/{spec['pdf_basename']}"
    docx_path   = f"{user_home}/{spec['docx_basename']}"
    xlsx_path   = f"{_DESKTOP}/{spec['xlsx_basename']}"
    expected_pdf = f"{_TMP}/expected_{template_id}.pdf"
    expected_docx = f"{_TMP}/expected_{template_id}.docx"

    # Build a 2-row xlsx (Title, URL) + a small gold PDF (txt → pdf via
    # pandoc) + a single-paragraph gold docx with the citing title.
    init_steps = [
        _build_xlsx_step(xlsx_path, spec["rows"]),
        # Use pandoc with pdf engine OR fall back to soffice. Use soffice
        # for determinism (already used by _build_oracle_lo).
        _python_step(textwrap.dedent(f"""\
            import os, subprocess, tempfile
            tmp_txt = tempfile.mktemp(suffix='.txt')
            with open(tmp_txt,'w') as f:
                f.write('Attention Is All You Need\\n\\nAbstract: a transformer paper.')
            tmp_dir = tempfile.mkdtemp()
            subprocess.run(['soffice','--headless','--norestore','--nofirststartwizard',
                            '--convert-to','pdf','--outdir',tmp_dir,tmp_txt],
                           check=True, env={{**os.environ,'DISPLAY':':1'}}, timeout=120)
            import shutil
            shutil.move(os.path.join(tmp_dir, os.path.basename(tmp_txt).replace('.txt','.pdf')),
                        {expected_pdf!r})
            """)),
        _build_docx_step(expected_docx, [(None, spec["citing_title"])]),
        *_multi_app_preopen(foreground=xlsx_path),
    ]
    oracle_steps = [
        _execute(f"cp '{expected_pdf}' '{pdf_path}'"),
        *_build_oracle_lo(docx_path, expected_docx, fmt="docx"),
    ]

    evaluator = {
        "func": ["compare_pdfs", "compare_docx_files"],
        "result": [
            {"type": "vm_file", "path": pdf_path,  "dest": spec["pdf_basename"]},
            {"type": "vm_file", "path": docx_path, "dest": spec["docx_basename"]},
        ],
        "expected": [
            {"type": "vm_file", "path": expected_pdf,  "dest": f"gold_{spec['pdf_basename']}"},
            {"type": "vm_file", "path": expected_docx, "dest": f"gold_{spec['docx_basename']}"},
        ],
    }
    first_title = spec["rows"][1][0]
    instructions = [
        f"I've compiled a list of papers with links in ~/Desktop/{spec['xlsx_basename']}. Please download the PDF of the first paper ('{first_title}') and save it as ~/{spec['pdf_basename']}. Also, identify which paper in the list cites the first one and record that paper's title in ~/{spec['docx_basename']} (one paragraph, the title only).",
        f"Using the bibliography spreadsheet at ~/Desktop/{spec['xlsx_basename']}, fetch the PDF of the first entry to ~/{spec['pdf_basename']}; then determine which paper from the table cites that first paper and write its title to ~/{spec['docx_basename']}.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", xlsx_path],
            "oracle_after_postconfig": False,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="compare_pdfs+compare_docx_files",
        setup_class="eval_emul_paper_pdf_docx",
    )


# ---- Cat U.8 — check_json+diff_text_file^4 (web-ext scaffold) -------------
# emulates osworld_multi_apps_74d5859f: agent scaffolds a browser web
# extension (manifest.json + background_script.js + browserAction/
# {index.html, style.css, script.js}); eval = check_json (manifest schema
# keys) + diff_text_file^4 (the 4 source files byte-equal to gold). The
# 4 diff_text_file atoms check 4 DIFFERENT files, so they're semantically
# distinct (anti-hacking principle satisfied).
def _make_webext_scaffold_template() -> SynthTemplate:
    template_id  = "multi_eval_webext_scaffold_happy_extension"
    proj_dir     = "/home/user/Projects/happy-extension"
    files = {
        "manifest.json":                       "manifest",
        "background_script.js":                "bg",
        "browserAction/index.html":            "idx",
        "browserAction/style.css":             "css",
        "browserAction/script.js":             "js",
    }
    gold_bg  = (
        "// background_script.js — happy-extension v0.0.1\n"
        "chrome.runtime.onInstalled.addListener(() => {\n"
        "  console.log('happy-extension installed.');\n"
        "});\n"
    )
    gold_idx = (
        "<!doctype html>\n"
        "<html><head><meta charset=\"utf-8\"><link rel=\"stylesheet\" href=\"style.css\"></head>\n"
        "<body><h1>happy-extension</h1><script src=\"script.js\"></script></body></html>\n"
    )
    gold_css = (
        "body { font-family: sans-serif; padding: 8px; }\n"
        "h1 { color: #2a7; }\n"
    )
    gold_js  = (
        "// script.js — popup logic for happy-extension.\n"
        "document.addEventListener('DOMContentLoaded', () => {\n"
        "  document.querySelector('h1').addEventListener('click', () => alert('hello'));\n"
        "});\n"
    )
    manifest_json = textwrap.dedent("""\
        {
          "manifest_version": 2,
          "name": "happy-extension",
          "version": "0.0.1",
          "description": "",
          "background": { "scripts": ["background_script.js"], "persistent": false },
          "browser_action": {
            "default_icon": { "64": "icons/icon.png" },
            "default_popup": "browserAction/index.html",
            "default_title": "happy-extension"
          }
        }
        """)
    # Gold copies in /tmp.
    exp = {k: f"{_TMP}/expected_{template_id}_{k.replace('/', '_')}" for k in files}

    init_steps = [
        _execute(f"mkdir -p '{proj_dir}/browserAction' '{proj_dir}/icons'"),
        _write_text_step(exp["manifest.json"], manifest_json),
        _write_text_step(exp["background_script.js"], gold_bg),
        _write_text_step(exp["browserAction/index.html"], gold_idx),
        _write_text_step(exp["browserAction/style.css"], gold_css),
        _write_text_step(exp["browserAction/script.js"], gold_js),
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [
        _execute(f"cp '{exp[k]}' '{proj_dir}/{k}'") for k in files
    ]

    evaluator = {
        "func": ["check_json", "diff_text_file", "diff_text_file", "diff_text_file", "diff_text_file"],
        "result": [
            {"type": "vm_file", "path": f"{proj_dir}/manifest.json",                "dest": "manifest.json"},
            {"type": "vm_file", "path": f"{proj_dir}/background_script.js",         "dest": "background_script.js"},
            {"type": "vm_file", "path": f"{proj_dir}/browserAction/index.html",     "dest": "index.html"},
            {"type": "vm_file", "path": f"{proj_dir}/browserAction/style.css",      "dest": "style.css"},
            {"type": "vm_file", "path": f"{proj_dir}/browserAction/script.js",      "dest": "script.js"},
        ],
        "expected": [
            {"type": "rule", "rules": {"expect": [
                {"key": ["name"],                          "method": "eq", "ref": "happy-extension"},
                {"key": ["version"],                       "method": "eq", "ref": "0.0.1"},
                {"key": ["background", "scripts"],         "method": "eq", "ref": ["background_script.js"]},
                {"key": ["browser_action", "default_icon"],"method": "eq", "ref": {"64": "icons/icon.png"}},
                {"key": ["browser_action", "default_popup"],"method": "eq", "ref": "browserAction/index.html"},
                {"key": ["browser_action", "default_title"],"method": "eq", "ref": "happy-extension"},
            ]}},
            {"type": "vm_file", "path": exp["background_script.js"],     "dest": "gold_bg.js"},
            {"type": "vm_file", "path": exp["browserAction/index.html"], "dest": "gold_idx.html"},
            {"type": "vm_file", "path": exp["browserAction/style.css"],  "dest": "gold_css.css"},
            {"type": "vm_file", "path": exp["browserAction/script.js"],  "dest": "gold_script.js"},
        ],
    }
    instructions = [
        "Scaffold a new browser web extension at ~/Projects/happy-extension/, tagged 'happy-extension v0.0.1' with an empty description. Include a background script (background_script.js) and a browser action with default_popup browserAction/index.html, default_title 'happy-extension', and default_icon {64: 'icons/icon.png'}; populate browserAction/ with index.html, style.css and script.js.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_json+diff_text_file+diff_text_file+diff_text_file+diff_text_file",
        setup_class="eval_emul_webext_scaffold",
    )


# ---- Cat U.9 — exact_match^3 (Desktop file triage) ------------------------
# emulates osworld_multi_apps_869de13e: 3-way file categorisation via
# `ls` on three Desktop subfolders. Each atom checks a different folder's
# listing, distinct file lists.
_DESK_TRIAGE_SPECS: list[dict] = [
    {
        "id": "multi_eval_desktop_triage_papers_projects_misc",
        "groups": {
            "Paper_reading": ["1706.03762.pdf", "1802.05365.pdf", "GLUE_benchmark.pdf"],
            "Projects":      ["flask-app", "ml-pipeline"],
            "Miscellaneous": ["budget-2024.xlsx", "meeting-notes.docx", "vacation.jpg"],
        },
    },
    {
        "id": "multi_eval_desktop_triage_books_code_other",
        "groups": {
            "Books":     ["alice-in-wonderland.epub", "moby-dick.epub", "war-of-the-worlds.epub"],
            "Code":      ["fizzbuzz.py", "merge-sort.cpp"],
            "Receipts":  ["amazon-2024-01.pdf", "office-supplies.pdf"],
        },
    },
]


def _make_desk_triage_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    groups = spec["groups"]
    group_names = list(groups.keys())
    folder_paths = {g: f"{_DESKTOP}/{g}" for g in group_names}

    init_steps = []
    for g, files in groups.items():
        init_steps.append(_execute(f"mkdir -p '{folder_paths[g]}'"))
        # Stage files as empty stubs in Desktop root (so agent must move them).
        for f in files:
            init_steps.append(_write_text_step(f"{_DESKTOP}/{f}", ""))
    init_steps.append({"type": "launch", "parameters": {"command": ["nautilus", _DESKTOP]}})

    oracle_steps = []
    for g, files in groups.items():
        for f in files:
            oracle_steps.append(_execute(f"mv '{_DESKTOP}/{f}' '{folder_paths[g]}/{f}'"))

    evaluator = {
        "func": ["exact_match"] * len(group_names),
        "result": [
            {"type": "vm_command_line", "command": ["ls", folder_paths[g]]}
            for g in group_names
        ],
        "expected": [
            {"type": "rule", "rules": {"expected": "\n".join(sorted(groups[g])) + "\n"}}
            for g in group_names
        ],
    }
    summary = "; ".join(f"{g}/ for {', '.join(files)}" for g, files in groups.items())
    instructions = [
        f"Tidy up my Desktop by sorting files into category folders: {summary}. Create the folders if needed and move every listed file into its category folder, leaving the Desktop root clean of those items.",
        f"Organise the loose items on the Desktop into the following category folders (create them if missing) and move each file accordingly: {summary}. After the move the Desktop root must no longer contain those items.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="exact_match+exact_match+exact_match",
        setup_class="eval_emul_desktop_triage",
    )


# ---- Cat U.10 — compare_csv+compare_table (contacts export+convert) -------
# emulates osworld_multi_apps_c867c42d: export contacts to csv, then save
# as xlsx in Calc. Apps: Thunderbird-style address book stub + Calc.
_CONTACTS_EXPORT_SPECS: list[dict] = [
    {
        "id":   "multi_eval_contacts_export_personal",
        "csv_basename":  "contacts.csv",
        "xlsx_basename": "contacts.xlsx",
        "rows": [
            ["First Name", "Last Name", "Email", "Phone"],
            ["Alice",      "Nguyen",    "alice.nguyen@example.com",  "+1-555-0101"],
            ["Bob",        "Patel",     "bob.patel@example.com",     "+1-555-0102"],
            ["Carol",      "Sato",      "carol.sato@example.com",    "+1-555-0103"],
            ["David",      "Klein",     "david.klein@example.com",   "+1-555-0104"],
        ],
    },
    {
        "id":   "multi_eval_contacts_export_work",
        "csv_basename":  "work_contacts.csv",
        "xlsx_basename": "work_contacts.xlsx",
        "rows": [
            ["First Name", "Last Name", "Email", "Phone"],
            ["Sergey",     "Ivanov",    "sergey@example.org", "+44-20-7946-0001"],
            ["Yuki",       "Tanaka",    "yuki@example.org",   "+81-3-1234-5678"],
            ["Pedro",      "Almeida",   "pedro@example.org",  "+351-21-987-6543"],
        ],
    },
]


def _make_contacts_export_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    csv_path     = f"{_DESKTOP}/{spec['csv_basename']}"
    xlsx_path    = f"{_DESKTOP}/{spec['xlsx_basename']}"
    expected_csv  = f"{_TMP}/expected_{template_id}.csv"
    expected_xlsx = f"{_TMP}/expected_{template_id}.xlsx"

    init_steps = [
        _write_text_step(expected_csv, _csv_text([[str(c) for c in r] for r in spec["rows"]])),
        _build_xlsx_step(expected_xlsx, spec["rows"]),
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [
        _execute(f"cp '{expected_csv}' '{csv_path}'"),
        _execute(f"cp '{expected_xlsx}' '{xlsx_path}'"),
    ]

    evaluator = {
        "func": ["compare_csv", "compare_table"],
        "conj": "and",
        "result": [
            {"type": "vm_file", "path": csv_path,  "dest": spec["csv_basename"]},
            {"type": "vm_file", "path": xlsx_path, "dest": spec["xlsx_basename"]},
        ],
        "expected": [
            {"type": "vm_file", "path": expected_csv,  "dest": f"gold_{spec['csv_basename']}"},
            {"type": "vm_file", "path": expected_xlsx, "dest": f"gold_{spec['xlsx_basename']}"},
        ],
        "options": [
            {},
            {"rules": [{"type": "sheet_data", "sheet_idx0": "RI0", "sheet_idx1": "EI0"}]},
        ],
    }
    instructions = [
        f"Export the entries of my Personal Address Book to a CSV at ~/Desktop/{spec['csv_basename']} (one header row + one row per contact with columns First Name, Last Name, Email, Phone), then open that CSV in Calc and save a copy as ~/Desktop/{spec['xlsx_basename']}.",
        f"Dump the Thunderbird Personal Address Book to ~/Desktop/{spec['csv_basename']} as a 4-column CSV (First Name, Last Name, Email, Phone), and then convert that CSV to xlsx via LibreOffice Calc at ~/Desktop/{spec['xlsx_basename']}.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="compare_csv+compare_table",
        setup_class="eval_emul_contacts_export",
    )


# ---- Cat U.11 — is_expected_tabs+compare_htmls (xlsx → html → chrome) -----
# emulates osworld_multi_apps_e135df7c: convert xlsx to html with Calc and
# open it in Chrome alongside other tabs. Apps: 3+ (Calc + Chrome + tabs).
_XLSX_TO_HTML_SPECS: list[dict] = [
    {
        "id":   "multi_eval_xlsx_html_enterprise_survey",
        "xlsx_basename": "annual-survey.xlsx",
        "html_basename": "annual-survey.html",
        "tabs": ["https://aclanthology.org/", "https://openai.com/", "https://www.linkedin.com/home/"],
        "rows": [
            ["Year", "Sector",         "RevenueM"],
            ["2021", "Manufacturing",  "120"],
            ["2021", "Construction",   "85"],
            ["2021", "Retail",         "150"],
        ],
    },
    {
        "id":   "multi_eval_xlsx_html_quarterly_summary",
        "xlsx_basename": "quarterly-summary.xlsx",
        "html_basename": "quarterly-summary.html",
        "tabs": ["https://news.ycombinator.com/", "https://arxiv.org/", "https://stackoverflow.com/"],
        "rows": [
            ["Quarter", "Revenue", "Profit"],
            ["Q1",      "1200",    "180"],
            ["Q2",      "1350",    "210"],
            ["Q3",      "1420",    "240"],
        ],
    },
]


def _make_xlsx_to_html_template(spec: dict) -> SynthTemplate:
    template_id   = spec["id"]
    xlsx_path     = f"{_DESKTOP}/{spec['xlsx_basename']}"
    html_path     = f"{_DESKTOP}/{spec['html_basename']}"
    expected_html = f"{_TMP}/expected_{template_id}.html"
    expected_xlsx = f"{_TMP}/expected_{template_id}.xlsx"
    file_url = f"file://{html_path}"

    init_steps = [
        _build_xlsx_step(expected_xlsx, spec["rows"]),
        _execute(f"cp '{expected_xlsx}' '{xlsx_path}'"),
        # Build a gold HTML the way LibreOffice Calc would: a single
        # <table> serialization mirroring Calc's "Save as HTML" output.
        _python_step(textwrap.dedent(f"""\
            import os
            rows = {spec["rows"]!r}
            parts = ['<html><body><table>']
            for r in rows:
                parts.append('<tr>' + ''.join('<td>'+str(c)+'</td>' for c in r) + '</tr>')
            parts.append('</table></body></html>')
            os.makedirs(os.path.dirname({expected_html!r}) or '.', exist_ok=True)
            with open({expected_html!r}, 'w') as f:
                f.write('\\n'.join(parts) + '\\n')
            """)),
        # Pre-open Chrome with the BG tabs (mirrors eval task expectations).
        {"type": "launch", "parameters": {"command": ["socat", "tcp-listen:9222,fork", "tcp:localhost:1337"]}},
        *_multi_app_preopen(foreground=xlsx_path),
    ]
    # Oracle: plant html, then open chrome with all tabs (3 BG + the file URL).
    chrome_cmd = ["google-chrome", "--remote-debugging-port=1337", "--no-first-run",
                  "--no-default-browser-check"] + list(spec["tabs"]) + [file_url]
    oracle_steps = [
        _execute(f"cp '{expected_html}' '{html_path}'"),
        {"type": "launch", "parameters": {"command": chrome_cmd}},
        {"type": "sleep", "parameters": {"seconds": 4}},
    ]

    evaluator = {
        "func": ["is_expected_tabs", "compare_htmls"],
        "result": [
            {"type": "open_tabs_info"},
            {"type": "vm_file", "path": html_path, "dest": spec["html_basename"]},
        ],
        "expected": [
            {"type": "rule", "rules": {"type": "url", "urls": list(spec["tabs"]) + [file_url]}},
            {"type": "vm_file", "path": expected_html, "dest": f"gold_{spec['html_basename']}"},
        ],
        "options": [{}, {"ignore_sdnum": True}],
    }
    tab_summary = ", ".join(spec["tabs"])
    instructions = [
        f"Open ~/Desktop/{spec['xlsx_basename']} in LibreOffice Calc, save (or export) it as an HTML file at ~/Desktop/{spec['html_basename']}, then view that HTML in Chrome alongside the reference tabs {tab_summary}.",
        f"Export the spreadsheet ~/Desktop/{spec['xlsx_basename']} to HTML at ~/Desktop/{spec['html_basename']} (using Calc), then load that file URL in Chrome together with the other reference tabs: {tab_summary}.",
    ]

    def _params(seed: int) -> dict:
        return {
            "instr": instructions[seed % len(instructions)],
            "pre_config_steps": init_steps,
            "open_command": ["libreoffice", "--calc", xlsx_path],
            "oracle_after_postconfig": False,
        }

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="is_expected_tabs+compare_htmls",
        setup_class="eval_emul_xlsx_to_html_chrome",
    )


# ---- Cat U.12 — check_list+check_list (install ext + apt package) ---------
# emulates osworld_multi_apps_e1fc0df3: install a LibreOffice extension AND
# its Java runtime dependency. Atom 1 = check_list on `unopkg list` output
# matching the .oxt regex; atom 2 = check_list on `dpkg -l | grep openjdk`.
# Two distinct atoms — same fn, different probe targets, different regexes.
def _make_lo_extension_install_template() -> SynthTemplate:
    template_id  = "multi_eval_lo_extension_languagetool_jdk"
    grep_out     = f"{_TMP}/grep_{template_id}.out"
    apt_out      = f"{_TMP}/apt_{template_id}.out"

    init_steps = [
        *_terminal_preopen_steps(),
    ]
    oracle_steps = [
        _write_text_step(grep_out, "org.openoffice.languagetool.oxt 5.9 user\n"),
        _write_text_step(apt_out,  "ii  openjdk-21-jre  21.0.4+7-1ubuntu3  amd64  OpenJDK Java runtime\n"),
    ]

    evaluator = {
        "func": ["check_list", "check_list"],
        "result": [
            {"type": "vm_file", "path": grep_out, "dest": "grep.out"},
            {"type": "vm_file", "path": apt_out,  "dest": "apt.out"},
        ],
        "expected": [
            {"type": "rule", "rules": {"expect": ["org\\.openoffice\\.languagetool\\.oxt"]}},
            {"type": "rule", "rules": {"expect": ["openjdk-\\d+-(jre|jdk)"]}},
        ],
    }
    instructions = [
        "Install the LanguageTool extension for LibreOffice on this machine. The extension requires a Java runtime, so install an OpenJDK package via apt first, then add the org.openoffice.languagetool.oxt extension via `unopkg add`. Save the output of `unopkg list` to /tmp/grep_multi_eval_lo_extension_languagetool_jdk.out and the matching `dpkg -l` line to /tmp/apt_multi_eval_lo_extension_languagetool_jdk.out so I can verify both are present.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_list+check_list",
        setup_class="eval_emul_lo_extension_install",
    )


# ---- Cat U.13 — check_structure_sim_with_threshold^2 (GIMP + Python) ------
# emulates osworld_multi_apps_e8172110: extract a pixel-art character via
# GIMP (manual) AND via a Python script — both outputs must pass SSIM ≥ 0.85
# against the same gold image. Apps: GIMP + Python/terminal.
def _make_pixel_art_extract_template() -> SynthTemplate:
    template_id = "multi_eval_pixel_art_extract_gimp_code"
    src_path    = f"{_DESKTOP}/character.png"
    gimp_path   = f"{_DESKTOP}/character_gimp.png"
    code_path   = f"{_DESKTOP}/character_code.png"
    gold_path   = f"{_TMP}/expected_{template_id}.png"

    init_steps = [
        # Build a deterministic pixel-art-style source: a 64×64 image with
        # a simple solid figure on white background, then the gold removes
        # the white background (alpha=0 there). Code build same as gold.
        _python_step(textwrap.dedent(f"""\
            import os
            from PIL import Image, ImageDraw
            os.makedirs(os.path.dirname({src_path!r}) or '.', exist_ok=True)
            img = Image.new('RGBA', (64, 64), (255,255,255,255))
            d = ImageDraw.Draw(img)
            d.rectangle([20,16, 44,48], fill=(80,160,80,255), outline=(0,0,0,255))
            d.rectangle([26,22, 30,26], fill=(0,0,0,255))
            d.rectangle([34,22, 38,26], fill=(0,0,0,255))
            d.line([26,40, 38,40], fill=(0,0,0,255))
            img.save({src_path!r})
            # Gold: same image with white background converted to alpha=0.
            gold = img.copy()
            data = gold.load()
            w,h = gold.size
            for y in range(h):
                for x in range(w):
                    r,g,b,a = data[x,y]
                    if (r,g,b) == (255,255,255):
                        data[x,y] = (0,0,0,0)
            gold.save({gold_path!r})
            """)),
        *_multi_app_preopen(foreground=src_path),
    ]
    oracle_steps = [
        _execute(f"cp '{gold_path}' '{gimp_path}'"),
        _execute(f"cp '{gold_path}' '{code_path}'"),
    ]

    evaluator = {
        "func": ["check_structure_sim_with_threshold", "check_structure_sim_with_threshold"],
        "result": [
            {"type": "vm_file", "path": gimp_path, "dest": "character_gimp.png"},
            {"type": "vm_file", "path": code_path, "dest": "character_code.png"},
        ],
        "expected": [
            {"type": "vm_file", "path": gold_path, "dest": "character_gold.png"},
            {"type": "vm_file", "path": gold_path, "dest": "character_gold.png"},
        ],
        "options": [{"ssim_threshold": 0.85}, {"ssim_threshold": 0.85}],
    }
    instructions = [
        "Open ~/Desktop/character.png in GIMP, extract the pixel-art character by making the white background transparent, and save the result as ~/Desktop/character_gimp.png. Additionally, write a Python script that automates the same background-removal pipeline and emits the result as ~/Desktop/character_code.png. Both outputs must be visually similar to the cleanly extracted character.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[0], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=1,
        eval_class="check_structure_sim_with_threshold+check_structure_sim_with_threshold",
        setup_class="eval_emul_pixel_art_extract",
    )


# ---- Cat U.14 — check_include_exclude+compare_archive (.doc -> .pdf tgz) --
# emulates osworld_multi_apps_f7dfbef3: agent batch-converts .doc files to
# .pdf via soffice, then archives the .pdfs as a tar.gz. atom 1 =
# check_include_exclude on bash_history probe; atom 2 = compare_archive
# (file_type=pdf) against gold tarball.
_DOC_BATCH_CONVERT_SPECS: list[dict] = [
    {
        "id":   "multi_eval_doc_batch_to_pdf_targz_meeting",
        "doc_basenames": ["minutes-q1.doc", "minutes-q2.doc", "minutes-q3.doc"],
        "doc_texts": [
            "Q1 meeting minutes. Decision: launch beta.",
            "Q2 meeting minutes. Decision: hire two engineers.",
            "Q3 meeting minutes. Decision: defer feature flag refactor.",
        ],
        "tgz_basename": "pdf.tar.gz",
    },
    {
        "id":   "multi_eval_doc_batch_to_pdf_targz_invoices",
        "doc_basenames": ["invoice-2024-01.doc", "invoice-2024-02.doc"],
        "doc_texts": [
            "Invoice January 2024. Total: 1240 USD.",
            "Invoice February 2024. Total: 980 USD.",
        ],
        "tgz_basename": "pdf.tar.gz",
    },
]


def _make_doc_batch_convert_template(spec: dict) -> SynthTemplate:
    template_id = spec["id"]
    work_dir    = f"/home/user/work_{template_id}"
    tgz_path    = f"{_DESKTOP}/{spec['tgz_basename']}"
    expected_tgz = f"{_TMP}/expected_{template_id}.tar.gz"
    hist_probe   = f"{_TMP}/probe_{template_id}.txt"

    init_steps = [
        _execute(f"mkdir -p '{work_dir}'"),
        # Stage N .doc files (use soffice to convert from .txt → .doc).
        _python_step(textwrap.dedent(f"""\
            import os, subprocess, tempfile
            work = {work_dir!r}
            texts = {spec["doc_texts"]!r}
            names = {spec["doc_basenames"]!r}
            for n, t in zip(names, texts):
                txt = os.path.join(work, n.replace('.doc', '.txt'))
                with open(txt, 'w') as f:
                    f.write(t)
                subprocess.run(['soffice','--headless','--norestore','--nofirststartwizard',
                                '--convert-to','doc','--outdir',work,txt],
                               check=True, env={{**os.environ, 'DISPLAY':':1'}}, timeout=180)
                os.remove(txt)
            """)),
        # Build the gold tar.gz of converted pdfs from the same texts.
        _python_step(textwrap.dedent(f"""\
            import os, subprocess, tempfile, tarfile
            stage = tempfile.mkdtemp()
            texts = {spec["doc_texts"]!r}
            names = {spec["doc_basenames"]!r}
            for n, t in zip(names, texts):
                txt = os.path.join(stage, n.replace('.doc', '.txt'))
                with open(txt, 'w') as f:
                    f.write(t)
                subprocess.run(['soffice','--headless','--norestore','--nofirststartwizard',
                                '--convert-to','pdf','--outdir',stage,txt],
                               check=True, env={{**os.environ, 'DISPLAY':':1'}}, timeout=180)
                os.remove(txt)
            with tarfile.open({expected_tgz!r}, 'w:gz') as tf:
                for fn in sorted(os.listdir(stage)):
                    tf.add(os.path.join(stage, fn), arcname=fn)
            """)),
        _execute("history -c 2>/dev/null || true; echo > ~/.bash_history"),
        *_terminal_preopen_steps(),
    ]
    # Oracle: simulate user running the desired soffice batch command,
    # plant tarball + history probe.
    oracle_steps = [
        _execute(f"echo 'cd {work_dir} && soffice --headless --convert-to pdf *.doc' >> ~/.bash_history"),
        _execute(
            "if cat ~/.bash_history | grep -E '(soffice|libreoffice).+--convert-to[[:space:]]+pdf.+\\*\\.doc' >/dev/null; "
            f"then echo 'catch the desired command' > '{hist_probe}'; else echo 'failed to complete this task' > '{hist_probe}'; fi"
        ),
        _execute(f"cp '{expected_tgz}' '{tgz_path}'"),
    ]

    evaluator = {
        "func": ["check_include_exclude", "compare_archive"],
        "result": [
            {"type": "vm_command_line", "command": f"cat {hist_probe}", "shell": True},
            {"type": "vm_file", "path": tgz_path,   "dest": spec["tgz_basename"]},
        ],
        "expected": [
            {"type": "rule", "rules": {
                "include": ["catch the desired command"],
                "exclude": ["failed to complete this task"],
            }},
            {"type": "vm_file", "path": expected_tgz, "dest": f"gold_{spec['tgz_basename']}"},
        ],
        "options": [{}, {"file_type": "pdf"}],
    }
    instructions = [
        f"From a single command line invocation in {work_dir}, batch-convert every .doc in the directory to .pdf (e.g. `soffice --headless --convert-to pdf *.doc`), then package the resulting PDFs into ~/Desktop/{spec['tgz_basename']}.",
        f"Inside {work_dir}, convert all .doc files to PDF at once using a single soffice command line (`soffice --headless --convert-to pdf *.doc`), and archive the resulting PDFs into ~/Desktop/{spec['tgz_basename']} as a gzip-compressed tarball.",
    ]

    def _params(seed: int) -> dict:
        return {"instr": instructions[seed % len(instructions)], "config_override": init_steps}

    return SynthTemplate(
        template_id=template_id,
        domain="multi_apps",
        instruction_fn=lambda p: p["instr"],
        evaluator_fn=lambda _p: evaluator,
        oracle_fn=lambda _p: oracle_steps,
        postconfig_fn=lambda _p: None,
        param_fn=_params,
        n_rows=2,
        eval_class="check_include_exclude+compare_archive",
        setup_class="eval_emul_doc_batch_convert_archive",
    )


# ---------------------------------------------------------------------------
# TEMPLATES export
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# DISABLED families (#116 parity + agent/env separation).
# These templates hard-code an instruction to use a CLI/library that the
# official OSWorld guest does NOT have and that lite.osworld deliberately keeps
# env-side under /opt/env (vendored to /opt/env/bin): ImageMagick (convert/magick/…),
# pandoc, pdftk, pymupdf(`import fitz`), mutagen. After the agent/env-facing
# separation the agent can neither discover (`which`/`whereis`/`find`) nor run
# these tools, so a task that mandates them is unsolvable-as-written and (pre-
# separation) drove /opt/env leakage. They were originally added "once the
# Dockerfile installs IM/pandoc" — an assumption #116 reversed.
# Commented out rather than deleted so they can be revived tool-neutrally: an
# image task → GIMP; mp3 tags → VLC Media Information; PDF page count/text →
# poppler (pdfinfo/pdftotext) or evince; PDF split/merge → poppler
# (pdfseparate/pdfunite). The pandoc/epub families have NO installed equivalent
# on the faithful guest and would also need their eval loosened, so they stay
# off until reworked.
# ---------------------------------------------------------------------------
TEMPLATES: list[SynthTemplate] = (
    [_make_txt_to_docx_template(s) for s in _TXT_TO_DOCX_SPECS]
    + [_make_csv_to_xlsx_template(s) for s in _CSV_TO_XLSX_SPECS]
    + [_make_xlsx_to_csv_template(s) for s in _XLSX_TO_CSV_SPECS]
    + [_make_xlsx_to_docx_table_template(s) for s in _XLSX_TO_DOCX_TABLE_SPECS]
    + _shell_pipeline_templates()
    + [
        _make_extract_sort_template(),
        _make_extract_targz_template(),
        _make_zip_create_template(),
    ]
    + [_make_asset_shell_row(s) for s in _ASSET_SHELL_SPECS]
    + [_make_asset_csv_to_xlsx_template()]
    # Cat I — cross-app real-asset rows (CSV→docx-table, photo→docx,
    # photo→pptx, CSV→xlsx-with-chart, multi-CSV concat, HTML extract).
    + [_make_csv_to_docx_table_template(s) for s in _CSV_TO_DOCX_TABLE_SPECS]
    # Cartesian (5 topics × 3 photo skills): photo→docx, photo→pptx, photo
    # gallery (4 slides). Replaces the legacy hardcoded `_PHOTO_TO_DOCX_SPECS`
    # / `_PHOTO_TO_PPTX_SPECS` / `_PHOTO_GALLERY_SPECS` lists.
    + [_make_topic_photo_to_docx_template(t) for t in TOPIC_FAMILIES]
    + [_make_topic_photo_to_docx_template(t, n_rows=2, _variant="cover") for t in TOPIC_FAMILIES]
    + [_make_topic_photo_to_pptx_template(t) for t in TOPIC_FAMILIES]
    # Gallery template uses 4 photos — only run on families with ≥4 photos.
    # Some topic families have only 2-3 photos
    # which can't satisfy gallery N=4; gracefully skip them.
    + [_make_topic_photo_gallery_template(t) for t in TOPIC_FAMILIES if len(TOPIC_FAMILIES[t]) >= 4]
    # Cat I.7b — `compare_docx_images` (PIL pixel-eq) bridges DROPPED in
    # compare_docx_images. PIL `tobytes()` byte-equality means the gold
    # docx (built by `add_picture`) differs in encoding from any agent
    # docx that goes through LO Insert > Picture (JPEG re-encode); the
    # only way to pass was the cp-only oracle path. Real-agent skill
    # signal was zero. Drop until shape/size eval infra exists.
    # + [_make_photo_to_docx_images_eval_template("astronomy")]
    # + [_make_dual_photo_to_docx_images_eval_template("astronomy")]
    # Validation fix (CRITICAL): chart templates DROPPED. openpyxl chart
    # XML round-trips imperfectly through LO save (chart.title is wiped when
    # no label container is set; openpyxl Series/NumDataSource serialization
    # varies). Gold post-LO-normalize has title=None; an agent UI insert
    # WITH title would fail compare_table chart-prop equality. Defer until
    # a deterministic build-via-Calc-UI gold path is wired.
    # + [_make_csv_to_xlsx_chart_template(s) for s in _CSV_TO_XLSX_CHART_SPECS]
    + [_make_asset_concat_3way_economic()]
    + [_make_asset_wikipedia_h2_extract()]
    # Cat I (third pass, 2026-05-10) — additional cross-app real-asset rows.
    + [_make_csv_to_docx_text_summary_template(s) for s in _CSV_TO_DOCX_TEXT_SUMMARY_SPECS]
    # REVIVED (Tier 1): arxiv PDF tasks now extract via poppler (pdftotext/
    # pdfinfo, on the agent PATH) instead of pymupdf fitz — gold+oracle+agent all
    # use poppler so the exact compare_text_file (+ignore_blanks) stays solvable.
    + [_make_arxiv_pdf_template(s) for s in _ARXIV_PDF_SPECS]
    + [_make_csv_concat_xlsx_template(s) for s in _CSV_CONCAT_XLSX_SPECS]
    + [_make_code_to_docx_template(s) for s in _CODE_TO_DOCX_SPECS]
    # Cat J — DISABLED (#116): pandoc is env-only; no installed markdown converter.
    # + [_pandoc_md_to_docx_template(s) for s in _PANDOC_MD_TO_DOCX_SPECS]
    # + [_pandoc_csv_to_md_template()]
    # + [_pandoc_html_to_md_template(s) for s in _PANDOC_HTML_TO_MD_SPECS]
    # + [_pandoc_docx_to_txt_template()]
    # Cat K — DISABLED (#116): ImageMagick is env-only; gold built by IM won't match a GUI result.
    # + [_make_im_template(s) for s in _IM_PHOTO_SPECS]
    # + [_make_im_montage_contact_sheet()]
    # Cat L — poppler-utils (pdftotext / pdfinfo) rows.
    + [_make_pdftotext_template(s) for s in _PDFTOTEXT_SPECS]
    + [_make_pdfinfo_pagecount_template()]
    # Cat M — pdftoppm (PDF page → PNG image).
    + [_make_pdftoppm_template(s) for s in _PDFTOPPM_SPECS]
    # Cat N — compare_pdfs rows (validation: docx/xlsx/pptx → pdf via soffice;
    # pdftk merge / extract pages).
    + [_make_docx_to_pdf_template(s) for s in _DOCX_TO_PDF_SPECS]
    + [_make_xlsx_to_pdf_template(s) for s in _XLSX_TO_PDF_SPECS]
    + [_make_pptx_to_pdf_template()]
    # REVIVED (Tier 1): instruction now uses poppler pdfseparate/pdfunite (agent
    # PATH); gold stays pdftk env-side (compare_pdfs is fuzzy-text, tool-neutral).
    + [_make_pdftk_merge_template()]
    + [_make_pdftk_extract_template()]
    # Cat O — exact_match rows (validation: csv → single-string answer).
    + [_make_exact_match_template(s) for s in _EXACT_MATCH_SPECS]
    # Cat O.2 — HTML→single-string-answer ( Chrome+terminal gap).
    + [_make_html_fact_extract_template(s) for s in _HTML_FACT_EXTRACT_SPECS]
    # Cat P — compare_image_list rows (validation: photo subset collect/rename).
    + [_make_compare_image_list_template(s) for s in _COMPARE_IMAGE_LIST_SPECS]
    # Cat Q — check_python_file_by_test_suite rows (validation: fix-bug).
    + [_make_py_test_template(s) for s in _PY_TEST_SPECS]
    # Cat R — Long-tail gap-fill ( 2026-05-10): under-covered eval
    # buckets brought from synth=0 to 1+ row each.
    + [_make_diff_text_file_template(s) for s in _DIFF_TEXT_FILE_SPECS]
    + [_make_check_list_template(s) for s in _CHECK_LIST_SPECS]
    # DISABLED (#116): mp3_meta mandates mutagen (env-only); use VLC Media Information if revived.
    # + [_make_check_mp3_meta_template(s) for s in _CHECK_MP3_META_SPECS]
    # DISABLED (#116): image_text renders text via ImageMagick (env-only, pixel-exact gold).
    # + [_make_compare_image_text_template(s) for s in _COMPARE_IMAGE_TEXT_SPECS]
    # DISABLED (#116): epub tasks mandate pandoc (env-only); no installed epub converter.
    # + [_make_compare_epub_template(s) for s in _COMPARE_EPUB_SPECS]
    + [_make_is_expected_tabs_template(s) for s in _IS_EXPECTED_TABS_SPECS]
    # REVIVED (Tier 1): instructions now use poppler pdfseparate/pdfunite (agent
    # PATH); gold stays pdftk env-side (compare_pdfs is fuzzy-text, tool-neutral).
    + [_make_compare_pdfs_template(s) for s in _COMPARE_PDFS_SPECS]
    + [_make_compare_table_extra_template(s) for s in _COMPARE_TABLE_EXTRA_SPECS]
    # Long-tail K=1 buckets — one template each.
    + [
        _make_check_json_template(),
        _make_check_direct_json_object_template(),
        _make_compare_zip_files_template(),
        # REVIVED (Tier 1): instruction now uses GIMP (agent PATH); gold stays
        # ImageMagick env-side (evals are rule-based: width/height, max_size).
        _make_check_image_size_template(),
        _make_check_image_file_size_template(),
        _make_compare_python_pure_text_template(),
        _make_check_line_number_template(),
        _make_file_contains_template(),
        _make_is_in_list_template(),
        _make_compare_docx_files_and_ignore_new_lines_template(),
        # DISABLED (#116): md→html via pandoc (env-only).
        # _make_compare_htmls_template(),
        _make_compare_audios_template(),
        # pdftoppm verified at /usr/bin/pdftoppm in the container (poppler-utils).
        _make_compare_pdf_images_template(),
        # DISABLED (#116): crop/transform gold built via ImageMagick (env-only).
        # _make_check_structure_sim_template(),
        # _make_check_structure_sim_with_threshold_template(),
        _make_literal_match_template(),
    ]
    # Cat S — 3-app cross-app workflow templates (multi_apps Axis A).
    # Genuine 3-app GUI hand-off (chrome HTML → calc xlsx → writer docx,
    # thunderbird-derived xlsx → calc → writer, impress → writer → pdf).
    + [
        _make_3app_chrome_calc_writer_market(),
        _make_3app_chrome_calc_writer_landmarks(),
        _make_3app_chrome_calc_writer_tech(),
        _make_3app_calc_thunderbird_writer_inbox(),
        _make_3app_impress_writer_pdf_brief(),
        _make_3app_chrome_writer_pdf_factsheet(),
        _make_3app_calc_writer_pdf_quarterly(),
    ]
    # Cat T — eval-emulation templates ( 2026-05-12). Each bridges
    # a compound no_fn-analog signature with 0% synth coverage; see header
    # docstring "Cat T" for citations of the osworld_multi_apps_<id> rows
    # each template emulates.
    + [
        # DISABLED (#116): mp3 batch mandates mutagen; epub mandates pandoc (both env-only).
        # _make_eval_mp3_meta_batch5_template(),               # 3f05f3b9
        # _make_eval_compare_epub_3path_or_template(),         # 42d25c08
        _make_eval_xcf_to_docx_image_paste_template(),       # 227d2f97
        _make_eval_notebook_extract_python_template(),       # dd60633f
        _make_eval_xlsx_ods_merge_csv_compound_template(),   # 3680a5ee
        _make_eval_tally_book_append_compound_template(),    # 415ef462
    ]
    # Cat U — validation compound-signature gap-closers (2026-05-12).
    # Each template emulates one or more osworld_multi_apps_<id> tasks with
    # a zero-coverage compound evaluator signature.
    + [_make_terminal_screenshot_template(s) for s in _TERMINAL_SCREENSHOT_SPECS]  # 02ce9a50
    + [_make_sysstat_template(s) for s in _SYSSTAT_SPECS]                          # 2373b66a
    + [_make_timetable_template(s) for s in _TIMETABLE_SPECS]                      # 3a93cae4
    + [_make_vscode_ext_img_template(s) for s in _VSCODE_EXT_IMG_SPECS]            # 42f4d1c7
    + [_make_workspace_setup_template(s) for s in _WORKSPACE_SETUP_SPECS]          # 48c46dc7
    + [_make_vscode_open_from_term_template(s) for s in _VSCODE_OPEN_FROM_TERM_SPECS]  # 510f64c8
    + [_make_paper_pdf_docx_template(s) for s in _PAPER_PDF_DOCX_SPECS]            # 68a25bd4
    + [_make_webext_scaffold_template()]                                           # 74d5859f
    + [_make_desk_triage_template(s) for s in _DESK_TRIAGE_SPECS]                  # 869de13e
    + [_make_contacts_export_template(s) for s in _CONTACTS_EXPORT_SPECS]          # c867c42d
    + [_make_xlsx_to_html_template(s) for s in _XLSX_TO_HTML_SPECS]                # e135df7c
    + [_make_lo_extension_install_template()]                                      # e1fc0df3
    + [_make_pixel_art_extract_template()]                                         # e8172110
    + [_make_doc_batch_convert_template(s) for s in _DOC_BATCH_CONVERT_SPECS]      # f7dfbef3
    # Note: 78aed49a (check_thunderbird_folder+compare_pdfs) intentionally
    # SKIPPED — requires Google Drive cloud auth + Thunderbird profile-level
    # message-move semantics that we can't oracle deterministically.
)
