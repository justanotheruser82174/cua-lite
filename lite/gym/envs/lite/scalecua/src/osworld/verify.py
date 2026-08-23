"""ScaleCUA-aware OSWorld evaluator."""

from __future__ import annotations

import asyncio
import base64
import contextlib
import copy
import hashlib
import inspect
import functools
import json
import logging
import math
import os
import re
import shlex
import statistics
import tempfile
import xml.etree.ElementTree as ET
from datetime import datetime
from typing import Any
from urllib.parse import parse_qs, unquote, urlparse, urlsplit
from urllib.request import urlopen

from lite.core.tools.calls import (
    tool_call_arguments,
    tool_call_name,
    validate_lite_tool_call,
)
from lite.gym.envs.lite.osworld.src.eval import runner as base_runner
from lite.gym.envs.lite.scalecua.src.osworld import judges
from lite.gym.envs.lite.scalecua.src.osworld.setup import dispatch_strict
from lite.gym.sandbox.types import SandboxTaskConfig

logger = logging.getLogger(__name__)


BASE_RUNNER_RESULT_TYPES = {
    "accessibility_tree",
    "audio_in_slide",
    "background_image_in_slide",
    "bookmarks",
    "cache_file",
    "check_dns_prefetch",
    "chrome_appearance_mode_ui",
    "chrome_color_scheme",
    "chrome_font_size",
    "cloud_file",
    "content_from_vm_file",
    "cookie_data",
    "default_search_engine",
    "default_video_player",
    "enable_safe_browsing",
    "find_installed_extension_name",
    "gimp_config_file",
    "googledrive_file",
    "history",
    "list_directory",
    "page_info",
    "profile_name",
    "rule",
    "shortcuts_on_desktop",
    "url_path_parse",
    "vlc_config",
    "vm_command_error",
    "vm_command_line",
    "vm_screen_size",
    "vm_terminal_output",
    "vm_wallpaper",
    "vm_window_size",
    "vscode_config",
}

SCALECUA_LOCAL_RESULT_TYPES = {
    "active_tab_html_parse",
    "active_tab_info",
    "active_tab_url_parse",
    "active_url_from_accessTree",
    "data_delete_automacally",
    "disable_safe_browsing",
    "enable_do_not_track",
    "enable_enhanced_safety_browsing",
    "enable_safe_browsing",
    "find_unpacked_extension_path",
    "new_startup_page",
    "open_tabs_info",
    "url_dashPart",
    "vlc_playing_info",
    "vm_file",
}

SCALECUA_LOCAL_HASHED_RESULT_TYPES = {
    "recreation_url_check",
    "recreation_devilsgarden_html",
}

SCALECUA_LOCAL_CHROME_PREF_RESULT_TYPES = {
    "block_third_party_cookies",
    "chrome_block_third_party_cookies",
    "chrome_third_party_cookies",
    "chrome_third_party_cookies_blocked",
    "data_delete_automacally",
    "disable_safe_browsing",
    "enable_do_not_track",
    "enable_enhanced_safety_browsing",
    "enable_safe_browsing",
    "new_startup_page",
    "password_manager_disabled",
    "password_manager_enabled",
    "third_party_cookies_blocked",
}

VM_FILE_TEXT_METRIC_PREFIXES = (
    "check_tb_triple_prefs__",
    "check_tb_dual_auto_quote_reply__",
    "check_tb_dual_auto_quote_sig__",
)

_RECREATION_WEEKDAY_DATE_RE = re.compile(
    r"\b(?:mon|tue|wed|thu|fri|sat|sun)\s+\d{1,2}\b",
    re.IGNORECASE,
)
_RECREATION_MONTH_DATE_RE = re.compile(
    r"\b(?:jan(?:uary)?|feb(?:ruary)?|mar(?:ch)?|apr(?:il)?|may|jun(?:e)?|"
    r"jul(?:y)?|aug(?:ust)?|sep(?:t(?:ember)?)?|oct(?:ober)?|"
    r"nov(?:ember)?|dec(?:ember)?)\s+\d{1,2}(?:,\s*\d{4})?\b",
    re.IGNORECASE,
)
_GIMP_ACTION_HISTORY = "/home/user/.config/GIMP/2.10/action-history"
_GIMP_WINDOW_MARKERS = ("gimp", "gnu image manipulation program")
_GIMP_CONFIG_FILE_MARKERS = (
    "/home/user/.config/GIMP/2.10/gimprc",
    "/home/user/.config/GIMP/2.10/sessionrc",
    "~/.config/GIMP/2.10/gimprc",
    "~/.config/GIMP/2.10/sessionrc",
)
_GIMP_CONFIG_RESULT_TYPES = {
    "gimp_config_file",
    "gimp_icon_theme",
    "gimp_theme_and_icons",
    "gimp_theme_setting",
}
_GIMP_CONFIG_RESULT_PREFIXES = (
    "gimp_icon_theme",
    "gimp_theme",
    "gimp_config",
)
_GIMP_ACTION_WINDOW_TOKENS: dict[str, tuple[str, ...]] = {
    "colors-brightness-contrast": ("brightness-contrast", "brightness contrast"),
    "colors-color-balance": ("color balance",),
    "colors-curves": ("curves", "color curves"),
    "colors-desaturate": ("desaturate",),
    "colors-hue-saturation": ("hue-saturation", "hue saturation"),
    "colors-levels": ("levels",),
    "colors-posterize": ("posterize",),
    "colors-threshold": ("threshold",),
    "filters-bloom": ("bloom",),
    "filters-blur": ("blur", "gaussian blur"),
    "filters-cartoon": ("cartoon",),
    "filters-despeckle": ("despeckle",),
    "filters-drop-shadow": ("drop shadow",),
    "filters-edge": ("edge", "edge-detect", "edge detect"),
    "filters-gaussian-blur": ("gaussian blur",),
    "filters-lens-distortion": ("lens distortion",),
    "filters-light-shadow-drop-shadow": ("drop shadow",),
    "filters-mblur": ("motion blur",),
    "filters-noise-reduction": ("noise reduction",),
    "filters-oilify": ("oilify",),
    "filters-pixelize": ("pixelize",),
    "filters-red-eye-removal": ("red eye removal", "red-eye removal"),
    "filters-sharpen": ("sharpen",),
    "filters-softglow": ("softglow", "soft glow"),
    "filters-unsharp-mask": ("unsharp mask", "sharpen (unsharp mask)", "sharpen"),
    "filters-vignette": ("vignette",),
}
# Legacy upstream task metadata may point at these original-author cache roots.
# Keep exact-prefix matching: this is a compatibility sentinel for upstream
# embedded paths, not a reference to the current checkout or host machine.
_UPSTREAM_LEGACY_CACHE_PREFIXES = (
    "/home/lvbowen/project/AutoGen/src/envs/osworld_env/cache/",
    "/home/lvbowen/project/SCALE-CUA/VeriGen/osworld_env/cache/",
)
_REFERENCE_METRIC_KEYS = {
    "source_path",
    "original_path",
    "source_cache_path",
    "src_path",
    "tgt_path",
}
_REFERENCE_RESULT_CONFIG_KEYS = {"source_cache_path"}
_CHROME_SETTINGS_SNAPSHOT_MARKER = "__SCALECUA_CHROME_SETTINGS_SNAPSHOT__"

BASE_RUNNER_EXPECTED_TYPES = {
    "accessibility_tree",
    "cloud_file",
    "content_from_vm_file",
    "gotoRecreationPage_and_get_html_content",
    "info_from_website",
    "list",
    "pdf_from_url",
    "rule",
    "rule_relativeTime",
    "vm_command_error",
    "vm_command_line",
    "vm_screen_size",
    "vm_window_size",
}

SCALECUA_LOCAL_EXPECTED_TYPES = {
    "vm_file",
}


async def evaluate_final_fn(
    task: SandboxTaskConfig,
    computer,
    actions: list | None = None,
    debug: bool = False,
) -> float | tuple[float, dict]:
    evaluator = copy.deepcopy(task.metadata.get("evaluator", {}))
    if not evaluator:
        return (0.0, {"error": "no evaluator"}) if debug else 0.0

    runtime_split = (task.metadata.get("scalecua") or {}).get("runtime_split")
    actions = actions or []
    pre_postconfig_state = await _capture_pre_postconfig_state(computer, evaluator)
    postconfig_done = bool(evaluator.get("_postconfig_done"))
    if evaluator.get("postconfig"):
        cache_dir = tempfile.mkdtemp(prefix="scalecua_eval_")
        os.makedirs(cache_dir, exist_ok=True)
        if not postconfig_done:
            try:
                await _run_postconfig(computer, evaluator, cache_dir)
            except Exception as exc:
                if debug:
                    return 0.0, {"postconfig_error": str(exc)}
                return 0.0
    else:
        cache_dir = None

    func = evaluator.get("func", "")
    if func == "infeasible":
        ok = _reported_infeasible(actions)
        return (1.0 if ok else 0.0, {"infeasible": ok}) if debug else (1.0 if ok else 0.0)

    forfeit_reason = _final_action_forfeit_reason(actions)
    if forfeit_reason:
        return (
            (0.0, {"terminal_failure": True, "reason": forfeit_reason})
            if debug
            else 0.0
        )

    if runtime_split in {"train", "rl"}:
        return await evaluate_scalecua_task(
            computer,
            evaluator,
            runtime_split=runtime_split,
            cache_dir=cache_dir,
            run_postconfig=cache_dir is None and not postconfig_done,
            pre_postconfig_state=pre_postconfig_state,
            reference_sources=_reference_sources_from_task(task),
            debug=debug,
        )
    return await base_runner.evaluate_osworld_task(
        computer,
        evaluator,
        cache_dir=cache_dir,
        debug=debug,
    )


def _reported_infeasible(actions: list[dict[str, Any]]) -> bool:
    if not actions:
        return False
    action = actions[-1]
    if _action_is_failure(action):
        return True
    name, args = _action_call(action)
    return name == "response" and "[infeasible]" in str(args.get("text", "")).lower()


def _final_action_forfeit_reason(actions: list[dict[str, Any]]) -> str | None:
    """Why the final action forfeits credit BEFORE metric evaluation, or ``None``.

    THE CONTRACT (``devs/envs/lite.scalecua/validate/rollout/analysis.md``):
    *a final FAIL action must return 0 before metric evaluation*. Two distinct
    ways to trip it, and the second is the whole point of this function:

    * ``"explicit_failure"`` — the agent gave up on the record:
      ``terminate(status="failure")`` / ``report_infeasible``.
    * ``"unreadable_final_action"`` — no name could be read off the final action
      at all. ``_action_call`` is deliberately TOTAL (see its docstring) so a
      malformed action cannot error the episode out of the eval denominator, but
      totality must land as a FAILURE, not as the ABSENCE of one. Reading
      ``("", {})`` and asking "is that a failure?" answers *no*, which falls
      through to normal metric evaluation and lets an unreadable terminal action
      score **1.0** — strictly worse than both the raise it replaced and the 0.0
      the contract demands. So the caller distinguishes "read a name" (``str``)
      from "could not read one" (``None``), and the latter forfeits.

    ``_action_is_failure`` stays a POSITIVE predicate (unreadable → ``False``)
    because ``_reported_infeasible`` reads it with the opposite polarity: there a
    match PAYS 1.0, so an unreadable action must not be mistaken for a correct
    infeasible report.
    """
    if not actions:
        return None
    action = actions[-1]
    if _action_call(action)[0] is None:
        return "unreadable_final_action"
    return "explicit_failure" if _action_is_failure(action) else None


def _action_is_failure(action: Any) -> bool:
    """True iff the action POSITIVELY declares failure. Unreadable → False."""
    name, args = _action_call(action)
    if name == "report_infeasible":
        return True
    return name == "terminate" and args.get("status") == "failure"


def _action_call(action: Any) -> tuple[str | None, dict[str, Any]]:
    """Read ``(name, arguments)`` off a final action. NEVER raises.

    Returns ``(None, {})`` — NOT ``("", {})`` — when no name could be read, so
    callers can tell "read a name" from "could not read one"; the reward path
    must treat the latter as a FAILURE rather than as the absence of one (see
    ``_final_action_forfeit_reason``).

    Deliberately total, unlike the rest of this module, because it is the only
    thing standing between a malformed final action and the reward: it is called
    unguarded from ``evaluate_final_fn``, so a raise here does not degrade the
    score, it ERRORS the episode out of the eval denominator entirely — strictly
    worse than the ``0.0`` a terminal-failure action is supposed to earn.

    ``env.step`` gates its rollout actions through
    ``lite.gym.utils.feedback.ingress.prepare_env_tool_calls``, so rollout
    callers usually arrive here as env-internal bare actions. Canonical callers
    may arrive as nested Lite tool-call envelopes; those are read only through
    ``lite.core.tools.calls`` canonical accessors. Top-level ``name`` /
    ``arguments`` reads below are for env-internal accepted actions and the
    direct oracle/replay harnesses under ``devs/envs/lite.scalecua/validate/``.
    An OpenAI-style provider passthrough can hand top-level ``arguments`` over
    as a JSON STRING — the one non-canonical oracle-private shape common enough
    to decode rather than discard.

    OSWorld's own ``"FAIL"`` / ``{"action_type": "FAIL"}`` sentinels are NOT
    accepted: they are the wire encoding of the ``osworld``/``osworld_2`` envs
    (which post ``{"cmd": "FAIL"}`` to their containers), unrepresentable as a
    Lite tool call, and nothing on this env's path emits them. Discarding them
    is safe only because discarding is not silent: a sentinel reads as ``None``
    here and ``_final_action_forfeit_reason`` turns that into a 0.0, which is
    what the contract in ``devs/envs/lite.scalecua/validate/rollout/analysis.md``
    ("a final FAIL action must return 0 before metric evaluation") asks for —
    the RECOGNIZED spellings ``terminate(status="failure")`` /
    ``report_infeasible`` and every unrecognized one alike.
    """
    if not isinstance(action, dict):
        return None, {}
    canonical = validate_lite_tool_call(action, "action", require_id=False) is None
    if canonical:
        return tool_call_name(action), tool_call_arguments(action)
    name = action.get("name")
    args = action.get("arguments")
    if isinstance(args, str):
        try:
            args = json.loads(args)
        except Exception:
            args = {}
    if not isinstance(name, str) or not name or not isinstance(args, dict):
        return None, {}
    return name, args


async def evaluate_scalecua_task(
    computer,
    evaluator: dict[str, Any],
    *,
    runtime_split: str,
    cache_dir: str | None = None,
    run_postconfig: bool = True,
    pre_postconfig_state: str | None = None,
    reference_sources: dict[str, list[str]] | None = None,
    debug: bool = False,
) -> float | tuple[float, dict]:
    if cache_dir is None:
        cache_dir = tempfile.mkdtemp(prefix="scalecua_eval_")
    os.makedirs(cache_dir, exist_ok=True)
    eval_env = judges.make_eval_env(computer, cache_dir)
    if reference_sources:
        setattr(eval_env, "_scalecua_reference_source_urls", reference_sources)

    if run_postconfig:
        await _run_postconfig(computer, evaluator, cache_dir)

    raw_func = evaluator.get("func", "")
    multi_metric = isinstance(raw_func, list)
    func_list = raw_func
    func_list = func_list if isinstance(func_list, list) else [func_list]
    result_list = _as_list(evaluator.get("result", {}))
    expected_list = _as_list(evaluator.get("expected", {}))
    options_list = _as_list(evaluator.get("options", {}))
    conj = evaluator.get("conj", "and")
    n = len(func_list)
    result_list += [{}] * (n - len(result_list))
    expected_list += [{}] * (n - len(expected_list))
    options_list += [{}] * (n - len(options_list))

    scores: list[float] = []
    details: list[dict[str, Any]] = []
    for i, fn_name in enumerate(func_list):
        try:
            result_data = await _get_result(
                eval_env, result_list[i] or {}, cache_dir, runtime_split
            )
            if i == len(func_list) - 1 and len(expected_list) > len(func_list):
                remaining = expected_list[i:]
                if len(remaining) > 1:
                    expected_data = [
                        await _get_expected(eval_env, cfg or {}, cache_dir, runtime_split)
                        for cfg in remaining
                    ]
                else:
                    expected_data = await _get_expected(
                        eval_env, expected_list[i] or {}, cache_dir, runtime_split
                    )
            else:
                expected_data = await _get_expected(
                    eval_env, expected_list[i] or {}, cache_dir, runtime_split
                )
            expected_data = _normalize_scalecua_expected_rules(
                str(fn_name), expected_data
            )
            result_data, expected_data = _normalize_chrome_url_metric_pair(
                result_list[i] or {},
                result_data,
                expected_data,
            )
            result_data = await _repair_accessibility_namespace_url_result(
                eval_env,
                expected_data,
                result_data,
            )
            result_data = _repair_chrome_settings_result_from_pre_postconfig_state(
                result_list[i] or {},
                result_data,
                pre_postconfig_state,
            )
            raw_result_data = result_data
            result_data = _augment_gimp_action_history_result(
                str(fn_name),
                result_list[i] or {},
                expected_data,
                result_data,
                pre_postconfig_state,
            )
            metric_fn = judges.resolve_metric(str(fn_name), runtime_split)
            opts = options_list[i] or {}
            metric_result = _prepare_metric_result(
                metric_fn, result_data, result_list[i] or {}
            )
            metric_result, metric_expected = _prepare_metric_args(
                str(fn_name), metric_result, expected_data
            )
            raw_score = await judges.call_metric(
                metric_fn,
                eval_env,
                metric_result,
                metric_expected,
                opts,
            )
            score = _coerce_score(raw_score)
            scores.append(score)
            if debug:
                detail = {
                    "func": fn_name,
                    "score": score,
                    "result_preview": _debug_preview(result_data),
                    "expected_preview": _debug_preview(expected_data),
                }
                if _is_gimp_action_history_result(result_list[i] or {}):
                    detail["result_preview"] = _debug_preview(raw_result_data)
                    detail["result_after_window_fallback_preview"] = _debug_preview(
                        result_data
                    )
                    detail["expected_preview"] = _debug_preview(expected_data)
                    detail["pre_postconfig_state_preview"] = _debug_preview(
                        pre_postconfig_state
                    )
                details.append(detail)
            if multi_metric and conj == "and" and score == 0.0:
                return _format_result(
                    0.0, conj, scores, details, debug, flush_stats=_flush_stats_snapshot(eval_env)
                )
            if multi_metric and conj == "or" and score == 1.0:
                return _format_result(
                    1.0, conj, scores, details, debug, flush_stats=_flush_stats_snapshot(eval_env)
                )
        except Exception as exc:
            scores.append(0.0)
            if debug:
                details.append({"func": fn_name, "error": str(exc), "score": 0.0})
            if multi_metric and conj == "and":
                return _format_result(
                    0.0, conj, scores, details, debug, flush_stats=_flush_stats_snapshot(eval_env)
                )

    final = _combine_scores(scores, conj, multi_metric)
    return _format_result(
        final, conj, scores, details, debug, flush_stats=_flush_stats_snapshot(eval_env)
    )


async def _run_postconfig(computer, evaluator: dict[str, Any], cache_dir: str) -> None:
    for index, step in enumerate(evaluator.get("postconfig", [])):
        if _is_chrome_pkill_postconfig_step(step):
            await _flush_chrome_profile(judges.make_eval_env(computer, cache_dir))
            continue
        await dispatch_strict(
            computer,
            step,
            phase="postconfig",
            index=index,
            cache_dir=cache_dir,
        )
    evaluator["_postconfig_done"] = True


def _is_chrome_pkill_postconfig_step(step: Any) -> bool:
    if not isinstance(step, dict) or step.get("type") != "launch":
        return False
    params = step.get("parameters") or {}
    command = params.get("command")
    if isinstance(command, str):
        try:
            tokens = shlex.split(command)
        except ValueError:
            tokens = command.split()
    elif isinstance(command, list):
        tokens = [str(part) for part in command]
    else:
        return False
    if not tokens:
        return False
    if os.path.basename(tokens[0]) != "pkill":
        return False
    return any("chrome" in token.lower() or "chromium" in token.lower() for token in tokens[1:])


async def _capture_pre_postconfig_state(computer, evaluator: dict[str, Any]) -> str | None:
    parts: list[str] = []
    if _is_gimp_action_history_evaluator(evaluator):
        try:
            result = await computer.interface.run_command(
                _build_gimp_window_state_command(),
                timeout=15,
            )
            stdout = getattr(result, "stdout", "") or ""
            if stdout:
                parts.append(stdout)
        except Exception:
            pass
        try:
            result = await computer.interface.run_command(
                _build_accessibility_tree_snapshot_command(),
                timeout=8,
            )
            stdout = getattr(result, "stdout", "") or ""
            if stdout:
                parts.append(stdout)
        except Exception:
            pass
        try:
            result = await computer.interface.run_command(
                _build_accessibility_tree_direct_command(),
                timeout=8,
            )
            stdout = getattr(result, "stdout", "") or ""
            if stdout:
                parts.append(stdout)
        except Exception:
            pass
    chrome_state = await _capture_pre_postconfig_chrome_settings_state(
        computer,
        evaluator,
    )
    if chrome_state:
        parts.append(
            _CHROME_SETTINGS_SNAPSHOT_MARKER
            + json.dumps(chrome_state, ensure_ascii=False, default=str)
        )
    return "\n".join(parts) if parts else None


async def _capture_pre_postconfig_chrome_settings_state(
    computer,
    evaluator: dict[str, Any],
) -> dict[str, dict[str, Any]]:
    if not evaluator.get("postconfig"):
        return {}
    result_configs = [
        config
        for config in _as_list(evaluator.get("result", {}))
        if isinstance(config, dict)
    ]
    wanted: dict[str, str] = {}
    for config in result_configs:
        result_type = str(config.get("type") or "")
        url = _chrome_settings_url_for_result(result_type)
        if url:
            wanted[_base_result_type(result_type)] = url
    if not wanted:
        return {}
    eval_env = judges.make_eval_env(computer, tempfile.gettempdir())
    snapshots_by_url: dict[str, dict[str, Any]] = {}
    snapshots_by_type: dict[str, dict[str, Any]] = {}
    for result_type, url in wanted.items():
        snapshot = snapshots_by_url.get(url)
        if snapshot is None:
            snapshot = await _get_chrome_settings_snapshot(eval_env, url)
            snapshots_by_url[url] = snapshot
        if snapshot:
            snapshots_by_type[result_type] = snapshot
    return snapshots_by_type


def _build_gimp_window_state_command() -> str:
    script = """
set +e
for display in "${DISPLAY:-}" :1 :0; do
  [ -z "$display" ] && continue
  DISPLAY="$display"; export DISPLAY
  printf '\\n__DISPLAY__%s\\n' "$DISPLAY"
  printf '__ACTIVE__\\n'
  active="$(timeout 1s xdotool getactivewindow 2>/dev/null || true)"
  if [ -n "$active" ]; then
    timeout 1s xdotool getwindowname "$active" 2>/dev/null || true
    timeout 1s xprop -id "$active" WM_NAME _NET_WM_NAME WM_CLASS 2>/dev/null || true
  fi
  printf '\\n__WMCTRL_L__\\n'
  timeout 1s wmctrl -l 2>/dev/null || true
  printf '\\n__WMCTRL_LX__\\n'
  timeout 1s wmctrl -lx 2>/dev/null || true
  printf '\\n__XWININFO_ROOT_TREE__\\n'
  timeout 2s xwininfo -root -tree 2>/dev/null || true
  printf '\\n__XPROP_ROOT_CLIENT_LIST__\\n'
  timeout 1s xprop -root _NET_CLIENT_LIST _NET_CLIENT_LIST_STACKING 2>/dev/null || true
  printf '\\n__XPROP_CLIENT_WINDOWS__\\n'
  for wid in $(timeout 1s xprop -root _NET_CLIENT_LIST _NET_CLIENT_LIST_STACKING 2>/dev/null | grep -o '0x[0-9a-fA-F]\\+' | sort -u || true); do
    printf '__WINDOW__%s\\n' "$wid"
    timeout 1s xwininfo -id "$wid" 2>/dev/null || true
    timeout 1s xprop -id "$wid" WM_NAME _NET_WM_NAME WM_CLASS 2>/dev/null || true
  done
  printf '\\n__XDO_TOOL_GIMP_WINDOWS__\\n'
  for wid in $(timeout 1s xdotool search --class 'Gimp' 2>/dev/null; timeout 1s xdotool search --class 'gimp' 2>/dev/null; timeout 1s xdotool search --class 'Gimp-2.10' 2>/dev/null || true); do
    printf '__WINDOW__%s\\n' "$wid"
    timeout 1s xdotool getwindowname "$wid" 2>/dev/null || true
    timeout 1s xprop -id "$wid" WM_NAME _NET_WM_NAME WM_CLASS 2>/dev/null || true
  done
  printf '\\n__XDO_TOOL_VISIBLE_WINDOWS__\\n'
  for wid in $(timeout 1s xdotool search --onlyvisible --name '.*' 2>/dev/null || true); do
    printf '__WINDOW__%s\\n' "$wid"
    timeout 1s xdotool getwindowname "$wid" 2>/dev/null || true
    timeout 1s xprop -id "$wid" WM_NAME _NET_WM_NAME WM_CLASS 2>/dev/null || true
  done
done
"""
    return "bash -lc " + shlex.quote(script)


def _build_accessibility_tree_snapshot_command() -> str:
    script = r"""
import json
import urllib.request

try:
    raw = urllib.request.urlopen("http://localhost:5000/accessibility", timeout=5).read()
    data = json.loads(raw.decode("utf-8", errors="replace") or "{}")
    text = data.get("AT", "")
    if isinstance(text, str):
        print("__ACCESSIBILITY_TREE__")
        print(text[:200000])
except Exception:
    pass
"""
    return "python3 -c " + shlex.quote(script)


def _build_accessibility_tree_direct_command() -> str:
    python_script = r"""
import sys

try:
    import pyatspi
except Exception:
    sys.exit(0)

interesting_roles = {
    "application",
    "dialog",
    "frame",
    "filler",
    "menu item",
    "page tab",
    "push button",
    "toggle button",
    "window",
}

try:
    desktop = pyatspi.Registry.getDesktop(0)
except Exception:
    sys.exit(0)

print("__ACCESSIBILITY_TREE_DIRECT__")
seen = 0

def walk(node, depth=0):
    global seen
    if seen >= 4000:
        return
    try:
        role = node.getRoleName()
    except Exception:
        role = ""
    try:
        name = node.name or ""
    except Exception:
        name = ""
    try:
        states = node.getState()
        showing = states.contains(pyatspi.STATE_SHOWING)
        visible = states.contains(pyatspi.STATE_VISIBLE)
        active = states.contains(pyatspi.STATE_ACTIVE)
    except Exception:
        showing = visible = active = False
    text = f"{role} {name}".strip()
    if text and (showing or visible or active or role in interesting_roles):
        print("  " * min(depth, 12) + f"{role}: {name} showing={showing} visible={visible} active={active}")
        seen += 1
    try:
        count = node.childCount
    except Exception:
        count = 0
    for idx in range(min(count, 200)):
        try:
            child = node.getChildAtIndex(idx)
        except Exception:
            continue
        walk(child, depth + 1)

walk(desktop)
"""
    shell_script = (
        "set +e\n"
        "export DISPLAY=${DISPLAY:-:1}\n"
        "export XDG_RUNTIME_DIR=${XDG_RUNTIME_DIR:-/run/user/1000}\n"
        "if [ -z \"${DBUS_SESSION_BUS_ADDRESS:-}\" ] "
        "&& [ -s /tmp/dbus-session-bus-address ]; then\n"
        "  export DBUS_SESSION_BUS_ADDRESS=\"$(cat /tmp/dbus-session-bus-address)\"\n"
        "fi\n"
        "python3 -c "
        + shlex.quote(python_script)
    )
    return "bash -lc " + shlex.quote(shell_script)


def _is_gimp_action_history_evaluator(evaluator: dict[str, Any]) -> bool:
    if "check_include_exclude" not in _as_list(evaluator.get("func", "")):
        return False
    return any(
        _is_gimp_action_history_result(config)
        for config in _as_list(evaluator.get("result", {}))
        if isinstance(config, dict)
    )


def _is_gimp_action_history_result(config: dict[str, Any]) -> bool:
    return (
        config.get("type") == "vm_command_line"
        and _GIMP_ACTION_HISTORY in str(config.get("command") or "")
    )


def _augment_gimp_action_history_result(
    fn_name: str,
    result_config: dict[str, Any],
    expected_data: Any,
    result_data: Any,
    pre_postconfig_state: str | None,
) -> Any:
    if (
        fn_name != "check_include_exclude"
        or not isinstance(result_config, dict)
        or not _is_gimp_action_history_result(result_config)
        or not isinstance(expected_data, dict)
        or not isinstance(result_data, str)
        or not pre_postconfig_state
    ):
        return result_data

    include = expected_data.get("include", [])
    if not isinstance(include, list):
        return result_data
    missing = [
        str(token)
        for token in include
        if str(token) not in result_data
        and _gimp_action_window_matches(str(token), pre_postconfig_state)
    ]
    if not missing:
        return result_data
    return result_data + "\n" + "\n".join(missing)


def _gimp_action_window_matches(action_id: str, window_state: str) -> bool:
    tokens = _GIMP_ACTION_WINDOW_TOKENS.get(action_id)
    if not tokens:
        return False
    normalized = re.sub(r"\s+", " ", window_state.lower())
    if not any(marker in normalized for marker in _GIMP_WINDOW_MARKERS):
        return False
    return any(token in normalized for token in tokens)


async def _repair_accessibility_namespace_url_result(
    eval_env,
    expected_data: Any,
    result_data: Any,
) -> Any:
    if not isinstance(expected_data, dict):
        return result_data
    expected_url = expected_data.get("expected_url") or expected_data.get("url")
    if not isinstance(expected_url, str) or not expected_url:
        return result_data
    if isinstance(result_data, str):
        if not _is_accessibility_namespace_url(result_data):
            return result_data
        return (
            await _get_cdp_expected_url_fallback(eval_env.computer, expected_url)
            or result_data
        )
    if not isinstance(result_data, dict):
        return result_data
    repaired = dict(result_data)
    changed = False
    for key in ("url", "active_url"):
        if _is_accessibility_namespace_url(repaired.get(key)):
            fallback = await _get_cdp_expected_url_fallback(
                eval_env.computer,
                expected_url,
            )
            if fallback:
                repaired[key] = fallback
                changed = True
    return repaired if changed else result_data


async def _get_result(eval_env, config: dict, cache_dir: str, runtime_split: str):
    config = _alias_chrome_profile_config(config)
    config = await _repair_reference_paths(
        eval_env,
        config,
        cache_dir,
        reference_keys=_REFERENCE_RESULT_CONFIG_KEYS,
        repair_home_user=False,
    )
    config = _normalize_scalecua_command_config(config)
    result_type = config.get("type", "")
    if _needs_chrome_profile_flush(result_type, config):
        _record_flush_needed(eval_env, "chrome", result_type)
        await _flush_chrome_profile(eval_env)
    if _needs_gimp_config_flush(result_type, config):
        _record_flush_needed(eval_env, "gimp", result_type)
        await _flush_gimp_config(eval_env)
    if _needs_thunderbird_config_flush(result_type, config):
        _record_flush_needed(eval_env, "thunderbird", result_type)
        await _flush_thunderbird(eval_env)
    if _needs_vlc_config_flush(result_type, config):
        _record_flush_needed(eval_env, "vlc", result_type)
        await _flush_vlc(eval_env)
    if _needs_libreoffice_config_flush(result_type, config):
        _record_flush_needed(eval_env, "libreoffice", result_type)
        await _flush_libreoffice(eval_env)
    # Reproduce the VS Code eval-helper extension's answer file
    # from on-disk state before any vscode getter downloads it — covers both the
    # base `vscode_config` and the hashed `vscode_*__<hash>` twins with one shared
    # reader. No extension, no image rebuild. Idempotent w.r.t. base_runner.
    if config.get("vscode_extension_command"):
        await base_runner._vscode_write_answer_file(eval_env.computer, config)
    if _is_django_repo_file_count_result(config):
        return await _get_django_repo_file_count(eval_env)
    if _base_result_type(str(result_type)) == "xls_cell_value":
        return await _get_generated_xls_cell_value(eval_env, config)
    if _is_scalecua_local_result_type(result_type):
        result = await _get_scalecua_local_result(eval_env, config, cache_dir)
        result = await _repair_clipboard_result(eval_env, config, result)
        result = await _repair_chrome_extension_result(eval_env, config, result)
        result = _normalize_scalecua_result(result_type, result)
        return await _repair_reference_paths(
            eval_env,
            result,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    if result_type in BASE_RUNNER_RESULT_TYPES:
        result = await base_runner._get_result(eval_env.computer, config, cache_dir)
        result = await _repair_clipboard_result(eval_env, config, result)
        result = await _repair_chrome_extension_result(eval_env, config, result)
        result = await _repair_default_search_engine_result(eval_env, config, result)
        result = await _repair_chrome_profile_name_result(eval_env, result_type, result)
        result = _normalize_scalecua_result(result_type, result)
        return await _repair_reference_paths(
            eval_env,
            result,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    getter = judges.resolve_getter(result_type, runtime_split)
    if getter is not None:
        result = await judges.call_overlay_getter(getter, eval_env, config, cache_dir)
        result = await _repair_generated_xlsx_formula_result(eval_env, config, result)
        result = await _repair_generated_xlsx_professor_notes_result(eval_env, config, result)
        result = await _repair_generated_terminal_gsettings_result(eval_env, config, result)
        result = await _repair_generated_pptx_result(eval_env, config, result)
        result = await _repair_generated_sar_report_result(eval_env, config, result)
        result = await _repair_clipboard_result(eval_env, config, result)
        result = await _repair_chrome_extension_result(eval_env, config, result)
        result = await _repair_chrome_experiments_result(eval_env, config, result)
        result = await _repair_chrome_third_party_cookies_result(eval_env, config, result)
        result = await _repair_chrome_settings_pref_result(eval_env, result_type, result)
        result = _normalize_scalecua_result(result_type, result)
        return await _repair_reference_paths(
            eval_env,
            result,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    result = await base_runner._get_result(eval_env.computer, config, cache_dir)
    result = await _repair_clipboard_result(eval_env, config, result)
    result = await _repair_chrome_extension_result(eval_env, config, result)
    result = await _repair_chrome_profile_name_result(eval_env, result_type, result)
    result = _normalize_scalecua_result(result_type, result)
    return await _repair_reference_paths(
        eval_env,
        result,
        cache_dir,
        reference_keys=_REFERENCE_METRIC_KEYS,
        repair_home_user=True,
    )


async def _repair_generated_pptx_result(eval_env, config: dict, result: Any) -> Any:
    raw_result_type = str(config.get("type") or "")
    result_type = _base_result_type(raw_result_type)
    if result_type not in {
        "pptx_specific_text",
        "pptx_slide_subtitle",
        "pptx_textbox_fonts",
        "pptx_slide_title_and_bg",
    }:
        return result
    path = config.get("path") or config.get("ppt_file_path")
    if not isinstance(path, str):
        return result
    try:
        file_bytes = await eval_env.computer.interface.read_bytes(path)
    except Exception:
        return result
    if not file_bytes:
        return result
    # pptx parsing is CPU/disk-bound — offload (see _extract_accessibility_visible_text).
    return await asyncio.to_thread(
        _repair_generated_pptx_result_sync, raw_result_type, result_type, config, result, file_bytes
    )


def _repair_generated_pptx_result_sync(
    raw_result_type: str, result_type: str, config: dict, result: Any, file_bytes: bytes
) -> Any:
    tmp_path = ""
    try:
        from pptx import Presentation

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        presentation = Presentation(tmp_path)
        slide_index = int(config.get("slide_index", 0) or 0)
        if slide_index < 0 or slide_index >= len(presentation.slides):
            return result
        slide = presentation.slides[slide_index]
        if result_type == "pptx_specific_text":
            repaired = dict(result) if isinstance(result, dict) else {"text": ""}
            repaired["all_text"] = _pptx_slide_all_text(slide)
            return repaired
        if result_type == "pptx_textbox_fonts":
            repaired = dict(result) if isinstance(result, dict) else {}
            for key, shape_key in (
                ("textbox1_font_size", "textbox1_shape_idx"),
                ("textbox2_font_size", "textbox2_shape_idx"),
            ):
                if repaired.get(key) is not None:
                    continue
                try:
                    shape = slide.shapes[int(config.get(shape_key))]
                except Exception:
                    continue
                size = _pptx_shape_first_font_size(shape)
                if size is not None:
                    repaired[key] = size
            return repaired
        if result_type == "pptx_slide_title_and_bg":
            repaired = dict(result) if isinstance(result, dict) else {}
            repaired.setdefault("all_text", _pptx_slide_all_text(slide))
            if not repaired.get("bg_color"):
                bg = _pptx_slide_bg_rgb(slide)
                if bg is not None:
                    repaired["bg_color"] = bg
            return repaired
        if (
            raw_result_type == "pptx_slide_subtitle__c9f628a8"
            and not str(result or "").strip()
        ):
            return _pptx_slide_first_non_placeholder_text(slide) or result
    except Exception:
        return result
    finally:
        if tmp_path:
            with contextlib.suppress(OSError):
                os.unlink(tmp_path)
    return result


def _pptx_slide_all_text(slide: Any) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        text = getattr(shape, "text", None)
        if isinstance(text, str) and text.strip():
            texts.append(text.strip())
    return texts


def _pptx_slide_first_non_placeholder_text(slide: Any) -> str:
    placeholder_text = {
        "click to add title",
        "click to add subtitle",
        "click to add text",
    }
    for text in _pptx_slide_all_text(slide):
        normalized = re.sub(r"\s+", " ", text).strip()
        if normalized.lower() not in placeholder_text:
            return normalized
    return ""


def _pptx_shape_first_font_size(shape: Any) -> float | None:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is not None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                size = getattr(getattr(run, "font", None), "size", None)
                if size is not None:
                    return float(size.pt)
            paragraph_font = getattr(paragraph, "font", None)
            size = getattr(paragraph_font, "size", None)
            if size is not None:
                return float(size.pt)
    return _pptx_shape_xml_font_size(shape)


def _pptx_shape_xml_font_size(shape: Any) -> float | None:
    element = getattr(shape, "_element", None)
    xml = getattr(element, "xml", None)
    if not isinstance(xml, str):
        return None
    try:
        root = ET.fromstring(xml)
    except ET.ParseError:
        return None
    for elem in root.iter():
        if not elem.tag.endswith(("}rPr", "}defRPr")):
            continue
        size = elem.attrib.get("sz")
        if size:
            try:
                return float(size) / 100.0
            except ValueError:
                continue
    return None


def _pptx_slide_bg_rgb(slide: Any) -> tuple[int, int, int] | None:
    fill = getattr(getattr(slide, "background", None), "fill", None)
    color = getattr(getattr(fill, "fore_color", None), "rgb", None)
    rgb = judges._rgb_tuple_from_color(color)
    if rgb is not None:
        return rgb
    element = getattr(getattr(slide, "background", None), "_element", None)
    xml = getattr(element, "xml", None)
    if not isinstance(xml, str):
        return None
    try:
        root = ET.fromstring(xml)
    except ET.ParseError:
        return None
    for elem in root.iter():
        if elem.tag.endswith("}srgbClr"):
            value = elem.attrib.get("val")
            rgb = judges._rgb_tuple_from_color(value)
            if rgb is not None:
                return rgb
    return None


async def _repair_generated_xlsx_formula_result(eval_env, config: dict, result: Any) -> Any:
    if result not in (None, ""):
        return result
    result_type = _base_result_type(str(config.get("type") or ""))
    if result_type not in {"xlsx_cell_value", "xlsx_percentage_value"}:
        return result
    path = config.get("path")
    cell_address = config.get("cell")
    if not isinstance(path, str) or not isinstance(cell_address, str):
        return result
    try:
        file_bytes = await eval_env.computer.interface.read_bytes(path)
    except Exception:
        return result
    if not file_bytes:
        return result
    # openpyxl load_workbook is CPU-heavy — offload.
    return await asyncio.to_thread(
        _repair_generated_xlsx_formula_result_sync, config, result, cell_address, file_bytes
    )


def _repair_generated_xlsx_formula_result_sync(
    config: dict, result: Any, cell_address: str, file_bytes: bytes
) -> Any:
    tmp_path = ""
    try:
        import openpyxl

        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
            tmp.write(file_bytes)
            tmp_path = tmp.name
        workbook = openpyxl.load_workbook(tmp_path, data_only=False)
        sheet = _select_xlsx_sheet(workbook, config.get("sheet", 0))
        formula = sheet[cell_address].value
        if not isinstance(formula, str) or not formula.lstrip().startswith("="):
            workbook.close()
            return result
        value = _evaluate_generated_xlsx_formula(workbook, sheet, formula)
        workbook.close()
        return value if value is not None else result
    except Exception:
        return result
    finally:
        if tmp_path:
            with contextlib.suppress(OSError):
                os.unlink(tmp_path)


async def _repair_generated_xlsx_professor_notes_result(
    eval_env,
    config: dict,
    result: Any,
) -> Any:
    if str(config.get("type") or "") != "xlsx_cells_dict__a9a82c07":
        return result
    path = config.get("path")
    if not isinstance(path, str):
        return result
    try:
        file_bytes = await eval_env.computer.interface.read_bytes(path)
    except Exception:
        return result
    if not file_bytes:
        return result
    # openpyxl load_workbook is CPU-heavy — offload.
    return await asyncio.to_thread(
        _repair_generated_xlsx_professor_notes_result_sync, config, result, file_bytes
    )


def _repair_generated_xlsx_professor_notes_result_sync(
    config: dict, result: Any, file_bytes: bytes
) -> Any:
    try:
        import openpyxl
        from io import BytesIO

        workbook = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)
        sheet = _select_xlsx_sheet(workbook, config.get("sheet", 0))
        header = sheet["G2"].value
        professors: list[dict[str, Any]] = []
        for row_num in range(3, 1001):
            name = sheet[f"B{row_num}"].value
            if name is None or str(name).strip() == "":
                break
            notes = sheet[f"G{row_num}"].value
            professors.append(
                {
                    "row": row_num,
                    "name": str(name).strip(),
                    "notes": str(notes).strip() if notes is not None else "",
                }
            )
        workbook.close()
        return {"header": header, "professors": professors}
    except Exception:
        return result


async def _repair_generated_terminal_gsettings_result(
    eval_env,
    config: dict,
    result: Any,
) -> Any:
    """Read generated terminal setting evaluators directly from gsettings.

    SCALE-CUA generated tasks run a gsettings command in a terminal during
    postconfig and then read terminal output. In the container this adds a UI
    synchronization dependency that is not part of the semantic check.
    """
    result_type = _base_result_type(str(config.get("type") or ""))
    key_by_type = {
        "terminal_cursor_shape": "cursor-shape",
        "terminal_scrollback": "scrollback-lines",
    }
    if result_type in key_by_type:
        stdout = await _read_default_terminal_profile_key(
            eval_env,
            key_by_type[result_type],
        )
        return stdout if stdout else result
    if result_type == "terminal_profile_name":
        stdout = await _read_default_terminal_profile_name(eval_env)
        return stdout if stdout else result
    if result_type == "terminal_color_scheme":
        stdout = await _read_default_terminal_color_scheme(eval_env)
        return stdout if stdout else result
    return result


async def _read_default_terminal_profile_key(eval_env, key: str) -> str | None:
    command = (
        "profile=$(gsettings get org.gnome.Terminal.ProfilesList default | tr -d \"'\"); "
        + f"gsettings get org.gnome.Terminal.Legacy.Profile:/org/gnome/terminal/legacy/profiles:/:$profile/ {shlex.quote(key)}"
    )
    return await _run_eval_shell_stdout(eval_env, command)


async def _read_default_terminal_profile_name(eval_env) -> str | None:
    command = (
        "profile=$(gsettings get org.gnome.Terminal.ProfilesList default | tr -d \"'\"); "
        + "printf '%s\n' \"$profile\"; "
        + "gsettings get org.gnome.Terminal.Legacy.Profile:/org/gnome/terminal/legacy/profiles:/:$profile/ visible-name"
    )
    return await _run_eval_shell_stdout(eval_env, command)


async def _read_default_terminal_color_scheme(eval_env) -> str | None:
    command = (
        "profile=$(gsettings get org.gnome.Terminal.ProfilesList default | tr -d \"'\"); "
        + "schema=org.gnome.Terminal.Legacy.Profile:/org/gnome/terminal/legacy/profiles:/:$profile/; "
        + "printf 'use-theme-colors:\n'; "
        + "gsettings get \"$schema\" use-theme-colors; "
        + "printf 'background-color:\n'; "
        + "gsettings get \"$schema\" background-color; "
        + "printf 'foreground-color:\n'; "
        + "gsettings get \"$schema\" foreground-color"
    )
    return await _run_eval_shell_stdout(eval_env, command)


async def _run_eval_shell_stdout(eval_env, command: str) -> str | None:
    try:
        proc = await eval_env.computer.interface.run_command(
            "bash -lc " + shlex.quote(command),
            timeout=10,
        )
    except Exception:
        return None
    if getattr(proc, "returncode", 0) not in (0, None):
        return None
    stdout = getattr(proc, "stdout", "") or ""
    return stdout.strip() or None


async def _get_generated_xls_cell_value(eval_env, config: dict) -> dict[str, Any]:
    path = config.get("path")
    if not isinstance(path, str) or not path:
        return {"error": "missing path"}
    sheet = int(config.get("sheet", 0) or 0)
    row = int(config.get("row", 0) or 0)
    col = int(config.get("col", 0) or 0)
    command = f"""
set -e
tmpdir="$(mktemp -d)"
cleanup() {{ rm -rf "$tmpdir"; }}
trap cleanup EXIT
soffice --headless --convert-to xlsx --outdir "$tmpdir" {shlex.quote(path)} >/dev/null 2>&1
python3 - "$tmpdir" <<'PY'
import json
import sys
from pathlib import Path

import openpyxl

tmpdir = Path(sys.argv[1])
paths = sorted(tmpdir.glob("*.xlsx"))
if not paths:
    print("__SCALECUA_XLS_CELL__" + json.dumps({{"error": "converted xlsx not found"}}))
    raise SystemExit(0)
wb = openpyxl.load_workbook(paths[0], data_only=True)
try:
    ws = wb.worksheets[{sheet}]
    value = ws.cell(row={row + 1}, column={col + 1}).value
finally:
    wb.close()
print("__SCALECUA_XLS_CELL__" + json.dumps({{"value": value}}))
PY
"""
    try:
        # This getter is the ONLY office-python evaluator that runs IN the
        # container (it needs `soffice` to convert legacy .xls → .xlsx before
        # `openpyxl` reads it; every other openpyxl/pptx path runs host-side in
        # the eval process). `openpyxl` lives in the /opt/env venv; a bare `python3`
        # resolves to /opt/env/venv/bin/python3 via the exec-stdio server's PATH (the
        # system python3 on /usr/bin has no openpyxl), and soffice on /usr/bin is
        # reachable too.
        result = await eval_env.computer.interface.run_command(
            "bash -lc " + shlex.quote(command),
            timeout=60,
        )
    except Exception as exc:
        return {"error": str(exc)}
    stdout = getattr(result, "stdout", "") or ""
    marker = "__SCALECUA_XLS_CELL__"
    if marker not in stdout:
        stderr = getattr(result, "stderr", "") or ""
        return {"error": (stderr or stdout or "xls conversion produced no output").strip()}
    payload = stdout.rsplit(marker, 1)[1].strip().splitlines()[0]
    try:
        parsed = json.loads(payload)
    except json.JSONDecodeError:
        return {"error": f"invalid xls cell payload: {payload[:200]}"}
    return parsed if isinstance(parsed, dict) else {"error": "invalid xls cell result"}


def _select_xlsx_sheet(workbook: Any, sheet: Any):
    if isinstance(sheet, int):
        return workbook.worksheets[sheet]
    if isinstance(sheet, str):
        return workbook[sheet]
    return workbook.worksheets[0]


def _evaluate_generated_xlsx_formula(workbook: Any, sheet: Any, formula: str) -> float | int | None:
    text = str(formula or "").strip()
    if text.startswith("="):
        text = text[1:].strip()
    text = text.replace(";", ",")
    division = _split_top_level_operator(text, "/")
    if division is not None:
        left = _evaluate_generated_xlsx_formula(workbook, sheet, division[0])
        right = _evaluate_generated_xlsx_formula(workbook, sheet, division[1])
        if isinstance(left, (int, float)) and isinstance(right, (int, float)) and right:
            return float(left) / float(right)
        return None
    match = re.fullmatch(r"SUM\((?P<args>.+)\)", text, flags=re.IGNORECASE)
    if match:
        return sum(
            _iter_xlsx_numeric_values(workbook, sheet, arg.strip())
            for arg in _split_formula_args(match.group("args"))
        )
    match = re.fullmatch(r"MEDIAN\((?P<args>.+)\)", text, flags=re.IGNORECASE)
    if match:
        values: list[float] = []
        for arg in _split_formula_args(match.group("args")):
            values.extend(_iter_xlsx_numeric_values(workbook, sheet, arg.strip(), as_list=True))
        if values:
            return float(statistics.median(values))
        return None
    match = re.fullmatch(
        r"SUMIF\((?P<range>.+?),(?P<criteria>.+?),(?P<sum_range>.+)\)",
        text,
        flags=re.IGNORECASE,
    )
    if match:
        return _evaluate_xlsx_sumif(
            workbook,
            sheet,
            match.group("range").strip(),
            match.group("criteria").strip(),
            match.group("sum_range").strip(),
        )
    return None


def _split_top_level_operator(text: str, operator: str) -> tuple[str, str] | None:
    depth = 0
    in_quote = False
    for index, char in enumerate(text):
        if char == '"':
            in_quote = not in_quote
            continue
        if in_quote:
            continue
        if char == "(":
            depth += 1
        elif char == ")":
            depth = max(0, depth - 1)
        elif char == operator and depth == 0:
            left = text[:index].strip()
            right = text[index + 1 :].strip()
            if left and right:
                return left, right
    return None


def _split_formula_args(text: str) -> list[str]:
    args: list[str] = []
    current: list[str] = []
    in_quote = False
    for char in text:
        if char == '"':
            in_quote = not in_quote
        if char == "," and not in_quote:
            args.append("".join(current).strip())
            current = []
        else:
            current.append(char)
    tail = "".join(current).strip()
    if tail:
        args.append(tail)
    return args


def _iter_xlsx_numeric_values(
    workbook: Any,
    current_sheet: Any,
    reference: str,
    *,
    as_list: bool = False,
) -> float | list[float]:
    values: list[float] = []
    for value in _iter_xlsx_reference_values(workbook, current_sheet, reference):
        if isinstance(value, (int, float)):
            values.append(float(value))
    return values if as_list else sum(values)


def _evaluate_xlsx_sumif(
    workbook: Any,
    current_sheet: Any,
    criteria_range: str,
    criteria: str,
    sum_range: str,
) -> float:
    expected = _unquote_formula_string(criteria)
    criteria_values = list(_iter_xlsx_reference_values(workbook, current_sheet, criteria_range))
    sum_values = list(_iter_xlsx_reference_values(workbook, current_sheet, sum_range))
    total = 0.0
    for value, amount in zip(criteria_values, sum_values):
        if str(value) == expected and isinstance(amount, (int, float)):
            total += float(amount)
    return total


def _unquote_formula_string(value: str) -> str:
    value = value.strip()
    if len(value) >= 2 and value[0] == value[-1] == '"':
        return value[1:-1]
    return value


def _iter_xlsx_reference_values(workbook: Any, current_sheet: Any, reference: str):
    from openpyxl.utils.cell import column_index_from_string, range_boundaries

    ref = reference.strip().replace("$", "")
    sheet = current_sheet
    if "!" in ref:
        sheet_name, ref = ref.split("!", 1)
        sheet = workbook[sheet_name.strip("'")]
    if re.fullmatch(r"[A-Za-z]+:[A-Za-z]+", ref):
        start_col, end_col = ref.split(":", 1)
        min_col = column_index_from_string(start_col)
        max_col = column_index_from_string(end_col)
        min_row = 1
        max_row = sheet.max_row
    elif re.fullmatch(r"[A-Za-z]+", ref):
        min_col = max_col = column_index_from_string(ref)
        min_row = 1
        max_row = sheet.max_row
    else:
        min_col, min_row, max_col, max_row = range_boundaries(ref)
        min_row = min_row or 1
        max_row = max_row or sheet.max_row
        min_col = min_col or 1
        max_col = max_col or sheet.max_column
    for row in sheet.iter_rows(
        min_row=min_row,
        max_row=max_row,
        min_col=min_col,
        max_col=max_col,
        values_only=True,
    ):
        for value in row:
            yield value




# ---------------------------------------------------------------------------
# ScaleCUA baked bashrc-alias check hardening (parity with the osworld os_05 fix).
#
# Several ScaleCUA verigen `vm_command_line` evaluators detect a `~/.bashrc`
# alias by raw-greping for an exact quoted literal, e.g.
#     grep -q "alias ll='ls -la'" ~/.bashrc && echo 1 || echo 0
# which FNs a correct agent writing the semantically identical
# `alias ll="ls -la"` (double quotes) or `alias ll='ls -al'` (bash does NOT
# reorder alias bodies, so a different short-flag ORDER is the same flag SET).
# A sibling shape greps with an optional-DOUBLE-quote regex
#     grep -qE "alias python=(\")?python3(\")? *$" ~/.bashrc && echo 1 || echo 0
# which conversely FNs the single-quote `alias python='python3'`.
#
# These are upstream DATA (baked into data/{rl,train}.jsonl), not reachable by
# the osworld synth patch, so we harden them at eval time. The rewrite sources
# the rc and lets the bash `alias` builtin canonicalize quoting before the
# comparison; for an `ls -<shortflags>` body it collapses the flags to a
# sorted-unique SET so any permutation matches while a wrong set (missing OR
# extra flag) still fails -- a conservative FN, never an FP. Predicate is over
# the command SHAPE (any `alias X='...'` grep), not per-hash.
# ---------------------------------------------------------------------------
_ALIAS_LITERAL_CHECK_RE = re.compile(
    r'^\s*grep\s+-q\w*\s+"alias\s+(?P<name>[A-Za-z_]\w*)=\'(?P<body>[^\']*)\'"\s+'
    r'(?P<path>\S+)\s*&&\s*echo\s+1\s*\|\|\s*echo\s+0\s*$'
)
_ALIAS_OPTQUOTE_CHECK_RE = re.compile(
    r'^\s*grep\s+-qE\s+"alias\s+(?P<name>[A-Za-z_]\w*)='
    r'\(\\"\)\?(?P<body>[\w.+-]+)\(\\"\)\?\s*\*\$"\s+'
    r'(?P<path>\S+)\s*&&\s*echo\s+1\s*\|\|\s*echo\s+0\s*$'
)
_ALIAS_LS_SHORTFLAGS_RE = re.compile(r"^ls((?:\s+-[A-Za-z]+)+)(?:\s+--\S+)*\s*$")


def _build_alias_check_command(name: str, body: str, path: str) -> str:
    """Quote-tolerant (and, for `ls -<shortflags>` bodies, flag-order-tolerant)
    replacement that still echoes exactly ``1``/``0``."""
    ls_flags = _ALIAS_LS_SHORTFLAGS_RE.match(body)
    if ls_flags:
        canon = "".join(sorted(set(re.findall(r"[A-Za-z]", ls_flags.group(1)))))
        inner = (
            f"LC_ALL=C; "
            # DON'T `source` the rc — Ubuntu's ~/.bashrc opens with the
            # `case $- in *i*) ;; *) return;; esac` non-interactive guard, which
            # `return`s before any alias line under `bash -c` ($- has no `i`) ->
            # the alias is never defined -> the check FALSE-NEGATIVES every task
            # (verified on gate-zzh). Extract just the alias line and `eval` it,
            # bypassing the guard while still letting bash canonicalize quoting.
            f'l=$(grep -hE "^[[:space:]]*alias {name}=" {path} 2>/dev/null '
            f'| tail -n1); eval "$l" 2>/dev/null; '
            f"a=$(alias {name} 2>/dev/null); "
            f'case "$a" in *ls*) '
            f"f=$(printf '%s\\n' \"$a\" | grep -oE '(^| )-[A-Za-z]+' | "
            f"grep -oE '[A-Za-z]' | LC_ALL=C sort -u | tr -d '\\n'); "
            f'[ "$f" = "{canon}" ] && echo 1 || echo 0 ;; '
            f"*) echo 0 ;; esac"
        )
    else:
        expected = f"alias {name}='{body}'"
        inner = (
            f"LC_ALL=C; "
            # DON'T `source` the rc — Ubuntu's ~/.bashrc opens with the
            # `case $- in *i*) ;; *) return;; esac` non-interactive guard, which
            # `return`s before any alias line under `bash -c` ($- has no `i`) ->
            # the alias is never defined -> the check FALSE-NEGATIVES every task
            # (verified on gate-zzh). Extract just the alias line and `eval` it,
            # bypassing the guard while still letting bash canonicalize quoting.
            f'l=$(grep -hE "^[[:space:]]*alias {name}=" {path} 2>/dev/null '
            f'| tail -n1); eval "$l" 2>/dev/null; '
            f"a=$(alias {name} 2>/dev/null); "
            f'[ "$a" = {shlex.quote(expected)} ] && echo 1 || echo 0'
        )
    return "bash -c " + shlex.quote(inner)


def _harden_alias_check_command(command: Any) -> Any:
    """Rewrite a brittle baked bashrc alias-grep evaluator to a canonicalizing
    form. Returns *command* unchanged when it is not an alias-check."""
    if not isinstance(command, str):
        return command
    for pattern in (_ALIAS_LITERAL_CHECK_RE, _ALIAS_OPTQUOTE_CHECK_RE):
        m = pattern.match(command)
        if m:
            return _build_alias_check_command(m["name"], m["body"], m["path"])
    return command


def _normalize_scalecua_command_config(config: dict) -> dict:
    if config.get("type") not in {"vm_command_line", "vm_command_error"}:
        return config
    command = config.get("command")
    fixed = _normalize_python_command(command, shell=bool(config.get("shell")))
    fixed = _repair_git_root_commit_diff_tree_command(fixed)
    fixed = _harden_alias_check_command(fixed)
    # A session-bus getter (gsettings/dconf/pactl) runs as the desktop user (the
    # exec-stdio server's identity) with the desktop-session env already available —
    # so no manual session prefix or VM-path rewrite here.
    if fixed is command:
        return config
    return {**config, "command": fixed}


def _normalize_python_command(command: Any, *, shell: bool) -> Any:
    """Match OSWorld VM's ``python -> python3`` shape inside ScaleCUA eval."""
    if isinstance(command, list):
        if command and str(command[0]) == "python":
            return ["python3", *command[1:]]
        return command
    if isinstance(command, tuple):
        if command and str(command[0]) == "python":
            return ("python3", *command[1:])
        return command
    if not isinstance(command, str) or shell:
        return command
    try:
        parts = shlex.split(command)
    except ValueError:
        return command
    if not parts or parts[0] != "python":
        return command
    return " ".join(shlex.quote(part) for part in ["python3", *parts[1:]])


def _repair_git_root_commit_diff_tree_command(command: Any) -> Any:
    """Repair a generated ScaleCUA git evaluator that inspects a root commit."""
    if isinstance(command, list):
        fixed = [_repair_git_root_commit_diff_tree_command(part) for part in command]
        return fixed if fixed != command else command
    if isinstance(command, tuple):
        fixed = tuple(_repair_git_root_commit_diff_tree_command(part) for part in command)
        return fixed if fixed != command else command
    if not isinstance(command, str):
        return command
    if "first_commit_only_readme" not in command or '"$first_hash"' not in command:
        return command
    target = 'git diff-tree --no-commit-id --name-only -r "$first_hash"'
    if target not in command:
        return command
    return command.replace(
        target,
        'git diff-tree --root --no-commit-id --name-only -r "$first_hash"',
    )


def _is_django_repo_file_count_result(config: dict) -> bool:
    return (
        config.get("type") == "dir_file_count__00db2192"
        and str(config.get("dir_path") or "").rstrip("/") == "/home/user/django"
        and str(config.get("pattern") or "") == "*.py"
    )


async def _get_django_repo_file_count(eval_env) -> int:
    command = r"""
dir=/home/user/django
if [ ! -d "$dir/.git" ] || [ ! -f "$dir/django/__init__.py" ]; then
  echo 0
  exit 0
fi
remote="$(git -C "$dir" config --get remote.origin.url 2>/dev/null || true)"
case "$remote" in
  *github.com/django/django*|*github.com:django/django*) ;;
  *) echo 0; exit 0 ;;
esac
find "$dir" -type f -name '*.py' 2>/dev/null | wc -l
"""
    result = await eval_env.computer.interface.run_command(
        "bash -lc " + shlex.quote(command),
        timeout=20,
    )
    try:
        return int(str(getattr(result, "stdout", "") or "0").strip())
    except ValueError:
        return 0


def _normalize_scalecua_result(result_type: str, result: Any) -> Any:
    if _base_result_type(result_type) in {
        "default_search_engine",
        "chrome_search_engine",
    }:
        return _normalize_default_search_engine(result)
    if _base_result_type(result_type).startswith("chrome_homepage") and isinstance(result, str):
        return judges._normalize_web_url_for_task(result)
    return result


def _normalize_default_search_engine(result: Any) -> Any:
    if not isinstance(result, str):
        return result
    normalized = re.sub(r"\s+", " ", result.strip()).lower()
    compact = re.sub(r"[^a-z0-9]+", "", normalized)
    if "startpage.com" in normalized or compact in {"startpage", "startpagecom"}:
        return "Startpage"
    if compact in {"duckduckgo", "duckduckgocom"}:
        return "DuckDuckGo"
    if "duckduckgo.com" in normalized:
        return "DuckDuckGo"
    if normalized.startswith("yahoo") or "search.yahoo.com" in normalized:
        return "Yahoo!"
    if compact in {"bing", "microsoftbing", "bingdefault", "microsoftbingdefault"}:
        return "Bing"
    if normalized.startswith(("bing", "microsoft bing")) or "bing.com" in normalized:
        return "Bing"
    return result


async def _repair_clipboard_result(eval_env, config: dict, result: Any) -> Any:
    if not _is_clipboard_result(config):
        return result
    if isinstance(result, str):
        if result.strip():
            return result
        clipboard = await _read_vm_clipboard(eval_env)
        return clipboard if clipboard is not None else result
    if isinstance(result, dict) and "clipboard" in result:
        if str(result.get("clipboard") or "").strip():
            return result
        clipboard = await _read_vm_clipboard(eval_env)
        if clipboard is None:
            return result
        fixed = dict(result)
        fixed["clipboard"] = clipboard
        return fixed
    return result


def _is_clipboard_result(config: dict) -> bool:
    result_type = _base_result_type(str(config.get("type") or ""))
    if result_type in {"clipboard_content", "file_copy_and_clipboard"}:
        return True
    if result_type != "vm_command_line":
        return False
    return _command_reads_clipboard(config.get("command"))


def _command_reads_clipboard(command: Any) -> bool:
    if isinstance(command, (list, tuple)):
        text = " ".join(str(part) for part in command)
    else:
        text = str(command or "")
    normalized = re.sub(r"\s+", " ", text.lower())
    return "clipboard" in normalized and (
        "xsel" in normalized
        or "xclip" in normalized
        or "wl-paste" in normalized
    )


async def _read_vm_clipboard(eval_env) -> str | None:
    command = r"""
set +e
emit_if_nonempty() {
  data="$1"
  if [ -n "$data" ]; then
    printf '__SCALECUA_CLIPBOARD_B64__'
    printf '%s' "$data" | base64 -w0
    printf '\n'
    exit 0
  fi
}
for display in "${DISPLAY:-}" :1 :0; do
  [ -z "$display" ] && continue
  DISPLAY="$display"; export DISPLAY
  emit_if_nonempty "$(timeout 2s xsel --clipboard --output 2>/dev/null || true)"
  emit_if_nonempty "$(timeout 2s xclip -selection clipboard -o 2>/dev/null || true)"
  emit_if_nonempty "$(timeout 2s wl-paste 2>/dev/null || true)"
  emit_if_nonempty "$(timeout 3s python3 - 2>/dev/null <<'PYEOF' || true
try:
    import gi
    gi.require_version("Gtk", "3.0")
    from gi.repository import Gdk, Gtk
    text = Gtk.Clipboard.get(Gdk.SELECTION_CLIPBOARD).wait_for_text()
    if text:
        print(text, end="")
except Exception:
    pass
PYEOF
)"
  emit_if_nonempty "$(timeout 3s python3 - 2>/dev/null <<'PYEOF' || true
try:
    import tkinter as tk
    root = tk.Tk()
    root.withdraw()
    text = root.clipboard_get()
    root.destroy()
    if text:
        print(text, end="")
except Exception:
    pass
PYEOF
)"
done
"""
    try:
        proc = await eval_env.computer.interface.run_command(
            "bash -lc " + shlex.quote(command),
            timeout=20,
        )
    except Exception:
        return None
    stdout = getattr(proc, "stdout", "") or ""
    marker = "__SCALECUA_CLIPBOARD_B64__"
    if marker not in stdout:
        return None
    encoded = stdout.split(marker, 1)[1].strip().splitlines()[0]
    try:
        return base64.b64decode(encoded).decode("utf-8", errors="replace")
    except Exception:
        return None


async def _repair_generated_sar_report_result(eval_env, config: dict, result: Any) -> Any:
    """Repair VeriGen sar getters whose generated Python `-c` command is invalid."""

    result_type = _base_result_type(str(config.get("type") or ""))
    if result_type == "sar_cpu_report_state":
        return await _read_sar_cpu_report_state(eval_env)
    if result_type == "sar_disk_report_state":
        return await _read_sar_disk_report_state(eval_env)
    return result


async def _read_sar_cpu_report_state(eval_env) -> dict[str, Any]:
    content = await _read_vm_text_file(eval_env, "/home/user/Desktop/System_Resources_Report.txt")
    if content is None:
        return {"exists": False, "cpu_lines": 0, "has_header": False}
    lines = content.splitlines()
    return {
        "exists": True,
        "cpu_lines": sum(1 for line in lines if " all " in line),
        "has_header": any("%user" in line for line in lines),
    }


async def _read_sar_disk_report_state(eval_env) -> dict[str, Any]:
    content = await _read_vm_text_file(eval_env, "/home/user/Desktop/Disk_IO_Report.txt")
    if content is None:
        return {"exists": False, "has_disk_data": False, "line_count": 0}
    lines = content.splitlines()
    return {
        "exists": True,
        "has_disk_data": any(
            "tps" in line.lower()
            or "rkb/s" in line.lower()
            or "wkb/s" in line.lower()
            or "DEV" in line
            for line in lines
        ),
        "line_count": len(lines),
    }


async def _read_vm_text_file(eval_env, path: str) -> str | None:
    try:
        data = await eval_env.computer.interface.read_bytes(path)
    except Exception:
        return None
    if not data:
        return None
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("utf-8", errors="replace")


_VSCODE_THEME_VALUE_ALIASES = {
    "Light+ (default light)": "Default Light+",
    "Dark+ (default dark)": "Default Dark+",
}


def _normalize_vscode_theme_aliases(value: Any) -> Any:
    if isinstance(value, dict):
        changed = False
        normalized: dict[Any, Any] = {}
        for key, item in value.items():
            if (
                key == "workbench.colorTheme"
                and isinstance(item, str)
                and item in _VSCODE_THEME_VALUE_ALIASES
            ):
                normalized[key] = _VSCODE_THEME_VALUE_ALIASES[item]
                changed = True
                continue
            normalized_item = _normalize_vscode_theme_aliases(item)
            if normalized_item is not item:
                changed = True
            normalized[key] = normalized_item
        return normalized if changed else value
    if isinstance(value, list):
        normalized_items = [_normalize_vscode_theme_aliases(item) for item in value]
        if any(new is not old for new, old in zip(normalized_items, value)):
            return normalized_items
    return value


def _normalize_scalecua_expected_rules(fn_name: str, expected_data: Any) -> Any:
    expected_data = _normalize_vscode_theme_aliases(expected_data)
    if fn_name == "check_ext_version__b28217ae" and isinstance(expected_data, dict):
        if set(expected_data) == {"expected"}:
            return expected_data["expected"]
    if fn_name == "check_pptx_shape_text__8b4cb395" and isinstance(expected_data, dict):
        if set(expected_data) == {"value"}:
            return expected_data["value"]
    if fn_name != "check_direct_json_object" or not isinstance(expected_data, dict):
        return expected_data
    expected = expected_data.get("expected")
    if not isinstance(expected, dict) or "ignore_list_order" not in expected:
        return expected_data

    fixed = copy.deepcopy(expected_data)
    fixed_expected = fixed.get("expected", {})
    if isinstance(fixed_expected, dict):
        fixed.setdefault("ignore_list_order", bool(fixed_expected["ignore_list_order"]))
        fixed_expected.pop("ignore_list_order", None)
    return fixed


def _normalize_chrome_url_metric_pair(
    result_config: dict[str, Any],
    result_data: Any,
    expected_data: Any,
) -> tuple[Any, Any]:
    result_type = _base_result_type(str(result_config.get("type") or ""))
    if not result_type.startswith("chrome_homepage"):
        return result_data, expected_data
    normalized_result = (
        judges._normalize_web_url_for_task(result_data)
        if isinstance(result_data, str)
        else result_data
    )
    if isinstance(expected_data, dict) and isinstance(expected_data.get("expected"), str):
        normalized_expected = dict(expected_data)
        normalized_expected["expected"] = judges._normalize_web_url_for_task(
            expected_data["expected"]
        )
        return normalized_result, normalized_expected
    if isinstance(expected_data, str):
        return normalized_result, judges._normalize_web_url_for_task(expected_data)
    return normalized_result, expected_data


def _alias_chrome_profile_config(value):
    if isinstance(value, str):
        return judges._alias_chrome_profile_path(value)
    if isinstance(value, list):
        return [_alias_chrome_profile_config(item) for item in value]
    if isinstance(value, tuple):
        return tuple(_alias_chrome_profile_config(item) for item in value)
    if isinstance(value, dict):
        return {
            key: _alias_chrome_profile_config(item)
            for key, item in value.items()
        }
    return value


_CHROME_PROFILE_RESULT_TYPES = {
    "bookmarks",
    "check_dns_prefetch",
    "chrome_appearance_mode_ui",
    "chrome_color_scheme",
    "chrome_font_size",
    "cookie_data",
    "data_delete_automacally",
    "default_search_engine",
    "enable_do_not_track",
    "enable_enhanced_safety_browsing",
    "enable_safe_browsing",
    "enabled_experiments",
    "find_installed_extension_name",
    "find_unpacked_extension_path",
    "history",
    "new_startup_page",
    "profile_name",
}

_CHROME_GENERATED_EXTENSION_RESULT_TYPES = {
    "ext_count",
    "ext_description",
    "ext_enabled",
    "ext_enabled_status",
    "ext_info",
    "ext_location",
    "ext_manifest",
    "ext_manifest_data",
    "ext_manifest_version",
    "ext_name_by_path",
    "ext_name_list",
    "ext_names",
    "ext_version",
}

_CHROME_PROFILE_RESULT_PREFIXES = (
    "block_third_party_cookies",
    "bookmark_",
    "bookmarks_",
    "chrome_appearance",
    "chrome_clear_on_exit",
    "chrome_color_scheme",
    "chrome_default_font_size",
    "chrome_default_search",
    "chrome_do_not_track",
    "chrome_download_",
    "chrome_enhanced_safe_browsing",
    "chrome_experiments",
    "chrome_ext_",
    "chrome_extension_",
    "chrome_extensions_",
    "chrome_fixed_font_size",
    "chrome_font_size",
    "chrome_homepage",
    "chrome_language",
    "chrome_location_blocked",
    "chrome_min_font",
    "chrome_newtab_startup",
    "chrome_notifications_blocked",
    "chrome_password_manager",
    "chrome_popups_blocked",
    "chrome_safe_browsing",
    "chrome_search_engine",
    "chrome_session_restore",
    "chrome_setting",
    "chrome_show_home_button",
    "chrome_standard_font_size",
    "chrome_startup",
    "chrome_third_party_cookies",
    "clear_cookies_on_exit",
    "cookies_clear_on_exit",
    "data_delete_automacally",
    "default_search_engine",
    "disable_safe_browsing",
    "do_not_track",
    "enable_enhanced_safety_browsing",
    "extension_",
    "new_startup_page",
    "password_manager_",
    "popup_blocker_",
    "safe_browsing_",
    "standard_safe_browsing",
    "third_party_cookies_",
    "unpacked_extension_",
    "webext_",
)

_CHROME_PROFILE_FILE_MARKERS = (
    "/home/user/chrome-data/",
    "/home/user/.config/google-chrome/",
    "/home/user/.config/chromium/",
    "/home/user/snap/chromium/common/chromium/",
)


def _base_result_type(result_type: str) -> str:
    return str(result_type or "").split("__", 1)[0]


def _is_scalecua_local_result_type(result_type: str) -> bool:
    base = _base_result_type(result_type)
    return (
        result_type in SCALECUA_LOCAL_RESULT_TYPES
        or base in SCALECUA_LOCAL_RESULT_TYPES
        or base in SCALECUA_LOCAL_HASHED_RESULT_TYPES
        or base in SCALECUA_LOCAL_CHROME_PREF_RESULT_TYPES
    )


def _flush_stats(eval_env) -> dict[str, dict[str, Any]]:
    stats = getattr(eval_env, "_scalecua_flush_stats", None)
    if not isinstance(stats, dict):
        stats = {}
        setattr(eval_env, "_scalecua_flush_stats", stats)
    return stats


def _flush_app_stats(eval_env, app: str) -> dict[str, Any]:
    stats = _flush_stats(eval_env)
    item = stats.get(app)
    if not isinstance(item, dict):
        item = {
            "needed": 0,
            "executed": 0,
            "skipped_already_flushed": 0,
            "result_types": [],
        }
        stats[app] = item
    return item


def _record_flush_needed(eval_env, app: str, result_type: str) -> None:
    item = _flush_app_stats(eval_env, app)
    item["needed"] += 1
    item["result_types"].append(str(result_type or ""))


def _record_flush_executed(eval_env, app: str) -> None:
    _flush_app_stats(eval_env, app)["executed"] += 1


def _record_flush_skipped(eval_env, app: str) -> None:
    _flush_app_stats(eval_env, app)["skipped_already_flushed"] += 1


def _flush_stats_snapshot(eval_env) -> dict[str, dict[str, Any]]:
    stats = getattr(eval_env, "_scalecua_flush_stats", None)
    if not isinstance(stats, dict):
        return {}
    return {
        app: {
            "needed": int(item.get("needed", 0)),
            "executed": int(item.get("executed", 0)),
            "skipped_already_flushed": int(item.get("skipped_already_flushed", 0)),
            "result_types": list(item.get("result_types", [])),
        }
        for app, item in sorted(stats.items())
        if isinstance(item, dict)
    }


def _flush_fired_counts_from_stats(stats: dict[str, dict[str, Any]]) -> dict[str, int]:
    return {
        app: int(item.get("needed", 0) or 0)
        for app, item in sorted(stats.items())
        if isinstance(item, dict)
    }


def _flush_counter_aliases(stats: dict[str, dict[str, Any]]) -> dict[str, int]:
    counts = _flush_fired_counts_from_stats(stats)
    return {
        "vlc_flush_fired": int(counts.get("vlc", 0)),
        "thunderbird_flush_fired": int(counts.get("thunderbird", 0)),
    }


def _needs_chrome_profile_flush(result_type: str, config: dict | None = None) -> bool:
    base = _base_result_type(result_type)
    if _chrome_settings_url_for_result(base):
        return False
    return (
        base in _CHROME_PROFILE_RESULT_TYPES
        or _is_chrome_extension_result(result_type)
        or base.startswith(_CHROME_PROFILE_RESULT_PREFIXES)
        or _mentions_chrome_profile_file(config)
    )


def _mentions_chrome_profile_file(value: Any) -> bool:
    if isinstance(value, str):
        lowered = value.lower()
        return any(marker in lowered for marker in _CHROME_PROFILE_FILE_MARKERS)
    if isinstance(value, dict):
        return any(_mentions_chrome_profile_file(item) for item in value.values())
    if isinstance(value, (list, tuple)):
        return any(_mentions_chrome_profile_file(item) for item in value)
    return False


def _is_chrome_extension_result(result_type: str) -> bool:
    result_type = _base_result_type(result_type)
    return (
        result_type in {"find_unpacked_extension_path", "find_installed_extension_name"}
        or result_type in _CHROME_GENERATED_EXTENSION_RESULT_TYPES
        or result_type.startswith("extension_")
        or result_type.startswith("chrome_extension_")
        or result_type.startswith("unpacked_extension_")
    )


def _needs_gimp_config_flush(result_type: str, config: dict | None = None) -> bool:
    base = _base_result_type(result_type)
    return (
        base in _GIMP_CONFIG_RESULT_TYPES
        or base.startswith(_GIMP_CONFIG_RESULT_PREFIXES)
        or _mentions_gimp_config_file(config)
    )


def _mentions_gimp_config_file(value: Any) -> bool:
    if isinstance(value, str):
        return any(marker in value for marker in _GIMP_CONFIG_FILE_MARKERS)
    if isinstance(value, dict):
        return any(_mentions_gimp_config_file(item) for item in value.values())
    if isinstance(value, (list, tuple)):
        return any(_mentions_gimp_config_file(item) for item in value)
    return False


def _best_effort_flush(fn):
    """Make a config-flush helper best-effort: a flush persists in-memory app
    state to disk before a getter reads it, but it is an OPTIMIZATION, not a
    correctness requirement — a flush that hangs/times out under load must NEVER
    fail the reward. (`_flush_vlc`'s `run_command(timeout=45)` has raised
    `TimeoutExpired`, which would otherwise propagate to reward 0.0 for a
    correct agent.) Swallow
    + log any error. Worst case the getter reads un-flushed (stale) state — exactly
    as if the flush had not run — which can only LOWER the score, never fabricate
    reward (no new FP). `Exception` only, so `asyncio.CancelledError` still
    propagates. Applies to all five `_flush_*` helpers, protecting every call site.
    """
    @functools.wraps(fn)
    async def _wrapped(eval_env, *args, **kwargs):
        try:
            return await fn(eval_env, *args, **kwargs)
        except Exception as exc:  # noqa: BLE001 — best-effort by design
            logger.warning(
                "config-flush %s failed (best-effort, ignored): %r", fn.__name__, exc
            )
    return _wrapped


@_best_effort_flush
async def _flush_gimp_config(eval_env) -> None:
    if getattr(eval_env, "_gimp_config_flushed", False):
        _record_flush_skipped(eval_env, "gimp")
        return
    command = r"""
sent_quit=0
for display in "${DISPLAY:-}" :1 :0; do
  [ -z "$display" ] && continue
  DISPLAY="$display"; export DISPLAY
  for wid in $(xdotool search --class 'Gimp' 2>/dev/null); do
    xdotool windowactivate --sync "$wid" 2>/dev/null || true
    xdotool key --window "$wid" ctrl+q 2>/dev/null || true
    sent_quit=1
    break 2
  done
done
if [ "$sent_quit" -eq 1 ]; then
  sleep 2
  for display in "${DISPLAY:-}" :1 :0; do
    [ -z "$display" ] && continue
    DISPLAY="$display"; export DISPLAY
    wid="$(xdotool search --name 'Quit GIMP' 2>/dev/null | head -1)"
    if [ -n "$wid" ]; then
      xdotool windowactivate "$wid" 2>/dev/null || true
      xdotool key alt+d 2>/dev/null || true
      break
    fi
  done
  sleep 1
  xdotool key Return 2>/dev/null || true
  for _ in $(seq 1 30); do
    if ! pgrep -x gimp >/dev/null 2>&1 && ! pgrep -x gimp-2.10 >/dev/null 2>&1; then
      break
    fi
    sleep 1
  done
  pkill -f gimp 2>/dev/null || true
  sleep 1
  sync
fi
"""
    await eval_env.computer.interface.run_command(
        "bash -lc " + shlex.quote(command),
        timeout=45,
    )
    _record_flush_executed(eval_env, "gimp")
    setattr(eval_env, "_gimp_config_flushed", True)


@_best_effort_flush
async def _flush_chrome_profile(eval_env) -> None:
    if getattr(eval_env, "_chrome_profile_flushed", False):
        _record_flush_skipped(eval_env, "chrome")
        return
    await eval_env.computer.interface.run_command(
        "xdotool key Tab 2>/dev/null || true; "
        "xdotool key Escape 2>/dev/null || true; "
        "sleep 1; "
        "for wid in $(xdotool search --class 'google-chrome' 2>/dev/null); do "
        "  xdotool windowactivate --sync \"$wid\" 2>/dev/null || true; "
        "  xdotool key --window \"$wid\" Alt+F4 2>/dev/null || true; "
        "done; "
        "for _ in $(seq 1 20); do "
        "  pgrep -x chrome >/dev/null 2>&1 || "
        "    pgrep -x google-chrome >/dev/null 2>&1 || "
        "    pgrep -x chromium >/dev/null 2>&1 || break; "
        "  sleep 0.25; "
        "done; "
        "pkill -TERM chrome || true; "
        "pkill -TERM google-chrome || true; "
        "pkill -TERM chromium || true; "
        "sleep 5; sync",
        timeout=25,
    )
    _record_flush_executed(eval_env, "chrome")
    setattr(eval_env, "_chrome_profile_flushed", True)


# --- Thunderbird / VLC config-flush scaffolding --------------------------------
#
# Mirror of the GIMP config-flush adapter (_flush_gimp_config): a graceful
# app-quit (ctrl+q -> wait-for-exit -> SIGTERM fallback, NEVER `kill -9`) that
# runs the app's own config-save path so a running-app GUI edit lands in the
# on-disk config (Thunderbird prefs.js / msgFilterRules.dat; VLC vlcrc) before
# the getter reads it.
#
# NOTE: during oracle replay the app has already been
# killed by `oracle_actions` (e.g. `pkill -9 thunderbird`), so `xdotool search`
# finds no window, `sent_quit` stays 0, and the whole body no-ops -- it MUST NOT
# disturb an already-passing fixture. The FN-recovery this enables is
# rollout-only.

_THUNDERBIRD_CONFIG_FILE_MARKERS = (
    "/home/user/.thunderbird/",
    "~/.thunderbird/",
    "prefs.js",
    "msgFilterRules.dat",
)
_THUNDERBIRD_CONFIG_RESULT_TYPES = {
    "check_thunderbird_prefs",
}
_THUNDERBIRD_CONFIG_RESULT_PREFIXES = (
    "check_tb_",
    "check_thunderbird",
    "thunderbird_config",
    "tb_prefs",
)

# VLC persists `vlcrc` only on a clean quit. `vlc_playing_info` (play/loop status
# read over the HTTP interface) is deliberately EXCLUDED: quitting VLC would tear
# down the very interface the getter reads (a launch-config issue, not a flush).
_VLC_CONFIG_FILE_MARKERS = (
    "/home/user/.config/vlc/vlcrc",
    ".config/vlc/vlcrc",
    "/vlcrc",
)
_VLC_CONFIG_RESULT_TYPES = {
    "vlc_config",
}
_VLC_CONFIG_RESULT_PREFIXES = (
    "vlc_config",
    "check_vlc_config",
    "check_vlcrc",
)


def _needs_thunderbird_config_flush(result_type: str, config: dict | None = None) -> bool:
    base = _base_result_type(result_type)
    return (
        base in _THUNDERBIRD_CONFIG_RESULT_TYPES
        or base.startswith(_THUNDERBIRD_CONFIG_RESULT_PREFIXES)
        or _mentions_thunderbird_config_file(config)
    )


def _mentions_thunderbird_config_file(value: Any) -> bool:
    if isinstance(value, str):
        return any(marker in value for marker in _THUNDERBIRD_CONFIG_FILE_MARKERS)
    if isinstance(value, dict):
        return any(_mentions_thunderbird_config_file(item) for item in value.values())
    if isinstance(value, (list, tuple)):
        return any(_mentions_thunderbird_config_file(item) for item in value)
    return False


def _needs_vlc_config_flush(result_type: str, config: dict | None = None) -> bool:
    base = _base_result_type(result_type)
    if base == "vlc_playing_info":
        return False
    return (
        base in _VLC_CONFIG_RESULT_TYPES
        or base.startswith(_VLC_CONFIG_RESULT_PREFIXES)
        or _mentions_vlc_config_file(config)
    )


def _mentions_vlc_config_file(value: Any) -> bool:
    if isinstance(value, str):
        return any(marker in value for marker in _VLC_CONFIG_FILE_MARKERS)
    if isinstance(value, dict):
        return any(_mentions_vlc_config_file(item) for item in value.values())
    if isinstance(value, (list, tuple)):
        return any(_mentions_vlc_config_file(item) for item in value)
    return False


@_best_effort_flush
async def _flush_thunderbird(eval_env) -> None:
    if getattr(eval_env, "_thunderbird_config_flushed", False):
        _record_flush_skipped(eval_env, "thunderbird")
        return
    command = r"""
sent_quit=0
for display in "${DISPLAY:-}" :1 :0; do
  [ -z "$display" ] && continue
  DISPLAY="$display"; export DISPLAY
  for wid in $(xdotool search --class 'Thunderbird' 2>/dev/null; xdotool search --class 'thunderbird' 2>/dev/null); do
    xdotool windowactivate --sync "$wid" 2>/dev/null || true
    xdotool key --window "$wid" ctrl+q 2>/dev/null || true
    sent_quit=1
    break 2
  done
done
if [ "$sent_quit" -eq 1 ]; then
  for _ in $(seq 1 30); do
    if ! pgrep -x thunderbird >/dev/null 2>&1 && ! pgrep -x thunderbird-bin >/dev/null 2>&1; then
      break
    fi
    sleep 1
  done
  pkill -TERM thunderbird 2>/dev/null || true
  sleep 2
  sync
fi
"""
    await eval_env.computer.interface.run_command(
        "bash -lc " + shlex.quote(command),
        timeout=45,
    )
    _record_flush_executed(eval_env, "thunderbird")
    setattr(eval_env, "_thunderbird_config_flushed", True)


@_best_effort_flush
async def _flush_vlc(eval_env) -> None:
    if getattr(eval_env, "_vlc_config_flushed", False):
        _record_flush_skipped(eval_env, "vlc")
        return
    # VLC persists `[qt]` keys (e.g. qt-continue) to vlcrc only on a CLEAN quit
    # through its own Qt event loop. A `ctrl+q` sent while a modal dialog (Simple
    # Preferences, the first-run privacy prompt, ...) is still open is swallowed
    # by the dialog and never reaches the main window's quit/save path, so the
    # edit is lost. We therefore (1) commit+close any lingering
    # VLC dialog with Return (default button of Simple Preferences is "Save", so
    # this persists rather than discards the agent's in-memory change), then
    # (2) send ctrl+q to every VLC window so the main window's clean-quit save
    # path runs, then (3) wait-for-exit with a SIGTERM (never -9) fallback + sync.
    command = r"""
sent_quit=0
for display in "${DISPLAY:-}" :1 :0; do
  [ -z "$display" ] && continue
  DISPLAY="$display"; export DISPLAY
  vlc_wids="$(xdotool search --class 'vlc' 2>/dev/null)"
  [ -z "$vlc_wids" ] && continue
  # 1) Commit + close any open modal dialog (Simple Preferences / privacy prompt)
  #    so the subsequent ctrl+q reaches the main window's config-save quit path.
  for wid in $(xdotool search --name 'Preferences' 2>/dev/null); do
    xdotool windowactivate --sync "$wid" 2>/dev/null || true
    xdotool key --window "$wid" Return 2>/dev/null || true
    sleep 1
  done
  # 2) Clean-quit every VLC window; the main window's ctrl+q writes vlcrc.
  for wid in $vlc_wids; do
    xdotool windowactivate --sync "$wid" 2>/dev/null || true
    xdotool key --window "$wid" ctrl+q 2>/dev/null || true
    sent_quit=1
  done
  break
done
if [ "$sent_quit" -eq 1 ]; then
  for _ in $(seq 1 30); do
    if ! pgrep -x vlc >/dev/null 2>&1; then
      break
    fi
    sleep 1
  done
  pkill -TERM vlc 2>/dev/null || true
  sleep 2
  sync
fi
"""
    await eval_env.computer.interface.run_command(
        "bash -lc " + shlex.quote(command),
        timeout=45,
    )
    _record_flush_executed(eval_env, "vlc")
    setattr(eval_env, "_vlc_config_flushed", True)


# --- LibreOffice registry config-flush (LO family) ------------------------------
#
# Structural clone of _flush_gimp_config. LibreOffice commits its user registry
# (registrymodifications.xcu) ONLY on a clean shutdown, so the writer_* registry
# getters (default font, zoom, margins, autosave, ...) read a stale/empty xcu
# while the agent's soffice.bin is still running -> FN. We send ctrl+q so the
# app's own quit->config-save path runs, accept the "Save changes?" modal so the
# quit is not blocked, then wait-for-exit with a SIGTERM fallback -- NEVER -9:
# soffice.bin's SIGTERM handler flushes the user registry, whereas SIGKILL skips
# it and re-introduces the stale-read FN -- then sync.
#
# During oracle replay the fixture has already run
# `killall -9 -q soffice soffice.bin libreoffice`, so `xdotool search` finds no
# window, `sent_quit` stays 0, and the whole body no-ops. It MUST NOT disturb an
# already-passing fixture. The FN-recovery this enables is rollout-only.

_LIBREOFFICE_CONFIG_FILE_MARKERS = (
    "registrymodifications.xcu",
    "/home/user/.config/libreoffice/",
    "~/.config/libreoffice/",
    ".config/libreoffice/",
)
_LIBREOFFICE_CONFIG_RESULT_TYPES = {
    "writer_default_font",
    "writer_default_font_size",
    "writer_heading_font",
    "writer_heading_font_size",
    "writer_list_font",
    "writer_caption_font",
    "writer_index_font",
    "writer_font_combo",
    "writer_user_data",
    "writer_multiple_settings",
    "writer_autosave_setting",
    "writer_zoom_level",
    "writer_page_margin",
    "writer_paragraph_spacing",
}
_LIBREOFFICE_CONFIG_RESULT_PREFIXES = (
    "writer_default_font",
    "writer_heading_font",
    "writer_user_data",
    "writer_autosave",
    "writer_zoom",
    "writer_page_margin",
    "writer_paragraph_spacing",
)


def _needs_libreoffice_config_flush(result_type: str, config: dict | None = None) -> bool:
    base = _base_result_type(result_type)
    return (
        base in _LIBREOFFICE_CONFIG_RESULT_TYPES
        or base.startswith(_LIBREOFFICE_CONFIG_RESULT_PREFIXES)
        or _mentions_libreoffice_config_file(config)
    )


def _mentions_libreoffice_config_file(value: Any) -> bool:
    if isinstance(value, str):
        lowered = value.lower()
        return any(marker in lowered for marker in _LIBREOFFICE_CONFIG_FILE_MARKERS)
    if isinstance(value, dict):
        return any(_mentions_libreoffice_config_file(item) for item in value.values())
    if isinstance(value, (list, tuple)):
        return any(_mentions_libreoffice_config_file(item) for item in value)
    return False


@_best_effort_flush
async def _flush_libreoffice(eval_env) -> None:
    if getattr(eval_env, "_libreoffice_config_flushed", False):
        _record_flush_skipped(eval_env, "libreoffice")
        return
    command = r"""
sent_quit=0
for display in "${DISPLAY:-}" :1 :0; do
  [ -z "$display" ] && continue
  DISPLAY="$display"; export DISPLAY
  for wid in $(xdotool search --class 'libreoffice' 2>/dev/null; xdotool search --class 'soffice' 2>/dev/null); do
    xdotool windowactivate --sync "$wid" 2>/dev/null || true
    xdotool key --window "$wid" ctrl+q 2>/dev/null || true
    sent_quit=1
    break 2
  done
done
if [ "$sent_quit" -eq 1 ]; then
  sleep 2
  # Accept the "Save Document?" modal (default button = Save) so ctrl+q's clean
  # quit path completes; harmless no-op when the document is not dirty.
  xdotool key Return 2>/dev/null || true
  for _ in $(seq 1 30); do
    if ! pgrep -x soffice.bin >/dev/null 2>&1; then
      break
    fi
    sleep 1
  done
  pkill -TERM soffice.bin 2>/dev/null || true
  pkill -TERM oosplash 2>/dev/null || true
  sleep 2
  sync
fi
"""
    await eval_env.computer.interface.run_command(
        "bash -lc " + shlex.quote(command),
        timeout=45,
    )
    _record_flush_executed(eval_env, "libreoffice")
    setattr(eval_env, "_libreoffice_config_flushed", True)


async def _read_vm_json(eval_env, path: str) -> dict[str, Any]:
    try:
        data = await eval_env.computer.interface.read_bytes(
            judges._alias_chrome_profile_path(path)
        )
        parsed = json.loads(data.decode("utf-8", errors="replace"))
    except Exception:
        return {}
    return parsed if isinstance(parsed, dict) else {}


async def _load_chrome_local_state(eval_env) -> dict[str, Any]:
    return await _read_vm_json(eval_env, "/home/user/chrome-data/Local State")


async def _load_chrome_secure_prefs(eval_env) -> dict[str, Any]:
    return await _read_vm_json(eval_env, "/home/user/chrome-data/Default/Secure Preferences")


async def _repair_chrome_experiments_result(eval_env, config: dict, result: Any) -> Any:
    result_type = str(config.get("type") or "")
    result_base = _base_result_type(result_type)
    if result_base != "enabled_experiments" and not result_base.startswith(
        "chrome_experiments"
    ):
        return result

    local_state = await _load_chrome_local_state(eval_env)
    raw_experiments = local_state.get("browser", {}).get("enabled_labs_experiments", [])
    if not isinstance(raw_experiments, list):
        raw_experiments = result if isinstance(result, list) else []
    if not raw_experiments:
        return result

    repaired = []
    for item in raw_experiments:
        if not isinstance(item, str) or not item:
            continue
        name, _, value = item.partition("@")
        if name == "enable-quic" and value == "2":
            repaired.append("disable-quic")
        elif value == "2":
            continue
        elif name:
            repaired.append(name)
    return repaired


def _chrome_third_party_cookies_blocked(prefs: dict[str, Any]) -> bool | None:
    profile = prefs.get("profile", {})
    if not isinstance(profile, dict):
        return None
    if profile.get("block_third_party_cookies") is True:
        return True
    mode = profile.get("cookie_controls_mode")
    with contextlib.suppress(TypeError, ValueError):
        return int(mode) == 1
    return False


def _is_chrome_third_party_cookies_result(result_type: str) -> bool:
    base = _base_result_type(result_type)
    return base in {
        "block_third_party_cookies",
        "chrome_block_third_party_cookies",
        "chrome_third_party_cookies",
        "chrome_third_party_cookies_blocked",
        "third_party_cookies_blocked",
    }


async def _repair_chrome_third_party_cookies_result(
    eval_env,
    config: dict,
    result: Any,
) -> Any:
    result_type = str(config.get("type") or "")
    if not _is_chrome_third_party_cookies_result(result_type):
        return result
    blocked = _chrome_third_party_cookies_blocked(await _load_chrome_prefs(eval_env))
    if blocked is None:
        return result
    if isinstance(result, dict):
        repaired = dict(result)
        base = _base_result_type(result_type)
        if "block_third_party_cookies" in repaired or base == "block_third_party_cookies":
            repaired["block_third_party_cookies"] = blocked
        if "block_third_party" in repaired or base == "chrome_third_party_cookies":
            repaired["block_third_party"] = blocked
        return repaired
    return "true" if blocked else "false"


async def _chrome_extension_entries(eval_env) -> list[dict[str, Any]]:
    settings: dict[str, Any] = {}
    for prefs in (await _load_chrome_prefs(eval_env), await _load_chrome_secure_prefs(eval_env)):
        candidate = prefs.get("extensions", {}).get("settings", {})
        if isinstance(candidate, dict):
            settings.update(candidate)
    if not settings:
        return []
    entries: list[dict[str, Any]] = []
    for extension_id, data in settings.items():
        if not isinstance(data, dict):
            continue
        path = str(data.get("path") or "")
        manifest = data.get("manifest") if isinstance(data.get("manifest"), dict) else {}
        manifest = dict(manifest)
        if path and not manifest.get("name"):
            disk_manifest = await _read_vm_json(eval_env, f"{path.rstrip('/')}/manifest.json")
            if disk_manifest:
                manifest = {**disk_manifest, **manifest}
        entries.append(
            {
                "id": extension_id,
                "path": path,
                "state": data.get("state"),
                "location": data.get("location"),
                "manifest": manifest,
            }
        )
    return entries


async def _repair_chrome_extension_result(eval_env, config: dict, result):
    result_type = str(config.get("type") or "")
    if not _is_chrome_extension_result(result_type):
        return result
    result_base = _base_result_type(result_type)

    if result_base in {
        "ext_names",
        "ext_name_list",
        "extension_names",
        "find_installed_extension_name",
    }:
        entries = await _chrome_extension_entries(eval_env)
        names = [
            entry.get("manifest", {}).get("name", "")
            for entry in entries
            if entry.get("manifest", {}).get("name", "")
        ]
        return names or result

    if not _is_empty_chrome_extension_result(result):
        return result

    entries = await _chrome_extension_entries(eval_env)
    if not entries:
        return result
    if result_base == "find_unpacked_extension_path":
        return [entry["path"] for entry in entries if entry.get("path")]
    if result_base == "ext_count" and result_type == "ext_count__dc791077":
        return sum(1 for entry in entries if _chrome_extension_is_unpacked(entry))
    if result_base == "ext_info":
        named = [entry for entry in entries if entry.get("manifest", {}).get("name")]
        return {
            "names": [entry.get("manifest", {}).get("name", "") for entry in named],
            "paths": [entry.get("path", "") for entry in entries if entry.get("path")],
            "enabled_states": [_chrome_extension_is_enabled(entry) for entry in named],
            "dev_mode_states": [_chrome_extension_is_dev_mode(entry) for entry in named],
        }
    if result_base == "ext_location":
        return [
            {
                "id": entry.get("id", ""),
                "path": entry.get("path", ""),
                "name": entry.get("manifest", {}).get("name", ""),
                "state": entry.get("state", 0),
            }
            for entry in entries
            if _chrome_extension_is_hello(entry)
            and "/desktop/" in str(entry.get("path") or "").lower()
        ]
    if result_base == "extension_details":
        return [
            {
                "name": entry.get("manifest", {}).get("name", ""),
                "version": entry.get("manifest", {}).get("version", ""),
                "description": entry.get("manifest", {}).get("description", ""),
                "manifest_version": entry.get("manifest", {}).get("manifest_version", 0),
            }
            for entry in entries
        ]
    if result_type == "extension_description__ae6416e4":
        extensions = {
            entry.get("manifest", {}).get("name", ""): entry.get("manifest", {}).get("description", "")
            for entry in entries
            if entry.get("manifest", {}).get("name", "")
        }
        return {"extensions": extensions, "all_names": list(extensions)}
    if result_type == "extension_manifest_version__ae6416e4":
        return {
            entry.get("manifest", {}).get("name", ""): {
                "manifest_version": entry.get("manifest", {}).get("manifest_version", 0),
                "is_unpacked": _chrome_extension_is_unpacked(entry),
            }
            for entry in entries
            if entry.get("manifest", {}).get("name", "")
        }
    if result_type.startswith("extension_source_type"):
        return {
            entry.get("manifest", {}).get("name", ""): (
                "webstore" if bool(entry.get("from_webstore", False)) else "unpacked"
            )
            for entry in entries
            if entry.get("manifest", {}).get("name", "")
        }

    matched = _find_chrome_extension_entry(entries, config)
    if matched is None:
        return result
    manifest = matched.get("manifest") or {}

    if result_base == "ext_count":
        if result_type == "ext_count__2cf92b2479d096498965f5b9ffc3704c":
            return {
                "found": True,
                "path": matched.get("path", ""),
                "name": manifest.get("name", ""),
                "enabled": _chrome_extension_is_enabled(matched),
            }
        return {
            "found_hello_extension": True,
            "extension_name": manifest.get("name", ""),
            "extension_path": matched.get("path", ""),
            "is_unpacked": _chrome_extension_is_unpacked(matched),
            "total_extensions": len(entries),
        }
    if result_base == "ext_manifest_data":
        return manifest
    if result_base == "ext_manifest":
        return {
            "path": matched.get("path", ""),
            "manifest_version": manifest.get("manifest_version", 0),
            "name": manifest.get("name", ""),
            "version": manifest.get("version", ""),
        }
    if result_base == "ext_manifest_version":
        return manifest.get("manifest_version") or result
    if result_base == "ext_version":
        return manifest.get("version") or result
    if result_base == "ext_description":
        return manifest.get("description") or result
    if result_base == "ext_enabled":
        return {
            "path": matched.get("path", ""),
            "enabled": _chrome_extension_is_enabled(matched),
            "state": matched.get("state", 0),
        }
    if result_base == "ext_enabled_status":
        return {
            "path": matched.get("path", ""),
            "enabled": _chrome_extension_is_enabled(matched),
            "found": True,
        }
    if result_base == "ext_name_by_path":
        return manifest.get("name") or result
    if result_type.startswith("extension_version"):
        return manifest.get("version") or result
    if result_type.startswith("extension_description"):
        return manifest.get("description") or result
    if result_type.startswith("extension_manifest_version"):
        return manifest.get("manifest_version") or result
    if result_type.startswith("extension_name"):
        return manifest.get("name") or result
    if result_type.startswith("extension_details"):
        return [
            {
                "name": entry.get("manifest", {}).get("name", ""),
                "version": entry.get("manifest", {}).get("version", ""),
                "description": entry.get("manifest", {}).get("description", ""),
                "manifest_version": entry.get("manifest", {}).get("manifest_version", 0),
            }
            for entry in entries
        ]
    return result


async def _repair_default_search_engine_result(eval_env, config: dict, result: Any) -> Any:
    if _base_result_type(str(config.get("type", ""))) != "default_search_engine":
        return result
    profile_result = await eval_env.computer.interface.run_command(
        "python3 - <<'PY'\n"
        "import json, os\n"
        "paths = [\n"
        "    '/home/user/chrome-data/Default/Preferences',\n"
        "    '/home/user/chrome-data/Default/Secure Preferences',\n"
        "    '/home/user/.config/google-chrome/Default/Preferences',\n"
        "    '/home/user/.config/google-chrome/Default/Secure Preferences',\n"
        "]\n"
        "records = []\n"
        "for path_index, path in enumerate(paths):\n"
        "    try:\n"
        "        with open(path, encoding='utf-8') as handle:\n"
        "            data = json.load(handle)\n"
        "    except Exception:\n"
        "        continue\n"
        "    if not isinstance(data, dict):\n"
        "        continue\n"
        "    candidates = []\n"
        "    provider_data = data.get('default_search_provider_data', {})\n"
        "    if isinstance(provider_data, dict):\n"
        "        template = provider_data.get('template_url_data', {})\n"
        "        if isinstance(template, dict):\n"
        "            for key in ('short_name', 'keyword', 'url', 'suggestions_url'):\n"
        "                value = template.get(key)\n"
        "                if value:\n"
        "                    candidates.append((3, str(value)))\n"
        "        for key in ('template_url', 'search_url', 'suggestions_url'):\n"
        "            value = provider_data.get(key)\n"
        "            if value:\n"
        "                candidates.append((3, str(value)))\n"
        "    provider = data.get('default_search_provider', {})\n"
        "    if isinstance(provider, dict):\n"
        "        for key in ('name', 'keyword', 'search_url', 'suggest_url'):\n"
        "            value = provider.get(key)\n"
        "            if value:\n"
        "                candidates.append((2, str(value)))\n"
        "    overrides = data.get('search_provider_overrides', [])\n"
        "    if isinstance(overrides, list):\n"
        "        for entry in overrides:\n"
        "            if not isinstance(entry, dict):\n"
        "                continue\n"
        "            for key in ('short_name', 'keyword', 'url', 'search_url'):\n"
        "                value = entry.get(key)\n"
        "                if value:\n"
        "                    candidates.append((1, str(value)))\n"
        "    path_rank = len(paths) - path_index\n"
        "    for index, (source_rank, value) in enumerate(candidates):\n"
        "        records.append((path_rank, source_rank, -index, value))\n"
        "records.sort(reverse=True)\n"
        "seen = set()\n"
        "values = []\n"
        "for _, _, _, value in records:\n"
        "    key = value.strip().lower()\n"
        "    if not key or key in seen:\n"
        "        continue\n"
        "    seen.add(key)\n"
        "    values.append(value)\n"
        "print('__SCALECUA_SEARCH_CANDIDATES__' + json.dumps(values))\n"
        "PY",
        timeout=10,
    )
    stdout = str(getattr(profile_result, "stdout", "") or "")
    marker = "__SCALECUA_SEARCH_CANDIDATES__"
    if marker in stdout:
        try:
            candidates = json.loads(stdout.split(marker, 1)[1].strip().splitlines()[0])
        except Exception:
            candidates = []
        if isinstance(candidates, list) and candidates:
            return {
                "__scalecua_candidates__": True,
                "value": candidates[0],
                "candidates": candidates,
            }
        ui_candidate = await _default_search_engine_from_settings_ui(eval_env)
        if ui_candidate:
            return ui_candidate
    candidate = stdout.strip()
    if candidate and marker not in candidate:
        return candidate
    ui_candidate = await _default_search_engine_from_settings_ui(eval_env)
    if ui_candidate:
        return ui_candidate
    return result


async def _default_search_engine_from_settings_ui(eval_env) -> str | None:
    snapshot = await _get_chrome_settings_snapshot(
        eval_env,
        "chrome://settings/searchEngines",
    )
    text = str(snapshot.get("text") or "")
    for engine in (
        "Microsoft Bing",
        "DuckDuckGo",
        "Startpage",
        "Yahoo! Hong Kong",
        "Yahoo!",
        "Google",
        "Baidu",
        "360",
    ):
        if re.search(rf"\b{re.escape(engine)}\s*\(Default\)", text, re.IGNORECASE):
            return _normalize_default_search_engine(engine)
    match = re.search(r"([A-Za-z0-9][A-Za-z0-9 !.+&-]{1,80}?)\s*\(Default\)", text)
    if match is None:
        return None
    return _normalize_default_search_engine(match.group(1).strip())


def _is_empty_chrome_extension_result(result: Any) -> bool:
    if result in ("", None, [], {}, 0):
        return True
    if isinstance(result, dict):
        if result.get("found") is False or result.get("found_hello_extension") is False:
            return True
        if result.get("extensions") == {} and result.get("all_names") == []:
            return True
        non_status_values = [
            value
            for key, value in result.items()
            if key not in {"enabled", "state", "is_unpacked"}
        ]
        return not any(non_status_values)
    return False


def _chrome_extension_name(entry: dict[str, Any]) -> str:
    return str((entry.get("manifest") or {}).get("name") or "")


def _chrome_extension_path(entry: dict[str, Any]) -> str:
    return str(entry.get("path") or "")


def _chrome_extension_is_enabled(entry: dict[str, Any]) -> bool:
    return entry.get("state", 0) == 1


def _chrome_extension_is_unpacked(entry: dict[str, Any]) -> bool:
    return bool(_chrome_extension_path(entry))


def _chrome_extension_is_dev_mode(entry: dict[str, Any]) -> bool:
    return not bool(entry.get("from_webstore", True)) and _chrome_extension_is_unpacked(entry)


def _chrome_extension_is_hello(entry: dict[str, Any]) -> bool:
    name = _chrome_extension_name(entry).lower()
    path = _chrome_extension_path(entry).lower()
    return ("hello" in name and "extension" in name) or "helloextension" in path


def _find_chrome_extension_entry(
    entries: list[dict[str, Any]],
    config: dict[str, Any],
) -> dict[str, Any] | None:
    target_name = str(config.get("extension_name") or "").strip().lower()
    target_path = str(config.get("extension_path") or config.get("path") or "").rstrip("/")
    path_contains = str(config.get("path_contains") or "").strip()
    name_pattern = str(config.get("name_pattern") or "").strip()

    def matches(entry: dict[str, Any]) -> bool:
        name = _chrome_extension_name(entry)
        path = _chrome_extension_path(entry).rstrip("/")
        if target_name and name.strip().lower() == target_name:
            return True
        if target_path and path == target_path:
            return True
        if path_contains and path_contains in path:
            return True
        if name_pattern and name:
            try:
                if re.search(name_pattern, name, re.IGNORECASE):
                    return True
            except re.error:
                if name_pattern.lower() in name.lower():
                    return True
        return False

    if any(
        value
        for value in (target_name, target_path, path_contains, name_pattern)
    ):
        return next((entry for entry in entries if matches(entry)), None)
    return next((entry for entry in entries if _chrome_extension_is_hello(entry)), None)


async def _get_scalecua_local_result(eval_env, config: dict, cache_dir: str):
    result_type = config.get("type", "")
    result_base = _base_result_type(result_type)
    if result_base == "active_url_from_accessTree":
        return await _get_active_url_from_access_tree(eval_env, config)
    if result_base == "active_tab_info":
        return await _get_active_tab_info(eval_env, config)
    if result_base == "active_tab_url_parse":
        return await _get_active_tab_url_parse(eval_env, config)
    if result_base == "active_tab_html_parse":
        return await _get_active_tab_html_parse(eval_env, config)
    if result_base == "open_tabs_info":
        return await _get_open_tabs_info(eval_env)
    if result_base == "url_dashPart":
        return await _get_url_dash_part(eval_env, config)
    if result_base in SCALECUA_LOCAL_CHROME_PREF_RESULT_TYPES:
        prefs = await _load_chrome_prefs(eval_env)
        result = _extract_scalecua_chrome_pref(result_base, prefs)
        return await _repair_chrome_settings_pref_result(eval_env, result_base, result)
    if result_base == "vlc_playing_info":
        return await _get_vlc_playing_info(eval_env, config, cache_dir)
    if result_base == "vm_file":
        return await _get_vm_file_official(eval_env, config, cache_dir)
    if result_base == "find_unpacked_extension_path":
        entries = await _chrome_extension_entries(eval_env)
        return [entry["path"] for entry in entries if entry.get("path")]
    if result_base == "recreation_devilsgarden_html":
        return await _get_recreation_devilsgarden_html(eval_env, config)
    if result_base == "recreation_url_check":
        return {
            "url": await _get_active_url_from_access_tree(
                eval_env,
                {"goto_prefix": config.get("goto_prefix", "https://")},
            )
        }
    raise KeyError(f"unsupported ScaleCUA local result type: {result_type}")


async def _get_recreation_devilsgarden_html(eval_env, config: dict) -> dict[str, Any] | None:
    pages = [
        tab
        for tab in await _get_cdp_pages(eval_env.computer)
        if tab.get("type") == "page"
        and not str(tab.get("url", "")).startswith(("chrome://", "devtools://"))
    ]
    if not pages:
        return None

    def rank(tab: dict[str, Any]) -> tuple[int, int]:
        url = str(tab.get("url", "")).lower()
        title = str(tab.get("title", "")).lower()
        text = f"{url} {title}"
        return (
            0 if "recreation.gov" in text else 1,
            0 if any(token in text for token in ("devil", "devils", "232449")) else 1,
        )

    target_url = str(sorted(pages, key=rank)[0].get("url") or "")
    if not target_url:
        return None
    page = await base_runner._get_page_info_via_cdp(eval_env.computer, target_url)
    html = str(page.get("content") or "")
    current_url = str(page.get("url") or target_url)
    result: dict[str, Any] = {
        "location_verified": False,
        "reservation_table_present": False,
        "url": current_url,
        "page_title": str(page.get("title") or ""),
        "has_availability_data": False,
        "dates_sorted": False,
        "earliest_reservation_identified": False,
        "reservation_dates": [],
    }
    try:
        min_count = max(1, int(config.get("order", 2)))
    except (TypeError, ValueError):
        min_count = 2

    if html:
        try:
            from bs4 import BeautifulSoup
        except Exception:
            soup = None
        else:
            soup = BeautifulSoup(html, "html.parser")

        if soup is not None:
            raw_page_text = soup.get_text(" ", strip=True)
            page_text = raw_page_text.lower()
            url_lower = current_url.lower()
            if (
                "devil" in url_lower
                or "devils" in url_lower
                or "232449" in url_lower
                or "devil's garden" in page_text
                or "devils garden" in page_text
                or "devil garden" in page_text
            ):
                result["location_verified"] = True

            title = soup.find("title")
            if title:
                result["page_title"] = title.get_text(strip=True)

            if config.get("selector", "class") == "class":
                class_name = config.get("class", "camp-sortable-column-header")
                if len(soup.find_all(class_=class_name)) >= min_count:
                    result["reservation_table_present"] = True
                elif _has_recreation_grid(soup, page_text, min_count):
                    result["reservation_table_present"] = True

            reservation_dates = _extract_recreation_reservation_dates(raw_page_text)
            if reservation_dates:
                result["reservation_dates"] = reservation_dates
                result["dates_sorted"] = True

            availability_indicators = [
                soup.find_all(class_="availability-status"),
                soup.find_all(class_="campsite-row"),
                soup.find_all(class_="rec-availability"),
                _find_recreation_available_cells(soup),
                soup.find_all(
                    "td",
                    class_=lambda value: bool(value and "available" in str(value).lower()),
                ),
                soup.find_all(
                    "div",
                    class_=lambda value: bool(value and "reservation" in str(value).lower()),
                ),
            ]
            if any(items for items in availability_indicators):
                result["has_availability_data"] = True
            else:
                date_elements = soup.find_all(
                    class_=lambda value: bool(value and "date" in str(value).lower())
                )
                result["has_availability_data"] = (
                    result["reservation_table_present"]
                    and bool(reservation_dates)
                    and (len(date_elements) >= 3 or "next available" in page_text)
                )
            result["earliest_reservation_identified"] = bool(
                result["has_availability_data"] and reservation_dates
            )

    if not _recreation_result_complete(result):
        at_text = await _get_recreation_accessibility_text(eval_env)
        _apply_recreation_text_result(result, at_text, min_count)

    return result


def _recreation_result_complete(result: dict[str, Any]) -> bool:
    return bool(
        result.get("location_verified")
        and result.get("reservation_table_present")
        and result.get("has_availability_data")
        and result.get("dates_sorted")
        and result.get("earliest_reservation_identified")
        and result.get("reservation_dates")
    )


async def _get_recreation_accessibility_text(eval_env) -> str:
    at_xml = await base_runner._get_at_xml(eval_env.computer)
    return await asyncio.to_thread(_extract_accessibility_visible_text, at_xml or "")


def _extract_accessibility_visible_text(at_xml: str) -> str:
    if not at_xml:
        return ""
    parts: list[str] = []
    try:
        import lxml.etree

        tree = lxml.etree.fromstring(
            at_xml.encode() if isinstance(at_xml, str) else at_xml
        )
        for elem in tree.iter():
            if elem.text and elem.text.strip():
                parts.append(elem.text.strip())
            name = elem.get("name")
            if name:
                parts.append(name)
            for attr_name, attr_value in elem.attrib.items():
                local_name = attr_name.rsplit("}", 1)[-1]
                if local_name in {"value", "description"} and attr_value:
                    parts.append(str(attr_value))
    except Exception:
        parts.append(re.sub(r"<[^>]+>", " ", str(at_xml)))
    return " ".join(" ".join(parts).split())


def _apply_recreation_text_result(
    result: dict[str, Any],
    raw_text: str,
    min_count: int,
) -> None:
    if not raw_text:
        return
    page_text = raw_text.lower()
    url_lower = str(result.get("url") or "").lower()
    if (
        "devil" in url_lower
        or "devils" in url_lower
        or "232449" in url_lower
        or "devil's garden" in page_text
        or "devils garden" in page_text
        or "devil garden" in page_text
    ):
        result["location_verified"] = True

    if _has_recreation_grid_text(page_text, min_count):
        result["reservation_table_present"] = True

    reservation_dates = _extract_recreation_reservation_dates(raw_text)
    if reservation_dates and not result.get("reservation_dates"):
        result["reservation_dates"] = reservation_dates
        result["dates_sorted"] = True

    if (
        result["reservation_table_present"]
        and result.get("reservation_dates")
        and _recreation_text_has_available_cell(raw_text)
    ):
        result["has_availability_data"] = True
    result["earliest_reservation_identified"] = bool(
        result["has_availability_data"] and result.get("reservation_dates")
    )


def _has_recreation_grid(soup, page_text: str, min_count: int) -> bool:
    tableish_classes = (
        "availability-grid",
        "campground-availability",
        "campsite-list",
        "rec-availability",
    )
    if soup.find_all(
        class_=lambda value: bool(
            value and any(token in str(value).lower() for token in tableish_classes)
        )
    ):
        return True
    return _has_recreation_grid_text(page_text, min_count)


def _has_recreation_grid_text(page_text: str, min_count: int) -> bool:
    grid_terms = ("campsite list", "sites", "loop", "next available")
    if not all(term in page_text for term in grid_terms):
        return False
    return len(_RECREATION_WEEKDAY_DATE_RE.findall(page_text)) >= min_count


def _find_recreation_available_cells(soup) -> list[Any]:
    available_cells = []
    for tag in soup.find_all(True):
        text = tag.get_text(" ", strip=True).lower()
        attrs = " ".join(
            str(value)
            for key, value in tag.attrs.items()
            if key in {"aria-label", "title", "data-status", "data-availability", "class"}
        ).lower()
        if text in {"a", "available"} or re.search(r"\bavailable\b", attrs):
            available_cells.append(tag)
    return available_cells


def _recreation_text_has_available_cell(text: str) -> bool:
    without_next_available = re.sub(
        r"\bnext\s+available\b",
        " ",
        str(text),
        flags=re.IGNORECASE,
    )
    if re.search(r"\bavailable\b", without_next_available, re.IGNORECASE):
        return True
    status_tokens = re.findall(r"(?<![A-Za-z])[ARX](?![A-Za-z])", str(text))
    return "A" in status_tokens and len(status_tokens) >= 2


def _extract_recreation_reservation_dates(text: str) -> list[str]:
    dates: list[str] = []
    seen: set[str] = set()
    for pattern in (_RECREATION_MONTH_DATE_RE, _RECREATION_WEEKDAY_DATE_RE):
        for match in pattern.finditer(text):
            value = " ".join(match.group(0).split())
            key = value.lower()
            if key not in seen:
                seen.add(key)
                dates.append(value)
    return dates


async def _get_active_url_from_access_tree(eval_env, config: dict) -> str | None:
    at_xml = await base_runner._get_at_xml(eval_env.computer)
    raw_url = None
    if not at_xml:
        if _allows_chrome_internal_fallback(config):
            return await _get_chrome_internal_url_fallback(eval_env.computer)
        return None
    extracted = await asyncio.to_thread(base_runner._extract_address_bar_url, at_xml)
    if extracted:
        raw_url = str(extracted).strip()
        fallback_url = await _get_accessibility_namespace_url_fallback(
            eval_env.computer,
            raw_url,
        )
        if fallback_url:
            return fallback_url
        if _looks_prefixed_url(raw_url):
            if raw_url.startswith("chrome://"):
                return await _get_chrome_internal_url_fallback(
                    eval_env.computer, raw_url
                ) or raw_url
            return raw_url
        prefix = str(config.get("goto_prefix", "https://") or "")
        if prefix:
            return _apply_goto_prefix(raw_url, prefix)

    # Chrome internal pages such as chrome://bookmarks can focus an in-page
    # search field, making AT-SPI expose non-URL text while CDP still has the
    # active tab URL. Keep this fallback narrow to internal Chrome pages.
    internal_url = None
    if _allows_chrome_internal_fallback(config):
        internal_url = await _get_chrome_internal_url_fallback(eval_env.computer)
    if internal_url:
        return internal_url
    return raw_url


async def _get_accessibility_namespace_url_fallback(
    computer,
    raw_url: str | None,
) -> str | None:
    if not _is_accessibility_namespace_url(raw_url):
        return None
    candidates = [
        url
        for url in (
            str(tab.get("url") or "")
            for tab in await _get_cdp_pages(computer)
            if tab.get("type") == "page"
        )
        if url.startswith(("http://", "https://"))
        and not _is_accessibility_namespace_url(url)
    ]
    return candidates[0] if len(candidates) == 1 else None


async def _get_cdp_expected_url_fallback(computer, expected_url: str) -> str | None:
    expected = _normalize_expected_url_fragment(expected_url)
    if not expected:
        return None
    matches: list[str] = []
    for tab in await _get_cdp_pages(computer):
        if tab.get("type") != "page":
            continue
        url = str(tab.get("url") or "")
        if not url.startswith(("http://", "https://")):
            continue
        if _is_accessibility_namespace_url(url):
            continue
        normalized = _normalize_expected_url_fragment(url)
        if expected in normalized:
            matches.append(url)
    return matches[0] if len(matches) == 1 else None


def _normalize_expected_url_fragment(url: str) -> str:
    text = unquote(str(url or "")).strip().lower()
    if not text:
        return ""
    parsed = urlsplit(text if "://" in text else f"https://{text}")
    host = parsed.netloc or parsed.path.split("/", 1)[0]
    path = parsed.path if parsed.netloc else "/" + text.split("/", 1)[1] if "/" in text else ""
    host = host[4:] if host.startswith("www.") else host
    return f"{host}{path}".rstrip("/")


def _is_accessibility_namespace_url(url: str | None) -> bool:
    text = str(url or "").strip()
    if text.startswith("https://accessibility.ubuntu.example.org/ns/"):
        return True
    return text.startswith("accessibility.ubuntu.example.org/ns/")


def _looks_prefixed_url(url: str) -> bool:
    return "://" in url or url.startswith(("about:", "chrome:", "file:", "mailto:"))


def _apply_goto_prefix(raw_url: str, prefix: str) -> str:
    if not prefix or _looks_prefixed_url(raw_url):
        return raw_url
    if prefix.endswith("www.") and _is_subdomain_like(raw_url):
        scheme = urlsplit(prefix).scheme or "https"
        return f"{scheme}://{raw_url}"
    return f"{prefix}{raw_url}"


def _is_subdomain_like(raw_url: str) -> bool:
    host = str(raw_url).split("/", 1)[0].split("?", 1)[0].split("#", 1)[0]
    return bool(re.match(r"^[A-Za-z0-9-]+(?:\.[A-Za-z0-9-]+){2,}$", host))


def _allows_chrome_internal_fallback(config: dict) -> bool:
    prefix = config.get("goto_prefix", "https://")
    return prefix == "" or str(prefix).startswith("chrome://")


async def _get_chrome_internal_url_fallback(
    computer, raw_url: str | None = None
) -> str | None:
    candidates: list[str] = []
    for tab in await _get_cdp_pages(computer):
        if tab.get("type") != "page":
            continue
        url = tab.get("url", "") or ""
        if not url.startswith("chrome://"):
            continue
        if url.startswith("chrome://omnibox-popup"):
            continue
        if url.rstrip("/") in {"chrome://newtab", "chrome://new-tab-page"}:
            continue
        candidates.append(url)
    if raw_url:
        raw_normalized = _normalize_url(raw_url)
        for url in candidates:
            if _normalize_url(url) == raw_normalized:
                return url
        return None
    return candidates[0] if len(candidates) == 1 else None


async def _get_active_tab_info(eval_env, config: dict):
    active_url = await _get_active_url_from_access_tree(eval_env, config)
    if active_url is None:
        return None
    return await base_runner._get_page_info_via_cdp(eval_env.computer, active_url)


async def _get_active_tab_url_parse(eval_env, config: dict):
    active_url = await _get_active_url_from_access_tree(eval_env, config)
    if active_url is None:
        return None
    parsed_url = urlparse(active_url)
    query_params = parse_qs(parsed_url.query)
    extracted = {key: query_params.get(key, [""])[0] for key in config["parse_keys"]}
    extracted = _repair_current_url_param_aliases(
        extracted,
        query_params,
        parsed_url,
        config,
    )
    for old_key, new_key in config.get("replace", {}).items():
        extracted[new_key] = extracted.pop(old_key, "")
    if config.get("split_list", False):
        extracted = {key: value.split(",") for key, value in extracted.items()}
    return extracted


def _repair_current_url_param_aliases(
    extracted: dict[str, str],
    query_params: dict[str, list[str]],
    parsed_url,
    config: dict,
) -> dict[str, str]:
    parse_keys = set(config.get("parse_keys") or ())
    host = str(getattr(parsed_url, "netloc", "") or "").lower()
    path = str(getattr(parsed_url, "path", "") or "").lower()
    if "ryanair." not in host or "/flights/" not in path:
        return extracted
    aliases = {
        "originIata": "originMac",
        "destinationIata": "destinationMac",
        "tpAdults": "adults",
        "tpChildren": "children",
        "tpTeens": "teens",
        "tpStartDate": "dateOut",
        "isReturn": "isReturn",
    }
    repaired = dict(extracted)
    for target_key, source_key in aliases.items():
        if target_key not in parse_keys or repaired.get(target_key):
            continue
        value = query_params.get(source_key, [""])[0]
        if value:
            repaired[target_key] = value
    return repaired


async def _get_active_tab_html_parse(eval_env, config: dict):
    active_url = await _get_active_url_from_access_tree(eval_env, config)
    if not isinstance(active_url, str):
        return None
    if not await _cdp_has_page_url(eval_env.computer, active_url):
        return {}
    result = await base_runner._get_html_parse_via_cdp(eval_env.computer, active_url, config)
    return await _repair_google_flights_html_parse(eval_env, active_url, config, result)


async def _get_vm_now(eval_env) -> datetime | None:
    """VM local wall clock (tz-aware). Relative-date getters must anchor on the clock the
    agent saw *inside the VM*, not the grader host (mirrors upstream desktop_env
    _get_vm_now_datetime). Returns None on any failure so callers fall back to host
    datetime.now() — byte-identical to prior behavior when clocks agree (zero regression)."""
    try:
        result = await eval_env.computer.interface.run_command(
            "date +%Y-%m-%dT%H:%M:%S%z", timeout=10
        )
        out = (getattr(result, "stdout", "") or "").strip()
        if out:
            return datetime.strptime(out, "%Y-%m-%dT%H:%M:%S%z")
    except Exception:
        pass
    return None


async def _repair_google_flights_html_parse(
    eval_env,
    active_url: str,
    config: dict,
    result: Any,
) -> Any:
    if not _is_google_flights_html_parse_config(config):
        return result
    if not isinstance(result, dict):
        result = {}
    if not any(not str(result.get(key) or "").strip() for key in ("start", "end", "time")):
        return result
    text = await _get_cdp_page_inner_text(eval_env.computer, active_url)
    parsed = _parse_google_flights_context_text(text, now=await _get_vm_now(eval_env))
    if not parsed:
        return result
    repaired = dict(result)
    for key, value in parsed.items():
        if key in {"start", "end", "time"} and not str(repaired.get(key) or "").strip():
            repaired[key] = value
    return repaired


def _is_google_flights_html_parse_config(config: dict) -> bool:
    if config.get("type") not in {None, "", "active_tab_html_parse"}:
        return False
    if config.get("category") != "class":
        return False
    class_single = config.get("class_singleObject") or {}
    class_child = config.get("class_multiObject_child") or {}
    return (
        class_single.get("mach-flight-context-info__wrapper--date") == "time"
        and "mach-flight-context-info__wrapper__info--separator" in class_child
    )


async def _get_cdp_page_inner_text(computer, target_url: str) -> str:
    script = r'''
import json
import requests
from urllib.parse import unquote

try:
    import websocket
except Exception:
    print("")
    raise SystemExit

target_url = ''' + json.dumps(target_url) + r'''

def norm(url):
    return unquote(str(url or "")).rstrip("/")

try:
    pages = requests.get("http://localhost:1337/json", timeout=5).json()
except Exception:
    print("")
    raise SystemExit

target = None
for page in pages:
    if page.get("type") == "page" and norm(page.get("url")) == norm(target_url):
        target = page
        break
if not target:
    print("")
    raise SystemExit

try:
    ws = websocket.create_connection(target["webSocketDebuggerUrl"], timeout=10)
    ws.send(json.dumps({
        "id": 1,
        "method": "Runtime.evaluate",
        "params": {
            "expression": "document.body ? document.body.innerText : document.documentElement.innerText",
            "returnByValue": True,
        },
    }))
    response = json.loads(ws.recv())
    ws.close()
    value = response.get("result", {}).get("result", {}).get("value", "")
    print(value if isinstance(value, str) else "")
except Exception:
    print("")
'''
    result = await computer.interface.run_command(f"python3 -c {shlex.quote(script)}")
    return getattr(result, "stdout", "") or ""


_GOOGLE_FLIGHTS_MONTHS = {
    "jan": 1,
    "feb": 2,
    "mar": 3,
    "apr": 4,
    "may": 5,
    "jun": 6,
    "jul": 7,
    "aug": 8,
    "sep": 9,
    "oct": 10,
    "nov": 11,
    "dec": 12,
}


def _parse_google_flights_context_text(
    text: str,
    *,
    now: datetime | None = None,
) -> dict[str, str] | None:
    normalized = re.sub(r"\s+", " ", text or "").strip()
    if not normalized:
        return None
    route_match = re.search(r"\b([A-Z]{3})\s*[-\u2010-\u2015]\s*([A-Z]{3})\b", normalized)
    if not route_match:
        return None
    date_match = _selected_google_flights_date_match(normalized)
    if not date_match:
        return None
    month_text, day_text = date_match.group(1), date_match.group(2)
    date_value = _format_google_flights_date(month_text, int(day_text), now=now)
    if not date_value:
        return None
    return {
        "start": route_match.group(1),
        "end": route_match.group(2),
        "time": date_value,
    }


def _selected_google_flights_date_match(text: str) -> re.Match[str] | None:
    month_re = r"(Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)"
    patterns = (
        rf"Track prices\s+{month_re}\s+(\d{{1,2}})\b",
        rf"{month_re}\s+(\d{{1,2}})\s+Any dates\b",
    )
    for pattern in patterns:
        match = re.search(pattern, text)
        if match:
            return match
    for match in re.finditer(rf"\b{month_re}\s+(\d{{1,2}})\b", text):
        prefix = text[max(0, match.start() - 16):match.start()].lower()
        if prefix.rstrip().endswith("travel on"):
            continue
        return match
    return None


def _format_google_flights_date(
    month_text: str,
    day: int,
    *,
    now: datetime | None = None,
) -> str | None:
    month = _GOOGLE_FLIGHTS_MONTHS.get(month_text[:3].lower())
    if month is None:
        return None
    today = (now or datetime.now()).date()
    candidates = []
    for year in (today.year - 1, today.year, today.year + 1):
        try:
            candidate = datetime(year, month, day).date()
        except ValueError:
            continue
        delta = (candidate - today).days
        if delta >= -1:
            candidates.append((delta, candidate))
    if not candidates:
        return None
    _, selected = min(candidates, key=lambda item: item[0])
    return selected.strftime("%a, %b %d, %Y")


async def _get_url_dash_part(eval_env, config: dict):
    active_url = await _get_active_url_from_access_tree(eval_env, config)
    if active_url is None:
        return None
    try:
        part_index = int(config["partIndex"])
    except (KeyError, TypeError, ValueError):
        return None
    url_parts = active_url.split("/")
    try:
        dash_part = url_parts[part_index]
    except IndexError:
        return None
    if config.get("needDeleteId", False):
        dash_part = dash_part.split("?")[0]
    if config.get("returnType") == "string":
        return dash_part
    if config.get("returnType") == "json":
        return {config["key"]: dash_part}
    return None


async def _cdp_has_page_url(computer, target_url: str) -> bool:
    for tab in await _get_cdp_pages(computer):
        if tab.get("type") == "page" and _normalize_url(tab.get("url", "")) == _normalize_url(target_url):
            return True
    return False


async def _get_cdp_pages(computer) -> list[dict[str, Any]]:
    result = await computer.interface.run_command("curl -s http://localhost:1337/json")
    try:
        pages = json.loads(getattr(result, "stdout", "") or "[]")
    except (json.JSONDecodeError, TypeError):
        return []
    return pages if isinstance(pages, list) else []


async def _get_chrome_settings_snapshot(eval_env, settings_url: str) -> dict[str, Any]:
    try:
        result = await eval_env.computer.interface.run_command(
            _chrome_settings_snapshot_command(settings_url),
            timeout=20,
        )
    except Exception:
        return {}
    try:
        data = json.loads(getattr(result, "stdout", "") or "{}")
    except (json.JSONDecodeError, TypeError):
        return {}
    return data if isinstance(data, dict) else {}


def _chrome_settings_snapshot_command(settings_url: str) -> str:
    return "python3 - " + shlex.quote(settings_url) + r""" <<'PY'
import base64
import json
import os
import socket
import struct
import sys
import time
import urllib.request
import urllib.parse

settings_url = sys.argv[1]

def ws_connect(url):
    parsed = urllib.parse.urlparse(url)
    sock = socket.create_connection((parsed.hostname, parsed.port), timeout=10)
    key = base64.b64encode(os.urandom(16)).decode()
    path = parsed.path or "/"
    if parsed.query:
        path += "?" + parsed.query
    handshake = (
        f"GET {path} HTTP/1.1\r\n"
        f"Host: {parsed.hostname}:{parsed.port}\r\n"
        "Upgrade: websocket\r\n"
        "Connection: Upgrade\r\n"
        f"Sec-WebSocket-Key: {key}\r\n"
        "Sec-WebSocket-Version: 13\r\n\r\n"
    )
    sock.sendall(handshake.encode())
    response = b""
    while b"\r\n\r\n" not in response:
        chunk = sock.recv(4096)
        if not chunk:
            raise RuntimeError("websocket handshake failed")
        response += chunk
    return sock

def ws_send(sock, payload):
    data = payload.encode()
    frame = bytearray([0x81])
    mask_key = os.urandom(4)
    length = len(data)
    if length < 126:
        frame.append(0x80 | length)
    elif length < 65536:
        frame.append(0x80 | 126)
        frame.extend(struct.pack(">H", length))
    else:
        frame.append(0x80 | 127)
        frame.extend(struct.pack(">Q", length))
    frame.extend(mask_key)
    frame.extend(bytearray(b ^ mask_key[i % 4] for i, b in enumerate(data)))
    sock.sendall(frame)

def ws_recv(sock, timeout=10):
    sock.settimeout(timeout)
    header = sock.recv(2)
    if len(header) < 2:
        return ""
    length = header[1] & 0x7f
    if length == 126:
        length = struct.unpack(">H", sock.recv(2))[0]
    elif length == 127:
        length = struct.unpack(">Q", sock.recv(8))[0]
    data = b""
    while len(data) < length:
        chunk = sock.recv(min(65536, length - len(data)))
        if not chunk:
            break
        data += chunk
    return data.decode(errors="replace")

def recv_id(sock, msg_id, timeout=10):
    deadline = time.time() + timeout
    while time.time() < deadline:
        try:
            message = ws_recv(sock, timeout=max(0.5, deadline - time.time()))
        except Exception:
            return {}
        if not message:
            continue
        try:
            data = json.loads(message)
        except Exception:
            continue
        if data.get("id") == msg_id:
            return data
    return {}

def pages():
    pages = json.loads(urllib.request.urlopen("http://localhost:1337/json", timeout=5).read())
    return [p for p in pages if p.get("type") == "page" and p.get("webSocketDebuggerUrl")]

def same_settings_url(current, target):
    if not current:
        return False
    current = current.split('?', 1)[0].rstrip('/')
    target = target.split('?', 1)[0].rstrip('/')
    return current == target

js = r'''
(() => {
  const ui = document.querySelector('settings-ui');
  const prefs = ui && ui.prefs;
  const selectedRadioTexts = [];
  const seen = new Set();
  const visit = (node) => {
    if (!node || seen.has(node)) return;
    seen.add(node);
    if (node.nodeType === Node.ELEMENT_NODE) {
      const tag = (node.tagName || '').toLowerCase();
      const role = node.getAttribute && node.getAttribute('role');
      const checked =
        node.checked === true ||
        node.hasAttribute && node.hasAttribute('checked') ||
        node.getAttribute && node.getAttribute('aria-checked') === 'true';
      if (checked && (role === 'radio' || tag.includes('radio'))) {
        const text = (node.innerText || node.textContent || '').replace(/\s+/g, ' ').trim();
        if (text) selectedRadioTexts.push(text);
      }
      if (node.shadowRoot) visit(node.shadowRoot);
    }
    for (const child of Array.from(node.children || [])) visit(child);
  };
  visit(document.documentElement);
  const get = (path) => {
    let cur = prefs;
    for (const key of path) {
      if (!cur || !(key in cur)) return null;
      cur = cur[key];
    }
    return cur === undefined ? null : cur;
  };
  const text = document.body && document.body.innerText ? document.body.innerText : '';
  return JSON.stringify({
    url: location.href,
    title: document.title,
    safebrowsing: {
      enhanced: get(['safebrowsing', 'enhanced', 'value']),
      enabled: get(['safebrowsing', 'enabled', 'value'])
    },
    profile: {
      name: get(['profile', 'name', 'value'])
    },
    cookies: {
      default_content_setting_values_cookies: get(['profile', 'default_content_setting_values', 'cookies', 'value']),
      cookie_controls_mode: get(['profile', 'cookie_controls_mode', 'value']),
      block_third_party_cookies: get(['profile', 'block_third_party_cookies', 'value'])
    },
    privacy: {
      do_not_track: get(['enable_do_not_track', 'value'])
    },
    security: {
      generated_https_first_mode_enabled: get(['generated', 'https_first_mode_enabled', 'value']),
      generated_https_only_mode_enabled: get(['generated', 'https_only_mode_enabled', 'value']),
      profile_https_only_mode_enabled: get(['profile', 'https_only_mode_enabled', 'value']),
      https_first_mode_enabled: get(['https_first_mode_enabled', 'value']),
      https_only_mode_enabled: get(['https_only_mode_enabled', 'value'])
    },
    startup: {
      restore_on_startup: get(['session', 'restore_on_startup', 'value'])
    },
    selected_radio_texts: selectedRadioTexts.slice(0, 20),
    text: text.slice(0, 4000)
  });
})()
'''

def eval_snapshot(page, navigate=False):
    sock = ws_connect(page["webSocketDebuggerUrl"])
    if navigate:
        ws_send(sock, json.dumps({"id": 1, "method": "Page.navigate", "params": {"url": settings_url}}))
        recv_id(sock, 1, timeout=5)
        time.sleep(2.5)
    ws_send(sock, json.dumps({"id": 2, "method": "Runtime.evaluate", "params": {"expression": js, "returnByValue": True}}))
    data = recv_id(sock, 2, timeout=10)
    value = data.get("result", {}).get("result", {}).get("value", "{}")
    return value if isinstance(value, str) else json.dumps(value)

def useful_snapshot(value):
    try:
        data = json.loads(value)
    except Exception:
        return False
    if not isinstance(data, dict):
        return False
    if not data.get("url", "").startswith("chrome://settings"):
        return False
    return any(
        section
        for section in (
            data.get("safebrowsing"),
            data.get("profile"),
            data.get("cookies"),
            data.get("privacy"),
            data.get("security"),
            data.get("startup"),
        )
        if isinstance(section, dict) and any(v is not None for v in section.values())
    )

try:
    page_list = pages()
    page = next((p for p in page_list if same_settings_url(p.get("url"), settings_url)), None)
    if page:
        value = eval_snapshot(page, navigate=False)
        if useful_snapshot(value):
            print(value)
            raise SystemExit
    page = page or (page_list[0] if page_list else None)
    if not page:
        print("{}")
        raise SystemExit
    print(eval_snapshot(page, navigate=True))
except Exception:
    print("{}")
PY"""


async def _repair_chrome_settings_pref_result(
    eval_env,
    result_type: str,
    result: Any,
) -> Any:
    if result == "true":
        return result
    settings_url = _chrome_settings_url_for_result(result_type)
    if not settings_url:
        return result
    snapshot = await _get_chrome_settings_snapshot(eval_env, settings_url)
    repaired = _extract_chrome_settings_snapshot_result(result_type, snapshot)
    return repaired if repaired is not None else result


async def _repair_chrome_profile_name_result(
    eval_env,
    result_type: str,
    result: Any,
) -> Any:
    if _base_result_type(result_type) != "profile_name":
        return result
    if not _should_fallback_to_chrome_profile_snapshot(result):
        return result
    snapshot = await _get_chrome_settings_snapshot(eval_env, "chrome://settings/manageProfile")
    name = _get_nested(snapshot, ("profile", "name"))
    if isinstance(name, str) and name.strip():
        return name.strip()
    return result


def _should_fallback_to_chrome_profile_snapshot(result: Any) -> bool:
    if result is None:
        return True
    if not isinstance(result, str):
        return False
    normalized = re.sub(r"\s+", " ", result).strip().lower()
    return normalized in {
        "",
        "{}",
        "null",
        "none",
        "default",
        "person 1",
        "profile 1",
        "user",
        "chrome",
        "chromium",
    }


def _chrome_settings_url_for_result(result_type: str) -> str | None:
    base = _base_result_type(result_type)
    if base in {
        "enable_enhanced_safety_browsing",
        "enable_safe_browsing",
        "disable_safe_browsing",
        "chrome_https_only_mode",
        "https_only_mode",
    }:
        return "chrome://settings/security"
    if base == "data_delete_automacally":
        return "chrome://settings/content/siteData"
    if base == "enable_do_not_track":
        return "chrome://settings/cookies"
    if base == "new_startup_page":
        return "chrome://settings/onStartup"
    if _is_chrome_third_party_cookies_result(base):
        return "chrome://settings/cookies"
    return None


def _extract_chrome_settings_snapshot_result(
    result_type: str,
    snapshot: dict[str, Any],
) -> str | None:
    if not snapshot:
        return None
    base = _base_result_type(result_type)
    if base in {
        "enable_enhanced_safety_browsing",
        "enable_safe_browsing",
        "disable_safe_browsing",
    }:
        enhanced = _coerce_bool_or_none(_get_nested(snapshot, ("safebrowsing", "enhanced")))
        enabled = _coerce_bool_or_none(_get_nested(snapshot, ("safebrowsing", "enabled")))
        if base == "enable_enhanced_safety_browsing":
            return _bool_text(enhanced) if enhanced is not None else None
        active = bool(enhanced) or bool(enabled)
        if enhanced is None and enabled is None:
            return None
        return "true" if active else "false"
    if base in {"chrome_https_only_mode", "https_only_mode"}:
        security = snapshot.get("security", {})
        if not isinstance(security, dict):
            return None
        values = [
            _coerce_bool_or_none(security.get(key))
            for key in (
                "generated_https_first_mode_enabled",
                "generated_https_only_mode_enabled",
                "profile_https_only_mode_enabled",
                "https_first_mode_enabled",
                "https_only_mode_enabled",
            )
        ]
        values = [value for value in values if value is not None]
        return _bool_text(any(values)) if values else None
    if base == "data_delete_automacally":
        cookies = _get_nested(snapshot, ("cookies", "default_content_setting_values_cookies"))
        with contextlib.suppress(TypeError, ValueError):
            return "true" if int(cookies) == 4 else "false"
        return None
    if base == "enable_do_not_track":
        dnt = _coerce_bool_or_none(_get_nested(snapshot, ("privacy", "do_not_track")))
        return _bool_text(dnt) if dnt is not None else None
    if base == "new_startup_page":
        restore = _get_nested(snapshot, ("startup", "restore_on_startup"))
        with contextlib.suppress(TypeError, ValueError):
            return "true" if int(restore) == 5 else "false"
        return None
    if _is_chrome_third_party_cookies_result(base):
        selected = _chrome_settings_selected_radio_texts(snapshot)
        for text in selected:
            normalized = re.sub(r"\s+", " ", text).strip().lower()
            if "allow third-party cookies" in normalized:
                return "false"
            if "block third-party cookies" in normalized:
                return "false" if "incognito" in normalized else "true"
        blocked = _chrome_third_party_cookies_blocked(
            {
                "profile": {
                    "block_third_party_cookies": _get_nested(
                        snapshot, ("cookies", "block_third_party_cookies")
                    ),
                    "cookie_controls_mode": _get_nested(
                        snapshot, ("cookies", "cookie_controls_mode")
                    ),
                }
            }
        )
        return _bool_text(blocked) if blocked is not None else None
    return None


def _chrome_settings_selected_radio_texts(snapshot: dict[str, Any]) -> list[str]:
    values = snapshot.get("selected_radio_texts")
    if not isinstance(values, list):
        return []
    return [str(value) for value in values if str(value).strip()]


def _repair_chrome_settings_result_from_pre_postconfig_state(
    result_config: dict[str, Any],
    result: Any,
    pre_postconfig_state: str | None,
) -> Any:
    if result == "true" or not isinstance(result_config, dict):
        return result
    result_type = str(result_config.get("type") or "")
    if not _chrome_settings_url_for_result(result_type):
        return result
    snapshots = _pre_postconfig_chrome_settings_snapshots(pre_postconfig_state)
    snapshot = snapshots.get(_base_result_type(result_type))
    if not snapshot:
        return result
    repaired = _extract_chrome_settings_snapshot_result(result_type, snapshot)
    return repaired if repaired is not None else result


def _pre_postconfig_chrome_settings_snapshots(
    pre_postconfig_state: str | None,
) -> dict[str, dict[str, Any]]:
    if not pre_postconfig_state:
        return {}
    for line in str(pre_postconfig_state).splitlines():
        if not line.startswith(_CHROME_SETTINGS_SNAPSHOT_MARKER):
            continue
        payload = line[len(_CHROME_SETTINGS_SNAPSHOT_MARKER):]
        try:
            data = json.loads(payload)
        except (json.JSONDecodeError, TypeError):
            return {}
        if not isinstance(data, dict):
            return {}
        return {
            str(key): value
            for key, value in data.items()
            if isinstance(value, dict)
        }
    return {}


def _get_nested(data: dict[str, Any], path: tuple[str, ...]) -> Any:
    cur: Any = data
    for key in path:
        if not isinstance(cur, dict):
            return None
        cur = cur.get(key)
    return cur


def _coerce_bool_or_none(value: Any) -> bool | None:
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in {"true", "1", "yes", "enabled"}:
            return True
        if normalized in {"false", "0", "no", "disabled"}:
            return False
    return None


def _bool_text(value: bool | None) -> str | None:
    if value is None:
        return None
    return "true" if value else "false"


def _normalize_url(url: str) -> str:
    return unquote(str(url or "")).rstrip("/")


async def _get_open_tabs_info(eval_env) -> list[dict[str, str]]:
    tabs: list[dict[str, str]] = []
    for tab in await _get_cdp_pages(eval_env.computer):
        if tab.get("type") != "page":
            continue
        url = tab.get("url", "") or ""
        if url.startswith("chrome://omnibox-popup") or url.startswith("devtools://"):
            continue
        tabs.append({"title": tab.get("title", "") or "", "url": url})
    return tabs


async def _load_chrome_prefs(eval_env) -> dict[str, Any]:
    result = await eval_env.computer.interface.run_command(
        _chrome_prefs_load_command()
    )
    try:
        data = json.loads(getattr(result, "stdout", "") or "{}")
    except (json.JSONDecodeError, TypeError):
        return {}
    return data if isinstance(data, dict) else {}


def _chrome_prefs_load_command() -> str:
    return r"""python3 - <<'PY'
import json
import os
import re
import subprocess
import sys

seen = set()
candidates = []

def add(path, rank):
    if not path or path in seen:
        return
    seen.add(path)
    candidates.append((rank, path))

try:
    ps = subprocess.check_output(
        ["ps", "-eo", "args="],
        text=True,
        stderr=subprocess.DEVNULL,
    )
except Exception:
    ps = ""

for line in ps.splitlines():
    lowered = line.lower()
    if "chrome" not in lowered and "chromium" not in lowered:
        continue
    match = re.search(r"--user-data-dir=(?:'([^']+)'|\"([^\"]+)\"|(\S+))", line)
    if match:
        user_data_dir = next(group for group in match.groups() if group)
        add(os.path.join(user_data_dir, "Default", "Preferences"), 100)

for path in (
    "/home/user/chrome-data/Default/Preferences",
    "/home/user/.config/google-chrome/Default/Preferences",
    "/home/user/.config/chromium/Default/Preferences",
):
    add(path, 0)

existing = []
for rank, path in candidates:
    try:
        stat = os.stat(path)
    except OSError:
        continue
    existing.append((rank, stat.st_mtime, stat.st_size, path))

for _, _, _, path in sorted(existing, reverse=True):
    try:
        with open(path, encoding="utf-8", errors="replace") as f:
            data = json.load(f)
    except Exception:
        continue
    if isinstance(data, dict):
        json.dump(data, sys.stdout)
        raise SystemExit(0)

sys.stdout.write("{}")
PY"""


def _extract_scalecua_chrome_pref(result_type: str, prefs: dict[str, Any]):
    profile = prefs.get("profile", {})
    if not isinstance(profile, dict):
        profile = {}
    if result_type == "enable_do_not_track":
        return "true" if prefs.get("enable_do_not_track", {}) else "false"
    if result_type == "new_startup_page":
        restore = prefs.get("session", {}).get("restore_on_startup", {})
        return "true" if restore == 5 else "false"
    if result_type == "data_delete_automacally":
        cookies = (
            profile
            .get("default_content_setting_values", {})
            .get("cookies")
        )
        content_settings = profile.get("content_settings", {})
        clear_on_exit = (
            content_settings.get("cookies", {}).get("clear_on_exit")
            if isinstance(content_settings, dict)
            else None
        )
        lifetime_enabled = (
            prefs.get("browser", {})
            .get("clear_data", {})
            .get("browsing_data_lifetime", {})
            .get("enabled")
        )
        return "true" if (
            cookies == 4
            or clear_on_exit is True
            or profile.get("clear_lso_data_enabled") is True
            or lifetime_enabled is True
        ) else "false"
    if result_type == "enable_enhanced_safety_browsing":
        return "true" if bool(prefs.get("safebrowsing", {}).get("enhanced", False)) else "false"
    if result_type == "enable_safe_browsing":
        safebrowsing = prefs.get("safebrowsing", {})
        enabled = bool(safebrowsing.get("enabled", False))
        enhanced = bool(safebrowsing.get("enhanced", False))
        return "true" if enabled or enhanced else "false"
    if result_type == "disable_safe_browsing":
        safebrowsing = prefs.get("safebrowsing", {})
        enabled = bool(safebrowsing.get("enabled", False))
        enhanced = bool(safebrowsing.get("enhanced", False))
        return "false" if not (enabled or enhanced) else "true"
    if result_type in {"chrome_https_only_mode", "https_only_mode"}:
        generated = prefs.get("generated", {})
        if not isinstance(generated, dict):
            generated = {}
        values = [
            generated.get("https_first_mode_enabled"),
            generated.get("https_only_mode_enabled"),
            profile.get("https_only_mode_enabled"),
            prefs.get("https_first_mode_enabled"),
            prefs.get("https_only_mode_enabled"),
        ]
        bool_values = [
            coerced
            for value in values
            if (coerced := _coerce_bool_or_none(value)) is not None
        ]
        return _bool_text(any(bool_values)) if bool_values else None
    if result_type in {
        "block_third_party_cookies",
        "chrome_block_third_party_cookies",
        "chrome_third_party_cookies",
        "chrome_third_party_cookies_blocked",
        "third_party_cookies_blocked",
    }:
        blocked = _chrome_third_party_cookies_blocked(prefs)
        if blocked is not None:
            return "true" if blocked else "false"
    if result_type == "password_manager_disabled":
        credentials_enabled = prefs.get("credentials_enable_service")
        password_manager_enabled = profile.get("password_manager_enabled")
        return "true" if (
            credentials_enabled is False or password_manager_enabled is False
        ) else "false"
    if result_type == "password_manager_enabled":
        credentials_enabled = prefs.get("credentials_enable_service", True)
        password_manager_enabled = profile.get("password_manager_enabled", True)
        return "true" if credentials_enabled and password_manager_enabled else "false"
    return None


async def _get_vlc_playing_info(eval_env, config: dict, cache_dir: str):
    dest = config.get("dest", "vlc_status.xml")
    local_path = os.path.join(cache_dir, dest)
    stdout = ""
    for command in (
        "curl -s http://localhost:8080/requests/status.xml",
        "curl -s --user :password http://localhost:8080/requests/status.xml",
        "curl -s --user :vlc http://localhost:8080/requests/status.xml",
        "curl -s --user :a http://localhost:8080/requests/status.xml",
    ):
        result = await eval_env.computer.interface.run_command(command)
        stdout = getattr(result, "stdout", "") or ""
        if "<root" in stdout:
            break
    else:
        return None
    os.makedirs(os.path.dirname(local_path) or cache_dir, exist_ok=True)
    with open(local_path, "wb") as f:
        f.write(stdout.encode("utf-8"))
    return local_path


async def _get_vm_file_official(eval_env, config: dict, cache_dir: str):
    if not config.get("multi", False):
        paths = [config["path"]]
        dests = [config["dest"]]
        if config.get("time_suffix", False):
            time_format = config.get("time_format", "%Y%m%d_%H%M%S")
            dests = [
                f"{os.path.splitext(dest)[0]}_{datetime.now().strftime(time_format)}"
                f"{os.path.splitext(dest)[1]}"
                for dest in dests
            ]
    else:
        paths = list(config["path"])
        dests = list(config["dest"])

    cache_paths: list[str | None] = []
    gives = set(config.get("gives", [0]))
    for index, (remote_path, dest) in enumerate(zip(paths, dests)):
        local_path = os.path.join(cache_dir, dest)
        data = await _read_container_file_bytes(eval_env.computer, remote_path)
        if data is None:
            if index in gives:
                cache_paths.append(None)
            continue
        os.makedirs(os.path.dirname(local_path) or cache_dir, exist_ok=True)
        with open(local_path, "wb") as f:
            f.write(data)
        if index in gives:
            cache_paths.append(local_path)
    return cache_paths[0] if len(cache_paths) == 1 else cache_paths


async def _read_container_file_bytes(computer, remote_path: str) -> bytes | None:
    try:
        data = await computer.interface.read_bytes(remote_path)
        if data is not None:
            return data
    except Exception:
        pass

    marker = "__SCALECUA_FILE_B64__"
    quoted = shlex.quote(remote_path)
    result = await computer.interface.run_command(
        f"if test -e {quoted}; then printf {shlex.quote(marker)}; base64 -w0 {quoted}; fi"
    )
    stdout = getattr(result, "stdout", "") or ""
    if marker not in stdout:
        return None
    payload = stdout.split(marker, 1)[1].strip()
    if not payload:
        return b""
    try:
        return base64.b64decode(payload)
    except Exception:
        return None


async def _repair_reference_paths(
    eval_env,
    value: Any,
    cache_dir: str,
    *,
    reference_keys: set[str],
    repair_home_user: bool,
) -> Any:
    """Materialize SCALE-CUA metric reference paths into local eval cache."""
    if isinstance(value, dict):
        fixed: dict[Any, Any] = {}
        changed = False
        for key, item in value.items():
            if (
                isinstance(key, str)
                and key in reference_keys
                and isinstance(item, str)
                and _is_materializable_reference_path(item, repair_home_user)
            ):
                repaired = await _materialize_reference_path(eval_env, item, cache_dir)
            else:
                repaired = await _repair_reference_paths(
                    eval_env,
                    item,
                    cache_dir,
                    reference_keys=reference_keys,
                    repair_home_user=repair_home_user,
                )
            fixed[key] = repaired
            changed = changed or repaired is not item
        return fixed if changed else value
    if isinstance(value, list):
        fixed_items = [
            await _repair_reference_paths(
                eval_env,
                item,
                cache_dir,
                reference_keys=reference_keys,
                repair_home_user=repair_home_user,
            )
            for item in value
        ]
        return fixed_items if any(a is not b for a, b in zip(fixed_items, value)) else value
    if isinstance(value, tuple):
        fixed_items = tuple(
            [
                await _repair_reference_paths(
                    eval_env,
                    item,
                    cache_dir,
                    reference_keys=reference_keys,
                    repair_home_user=repair_home_user,
                )
                for item in value
            ]
        )
        return fixed_items if any(a is not b for a, b in zip(fixed_items, value)) else value
    return value


def _is_materializable_reference_path(path: str, repair_home_user: bool) -> bool:
    if _is_author_cache_path(path):
        return True
    return repair_home_user and path.startswith("/home/user/")


def _is_author_cache_path(path: str) -> bool:
    return path.startswith(_UPSTREAM_LEGACY_CACHE_PREFIXES)


def _reference_sources_from_task(task: SandboxTaskConfig) -> dict[str, list[str]]:
    config = task.metadata.get("config", []) or []
    sources: dict[str, list[str]] = {}
    for step in config:
        if not isinstance(step, dict) or step.get("type") != "download":
            continue
        for file_info in (step.get("parameters") or {}).get("files", []) or []:
            if not isinstance(file_info, dict):
                continue
            url = file_info.get("url")
            path = file_info.get("path")
            if not isinstance(url, str) or not url or not isinstance(path, str):
                continue
            for key in _reference_source_keys(path, url):
                sources.setdefault(key, [])
                if url not in sources[key]:
                    sources[key].append(url)
    return sources


def _reference_source_keys(path: str, url: str | None = None) -> list[str]:
    keys: list[str] = []
    basename = os.path.basename(path)
    if basename:
        keys.append(basename)
    for value in (path, url or ""):
        parts = [part for part in urlsplit(value).path.split("/") if part]
        if len(parts) >= 2:
            key = "/".join(parts[-2:])
            if key not in keys:
                keys.append(key)
    return keys


async def _materialize_reference_path(eval_env, path: str, cache_dir: str) -> str:
    memo = getattr(eval_env, "_scalecua_reference_path_cache", None)
    if memo is None:
        memo = {}
        setattr(eval_env, "_scalecua_reference_path_cache", memo)
    if path in memo:
        return memo[path]

    if _is_author_cache_path(path):
        source_urls = getattr(eval_env, "_scalecua_reference_source_urls", {}) or {}
        for key in _reference_source_keys(path):
            for url in source_urls.get(key, []):
                # Sync urlopen — offload so a slow URL cannot stall the event loop.
                data = await asyncio.to_thread(_read_reference_source_url, url)
                if not data:
                    continue
                local_path = _reference_cache_path(cache_dir, path)
                os.makedirs(os.path.dirname(local_path), exist_ok=True)
                with open(local_path, "wb") as f:
                    f.write(data)
                memo[path] = local_path
                memo[url] = local_path
                return local_path

    for candidate in _reference_vm_candidates(path):
        data = await _read_container_file_bytes(eval_env.computer, candidate)
        if data is None:
            continue
        local_path = _reference_cache_path(cache_dir, path)
        os.makedirs(os.path.dirname(local_path), exist_ok=True)
        with open(local_path, "wb") as f:
            f.write(data)
        memo[path] = local_path
        memo[candidate] = local_path
        return local_path

    memo[path] = path
    return path


def _read_reference_source_url(url: str) -> bytes | None:
    try:
        with urlopen(url, timeout=30) as response:
            return response.read()
    except Exception:
        return None


def _reference_vm_candidates(path: str) -> list[str]:
    # Only exact VM paths are trusted. A basename search over agent-writable
    # dirs (Desktop/Downloads/...) would let a policy plant a same-name file
    # equal to its own output and trivially pass compare-against-original
    # metrics, so author-cache paths whose pinned URL fetch failed stay
    # unmaterialized (the metric then scores 0 instead of being spoofable).
    if path.startswith("/home/user/"):
        return [path]
    return []


def _reference_cache_path(cache_dir: str, path: str) -> str:
    basename = os.path.basename(path) or "reference"
    stem, ext = os.path.splitext(basename)
    safe_stem = re.sub(r"[^A-Za-z0-9._-]+", "_", stem).strip("._-") or "reference"
    safe_ext = re.sub(r"[^A-Za-z0-9.]+", "", ext)
    digest = hashlib.sha1(path.encode("utf-8", errors="replace")).hexdigest()[:10]
    return os.path.join(cache_dir, "_reference_assets", f"{safe_stem}.{digest}{safe_ext}")


def _resolve_relative_time(config: dict) -> dict:
    return base_runner._resolve_relative_time(config)


async def _get_expected(eval_env, config: dict, cache_dir: str, runtime_split: str):
    expected_type = config.get("type", "")
    if expected_type in SCALECUA_LOCAL_EXPECTED_TYPES:
        expected = await _get_scalecua_local_expected(eval_env, config, cache_dir)
        return await _repair_reference_paths(
            eval_env,
            expected,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    if (
        expected_type == "rule"
        and isinstance(config.get("rules"), dict)
        and isinstance(config["rules"].get("relativeTime"), dict)
    ):
        expected = _resolve_relative_time(config)
        return await _repair_reference_paths(
            eval_env,
            expected,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    if expected_type in BASE_RUNNER_EXPECTED_TYPES:
        expected = await base_runner._get_expected(eval_env.computer, config, cache_dir)
        return await _repair_reference_paths(
            eval_env,
            expected,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    getter = judges.resolve_getter(expected_type, runtime_split)
    if getter is not None:
        expected = await judges.call_overlay_getter(getter, eval_env, config, cache_dir)
        return await _repair_reference_paths(
            eval_env,
            expected,
            cache_dir,
            reference_keys=_REFERENCE_METRIC_KEYS,
            repair_home_user=True,
        )
    expected = await base_runner._get_expected(eval_env.computer, config, cache_dir)
    return await _repair_reference_paths(
        eval_env,
        expected,
        cache_dir,
        reference_keys=_REFERENCE_METRIC_KEYS,
        repair_home_user=True,
    )


async def _get_scalecua_local_expected(eval_env, config: dict, cache_dir: str):
    expected_type = config.get("type", "")
    if expected_type == "vm_file":
        return await _get_vm_file_official(eval_env, config, cache_dir)
    raise KeyError(f"unsupported ScaleCUA local expected type: {expected_type}")


def _as_list(value):
    return value if isinstance(value, list) else [value]


def _coerce_score(raw_score: Any) -> float:
    if isinstance(raw_score, bool):
        return 1.0 if raw_score else 0.0
    if isinstance(raw_score, (int, float)):
        score = float(raw_score)
        if not math.isfinite(score):
            return 0.0
        if score >= 1.0 - 1e-9:
            return 1.0
        if score <= 1e-9:
            return 0.0
        return min(1.0, max(0.0, score))
    # Generated metrics have no return-type guarantee. A truthy non-numeric
    # return (error string, diagnostics dict, the string "0") must not count
    # as success — an RL policy would learn to trigger exactly those paths.
    if raw_score:
        logger.warning(
            "non-numeric metric return %s treated as 0.0: %s",
            type(raw_score).__name__, _debug_preview(raw_score, limit=200),
        )
    return 0.0


def _debug_preview(value: Any, *, limit: int = 1200) -> str | None:
    if value is None:
        return None
    text = value if isinstance(value, str) else json.dumps(value, ensure_ascii=False, default=str)
    text = str(text)
    if len(text) <= limit:
        return text
    return text[:limit] + "...<truncated>"


def _combine_scores(scores: list[float], conj: str, multi_metric: bool) -> float:
    if not scores:
        return 0.0
    if not multi_metric:
        return scores[0]
    if conj == "or":
        return max(scores)
    return sum(scores) / len(scores)


def _prepare_metric_result(metric_fn, result_data: Any, result_config: dict[str, Any]) -> Any:
    if result_config.get("type") != "vm_file":
        return result_data
    if not isinstance(result_data, str) or not os.path.isfile(result_data):
        return result_data
    if _metric_result_wants_text(metric_fn):
        with open(result_data, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    if not _metric_result_wants_bytes(metric_fn):
        return result_data
    with open(result_data, "rb") as f:
        return f.read()


def _prepare_metric_args(
    fn_name: str,
    result_data: Any,
    expected_data: Any,
) -> tuple[Any, Any]:
    # Official ScaleCUA/OSWorld names this metric "clickboard" and defines the
    # signature as (config, terminal_output), unlike the usual (result, expected).
    if fn_name == "is_in_vm_clickboard" and isinstance(expected_data, dict):
        return expected_data, result_data
    return result_data, expected_data


def _metric_result_wants_text(metric_fn) -> bool:
    name = getattr(metric_fn, "__name__", "")
    return any(name.startswith(prefix) for prefix in VM_FILE_TEXT_METRIC_PREFIXES)


def _metric_result_wants_bytes(metric_fn) -> bool:
    try:
        params = list(inspect.signature(metric_fn).parameters.values())
    except (TypeError, ValueError):
        return False
    if not params:
        return False
    ann = params[0].annotation
    return ann is bytes or ann == "bytes" or getattr(ann, "__name__", None) == "bytes"


def _format_result(
    final: float,
    conj: str,
    scores: list[float],
    details: list[dict[str, Any]],
    debug: bool,
    *,
    flush_stats: dict[str, dict[str, Any]] | None = None,
):
    if debug:
        payload = {"conj": conj, "scores": scores, "details": details}
        if flush_stats:
            payload["flush_stats"] = flush_stats
            payload["flush_fired_counts"] = _flush_fired_counts_from_stats(flush_stats)
            payload["flush_counters"] = _flush_counter_aliases(flush_stats)
        return final, payload
    return final
