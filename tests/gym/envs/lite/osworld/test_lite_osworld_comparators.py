"""Unit tests for the lite.osworld eval comparator OVERRIDES in
``lite/gym/envs/lite/osworld/src/eval/metrics.py``.

These guard the hand-written delegate-then-widen comparators that a prior PR
review flagged as untested: each must ACCEPT a provably-equivalent result and
still REJECT a genuinely-wrong one. The data files are byte-locked elsewhere;
this protects the comparator LOGIC against future refactors.

Run:
    uv run pytest tests/gym/envs/lite/osworld/test_lite_osworld_comparators.py -q
"""

from __future__ import annotations

import importlib.util
import zipfile

import pytest
from pptx import Presentation
from pptx.util import Emu

from lite.gym.envs.lite.osworld.src.eval import metrics

# The docx/table/line-spacing comparators delegate to the upstream OSWorld
# eval package (`desktop_env`, an install.sh-only dep that `uv sync` evicts);
# the pure-python tests below run without it.
needs_desktop_env = pytest.mark.skipif(
    importlib.util.find_spec("desktop_env") is None,
    reason="desktop_env not installed (lite.osworld install.sh)",
)


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------


def _xlsx(path, freeze=None):
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    for r in range(1, 6):
        for c in range(1, 6):
            ws.cell(r, c, f"{r}{c}")
    if freeze is not None:
        ws.freeze_panes = freeze  # real LO-equivalent setter
    wb.save(path)
    wb.close()
    return str(path)


def _docx(path, paras, normal_font=None, alignments=None, line_spacings=None):
    from docx import Document

    doc = Document()
    if normal_font is not None:
        doc.styles["Normal"].font.name = normal_font
    for i, text in enumerate(paras):
        p = doc.add_paragraph()
        p.add_run(text)
        # leave run.font.name = None so the Normal-style cascade is exercised
        if alignments is not None and alignments[i] is not None:
            p.alignment = alignments[i]
        if line_spacings is not None and line_spacings[i] is not None:
            p.paragraph_format.line_spacing = line_spacings[i]
    doc.save(path)
    return str(path)


def _epub(path, html_names, chapters):
    """Minimal epub: opf + ncx + chapter html files (with XML encoding decl)."""
    decl = '<?xml version="1.0" encoding="utf-8"?>'
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("content.opf", decl + "<package><metadata/></package>")
        z.writestr("toc.ncx", decl + "<ncx><docTitle><text>B</text></docTitle></ncx>")
        for name, body in zip(html_names, chapters):
            z.writestr(name, f"{decl}<html><body><p>{body}</p></body></html>")
    return str(path)


# ---------------------------------------------------------------------------
# check_xlsx_freeze_pane — row-only relaxation must NOT pass an extra column freeze
# ---------------------------------------------------------------------------


def test_freeze_pane_exact_a2_passes(tmp_path):
    p = _xlsx(tmp_path / "a2.xlsx", freeze="A2")
    assert (
        metrics.check_xlsx_freeze_pane(p, expected_first_row=2, expected_first_col_letter="A")
        == 1.0
    )


def test_freeze_pane_overfreeze_column_fails(tmp_path):
    # freeze_panes="B2" freezes row 1 AND column A (xSplit=1) — a row-only task
    # must NOT accept it even though openpyxl's freeze_panes reads "B2".
    p = _xlsx(tmp_path / "b2.xlsx", freeze="B2")
    assert (
        metrics.check_xlsx_freeze_pane(p, expected_first_row=2, expected_first_col_letter="A")
        == 0.0
    )


def test_freeze_pane_wrong_row_fails(tmp_path):
    p = _xlsx(tmp_path / "a3.xlsx", freeze="A3")
    assert (
        metrics.check_xlsx_freeze_pane(p, expected_first_row=2, expected_first_col_letter="A")
        == 0.0
    )


def test_freeze_pane_no_freeze_fails(tmp_path):
    p = _xlsx(tmp_path / "none.xlsx", freeze=None)
    assert (
        metrics.check_xlsx_freeze_pane(p, expected_first_row=2, expected_first_col_letter="A")
        == 0.0
    )


# ---------------------------------------------------------------------------
# compare_table_numfmt_tolerant — expected=None must not TypeError
# ---------------------------------------------------------------------------


@needs_desktop_env
def test_compare_table_expected_none_no_typeerror(tmp_path):
    # A check_cell-only task has evaluator.expected = null → metric called with
    # no `expected` arg. Must return a float (delegated to upstream), not raise.
    p = _xlsx(tmp_path / "r.xlsx")
    rules = [{"type": "sheet_name", "sheet_idx": 0, "value": "Sheet"}]
    out = metrics.compare_table_numfmt_tolerant(p, rules=rules)
    assert isinstance(out, (int, float))


@needs_desktop_env
def test_compare_table_nf_rule_without_gold_fails(tmp_path):
    # number_format rule needs a gold; expected=None → explicit 0.0, not a crash.
    p = _xlsx(tmp_path / "r2.xlsx")
    rules = [{"type": "style", "props": ["number_format"], "sheet_idx0": 0}]
    assert metrics.compare_table_numfmt_tolerant(p, rules=rules) == 0.0


# ---------------------------------------------------------------------------
# compare_font_names (aliased to loose) — effective-font cascade
# ---------------------------------------------------------------------------


def test_font_names_loose_hoisted_to_normal_passes(tmp_path):
    # LO hoists the font to the Normal style; per-run font.name is None. The
    # effective font is still Times New Roman everywhere.
    p = _docx(tmp_path / "tnr.docx", ["Body one.", "Body two."], normal_font="Times New Roman")
    assert metrics.compare_font_names(p, {"font_name": "Times New Roman"}) == 1.0


def test_font_names_loose_wrong_font_fails(tmp_path):
    p = _docx(tmp_path / "arial.docx", ["Body one."], normal_font="Arial")
    assert metrics.compare_font_names(p, {"font_name": "Times New Roman"}) == 0.0


# ---------------------------------------------------------------------------
# compare_docx_files — trailing-space rstrip fallback
# ---------------------------------------------------------------------------


@needs_desktop_env
def test_docx_files_trailing_space_passes(tmp_path):
    gold = _docx(tmp_path / "gold.docx", ["The lecture. "])  # trailing space
    result = _docx(tmp_path / "res.docx", ["The lecture."])  # stripped
    assert metrics.compare_docx_files(result, gold, ignore_blanks=False) == 1


@needs_desktop_env
def test_docx_files_real_diff_fails(tmp_path):
    gold = _docx(tmp_path / "g2.docx", ["Hello world."])
    result = _docx(tmp_path / "r3.docx", ["Goodbye world."])
    assert metrics.compare_docx_files(result, gold, ignore_blanks=False) == 0


# ---------------------------------------------------------------------------
# compare_epub — layout-independent chapter-text fallback
# ---------------------------------------------------------------------------

_CH1 = "The quick brown fox jumps over the lazy dog repeatedly today."
_CH2 = "Pack my box with five dozen liquor jugs of fine quality wine."


def test_epub_misaligned_layout_same_prose_passes(tmp_path):
    agent = _epub(tmp_path / "agent.epub", ["ch001.xhtml", "ch002.xhtml"], [_CH1, _CH2])
    gold = _epub(tmp_path / "gold.epub", ["file0.html", "file1.html"], [_CH1, _CH2])
    assert metrics.compare_epub(agent, gold) >= 0.5


def test_epub_different_book_fails(tmp_path):
    agent = _epub(tmp_path / "a2.epub", ["ch001.xhtml", "ch002.xhtml"], [_CH1, _CH2])
    gold = _epub(
        tmp_path / "g3.epub",
        ["file0.html", "file1.html"],
        [
            "Completely unrelated content about astronomy.",
            "Nothing here matches the agent book at all.",
        ],
    )
    assert metrics.compare_epub(agent, gold) < 0.5


# ---------------------------------------------------------------------------
# compare_line_spacing — None (default single) normalized to 1.0
# ---------------------------------------------------------------------------


@needs_desktop_env
def test_line_spacing_none_equals_single(tmp_path):
    from docx.enum.text import WD_LINE_SPACING  # noqa: F401 (import sanity)

    a = _docx(tmp_path / "ls_none.docx", ["Same text."], line_spacings=[None])
    b = _docx(tmp_path / "ls_one.docx", ["Same text."], line_spacings=[1.0])
    assert metrics.compare_line_spacing(a, b) == 1


@needs_desktop_env
def test_line_spacing_double_still_required(tmp_path):
    a = _docx(tmp_path / "ls_none2.docx", ["Same text."], line_spacings=[None])
    b = _docx(tmp_path / "ls_two.docx", ["Same text."], line_spacings=[2.0])
    assert metrics.compare_line_spacing(a, b) == 0


# ---------------------------------------------------------------------------
# compare_docx_strict — effective alignment cascade (LEFT == None default)
# ---------------------------------------------------------------------------


def test_docx_strict_alignment_default_equiv(tmp_path):
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    a = _docx(tmp_path / "al_none.docx", ["Para."], alignments=[None])
    b = _docx(tmp_path / "al_left.docx", ["Para."], alignments=[WD_ALIGN_PARAGRAPH.LEFT])
    # explicit LEFT and inherit-default-LEFT are the same visual result
    assert metrics.compare_docx_strict(a, b) == 1


def test_docx_strict_alignment_center_differs(tmp_path):
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    a = _docx(tmp_path / "al_none2.docx", ["Para."], alignments=[None])
    b = _docx(tmp_path / "al_center.docx", ["Para."], alignments=[WD_ALIGN_PARAGRAPH.CENTER])
    assert metrics.compare_docx_strict(a, b) == 0


def _docx_styled(path, *, indent=None, size=None, color=None):
    # A doc whose paragraph/run inherit indent/size/color from a shared style "S".
    # If indent/size/color are given, they're ALSO set EXPLICITLY (the gold shape,
    # redundant with the style); if None, the value is inherited (the LO-saved shape).
    import docx
    from docx.enum.style import WD_STYLE_TYPE
    from docx.shared import Pt, RGBColor

    d = docx.Document()
    st = d.styles.add_style("S", WD_STYLE_TYPE.PARAGRAPH)
    st.paragraph_format.first_line_indent = Pt(36)
    st.font.size = Pt(18)
    st.font.color.rgb = RGBColor.from_string("FF0000")
    p = d.add_paragraph()
    p.style = st
    if indent is not None:
        p.paragraph_format.first_line_indent = Pt(indent)
    r = p.add_run("hello world")
    if size is not None:
        r.font.size = Pt(size)
    if color is not None:
        r.font.color.rgb = RGBColor.from_string(color)
    d.save(path)
    return str(path)


def test_docx_strict_style_inheritance_parity_size_color_indent(tmp_path):
    # CLASS (same root as effective font-name/alignment/line_spacing): LO save drops a
    # run/paragraph's explicit size/color/indent to None when it equals the style value,
    # while the python-docx gold keeps it explicit. compare_docx_strict must resolve the
    # EFFECTIVE value through the style cascade for size/color/indent too — else a correct
    # styled paragraph false-fails. (De-narrowing: one class, not the 3 attrs individually.)
    opts = dict(
        examine_font_name=False, examine_font_size=True, examine_color=True, examine_highlight=False
    )
    gold = _docx_styled(tmp_path / "gold.docx", indent=36, size=18, color="FF0000")
    inherited = _docx_styled(tmp_path / "inh.docx")  # all values inherited (raw None)
    wrong = _docx_styled(tmp_path / "wrong.docx", indent=72, size=24, color="00FF00")
    assert metrics.compare_docx_strict(inherited, gold, **opts) == 1  # 0->1 fixed
    assert metrics.compare_docx_strict(wrong, gold, **opts) == 0  # neg control (no FP)


# ---------------------------------------------------------------------------
# Preventive class tests for landed comparator fixes: assert the general rule
# plus a negative control, not the single instance.
# ---------------------------------------------------------------------------
def test_check_list_msf_sidecar_case_insensitive(tmp_path):
    # Class rule, not the single Drafts.msf instance: ANY `\.msf`
    # sidecar pattern matches case-INSENSITIVELY (Thunderbird auto-generates the
    # sidecar title-case while the eval regex is uppercase).
    listing = tmp_path / "ls.txt"
    listing.write_text("Drafts.msf\nSent.msf\nArchives.msf\n")
    rules = {"expect": [r"\bDRAFTS\.msf\b", r"\bSENT\.msf\b", r"\bARCHIVES\.msf\b"]}
    assert metrics.check_list(str(listing), rules) == 1.0


def test_check_list_mbox_name_stays_case_sensitive(tmp_path):
    # NEG control: a folder/mbox pattern carries `.msf` only inside a negative
    # lookahead `(?!\.msf)`, so it must NOT be relaxed — a lowercase mbox name
    # must still FAIL an uppercase pattern.
    listing = tmp_path / "ls2.txt"
    listing.write_text("drafts\n")
    rules = {"expect": [r"\bDRAFTS/?(?!\.msf)"]}
    assert metrics.check_list(str(listing), rules) == 0.0


def test_check_list_msf_wrong_name_still_fails(tmp_path):
    # NEG control 2: case-insensitivity must not become "matches anything".
    listing = tmp_path / "ls3.txt"
    listing.write_text("Inbox.msf\n")
    rules = {"expect": [r"\bDRAFTS\.msf\b"]}
    assert metrics.check_list(str(listing), rules) == 0.0


def test_check_config_status_memsize_unit_equivalence(tmp_path):
    # GIMP-memsize CLASS: ANY two memsize tokens resolving to the same byte
    # count are equal, regardless of unit spelling ('2G' == raw bytes).
    cfg = tmp_path / "gimprc"
    cfg.write_text("(tile-cache-size 2G)\n")
    rule = {"key": "tile-cache-size", "value": str(2 * 1024**3)}
    assert metrics.check_config_status(str(cfg), rule) == 1.0


def test_check_config_status_memsize_different_value_fails(tmp_path):
    # NEG control: distinct memsizes must still score 0.
    cfg = tmp_path / "gimprc2"
    cfg.write_text("(tile-cache-size 512M)\n")
    assert (
        metrics.check_config_status(str(cfg), {"key": "tile-cache-size", "value": "1024M"}) == 0.0
    )


def test_char_format_trailing_whitespace_neutralized(tmp_path):
    # Class rule: an unstruck TRAILING-whitespace run is format-neutralized,
    # so a fully-struck gold still matches a result whose trailing spaces are
    # unstruck. (Interior-unstruck still mismatches — covered by the FN direction.)
    from docx import Document

    def _make(path, strike_trailing):
        d = Document()
        p = d.add_paragraph()
        if strike_trailing:
            r = p.add_run("hello  ")
            r.font.strike = True
        else:
            w = p.add_run("hello")
            w.font.strike = True
            p.add_run("  ")  # unstruck trailing whitespace
        d.save(path)
        return str(path)

    gold = _make(tmp_path / "g.docx", strike_trailing=True)
    res = _make(tmp_path / "r.docx", strike_trailing=False)
    assert metrics.compare_docx_strict(res, gold) == 1


@needs_desktop_env
def test_chart_range_aware_tolerates_lo_prefix_drop_rejects_wrong_range(tmp_path):
    # require_value_range accepts a sheet-prefix-dropped value ref
    # (LO GUI normalization) but rejects a genuinely-different value column.
    import openpyxl
    from openpyxl.chart import BarChart
    from openpyxl.chart.data_source import AxDataSource, NumDataSource, NumRef
    from openpyxl.chart.series import Series

    def _mk(p, val_ref, cat_ref):
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "S"
        for r in range(1, 6):
            for c in range(1, 4):
                ws.cell(r, c, r * c)
        ch = BarChart()
        ch.title = "T"
        s = Series()
        s.val = NumDataSource(numRef=NumRef(f=val_ref))
        s.cat = AxDataSource(numRef=NumRef(f=cat_ref))
        ch.series.append(s)
        ws.add_chart(ch, "E2")
        wb.save(str(p))
        wb.close()
        return str(p)

    gold = _mk(tmp_path / "g.xlsx", "S!$C$2:$C$5", "S!$A$2:$A$5")
    ok = _mk(tmp_path / "ok.xlsx", "$C$2:$C$5", "$A$2:$A$5")  # LO prefix-dropped, correct
    bad = _mk(tmp_path / "bad.xlsx", "S!$B$2:$B$5", "S!$A$2:$A$5")  # wrong value column
    rule = {
        "type": "chart",
        "sheet_idx0": 0,
        "sheet_idx1": "EI0",
        "chart_props": ["title", "type", "direction"],
        "require_value_range": True,
    }
    assert metrics.compare_calc_chart_type(gold, gold, rules=[rule]) == 1.0
    assert metrics.compare_calc_chart_type(ok, gold, rules=[rule]) == 1.0  # 0->1
    assert metrics.compare_calc_chart_type(bad, gold, rules=[rule]) == 0.0  # neg control


@needs_desktop_env
def test_check_json_settings_default_aware_and_neg_controls(tmp_path):
    # an absent expected key passes iff its shipped VS Code default
    # == expected (strict superset); present keys exact; everything else 0.
    import json

    def _s(obj):
        p = tmp_path / "settings.json"
        p.write_text(json.dumps(obj))
        return str(p)

    def _r(d):
        return {"expected": d}

    # 0->1 recovery: key omitted, expected == shipped default (focus=True, wrap=False)
    assert metrics.check_json_settings(_s({}), _r({"debug.focusEditorOnBreak": True})) == 1.0
    assert metrics.check_json_settings(_s({}), _r({"workbench.editor.wrapTabs": False})) == 1.0
    # already-passing: present & correct stays 1.0
    assert (
        metrics.check_json_settings(
            _s({"debug.focusEditorOnBreak": True}), _r({"debug.focusEditorOnBreak": True})
        )
        == 1.0
    )
    # NEG: present but wrong -> 0
    assert (
        metrics.check_json_settings(
            _s({"debug.focusEditorOnBreak": False}), _r({"debug.focusEditorOnBreak": True})
        )
        == 0.0
    )
    # NEG: absent but expected != default -> 0
    assert metrics.check_json_settings(_s({}), _r({"debug.focusEditorOnBreak": False})) == 0.0
    assert metrics.check_json_settings(_s({}), _r({"workbench.editor.wrapTabs": True})) == 0.0
    # NEG: absent unknown (non-whitelisted) key -> 0
    assert metrics.check_json_settings(_s({}), _r({"editor.fontSize": 14})) == 0.0
    # NEG: missing file -> 0
    assert (
        metrics.check_json_settings(
            str(tmp_path / "nope.json"), _r({"debug.focusEditorOnBreak": True})
        )
        == 0.0
    )


# region options a move_object(table, right) row carries (generator _op_flags).
_REGION_RIGHT = {
    "examine_shape": False,
    "position_mode": "region",
    "position_region": {"slide": 0, "shape_type": 19, "edge": "right"},
    "position_tolerance_emu": 50000,
}
_W, _H, _MARGIN = 9144000, 6858000, 457200


def _pptx_table(path, *, left, width, rows=3, cols=3, top=2000000, height=1500000):
    prs = Presentation()
    prs.slide_width, prs.slide_height = _W, _H
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_table(rows, cols, Emu(left), Emu(top), Emu(width), Emu(height))
    prs.save(str(path))
    return str(path)


def _gold_right(tmp, width):
    return _pptx_table(tmp / "gold.pptx", left=_W - width - _MARGIN, width=width)


# ---- #23: a NEAR-FULL-WIDTH table (center pinned in the middle third) ----------
_WIDE = 7886520  # gold right disp = +171540 EMU (< H/3, center stays middle third)


def test_geom_region_right_exact_move_passes(tmp_path):
    gold = _gold_right(tmp_path, _WIDE)
    res = _pptx_table(tmp_path / "r.pptx", left=_W - _WIDE - _MARGIN, width=_WIDE)
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 1


def test_geom_region_right_realistic_drag_passes(tmp_path):
    # A GUI drag lands 120000 EMU off gold but still biased right past tol
    # (disp 51540 >= 50000). The OLD center-in-third predicate scored this 0
    # (full-width center can't reach the right third) -> this is the #23 FN guard.
    gold = _gold_right(tmp_path, _WIDE)
    res = _pptx_table(tmp_path / "r.pptx", left=_W - _WIDE - _MARGIN - 120000, width=_WIDE)
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 1


def test_geom_region_right_far_move_passes(tmp_path):
    gold = _gold_right(tmp_path, _WIDE)
    res = _pptx_table(tmp_path / "r.pptx", left=_W - _WIDE - _MARGIN + 120000, width=_WIDE)
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 1


def test_geom_region_centered_fails(tmp_path):
    gold = _gold_right(tmp_path, _WIDE)
    res = _pptx_table(tmp_path / "r.pptx", left=(_W - _WIDE) // 2, width=_WIDE)  # disp 0
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 0


def test_geom_region_moved_left_fails(tmp_path):
    gold = _gold_right(tmp_path, _WIDE)
    res = _pptx_table(tmp_path / "r.pptx", left=_MARGIN, width=_WIDE)  # left region, disp -171540
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 0


def test_geom_region_wrong_size_fails(tmp_path):
    # correct RIGHT position but width halved -> upstream examine_shape size gate fires.
    gold = _gold_right(tmp_path, _WIDE)
    narrow = _WIDE // 2
    res = _pptx_table(tmp_path / "r.pptx", left=_W - narrow - _MARGIN, width=narrow)
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 0


# ---- regression guard: a COMPACT table (center CAN leave the third) still works,
#      i.e. the bias predicate is a strict superset of the old third test (#22). ----
_NARROW = 1200000


def test_geom_region_compact_shape_edge_passes(tmp_path):
    gold = _gold_right(tmp_path, _NARROW)  # center deep in right third
    res = _pptx_table(tmp_path / "r.pptx", left=_W - _NARROW - _MARGIN, width=_NARROW)
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 1


def test_geom_region_compact_shape_centered_fails(tmp_path):
    gold = _gold_right(tmp_path, _NARROW)
    res = _pptx_table(tmp_path / "r.pptx", left=(_W - _NARROW) // 2, width=_NARROW)
    assert metrics.compare_pptx_files(res, gold, **dict(_REGION_RIGHT)) == 0


# ---------------------------------------------------------------------------
# compare_pptx_appended_titles  (#11 multi_apps 47f7c0ce Chrome->pptx append)
# ---------------------------------------------------------------------------


def _pptx_deck(path, n_prefix, appended, append_with_placeholders):
    """Build a deck: n_prefix pre-existing slides + one appended slide per title.

    append_with_placeholders=True  -> GOLD shape: python-pptx blank/title layout
      that carries placeholder shapes + a title textbox (several shapes).
    append_with_placeholders=False -> AGENT shape: truly blank layout + 1 textbox
      (one shape) — the LO-Impress-style append whose shape count differs.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # carries 2 placeholders
    blank_layout = prs.slide_layouts[6]  # 0 placeholders
    for i in range(n_prefix):
        s = prs.slides.add_slide(title_layout)
        s.shapes.title.text = f"Existing {i}"
    for t in appended:
        s = prs.slides.add_slide(title_layout if append_with_placeholders else blank_layout)
        tx = s.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
        tx.text_frame.text = t
    prs.save(str(path))
    return str(path)


def test_pptx_appended_titles_shape_count_divergence_passes(tmp_path):
    """Gold's appended slide has several shapes, the agent's has one; the visible
    titles match. The new comparator accepts it (upstream's unconditional shape-count
    gate would reject — that is the oracle-masked FN this fix targets)."""
    titles = [
        "Hands-on: Camera Calibration",
        "Hands-on: Path Planning",
        "Hands-on: Gripper Control",
    ]
    gold = _pptx_deck(tmp_path / "gold.pptx", 3, titles, append_with_placeholders=True)
    agent = _pptx_deck(tmp_path / "agent.pptx", 3, titles, append_with_placeholders=False)
    assert metrics.compare_pptx_appended_titles(agent, gold, num_appended=3) == 1.0


def test_pptx_appended_titles_wrong_title_fails(tmp_path):
    titles = ["Alpha", "Beta", "Gamma"]
    gold = _pptx_deck(tmp_path / "gold.pptx", 2, titles, append_with_placeholders=True)
    bad = _pptx_deck(
        tmp_path / "bad.pptx", 2, ["Alpha", "WRONG", "Gamma"], append_with_placeholders=False
    )
    assert metrics.compare_pptx_appended_titles(bad, gold, num_appended=3) == 0.0


def test_pptx_appended_titles_wrong_order_fails(tmp_path):
    titles = ["Alpha", "Beta", "Gamma"]
    gold = _pptx_deck(tmp_path / "gold.pptx", 2, titles, append_with_placeholders=True)
    swapped = _pptx_deck(
        tmp_path / "swap.pptx", 2, ["Beta", "Alpha", "Gamma"], append_with_placeholders=False
    )
    assert metrics.compare_pptx_appended_titles(swapped, gold, num_appended=3) == 0.0


def test_pptx_appended_titles_wrong_count_fails(tmp_path):
    titles = ["Alpha", "Beta", "Gamma"]
    gold = _pptx_deck(tmp_path / "gold.pptx", 2, titles, append_with_placeholders=True)
    short = _pptx_deck(
        tmp_path / "short.pptx", 2, ["Alpha", "Beta"], append_with_placeholders=False
    )
    assert metrics.compare_pptx_appended_titles(short, gold, num_appended=3) == 0.0


@needs_desktop_env
def test_pptx_appended_titles_guards_upstream_shape_count_fn(tmp_path):
    """Regression witness: upstream compare_pptx_files (even examine_shape=False)
    scores the shape-count-divergent-but-correct append 0 — the FN the new
    comparator flips to 1.0."""
    from desktop_env.evaluators.metrics.slides import compare_pptx_files as upstream

    titles = ["Alpha", "Beta"]
    gold = _pptx_deck(tmp_path / "gold.pptx", 2, titles, append_with_placeholders=True)
    agent = _pptx_deck(tmp_path / "agent.pptx", 2, titles, append_with_placeholders=False)
    assert (
        upstream(
            agent,
            gold,
            examine_shape=False,
            examine_font_name=False,
            examine_background_color=False,
            enable_debug=False,
        )
        == 0
    )
    assert metrics.compare_pptx_appended_titles(agent, gold, num_appended=2) == 1.0
