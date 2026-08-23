"""ScaleCUA generated judge getter and metric repair tests."""

from __future__ import annotations

import json
import sys
from types import SimpleNamespace

import pytest

from lite.gym.envs.lite.scalecua.src.osworld import judges
from lite.gym.envs.lite.scalecua.src.osworld import verify as scalecua_verify
from lite.gym.envs.lite.scalecua.src.utils import dataset


def _cache_ready() -> bool:
    return all(dataset.catalog_path(split).is_file() for split in dataset.RUNTIME_SPLITS)


def _overlays_ready() -> bool:
    """Judge-overlay tests need the imported getters/metrics modules too."""
    return all(
        (root / f"{name}.py").is_file()
        for split in dataset.RUNTIME_SPLITS
        if (root := judges.overlay_dir(split)) is not None
        for name in ("getters", "metrics")
    )


class _FakeInterface:
    def __init__(
        self,
        stdout: str | list[str] = "stdout",
        files: dict[str, bytes] | None = None,
    ):
        self.commands: list[str] = []
        self.command_calls: list[dict[str, object]] = []
        self.hotkeys: list[tuple[str, ...]] = []
        self.typed_text: list[str] = []
        self.stdout = stdout
        self.files = files or {}

    async def read_bytes(self, path: str) -> bytes:
        if path in self.files:
            return self.files[path]
        return f"bytes:{path}".encode()

    async def screenshot(self) -> bytes:
        return b"png"

    async def get_screen_size(self):
        return {"width": 800, "height": 600}

    async def run_command(self, command: str, timeout=None):
        self.commands.append(command)
        self.command_calls.append({"command": command, "timeout": timeout})
        if isinstance(self.stdout, list):
            stdout = self.stdout.pop(0) if self.stdout else ""
        else:
            stdout = self.stdout
        return SimpleNamespace(stdout=stdout, stderr="", returncode=0)

    async def hotkey(self, *keys: str):
        self.hotkeys.append(tuple(keys))

    async def type_text(self, text: str):
        self.typed_text.append(text)


class _FakeComputer:
    def __init__(
        self,
        stdout: str = "stdout",
        files: dict[str, bytes] | None = None,
    ):
        self.interface = _FakeInterface(stdout=stdout, files=files)


def test_scalecua_judge_aliases_windows_vlc_path():
    from lite.gym.envs.lite.scalecua.src.osworld import judges

    assert (
        judges._alias_chrome_profile_path(r"C:\Users\user\AppData\Roaming\vlc\vlcrc")
        == "/home/user/.config/vlc/vlcrc"
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_judge_overlay_loads_relative_imports_without_sys_path_mutation():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    before = list(sys.path)
    mod = judges._load_overlay_module("train", "metrics")
    assert mod is not None
    assert hasattr(mod, "__all__")
    assert sys.path == before


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_getter_helpers_are_injected():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    mod = judges._load_overlay_module("train", "getters")
    getter = getattr(mod, "get_file_exists__3b8e423e430323c0078f4425aded05b9")
    helpers = getter.__globals__

    assert helpers["_get_video_rotation"] is judges._get_video_rotation
    assert judges._get_video_rotation(b"") == (None, False)

    root = helpers["etree"].fromstring(b"<root><child /></root>")
    assert root.xpath("count(//child)") == 1.0


def test_scalecua_video_rotation_helper_parses_ffprobe_shapes():
    assert judges._rotation_from_ffprobe_stream({"tags": {"rotate": "180"}}) == 180
    assert judges._rotation_from_ffprobe_stream({"side_data_list": [{"rotation": -90.0}]}) == -90
    assert judges._rotation_from_ffprobe_stream({"side_data_list": [{}]}) is None


def test_scalecua_bool_existence_guard_preserves_name_returned_dict_getters():
    # The file-existence override flattens a get_file_exists__
    # getter to a bare `test -f` bool ONLY when it is genuinely bool-returning. A
    # rich getter that builds `result = {...}; return result` (return node is an
    # ast.Name, not an ast.Dict literal) must be PRESERVED — flattening it drops
    # the has_content/row_count/... fields the paired metric grades -> reward FN.
    def _bare_bool(env, config):
        return _controller_path_is_file(env, config["path"])  # noqa: F821

    def _dict_literal(env, config):
        return {"exists": True, "has_content": False}

    def _name_returned_dict(env, config):  # the R1 victim shape
        result = {"exists": True, "is_csv": True, "row_count": 0, "has_content": False}
        result["exists"] = _controller_path_is_file  # noqa: F821 (mutate, still dict)
        return result

    assert judges._is_bool_existence_getter(_bare_bool) is True  # flatten (safe)
    assert judges._is_bool_existence_getter(_dict_literal) is False  # keep
    assert judges._is_bool_existence_getter(_name_returned_dict) is False  # keep (R1)


def test_scalecua_generic_rename_status_emits_old_gone_superset(monkeypatch):
    # #154 R3: the file-rename override replaces getters UNCONDITIONALLY, so it
    # must emit a SUPERSET of the keys any paired metric reads — new_exists/
    # old_exists (train metrics) AND old_gone (the RL metric's +0.5 key). Before
    # the fix the flat {new_exists, old_exists} dropped old_gone -> RL correct
    # renames scored 0.5 not 1.0. old_gone = not old_exists cannot fabricate
    # reward (same real path probe, inverted).
    present = {"/new": True, "/old": True}
    monkeypatch.setattr(judges, "_controller_path_is_file", lambda env, p: present.get(p, False))
    cfg = {"new_path": "/new", "old_path": "/old"}
    present["/old"] = False  # correct rename: new exists, old gone
    r = judges._generic_file_rename_status(None, cfg)
    assert r == {"new_exists": True, "old_exists": False, "old_gone": True}
    present["/old"] = True  # copy-not-rename: old still there -> NOT gone (no FP)
    assert judges._generic_file_rename_status(None, cfg)["old_gone"] is False


def test_scalecua_xlsx_effective_fill_uses_highest_precedence_cf_rule():
    # #154 CF-precedence: Excel renders the matching CF rule with the LOWEST
    # priority NUMBER (highest precedence). The old sorted-asc + last-wins loop
    # returned the highest-number = lowest-precedence fill -> inverted color when
    # >=2 overlapping rules match. The min-priority accumulator fixes it.
    import importlib.util
    import os
    import tempfile

    import openpyxl
    from openpyxl.formatting.rule import CellIsRule
    from openpyxl.styles import PatternFill

    spec = importlib.util.spec_from_file_location(
        "_jh", "lite/gym/envs/lite/scalecua/src/gen/eval/judge_helpers.py"
    )
    jh = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(jh)

    def build(rules):  # rules: list of (color, priority)
        wb = openpyxl.Workbook()
        ws = wb.active
        ws["A1"] = 50
        for color, prio in rules:
            r = CellIsRule(
                operator="greaterThan",
                formula=["10"],
                fill=PatternFill(start_color=color, end_color=color, fill_type="solid"),
            )
            r.priority = prio
            ws.conditional_formatting.add("A1", r)
        f = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        wb.save(f.name)
        return f.name

    RED, GREEN, BLUE = "FFFF0000", "FF00FF00", "FF0000FF"
    p = build([(RED, 1), (GREEN, 2)])  # RED is highest precedence (prio 1)
    assert jh.scalecua_xlsx_effective_fill(p, None, "A1", 50) == RED
    os.unlink(p)
    p = build([(GREEN, 2), (RED, 1)])  # order-independent -> still RED
    assert jh.scalecua_xlsx_effective_fill(p, None, "A1", 50) == RED
    os.unlink(p)
    p = build([(BLUE, 1)])  # single rule (the norm) -> unchanged
    assert jh.scalecua_xlsx_effective_fill(p, None, "A1", 50) == BLUE
    os.unlink(p)


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_gimp_metric_helpers_are_injected():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    mod = judges._load_overlay_module("train", "metrics")
    metric = getattr(mod, "check_brightness_increase__677c871105619c52013adb82fa7e5d28")
    helpers = metric.__globals__

    assert callable(helpers.get("calculate_brightness"))
    assert callable(helpers.get("normalize_brightness"))
    assert callable(helpers.get("measure_saturation"))
    assert callable(helpers.get("structure_check_by_mse"))


def test_scalecua_measure_image_saturation_orders_gray_and_red():
    Image = pytest.importorskip("PIL.Image")

    gray = Image.new("RGB", (8, 8), (120, 120, 120))
    red = Image.new("RGB", (8, 8), (220, 20, 20))

    assert judges._measure_image_saturation(gray) == 0.0
    assert judges._measure_image_saturation(red) > 0.5


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_gimp_triangle_helpers_are_injected(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    import numpy as np
    from PIL import Image, ImageDraw, ImageOps

    original = Image.new("RGB", (120, 80), "white")
    draw = ImageDraw.Draw(original)
    draw.polygon([(12, 10), (12, 70), (52, 40)], fill=(255, 255, 0))
    result = ImageOps.mirror(original)
    original_path = tmp_path / "original.png"
    result_path = tmp_path / "result.png"
    original.save(original_path)
    result.save(result_path)

    mod = judges._load_overlay_module("train", "metrics")
    metric = getattr(mod, "check_triangle_flipped__1aaf02638da71a4a84f47e22e2395da3")
    helpers = metric.__globals__

    assert callable(helpers.get("detect_yellow_triangle"))
    assert callable(helpers.get("_detect_yellow_triangle"))
    assert callable(helpers.get("verify_horizontal_flip"))
    assert callable(helpers.get("_find_triangle_color"))
    assert callable(helpers.get("_calculate_centroid"))
    assert callable(helpers.get("is_yellow_color"))
    detection = helpers["_detect_yellow_triangle"](np.array(original))
    mask, angle, contour = detection
    assert np.sum(detection) > 0
    assert mask.any()
    assert angle is not None or contour is not None
    assert helpers["_find_triangle_color"](np.array(original)) is not None
    assert helpers["_calculate_centroid"](mask) is not None
    assert helpers["is_yellow_color"]("FFFFFF00")
    assert (
        metric(
            {"original_path": str(original_path), "result_path": str(result_path)},
            {},
        )
        == 1.0
    )
    top_left = Image.new("RGB", (120, 80), "white")
    top_left_draw = ImageDraw.Draw(top_left)
    top_left_draw.polygon([(2, 2), (2, 22), (22, 2)], fill=(255, 255, 0))
    top_left_path = tmp_path / "top_left.png"
    top_left.save(top_left_path)
    top_left_metric = getattr(mod, "check_triangle_top_left__86e46240")
    assert top_left_metric(str(top_left_path), {}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_excel_color_metric_helpers_are_injected():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    mod = judges._load_overlay_module("train", "metrics")
    metric = getattr(mod, "check_xlsx_conditional_colors__5ae23d3f4f11b8f1328416ec40bcd398")
    helpers = metric.__globals__

    assert callable(helpers.get("_colors_similar"))
    assert helpers["_colors_similar"]("FFFFFF00", "FFFFF10A", 20)
    assert not helpers["_colors_similar"]("FFFF0000", "FF00FF00", 20)
    assert (
        metric(
            {2: "FFFFF10A"},
            {"color_map": {"high": "FFFFFF00"}, "marks_data": {"2": 90}},
            color_tolerance=20,
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_excel_rgb_and_yellow_helpers_are_injected():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill

    mod = judges._load_overlay_module("train", "metrics")
    header_metric = getattr(mod, "check_xlsx_header_format__5c61a707")
    yellow_metric = getattr(mod, "check_xlsx_highlighted_month__f071a3f381bbaeb7d60101b95a35963a")
    top_n_metric = getattr(mod, "check_xlsx_top_n_highlight__5ba4beb6")

    assert callable(header_metric.__globals__.get("_parse_rgb_color"))
    assert callable(header_metric.__globals__.get("_color_distance"))
    assert callable(yellow_metric.__globals__.get("_is_yellow_color"))
    assert callable(top_n_metric.__globals__.get("is_yellow_color"))

    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"].font = Font(bold=True)
    sheet["A1"].fill = PatternFill(fill_type="solid", fgColor="ADD8E6")
    assert (
        header_metric(
            workbook,
            {
                "sheet": 0,
                "header_cells": ["A1"],
                "expected_color": (173, 216, 230),
            },
        )
        == 1.0
    )
    assert (
        yellow_metric(
            {"highlighted_column": "B", "month_name": "March", "color_rgb": "FFFFFF00"},
            {"expected_month": "March", "expected_column": "B"},
        )
        == 1.0
    )
    assert (
        top_n_metric(
            {
                "column_values": {"A2": 10, "A3": 8},
                "highlighted_cells": ["A2"],
                "cell_colors": {"A2": "FFFFFF00"},
            },
            {
                "top_n": 1,
                "values": [10],
                "strict_color_match": True,
                "color_tolerance": 20,
            },
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_docx_italic_color_getter_accepts_rgbcolor(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    from docx import Document
    from docx.shared import RGBColor

    path = tmp_path / "italic-red.docx"
    document = Document()
    run = document.add_paragraph().add_run("Red italic text")
    run.italic = True
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    document.save(path)
    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda vm_path: path.read_bytes()))

    getter = judges.resolve_getter(
        "docx_italic_color_info__96c1a3698fda406b269dcc52d1b2d9c9",
        "train",
    )
    result = getter(env, {"path": "/home/user/Desktop/italic-red.docx"})

    assert getter is judges._get_docx_italic_color_info_96c1
    assert result["italic_runs"][0]["color_rgb"] == "FF0000"
    metric = judges.resolve_metric(
        "check_italic_text_color__96c1a3698fda406b269dcc52d1b2d9c9",
        "train",
    )
    assert metric(result, {"expected_color_rgb": "FF0000"}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_docx_highlight_getter_accepts_ooxml_shading(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    from docx import Document
    from docx.oxml import parse_xml
    from docx.oxml.ns import nsdecls

    path = tmp_path / "highlight-shading.docx"
    target = '"Chocolate is a healthy food." Discuss.'
    document = Document()
    run = document.add_paragraph().add_run(target)
    run._r.get_or_add_rPr().append(
        parse_xml(f'<w:shd {nsdecls("w")} w:val="clear" w:color="auto" w:fill="FFFF00"/>')
    )
    document.save(path)
    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda vm_path: path.read_bytes()))

    getter = judges.resolve_getter(
        "text_highlighting__1f6188764666903805a2bbde08d45ff2",
        "train",
    )
    result = getter(env, {"path": "/home/user/Desktop/highlight-shading.docx"})

    assert getter is judges._get_text_highlighting_1f618
    assert result["found"] is True
    assert result["is_highlighted"] is True
    metric = judges.resolve_metric(
        "check_text_highlighting__1f6188764666903805a2bbde08d45ff2",
        "train",
    )
    assert (
        metric(
            result,
            {"should_be_highlighted": True, "expected_color": "yellow"},
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_csv_merge_getter_tolerates_missing_source_helpers():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    merged = "First Name\nDulce\nMara\nPhilip\n"
    files = {"/home/user/Desktop/merged_data.csv": merged.encode("utf-8")}
    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda vm_path: files.get(vm_path)))

    getter = judges.resolve_getter(
        "csv_merge_data__c66369e707b97de3ccd6da4699663fe6",
        "train",
    )
    result = getter(env, {"path": "/home/user/Desktop/merged_data.csv"})

    assert getter is judges._get_csv_merge_data_c663
    assert result["row_count"] == 3
    assert result["has_single_header"] is True
    assert result["source_file1_unique_values"] == set()
    metric = judges.resolve_metric(
        "check_csv_merge__c66369e707b97de3ccd6da4699663fe6",
        "train",
    )
    assert metric(result, {"row_count": 3, "has_single_header": True}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_calc_initials_metric_accepts_header_row_offset():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric(
        "check_xlsx_column_g__cae1c5edda471ddd033496e2d51a93a2",
        "train",
    )

    assert (
        metric(
            {2: "Initials", 3: "QL", 4: "TY", 5: "LK"},
            {"rows": {"2": "QL", "3": "TY", "4": "LK"}},
        )
        == 1.0
    )
    assert (
        metric(
            {2: "Wrong", 3: "QL", 4: "TY", 5: "LK"},
            {"rows": {"2": "QL", "3": "TY", "4": "LK"}},
        )
        == 0.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_calc_sorted_column_metric_accepts_top_value_prefix():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric("check_xlsx_column_sorted__b22cf776", "train")

    assert (
        metric(
            {"is_sorted": True, "values": [97, 95, 93, 93, 91, 83, 81]},
            {"is_sorted": True, "expected_values": [97.0, 95.0, 93.0, 93.0, 91.0]},
        )
        == 1.0
    )
    assert (
        metric(
            {"is_sorted": True, "values": [97, 95, 92, 93, 91, 83, 81]},
            {"is_sorted": True, "expected_values": [97.0, 95.0, 93.0, 93.0, 91.0]},
        )
        == 0.5
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_gross_profit_getter_uses_first_data_row(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")
    judges._load_overlay_module.cache_clear()

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    net_sales = [70000, 60000, 65000, 74000, 76000, 71000, 72000, 63000, 57000]
    cogs = [15000, 12338, 11549, 13181, 13136, 13076, 12366, 12148, 13351]
    gross_profit = [ns - cg for ns, cg in zip(net_sales, cogs)]
    for offset, (ns, cg, gp) in enumerate(zip(net_sales, cogs, gross_profit), start=2):
        sheet[f"E{offset}"] = ns
        sheet[f"I{offset}"] = cg
        sheet[f"J{offset}"] = gp
    sheet["J12"] = sum(gross_profit) / len(gross_profit)
    path = tmp_path / "income.xlsx"
    workbook.save(path)
    workbook.close()

    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda _: path.read_bytes()))
    getter = judges.resolve_getter(
        "xlsx_gross_profit_calculations__87606b1393cde1435f28787ed7fdd477",
        "train",
    )
    result = getter(env, {"path": "/home/user/IncomeStatement2.xlsx", "sheet": 0})

    assert getter is judges._get_xlsx_gross_profit_calculations_87606
    assert result["gross_profit_values"] == [float(value) for value in gross_profit]
    assert result["net_sales_values"] == [float(value) for value in net_sales]
    assert result["cogs_values"] == [float(value) for value in cogs]
    metric = judges.resolve_metric(
        "check_xlsx_gross_profit_calculations__87606b1393cde1435f28787ed7fdd477",
        "train",
    )
    assert (
        metric(
            result,
            {
                "expected_gross_profit_values": [float(value) for value in gross_profit],
                "expected_average": round(sum(gross_profit) / len(gross_profit), 2),
            },
            tolerance=1.0,
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_sales_pie_getter_derives_title_and_data_points(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")
    judges._load_overlay_module.cache_clear()

    from openpyxl.chart import PieChart, Reference
    from openpyxl.styles import PatternFill

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Rep", "Jan", "Feb", "Mar", "Apr", "May", "Jun"])
    sheet.append(["Kandice Hussey", 500, 500, 500, 500, 500, 500])
    sheet["A12"] = "Top Performer"
    sheet["B12"] = "Kandice Hussey"
    sheet["C12"] = 3000
    yellow = PatternFill(fill_type="solid", fgColor="FFFF00")
    for col in range(1, 8):
        sheet.cell(row=12, column=col).fill = yellow

    chart = PieChart()
    chart.title = "Top Performer Monthly Breakdown"
    chart.add_data(Reference(sheet, min_col=2, max_col=7, min_row=2, max_row=2))
    sheet.add_chart(chart, "E12")

    path = tmp_path / "SalesRep.xlsx"
    workbook.save(path)
    workbook.close()

    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda _: path.read_bytes()))
    getter = judges.resolve_getter(
        "xlsx_sales_data__68a2bc639d62cfd83c2dd75c92c0cd2c",
        "train",
    )
    result = getter(env, {"path": "/home/user/SalesRep.xlsx"})

    assert getter is judges._get_xlsx_sales_data_68a2
    assert result["A12"] == "Top Performer"
    assert result["B12"] == "Kandice Hussey"
    assert result["C12"] == 3000
    assert result["_charts"][0]["type"] == "PieChart"
    assert result["_charts"][0]["title"] == "Top Performer Monthly Breakdown"
    assert result["_charts"][0]["data_points"] == 6
    metric = judges.resolve_metric(
        "check_xlsx_sales_task__68a2bc639d62cfd83c2dd75c92c0cd2c",
        "train",
    )
    assert (
        metric(
            result,
            {
                "A12": "Top Performer",
                "B12": "Kandice Hussey",
                "C12": 3000,
            },
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_website_qty_getter_accepts_descriptive_summary_sheet(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")
    judges._load_overlay_module.cache_clear()

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet.append(
        ["Date Time", "Web Site", "Product", "Type", "Quantity", "Discount", "Total Quantity"]
    )
    rows = [
        ("amazon.com", 10),
        ("amazon.com", 5),
        ("ebay.com", 7),
        ("gel-boomerang.com", 3),
        ("coloradoboomerangs.com", 2),
    ]
    expected = {}
    for website, quantity in rows:
        expected[website] = expected.get(website, 0) + quantity
        sheet.append(["", website, "Product", "Retail", quantity, 0, expected[website]])
    summary = workbook.create_sheet("Website Quantity Summary")
    summary.append(["Web Site", "Total Quantity"])
    for website, quantity in expected.items():
        summary.append([website, quantity])

    path = tmp_path / "BoomerangSales.xlsx"
    workbook.save(path)
    workbook.close()

    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda _: path.read_bytes()))
    getter = judges.resolve_getter(
        "xlsx_website_qty__57e581555e19ceb3db2669015c9e00b9",
        "train",
    )
    result = getter(env, {"path": "/home/user/BoomerangSales.xlsx"})

    assert getter is judges._get_xlsx_website_qty_57e
    assert result["has_summary_sheet"] is True
    assert result["summary_data"] == expected
    assert result["has_total_column"] is True
    metric = judges.resolve_metric(
        "check_xlsx_website_qty__57e581555e19ceb3db2669015c9e00b9",
        "train",
    )
    assert metric(result, {"website_quantities": expected}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_sheet_cells_metric_treats_missing_blank_boundary_as_unmerged():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric(
        "check_xlsx_sheet_cells__7001880f97ce5d91e9e11e6dd84fdc47",
        "train",
    )
    result = {
        "data": {
            "Sheet3": {
                "cells": {
                    "A1": {"value": "Financial Data", "merged": True},
                    "B1": {"value": None, "merged": True},
                    "C1": {"value": None, "merged": True},
                    "D1": {"value": None, "merged": True},
                    "A2": {"value": "Principal Amount", "merged": True},
                    "B2": {"value": None, "merged": True},
                },
                "merged_ranges": ["A1:D1", "A2:B2"],
            }
        }
    }
    expected = {
        "sheet_name": "Sheet3",
        "checks": [
            {"coord": "A1", "value": "Financial Data", "merged": True},
            {"coord": "B1", "merged": True},
            {"coord": "C1", "merged": True},
            {"coord": "D1", "merged": True},
            {"coord": "E1", "merged": False},
            {"coord": "A2", "value": "Principal Amount", "merged": True},
            {"coord": "B2", "merged": True},
            {"coord": "C2", "merged": False},
        ],
    }

    assert metric(result, expected) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_percentage_metric_accepts_percent_scale():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric("check_percentage__3c3385e3", "train")

    assert metric(0.6939, {"value": 69.39}, tolerance=0.5, is_decimal=True) == 1.0
    assert metric(0.51, {"value": 69.39}, tolerance=0.5, is_decimal=True) == 0.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_numeric_metric_accepts_expected_value_dict():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric("check_numeric_value__28be6047", "train")

    assert metric(649.16, {"value": 649.16}, tolerance=0.01) == 1.0
    assert metric("649.16", {"value": 649.16}, tolerance=0.01) == 1.0
    assert metric(650.0, {"value": 649.16}, tolerance=0.01) == 0.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_decimal_text_metric_accepts_text_formula_fallback():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric(
        "check_calc_cell_text__a8896662f47f9f0431d9d6ec0b20825a",
        "train",
    )

    assert (
        metric(
            {
                "value": None,
                "formula": '=TEXT(C1/10;"0.0000")',
                "source_cell": "C1",
                "source_value": 19.5,
                "numeric_value": None,
            },
            {"text": "1.9500"},
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_decimal_text_metric_accepts_numeric_formula_fallback():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric(
        "check_calc_cell_text__154194afadf90d869e3ab3e594cc44ae",
        "train",
    )

    assert (
        metric(
            {
                "value": None,
                "number_format": "0.000",
                "formula": "=C1*2",
                "source_value": 19.5,
            },
            {"text": "39.000"},
        )
        == 1.0
    )


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_total_label_getter_uses_product_header(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")
    judges._load_overlay_module.cache_clear()

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Date", "ID", "Product", "Amount"])
    sheet["C20"] = "Total"
    sheet["D20"] = 76079
    path = tmp_path / "Arrang_Value_min_to_max.xlsx"
    workbook.save(path)
    workbook.close()

    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda _: path.read_bytes()))
    getter = judges.resolve_getter(
        "xlsx_cell_value__f211fc504e7a50f585ecb0b30c1d4299",
        "train",
    )
    result = getter(
        env,
        {
            "path": "/home/user/Arrang_Value_min_to_max.xlsx",
            "cell": "D20",
            "label_cell": "B20",
        },
    )

    assert getter is judges._get_xlsx_cell_value_f211
    assert result == {"product_label": "Total", "amount_value": 76079.0}
    metric = judges.resolve_metric(
        "check_xlsx_cell_value__f211fc504e7a50f585ecb0b30c1d4299",
        "train",
    )
    assert metric(result, {"product_label": "Total", "amount_value": 76079}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_2020_cost_getter_finds_lower_table(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")
    judges._load_overlay_module.cache_clear()

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet["A18"] = "Personal Costs - 2020"
    sheet["A19"] = "Month"
    sheet["B19"] = "Total Cost"
    values = [900, 800, 700, 650]
    for row, (month, value) in enumerate(zip(["Jan", "Feb", "Mar", "Apr"], values), start=20):
        sheet[f"A{row}"] = month
        sheet[f"B{row}"] = value
    sheet["A23"].font = openpyxl.styles.Font(bold=True)
    path = tmp_path / "Create_column_charts_using_statistics.xlsx"
    workbook.save(path)
    workbook.close()

    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda _: path.read_bytes()))
    getter = judges.resolve_getter(
        "calc_cell_format__a3c8817989e00a27454e516ffd39d47e",
        "train",
    )
    result = getter(env, {"path": "/home/user/Create_column_charts_using_statistics.xlsx"})

    assert getter is judges._get_calc_cell_format_a3c
    assert result["cost_2020_column"][-1] == (23, "B23", 650.0)
    metric = judges.resolve_metric(
        "check_calc_cell_format__a3c8817989e00a27454e516ffd39d47e",
        "train",
    )
    assert metric(result, {"cell": "A23", "bold": True}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_row_count_getter_detects_type_header(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")
    judges._load_overlay_module.cache_clear()

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.append(["Date", "Website", "Product", "Type", "Amount"])
    for row_idx in range(2, 15):
        sheet.append(["", "", "", "Wholesale", 1])
    for row_idx in range(15, 37):
        sheet.append(["", "", "", "Retail", 1])
        sheet.row_dimensions[row_idx].hidden = True
    sheet.auto_filter.ref = "A1:E36"
    sheet.auto_filter.add_filter_column(3, ["Wholesale"])
    path = tmp_path / "BoomerangSales.xlsx"
    workbook.save(path)
    workbook.close()

    env = SimpleNamespace(controller=SimpleNamespace(get_file=lambda _: path.read_bytes()))
    getter = judges.resolve_getter("xlsx_row_count__6c4a1081", "train")
    result = getter(
        env,
        {
            "path": "/home/user/BoomerangSales.xlsx",
            "dest": "BoomerangSales.xlsx",
            "sheet": 0,
        },
    )

    assert getter is judges._get_xlsx_row_count_6c4a
    assert result["type_column_filtered"] is True
    assert result["visible_row_types"] == ["Wholesale"] * 13
    assert result["hidden_row_types"] == ["Retail"] * 22
    metric = judges.resolve_metric("check_row_count__6c4a1081", "train")
    assert metric(result, {"count": 13, "tolerance": 1}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_generated_period_rate_metric_allows_adjacent_total_label():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    metric = judges.resolve_metric(
        "check_xlsx_period_rate_sum__da315cf2edc339756d9b2f4defe564ad",
        "train",
    )
    result = {
        "header": "Period Rate (%)",
        "total_data_rows": 24,
        "period_rate_values": [
            {
                "annual_rate": 12.0,
                "period_per_year": 12.0,
                "period_rate": 1.0,
                "is_numeric": True,
            }
            for _ in range(24)
        ],
        "sum_value": 24.0,
        "sum_row_label": "",
    }

    assert metric(result, {"expected_sum": 24.0}) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_calc_chart_helpers_are_injected_and_titles_normalized():
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    judges._load_overlay_module.cache_clear()

    mod = judges._load_overlay_module("train", "getters")
    getter = getattr(mod, "get_calc_chart_info__6d53c98d6c4b8d6ad256907b7e9092b3")
    helper = getter.__globals__.get("_parse_cell_range")

    assert callable(helper)
    workbook = SimpleNamespace(sheetnames=["WeeklySales"])
    assert helper("'WeeklySales'!$C$2:$C$11", workbook) == {
        "sheet_name": "WeeklySales",
        "columns": ["C"],
        "row_count": 10,
    }

    title = SimpleNamespace(
        tx=SimpleNamespace(
            rich=SimpleNamespace(
                p=[
                    SimpleNamespace(
                        r=[SimpleNamespace(t="COGS "), SimpleNamespace(t="Distribution")]
                    )
                ]
            )
        )
    )
    metric = judges.resolve_metric(
        "check_calc_chart__6d53c98d6c4b8d6ad256907b7e9092b3",
        "train",
    )
    result = {
        "charts": [
            {
                "sheet_name": "Sheet2",
                "chart_type": "PieChart",
                "chart_title": title,
                "data_column_names": ["COGS"],
                "data_row_count": 10,
            }
        ]
    }
    expected = {
        "sheet_name": "Sheet2",
        "chart_type": "PieChart",
        "chart_title": "COGS Distribution",
        "data_column_name": "COGS",
        "data_row_count": 10,
    }

    assert metric(result, expected) == 1.0

    area_metric = judges.resolve_metric(
        "check_calc_chart__d89ae1ac878087c6053fc4182b0f11d8",
        "train",
    )
    area_result = {
        "charts": [
            {
                "sheet_name": "Sheet2",
                "chart_type": "AreaChart",
                "chart_title": SimpleNamespace(text="COGS Over Time"),
                "data_column_names": ["COGS"],
                "data_row_count": 10,
            }
        ]
    }
    area_expected = {
        "sheet_name": "Sheet2",
        "chart_type": "AreaChart",
        "chart_title": "COGS Over Time",
        "data_column": "COGS",
        "min_data_rows": 10,
    }

    assert area_metric(area_result, area_expected) == 1.0


@pytest.mark.skipif(not _cache_ready(), reason="lite.scalecua cache not imported")
def test_scalecua_xlsx_formula_metric_accepts_correct_sum_formula(tmp_path):
    pytest.importorskip("desktop_env", reason="desktop_env not installed")
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    value = 1
    expected_total = 0
    for row in range(2, 10):
        for col in range(3, 16):
            sheet.cell(row=row, column=col).value = value
            expected_total += value
            value += 1
    sheet["C10"] = "=SUM(C2:O9)"
    path = tmp_path / "formula.xlsx"
    workbook.save(path)
    workbook.close()

    metric = judges.resolve_metric("check_xlsx_cell_value__38c661cc", "train")

    assert (
        metric(
            str(path),
            {
                "sheet": 0,
                "cell": "C10",
                "expected_value": expected_total,
                "tolerance": 0.01,
            },
        )
        == 1.0
    )
    assert (
        metric(
            str(path),
            {
                "sheet": 0,
                "cell": "C10",
                "expected_value": expected_total + 10,
                "tolerance": 0.01,
            },
        )
        == 0.0
    )


def test_scalecua_getter_falls_back_to_upstream(monkeypatch):
    upstream_getters = pytest.importorskip("desktop_env.evaluators.getters")
    monkeypatch.setattr(judges, "_load_overlay_module", lambda *_: None)

    getter = judges.resolve_getter("enabled_experiments", "train")

    assert getter is upstream_getters.get_enabled_experiments


@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
def test_scalecua_notes_panel_metric_accepts_selected_notes_tab():
    metric = judges.resolve_metric("check_notes_panel__be7f4bd7", "train")
    accessibility_tree = """
    <desktop xmlns:st="uri:deskat:state.at-spi.gnome.org">
      <application name="LibreOffice Impress">
        <object role="page tab list" name="View tabs">
          <object role="page tab" name="Notes">
            <st:state name="showing" />
            <st:state name="selected" />
          </object>
        </object>
        <object role="paragraph" name="Click to add Notes">
          <st:state name="showing" />
        </object>
      </application>
    </desktop>
    """
    hidden_tree = """
    <desktop xmlns:st="uri:deskat:state.at-spi.gnome.org">
      <application name="LibreOffice Impress">
        <object role="page tab" name="Notes">
          <st:state name="showing" />
        </object>
      </application>
    </desktop>
    """

    assert metric(accessibility_tree, {}) == 1.0
    assert metric(hidden_tree, {}) == 0.0


@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
def test_scalecua_speedtest_report_metric_accepts_idle_latency_label():
    metric = judges.resolve_metric("check_speedtest_report__26660ad1", "train")
    report = """
    Speedtest.net Network Speed Test Report
    Test provider: Ookla Speedtest
    Result URL: https://www.speedtest.net/result/19432281822
    Server information:
      Server: Tekify Fiber & Wireless
      Location: Fremont, CA
    Measured results:
      Download: 275.12 Mbps
      Upload: 86.95 Mbps
      Ping / idle latency: 4 ms
      Download latency: 119 ms
      Upload latency: 297 ms
    """

    assert metric(report, {}) == 1.0
    assert metric(report.replace("Upload: 86.95 Mbps", ""), {}) == 0.0


def test_scalecua_vscode_keybinding_metric_accepts_negative_unbinding_without_when(tmp_path):
    metric = judges.resolve_metric("check_json_keybindings", "train")
    keybindings = tmp_path / "keybindings.json"
    expected = {
        "expected": {
            "key": "ctrl+h",
            "command": "-editor.action.startFindReplaceAction",
            "when": "editorTextFocus && !editorReadonly",
        }
    }

    keybindings.write_text(
        json.dumps(
            [
                {
                    "key": "ctrl+h",
                    "command": "-editor.action.startFindReplaceAction",
                }
            ]
        ),
        encoding="utf-8",
    )
    assert metric(str(keybindings), expected) == 1.0

    keybindings.write_text(
        json.dumps(
            [
                {
                    "key": "ctrl+h",
                    "command": "-editor.action.startFindReplaceAction",
                    "when": "terminalFocus",
                }
            ]
        ),
        encoding="utf-8",
    )
    assert metric(str(keybindings), expected) == 0.0

    keybindings.write_text(
        json.dumps(
            [
                {
                    "key": "ctrl+h",
                    "command": "editor.action.startFindReplaceAction",
                }
            ]
        ),
        encoding="utf-8",
    )
    assert metric(str(keybindings), expected) == 0.0


@pytest.mark.asyncio
@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
async def test_scalecua_generated_xlsx_cell_value_recomputes_uncached_sumif_formula(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    rows = [
        ("Channel", "Revenue"),
        ("Web Site Sales", 6939),
        ("Retail", 3061),
    ]
    for row_index, row in enumerate(rows, start=1):
        sheet.cell(row=row_index, column=5).value = row[0]
        sheet.cell(row=row_index, column=7).value = row[1]
    sheet["H2"] = '=SUMIF(E:E,"Web Site Sales",G:G)/SUM(G:G)'
    path = tmp_path / "summer_sales.xlsx"
    workbook.save(path)
    workbook.close()

    remote_path = "/home/user/SummerSales.xlsx"
    env = judges.make_eval_env(_FakeComputer(files={remote_path: path.read_bytes()}), str(tmp_path))

    result = await scalecua_verify._get_result(
        env,
        {
            "type": "xlsx_cell_value__3c3385e3",
            "path": remote_path,
            "sheet": 0,
            "cell": "H2",
        },
        str(tmp_path),
        "train",
    )

    assert result == pytest.approx(0.6939)


@pytest.mark.asyncio
async def test_scalecua_repairs_generated_professor_notes_getter_name_column(tmp_path):
    openpyxl = pytest.importorskip("openpyxl")

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet["A2"] = "No."
    sheet["B2"] = "Professor"
    sheet["G2"] = "Notes"
    sheet["A3"] = 1
    sheet["B3"] = "Qi Liu"
    sheet["G3"] = "PhD advisor"
    sheet["A4"] = 2
    sheet["B4"] = "Tao Yu"
    path = tmp_path / "Professor_Contact.xlsx"
    workbook.save(path)
    workbook.close()

    remote_path = "/home/user/Desktop/Professor_Contact.xlsx"
    env = judges.make_eval_env(
        _FakeComputer(files={remote_path: path.read_bytes()}),
        str(tmp_path),
    )

    out = await scalecua_verify._repair_generated_xlsx_professor_notes_result(
        env,
        {"type": "xlsx_cells_dict__a9a82c07", "path": remote_path, "sheet": 0},
        {
            "header": "Notes",
            "professors": [{"row": 3, "name": "1", "notes": "PhD advisor"}],
        },
    )

    assert out == {
        "header": "Notes",
        "professors": [
            {"row": 3, "name": "Qi Liu", "notes": "PhD advisor"},
            {"row": 4, "name": "Tao Yu", "notes": ""},
        ],
    }


@pytest.mark.asyncio
async def test_scalecua_repairs_generated_terminal_gsettings_getters(tmp_path):
    computer = _FakeComputer(
        stdout=[
            "'ibeam'\n",
            "10000\n",
            "profile-uuid\n'MyWork'\n",
            (
                "use-theme-colors:\n"
                "false\n"
                "background-color:\n"
                "'#002b36'\n"
                "foreground-color:\n"
                "'#839496'\n"
            ),
        ]
    )
    env = judges.make_eval_env(computer, str(tmp_path))

    assert (
        await scalecua_verify._repair_generated_terminal_gsettings_result(
            env,
            {"type": "terminal_cursor_shape__abfc248d5ef921e3773f3d8ab3492012"},
            "old-terminal-output",
        )
        == "'ibeam'"
    )
    assert (
        await scalecua_verify._repair_generated_terminal_gsettings_result(
            env,
            {"type": "terminal_scrollback__8e9fd232f4307bd96cd64b8e7a1f9389"},
            "old-terminal-output",
        )
        == "10000"
    )
    assert (
        await scalecua_verify._repair_generated_terminal_gsettings_result(
            env,
            {"type": "terminal_profile_name__dfa021d620bcc8024aed513c1adfaad3"},
            "old-terminal-output",
        )
        == "profile-uuid\n'MyWork'"
    )
    assert await scalecua_verify._repair_generated_terminal_gsettings_result(
        env,
        {"type": "terminal_color_scheme__ad0133b72746f22f133946cc65aa7f55"},
        "old-terminal-output",
    ) == ("use-theme-colors:\nfalse\nbackground-color:\n'#002b36'\nforeground-color:\n'#839496'")

    assert all(command.startswith("bash -lc ") for command in computer.interface.commands)
    assert "cursor-shape" in computer.interface.commands[0]
    assert "scrollback-lines" in computer.interface.commands[1]
    assert "visible-name" in computer.interface.commands[2]
    assert "background-color" in computer.interface.commands[3]


@pytest.mark.asyncio
async def test_scalecua_generated_terminal_gsettings_repair_keeps_original_on_empty_stdout(
    tmp_path,
):
    env = judges.make_eval_env(_FakeComputer(stdout=""), str(tmp_path))

    assert (
        await scalecua_verify._repair_generated_terminal_gsettings_result(
            env,
            {"type": "terminal_cursor_shape__abfc248d5ef921e3773f3d8ab3492012"},
            "old-terminal-output",
        )
        == "old-terminal-output"
    )


@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
def test_scalecua_pptx_specific_text_metric_accepts_target_substring():
    metric = judges.resolve_metric(
        "check_pptx_specific_text__1a27b3b659a4498b5b5ec44abc0ec111",
        "train",
    )

    assert (
        metric(
            {"text": "The winning Team is The Eagles"},
            {"target_text": "The Eagles"},
        )
        == 1.0
    )
    assert metric({"text": "The winning Team is The Bears"}, {"target_text": "The Eagles"}) == 0.0


@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
def test_scalecua_pptx_image_size_metric_ignores_unstable_embedded_name(tmp_path):
    pytest.importorskip("pptx")
    pytest.importorskip("PIL")
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Cm

    image_path = tmp_path / "none.png"
    Image.new("RGB", (32, 32), "gray").save(image_path)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_picture(
        str(image_path),
        Cm(1),
        Cm(1),
        width=Cm(5),
        height=Cm(4),
    )
    shape.name = "Picture 1"
    path = tmp_path / "image-size.pptx"
    presentation.save(path)

    metric = judges.resolve_metric("check_pptx_image_size__ca0cb440", "train")
    assert (
        metric(
            str(path),
            {
                "slide_index": 0,
                "expected_width_cm": 5.0,
                "expected_height_cm": 4.0,
                "expected_image_name": "none.png",
                "tolerance_cm": 0.05,
            },
        )
        == 1.0
    )


@pytest.mark.asyncio
@pytest.mark.skipif(not _overlays_ready(), reason="lite.scalecua judge overlays not imported")
async def test_scalecua_generated_pptx_textbox_fonts_read_paragraph_default_size(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation
    from pptx.util import Pt

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox1 = slide.shapes.add_textbox(0, 0, 1000, 1000)
    p1 = textbox1.text_frame.paragraphs[0]
    p1.font.size = Pt(110)
    p1.text = "Market Research"
    textbox2 = slide.shapes.add_textbox(0, 1200, 1000, 1000)
    p2 = textbox2.text_frame.paragraphs[0]
    p2.font.size = Pt(70)
    p2.text = "43%"
    path = tmp_path / "fonts.pptx"
    presentation.save(path)

    remote_path = "/home/user/Desktop/45_1.pptx"
    env = judges.make_eval_env(_FakeComputer(files={remote_path: path.read_bytes()}), str(tmp_path))

    result = await scalecua_verify._get_result(
        env,
        {
            "type": "pptx_textbox_fonts__e888e523",
            "ppt_file_path": remote_path,
            "slide_index": 0,
            "textbox1_shape_idx": 0,
            "textbox2_shape_idx": 1,
        },
        str(tmp_path),
        "train",
    )

    assert result["textbox1_font_size"] == pytest.approx(110)
    assert result["textbox2_font_size"] == pytest.approx(70)


@pytest.mark.asyncio
async def test_scalecua_generated_pptx_slide_subtitle_accepts_textbox_fallback(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(0, 0, 4000000, 800000)
    textbox.text = "A comprehensive guide"
    path = tmp_path / "subtitle.pptx"
    presentation.save(path)

    remote_path = "/home/user/Desktop/189_4.pptx"
    env = judges.make_eval_env(
        _FakeComputer(files={remote_path: path.read_bytes()}),
        str(tmp_path),
    )

    result = await scalecua_verify._repair_generated_pptx_result(
        env,
        {
            "type": "pptx_slide_subtitle__c9f628a8",
            "path": remote_path,
            "slide_index": 1,
        },
        "",
    )

    assert result == "A comprehensive guide"
